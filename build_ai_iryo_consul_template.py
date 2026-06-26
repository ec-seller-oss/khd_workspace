"""
KHD AI医療コンサル 全体像テンプレ提案資料（院横断・展開用）。
LINEは数あるソリューションの1つとして配置。クリーム白×レンガ赤。出力: ai_iryo_consul_template.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.6),sz=22,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.5),Inches(12.1),Inches(0.35),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"KHD ｜ AI医療コンサル",Inches(0.5),H-Inches(0.42),Inches(8),Inches(0.32),sz=9,col=GRY)
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.06); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=(TEALBG if is_hi else (CARD if ri%2==1 else BG))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or (ci==0) or is_hi
                    r.font.color.rgb=(WHT if ri==0 else (TEALD if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A))))
    return tb

# ════ 1 表紙 ════
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"AI医療コンサル（KHD）",Inches(0.9),Inches(1.1),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"クリニックの「人手不足 × IT苦手 × バラバラなツール」を、\nAIで小さく解く。",Inches(0.88),Inches(1.7),Inches(11.7),Inches(1.6),sz=30,bold=True,col=INK,line_sp=1.12)
t(s,"院ごとに課題を診断し、刺さる1つから。0円ではじめて、効いた分だけ広げます。",Inches(0.92),Inches(3.7),Inches(11.4),Inches(0.5),sz=14,col=GRY)
bx(s,Inches(0.9),Inches(4.6),Inches(11.5),Inches(1.2),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(4.6),Inches(0.1),Inches(1.2),RED)
t(s,"提案だけで終わらせません。AIで「実物」を0円で作って、動かしてお見せします。",Inches(1.2),Inches(4.6),Inches(11),Inches(1.2),sz=15,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)
t(s,"ご提案先：＿＿＿＿＿＿＿＿ クリニック 御中　／　KHD ｜ AI医療コンサル",Inches(0.9),H-Inches(0.7),Inches(11),Inches(0.3),sz=11,col=GRY)

# ════ 2 共通の痛み ════
s=sl(); ft(s)
hdr(s,"いま起きていること","どのクリニックも、だいたい同じ所でつまずく","フェーズは違っても、負担の「型」は共通しています")
pains=[("① 受付・電話の波","問い合わせ・予約・順番の電話でスタッフの手が止まる"),
       ("② 入力・転記の二度手間","紙→電子カルテへ書き写し。同じことを二度"),
       ("③ 集患・認知が弱い","新患が増えない。HP・SNSに手が回らない"),
       ("④ 立地・開業の判断","開業・分院、「ここで食えるか」が読めない"),
       ("⑤ 承継・補助金の情報不足","承継やIT補助金、何が使えるか分からない")]
cw,ch,gx,gy=Inches(3.96),Inches(1.55),Inches(0.2),Inches(0.3); x0,y0=Inches(0.55),Inches(2.0)
for i,(ti,ds) in enumerate(pains):
    cx=x0+(cw+gx)*(i%3); cy=y0+(ch+gy)*(i//3)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.25),cy+Inches(0.18),cw-Inches(0.4),Inches(0.5),sz=14,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.27),cy+Inches(0.7),cw-Inches(0.45),Inches(0.7),sz=11,col=INK,line_sp=1.2)
bx(s,Inches(4.71),Inches(3.85),Inches(7.62),Inches(1.55),REDBG,line=RED,lw=1.0); bx(s,Inches(4.71),Inches(3.85),Inches(0.1),Inches(1.55),RED)
t(s,"でも、院内に「DX担当者」はいない。\n業者は高い・縦割り・自社製品売り＝院の課題は二の次。",Inches(4.95),Inches(3.85),Inches(7.3),Inches(1.55),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.3)

# ════ 3 なぜ隙間か ════
s=sl(); ft(s)
hdr(s,"我々の立ち位置","ここに、ぽっかり空いた「隙間」がある","院の現実と、既存業者のあいだ。そこに中立で入る。")
cols=[("院の現実","・院長は診療で多忙、DXに手が回らない\n・院内にIT/DX担当者がいない（紙・電話文化）\n・何から手をつけるか分からない",CARD,INK),
      ("既存の業者","・電カル/予約会社は高額・縦割り\n・「自社製品を買って」が出発点\n・院ごとの課題に合わせてくれない",CARD,INK)]
for i,(ti,ds,fill,tc) in enumerate(cols):
    cx=Inches(0.55)+(Inches(6.0)+Inches(0.23))*i
    bx(s,cx,Inches(2.0),Inches(6.0),Inches(2.05),fill,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.0),Inches(6.0),Inches(0.06),RED)
    t(s,ti,cx+Inches(0.3),Inches(2.15),Inches(5.5),Inches(0.45),sz=15,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.32),Inches(2.7),Inches(5.5),Inches(1.25),sz=11.5,col=INK,line_sp=1.3)
bx(s,Inches(0.55),Inches(4.4),Inches(12.23),Inches(1.7),TEALBG,line=TEALD,lw=1.5); bx(s,Inches(0.55),Inches(4.4),Inches(0.12),Inches(1.7),TEALD)
t(s,"その「あいだ」に、中立で入る。",Inches(0.85),Inches(4.55),Inches(11.6),Inches(0.5),sz=18,bold=True,col=TEALD)
t(s,"特定ツールを売らず、課題を診断して「刺さる1つ」から。AIで実物を0円で作り、効いた分だけ成果報酬。\n御用聞きでなく、院が少し楽になることだけを一緒に積み上げます。",Inches(0.87),Inches(5.05),Inches(11.6),Inches(1.0),sz=12.5,col=INK,line_sp=1.3)

# ════ 4 進め方（型） ════
s=sl(); ft(s)
hdr(s,"進め方（当社の型）","進め方は、いつも同じ「型」","この5ステップを、院ごとの課題に当てはめるだけ")
steps=[("①","無料診断","課題を一緒に棚卸し"),("②","刺さる1つを選ぶ","効果が大きい所から"),
       ("③","0円で小さく試す","AIで実物を作って動かす"),("④","数値で測る","導入前後を比べる"),
       ("⑤","効いた所を拡大","合わなければ見直す")]
cw,gx,x0=Inches(2.3),Inches(0.18),Inches(0.55)
for i,(no,ti,ds) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.3),cw,Inches(2.4),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.3),cw,Inches(0.06),RED)
    bx(s,cx+Inches(0.85),Inches(2.5),Inches(0.6),Inches(0.6),RED,shape=MSO_SHAPE.OVAL)
    t(s,no,cx+Inches(0.85),Inches(2.5),Inches(0.6),Inches(0.6),sz=18,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ti,cx+Inches(0.1),Inches(3.3),cw-Inches(0.2),Inches(0.6),sz=13.5,bold=True,col=INK,align=PP_ALIGN.CENTER,line_sp=1.05)
    t(s,ds,cx+Inches(0.12),Inches(3.95),cw-Inches(0.24),Inches(0.7),sz=10.5,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.15)
    if i<4: t(s,"→",cx+cw-Inches(0.02),Inches(3.1),Inches(0.3),Inches(0.5),sz=18,bold=True,col=RED,align=PP_ALIGN.CENTER)
bx(s,Inches(0.55),Inches(5.05),Inches(12.23),Inches(0.9),REDBG); bx(s,Inches(0.55),Inches(5.05),Inches(0.1),Inches(0.9),RED)
t(s,"※「刺さる1つ」は院ごとに違う。公式LINEとは限りません。まず診断ありき。",Inches(0.82),Inches(5.05),Inches(11.8),Inches(0.9),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 5 ソリューションメニュー ════
s=sl(); ft(s)
hdr(s,"打ち手の引き出し","課題に合わせて選ぶ ── ソリューション・メニュー","全部やるのでなく、課題に合うものから1つ。LINEはこの中の1つ。")
rows=[("院の課題","主な打ち手（ソリューション）"),
      ("受付・問い合わせ・電話の負担","公式LINE（FAQ自動応答／Web問診／予約／お知らせ・リマインド）"),
      ("電子カルテの入力・転記の負担","Stream Deck・AI入力／伝達・音声の自動抽出・入力連携"),
      ("集患・認知が弱い","オウンドメディア・HP・SNS × AI（発信の量産）"),
      ("開業・分院・立地の判断","診療圏調査（AIで高速・量産。ここで食えるかを数値化）"),
      ("承継・M&A・補助金","テナント承継コンサル＋IT導入補助金の活用支援")]
light_table(s,rows,Inches(0.55),Inches(2.0),Inches(12.23),Inches(3.5),[Inches(4.6),Inches(7.63)],hi_col=None,sz=12,header_sz=12.5)
t(s,"→ 今回ご縁のあった院では「公式LINE」が刺さりましたが、御院では別の打ち手が主役になるかもしれません。",Inches(0.55),Inches(5.7),Inches(12.2),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 6 実物を見せる（例：公式LINE） ════
s=sl(); ft(s)
hdr(s,"実物を作って、見せる","「提案」で終わらせない ── 例：公式LINE","AIで実物を0円で作って動かせる。これは数ある打ち手の一例です。")
SC="sc1.jpg"
if os.path.exists(SC):
    bx(s,Inches(8.7),Inches(1.95),Inches(4.1),Inches(4.7),CARD,line=CARDLN,lw=1.0)
    s.shapes.add_picture(SC,Inches(9.05),Inches(2.15),height=Inches(4.3))
    t(s,"実装例：京橋クリニックの公式LINE（本番稼働）",Inches(8.7),Inches(6.55),Inches(4.1),Inches(0.3),sz=9.5,col=GRY,align=PP_ALIGN.CENTER)
feats=[("提案だけでなく「実物」","他社は資料止まり。我々はAIで動く実物を0円で作る"),
       ("0円・中立","初期/月額0円（OSS×無料枠）。特定ツールを売らない"),
       ("現場目線で作る","スタッフの不安に配慮。仕事を奪わず、面倒を引き受ける助手"),
       ("これは「一例」","他院では、刺さる課題に応じて別の打ち手が主役に")]
y=Inches(2.05)
for ti,ds in feats:
    bx(s,Inches(0.55),y,Inches(7.9),Inches(1.05),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),y,Inches(0.1),Inches(1.05),RED)
    t(s,ti,Inches(0.8),y+Inches(0.13),Inches(7.5),Inches(0.4),sz=14,bold=True,col=REDD)
    t(s,ds,Inches(0.82),y+Inches(0.55),Inches(7.4),Inches(0.45),sz=11,col=INK,line_sp=1.15)
    y=y+Inches(1.16)

# ════ 7 幅を見せる（診療圏・承継） ════
s=sl(); ft(s)
hdr(s,"LINEだけじゃない","院のフェーズで、主役は変わる","開業前／運営中／承継 ── それぞれに打ち手があります")
phases=[("開業・分院を考える院","診療圏調査（AI量産）","「ここで本当に食えるか」を、年齢・男女別の受療率からAIで高速レポート化。出店判断の精度を上げる。"),
        ("運営中で忙しい院","業務のAI効率化","受付（公式LINE）・電カル入力（Stream Deck/音声）・発信（オウンドメディア）を、課題に合わせて。"),
        ("承継・出口を考える院","承継コンサル＋補助金","テナント承継・居抜きの座組み。IT導入補助金で導入コストを圧縮。次の一手まで伴走。")]
cw,gx,x0=Inches(3.96),Inches(0.18),Inches(0.55)
for i,(ph,sol,ds) in enumerate(phases):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.1),cw,Inches(3.4),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.1),cw,Inches(0.7),RED)
    t(s,ph,cx+Inches(0.15),Inches(2.2),cw-Inches(0.3),Inches(0.5),sz=12.5,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,sol,cx+Inches(0.2),Inches(3.0),cw-Inches(0.4),Inches(0.6),sz=15,bold=True,col=REDD,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.25),Inches(3.75),cw-Inches(0.5),Inches(1.6),sz=11,col=INK,align=PP_ALIGN.CENTER,line_sp=1.3)
t(s,"※ 一度きりでなく、院の成長に合わせて長くご一緒できるのが、AI医療コンサルの形です。",Inches(0.55),Inches(5.7),Inches(12.2),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 8 効果の測り方 ════
s=sl(); ft(s)
hdr(s,"効果の測り方","効果は、一緒に「数字」で確かめる","打ち手ごとに、無理なく測れる指標だけを選びます")
rows=[("見る指標（例）","どう良くなるか"),
      ("受付の電話本数 / 日","自動応答が一次対応 → 件数が減る"),
      ("入力・転記の時間","Web問診・AI入力 → 二度手間が減る"),
      ("新患数 / 問い合わせ数","発信（HP・SNS）→ 認知・来院が増える"),
      ("定期受診の来院率","リマインド → 取りこぼし防止で上がる"),
      ("月末の残業","業務平準化 → 月末の山が減る")]
light_table(s,rows,Inches(0.55),Inches(2.0),Inches(12.23),Inches(3.4),[Inches(4.6),Inches(7.63)],hi_col=None,sz=12,header_sz=12.5)
bx(s,Inches(0.55),Inches(5.6),Inches(12.23),Inches(0.7),GRYBG); bx(s,Inches(0.55),Inches(5.6),Inches(0.1),Inches(0.7),RED)
t(s,"どの指標を見るかは、スタッフの皆さまと相談して決めます。評価のためでなく「楽になったか」を確かめるためです。",Inches(0.82),Inches(5.6),Inches(11.8),Inches(0.7),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 9 なぜ我々か ════
s=sl(); ft(s)
hdr(s,"選ばれる理由","なぜ、KHDのAI医療コンサルか","売り込みでなく、信頼を積み上げる関係を目指します")
whys=[("中立","特定ツールを売らない。院の課題から最適な打ち手を"),
      ("0円・成果報酬","初期/月額0円ではじめ、効いた分だけの成果報酬も可"),
      ("AIで速く・柔軟","従来20〜250万のものを0円〜。その場で直せる"),
      ("現場目線","スタッフの不安に配慮。仕事を奪わず、面倒を引き受ける"),
      ("実物で見せる","資料でなく、動く実物を作って判断材料にできる"),
      ("GIVE先行","売り込まず、まず役に立つ。信頼の対価として")]
cw,ch,gx,gy=Inches(3.96),Inches(1.4),Inches(0.2),Inches(0.25); x0,y0=Inches(0.55),Inches(2.0)
for i,(ti,ds) in enumerate(whys):
    cx=x0+(cw+gx)*(i%3); cy=y0+(ch+gy)*(i//3)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.25),cy+Inches(0.15),cw-Inches(0.4),Inches(0.45),sz=15,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.27),cy+Inches(0.62),cw-Inches(0.45),Inches(0.7),sz=11,col=INK,line_sp=1.2)

# ════ 10 はじめ方 ════
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"はじめ方",Inches(0.9),Inches(0.85),Inches(11),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"まずは「無料診断」から。",Inches(0.88),Inches(1.35),Inches(11.6),Inches(0.8),sz=30,bold=True,col=INK)
t(s,"いきなり契約ではありません。課題を一緒に棚卸しして、刺さる1つを見つけます。",Inches(0.92),Inches(2.25),Inches(11.4),Inches(0.5),sz=14,col=GRY)
steps=[("STEP 1","無料診断（約30分）","院の課題を一緒に棚卸し。優先順位をつける"),
       ("STEP 2","刺さる1つを0円で","AIで実物を作って試す。負担なく小さく"),
       ("STEP 3","測って、広げる","効いた所だけ拡大。合わなければ見直す")]
cw,gx,x0=Inches(3.95),Inches(0.24),Inches(0.9)
for i,(st,ti,ds) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(3.05),cw,Inches(2.2),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(3.05),cw,Inches(0.7),RED)
    t(s,st,cx,Inches(3.17),cw,Inches(0.45),sz=18,bold=True,col=WHT,align=PP_ALIGN.CENTER)
    t(s,ti,cx+Inches(0.15),Inches(3.9),cw-Inches(0.3),Inches(0.5),sz=14,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.25),Inches(4.45),cw-Inches(0.5),Inches(0.8),sz=11.5,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.2)
bx(s,Inches(0.9),Inches(5.65),Inches(11.5),Inches(0.9),REDBG); bx(s,Inches(0.9),Inches(5.65),Inches(0.1),Inches(0.9),RED)
t(s,"御院に合うかどうか、一緒に確かめながら。まず小さく、いつでも見直せます。",Inches(1.2),Inches(5.65),Inches(11),Inches(0.9),sz=15,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("ai_iryo_consul_template.pptx")
print("saved ai_iryo_consul_template.pptx slides:",len(prs.slides._sldIdLst))

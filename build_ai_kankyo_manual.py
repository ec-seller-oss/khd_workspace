# -*- coding: utf-8 -*-
"""KHD AI仕事環境 整備マニュアル v1（横文字ひかえめ・目次型で厚くしていける土台）
KHDクリーム白×レンガ赤。営業・コンサルでも使える"型"として構築。
出力: KHD_AI仕事環境整備マニュアル_v1_260628.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); GRN=RGBColor(0x2E,0x7D,0x46)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold
        r.font.color.rgb=col; r.font.name=FONT
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.6),sz=24,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.5),Inches(12.1),Inches(0.35),sz=12.5,col=GRY)

def ft(slide,n):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"KHD AI仕事環境 整備マニュアル  ｜  v1  ｜  2026-06-28",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
    t(slide,str(n),Inches(12.4),H-Inches(0.42),Inches(0.5),Inches(0.32),sz=9,col=GRY,align=PP_ALIGN.RIGHT)

def panel(s,x,y,w,h,title,lines,tcol=RED,lh=0.56,tsz=16,bsz=13.5):
    bx(s,x,y,w,h,CARD,line=CARDLN,lw=1.0); bx(s,x,y,w,Inches(0.07),tcol)
    t(s,title,x+Inches(0.3),y+Inches(0.18),w-Inches(0.5),Inches(0.5),sz=tsz,bold=True,col=tcol)
    yy=y+Inches(0.85)
    for ln,col in lines:
        t(s,ln,x+Inches(0.32),yy,w-Inches(0.6),Inches(lh),sz=bsz,col=col,line_sp=1.13); yy=yy+Inches(lh)

def band(s,y,text,sub="",h=0.95):
    bx(s,Inches(0.55),y,Inches(12.23),Inches(h),REDBG); bx(s,Inches(0.55),y,Inches(0.1),Inches(h),RED)
    t(s,text,Inches(0.9),y+Inches(0.12),Inches(11.6),Inches(0.45),sz=15,bold=True,col=REDD,line_sp=1.1)
    if sub: t(s,sub,Inches(0.9),y+Inches(0.56),Inches(11.6),Inches(0.35),sz=12,col=INK,line_sp=1.1)

def table(s,x,y,rows,cw,rh=Inches(0.56),hot_row=None,hsz=12,bsz=11.5):
    for ri,row in enumerate(rows):
        cx=x; head=(ri==0); hot=(hot_row==ri)
        for ci,cell in enumerate(row):
            bg=RED if head else (REDBG if hot else (CARD if ri%2 else GRYBG))
            bx(s,cx,y+rh*ri,cw[ci],rh,bg,line=CARDLN,lw=0.8)
            tc=WHT if head else (REDD if hot and ci==0 else INK)
            t(s,cell,cx+Inches(0.16),y+rh*ri,cw[ci]-Inches(0.28),rh,sz=hsz if head else bsz,
              bold=head or ci==0,col=tc,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.0)
            cx=cx+cw[ci]

# ── 1 表紙 ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"KHD ｜ 仕事をAIで回すための土台づくり",Inches(0.9),Inches(1.5),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"AI仕事環境 整備マニュアル",Inches(0.88),Inches(2.15),Inches(11.7),Inches(0.95),sz=42,bold=True,col=INK)
t(s,"― はじめての人でも、同じ環境を迷わず作れる地図 ―",Inches(0.9),Inches(3.15),Inches(11.7),Inches(0.6),sz=20,bold=True,col=RED)
t(s,"パソコンの選び方・役割分担、データの置き場、外からの使い方、付属品まで。\nこの目次に章を足しながら“厚く”していきます。営業・コンサルでも使える型に。",
  Inches(0.9),Inches(4.0),Inches(11.6),Inches(0.9),sz=14,col=GRY,line_sp=1.3)
bx(s,Inches(0.9),Inches(5.6),Inches(11.5),Pt(1.2),LINE)
t(s,"KHD  ｜  菊池 研太  ｜  v1  2026-06-28",Inches(0.9),Inches(5.75),Inches(11),Inches(0.4),sz=13,bold=True,col=INK)

# ── 2 このマニュアルの使い方 ──
s=sl(); ft(s,2)
hdr(s,"はじめに","このマニュアルの“ねらい”と使い方","「人は打席（会って話す）だけ。残りはAIが全部やる」を支える“道具と環境”の地図")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(2.6),"何のため？",[
 ("AIで仕事を回すには、まず“環境”が要る",INK),
 ("（どのパソコンで・どこに保存し・",GRY),
 ("　外からどう使うか）",GRY),
 ("ここが整うと、菊池さんは“打席”に",INK),
 ("　集中でき、空き時間が増える",GRN),
],lh=0.5)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(2.6),"どう使う？",[
 ("困ったら、次ページの“目次”から開く",INK),
 ("章は、これから足して厚くしていく",INK),
 ("（今日できた所は ✓、これからは ＋）",GRY),
 ("そのまま 他の営業・コンサル先にも",INK),
 ("　配れる“型”として育てる",GRN),
],tcol=GRN,lh=0.5)
band(s,Inches(4.85),"大原則：むずかしい横文字は使わない。「机の広さ」「倉庫」「窓口」など、身近な言葉で。",
     "新しく関わる人（ゆーし・スタッフ・お客様）が、読んだだけで動ける状態を保つ。",h=1.0)

# ── 3 目次 ──
s=sl(); ft(s,3)
hdr(s,"CONTENTS","目次 ── この順番で“環境”が整う","✓＝今回入れた章／＋＝これから厚くする章")
chaps=[
 ("第1章","全体像（主役の1台＋のぞき窓）","✓"),
 ("第2章","速さの正体（机の広さ と 倉庫）","✓"),
 ("第3章","データの置き場（全部クラウドへ）","✓"),
 ("第4章","パソコンの役割分担（何台でどう使う）","✓"),
 ("第5章","外から使う（遠くから主役に入る）","✓"),
 ("第6章","付属品・周辺機器（ケーブル・画面 等）","✓"),
 ("第7章","新しい機械の初期設定（同じ状態に）","＋"),
 ("第8章","毎日の回し方（朝・日中・夜）","＋"),
 ("第9章","安全とバックアップ（失わない・分ける）","＋"),
 ("第10章","困った時（つまずきと直し方）","＋"),
]
colx=[Inches(0.55),Inches(6.75)]; cw=Inches(5.95)
for i,(no,title,mk) in enumerate(chaps):
    col=i//5; row=i%5
    x=colx[col]; y=Inches(1.95)+Inches(0.92)*row
    done=(mk=="✓")
    bx(s,x,y,cw,Inches(0.78),CARD if done else GRYBG,line=CARDLN,lw=0.8); bx(s,x,y,Inches(0.07),Inches(0.78),RED if done else GRY)
    t(s,no,x+Inches(0.25),y,Inches(1.3),Inches(0.78),sz=14,bold=True,col=RED if done else GRY,anchor=MSO_ANCHOR.MIDDLE)
    t(s,title,x+Inches(1.5),y,cw-Inches(2.1),Inches(0.78),sz=12.5,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.0)
    mc=GRN if done else GRY
    t(s,mk,x+cw-Inches(0.55),y,Inches(0.5),Inches(0.78),sz=18,bold=True,col=mc,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ── 4 第1章 全体像 ──
s=sl(); ft(s,4)
hdr(s,"第1章｜全体像","主役の1台に、みんなが“入って”使う","“主役パソコン（母艦）”がAIとデータを持ち、他は窓からのぞくだけ")
# 中央=母艦
bx(s,Inches(5.0),Inches(2.3),Inches(3.3),Inches(2.0),REDBG,line=RED,lw=1.5)
t(s,"主役パソコン\n（母艦）",Inches(5.0),Inches(2.5),Inches(3.3),Inches(0.9),sz=18,bold=True,col=REDD,align=PP_ALIGN.CENTER)
t(s,"AI（クロード）と\n会社のデータが\nここに集約・常時稼働",Inches(5.0),Inches(3.35),Inches(3.3),Inches(0.9),sz=11.5,col=INK,align=PP_ALIGN.CENTER,line_sp=1.1)
# 周りの窓口
clients=[("自分のスマホ",Inches(0.7),Inches(2.1)),("持ち出しパソコン",Inches(0.7),Inches(4.4)),
         ("ゆーしのパソコン",Inches(9.9),Inches(2.1)),("自宅の予備機",Inches(9.9),Inches(4.4))]
for lab,x,y in clients:
    bx(s,x,y,Inches(2.7),Inches(1.1),CARD,line=CARDLN,lw=1.0)
    t(s,lab,x,y+Inches(0.15),Inches(2.7),Inches(0.4),sz=13,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,"＝のぞき窓\n（入って使うだけ）",x,y+Inches(0.5),Inches(2.7),Inches(0.5),sz=10.5,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.0)
band(s,Inches(5.95),"だから“窓口”側は軽い。重い処理は全部“主役”が引き受ける。",
     "覚えること：強い1台を主役にして、他はそこに入る。これがAI環境の背骨。",h=1.0)

# ── 5 第2章 速さの正体 ──
s=sl(); ft(s,5)
hdr(s,"第2章｜速さの正体","“机の広さ”と“倉庫の大きさ”は別もの","パソコンが遅いのは「机が狭い」から。「倉庫が一杯」とは別の話")
table(s,Inches(0.55),Inches(2.0),[
 ("","机の広さ（速さ）","倉庫の大きさ（保存）"),
 ("たとえ","一度に広げて作業できる量","しまっておける量"),
 ("足りないと","固まる・もたつく・遅い","保存できない（速さは無関係）"),
 ("増やすには","強いパソコン（主役）に替える","クラウド（ネット倉庫）を使う"),
 ("今の旧Mac","机が狭い（8GB）＝ここが遅さの正体","倉庫はクラウドで足りる"),
],cw=[Inches(2.6),Inches(4.8),Inches(4.83)],hot_row=4)
band(s,Inches(5.7),"大事：データを全部クラウドに移しても“速さ”は直らない。",
     "速さ＝机の広さ。倉庫（クラウド）をいくら広げても、机は広がらない。遅さは“主役を強い機械にする”で直す。",h=1.05)

# ── 6 第3章 データの置き場 ──
s=sl(); ft(s,6)
hdr(s,"第3章｜データの置き場","データは全部クラウド（ネット倉庫）に","机の上に物を積まない＝パソコンが軽くなる")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.0),"こうする",[
 ("ファイルは Google ドライブ に置く",INK),
 ("「開く時だけ取り寄せ」の設定にする",INK),
 ("→ パソコンの中に物をためない＝軽い",GRN),
 ("3台のどれからでも同じファイルが見える",INK),
 ("倉庫はたっぷり（容量で困らない）",GRY),
],tcol=GRN,lh=0.52)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.0),"2つだけ例外（クラウド任せにしない）",[
 ("① 作業中の“記録の金庫”",INK),
 ("　＝同時に書き換えると壊れるので、",GRY),
 ("　　主役パソコンの中で扱う",GRY),
 ("② AIやアプリ本体（ソフト）",INK),
 ("　＝そもそも機械の中で動くもの",GRY),
],tcol=REDD,lh=0.52)
band(s,Inches(5.25),"合言葉：「ファイルはクラウド、机の上はいつも空っぽ」。",
     "これで保存に困らず、パソコンも軽く保てる。",h=1.0)

# ── 7 第4章 役割分担 ──
s=sl(); ft(s,7)
hdr(s,"第4章｜パソコンの役割分担","何台を、どう使い分けるか","強い1台を“主役”に。他は軽い“窓口”か“予備”")
table(s,Inches(0.55),Inches(2.0),[
 ("パソコン","役割","使い方"),
 ("新しい強いパソコン","主役（母艦）","AIとデータを持ち、ずっと動かす。重い処理はここ"),
 ("今までの旧パソコン","のぞき窓 か 予備","主役に入って使うだけ＝軽い。重い仕事はさせない"),
 ("ゆーしのパソコン","のぞき窓","主役に入って使う（くわしくは別マニュアル）"),
 ("スマホ","のぞき窓","外出先から主役に入る。アプリ1つで軽い"),
],cw=[Inches(3.0),Inches(2.6),Inches(6.63)],hot_row=1)
band(s,Inches(5.7),"今日のおすすめ：遅い旧パソコンを“主役”からはずし、新しい強いパソコンを主役にする。",
     "旧パソコンは「のぞき窓」か「予備」へ。これだけで旧パソコンがぐっと軽くなる。",h=1.05)

# ── 8 第5章 外から使う ──
s=sl(); ft(s,8)
hdr(s,"第5章｜外から使う","出先のスマホから、自宅の主役に“入る”","主役は家でずっと動かし、外からはそこへ入って使う")
panel(s,Inches(0.55),Inches(1.95),Inches(7.7),Inches(3.1),"やり方（くわしくは別マニュアル）",[
 ("主役パソコンは 家に置いて、ずっと起動",INK),
 ("外からは スマホ／別パソコンで そこへ入る",INK),
 ("画面ごと見る方法、AIだけ操作する方法 など",GRY),
 ("入る側は“窓”なので、通信さえあれば軽い",GRN),
 ("電源を切ると入れない→主役は つけっぱなしに",REDD),
],lh=0.54,bsz=14)
panel(s,Inches(8.45),Inches(1.95),Inches(4.33),Inches(3.1),"良いところ",[
 ("打席（商談）の直後に、その場で",INK),
 ("　一言だけ指示できる",GRY),
 ("帰宅前に“残りの仕事”が",INK),
 ("　片付いている",GRN),
],tcol=GRN,lh=0.54,bsz=14)
band(s,Inches(5.35),"※ くわしい手順は別冊「外から使う 設定ガイド」に。この章は“考え方”だけ。",h=0.62)

# ── 9 第6章 付属品・周辺機器 ──
s=sl(); ft(s,9)
hdr(s,"第6章｜付属品・周辺機器","AI環境に“付随して整える物”リスト","本体だけでなく、周りの道具も揃えて初めて回る")
table(s,Inches(0.55),Inches(1.95),[
 ("品目","何のため","要否の目安"),
 ("高速ケーブル（サンダーボルト等）","新しいパソコンへ“丸ごと引っ越し”を最速で","引っ越し時に必須"),
 ("外付けの大きい画面","主役パソコンは画面が広いほど作業が速い","主役にあると良い"),
 ("キーボード・マウス","長時間の入力を楽に・正確に","主役にあると良い"),
 ("有線LAN／安定したネット","常時起動の主役は通信が命綱","主役は必須級"),
 ("電源・電源タップ","主役はつけっぱなし＝電源を確保","必須"),
 ("バックアップ用の外付け保存","万一に備えてデータの控えを取る","あると安心"),
],cw=[Inches(3.7),Inches(5.5),Inches(3.03)],rh=Inches(0.6),hot_row=1)
t(s,"※ この章は、新しい道具を買う・もらうたびに行を足して“買い物リスト＆理由”として育てる。",
  Inches(0.6),Inches(6.7),Inches(12),Inches(0.3),sz=10,col=GRY)

# ── 10 これから厚くする章（予告）──
s=sl(); ft(s,10)
hdr(s,"これから厚くする章","第7〜10章 ── 足しながら育てる","今は見出しだけ。運用で分かった事を、章に足していく")
nexts=[
 ("第7章","新しい機械の初期設定","新しいパソコンを“いつもと同じ状態”にする段取り（引っ越し→ログイン→データ取り込み）"),
 ("第8章","毎日の回し方","朝（今日の打席を決める）・日中（その場で一言）・夜（結果を一言）の流れ"),
 ("第9章","安全とバックアップ","データを失わない／第三者（外の人）と“見える範囲”を分ける考え方"),
 ("第10章","困った時","つまずきやすい所と直し方。エラーは触らず写真を撮って相談、など"),
]
yy=Inches(2.0)
for no,title,desc in nexts:
    bx(s,Inches(0.55),yy,Inches(12.23),Inches(1.0),GRYBG,line=CARDLN,lw=0.8); bx(s,Inches(0.55),yy,Inches(0.1),Inches(1.0),GRY)
    t(s,no,Inches(0.8),yy,Inches(1.4),Inches(1.0),sz=15,bold=True,col=GRY,anchor=MSO_ANCHOR.MIDDLE)
    t(s,title,Inches(2.2),yy+Inches(0.14),Inches(10),Inches(0.4),sz=15,bold=True,col=INK)
    t(s,desc,Inches(2.2),yy+Inches(0.54),Inches(10.2),Inches(0.4),sz=11.5,col=GRY)
    yy=yy+Inches(1.12)

# ── 11 締め ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"この1冊が、AI環境の“標準”になる",Inches(0.9),Inches(2.2),Inches(11.7),Inches(0.8),sz=30,bold=True,col=INK)
t(s,"足しながら、営業・コンサルの“武器”に。",Inches(0.9),Inches(3.05),Inches(11.7),Inches(0.8),sz=30,bold=True,col=RED)
t(s,"むずかしい横文字を使わず、誰が読んでも動ける。\nこの型をそのままお客様にも渡せる＝“環境づくり”自体が、商品になる。",
  Inches(0.9),Inches(4.15),Inches(11.6),Inches(0.9),sz=15,col=INK,line_sp=1.3)
bx(s,Inches(0.9),Inches(5.6),Inches(11.5),Pt(1.2),LINE)
t(s,"KHD AI仕事環境 整備マニュアル  ｜  v1  2026-06-28  ｜  以後、章を追加",Inches(0.9),Inches(5.72),Inches(11),Inches(0.4),sz=12,bold=True,col=INK)

out="KHD_AI仕事環境整備マニュアル_v1_260628.pptx"
prs.save(out)
print("saved:",out,"slides:",len(prs.slides._sldIdLst))

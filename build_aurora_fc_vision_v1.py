"""
オーロラ訪問マッサージ 事業立ち上げ構想 ── 面談用「当社方針」説明デッキ v1
用途: 採用面談（石原さん等・即戦力候補）で見せる、KHDの事業ビジョン共有資料。
立て付け: 口説き／FC加盟確定の事業計画書 ではなく「一緒に創る仲間像」を共有する相互見極め資料。
中核信条: 相手目線の商売人・売り込まない・信頼の対価で稼ぐ（feedback_aite_mokusen）。
デザイン: クリニックDXデッキ準拠（クリーム白 #F9F6EF × レンガ赤 #AA2E26 / Hiragino Sans / 16:9）。
出力: aurora_fc_vision_v1.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 配色（参考デッキから抽出した実配色を踏襲）──
BG     = RGBColor(0xF9, 0xF6, 0xEF)   # クリーム背景
RED    = RGBColor(0xAA, 0x2E, 0x26)   # レンガ赤・主アクセント
REDD   = RGBColor(0x8C, 0x24, 0x1D)   # 濃赤
INK    = RGBColor(0x1A, 0x1A, 0x1A)   # 本文黒
GRY    = RGBColor(0x6E, 0x6E, 0x6E)   # サブ灰
LINE   = RGBColor(0xDA, 0xD6, 0xCF)   # 罫線
CARD   = RGBColor(0xF1, 0xEC, 0xE1)   # ベージュカード
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)   # カード枠
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)   # 薄赤（強調セル）
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)   # 薄灰ベージュ
WHT    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Hiragino Sans"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "オーロラ訪問マッサージ  ｜  KHD 事業立ち上げ構想", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)


def light_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = REDD if is_hi else RED
            else:
                cell.fill.fore_color.rgb = REDBG if is_hi else (CARD if ri % 2 == 1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = WHT
                    elif is_hi:
                        r.font.color.rgb = RED
                    elif ci == 0:
                        r.font.color.rgb = INK
                    else:
                        r.font.color.rgb = RGBColor(0x3A,0x3A,0x3A)
    return tb


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "AURORA ｜ 訪問マッサージ事業 立ち上げ構想", Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.45), sz=15, bold=True, col=RED)
t(s, "在宅に、医療の手を。", Inches(0.88), Inches(2.15), Inches(11.6), Inches(0.95), sz=44, bold=True, col=INK)
t(s, "船橋から始める、訪問マッサージ事業。", Inches(0.88), Inches(3.15), Inches(11.6), Inches(0.9), sz=30, bold=True, col=RED)
t(s, "通院が難しい高齢の方のもとへ、こちらから出向く。\nその第一歩を、一緒に踏み出してくれる仲間を探しています。",
  Inches(0.9), Inches(4.25), Inches(11.3), Inches(1.0), sz=15, col=GRY, line_sp=1.3)
bx(s, Inches(0.9), Inches(6.6), Inches(11.5), Pt(1.2), LINE)
t(s, "KHD（株式会社）  ｜  代表  菊池 研太", Inches(0.9), Inches(6.75), Inches(11), Inches(0.4), sz=14, bold=True, col=INK)

# ════════ SLIDE 2 — KHD / 菊池の自己紹介・資質 ════════
s = sl(); ft(s)
hdr(s, "WHO WE ARE", "私たちは何者か ── KHDと、代表・菊池研太", "「どんな人と、どんな会社で働くのか」をまず正直にお伝えします")
quals = [
    ("不動産デベロッパー", "現場で稼いできた実績", "仕入れ〜加工〜売却を一人称で回す事業家。数字とリスクで判断し、机上論で終わらせない。"),
    ("医療×事業の知見", "クリニック承継・医療コンサル", "医療クリニックの承継・開業支援を手がける。医療保険・在宅医療の世界観を理解している。"),
    ("仕組み化・AI活用", "少人数でも回る設計", "AIと仕組みで雑務を圧縮し、人は人にしかできない仕事に集中させる。立ち上げを軽くする。"),
    ("相手目線・誠実", "売り込まない商売人", "妻と幼い子を持つ一人の生活者。長く安心して働ける組織を、信頼を土台に創りたい。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.4), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.85)
for i, (ti, sub, body) in enumerate(quals):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    t(s, ti, cx+Inches(0.35), cy+Inches(0.28), cw-Inches(0.6), Inches(0.5), sz=18, bold=True, col=INK)
    t(s, sub, cx+Inches(0.37), cy+Inches(0.85), cw-Inches(0.6), Inches(0.4), sz=12.5, bold=True, col=RED)
    t(s, body, cx+Inches(0.37), cy+Inches(1.32), cw-Inches(0.7), Inches(0.95), sz=12, col=GRY, line_sp=1.18)

# ════════ SLIDE 3 — なぜ訪問マッサージか ════════
s = sl(); ft(s)
hdr(s, "WHY THIS BUSINESS", "なぜ今、訪問マッサージ事業なのか", "「儲かりそうだから」ではなく、伸びる必然と、やる意味の両方があるから")
# 左：3つの理由
reasons = [
    ("① 需要が伸び続ける", "後期高齢者は増え続け、通院が難しい在宅の方が増えていく。訪問のニーズは構造的に右肩上がり。"),
    ("② 医療保険ベースの安定収益", "医療保険適用の訪問マッサージは、単価が制度で支えられ、景気に左右されにくい安定したストック型。"),
    ("③ 社会的に意味がある", "「動けないから諦める」をなくす。手を必要としている人のもとへ、こちらから出向く仕事。"),
]
y = Inches(1.95)
for ti, body in reasons:
    bx(s, Inches(0.55), y, Inches(7.4), Inches(1.45), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.55), y, Inches(0.1), Inches(1.45), RED)
    t(s, ti, Inches(0.8), y+Inches(0.18), Inches(7.0), Inches(0.45), sz=16, bold=True, col=REDD)
    t(s, body, Inches(0.82), y+Inches(0.68), Inches(7.0), Inches(0.7), sz=12, col=INK, line_sp=1.15)
    y = y + Inches(1.6)
# 右：強調パネル
bx(s, Inches(8.25), Inches(1.95), Inches(4.53), Inches(4.65), RED)
t(s, "つまり", Inches(8.25), Inches(2.3), Inches(4.53), Inches(0.4), sz=14, bold=True, col=RGBColor(0xF2,0xD8,0xD6), align=PP_ALIGN.CENTER)
t(s, "伸びる市場 ×\n安定収益 ×\n社会的意義", Inches(8.25), Inches(2.85), Inches(4.53), Inches(1.6), sz=26, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.2)
bx(s, Inches(8.6), Inches(4.75), Inches(3.83), Pt(1.2), RGBColor(0xD9,0x8A,0x84))
t(s, "この3つが揃う事業は\nそう多くありません。", Inches(8.25), Inches(5.0), Inches(4.53), Inches(1.0), sz=15, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.25)

# ════════ SLIDE 4 — なぜオーロラの仕組みか ════════
s = sl(); ft(s)
hdr(s, "WHY AURORA", "なぜ“自前ゼロから”ではなく、オーロラの仕組みを使うのか", "立ち上げの遠回り・失敗を避け、現場の価値提供に最短で集中するため")
rows = [
    ("論点", "ゼロから自前で立ち上げ", "オーロラの仕組みを活用"),
    ("立ち上げ速度", "制度理解・体制づくりに時間", "確立した型に乗って最短で開始"),
    ("レセプト・制度対応", "手探りで算定ミスのリスク", "実績あるオペレーションで安心"),
    ("集客・患者獲得", "ケアマネ営業もゼロから", "立ち上げ支援・ノウハウあり"),
    ("私たちが集中すべき所", "事務に追われ現場が手薄に", "良い施術と信頼づくりに集中"),
]
light_table(s, rows, Inches(0.55), Inches(1.95), Inches(12.23), Inches(3.4),
            [Inches(3.0), Inches(4.6), Inches(4.63)], hi_col=2, sz=13, header_sz=13)
bx(s, Inches(0.55), Inches(5.6), Inches(12.23), Inches(1.05), REDBG)
bx(s, Inches(0.55), Inches(5.6), Inches(0.1), Inches(1.05), RED)
t(s, "仕組みは借りる。でも、現場の信頼と質は自分たちで創る。", Inches(0.85), Inches(5.72), Inches(11.6), Inches(0.4), sz=15, bold=True, col=REDD)
t(s, "フランチャイズに“やってもらう”のではなく、土台を借りて、私たちの色で地域一番の訪問チームを育てる発想です。",
  Inches(0.85), Inches(6.18), Inches(11.6), Inches(0.4), sz=12, col=INK)

# ════════ SLIDE 5 — 事業構想（船橋） ════════
s = sl(); ft(s)
hdr(s, "OUR PLAN", "どこで、どう始めるか ── 船橋を拠点にした事業構想", "広げる前に、まず一拠点で“地域で一番頼られる訪問チーム”を創る")
# 左：数字パネル3つ
stats = [
    ("船橋市", "テリトリー承認済", "事業を始められる地域として確保済み"),
    ("後期高齢者", "8万人超", "通院困難な在宅ニーズが厚いエリア"),
    ("起点", "3駅エリア", "無理に広げず、まず手の届く範囲から"),
]
y = Inches(1.95)
for big, val, body in stats:
    bx(s, Inches(0.55), y, Inches(5.7), Inches(1.4), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.55), y, Inches(0.1), Inches(1.4), RED)
    t(s, big, Inches(0.82), y+Inches(0.2), Inches(2.3), Inches(0.9), sz=15, bold=True, col=GRY, anchor=MSO_ANCHOR.MIDDLE)
    t(s, val, Inches(2.7), y+Inches(0.12), Inches(3.4), Inches(0.65), sz=26, bold=True, col=RED)
    t(s, body, Inches(2.72), y+Inches(0.82), Inches(3.4), Inches(0.5), sz=10.5, col=GRY, line_sp=1.05)
    y = y + Inches(1.55)
# 右：ステップ
bx(s, Inches(6.55), Inches(1.95), Inches(6.23), Inches(4.65), GRYBG, line=CARDLN, lw=1.0)
t(s, "立ち上げのステップ", Inches(6.85), Inches(2.15), Inches(5.6), Inches(0.4), sz=15, bold=True, col=REDD)
steps = [
    ("STEP 1", "船橋1拠点で立ち上げ", "良い施術 × ケアマネ・施設との信頼で、地域の患者基盤を固める"),
    ("STEP 2", "チーム化・育成", "施術者を増やし、運営を仕組み化。立ち上げメンバーは中核へ"),
    ("STEP 3", "近隣エリアへ展開", "1拠点の型ができたら、無理のない範囲で面を広げる"),
]
yy = Inches(2.7)
for st, ti, body in steps:
    bx(s, Inches(6.85), yy, Inches(1.25), Inches(0.5), RED)
    t(s, st, Inches(6.85), yy+Inches(0.04), Inches(1.25), Inches(0.42), sz=12, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, Inches(8.25), yy, Inches(4.3), Inches(0.4), sz=14, bold=True, col=INK)
    t(s, body, Inches(8.27), yy+Inches(0.42), Inches(4.3), Inches(0.7), sz=10.5, col=GRY, line_sp=1.1)
    yy = yy + Inches(1.25)

# ════════ SLIDE 6 — 中核信条（商売観） ════════
s = sl(); ft(s)
hdr(s, "OUR CREED", "私たちの仕事観 ── 信頼を土台に、長く続ける", "短期で刈り取る商売はしない。患者にも、仲間にも、誠実であり続ける")
creeds = [
    ("患者さんへ", "売り込まない・盛らない", "必要な施術を、必要なだけ。数字のために過剰な訪問はしない。信頼が次の紹介を生む。"),
    ("仲間へ", "無理をさせない", "短期の数字で人をすり潰さない。家族との時間も大事にしながら、長く働ける環境をつくる。"),
    ("地域へ", "ケアマネ・施設と誠実に", "目先の囲い込みより、地域で一番信頼される存在に。それが結局いちばん強い集客になる。"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.95)
for i, (who, head, body) in enumerate(creeds):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(3.6), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.06), RED)
    t(s, who, cx, y0+Inches(0.3), cw, Inches(0.4), sz=13, bold=True, col=GRY, align=PP_ALIGN.CENTER)
    t(s, head, cx+Inches(0.3), y0+Inches(0.85), cw-Inches(0.6), Inches(0.9), sz=19, bold=True, col=REDD, align=PP_ALIGN.CENTER, line_sp=1.1)
    bx(s, cx+Inches(0.6), y0+Inches(1.85), cw-Inches(1.2), Pt(1.2), LINE)
    t(s, body, cx+Inches(0.35), y0+Inches(2.05), cw-Inches(0.7), Inches(1.4), sz=12, col=INK, line_sp=1.2)
bx(s, Inches(0.55), Inches(5.85), Inches(12.23), Inches(0.8), REDBG)
t(s, "信頼が先、収益は結果。── これは不動産でも医療でも、私が一貫して大事にしてきた軸です。",
  Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.5), sz=14, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)

# ════════ SLIDE 7 — 一緒に創りたい仲間像 ════════
s = sl(); ft(s)
hdr(s, "WHO WE NEED", "一緒に第一歩を踏み出す、仲間を探しています", "立ち上げは「誰とやるか」が9割。だから、正直にお話しします")
# 左：求める人物像
bx(s, Inches(0.55), Inches(1.95), Inches(6.1), Inches(4.65), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(0.55), Inches(1.95), Inches(6.1), Inches(0.06), RED)
t(s, "こんな方と創りたい", Inches(0.85), Inches(2.18), Inches(5.5), Inches(0.4), sz=16, bold=True, col=REDD)
wants = [
    "国家資格を持ち、訪問の現場を任せられる方",
    "施術だけでなく、ケアマネ・施設との関係づくりも一緒に考えられる方",
    "「立ち上げ」をゼロイチで創るおもしろさに、わくわくできる方",
    "目先の数字より、患者・地域の信頼を大事にできる方",
    "将来、拠点長・幹部として中核を担う意志のある方",
]
yy = Inches(2.7)
for wq in wants:
    t(s, "●", Inches(0.9), yy, Inches(0.4), Inches(0.4), sz=13, bold=True, col=RED)
    t(s, wq, Inches(1.3), yy-Inches(0.02), Inches(5.2), Inches(0.7), sz=12.5, col=INK, line_sp=1.12)
    yy = yy + Inches(0.76)
# 右：メッセージパネル
bx(s, Inches(6.95), Inches(1.95), Inches(5.83), Inches(4.65), RED)
t(s, "正直に言うと", Inches(6.95), Inches(2.3), Inches(5.83), Inches(0.4), sz=14, bold=True, col=RGBColor(0xF2,0xD8,0xD6), align=PP_ALIGN.CENTER)
t(s, "この事業を本気でやる確信は、\n“いい人に出会えるか”に\nかかっています。",
  Inches(7.2), Inches(2.85), Inches(5.4), Inches(1.7), sz=20, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.3)
bx(s, Inches(7.4), Inches(4.75), Inches(4.93), Pt(1.2), RGBColor(0xD9,0x8A,0x84))
t(s, "だから今日は、採用の合否を決める場ではなく、\nお互いに「この人と組めるか」を確かめる場にしたい。",
  Inches(7.2), Inches(5.0), Inches(5.4), Inches(1.3), sz=13, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.3)

# ════════ SLIDE 8 — 働き方・処遇の考え方 ════════
s = sl(); ft(s)
hdr(s, "HOW WE WORK", "働き方と、処遇の考え方", "細かい条件は対話で。でも考え方は最初に揃えておきたい")
ways = [
    ("立ち上げメンバーへの還元", "事業が育てば、その価値は一緒に創った人に返す。立ち上げの貢献は処遇・ポジションで報いる。"),
    ("成長に応じた役割", "施術者 → リーダー → 拠点長・幹部。本人の意志と実力に応じて、任せる範囲を広げていく。"),
    ("無理をさせない設計", "短期の数字で人をすり減らさない。AI・仕組みで雑務を減らし、現場と生活の両立を守る。"),
    ("条件はフェアに、対話で", "数字ありきの押し付けはしない。お互いが納得できる形を、これから一緒に決めていきたい。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.0), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.95)
for i, (ti, body) in enumerate(ways):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    t(s, str(i+1), cx+Inches(0.32), cy+Inches(0.25), Inches(0.8), Inches(0.8), sz=30, bold=True, col=RED)
    t(s, ti, cx+Inches(1.2), cy+Inches(0.28), cw-Inches(1.4), Inches(0.5), sz=15, bold=True, col=INK)
    t(s, body, cx+Inches(1.22), cy+Inches(0.85), cw-Inches(1.45), Inches(1.0), sz=11.5, col=GRY, line_sp=1.15)
t(s, "※ 給与・雇用形態などの具体条件は、お互いの希望をすり合わせながら決めていきます。",
  Inches(0.6), Inches(6.55), Inches(12), Inches(0.4), sz=10.5, col=GRY)

# ════════ SLIDE 9 — 3年後のビジョン ════════
s = sl(); ft(s)
hdr(s, "VISION", "3年後、こうなっていたい", "大きく広げることより、地域で一番信頼されることを先に取りに行く")
vis = [
    ("1年目", "船橋で確かな一歩", "良い施術と信頼で、地域のケアマネ・施設から「あそこに頼めば安心」と言われる拠点に。"),
    ("2年目", "チームとして回る", "施術者が増え、運営が仕組みで回る。立ち上げメンバーが後輩を育てる中核に。"),
    ("3年目", "面で支える存在へ", "船橋を起点に近隣エリアへ。地域の在宅高齢者を面で支える訪問チームへ成長。"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(2.0)
for i, (yr, ti, body) in enumerate(vis):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(3.9), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.85), RED)
    t(s, yr, cx, y0+Inches(0.16), cw, Inches(0.55), sz=24, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.3), y0+Inches(1.1), cw-Inches(0.6), Inches(0.8), sz=17, bold=True, col=REDD, align=PP_ALIGN.CENTER, line_sp=1.1)
    bx(s, cx+Inches(0.6), y0+Inches(1.95), cw-Inches(1.2), Pt(1.2), LINE)
    t(s, body, cx+Inches(0.35), y0+Inches(2.15), cw-Inches(0.7), Inches(1.5), sz=12, col=INK, line_sp=1.25)

# ════════ SLIDE 10 — クロージング ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "まずは、お互いを知るところから。", Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5), sz=16, bold=True, col=RED)
t(s, "良ければ、一緒に第一歩を。", Inches(0.88), Inches(2.0), Inches(11.6), Inches(1.0), sz=40, bold=True, col=INK)
t(s, "今日この場は、合否を決める面接ではありません。\nあなたの経験と思いを聞かせてください。そして、私たちの構想に正直な意見をください。",
  Inches(0.92), Inches(3.15), Inches(11.5), Inches(1.0), sz=15, col=GRY, line_sp=1.35)
# 3つの約束カード
items = [
    ("売り込みません", "あなたを口説き落とす場ではなく、合うかを確かめ合う場にします。"),
    ("正直にお話しします", "良いことも、これからの課題も、包み隠さずお伝えします。"),
    ("一緒に決めます", "働き方も条件も、お互いが納得できる形を対話で決めます。"),
]
for i, (ti, ds) in enumerate(items):
    cx = Inches(0.9) + (Inches(3.95) + Inches(0.24)) * i
    bx(s, cx, Inches(4.5), Inches(3.95), Inches(1.85), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, Inches(4.5), Inches(3.95), Inches(0.06), RED)
    t(s, "✓ " + ti, cx+Inches(0.3), Inches(4.72), Inches(3.4), Inches(0.45), sz=15, bold=True, col=REDD)
    t(s, ds, cx+Inches(0.32), Inches(5.25), Inches(3.4), Inches(1.0), sz=11.5, col=GRY, line_sp=1.2)
t(s, "KHD（株式会社）  ｜  代表  菊池 研太", Inches(0.9), Inches(6.7), Inches(11), Inches(0.4), sz=13, bold=True, col=INK)

prs.save("aurora_fc_vision_v1.pptx")
print("saved aurora_fc_vision_v1.pptx  /  slides:", len(prs.slides._sldIdLst))

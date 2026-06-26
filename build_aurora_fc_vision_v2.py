"""
オーロラ訪問マッサージ 事業立ち上げ構想 ── 面談用「当社方針」説明デッキ v2（5枚・写真版）
用途: 採用面談（石原さん＝ベテラン即戦力）で見せる、KHDの方針共有資料。
立て付け: 口説き／FC加盟確定の事業計画書ではない。大風呂敷も卑屈さも出さず、
          「現場のプロにお任せ × 経営の仕組みは私が引く」対等な座組みを写真で共有する。
中核信条: 相手目線・売り込まない・信頼の対価で稼ぐ（feedback_aite_mokusen）。
デザイン: クリーム白 #F9F6EF × レンガ赤 #AA2E26 / Hiragino Sans / 16:9。
出力: aurora_fc_vision_v2.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Hiragino Sans"
W = Inches(13.33)
H = Inches(7.5)

IMGDIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "aurora_imgs")
P_MASSAGE = os.path.join(IMGDIR, "p_massage.jpg")   # 512x344
P_BIKE    = os.path.join(IMGDIR, "p_bike.jpg")       # 366x197
P_GEAR    = os.path.join(IMGDIR, "p_gear.jpg")       # 266x158
P_TEAM    = os.path.join(IMGDIR, "p_team.jpg")       # 326x186
P_OFFICE  = os.path.join(IMGDIR, "p_office.jpg")     # 306x252

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def photo(slide, path, x, y, w_in, ratio):
    """幅固定・アスペクト維持で写真を配置し、細枠を付ける。ratio=高さ/幅。"""
    h_in = w_in * ratio
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w_in))
    pic.line.color.rgb = WHT
    pic.line.width = Pt(4)
    pic.shadow.inherit = False
    return h_in


def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide, page):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "KHD ｜ 船橋・訪問マッサージ事業 立ち上げ構想", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)
    t(slide, str(page), W-Inches(1.0), H-Inches(0.42), Inches(0.5), Inches(0.32), sz=9, col=GRY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 1. 表紙
# ════════════════════════════════════════════════════════════
s = sl()
# 右：施術風景の大判
photo(s, P_MASSAGE, 7.15, 0.0, 6.18, 344/512)  # 高さ約4.15in（上端から）
bx(s, Inches(0), Inches(0), Inches(7.2), H, BG)  # 左テキスト面の地色（写真とのかぶり防止用・背景同色）
# テキストは上の地色boxの上に置く
t(s, "KHD ｜ 船橋・訪問マッサージ事業", Inches(0.75), Inches(1.3), Inches(6.2), Inches(0.4), sz=14, bold=True, col=RED)
bx(s, Inches(0.78), Inches(1.78), Inches(1.5), Pt(3), RED)
t(s, "在宅に、医療の手を。", Inches(0.72), Inches(2.0), Inches(6.4), Inches(1.0), sz=40, bold=True, col=INK)
t(s, "船橋エリアで、訪問マッサージの新しい拠点を\n立ち上げます。", Inches(0.78), Inches(3.25), Inches(6.2), Inches(1.0), sz=16, col=INK, line_sp=1.25)
bx(s, Inches(0.78), Inches(5.45), Inches(5.9), Inches(1.1), CARD, CARDLN, 1.0)
t(s, "面談用資料", Inches(1.0), Inches(5.62), Inches(3), Inches(0.35), sz=11, bold=True, col=RED)
t(s, "KHD 代表  菊池 研太", Inches(1.0), Inches(5.95), Inches(5), Inches(0.4), sz=15, bold=True, col=INK)
ft(s, 1)

# ════════════════════════════════════════════════════════════
# 2. 何を立ち上げるか（座組み・端的に）
# ════════════════════════════════════════════════════════════
s = sl()
hdr(s, "WHAT WE BUILD", "立ち上げる体制",
    "細かい説明は省きます。要は、実績ある仕組み × 承認済みエリア × 経営は私が回す、の3つです。")
rows = [
    ("仕組み", "オーロラの訪問マッサージFCの仕組みを活用。\n医療保険×在宅のオペレーションは確立済み。"),
    ("エリア", "船橋市はテリトリー承認済み。\n後期高齢者8万人超／3駅圏の需要を狙える。"),
    ("経営", "集客・ケアマネ連携・資金・事務はKHDが担う。\n現場が施術に集中できる土台をつくる。"),
]
y = 2.15
for label, body in rows:
    bx(s, Inches(0.62), Inches(y), Inches(7.2), Inches(1.15), CARD, CARDLN, 1.0)
    bx(s, Inches(0.62), Inches(y), Inches(0.16), Inches(1.15), RED)
    t(s, label, Inches(0.95), Inches(y+0.12), Inches(1.6), Inches(0.9), sz=17, bold=True, col=RED, anchor=MSO_ANCHOR.MIDDLE)
    t(s, body, Inches(2.55), Inches(y+0.1), Inches(5.1), Inches(0.95), sz=12.5, col=INK, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.12)
    y += 1.33
# 右：店舗外観
ph = photo(s, P_OFFICE, 8.25, 2.3, 4.4, 252/306)
t(s, "船橋に構える事業拠点", Inches(8.25), Inches(2.3)+Inches(4.4*252/306)+Inches(0.08),
  Inches(4.4), Inches(0.3), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)
ft(s, 2)

# ════════════════════════════════════════════════════════════
# 3. 役割分担（プロにお任せ × 経営は私が引く）── 距離感の核
# ════════════════════════════════════════════════════════════
s = sl()
hdr(s, "HOW WE TEAM UP", "お任せしたいこと、私が引き受けること",
    "現場の質はプロに敵いません。だから役割を分けて、お互いの強みで長く続く拠点をつくりたい。")
# 左カード：現場（石原さんへ）
bx(s, Inches(0.62), Inches(2.2), Inches(5.85), Inches(3.5), WHT, CARDLN, 1.2)
bx(s, Inches(0.62), Inches(2.2), Inches(5.85), Inches(0.7), RED)
t(s, "現場 ──  あなたにお任せしたい", Inches(0.85), Inches(2.32), Inches(5.5), Inches(0.45), sz=15, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
for i, line in enumerate(["施術の質・臨床判断", "患者さんとの信頼関係", "現場のオペレーション最適化"]):
    yy = 3.15 + i*0.72
    t(s, "●", Inches(0.95), Inches(yy), Inches(0.4), Inches(0.4), sz=12, col=RED)
    t(s, line, Inches(1.35), Inches(yy-0.03), Inches(4.9), Inches(0.5), sz=14, bold=True, col=INK)
t(s, "プロの領域は、信頼してお任せします。", Inches(0.85), Inches(5.25), Inches(5.4), Inches(0.4), sz=11.5, italic=True, col=GRY)
# 右カード：経営（私が引く）
bx(s, Inches(6.85), Inches(2.2), Inches(5.85), Inches(3.5), CARD, CARDLN, 1.2)
bx(s, Inches(6.85), Inches(2.2), Inches(5.85), Inches(0.7), INK)
t(s, "経営 ──  私が引き受ける", Inches(7.08), Inches(2.32), Inches(5.5), Inches(0.45), sz=15, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
for i, line in enumerate(["集客導線・ケアマネ連携の仕組み", "資金・制度・許認可まわり", "事務・バックオフィス全般"]):
    yy = 3.15 + i*0.72
    t(s, "●", Inches(7.18), Inches(yy), Inches(0.4), Inches(0.4), sz=12, col=INK)
    t(s, line, Inches(7.58), Inches(yy-0.03), Inches(4.9), Inches(0.5), sz=14, bold=True, col=INK)
t(s, "テリトリー・資金・仕組みは、私の責任で用意します。", Inches(7.08), Inches(5.25), Inches(5.5), Inches(0.4), sz=11.5, italic=True, col=GRY)
# 下帯：まとめ一文
bx(s, Inches(0.62), Inches(5.95), Inches(12.1), Inches(0.7), REDBG, None)
t(s, "現場のプロ  ×  経営の仕組み  ──  どちらが欠けても続かない。だから組みたい。",
  Inches(0.62), Inches(5.95), Inches(12.1), Inches(0.7), sz=14, bold=True, col=REDD,
  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ft(s, 3)

# ════════════════════════════════════════════════════════════
# 4. 働き方・処遇・大切にしたい考え方
# ════════════════════════════════════════════════════════════
s = sl()
hdr(s, "HOW WE WORK", "働き方と、大切にしたい考え方",
    "条件の細部は対話で決めたい。でも“考え方”だけは最初に揃えておきたいです。")
# 左：仕事観（中核信条）
bx(s, Inches(0.62), Inches(2.2), Inches(7.0), Inches(3.9), CARD, CARDLN, 1.0)
t(s, "事業の芯にしたい考え方", Inches(0.88), Inches(2.4), Inches(6.5), Inches(0.4), sz=13, bold=True, col=RED)
t(s, "売り込まない。相手目線でGIVEし、\n信頼の対価としていただく。",
  Inches(0.88), Inches(2.85), Inches(6.5), Inches(1.1), sz=19, bold=True, col=INK, line_sp=1.15)
for i, line in enumerate([
    "患者さん・ケアマネさんに、まず誠実であること",
    "数を追うより、地域で長く信頼される拠点にすること",
    "働く人が、正当に報われる仕組みにすること",
]):
    yy = 4.35 + i*0.56
    t(s, "—", Inches(0.95), Inches(yy), Inches(0.4), Inches(0.4), sz=13, col=RED)
    t(s, line, Inches(1.35), Inches(yy-0.02), Inches(6.0), Inches(0.45), sz=13, col=INK)
# 右：処遇の考え方（数字は伏せる）＋写真
ph = photo(s, P_BIKE, 8.0, 2.2, 4.7, 197/366)
t(s, "処遇は、対話で決めます", Inches(8.0), Inches(4.85), Inches(4.7), Inches(0.4), sz=13, bold=True, col=RED)
t(s, "経験・役割に正当に報いる前提で、\nお互い納得できる形を一緒に決めたい。\n固定の枠に当てはめるつもりはありません。",
  Inches(8.0), Inches(5.25), Inches(4.7), Inches(1.2), sz=12, col=INK, line_sp=1.2)
ft(s, 4)

# ════════════════════════════════════════════════════════════
# 5. クロージング
# ════════════════════════════════════════════════════════════
s = sl()
photo(s, P_TEAM, 7.0, 0.0, 6.33, 186/326)  # 右上
bx(s, Inches(0), Inches(0), Inches(7.05), H, BG)
t(s, "FROM HERE", Inches(0.75), Inches(1.5), Inches(6), Inches(0.4), sz=13, bold=True, col=RED)
bx(s, Inches(0.78), Inches(1.9), Inches(1.5), Pt(3), RED)
t(s, "まずは、お互いを\n知るところから。", Inches(0.72), Inches(2.15), Inches(6.3), Inches(1.8), sz=34, bold=True, col=INK, line_sp=1.1)
t(s, "今日は口説く場ではありません。\nお互いに見極める場だと思っています。\n気になる点は、何でも遠慮なく聞いてください。",
  Inches(0.78), Inches(4.15), Inches(6.1), Inches(1.4), sz=14, col=INK, line_sp=1.3)
bx(s, Inches(0.78), Inches(5.95), Inches(5.9), Inches(0.8), CARD, CARDLN, 1.0)
t(s, "KHD 代表  菊池 研太", Inches(1.0), Inches(6.12), Inches(5.5), Inches(0.45), sz=15, bold=True, col=INK)
ft(s, 5)

prs.save("aurora_fc_vision_v2.pptx")
print("saved: aurora_fc_vision_v2.pptx  /  slides =", len(prs.slides._sldIdLst))

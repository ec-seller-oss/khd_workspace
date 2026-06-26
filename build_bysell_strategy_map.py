#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""バイセル連携 事業化 戦略設計図スライド（業務一覧スプシとセット）"""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent / "scripts"))
import loan_deck as L
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM, BRICK, INK, GRAY, GREEN, WHITE = L.CREAM, L.BRICK, L.INK, L.GRAY, L.GREEN, L.WHITE
BLUE = RGBColor(0x2B, 0x57, 0x9A)
ORANGE = RGBColor(0xD9, 0x7B, 0x2B)


def chip(s, x, y, w, h, text, fill=WHITE, fg=INK, size=12, bold=True, line=BRICK):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.color.rgb = line; r.line.width = Pt(1.25); r.shadow.inherit = False
    L.box(s, x, y + 0.02, w, h, text, size=size, color=fg, bold=bold,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)


def arrow(s, x, y, txt="→"):
    L.box(s, x, y, 0.5, 0.4, txt, size=22, color=BRICK, bold=True, align=PP_ALIGN.CENTER)


def build():
    prs = L.Presentation(); prs.slide_width = L.EMU_W; prs.slide_height = L.EMU_H

    # 1 表紙
    s = L.add_slide(prs)
    L.bar(s, 0, 2.5, 13.333, 0.10, BRICK)
    L.box(s, 0.9, 1.3, 11, 0.5, "戦略設計図（業務一覧スプシとセット）", size=14, color=GRAY)
    L.box(s, 0.9, 2.7, 11.8, 1.1, "🏗 バイセル資産 × 自社営業 × ツクビト", size=38, color=BRICK, bold=True)
    L.box(s, 0.95, 3.85, 11.8, 0.7, "過去案件を「現金」に変える 5フェーズの全体図", size=19, color=INK)
    L.box(s, 0.95, 6.2, 11.8, 0.6, "制作：テナントアシスト・ウイン 菊池研太 ／ 2026-06-10", size=13, color=GRAY)

    # 2 全体図（5フェーズ）
    s = L.add_slide(prs); L.header(s, 1, "全体図 ― データを現金に変える流れ")
    y = 1.9
    chip(s, 0.4, y, 2.3, 1.15, "①データ資産\n救出(6月末迄)\n使う分だけPDF", fill=RGBColor(0xEC,0xE0,0xDD))
    arrow(s, 2.75, y + 0.35)
    chip(s, 3.3, y, 2.3, 1.15, "②整理\n再アプローチ37\nルーティング90", fill=WHITE)
    arrow(s, 5.65, y + 0.35)
    chip(s, 6.2, y, 2.5, 1.15, "③営業(本体)\n毎朝荷電5件\n業者接触1件", fill=RGBColor(0xF7,0xE9,0xE7), line=BRICK)
    arrow(s, 8.75, y + 0.35)
    chip(s, 9.3, y, 1.7, 1.15, "④出口\n自社2/3\nよくばり1/3", fill=WHITE)
    arrow(s, 11.05, y + 0.35)
    chip(s, 11.6, y, 1.4, 1.15, "⑤現金\n+実績", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN)
    L.box(s, 0.5, 3.5, 12.3, 0.5, "▼ 支える仕組み（できてる）", size=14, color=GREEN, bold=True)
    chip(s, 0.5, 4.0, 3.9, 0.8, "AI査定エンジン(無人)\n反応案件を即査定→買付", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN, size=11)
    chip(s, 4.6, 4.0, 3.9, 0.8, "込山(師匠)との連携\n松戸つよつよ×千葉+AI", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN, size=11)
    chip(s, 8.7, 4.0, 4.1, 0.8, "宅建士+KHD宅建業\n込山6割=皆ハッピー構想", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN, size=11)
    L.box(s, 0.5, 5.2, 12.3, 1.3,
          "■ 戦略の一言：救出は『全部』でなく『今日荷電する分だけ』。③毎日の営業が本体で、他は全部その支援。\n"
          "■ 数字：現実シナリオ(週1件確保)で月次粗利499万 ＞ 必達ライン。ただし前提は楽観＝まず実績1件。",
          size=14, color=INK, ls=1.3)
    L.footer(s)

    # 3 業務一覧（スプシ連携サマリ）
    s = L.add_slide(prs); L.header(s, 2, "業務一覧 ― 何を・いつまでに（詳細はスプシ）")
    L.table(s, 0.5, 1.55, 12.3, [
        ("🔴 今週必須", "荷電前PDF救出(6月末締切)／毎朝荷電5件＋業者接触1件／ツクビト確認(帰属・競業)／込山と座組"),
        ("🟡 今月", "シーズネタ追加救出／別ファイル分析(入手次第)／事業計画の予実"),
        ("🟢 実績1件後", "羽鳥ら都内業者ルート拡販／よくばり1/3運用／込山案件のKHD受託(6割)"),
        ("❄️ 不要・完了", "過去全メール一括救出(不要と判断済)／APチーム救出(済)／在庫表53件(済)"),
    ], col1=0.22, rh=0.95, fs=13)
    L.box(s, 0.5, 5.8, 12.3, 0.8,
          "■ 期限が効くのは①救出だけ（Gmail 6月末閉鎖・あと約20日）。他は『毎日の営業』に従属。",
          size=14, color=BRICK, bold=True)
    L.footer(s)

    # 4 1日の流れ
    s = L.add_slide(prs); L.header(s, 3, "毎日なにをやるか（1日の型）")
    L.table(s, 0.5, 1.55, 12.3, [
        ("5:00-7:30", "調査士（別枠・死守）"),
        ("8:45-9:00", "今日荷電する5案件のメール詳細をPDF救出（①を15分だけ）"),
        ("9:00-12:00", "🔥 荷電5件「今いくらなら？」＋業者接触1件（紹介依頼）"),
        ("13:00-17:00", "対面・内見 ／ 反応案件→AI査定→買付 or 込山・ツクビトへパス"),
        ("18:00-", "家族（死守）"),
        ("夜", "日報：荷電結果をシートに記入＋黙る力5問＋温度感1語"),
    ], col1=0.18, rh=0.72, fs=14)
    L.box(s, 0.5, 6.3, 12.3, 0.6,
          "KPI＝荷電5件/日＋業者接触1件/日 → 月4件確保（現実シナリオ）に乗る",
          size=15, color=GREEN, bold=True)
    L.footer(s)

    # 5 結び
    s = L.add_slide(prs)
    L.bar(s, 0, 2.4, 13.333, 0.08, BRICK)
    L.box(s, 0.9, 2.7, 11.8, 1.5,
          "データは腐る、営業は実る。\n救出は使う分だけ、毎日の荷電が本体。実績1件がすべての扉。",
          size=22, color=INK, ls=1.35)
    L.box(s, 0.95, 5.6, 11.5, 0.7, "制作：テナントアシスト・ウイン 菊池研太 ／ 業務一覧スプシとセット運用", size=13, color=GRAY)

    out = Path.home() / "01_honbu_docs_automation" / "out_screener" / "05_バイセル連携_戦略設計図.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    print("✅", build())

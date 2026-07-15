#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
KHD 車両転売事業マニュアル v1 — 公用車入札の仕入〜売却 一気通貫フロー
不動産査定と同じ「発掘→資格→机上査定→現物→収支→入札→決済→売却」の型に落とし込み、
各STEPのチェックリストを同梱。実例＝岩手中部水道企業団ランクル案件(2026)。
デザインシステム: クリーム白 #F9F6EF × レンガ赤 #AA2E26 ＋ ゴールド差し色（KHD標準）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CREAM=RGBColor(0xF9,0xF6,0xEF); BRICK=RGBColor(0xAA,0x2E,0x26); DBRICK=RGBColor(0x82,0x21,0x1B)
INK=RGBColor(0x2B,0x24,0x22); GRAY=RGBColor(0x6B,0x60,0x5C); MUTE=RGBColor(0x9A,0x90,0x8A)
WHITE=RGBColor(0xFF,0xFF,0xFF); GOLD=RGBColor(0xC3,0x9B,0x4E); LGOLD=RGBColor(0xE7,0xD9,0xB7)
LBRICK=RGBColor(0xEC,0xDA,0xD7); PANEL=RGBColor(0xF3,0xEC,0xEB); SHADOW=RGBColor(0xE2,0xDB,0xD1)
GREEN=RGBColor(0x3F,0x7D,0x4E); LGREEN=RGBColor(0xDD,0xEA,0xDF); BLUE=RGBColor(0x2A,0x6E,0xA0)
FONT='Hiragino Kaku Gothic ProN'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
_pages={'n':0}

def E(v): return Emu(int(v))
def add(v,d): return Emu(int(v)+int(d))
def bg(s,c=CREAM):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,l,t,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def rrect(s,l,t,w,h,fill=None,line=None,lw=1.0):
    return box(s,l,t,w,h,fill,line,lw,MSO_SHAPE.ROUNDED_RECTANGLE)
def card(s,l,t,w,h,fill=WHITE,line=None,lw=1.0,depth=Inches(0.06)):
    rrect(s,add(l,depth),add(t,depth),w,h,fill=SHADOW)
    return rrect(s,l,t,w,h,fill=fill,line=line,lw=lw)
def seg(s,l,t,w,h,color=GOLD,lw=2.0):
    c=s.shapes.add_connector(2,l,t,add(l,w),add(t,h))
    c.line.color.rgb=color; c.line.width=Pt(lw); c.shadow.inherit=False; return c
def txt(s,l,t,w,h,text,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if sp is not None: p.line_spacing=sp
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=FONT
    return tb
def footer(s):
    _pages['n']+=1
    seg(s,Inches(0.7),Inches(7.02),Inches(11.93),0,color=LBRICK,lw=1.0)
    txt(s,Inches(0.7),Inches(7.04),Inches(9),Inches(0.35),'KHD ｜ 車両転売事業マニュアル v1（公用車入札）',size=9,color=MUTE)
    txt(s,Inches(11.4),Inches(7.04),Inches(1.23),Inches(0.35),f'{_pages["n"]:02d}',size=10,color=BRICK,bold=True,align=PP_ALIGN.RIGHT)
def header(s,kicker,title):
    box(s,0,0,Inches(0.22),SH,fill=BRICK)
    txt(s,Inches(0.7),Inches(0.40),Inches(11.8),Inches(0.38),kicker,size=12,color=GOLD,bold=True)
    txt(s,Inches(0.7),Inches(0.70),Inches(12.2),Inches(0.8),title,size=25,color=INK,bold=True)
    seg(s,Inches(0.72),Inches(1.44),Inches(2.0),0,color=GOLD,lw=2.5)
def chip(s,l,t,text,fill=LGOLD,fg=DBRICK,w=Inches(1.9),h=Inches(0.4),size=12):
    rrect(s,l,t,w,h,fill=fill); txt(s,l,t,w,h,text,size=size,color=fg,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def checklist(s,l,t,w,items,size=12.5,gap=0.36,color=INK):
    y=t
    for it in items:
        box(s,l,add(y,Inches(0.04)),Inches(0.17),Inches(0.17),fill=None,line=BRICK,lw=1.3)
        txt(s,add(l,Inches(0.30)),y,add(w,-Inches(0.30)),Inches(0.4),it,size=size,color=color,sp=1.0)
        y=add(y,Inches(gap))
    return y

# ============ 1. COVER ============
s=prs.slides.add_slide(BLANK); bg(s,INK)
box(s,0,0,Inches(0.45),SH,fill=BRICK); box(s,Inches(0.45),0,Inches(0.12),SH,fill=GOLD)
txt(s,Inches(1.1),Inches(1.15),Inches(11),Inches(0.5),'KHD ｜ 新規事業 ｜ 社内マニュアル',size=16,color=GOLD,bold=True)
txt(s,Inches(1.1),Inches(1.95),Inches(11.4),Inches(1.9),'車両転売事業マニュアル v1\n公用車入札 仕入〜売却 一気通貫フロー',size=40,color=WHITE,bold=True,sp=1.1)
txt(s,Inches(1.1),Inches(4.35),Inches(11),Inches(1.6),'不動産査定と同じ「型」で回す：発掘 → 資格 → 机上査定 → 現物 → 収支 → 入札 → 決済 → 売却\n全8STEPチェックリスト同梱 ／ 実例：岩手中部水道企業団 ランドクルーザー案件（2026）',size=15,color=LGOLD,sp=1.3)
txt(s,Inches(1.1),Inches(6.6),Inches(10),Inches(0.4),'KIKUCHIホールディングス株式会社 ｜ 2026-07-15 作成',size=12,color=MUTE)

# ============ 2. 全体フロー ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'OVERVIEW','全体フロー：不動産と同じ8ステップの「型」')
steps=[('S1','案件発掘','官公庁の売払い\n公告を拾う'),('S2','参加資格','地域要件・書類\n＝門前払い回避'),
       ('S3','机上査定','仕様書精読＋\n相見積りで相場'),('S4','現物確認','下回り・始動\n＝内見と同じ'),
       ('S5','収支確定','逆算で入札\n上限を決める'),('S6','入札','一発勝負\n端数を付ける'),
       ('S7','決済・搬出','全額前払い\n期限厳守'),('S8','売却・記録','最高値で売却\n担当者と関係構築')]
x=Inches(0.7); w=Inches(1.44); gapx=Inches(0.145)
for i,(no,name,desc) in enumerate(steps):
    l=add(x,int(add(w,gapx))*i)
    c=card(s,l,Inches(2.0),w,Inches(2.5),fill=WHITE)
    box(s,l,Inches(2.0),w,Inches(0.5),fill=BRICK if i not in (4,5) else DBRICK)
    txt(s,l,Inches(2.0),w,Inches(0.5),no,size=13,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,l,Inches(2.6),w,Inches(0.5),name,size=13.5,color=INK,bold=True,align=PP_ALIGN.CENTER)
    txt(s,l,Inches(3.15),w,Inches(1.3),desc,size=10,color=GRAY,align=PP_ALIGN.CENTER,sp=1.1)
card(s,Inches(0.7),Inches(5.0),Inches(12.0),Inches(1.7),fill=PANEL)
txt(s,Inches(1.0),Inches(5.2),Inches(11.5),Inches(0.4),'不動産との対応関係（同じ筋肉で回せる）',size=14,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(5.65),Inches(11.5),Inches(1.0),
    'レインズ発掘=売払い公告ウォッチ ／ 謄本・重説=物件仕様書・入札公告 ／ 内見・現調=現物公開日 ／ 買付=入札書（ただし一発・撤回不可）\n決済=契約締結・全額前払い ／ 転売先=買取・輸出業者（レインズの代わりに相見積りで最高値を作る）',size=12.5,color=INK,sp=1.25)
footer(s)

# ============ 3. S1 案件発掘 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 1','案件発掘：官公庁の売払い公告を拾う')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'探す場所',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '・市町村/水道企業団/一部事務組合のHP「入札情報」\n・「公用車 売払い 入札」「不用物品 売払い」で検索\n・官公庁オークション（KSI官公庁オークション等）\n・地域要件が緩い案件＝競争が激しい\n　地域限定案件＝参入障壁が高い分、安く取れる\n\n狙い目の型：\n・地方の水道/消防/土木系＝ランクル・ハイエース等\n　海外needsの高い車が出やすい\n・「最低売却価格が明らかに安い」＝役所は相場を\n　取りに行かない。ここが利益の源泉',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '公告日・入札方式（一般競争/条件付/せり）を確認した',
 '最低売却価格と車種で「勝ち筋」があるか10分で仮判定',
 '参加資格（地域要件・法人/個人）を最初に読んだ',
 '締切から逆算したスケジュール表を作った',
 '設計図書（仕様書）・様式類を全てDLして保存した',
 '現物公開日をカレンダー登録（事前連絡の要否も）',
 '過去の同種案件の落札結果を検索した（相場観）',
],size=12)
footer(s)

# ============ 4. S2 参加資格 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 2','参加資格・書類：門前払いを回避する')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'最初に確認する3点',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '① 地域要件：住民登録 or 本店/支店/営業所の所在地\n　→ KHDは花巻支店登記で岩手中部圏域をクリア済み\n　→ 他地域の案件は「支店登記の追加」も戦略になる\n② 欠格要件：破産/更生/税滞納/反社でないこと\n③ 添付書類の有効期限：発行後3ヶ月以内が通例\n　→ 締切から逆算して取得日を決める（早取りは失効リスク）\n\n法人申込の標準セット：\n全部事項証明書／印鑑登録証明書／市町村税納税証明書\n／誓約書（様式）／（代理人なら委任状＋身分証写し）',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '地域要件を自社の登記で満たすことを謄本で確認した',
 '税の未納がないこと（誓約書と矛盾しない）を確認した',
 '添付書類の取得日を締切から逆算して決めた',
 '様式（申込書/誓約書）を正確な商号・住所で下書きした',
 '不明点は質問書で企業団に事前照会した（期限内）',
 '提出は持参 or 簡易書留等の記録が残る方法にした',
 '受付期間内「必着」を確認した（消印有効ではない）',
],size=12)
footer(s)

# ============ 5. S3 机上査定 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 3','机上査定：仕様書精読 × 相見積りで「売値」を作る')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'仕様書から読み取る査定因子',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '・型式→グレードを特定（MT/AT・バン/ワゴンで絞れる）\n・走行距離／初度登録＝国内価値。ランクル等は\n　過走行でも海外価値が残る（50万kmからが本番）\n・車検/自賠責の残＝自走可否・搬出コストに直結\n・特殊装備（赤色灯等）＝撤去・用途変更の手間\n\n相見積りの鉄則：\n・最低3〜5社。国内買取/旧車専門/輸出系/知人ルート\n　と「タイプの違う」先に当てる\n・提示額は相互に伏せる（アンカリング回避）\n・車台番号は概算段階では渡さない（本命1社のみ）',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '型式・年式・走行・装備からグレードを特定した',
 '排ガス規制(NOx・PM法)の適合を確認した※重要',
 '査定先を3〜5社リストアップ（タイプ分散）した',
 '競合化リスクを精査した（地域要件を満たす業者は除外）',
 '概算査定を取り「保守/中央/強気」の3シナリオにした',
 '査定根拠（メール・査定書）を証跡として保存した',
 '想定vs実際の比較表を作り、取得のたび埋めている',
],size=12)
footer(s)

# ============ 6. S4 現物確認 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 4','現物確認：不動産の内見と同じ「唯一の実査機会」')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'当日の段取り',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '・公開日は事前予約制が通例。必ず電話で時間を確保\n・現状渡し＋契約不適合責任なし＝ここで見落とすと\n　全部自分持ち。写真は「多すぎるくらい」撮る\n・査定業者に送る前提で撮る：外装4方向/下回り/\n　エンジンルーム/内装/メーター/タイヤ/装備品\n\n最重要は下回り：\n・降雪地の公用車は融雪剤による錆が最大の減点\n・フレームの貫通錆の有無で輸出価値が激変\n・その場で判断できなければ写真を整備工場・買取業者\n　に即送して所見をもらう',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '事前予約の電話を入れ、訪問時間を確定した',
 '下回り（フレーム・足回り）の錆を撮影・打診した',
 '始動・アイドリング・白煙黒煙・オイル漏れを動画で記録',
 '車台番号と書類（車検証等）の一致を確認した',
 '特殊装備の脱着可否を担当者に直接質問した',
 '付属品（鍵・取説・リサイクル券・タイヤ等）を実見した',
 '写真20枚以上を撮り、当日中に査定先へ送付した',
],size=12)
footer(s)

# ============ 7. S5 収支確定 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 5','収支確定：売値から逆算して入札上限を決める')
card(s,Inches(0.7),Inches(1.8),Inches(7.4),Inches(2.9),fill=WHITE)
txt(s,Inches(1.0),Inches(1.95),Inches(6.8),Inches(0.4),'逆算式（転売前提・車検不要が基本）',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.45),Inches(6.9),Inches(2.1),
    '入札上限 ＝（保守シナリオ売却額 − 目標利益 − 諸費用）÷ 1.10\n\n・契約金額＝入札額×1.10（消費税相当加算）が通例\n・諸費用＝陸送(0〜10万・業者引取無料もある)＋名義変更等(1〜3万)\n・保有するなら＋車検整備(15〜40万)＋排ガス規制の登録可否を確認',size=13,color=INK,sp=1.3)
card(s,Inches(0.7),Inches(4.9),Inches(7.4),Inches(1.8),fill=LGREEN)
txt(s,Inches(1.0),Inches(5.05),Inches(6.8),Inches(0.4),'実例：岩手中部ランクル（旧車王 実査定100〜200万円）',size=13.5,color=GREEN,bold=True)
txt(s,Inches(1.0),Inches(5.5),Inches(6.9),Inches(1.1),
    '保守100万で利益30万狙い → 上限52万 ／ 中央150万なら → 上限98万\n→ 推奨入札 55〜65万円・絶対上限80万円（保守の損益分岐）',size=13,color=INK,sp=1.3)
card(s,Inches(8.4),Inches(1.8),Inches(4.3),Inches(4.9),fill=WHITE)
chip(s,Inches(8.7),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(8.7),Inches(2.6),Inches(3.7),[
 '売却3シナリオ（保守/中央/強気）を実査定で裏付けた',
 '諸費用を業者別の実見積りで確定した',
 '入札締切1週間前までに全見積りを回収した',
 '「絶対上限」を先に決め、書面に残した',
 '端数を付けた入札額にした（同額くじ回避）',
 '資金（契約金額の全額前払い）を確保した',
],size=11.5,gap=0.44)
footer(s)

# ============ 8. S6 入札 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 6','入札：一発勝負のルールを間違えない')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'絶対に落とせない実務ルール',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '・入札回数1回／提出後の書換・撤回不可が通例\n・金額は税抜で記入（契約金額は×1.10）\n・「￥」記入・アラビア数字・訂正不可ペン\n・封筒に件名明記・封印\n・郵送は「必着」。簡易書留/特定記録で記録を残す\n\n無効になる典型：\n記名押印漏れ／金額訂正／指定様式以外／\n資本・人的関係のある複数者の入札（連合とみなし）',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '入札書は指定様式・ボールペン・訂正なしで記入した',
 '金額の頭に￥、桁ズレがないか2人ダブルチェックした',
 '記名・押印（法人は代表者印）を確認した',
 '封筒に件名を明記し封印した',
 '提出期限の前日までに到達する手段で発送した',
 '開札日・結果公表の確認方法を控えた',
 '落札できなかった場合の次案件リストを用意した',
],size=12)
footer(s)

# ============ 9. S7 決済・搬出 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 7','落札後：契約・全額前払い・搬出を期限内に')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'期限が3つ走る（1つでも落とすと決定取消も）',size=14.5,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '① 契約締結期限（例：落札から1週間）\n② 代金全額前払い（納入通知書で契約締結までに）\n③ 搬出期限（例：月末まで・費用は落札者負担）\n\n搬出の実務：\n・車検切れ＝公道自走不可 → 積載車を事前手配\n　（仮ナンバー自走は老朽車ではリスク高・非推奨）\n・買取業者への直送も検討（引取無料の業者なら\n　陸送費が丸ごと浮く＝実質利益+5〜10万）\n・名義変更登録は速やかに（自動車税・責任の切替）',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '契約締結・支払・搬出の3期限をカレンダー登録した',
 '納入通知書での支払方法・期日を確認した',
 '陸送 or 業者引取を確定し、搬出日を予約した',
 '売却先と引取場所・日程を握った（直送できれば最良）',
 '名義変更（または抹消）の段取りを決めた',
 '特殊装備の撤去有無・タイミングを確定した',
 '実費の領収書を全て保存した（比較表に反映）',
],size=12)
footer(s)

# ============ 10. S8 売却・記録 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'STEP 8','売却・関係構築・記録：次の案件につなげる')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'売却と事業化',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '・相見積りの最高値＋対応の質で売却先を決定\n・取引完了時に「今後も官公庁売払いを継続的に扱う」\n　と伝え、担当者の直通連絡先を確保する\n・担当者は名刺DB/顧客マスターに登録（次回は\n　電話1本で査定が回る状態を作る）\n\n事業化の法規：\n・古物商許可＝反復継続で転売するなら必須\n　（無許可営業は3年以下の懲役/100万円以下の罰金）\n・2台目に着手する前に取得完了（申請〜許可1〜2ヶ月）',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=WHITE)
chip(s,Inches(7.3),Inches(2.0),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
checklist(s,Inches(7.3),Inches(2.6),Inches(5.1),[
 '本番査定（車台番号開示）は最有力1社に絞った',
 '売買契約書・振込記録を保存した',
 '担当者名・直通連絡先を名刺DBに登録した',
 '想定vs実際の比較表を完成させ、差異を分析した',
 '案件全体の学び（次回の入札精度向上）をnotesに記録',
 '古物商許可の取得状況を確認した（2台目の前提）',
 '次の売払い案件のウォッチを再開した',
],size=12)
footer(s)

# ============ 11. リスク・法規 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'RISK & LAW','リスクと法規：先に知らないと即死する4つ')
items=[('NOx・PM法','旧年式ディーゼルは首都圏・大阪・愛知等の対策地域で登録不可。\n「東京で保有」が物理的に不可能な車がある。転売(輸出)なら無関係。','最重要'),
('現状渡し・責任なし','引渡し後の故障・瑕疵は全て自己責任。現物確認が唯一の防衛線。\n下回り腐食は売値を数十万単位で変える。','高'),
('一発入札','書換・撤回・再入札なし。相場を外すと高値掴みか未落札。\n「絶対上限」を先に書面化してから入札書を書く。','高'),
('古物商許可','反復継続の転売は許可必須。1台の単発は直ちに違法ではないが、\n事業化宣言済みなら2台目までに取得完了が必須。','中')]
y=Inches(1.8)
for name,desc,tag in items:
    c=card(s,Inches(0.7),y,Inches(12.0),Inches(1.14),fill=WHITE)
    box(s,Inches(0.7),y,Inches(0.14),Inches(1.14),fill=BRICK)
    txt(s,Inches(1.05),add(y,Inches(0.10)),Inches(2.6),Inches(0.9),name,size=14.5,color=DBRICK,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.8),add(y,Inches(0.09)),Inches(7.6),Inches(1.0),desc,size=11.5,color=INK,sp=1.15)
    chip(s,Inches(11.55),add(y,Inches(0.36)),tag,fill=LBRICK,fg=DBRICK,w=Inches(0.95),h=Inches(0.4),size=11)
    y=add(y,Inches(1.30))
footer(s)

# ============ 12. 巻末：実例サマリ ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'CASE STUDY','実例：岩手中部水道企業団 ランドクルーザー案件（2026）')
card(s,Inches(0.7),Inches(1.8),Inches(6.0),Inches(4.9),fill=WHITE)
txt(s,Inches(1.0),Inches(2.0),Inches(5.4),Inches(0.4),'案件データ',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.5),Inches(5.5),Inches(3.9),
    '車両：ランクル100 バンVX ディーゼル5MT(KG-HDJ101K)\n平成10年式／28万km／車検切れ／公共応急作業車\n最低売却価格：15万円\n実査定：旧車王 概算100〜200万円\n推奨入札：55〜65万円（上限80万円）\n\nスケジュール：8/20現物確認 → 8/21書類締切 →\n9/7入札 → 9/8開札 → 9/15契約・支払 → 9/30搬出',size=12.5,color=INK,sp=1.25)
card(s,Inches(7.0),Inches(1.8),Inches(5.7),Inches(4.9),fill=PANEL)
txt(s,Inches(7.3),Inches(2.0),Inches(5.1),Inches(0.4),'この案件で確立した「型」',size=15,color=DBRICK,bold=True)
txt(s,Inches(7.3),Inches(2.5),Inches(5.1),Inches(3.9),
    '・地域要件は支店登記でクリアできる（参入障壁を\n　自分だけ越える＝安く仕入れる構造）\n・役所の最低売却価格は市場価格と無関係。\n　仕様書の精読とプロ査定で「本当の値段」を掴む\n・買取業者は地域要件を満たせない＝情報を渡しても\n　入札で競合しない。ただし買い叩きには注意\n・全記録は .company/secretary/notes/ に集約。\n　想定vs実際の差異が次回案件の精度になる',size=12.5,color=INK,sp=1.25)
footer(s)

prs.save('/home/user/khd_workspace/KHD_車両転売_仕入売却マニュアル_v1.pptx')
print('saved')

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
KHD 車両転売事業マニュアル v2 — 目利き図解＋実車写真 完全版
------------------------------------------------------------
v1からの強化:
- 目利き（インスペクション）章を新設・最重要扱い（不動産の内見と同格）
- 車両側面/下回りの模式図（KHD自作）に番号付きチェックポイント
- 岩手中部ランクル実車写真9枚（物件仕様書より）を注記付きで掲載
- 現物確認当日の撮影構図テンプレ・持ち物リスト・全チェックリスト総覧を追加
デザイン: 全社標準 クリーム白 #F9F6EF × レンガ赤 #AA2E26 ＋ ゴールド差し色
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CREAM=RGBColor(0xF9,0xF6,0xEF); BRICK=RGBColor(0xAA,0x2E,0x26); DBRICK=RGBColor(0x82,0x21,0x1B)
INK=RGBColor(0x2B,0x24,0x22); GRAY=RGBColor(0x6B,0x60,0x5C); MUTE=RGBColor(0x9A,0x90,0x8A)
WHITE=RGBColor(0xFF,0xFF,0xFF); GOLD=RGBColor(0xC3,0x9B,0x4E); LGOLD=RGBColor(0xE7,0xD9,0xB7)
LBRICK=RGBColor(0xEC,0xDA,0xD7); PANEL=RGBColor(0xF3,0xEC,0xEB); SHADOW=RGBColor(0xE2,0xDB,0xD1)
GREEN=RGBColor(0x3F,0x7D,0x4E); LGREEN=RGBColor(0xDD,0xEA,0xDF)
STEEL=RGBColor(0xC9,0xD2,0xD8); DSTEEL=RGBColor(0x8F,0x9D,0xA6)
FONT='Hiragino Kaku Gothic ProN'
ASSET='/home/user/khd_workspace/assets_landcruiser_148'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
_pages={'n':0}

def add(v,d): return Emu(int(v)+int(d))
def bg(s,c=CREAM):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,l,t,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def rrect(s,l,t,w,h,fill=None,line=None,lw=1.0):
    return box(s,l,t,w,h,fill,line,lw,MSO_SHAPE.ROUNDED_RECTANGLE)
def oval(s,l,t,w,h,fill=None,line=None,lw=1.0):
    return box(s,l,t,w,h,fill,line,lw,MSO_SHAPE.OVAL)
def card(s,l,t,w,h,fill=WHITE,line=None,lw=1.0,depth=Inches(0.06)):
    rrect(s,add(l,depth),add(t,depth),w,h,fill=SHADOW)
    return rrect(s,l,t,w,h,fill=fill,line=line,lw=lw)
def seg(s,l,t,w,h,color=GOLD,lw=2.0):
    c=s.shapes.add_connector(2,l,t,add(l,w),add(t,h))
    c.line.color.rgb=color; c.line.width=Pt(lw); c.shadow.inherit=False; return c
def txt(s,l,t,w,h,text,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if sp is not None: p.line_spacing=sp
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=FONT
    return tb
def footer(s):
    _pages['n']+=1
    seg(s,Inches(0.7),Inches(7.02),Inches(11.93),0,color=LBRICK,lw=1.0)
    txt(s,Inches(0.7),Inches(7.04),Inches(9),Inches(0.35),'KHD ｜ 車両転売事業マニュアル v2（公用車入札・目利き完全版）',size=9,color=MUTE)
    txt(s,Inches(11.4),Inches(7.04),Inches(1.23),Inches(0.35),f'{_pages["n"]:02d}',size=10,color=BRICK,bold=True,align=PP_ALIGN.RIGHT)
def header(s,kicker,title):
    box(s,0,0,Inches(0.22),SH,fill=BRICK)
    txt(s,Inches(0.7),Inches(0.40),Inches(11.8),Inches(0.38),kicker,size=12,color=GOLD,bold=True)
    txt(s,Inches(0.7),Inches(0.70),Inches(12.2),Inches(0.8),title,size=24,color=INK,bold=True)
    seg(s,Inches(0.72),Inches(1.42),Inches(2.0),0,color=GOLD,lw=2.5)
def chip(s,l,t,text,fill=LGOLD,fg=DBRICK,w=Inches(1.9),h=Inches(0.4),size=12):
    rrect(s,l,t,w,h,fill=fill); txt(s,l,t,w,h,text,size=size,color=fg,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def checklist(s,l,t,w,items,size=12.5,gap=0.36,color=INK):
    y=t
    for it in items:
        box(s,l,add(y,Inches(0.04)),Inches(0.17),Inches(0.17),fill=None,line=BRICK,lw=1.3)
        txt(s,add(l,Inches(0.30)),y,add(w,-Inches(0.30)),Inches(0.4),it,size=size,color=color,sp=1.0)
        y=add(y,Inches(gap))
    return y
def numdot(s,l,t,n,d=Inches(0.34),fill=BRICK,fg=WHITE,size=12):
    oval(s,l,t,d,d,fill=fill)
    txt(s,l,add(t,-Inches(0.015)),d,d,str(n),size=size,color=fg,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def pic(s,path,l,t,w=None,h=None,label=None,note=None,notew=None):
    p=s.shapes.add_picture(os.path.join(ASSET,path),l,t,width=w,height=h)
    if label:
        rrect(s,l,add(t,-Inches(0.02)),Inches(2.1),Inches(0.36),fill=BRICK)
        txt(s,l,add(t,-Inches(0.02)),Inches(2.1),Inches(0.36),label,size=11,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note:
        nw=notew or p.width
        txt(s,l,add(add(t,p.height),Inches(0.03)),nw,Inches(0.6),note,size=10.5,color=GRAY,sp=1.05)
    return p

# ============ 1. COVER ============
s=prs.slides.add_slide(BLANK); bg(s,INK)
box(s,0,0,Inches(0.45),SH,fill=BRICK); box(s,Inches(0.45),0,Inches(0.12),SH,fill=GOLD)
txt(s,Inches(1.1),Inches(1.0),Inches(11),Inches(0.5),'KHD ｜ 新規事業 ｜ 社内マニュアル（全社標準デザイン）',size=15,color=GOLD,bold=True)
txt(s,Inches(1.1),Inches(1.7),Inches(11.6),Inches(1.9),'車両転売事業マニュアル v2\n公用車入札 仕入〜売却 ＋ 目利き完全版',size=38,color=WHITE,bold=True,sp=1.1)
txt(s,Inches(1.1),Inches(4.0),Inches(11.2),Inches(1.9),'不動産と同じ8STEPの型 ＋ 最重要の「目利き」を図解と実車写真で徹底\n目利き章（模式図×12ポイント／実車写真9枚の見方）／ 8STEPチェックリスト ／ 撮影構図テンプレ ／ 総覧1枚\n実例：岩手中部水道企業団 ランドクルーザー100 バンVX（2026・公告148号）',size=14,color=LGOLD,sp=1.35)
txt(s,Inches(1.1),Inches(6.55),Inches(10),Inches(0.4),'KIKUCHIホールディングス株式会社 ｜ 2026-07-15 v2',size=12,color=MUTE)

# ============ 2. 全体フロー ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'OVERVIEW','全体フロー：不動産と同じ8ステップ ＋ 目利きが背骨')
steps=[('S1','案件発掘','公告を拾う'),('S2','参加資格','門前払い回避'),('S3','机上査定','相場を作る'),
       ('S4','現物確認','目利き＝内見'),('S5','収支確定','逆算で上限'),('S6','入札','一発勝負'),
       ('S7','決済・搬出','期限3本'),('S8','売却・記録','最高値＋関係')]
x=Inches(0.7); w=Inches(1.44); gapx=Inches(0.145)
for i,(no,name,desc) in enumerate(steps):
    l=add(x,int(add(w,gapx))*i)
    card(s,l,Inches(1.9),w,Inches(2.1),fill=WHITE)
    box(s,l,Inches(1.9),w,Inches(0.5),fill=DBRICK if i==3 else BRICK)
    txt(s,l,Inches(1.9),w,Inches(0.5),no,size=13,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,l,Inches(2.5),w,Inches(0.5),name,size=13,color=INK,bold=True,align=PP_ALIGN.CENTER)
    txt(s,l,Inches(3.05),w,Inches(0.9),desc,size=10,color=GRAY,align=PP_ALIGN.CENTER,sp=1.1)
card(s,Inches(0.7),Inches(4.35),Inches(12.0),Inches(1.1),fill=LGOLD)
txt(s,Inches(1.0),Inches(4.5),Inches(11.4),Inches(0.85),'★ 目利き（S4）が全STEPの背骨：不動産で「内見・現調しない買付はない」のと同じ。\n机上査定(S3)は仮説、現物確認(S4)が検証、収支(S5)は検証済みの数字だけで組む。',size=13.5,color=DBRICK,bold=True,sp=1.25)
card(s,Inches(0.7),Inches(5.65),Inches(12.0),Inches(1.15),fill=PANEL)
txt(s,Inches(1.0),Inches(5.78),Inches(11.5),Inches(0.95),'不動産との対応：レインズ発掘=公告ウォッチ ／ 謄本・重説=仕様書・公告 ／ 内見・現調=現物公開日 ／ 買付=入札書（一発・撤回不可）\n決済=契約・全額前払い ／ 転売先=買取・輸出業者（相見積りで最高値を作る）',size=12,color=INK,sp=1.25)
footer(s)

# ============ 3. 目利き全体マップ（側面模式図） ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'INSPECTION MAP ｜ 最重要','車両目利きマップ：見るべき12ポイント（側面模式図）')
# --- car schematic (side view) ---
CX=Inches(1.4); CY=Inches(2.7); # base
card(s,Inches(0.7),Inches(1.7),Inches(8.2),Inches(3.6),fill=WHITE)
# body
rrect(s,Inches(1.2),Inches(3.15),Inches(6.9),Inches(1.15),fill=STEEL,line=DSTEEL,lw=1.2)         # lower body
rrect(s,Inches(2.3),Inches(2.45),Inches(4.2),Inches(0.95),fill=STEEL,line=DSTEEL,lw=1.2)         # cabin
box(s,Inches(2.55),Inches(2.6),Inches(1.6),Inches(0.6),fill=WHITE,line=DSTEEL,lw=1.0)            # windows
box(s,Inches(4.35),Inches(2.6),Inches(1.85),Inches(0.6),fill=WHITE,line=DSTEEL,lw=1.0)
box(s,Inches(2.5),Inches(2.28),Inches(3.8),Inches(0.16),fill=DSTEEL)                              # roof rack
oval(s,Inches(5.6),Inches(2.02),Inches(0.3),Inches(0.3),fill=BRICK)                               # beacon
oval(s,Inches(2.05),Inches(3.95),Inches(1.05),Inches(1.05),fill=INK)                              # wheels
oval(s,Inches(6.35),Inches(3.95),Inches(1.05),Inches(1.05),fill=INK)
oval(s,Inches(2.33),Inches(4.23),Inches(0.5),Inches(0.5),fill=STEEL)
oval(s,Inches(6.63),Inches(4.23),Inches(0.5),Inches(0.5),fill=STEEL)
seg(s,Inches(1.2),Inches(4.62),Inches(6.9),0,color=DSTEEL,lw=1.5)                                 # ground line hint
# frame rail highlight (under body)
box(s,Inches(1.5),Inches(4.42),Inches(6.3),Inches(0.14),fill=LBRICK)
# numbered points on schematic
pts=[(1,Inches(3.9),Inches(4.55)),(2,Inches(2.25),Inches(4.05)),(3,Inches(1.35),Inches(3.35)),
     (4,Inches(4.9),Inches(2.62)),(5,Inches(5.55),Inches(1.95)),(6,Inches(7.5),Inches(3.1)),
     (7,Inches(4.0),Inches(3.5)),(8,Inches(3.1),Inches(2.9))]
for n,l,t in pts: numdot(s,l,t,n)
txt(s,Inches(1.2),Inches(5.02),Inches(7.4),Inches(0.3),'※KHD自作の模式図（車両を横から見た図）。番号=右の目利きポイントに対応。',size=10,color=MUTE)
# --- right: 12 point list ---
card(s,Inches(9.15),Inches(1.7),Inches(3.55),Inches(5.0),fill=WHITE)
txt(s,Inches(9.4),Inches(1.85),Inches(3.1),Inches(0.4),'12の目利きポイント',size=14,color=DBRICK,bold=True)
pts12=['① 下回り・フレーム錆（最重要）','② タイヤ・足回り・ブッシュ','③ ボディ下部パネルの錆・膨れ','④ ガラス・モール・ドア開閉','⑤ 特殊装備（赤色灯・ラック）','⑥ エンジンルーム（漏れ・煙）','⑦ 室内・シート・floor錆','⑧ メーター（ODOと書類一致）','⑨ シフト・クラッチ（MT）','⑩ 付属品（鍵・タイヤ・書類）','⑪ 車台番号打刻と書類一致','⑫ 排気（白煙/黒煙・異音）']
y=Inches(2.3)
for p in pts12:
    txt(s,Inches(9.4),y,Inches(3.2),Inches(0.35),p,size=11.5,color=INK)
    y=add(y,Inches(0.355))
card(s,Inches(0.7),Inches(5.5),Inches(8.2),Inches(1.2),fill=LGOLD)
txt(s,Inches(1.0),Inches(5.62),Inches(7.7),Inches(1.0),'配点イメージ：①下回り40点 ／ ⑥機関20点 ／ ③外装15点 ／ ⑧⑪整合10点 ／ その他15点。\n①が崩れたら（貫通錆）他が満点でも降りる。不動産の「構造・傾き」と同じ扱い。',size=12.5,color=DBRICK,bold=True,sp=1.25)
footer(s)

# ============ 4. 目利き詳細：下回り（模式図） ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'INSPECTION 1/4 ｜ 配点40','下回り・フレーム：ここで買値が数十万動く（下面模式図）')
card(s,Inches(0.7),Inches(1.7),Inches(6.4),Inches(4.9),fill=WHITE)
# under-body schematic (top-down)
fx=Inches(1.55); fy=Inches(2.2); fw=Inches(4.6); fh=Inches(3.6)
box(s,fx,fy,Inches(0.3),fh,fill=STEEL,line=DSTEEL)                          # left rail
box(s,add(fx,Inches(4.3)),fy,Inches(0.3),fh,fill=STEEL,line=DSTEEL)         # right rail
for i in range(5):                                                           # cross members
    box(s,add(fx,Inches(0.3)),add(fy,Inches(0.18+0.8*i)),Inches(4.0),Inches(0.22),fill=STEEL,line=DSTEEL)
oval(s,add(fx,-Inches(0.42)),add(fy,Inches(0.25)),Inches(0.42),Inches(0.95),fill=INK)  # wheels
oval(s,add(fx,Inches(4.6)),add(fy,Inches(0.25)),Inches(0.42),Inches(0.95),fill=INK)
oval(s,add(fx,-Inches(0.42)),add(fy,Inches(2.45)),Inches(0.42),Inches(0.95),fill=INK)
oval(s,add(fx,Inches(4.6)),add(fy,Inches(2.45)),Inches(0.42),Inches(0.95),fill=INK)
# hot zones
box(s,fx,add(fy,Inches(2.6)),Inches(0.3),Inches(1.0),fill=BRICK)             # rear rail hot
box(s,add(fx,Inches(4.3)),add(fy,Inches(2.6)),Inches(0.3),Inches(1.0),fill=BRICK)
box(s,add(fx,Inches(0.3)),add(fy,Inches(3.38)),Inches(4.0),Inches(0.22),fill=BRICK)  # rear cross
numdot(s,add(fx,Inches(1.9)),add(fy,Inches(3.28)),'A',fill=DBRICK)
numdot(s,add(fx,-Inches(0.15)),add(fy,Inches(1.5)),'B',fill=GOLD,fg=INK)
numdot(s,add(fx,Inches(1.9)),add(fy,Inches(0.1)),'C',fill=GOLD,fg=INK)
txt(s,Inches(1.0),Inches(6.05),Inches(5.8),Inches(0.5),'上から見たラダーフレーム模式図。■赤=融雪剤地域で錆が集中する「後部」ホットゾーン',size=10.5,color=MUTE,sp=1.1)
card(s,Inches(7.4),Inches(1.7),Inches(5.3),Inches(4.9),fill=WHITE)
txt(s,Inches(7.7),Inches(1.9),Inches(4.8),Inches(0.4),'見方と判定基準',size=14.5,color=DBRICK,bold=True)
txt(s,Inches(7.7),Inches(2.4),Inches(4.8),Inches(3.0),
    'A 後部フレーム・リアクロスメンバー（最頻発）\n　→ドライバーの柄などで軽く叩く。「ボコッ」と\n　　鈍い音・ボロボロ剥がれる＝内部腐食進行\nB サイドレール全長（下に潜って目視）\n　→表面サビ(茶色)はOK。「膨れ・層状剥離」はNG\nC 前部クロス・足回り取付部\n　→ブッシュ/ボルト周りの固着と亀裂\n\n判定：表面サビのみ=減点小（輸出値ほぼ維持）\n　　　膨れ・剥離=要減額交渉材料\n　　　貫通穴=入札額を大幅に下げるか撤退',size=12,color=INK,sp=1.2)
chip(s,Inches(7.7),Inches(5.75),'岩手=融雪剤地域。ここが本案件最大の変数',fill=LBRICK,w=Inches(4.7),size=11.5)
footer(s)

# ============ 5. 目利き詳細：機関（実写真） ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'INSPECTION 2/4 ｜ 配点20','機関（エンジン・排気）：実車写真で見る');
pic(s,'engine.jpg',Inches(0.7),Inches(1.85),w=Inches(6.2),label='実車：エンジンルーム')
txt(s,Inches(0.7),Inches(6.15),Inches(6.2),Inches(0.7),'1HD-FTE 4.2Lディーゼルターボ。ツインバッテリー（寒冷地/ディーゼル標準）。\n写真上：大きなオイル飛散・冷却水漏れ痕は見えず、年式比で良好な印象。',size=11,color=GRAY,sp=1.15)
card(s,Inches(7.2),Inches(1.85),Inches(5.5),Inches(4.8),fill=WHITE)
txt(s,Inches(7.5),Inches(2.05),Inches(5.0),Inches(0.4),'現地で確認する7点（写真では分からない）',size=13.5,color=DBRICK,bold=True)
checklist(s,Inches(7.5),Inches(2.55),Inches(4.9),[
 '冷間始動の一発始動性（ディーゼルの健康診断）',
 'アイドリング安定・異音（カラカラ/ガラガラ）',
 '排気色：白煙(オイル下がり)/黒煙(燃調・詰まり)',
 'オイルにじみ：ヘッドカバー・タービン周り',
 '冷却水の色（サビ色=管理不良のサイン）',
 'バッテリー端子の腐食・補機ベルトの亀裂',
 'ターボの効き（可能なら短距離の構内試走）',
],size=11.5,gap=0.42)
chip(s,Inches(7.5),Inches(5.9),'28万km級1HDは「壊れる」より「漏れ・煙」を見る',fill=LGOLD,w=Inches(4.9),size=11)
footer(s)

# ============ 6. 目利き詳細：内装・整合（実写真2枚） ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'INSPECTION 3/4 ｜ 配点10+α','内装・書類整合：メーターと書類の一致は「登記と現況の照合」')
pic(s,'meter_odo.jpg',Inches(0.7),Inches(1.85),w=Inches(5.6),label='実車：メーター')
txt(s,Inches(0.7),Inches(5.75),Inches(5.6),Inches(1.0),'★ODO実写 280,113km ＝ 物件仕様書の記載(280,113km/R8.7.9時点)と完全一致。\nメーター整合OKの現物証拠。タコメーター6,000スケール＝ディーゼルの証明。\n警告灯はエンジン停止状態の点灯（油圧・充電）で正常。',size=11,color=GRAY,sp=1.2)
pic(s,'cockpit_mt.jpg',Inches(6.7),Inches(1.85),w=Inches(5.6),label='実車：運転席')
txt(s,Inches(6.7),Inches(5.75),Inches(5.6),Inches(1.0),'★5速MTシフト＋トランスファーレバーを実写確認 ＝ バンVXディーゼル5MTの現物証拠。\n業務無線マイク残置（撤去対象）。シート・内装は業務車として並〜良好。\n現地では：クラッチの繋がり位置（奥すぎ=摩耗）、シートレール下の床錆も確認。',size=11,color=GRAY,sp=1.2)
footer(s)

# ============ 7. 目利き詳細：装備・付属品（実写真2枚） ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'INSPECTION 4/4 ｜ 配点15','装備・付属品：足す価値と引く手間を仕分ける')
pic(s,'cluster.jpg',Inches(0.7),Inches(1.85),w=Inches(5.6),label='実車：センター')
txt(s,Inches(0.7),Inches(5.75),Inches(5.6),Inches(1.0),'業務無線・広報アンプ類が残置（「さぎょうだん18」ラベル＝作業団車両）。\n→引く手間：無線・広報設備・赤色灯は転売前に撤去が基本（企業団に外して引渡し\n可能かを質問書で確認済みの論点）。エアコン・オーディオは作動確認済(仕様書)。',size=11,color=GRAY,sp=1.2)
pic(s,'trunk_tires.jpg',Inches(6.7),Inches(1.85),w=Inches(5.6),label='実車：荷室')
txt(s,Inches(6.7),Inches(5.75),Inches(5.6),Inches(1.0),'★スタッドレス4本（2021年12月購入・比較的新しい）が付属 ＝ 足す価値（+2〜5万円相当）。\nテールゲート開閉・ダンパー保持OK。荷室フロアの錆・水染みを現地で確認。\n付属品：鍵3本・車検証・取説・リサイクル券（現地で全点実物照合）。',size=11,color=GRAY,sp=1.2)
footer(s)

# ============ 8. 写真評価まとめ：外観2枚＋現時点の所見 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'PHOTO REVIEW','実車写真の現時点評価：写真で分かったこと／8/20で埋めること')
pic(s,'side_left.jpg',Inches(0.7),Inches(1.85),w=Inches(5.0),label='実車：左側面')
pic(s,'rear_vx.jpg',Inches(0.7),Inches(4.6),w=Inches(3.6),label='実車：後方')
txt(s,Inches(4.45),Inches(4.65),Inches(1.9),Inches(2.0),'← VXバッジ\n実写確認。\nグレード確定\nの現物証拠',size=11,color=DBRICK,bold=True,sp=1.2)
card(s,Inches(6.2),Inches(1.85),Inches(6.5),Inches(4.85),fill=WHITE)
txt(s,Inches(6.5),Inches(2.05),Inches(6.0),Inches(0.4),'写真9枚から言えること（2026-07-15時点）',size=14,color=DBRICK,bold=True)
txt(s,Inches(6.5),Inches(2.55),Inches(5.9),Inches(2.3),
    '◎ 外装：ツヤあり・大きな凹み/事故痕は視認できず。年式比で良好\n◎ グレード：後方VXバッジ・5MT・貨物登録 → バンVX確定\n◎ 整合：ODO実写=仕様書=280,113km一致。メーター信頼できる\n◎ 装備：ルーフラック・リアラダー・赤色灯・無線（撤去要否は判断）\n○ 機関：エンジンルーム表層はきれい。ただし写真では音・煙・漏れ不明\n△ 内装：業務使用感あり（減点小）',size=12,color=INK,sp=1.3)
box(s,Inches(6.5),Inches(5.0),Inches(5.9),Inches(0.02),fill=LBRICK)
txt(s,Inches(6.5),Inches(5.12),Inches(5.9),Inches(1.4),
    '✖ 写真に「下回り」が1枚もない ＝ 最大配点(40点)が未評価のまま。\n企業団の公開写真は上物だけ。だから8/20の現物確認が全て。\n下回りが表面サビ止まりなら、旧車王レンジ(100〜200万)の中央〜上限が狙える。',size=12.5,color=DBRICK,bold=True,sp=1.3)
footer(s)

# ============ 8b. 外観4方向（残り写真） ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'PHOTO REVIEW 2/2','外観の記録：前方・右側面（公式写真）と見るポイント')
pic(s,'front.jpg',Inches(0.7),Inches(1.85),w=Inches(5.6),label='実車：前方')
txt(s,Inches(0.7),Inches(5.75),Inches(5.6),Inches(1.0),'岩手800「さ・134」＝8ナンバー特種登録の現物確認。赤色灯・ルーフラック・フォグ付き。\nバンパー・グリルに大きな損傷なし。ヘッドライトのくもりは年式なり（研磨で改善可）。\n現地では：ボンネット先端・ルーフ前縁の飛び石サビ、ガラスの飛び石傷を確認。',size=11,color=GRAY,sp=1.2)
pic(s,'side_right.jpg',Inches(6.7),Inches(1.85),w=Inches(5.6),label='実車：右側面')
txt(s,Inches(6.7),Inches(5.75),Inches(5.6),Inches(1.0),'右側面もライン通り・パネル波打ちなし（写真上）。サイドステップ下端は錆やすい部位。\n現地では：ドア下端・リアフェンダーアーチ内側を指で触って膨れを確認。\nミラー・モール・ドアハンドルの欠品/割れも減点対象なので個別撮影する。',size=11,color=GRAY,sp=1.2)
footer(s)

# ============ 9-16. S1〜S8 ============
def step_slide(kicker,title,left_title,left_body,check_items):
    s=prs.slides.add_slide(BLANK); bg(s)
    header(s,kicker,title)
    card(s,Inches(0.7),Inches(1.7),Inches(6.0),Inches(5.0),fill=WHITE)
    txt(s,Inches(1.0),Inches(1.9),Inches(5.4),Inches(0.4),left_title,size=15,color=DBRICK,bold=True)
    txt(s,Inches(1.0),Inches(2.4),Inches(5.5),Inches(4.1),left_body,size=12,color=INK,sp=1.25)
    card(s,Inches(7.0),Inches(1.7),Inches(5.7),Inches(5.0),fill=WHITE)
    chip(s,Inches(7.3),Inches(1.9),'☑ チェックリスト',fill=LGOLD,w=Inches(2.4))
    checklist(s,Inches(7.3),Inches(2.5),Inches(5.1),check_items,size=12,gap=0.40)
    footer(s); return s

step_slide('STEP 1','案件発掘：官公庁の売払い公告を拾う','探す場所と狙い目',
 '・市町村/水道企業団/一部事務組合HPの「入札情報」\n・「公用車 売払い 入札」「不用物品 売払い」で定期検索\n・官公庁オークション（KSI等）も併読\n\n狙い目の型：\n・地方の水道/消防/土木系＝ランクル/ハイエース等\n　海外需要の高い車が出やすい\n・地域限定案件＝参入障壁が高い分、安く取れる\n・最低売却価格は市場価格と無関係（役所は相場を\n　取りに行かない）。ここが利益の源泉',
 ['公告日・入札方式（一般競争/条件付/せり）を確認した',
  '最低売却価格と車種で「勝ち筋」を10分で仮判定した',
  '参加資格（地域要件・法人/個人）を最初に読んだ',
  '締切から逆算したスケジュール表を作った',
  '設計図書・様式類を全てDLして保存した',
  '現物公開日をカレンダー登録した（事前連絡要否も）',
  '過去の同種案件の落札結果を検索した（相場観）'])

step_slide('STEP 2','参加資格・書類：門前払いを回避する','最初に確認する3点',
 '① 地域要件：住民登録 or 本店/支店/営業所の所在地\n　→ KHDは花巻支店登記で岩手中部圏域クリア済み\n　→ 他地域案件は「支店登記の追加」も戦略になる\n② 欠格要件：破産/更生/税滞納/反社でないこと\n③ 添付書類の有効期限：発行後3ヶ月以内が通例\n　→ 締切から逆算して取得（早取りは失効リスク）\n\n法人申込の標準セット：\n全部事項証明書／印鑑登録証明書／市町村税納税証明書\n／誓約書／（代理人なら委任状＋身分証写し）',
 ['地域要件を自社の登記で満たすことを謄本で確認した',
  '税の未納なし（誓約書と矛盾しない）を確認した',
  '添付書類の取得日を締切から逆算して決めた',
  '様式を正確な商号・住所で下書きした',
  '不明点は質問書で事前照会した（期限内）',
  '提出は持参 or 簡易書留等の記録が残る方法にした',
  '受付期間内「必着」を確認した（消印有効ではない）'])

step_slide('STEP 3','机上査定：仕様書精読 × 相見積りで「売値」を作る','仕様書から読み取る査定因子',
 '・型式→グレード特定（MT/AT・バン/ワゴンで絞れる）\n・走行距離/初度登録＝国内価値。ランクル等は過走行\n　でも海外価値が残る（50万kmからが本番）\n・車検/自賠責の残＝自走可否・搬出コストに直結\n・特殊装備（赤色灯等）＝撤去・用途変更の手間\n\n相見積りの鉄則：\n・3〜5社、タイプの違う先（国内買取/旧車専門/輸出\n　/知人ルート）に当てる\n・提示額は相互に伏せる（アンカリング回避）\n・車台番号は概算段階では渡さない（本命1社のみ）',
 ['型式・年式・走行・装備からグレードを特定した',
  '排ガス規制(NOx・PM法)の適合を確認した ※重要',
  '査定先を3〜5社リストアップ（タイプ分散）した',
  '競合化リスクを精査した（地域要件を満たす業者は除外）',
  '概算査定で「保守/中央/強気」の3シナリオを作った',
  '査定根拠（メール・査定書）を証跡保存した',
  '想定vs実際の比較表を作り、取得のたび埋めている'])

step_slide('STEP 4','現物確認：目利きマップ(p03〜08)を現場で回す','当日の動線（60〜90分）',
 '① 書類照合(10分)：車台番号打刻・車検証・付属品\n② 外周(15分)：外装4方向→パネル隙間→ガラス→装備\n③ 下回り(20分)：フレーム後部→サイドレール→足回り\n　※ここに最長時間。マット・ライト・軍手持参\n④ 機関(15分)：冷間始動→アイドル→排気色→漏れ\n⑤ 内装(10分)：ODO→クラッチ→床錆→装備作動\n⑥ 質問(10分)：装備撤去可否・整備記録の有無・\n　過去の修理歴を担当者に直接確認\n\n撮影は次ページの構図テンプレ20枚に沿って撮る',
 ['事前予約の電話を入れ、訪問時間を確定した',
  '持ち物を揃えた（ライト/軍手/マット/磁石/脚立）',
  '目利き12ポイントを順番どおり全て見た',
  '下回りは叩き＋目視で3ゾーン(A/B/C)判定した',
  '冷間始動と排気色を動画で記録した',
  '車台番号打刻と書類の一致を確認した',
  '写真20枚を撮り、当日中に査定先へ送付した'])

step_slide('STEP 5','収支確定：売値から逆算して入札上限を決める','逆算式と本案件の当てはめ',
 '入札上限 ＝（保守売却額 − 目標利益 − 諸費用）÷ 1.10\n\n・契約金額＝入札額×1.10（消費税相当）が通例\n・諸費用＝陸送(0〜10万・業者引取無料なら0)\n　＋名義変更等(1〜3万)\n・保有なら＋車検整備(15〜40万)＋NOx・PM登録可否\n\n本案件（旧車王実査定100〜200万円）：\n保守100万・利益30万狙い → 上限52万\n中央150万 → 上限98万\n→ 推奨入札 55〜65万円／絶対上限80万円',
 ['売却3シナリオを実査定で裏付けた',
  '諸費用を業者別の実見積りで確定した',
  '入札締切1週間前までに全見積りを回収した',
  '「絶対上限」を先に決め、書面に残した',
  '端数を付けた入札額にした（同額くじ回避）',
  '資金（契約金額の全額前払い）を確保した',
  '現物確認の結果を入札額に反映した（下回り減点等）'])

step_slide('STEP 6','入札：一発勝負のルールを間違えない','絶対に落とせない実務ルール',
 '・入札回数1回／提出後の書換・撤回不可\n・金額は税抜で記入（契約金額は×1.10）\n・「￥」記入・アラビア数字・訂正不可のペン\n・封筒に件名明記・封印\n・郵送は「必着」。簡易書留/特定記録で記録を残す\n\n無効になる典型：\n記名押印漏れ／金額訂正／指定様式以外／\n資本・人的関係のある複数者の入札（連合とみなし）',
 ['入札書は指定様式・訂正なしで記入した',
  '金額の頭に￥、桁ズレがないかダブルチェックした',
  '記名・押印（法人は代表者印）を確認した',
  '封筒に件名を明記し封印した',
  '提出期限の前日までに到達する手段で発送した',
  '開札日・結果公表の確認方法を控えた',
  '落札できなかった場合の次案件リストを用意した'])

step_slide('STEP 7','落札後：契約・全額前払い・搬出を期限内に','期限が3本走る（1本落とすと決定取消も）',
 '① 契約締結期限（例：落札から約1週間）\n② 代金全額前払い（納入通知書で契約締結までに）\n③ 搬出期限（例：月末・費用は落札者負担）\n\n搬出の実務：\n・車検切れ＝公道自走不可 → 積載車を事前手配\n　（仮ナンバー自走は老朽車ではリスク高・非推奨）\n・買取業者への直送も検討：引取無料の業者なら\n　陸送費が丸ごと浮く＝実質利益+5〜10万\n・名義変更/抹消は速やかに（税・責任の切替）',
 ['契約・支払・搬出の3期限をカレンダー登録した',
  '納入通知書での支払方法・期日を確認した',
  '陸送 or 業者引取を確定し、搬出日を予約した',
  '売却先と引取場所・日程を握った（直送が最良）',
  '名義変更（または抹消）の段取りを決めた',
  '特殊装備の撤去有無・タイミングを確定した',
  '実費の領収書を全て保存した（比較表に反映）'])

step_slide('STEP 8','売却・関係構築・記録：次の案件につなげる','売却と事業化',
 '・相見積りの最高値＋対応の質で売却先を決定\n・取引完了時に「今後も官公庁売払いを継続的に扱う」\n　と伝え、担当者の直通連絡先を確保\n・担当者は名刺DB/顧客マスターに登録（次回は電話\n　1本で査定が回る状態を作る）\n\n事業化の法規：\n・古物商許可＝反復継続の転売に必須\n　（無許可営業は3年以下の懲役/100万円以下の罰金）\n・2台目に着手する前に取得完了（申請〜許可1〜2ヶ月）',
 ['本番査定（車台番号開示）は最有力1社に絞った',
  '売買契約書・振込記録を保存した',
  '担当者名・直通連絡先を名刺DBに登録した',
  '想定vs実際の比較表を完成させ、差異を分析した',
  '案件全体の学びをnotesに記録した',
  '古物商許可の取得状況を確認した（2台目の前提）',
  '次の売払い案件のウォッチを再開した'])

# ============ 17. 撮影構図テンプレ ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'FIELD KIT','現地の持ち物 ＋ 撮影構図テンプレ20枚（査定送付用）')
card(s,Inches(0.7),Inches(1.7),Inches(4.4),Inches(5.0),fill=WHITE)
txt(s,Inches(1.0),Inches(1.9),Inches(3.9),Inches(0.4),'持ち物リスト',size=14.5,color=DBRICK,bold=True)
checklist(s,Inches(1.0),Inches(2.4),Inches(3.8),[
 'LEDライト（下回り照射用）','軍手・汚れてよい服装','地面用マット or 段ボール',
 '磁石（パテ盛り判定・ボディ用）','スマホ＋モバイルバッテリー','物件仕様書のコピー（照合用）',
 'メジャー（荷室採寸）','質問メモ（装備撤去可否等）'],size=11.5,gap=0.42)
card(s,Inches(5.4),Inches(1.7),Inches(7.3),Inches(5.0),fill=WHITE)
txt(s,Inches(5.7),Inches(1.9),Inches(6.8),Inches(0.4),'撮影構図20枚テンプレ（この順で撮ればそのまま査定に送れる）',size=13.5,color=DBRICK,bold=True)
cols=[['① 前方全景（ナンバー入り）','② 後方全景（バッジ入り）','③ 左側面全景','④ 右側面全景','⑤ 左前45度','⑥ 右後45度','⑦ 下回り：後部フレーム','⑧ 下回り：左レール','⑨ 下回り：右レール','⑩ 下回り：前部クロス'],
      ['⑪ タイヤ4本（溝・製造年）','⑫ エンジンルーム全景','⑬ エンジン：オイル周り接写','⑭ 冷間始動の動画＋排気','⑮ メーター（ODO読める解像度）','⑯ 運転席全景（シフト入り）','⑰ シート・内装の傷汚れ','⑱ 荷室＋付属タイヤ','⑲ 車台番号打刻の接写','⑳ 傷・錆・凹みの個別接写(全部)']]
x0=[Inches(5.7),Inches(9.3)]
for ci,col in enumerate(cols):
    y=Inches(2.45)
    for item in col:
        txt(s,x0[ci],y,Inches(3.5),Inches(0.35),item,size=11,color=INK)
        y=add(y,Inches(0.40))
footer(s)

# ============ 18. リスク・法規 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'RISK & LAW','リスクと法規：先に知らないと即死する4つ')
items=[('NOx・PM法','旧年式ディーゼルは首都圏・大阪・愛知等の対策地域で登録不可。「東京で保有」が\n物理的に不可能な車がある。転売(輸出)なら無関係。本案件は保有なら花巻登録前提。','最重要'),
('現状渡し・責任なし','引渡し後の故障・瑕疵は全て自己責任。現物確認が唯一の防衛線。\n下回り腐食は売値を数十万単位で変える（配点40点の理由）。','高'),
('一発入札','書換・撤回・再入札なし。相場を外すと高値掴みか未落札。\n「絶対上限」を先に書面化してから入札書を書く。','高'),
('古物商許可','反復継続の転売は許可必須。1台の単発は直ちに違法ではないが、\n事業化するなら2台目までに取得完了（申請〜許可1〜2ヶ月）。','中')]
y=Inches(1.75)
for name,desc,tag in items:
    card(s,Inches(0.7),y,Inches(12.0),Inches(1.14),fill=WHITE)
    box(s,Inches(0.7),y,Inches(0.14),Inches(1.14),fill=BRICK)
    txt(s,Inches(1.05),add(y,Inches(0.10)),Inches(2.6),Inches(0.9),name,size=14.5,color=DBRICK,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.8),add(y,Inches(0.09)),Inches(7.6),Inches(1.0),desc,size=11.5,color=INK,sp=1.15)
    chip(s,Inches(11.55),add(y,Inches(0.36)),tag,fill=LBRICK,fg=DBRICK,w=Inches(0.95),h=Inches(0.4),size=11)
    y=add(y,Inches(1.30))
footer(s)

# ============ 19. ケーススタディ ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'CASE STUDY','実例：岩手中部水道企業団 公告148号（2026）')
card(s,Inches(0.7),Inches(1.7),Inches(6.0),Inches(5.0),fill=WHITE)
txt(s,Inches(1.0),Inches(1.9),Inches(5.4),Inches(0.4),'案件データ（確定情報のみ）',size=15,color=DBRICK,bold=True)
txt(s,Inches(1.0),Inches(2.4),Inches(5.5),Inches(4.1),
    '車両：ランクル100 バンVX ディーゼル5MT(KG-HDJ101K)\n　　　※VXバッジ・5MT・貨物登録の実写で確定\n平成10年6月初度登録／280,113km（ODO実写一致）\n車検切れ(R8.6.17満了)／公共応急作業車(8ナンバー)\n最低売却価格：15万円\n実査定：旧車王 概算100〜200万円（不成立・伏せて相見積り中）\n推奨入札：55〜65万円（絶対上限80万円）\n\n8/20現物確認 → 8/21書類締切 → 9/7入札 →\n9/8開札 → 9/15契約・支払 → 9/30搬出',size=12,color=INK,sp=1.3)
card(s,Inches(7.0),Inches(1.7),Inches(5.7),Inches(5.0),fill=PANEL)
txt(s,Inches(7.3),Inches(1.9),Inches(5.1),Inches(0.4),'この案件で確立した「型」',size=15,color=DBRICK,bold=True)
txt(s,Inches(7.3),Inches(2.4),Inches(5.1),Inches(4.1),
    '・地域要件は支店登記でクリア（参入障壁を自分だけ\n　越える＝安く仕入れる構造）\n・役所の最低売却価格は市場価格と無関係。仕様書の\n　精読とプロ査定で「本当の値段」を掴む\n・買取業者は地域要件を満たせない＝情報を渡しても\n　入札で競合しない。ただし買い叩きには注意\n・車台番号は本命1社の本番査定まで温存\n・相見積りは提示額を相互に伏せる\n・全記録は .company/secretary/notes/ に集約。\n　想定vs実際の差異が次回案件の精度になる',size=12,color=INK,sp=1.3)
footer(s)

# ============ 20. 総覧 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'MASTER CHECKLIST','全チェックリスト総覧（印刷してこの1枚で回せる）')
groups=[('S1 発掘','公告確認/勝ち筋判定/資格先読み/逆算表/資料DL/公開日登録/過去落札'),
('S2 資格','登記で要件確認/税未納なし/書類逆算取得/様式下書き/質問書/記録郵送/必着'),
('S3 査定','グレード特定/排ガス規制/3〜5社分散/競合精査/3シナリオ/証跡保存/比較表'),
('S4 現物','予約/持ち物/12ポイント/下回りABC/冷間始動動画/打刻照合/写真20枚'),
('S5 収支','3シナリオ裏付け/諸費用実額/1週間前回収/絶対上限書面化/端数/資金/現物反映'),
('S6 入札','指定様式/金額ダブルチェック/記名押印/封筒件名/前日到達/結果確認/次案件'),
('S7 決済','3期限登録/納入通知書/搬出予約/売却先と握り/名義変更/装備撤去/領収書'),
('S8 売却','本番査定1社/契約書保存/担当者DB登録/比較表完成/学び記録/古物商/次ウォッチ')]
y=Inches(1.75)
for name,desc in groups:
    card(s,Inches(0.7),y,Inches(12.0),Inches(0.56),fill=WHITE,depth=Inches(0.03))
    box(s,Inches(0.7),y,Inches(1.5),Inches(0.56),fill=BRICK)
    txt(s,Inches(0.7),y,Inches(1.5),Inches(0.56),name,size=12,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(2.4),y,Inches(10.2),Inches(0.56),desc,size=11.5,color=INK,anchor=MSO_ANCHOR.MIDDLE)
    y=add(y,Inches(0.64))
txt(s,Inches(0.7),add(y,Inches(0.05)),Inches(12),Inches(0.4),'詳細は各STEPページ参照。目利き12ポイント＝p03、撮影テンプレ20枚＝p17。',size=11,color=MUTE)
footer(s)

prs.save('/home/user/khd_workspace/KHD_車両転売_仕入売却マニュアル_v2_目利き完全版.pptx')
print('saved pages:',_pages['n']+1)

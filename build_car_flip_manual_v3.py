#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
KHD 車両転売事業マニュアル v3 — H037テンプレ準拠版
------------------------------------------------------------
デザイン: H037_AI医療コンサル(260627)のデザイン言語に完全準拠
- フラット（影なし・角丸なし）／ベージュカード#F1ECE1に赤天冠バー
- Hiragino Sans／本文#1A1A1A／グレー#6E6E6E／赤#AA2E26／見出し赤#8C241D
- ポジティブ差し色グリーン#0F6E56/#E1F5EE
内容: v2（目利き完全版・実車写真9枚・全21ページ構成）を踏襲
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK=RGBColor(0x1A,0x1A,0x1A); GRAY=RGBColor(0x6E,0x6E,0x6E)
RED=RGBColor(0xAA,0x2E,0x26); DRED=RGBColor(0x8C,0x24,0x1D)
BEIGE=RGBColor(0xF1,0xEC,0xE1); LINEC=RGBColor(0xDA,0xD6,0xCF); LRED=RGBColor(0xF4,0xE4,0xE2)
GREEN=RGBColor(0x0F,0x6E,0x56); LGREEN=RGBColor(0xE1,0xF5,0xEE); WHITE=RGBColor(0xFF,0xFF,0xFF)
STEEL=RGBColor(0xC9,0xD2,0xD8); DSTEEL=RGBColor(0x8F,0x9D,0xA6)
FONT='Hiragino Sans'
ASSET='/home/user/khd_workspace/assets_landcruiser_148'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
_pages={'n':0}

def add(v,d): return Emu(int(v)+int(d))
def bg(s,c=WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def oval(s,l,t,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def txt(s,l,t,w,h,text,size=12,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if sp is not None: p.line_spacing=sp
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=FONT
    return tb
def footer(s):
    _pages['n']+=1
    box(s,Inches(0.5),Inches(7.05),Inches(12.33),Inches(0.02),fill=LINEC)
    txt(s,Inches(0.5),Inches(7.1),Inches(10.0),Inches(0.3),'KHD 車両転売事業マニュアル v3（公用車入札・目利き完全版）',size=8,color=GRAY)
    txt(s,Inches(12.45),Inches(7.1),Inches(0.4),Inches(0.3),str(_pages['n']+1),size=8,color=GRAY)
def header(s,kicker,title,sub=None):
    txt(s,Inches(0.55),Inches(0.35),Inches(12.0),Inches(0.4),kicker,size=11,color=GRAY,bold=True)
    box(s,Inches(0.57),Inches(0.72),Inches(1.6),Inches(0.04),fill=RED)
    txt(s,Inches(0.55),Inches(0.82),Inches(12.2),Inches(0.55),title,size=20,color=INK,bold=True)
    if sub: txt(s,Inches(0.57),Inches(1.36),Inches(12.2),Inches(0.3),sub,size=11,color=GRAY)
def bcard(s,l,t,w,h,title=None,fill=BEIGE,bar=RED):
    box(s,l,t,w,h,fill=fill)
    box(s,l,t,w,Inches(0.06),fill=bar)
    if title: txt(s,add(l,Inches(0.22)),add(t,Inches(0.18)),add(w,-Inches(0.44)),Inches(0.5),title,size=14,color=INK,bold=True)
def checklist(s,l,t,w,items,size=11.5,gap=0.375,color=INK):
    y=t
    for it in items:
        box(s,l,add(y,Inches(0.045)),Inches(0.15),Inches(0.15),fill=None,line=RED,lw=1.2)
        txt(s,add(l,Inches(0.28)),y,add(w,-Inches(0.28)),Inches(0.4),it,size=size,color=color,sp=1.0)
        y=add(y,Inches(gap))
    return y
def numdot(s,l,t,n,d=Inches(0.32),fill=RED,fg=WHITE,size=11):
    oval(s,l,t,d,d,fill=fill)
    txt(s,l,add(t,-Inches(0.015)),d,d,str(n),size=size,color=fg,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def pic(s,path,l,t,w=None,label=None,note=None):
    p=s.shapes.add_picture(os.path.join(ASSET,path),l,t,width=w)
    if label:
        box(s,l,t,Inches(1.9),Inches(0.34),fill=RED)
        txt(s,l,t,Inches(1.9),Inches(0.34),label,size=10.5,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note:
        txt(s,l,add(add(t,p.height),Inches(0.04)),p.width,Inches(0.9),note,size=10,color=GRAY,sp=1.1)
    return p

# ============ 1. COVER（テンプレの表紙文法） ============
s=prs.slides.add_slide(BLANK); bg(s)
box(s,Inches(0.5),Inches(0.5),Inches(0.06),Inches(6.5),fill=RED)
txt(s,Inches(0.9),Inches(1.1),Inches(11.0),Inches(0.45),'車両転売事業（公用車入札）',size=14,color=DRED,bold=True)
txt(s,Inches(0.88),Inches(1.65),Inches(11.7),Inches(1.7),'仕入〜売却マニュアル v3\n── 不動産と同じ8STEPの型 ＋ 目利き完全版',size=26,color=INK,bold=True,sp=1.15)
txt(s,Inches(0.92),Inches(3.8),Inches(11.4),Inches(0.5),'目利き（現物確認）が背骨。机上査定は仮説、現物が検証、収支は検証済みの数字だけで組む。',size=13,color=GRAY)
box(s,Inches(0.9),Inches(4.6),Inches(11.5),Inches(1.4),fill=BEIGE)
box(s,Inches(0.9),Inches(4.6),Inches(0.1),Inches(1.4),fill=RED)
txt(s,Inches(1.2),Inches(4.75),Inches(11.0),Inches(1.15),'発掘→資格→机上査定→現物→収支→入札→決済→売却。全STEPチェックリスト・目利き12ポイント図解・\n実車写真9枚（岩手中部水道企業団 ランクル100 バンVX・公告148号）・撮影テンプレ20枚・総覧1枚を収録。',size=12,color=INK,sp=1.3)
txt(s,Inches(0.92),Inches(6.45),Inches(11),Inches(0.4),'KIKUCHIホールディングス株式会社 ｜ 2026-07-15 ｜ v3（H037テンプレ準拠）',size=10,color=GRAY)

# ============ 2. 全体フロー ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'全体設計','仕入〜売却の8ステップ（不動産と同じ型）','赤＝目利き（S4）が全STEPの背骨')
steps=[('S1','案件発掘','公告を拾う'),('S2','参加資格','門前払い回避'),('S3','机上査定','相場を作る'),
       ('S4','現物確認','目利き＝内見'),('S5','収支確定','逆算で上限'),('S6','入札','一発勝負'),
       ('S7','決済・搬出','期限3本'),('S8','売却・記録','最高値＋関係')]
x=Inches(0.55); w=Inches(2.96); gx=Inches(0.13); gy=Inches(0.18); hh=Inches(1.18)
for i,(no,name,desc) in enumerate(steps):
    col=i%4; row=i//4
    l=add(x,int(add(w,gx))*col); t=add(Inches(1.9),int(add(hh,gy))*row)
    bcard(s,l,t,w,hh,bar=DRED if i==3 else RED)
    txt(s,add(l,Inches(0.22)),add(t,Inches(0.16)),Inches(2.5),Inches(0.4),f'{no}　{name}',size=14,color=DRED if i==3 else INK,bold=True)
    txt(s,add(l,Inches(0.24)),add(t,Inches(0.62)),Inches(2.5),Inches(0.5),desc,size=11,color=GRAY)
box(s,Inches(0.55),Inches(4.85),Inches(11.86),Inches(0.9),fill=LRED)
txt(s,Inches(0.85),Inches(4.97),Inches(11.3),Inches(0.7),'目利き（S4）が背骨：不動産で「内見しない買付」がないのと同じ。S3は仮説、S4が検証、S5は検証済みの数字だけで組む。',size=12.5,color=DRED,bold=True,sp=1.25)
box(s,Inches(0.55),Inches(5.95),Inches(11.86),Inches(0.95),fill=BEIGE)
txt(s,Inches(0.85),Inches(6.06),Inches(11.3),Inches(0.8),'不動産との対応：レインズ発掘=公告ウォッチ／謄本・重説=仕様書・公告／内見=現物公開日／買付=入札書（一発・撤回不可）／決済=契約・全額前払い／転売先=相見積りで最高値',size=11,color=INK,sp=1.25)
footer(s)

# ============ 3. 目利きマップ ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'目利き ｜ 最重要','車両目利きマップ：見るべき12ポイント（側面模式図）','番号＝右のリストに対応。配点は下段')
box(s,Inches(0.55),Inches(1.75),Inches(8.3),Inches(3.5),fill=BEIGE)
box(s,Inches(0.55),Inches(1.75),Inches(8.3),Inches(0.06),fill=RED)
# car schematic
box(s,Inches(1.15),Inches(3.15),Inches(6.9),Inches(1.1),fill=STEEL,line=DSTEEL,lw=1.2)
box(s,Inches(2.25),Inches(2.5),Inches(4.2),Inches(0.9),fill=STEEL,line=DSTEEL,lw=1.2)
box(s,Inches(2.5),Inches(2.64),Inches(1.6),Inches(0.58),fill=WHITE,line=DSTEEL)
box(s,Inches(4.3),Inches(2.64),Inches(1.85),Inches(0.58),fill=WHITE,line=DSTEEL)
box(s,Inches(2.45),Inches(2.33),Inches(3.8),Inches(0.14),fill=DSTEEL)
oval(s,Inches(5.55),Inches(2.08),Inches(0.28),Inches(0.28),fill=RED)
oval(s,Inches(2.0),Inches(3.9),Inches(1.0),Inches(1.0),fill=INK)
oval(s,Inches(6.3),Inches(3.9),Inches(1.0),Inches(1.0),fill=INK)
oval(s,Inches(2.27),Inches(4.17),Inches(0.46),Inches(0.46),fill=STEEL)
oval(s,Inches(6.57),Inches(4.17),Inches(0.46),Inches(0.46),fill=STEEL)
box(s,Inches(1.45),Inches(4.33),Inches(6.3),Inches(0.12),fill=LRED)
pts=[(1,Inches(3.85),Inches(4.42)),(2,Inches(2.2),Inches(3.98)),(3,Inches(1.28),Inches(3.3)),
     (4,Inches(4.85),Inches(2.66)),(5,Inches(5.5),Inches(2.0)),(6,Inches(7.42),Inches(3.05)),
     (7,Inches(3.95),Inches(3.45)),(8,Inches(3.05),Inches(2.9))]
for n,l,t in pts: numdot(s,l,t,n)
txt(s,Inches(1.15),Inches(4.85),Inches(7.4),Inches(0.3),'KHD自作の模式図（側面）。実車写真の見方は次ページ以降。',size=9.5,color=GRAY)
# right list
box(s,Inches(9.1),Inches(1.75),Inches(3.6),Inches(4.95),fill=BEIGE)
box(s,Inches(9.1),Inches(1.75),Inches(3.6),Inches(0.06),fill=RED)
txt(s,Inches(9.32),Inches(1.93),Inches(3.2),Inches(0.4),'12の目利きポイント',size=13,color=DRED,bold=True)
pts12=['① 下回り・フレーム錆（最重要）','② タイヤ・足回り・ブッシュ','③ ボディ下部パネルの錆・膨れ','④ ガラス・モール・ドア開閉','⑤ 特殊装備（赤色灯・ラック）','⑥ エンジンルーム（漏れ・煙）','⑦ 室内・シート・床の錆','⑧ メーター（ODO書類一致）','⑨ シフト・クラッチ（MT）','⑩ 付属品（鍵・タイヤ・書類）','⑪ 車台番号打刻と書類一致','⑫ 排気（白煙/黒煙・異音）']
y=Inches(2.38)
for p_ in pts12:
    txt(s,Inches(9.32),y,Inches(3.25),Inches(0.33),p_,size=11,color=INK)
    y=add(y,Inches(0.345))
box(s,Inches(0.55),Inches(5.5),Inches(8.3),Inches(1.2),fill=LRED)
txt(s,Inches(0.85),Inches(5.62),Inches(7.8),Inches(1.0),'配点：①下回り40点／⑥機関20点／③外装15点／⑧⑪整合10点／その他15点。\n①が崩れたら（貫通錆）他が満点でも降りる＝不動産の「構造・傾き」と同じ扱い。',size=12,color=DRED,bold=True,sp=1.3)
footer(s)

# ============ 4. 下回り ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'目利き 1/4 ｜ 配点40','下回り・フレーム：買値が数十万動く場所（下面模式図）','赤＝融雪剤地域で錆が集中する後部ホットゾーン')
box(s,Inches(0.55),Inches(1.75),Inches(6.3),Inches(4.95),fill=BEIGE)
box(s,Inches(0.55),Inches(1.75),Inches(6.3),Inches(0.06),fill=RED)
fx=Inches(1.5); fy=Inches(2.25); fh=Inches(3.55)
box(s,fx,fy,Inches(0.3),fh,fill=STEEL,line=DSTEEL)
box(s,add(fx,Inches(4.0)),fy,Inches(0.3),fh,fill=STEEL,line=DSTEEL)
for i in range(5):
    box(s,add(fx,Inches(0.3)),add(fy,Inches(0.18+0.78*i)),Inches(3.7),Inches(0.2),fill=STEEL,line=DSTEEL)
for wl,wt in [(-0.42,0.25),(4.3,0.25),(-0.42,2.4),(4.3,2.4)]:
    oval(s,add(fx,Inches(wl)),add(fy,Inches(wt)),Inches(0.4),Inches(0.92),fill=INK)
box(s,fx,add(fy,Inches(2.55)),Inches(0.3),Inches(1.0),fill=RED)
box(s,add(fx,Inches(4.0)),add(fy,Inches(2.55)),Inches(0.3),Inches(1.0),fill=RED)
box(s,add(fx,Inches(0.3)),add(fy,Inches(3.3)),Inches(3.7),Inches(0.2),fill=RED)
numdot(s,add(fx,Inches(1.8)),add(fy,Inches(3.22)),'A',fill=DRED)
numdot(s,add(fx,-Inches(0.13)),add(fy,Inches(1.45)),'B',fill=GRAY)
numdot(s,add(fx,Inches(1.8)),add(fy,Inches(0.08)),'C',fill=GRAY)
txt(s,Inches(0.85),Inches(6.05),Inches(5.8),Inches(0.5),'上から見たラダーフレーム模式図（KHD自作）',size=9.5,color=GRAY)
box(s,Inches(7.15),Inches(1.75),Inches(5.55),Inches(4.95),fill=BEIGE)
box(s,Inches(7.15),Inches(1.75),Inches(5.55),Inches(0.06),fill=RED)
txt(s,Inches(7.4),Inches(1.95),Inches(5.0),Inches(0.4),'見方と判定基準',size=14,color=DRED,bold=True)
txt(s,Inches(7.4),Inches(2.45),Inches(5.05),Inches(3.2),
    'A 後部フレーム・リアクロス（最頻発）\n　→柄などで軽く叩く。「ボコッ」と鈍い音・剥落＝内部腐食\nB サイドレール全長（潜って目視）\n　→表面サビ(茶)はOK。「膨れ・層状剥離」はNG\nC 前部クロス・足回り取付部\n　→ブッシュ/ボルト周りの固着と亀裂\n\n判定：\n・表面サビのみ＝減点小（輸出値ほぼ維持）\n・膨れ・剥離＝減額交渉材料\n・貫通穴＝入札額を大幅に下げるか撤退',size=11.5,color=INK,sp=1.25)
box(s,Inches(7.4),Inches(5.95),Inches(5.05),Inches(0.55),fill=LRED)
txt(s,Inches(7.5),Inches(5.98),Inches(4.9),Inches(0.5),'岩手＝融雪剤地域。本案件最大の変数はここ',size=11.5,color=DRED,bold=True,anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ============ 5. 機関 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'目利き 2/4 ｜ 配点20','機関（エンジン・排気）：実車写真で見る','写真で分かること／現地でしか分からないことを分ける')
pic(s,'engine.jpg',Inches(0.55),Inches(1.8),w=Inches(6.2),label='実車：エンジンルーム',
    note='1HD-FTE 4.2Lディーゼルターボ。ツインバッテリー（寒冷地/ディーゼル標準）。写真上は大きなオイル飛散・冷却水漏れ痕なし、年式比で良好な印象。')
box(s,Inches(7.05),Inches(1.8),Inches(5.65),Inches(4.85),fill=BEIGE)
box(s,Inches(7.05),Inches(1.8),Inches(5.65),Inches(0.06),fill=RED)
txt(s,Inches(7.3),Inches(2.0),Inches(5.1),Inches(0.4),'現地で確認する7点（写真では不明）',size=13.5,color=DRED,bold=True)
checklist(s,Inches(7.3),Inches(2.5),Inches(5.1),[
 '冷間始動の一発始動性（ディーゼルの健康診断）',
 'アイドリング安定・異音（カラカラ/ガラガラ）',
 '排気色：白煙(オイル下がり)/黒煙(燃調・詰まり)',
 'オイルにじみ：ヘッドカバー・タービン周り',
 '冷却水の色（サビ色＝管理不良のサイン）',
 'バッテリー端子の腐食・補機ベルトの亀裂',
 'ターボの効き（可能なら構内で短距離試走）'],size=11.5,gap=0.42)
box(s,Inches(7.3),Inches(5.85),Inches(5.1),Inches(0.55),fill=LRED)
txt(s,Inches(7.42),Inches(5.88),Inches(4.9),Inches(0.5),'28万km級1HDは「壊れる」より「漏れ・煙」を見る',size=11.5,color=DRED,bold=True,anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ============ 6. 内装・整合 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'目利き 3/4 ｜ 配点10＋α','内装・書類整合：「登記と現況の照合」に相当','ODO・車台番号・グレードの3整合が取れて初めて数字を信じる')
pic(s,'meter_odo.jpg',Inches(0.55),Inches(1.85),w=Inches(5.6),label='実車：メーター',
    note='★ODO実写 280,113km＝物件仕様書(R8.7.9時点)と完全一致。メーター整合OKの現物証拠。タコ6,000スケール＝ディーゼル。警告灯は停止中の点灯で正常。')
pic(s,'cockpit_mt.jpg',Inches(6.65),Inches(1.85),w=Inches(5.6),label='実車：運転席',
    note='★5速MT＋トランスファーレバー実写＝バンVXディーゼル5MTの現物証拠。業務無線マイク残置（撤去対象）。現地ではクラッチの繋がり位置・シート下の床錆を確認。')
footer(s)

# ============ 7. 装備・付属品 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'目利き 4/4 ｜ 配点15','装備・付属品：足す価値と引く手間を仕分ける','スタッドレス＝足す価値／無線・赤色灯＝引く手間')
pic(s,'cluster.jpg',Inches(0.55),Inches(1.85),w=Inches(5.6),label='実車：センター',
    note='業務無線・広報アンプ残置（「さぎょうだん18」ラベル＝作業団車両）。無線・広報設備・赤色灯は転売前に撤去が基本（外して引渡し可能かは質問書で確認）。')
pic(s,'trunk_tires.jpg',Inches(6.65),Inches(1.85),w=Inches(5.6),label='実車：荷室',
    note='★スタッドレス4本（2021年12月購入）付属＝＋2〜5万円相当の価値。テールゲート開閉OK。荷室フロアの錆・水染み、付属品（鍵3本・車検証・取説・リサイクル券）を現地で全点照合。')
footer(s)

# ============ 8. 写真評価1 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'写真評価 1/2','実車写真の現時点評価：分かったこと／8/20で埋めること','公式写真は上物だけ。最大配点（下回り40点）が未評価')
pic(s,'side_left.jpg',Inches(0.55),Inches(1.8),w=Inches(4.9),label='実車：左側面')
pic(s,'rear_vx.jpg',Inches(0.55),Inches(4.55),w=Inches(3.5),label='実車：後方')
txt(s,Inches(4.2),Inches(4.6),Inches(1.8),Inches(2.0),'← VXバッジ実写\nグレード確定の\n現物証拠',size=11,color=DRED,bold=True,sp=1.25)
box(s,Inches(6.1),Inches(1.8),Inches(6.6),Inches(4.85),fill=BEIGE)
box(s,Inches(6.1),Inches(1.8),Inches(6.6),Inches(0.06),fill=RED)
txt(s,Inches(6.4),Inches(2.0),Inches(6.0),Inches(0.4),'写真9枚から言えること（2026-07-15時点）',size=13.5,color=DRED,bold=True)
txt(s,Inches(6.4),Inches(2.5),Inches(6.0),Inches(2.4),
    '◎ 外装：ツヤあり・大きな凹み/事故痕は視認できず\n◎ グレード：後方VXバッジ・5MT・貨物登録 → バンVX確定\n◎ 整合：ODO実写=仕様書=280,113km一致\n◎ 装備：ラック・ラダー・赤色灯・無線（撤去要否を判断）\n○ 機関：表層きれい。ただし音・煙・漏れは写真で不明\n△ 内装：業務使用感あり（減点小）',size=11.5,color=INK,sp=1.35)
box(s,Inches(6.4),Inches(5.1),Inches(6.0),Inches(1.3),fill=LRED)
txt(s,Inches(6.55),Inches(5.2),Inches(5.8),Inches(1.1),'✖ 写真に「下回り」が1枚もない＝配点40点が未評価。\n下回りが表面サビ止まりなら、旧車王レンジ(100〜200万)の中央〜上限が狙える。',size=12,color=DRED,bold=True,sp=1.3)
footer(s)

# ============ 9. 写真評価2 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'写真評価 2/2','外観の記録：前方・右側面と見るポイント','ナンバー=岩手800「さ・134」＝8ナンバー特種登録の現物確認')
pic(s,'front.jpg',Inches(0.55),Inches(1.85),w=Inches(5.6),label='実車：前方',
    note='赤色灯・ルーフラック・フォグ付き。バンパー/グリルに大きな損傷なし。ライトのくもりは年式なり（研磨で改善可）。現地では飛び石サビ・ガラス傷を確認。')
pic(s,'side_right.jpg',Inches(6.65),Inches(1.85),w=Inches(5.6),label='実車：右側面',
    note='ライン通り・パネル波打ちなし（写真上）。サイドステップ下端・ドア下端・フェンダーアーチ内側は錆やすい部位＝指で触って膨れを確認。欠品/割れは個別撮影。')
footer(s)

# ============ 10-17. S1〜S8 ============
def step_slide(kicker,title,sub,left_title,left_body,check_items):
    s=prs.slides.add_slide(BLANK); bg(s)
    header(s,kicker,title,sub)
    box(s,Inches(0.55),Inches(1.75),Inches(6.1),Inches(4.95),fill=BEIGE)
    box(s,Inches(0.55),Inches(1.75),Inches(6.1),Inches(0.06),fill=RED)
    txt(s,Inches(0.85),Inches(1.95),Inches(5.5),Inches(0.4),left_title,size=14,color=DRED,bold=True)
    txt(s,Inches(0.85),Inches(2.45),Inches(5.55),Inches(4.1),left_body,size=11.5,color=INK,sp=1.28)
    box(s,Inches(6.95),Inches(1.75),Inches(5.75),Inches(4.95),fill=BEIGE)
    box(s,Inches(6.95),Inches(1.75),Inches(5.75),Inches(0.06),fill=RED)
    txt(s,Inches(7.2),Inches(1.95),Inches(5.0),Inches(0.4),'チェックリスト',size=13.5,color=DRED,bold=True)
    checklist(s,Inches(7.2),Inches(2.45),Inches(5.2),check_items,size=11.5,gap=0.44)
    footer(s); return s

step_slide('STEP 1','案件発掘：官公庁の売払い公告を拾う','狙い目＝地域限定×海外需要車種×安すぎる最低価格','探す場所と狙い目',
 '・市町村/水道企業団/一部事務組合HPの「入札情報」\n・「公用車 売払い 入札」「不用物品 売払い」で定期検索\n・官公庁オークション（KSI等）も併読\n\n狙い目の型：\n・地方の水道/消防/土木系＝ランクル/ハイエース等\n　海外需要の高い車が出やすい\n・地域限定案件＝参入障壁が高い分、安く取れる\n・最低売却価格は市場価格と無関係（役所は相場を\n　取りに行かない）。ここが利益の源泉',
 ['公告日・入札方式（一般競争/条件付/せり）を確認した',
  '最低売却価格と車種で「勝ち筋」を10分で仮判定した',
  '参加資格（地域要件・法人/個人）を最初に読んだ',
  '締切から逆算したスケジュール表を作った',
  '設計図書・様式類を全てDLして保存した',
  '現物公開日をカレンダー登録した（事前連絡要否も）',
  '過去の同種案件の落札結果を検索した（相場観）'])

step_slide('STEP 2','参加資格・書類：門前払いを回避する','地域要件→欠格要件→書類期限の順に読む','最初に確認する3点',
 '① 地域要件：住民登録 or 本店/支店/営業所の所在地\n　→ KHDは花巻支店登記で岩手中部圏域クリア済み\n　→ 他地域案件は「支店登記の追加」も戦略になる\n② 欠格要件：破産/更生/税滞納/反社でないこと\n③ 添付書類の有効期限：発行後3ヶ月以内が通例\n　→ 締切から逆算して取得（早取りは失効リスク）\n\n法人申込の標準セット：\n全部事項証明書／印鑑登録証明書／市町村税納税証明書\n／誓約書／（代理人なら委任状＋身分証写し）',
 ['地域要件を自社の登記で満たすことを謄本で確認した',
  '税の未納なし（誓約書と矛盾しない）を確認した',
  '添付書類の取得日を締切から逆算して決めた',
  '様式を正確な商号・住所で下書きした',
  '不明点は質問書で事前照会した（期限内）',
  '提出は持参 or 簡易書留等の記録が残る方法にした',
  '受付期間内「必着」を確認した（消印有効ではない）'])

step_slide('STEP 3','机上査定：仕様書精読 × 相見積りで「売値」を作る','車台番号は概算段階では渡さない（本命1社のみ）','仕様書から読み取る査定因子',
 '・型式→グレード特定（MT/AT・バン/ワゴンで絞れる）\n・走行距離/初度登録＝国内価値。ランクル等は過走行\n　でも海外価値が残る（50万kmからが本番）\n・車検/自賠責の残＝自走可否・搬出コストに直結\n・特殊装備（赤色灯等）＝撤去・用途変更の手間\n\n相見積りの鉄則：\n・3〜5社、タイプの違う先（国内買取/旧車専門/輸出\n　/知人ルート）に当てる\n・提示額は相互に伏せる（アンカリング回避）',
 ['型式・年式・走行・装備からグレードを特定した',
  '排ガス規制(NOx・PM法)の適合を確認した ※重要',
  '査定先を3〜5社リストアップ（タイプ分散）した',
  '競合化リスクを精査した（地域要件を満たす業者は除外）',
  '概算査定で「保守/中央/強気」の3シナリオを作った',
  '査定根拠（メール・査定書）を証跡保存した',
  '想定vs実際の比較表を作り、取得のたび埋めている'])

step_slide('STEP 4','現物確認：目利きマップを現場で回す','60〜90分の動線。下回りに最長時間を割く','当日の動線',
 '① 書類照合(10分)：車台番号打刻・車検証・付属品\n② 外周(15分)：外装4方向→パネル隙間→ガラス→装備\n③ 下回り(20分)：フレーム後部→サイドレール→足回り\n　※最長時間。マット・ライト・軍手持参\n④ 機関(15分)：冷間始動→アイドル→排気色→漏れ\n⑤ 内装(10分)：ODO→クラッチ→床錆→装備作動\n⑥ 質問(10分)：装備撤去可否・整備記録・修理歴を\n　担当者に直接確認\n\n撮影は「撮影テンプレ20枚」の順で撮る',
 ['事前予約の電話を入れ、訪問時間を確定した',
  '持ち物を揃えた（ライト/軍手/マット/磁石/脚立）',
  '目利き12ポイントを順番どおり全て見た',
  '下回りは叩き＋目視で3ゾーン(A/B/C)判定した',
  '冷間始動と排気色を動画で記録した',
  '車台番号打刻と書類の一致を確認した',
  '写真20枚を撮り、当日中に査定先へ送付した'])

step_slide('STEP 5','収支確定：売値から逆算して入札上限を決める','入札上限＝(保守売却額−目標利益−諸費用)÷1.10','逆算式と本案件の当てはめ',
 '・契約金額＝入札額×1.10（消費税相当）が通例\n・諸費用＝陸送(0〜10万・業者引取無料なら0)\n　＋名義変更等(1〜3万)\n・保有なら＋車検整備(15〜40万)＋NOx・PM登録可否\n\n本案件（旧車王実査定100〜200万円）：\n・保守100万・利益30万狙い → 上限52万\n・中央150万 → 上限98万\n→ 推奨入札 55〜65万円／絶対上限80万円\n\n※現物確認の減点（下回り等）を必ず反映してから確定',
 ['売却3シナリオを実査定で裏付けた',
  '諸費用を業者別の実見積りで確定した',
  '入札締切1週間前までに全見積りを回収した',
  '「絶対上限」を先に決め、書面に残した',
  '端数を付けた入札額にした（同額くじ回避）',
  '資金（契約金額の全額前払い）を確保した',
  '現物確認の結果を入札額に反映した'])

step_slide('STEP 6','入札：一発勝負のルールを間違えない','無効事由に1つでも触れたら終わり','絶対に落とせない実務ルール',
 '・入札回数1回／提出後の書換・撤回不可\n・金額は税抜で記入（契約金額は×1.10）\n・「￥」記入・アラビア数字・訂正不可のペン\n・封筒に件名明記・封印\n・郵送は「必着」。簡易書留/特定記録で記録を残す\n\n無効になる典型：\n記名押印漏れ／金額訂正／指定様式以外／\n資本・人的関係のある複数者の入札（連合とみなし）',
 ['入札書は指定様式・訂正なしで記入した',
  '金額の頭に￥、桁ズレをダブルチェックした',
  '記名・押印（法人は代表者印）を確認した',
  '封筒に件名を明記し封印した',
  '提出期限の前日までに到達する手段で発送した',
  '開札日・結果公表の確認方法を控えた',
  '落札できなかった場合の次案件リストを用意した'])

step_slide('STEP 7','落札後：契約・全額前払い・搬出を期限内に','期限3本。1本落とすと落札決定取消もある','搬出の実務',
 '① 契約締結期限（例：落札から約1週間）\n② 代金全額前払い（納入通知書で契約締結までに）\n③ 搬出期限（例：月末・費用は落札者負担）\n\n・車検切れ＝公道自走不可 → 積載車を事前手配\n　（仮ナンバー自走は老朽車ではリスク高・非推奨）\n・買取業者への直送も検討：引取無料の業者なら\n　陸送費が丸ごと浮く＝実質利益+5〜10万\n・名義変更/抹消は速やかに（税・責任の切替）',
 ['契約・支払・搬出の3期限をカレンダー登録した',
  '納入通知書での支払方法・期日を確認した',
  '陸送 or 業者引取を確定し、搬出日を予約した',
  '売却先と引取場所・日程を握った（直送が最良）',
  '名義変更（または抹消）の段取りを決めた',
  '特殊装備の撤去有無・タイミングを確定した',
  '実費の領収書を全て保存した（比較表に反映）'])

step_slide('STEP 8','売却・関係構築・記録：次の案件につなげる','担当者の直通確保→名刺DB登録→次回は電話1本','売却と事業化',
 '・相見積りの最高値＋対応の質で売却先を決定\n・取引完了時に「今後も官公庁売払いを継続的に扱う」\n　と伝え、担当者の直通連絡先を確保\n・担当者は名刺DB/顧客マスターに登録\n\n事業化の法規：\n・古物商許可＝反復継続の転売に必須\n　（無許可営業は3年以下の懲役/100万円以下の罰金）\n・2台目に着手する前に取得完了（申請〜許可1〜2ヶ月）',
 ['本番査定（車台番号開示）は最有力1社に絞った',
  '売買契約書・振込記録を保存した',
  '担当者名・直通連絡先を名刺DBに登録した',
  '想定vs実際の比較表を完成させ、差異を分析した',
  '案件全体の学びをnotesに記録した',
  '古物商許可の取得状況を確認した（2台目の前提）',
  '次の売払い案件のウォッチを再開した'])

# ============ 18. 現地キット ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'現地キット','持ち物 ＋ 撮影構図テンプレ20枚（この順で撮ればそのまま査定に送れる）')
box(s,Inches(0.55),Inches(1.75),Inches(4.4),Inches(4.95),fill=BEIGE)
box(s,Inches(0.55),Inches(1.75),Inches(4.4),Inches(0.06),fill=RED)
txt(s,Inches(0.85),Inches(1.95),Inches(3.9),Inches(0.4),'持ち物リスト',size=13.5,color=DRED,bold=True)
checklist(s,Inches(0.85),Inches(2.45),Inches(3.8),[
 'LEDライト（下回り照射用）','軍手・汚れてよい服装','地面用マット or 段ボール',
 '磁石（パテ盛り判定用）','スマホ＋モバイルバッテリー','物件仕様書のコピー（照合用）',
 'メジャー（荷室採寸）','質問メモ（装備撤去可否等）'],size=11.5,gap=0.44)
box(s,Inches(5.25),Inches(1.75),Inches(7.45),Inches(4.95),fill=BEIGE)
box(s,Inches(5.25),Inches(1.75),Inches(7.45),Inches(0.06),fill=RED)
txt(s,Inches(5.55),Inches(1.95),Inches(6.8),Inches(0.4),'撮影構図テンプレ20枚',size=13.5,color=DRED,bold=True)
cols=[['① 前方全景（ナンバー入り）','② 後方全景（バッジ入り）','③ 左側面全景','④ 右側面全景','⑤ 左前45度','⑥ 右後45度','⑦ 下回り：後部フレーム','⑧ 下回り：左レール','⑨ 下回り：右レール','⑩ 下回り：前部クロス'],
      ['⑪ タイヤ4本（溝・製造年）','⑫ エンジンルーム全景','⑬ オイル周り接写','⑭ 冷間始動の動画＋排気','⑮ メーター（ODO読める解像度）','⑯ 運転席全景（シフト入り）','⑰ シート・内装の傷汚れ','⑱ 荷室＋付属タイヤ','⑲ 車台番号打刻の接写','⑳ 傷・錆・凹みの個別接写(全部)']]
x0=[Inches(5.55),Inches(9.15)]
for ci,col in enumerate(cols):
    y=Inches(2.45)
    for item in col:
        txt(s,x0[ci],y,Inches(3.5),Inches(0.35),item,size=11,color=INK)
        y=add(y,Inches(0.40))
footer(s)

# ============ 19. リスク法規 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'リスクと法規','先に知らないと即死する4つ','赤タグ＝重要度')
items=[('NOx・PM法','旧年式ディーゼルは首都圏・大阪・愛知等の対策地域で登録不可。「東京で保有」が物理的に不可能な車がある。\n転売(輸出)なら無関係。本案件は保有なら花巻登録が前提。','最重要'),
('現状渡し・責任なし','引渡し後の故障・瑕疵は全て自己責任。現物確認が唯一の防衛線。\n下回り腐食は売値を数十万単位で変える（配点40点の理由）。','高'),
('一発入札','書換・撤回・再入札なし。相場を外すと高値掴みか未落札。\n「絶対上限」を先に書面化してから入札書を書く。','高'),
('古物商許可','反復継続の転売は許可必須。1台の単発は直ちに違法ではないが、\n事業化するなら2台目までに取得完了（申請〜許可1〜2ヶ月）。','中')]
y=Inches(1.75)
for name,desc,tag in items:
    box(s,Inches(0.55),y,Inches(11.9),Inches(1.12),fill=BEIGE)
    box(s,Inches(0.55),y,Inches(0.1),Inches(1.12),fill=RED)
    txt(s,Inches(0.9),add(y,Inches(0.1)),Inches(2.5),Inches(0.9),name,size=14,color=DRED,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.6),add(y,Inches(0.1)),Inches(7.6),Inches(0.95),desc,size=11,color=INK,sp=1.2)
    box(s,Inches(11.35),add(y,Inches(0.36)),Inches(0.95),Inches(0.4),fill=LRED)
    txt(s,Inches(11.35),add(y,Inches(0.36)),Inches(0.95),Inches(0.4),tag,size=10.5,color=DRED,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    y=add(y,Inches(1.28))
footer(s)

# ============ 20. ケーススタディ ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'ケーススタディ','実例：岩手中部水道企業団 公告148号（2026）','確定情報のみ記載。数字の根拠は notes/ の各記録')
box(s,Inches(0.55),Inches(1.75),Inches(6.1),Inches(4.95),fill=BEIGE)
box(s,Inches(0.55),Inches(1.75),Inches(6.1),Inches(0.06),fill=RED)
txt(s,Inches(0.85),Inches(1.95),Inches(5.4),Inches(0.4),'案件データ',size=14,color=DRED,bold=True)
txt(s,Inches(0.85),Inches(2.45),Inches(5.55),Inches(4.1),
    '車両：ランクル100 バンVX ディーゼル5MT(KG-HDJ101K)\n　※VXバッジ・5MT・貨物登録の実写で確定\n平成10年6月初度登録／280,113km（ODO実写一致）\n車検切れ(R8.6.17)／公共応急作業車(8ナンバー)\n最低売却価格：15万円\n実査定：旧車王 概算100〜200万円（不成立・伏せて相見積り中）\n推奨入札：55〜65万円（絶対上限80万円）\n\n8/20現物確認 → 8/21書類締切 → 9/7入札 →\n9/8開札 → 9/15契約・支払 → 9/30搬出',size=11.5,color=INK,sp=1.3)
box(s,Inches(6.95),Inches(1.75),Inches(5.75),Inches(4.95),fill=LGREEN)
box(s,Inches(6.95),Inches(1.75),Inches(5.75),Inches(0.06),fill=GREEN)
txt(s,Inches(7.2),Inches(1.95),Inches(5.1),Inches(0.4),'この案件で確立した「型」',size=14,color=GREEN,bold=True)
txt(s,Inches(7.2),Inches(2.45),Inches(5.2),Inches(4.1),
    '・地域要件は支店登記でクリア（参入障壁を自分だけ\n　越える＝安く仕入れる構造）\n・役所の最低売却価格は市場価格と無関係。仕様書の\n　精読とプロ査定で「本当の値段」を掴む\n・買取業者は地域要件を満たせない＝情報を渡しても\n　入札で競合しない。ただし買い叩きには注意\n・車台番号は本命1社の本番査定まで温存\n・相見積りは提示額を相互に伏せる\n・全記録は .company/secretary/notes/ に集約。\n　想定vs実際の差異が次回案件の精度になる',size=11.5,color=INK,sp=1.3)
footer(s)

# ============ 21. 総覧 ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'マスターチェックリスト','全チェックリスト総覧（印刷してこの1枚で回せる）')
groups=[('S1 発掘','公告確認／勝ち筋判定／資格先読み／逆算表／資料DL／公開日登録／過去落札'),
('S2 資格','登記で要件確認／税未納なし／書類逆算取得／様式下書き／質問書／記録郵送／必着'),
('S3 査定','グレード特定／排ガス規制／3〜5社分散／競合精査／3シナリオ／証跡保存／比較表'),
('S4 現物','予約／持ち物／12ポイント／下回りABC／冷間始動動画／打刻照合／写真20枚'),
('S5 収支','3シナリオ裏付け／諸費用実額／1週間前回収／絶対上限書面化／端数／資金／現物反映'),
('S6 入札','指定様式／金額ダブルチェック／記名押印／封筒件名／前日到達／結果確認／次案件'),
('S7 決済','3期限登録／納入通知書／搬出予約／売却先と握り／名義変更／装備撤去／領収書'),
('S8 売却','本番査定1社／契約書保存／担当者DB登録／比較表完成／学び記録／古物商／次ウォッチ')]
y=Inches(1.75)
for name,desc in groups:
    box(s,Inches(0.55),y,Inches(11.9),Inches(0.56),fill=BEIGE)
    box(s,Inches(0.55),y,Inches(1.5),Inches(0.56),fill=RED)
    txt(s,Inches(0.55),y,Inches(1.5),Inches(0.56),name,size=11.5,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(2.25),y,Inches(10.1),Inches(0.56),desc,size=11,color=INK,anchor=MSO_ANCHOR.MIDDLE)
    y=add(y,Inches(0.63))
txt(s,Inches(0.55),add(y,Inches(0.06)),Inches(12),Inches(0.4),'詳細は各STEPページ。目利き12ポイント＝p3／下回り判定＝p4／撮影テンプレ20枚＝p18。',size=10,color=GRAY)
footer(s)

prs.save('/home/user/khd_workspace/KHD_車両転売_仕入売却マニュアル_v3_テンプレ準拠.pptx')
print('saved pages:',_pages['n']+1)

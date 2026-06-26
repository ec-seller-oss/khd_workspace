"""
クリニックDX「My AI」 v5 ── 参考デッキ「260526_医療専門不動産のビジネスモデル」に
デザイン・色を厳密一致（クリーム白×レンガ赤）。実物PDFから配色抽出。
BG #F9F6EF / RED #AA2E26 / INK #1A1A1A / LINE #DAD6CF / CARD #F1ECE1。
SSoT: 提案スライド完全設計書(Google Doc 1QfTxY6...)。保存先: Drive「260526_AI医療コンサル」固定。
出力: clinic_dx_v5.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 参考デッキから抽出した実配色 ──
BG     = RGBColor(0xF9, 0xF6, 0xEF)   # クリーム背景
RED    = RGBColor(0xAA, 0x2E, 0x26)   # レンガ赤・主アクセント
REDD   = RGBColor(0x8C, 0x24, 0x1D)   # 濃赤
INK    = RGBColor(0x1A, 0x1A, 0x1A)   # 本文黒
GRY    = RGBColor(0x6E, 0x6E, 0x6E)   # サブ灰
LINE   = RGBColor(0xDA, 0xD6, 0xCF)   # 罫線
CARD   = RGBColor(0xF1, 0xEC, 0xE1)   # ベージュカード
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)   # カード枠
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)   # 薄赤（強調セル）
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)   # 薄灰ベージュ
WHT    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Hiragino Sans"
_MK = "/Users/kikuchikenta/01_honbu_docs_automation/myai_mockups"
IMG1 = _MK + "/screen1_doc.png"
IMG2 = _MK + "/screen2_dash.png"
IMG3 = _MK + "/screen3_line.png"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, eyebrow, main, sub=""):
    """参考デッキ準拠：赤のスモールキャップ＋赤下線→黒い主題→灰サブ。"""
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "クリニックDX「My AI」  ｜  AI医療コンサル", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)


def light_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = REDD if is_hi else RED
            else:
                cell.fill.fore_color.rgb = REDBG if is_hi else (CARD if ri % 2 == 1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = WHT
                    elif is_hi:
                        r.font.color.rgb = RED
                    elif ci == 0:
                        r.font.color.rgb = INK
                    else:
                        r.font.color.rgb = RGBColor(0x3A,0x3A,0x3A)
    return tb


def myai_mock(slide, x, y, w, h):
    """MyAI 紹介状ジェネレーター（白基調・赤アクセント）。"""
    bx(slide, x, y, w, h, CARD, line=CARDLN, lw=1.0)
    bx(slide, x, y, w, Inches(0.06), RED)
    bx(slide, x+Inches(0.18), y+Inches(0.2), w-Inches(0.36), Inches(0.42), RED)
    t(slide, "マイAI ｜ 紹介状ジェネレーター", x+Inches(0.32), y+Inches(0.21), w-Inches(0.5), Inches(0.4),
      sz=11, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
    bx(slide, x+Inches(0.18), y+Inches(0.74), w-Inches(0.36), Inches(0.5), GRYBG)
    t(slide, "カルテから抽出中… 患者ID / 主訴 / 既往歴 / 処方", x+Inches(0.3), y+Inches(0.76),
      w-Inches(0.5), Inches(0.46), sz=9.5, col=GRY, anchor=MSO_ANCHOR.MIDDLE)
    bx(slide, x+Inches(0.18), y+Inches(1.36), w-Inches(0.36), h-Inches(2.1), WHT, line=CARDLN, lw=0.75)
    t(slide, "AI生成プレビュー", x+Inches(0.34), y+Inches(1.44), w-Inches(0.6), Inches(0.3), sz=9, bold=True, col=RED)
    for i, ln in enumerate(["──────────────────",
                            "拝啓  時下ますますご清栄のことと…",
                            "下記の患者をご紹介申し上げます。",
                            "診断: ＿＿＿  /  所見: ＿＿＿＿",
                            "──────────────────"]):
        t(slide, ln, x+Inches(0.34), y+Inches(1.74)+Emu(int(Inches(0.28))*i), w-Inches(0.66), Inches(0.28),
          sz=9, col=GRY)
    bx(slide, x+Inches(0.18), y+h-Inches(0.6), w-Inches(0.36), Inches(0.42), RED)
    t(slide, "⚡ 8分 → 数秒で下書き完了", x+Inches(0.18), y+h-Inches(0.6), w-Inches(0.36), Inches(0.42),
      sz=11, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def dash_mock(slide, x, y, w, h):
    bx(slide, x, y, w, h, CARD, line=CARDLN, lw=1.0)
    bx(slide, x, y, w, Inches(0.06), RED)
    t(slide, "📊 削減効果ダッシュボード", x+Inches(0.25), y+Inches(0.18), w-Inches(0.4), Inches(0.35), sz=11, bold=True, col=RED)
    bars = [("問診", 0.70), ("電話", 0.55), ("書類", 0.62), ("レセプト", 0.90)]
    base_x = x+Inches(0.3); base_y = y+Inches(0.7); bw = Inches(0.55); maxh = Inches(1.5); gap = Inches(0.45)
    for i, (lab, v) in enumerate(bars):
        bxx = base_x + (bw+gap)*i
        bh = Emu(int(maxh*v))
        bx(slide, bxx, base_y+maxh-bh, bw, bh, RED if i == 3 else RGBColor(0xC9,0x6A,0x62))
        t(slide, str(int(v*100))+"%", bxx-Inches(0.1), base_y+maxh-bh-Inches(0.28), bw+Inches(0.2), Inches(0.26), sz=9, bold=True, col=INK, align=PP_ALIGN.CENTER)
        t(slide, lab, bxx-Inches(0.15), base_y+maxh+Inches(0.04), bw+Inches(0.3), Inches(0.26), sz=9, col=GRY, align=PP_ALIGN.CENTER)


def line_mock(slide, x, y, w, h):
    bx(slide, x, y, w, h, RGBColor(0x06,0xC7,0x55))
    bx(slide, x+Inches(0.12), y+Inches(0.12), w-Inches(0.24), h-Inches(0.24), WHT)
    t(slide, "LINE 問診・予約", x+Inches(0.28), y+Inches(0.2), w-Inches(0.5), Inches(0.3), sz=10, bold=True, col=RGBColor(0x06,0x7A,0x35))
    bubbles = [("来院前にスマホで問診✓", True), ("予約日をお選びください", False), ("○/○ 10:00 で予約完了", True)]
    by = y+Inches(0.6)
    for txt, right in bubbles:
        bw = w-Inches(1.0)
        bxx = x+(w-bw-Inches(0.28)) if right else x+Inches(0.28)
        bx(slide, bxx, by, bw, Inches(0.4), RGBColor(0x9A,0xE5,0xB0) if right else GRYBG)
        t(slide, txt, bxx+Inches(0.1), by+Inches(0.02), bw-Inches(0.2), Inches(0.36), sz=8.5, col=INK, anchor=MSO_ANCHOR.MIDDLE)
        by = by + Inches(0.48)


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "AI × CLINIC DX", Inches(0.9), Inches(1.7), Inches(7.2), Inches(0.45), sz=15, bold=True, col=RED)
t(s, "クリニックのAI・LINE導入", Inches(0.88), Inches(2.25), Inches(7.4), Inches(0.9), sz=40, bold=True, col=INK)
t(s, "まるごとおまかせ", Inches(0.88), Inches(3.05), Inches(7.4), Inches(0.9), sz=40, bold=True, col=RED)
t(s, "独自カスタマイズの「マイAI」×「Lステップ」×AI音声(IVR)で実現する、\n初期費用0円・完全成果報酬型のクリニックDX。",
  Inches(0.9), Inches(4.05), Inches(7.3), Inches(0.9), sz=13.5, col=GRY, line_sp=1.25)
# オファー3カード
ox, oy, ow, og = Inches(0.9), Inches(5.15), Inches(2.3), Inches(0.16)
offers = [("初期費用", "0円"), ("月額基本料", "0円"), ("成果報酬", "30%")]
for i, (lab, val) in enumerate(offers):
    cx = ox + (ow + og) * i
    bx(s, cx, oy, ow, Inches(1.3), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, oy, ow, Inches(0.06), RED)
    t(s, lab, cx, oy+Inches(0.18), ow, Inches(0.35), sz=12, col=GRY, align=PP_ALIGN.CENTER)
    t(s, val, cx, oy+Inches(0.48), ow, Inches(0.7), sz=34, bold=True, col=RED, align=PP_ALIGN.CENTER)
bx(s, Inches(0.9), Inches(6.78), Inches(11.5), Pt(1.2), LINE)
t(s, "チームてっかん  ｜  菊池 研太  ｜  @khd_medical01", Inches(0.9), Inches(6.9), Inches(11), Inches(0.4), sz=13, bold=True, col=INK)
# 右：MyAI画面（実写真）
bx(s, Inches(8.45), Inches(1.7), Inches(4.5), Inches(3.45), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(8.45), Inches(1.7), Inches(4.5), Inches(0.06), RED)
s.shapes.add_picture(IMG1, Inches(8.7), Inches(2.1), width=Inches(4.0))
t(s, "実際のMyAI画面：カルテ → 紹介状を数秒で自動下書き", Inches(8.45), Inches(5.28), Inches(4.5), Inches(0.5), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)

# ════════ SLIDE 2 — 4大課題 ════════
s = sl(); ft(s)
hdr(s, "THE PROBLEM", "現場はもう限界では？ ── 御院でも起きていませんか？", "ネット上の生の声が示す、クリニックのリアルな4大お悩み")
issues = [
    ("01", "採用しても、定着しない", "70万〜150万円", "医療事務の採用単価。保険制度の複雑さで即戦力にならず早期離職の悪循環。"),
    ("02", "カルテ入力で患者と目が合わない", "診察時間の半分", "が電子カルテ入力。診療後も紹介状・診断書作成で2時間超の残業。"),
    ("03", "月初10日のレセプトが毎月重い", "月10〜20時間", "の残業が常態化。目視チェックによる算定漏れの不安も重なる。"),
    ("04", "鳴り止まない電話で受付がパンク", "月1,200件超", "の電話対応。「今空いてますか？」で目の前の患者対応が滞る。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.4), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.85)
for i, (no, ti, num, desc) in enumerate(issues):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    t(s, no, cx+Inches(0.3), cy+Inches(0.2), Inches(1.0), Inches(0.6), sz=28, bold=True, col=RED)
    t(s, ti, cx+Inches(1.05), cy+Inches(0.28), cw-Inches(1.25), Inches(0.6), sz=16, bold=True, col=INK)
    t(s, num, cx+Inches(0.32), cy+Inches(1.0), cw-Inches(0.6), Inches(0.7), sz=29, bold=True, col=RED)
    t(s, desc, cx+Inches(0.34), cy+Inches(1.7), cw-Inches(0.6), Inches(0.6), sz=11.5, col=GRY, line_sp=1.1)

# ════════ SLIDE 3 — 自動化マップ ════════
s = sl(); ft(s)
hdr(s, "THE SOLUTION", "既存の業務をどう変えるか？ ── 完全自動化マップ", "医者の悩み（左）を、AI商品（右）がどう解決するか。クリニック全体を1つのエコシステムで最適化")
rows = [
    ("業務領域", "現状のアナログ業務", "マイAI・Lステップ・IVR導入後", "削減効果"),
    ("問診・予約", "電話受付＋紙問診をカルテへ手入力", "Lステップで来院前にスマホ上で完結", "問診業務 40〜70%減"),
    ("電話対応", "全件スタッフが受話", "AI音声(IVR)が一次対応・LINE誘導", "電話件数 最大40〜80%減"),
    ("書類作成", "過去カルテを探しつつ手作業", "自院学習「マイAI」が下書き自動生成", "書類時間 月30h以上減"),
    ("レセプト点検", "月末月初に目視で残業", "AIが4,000ルールで算定漏れ自動チェック", "点検時間 最大1/20"),
]
light_table(s, rows, Inches(0.55), Inches(1.85), Inches(12.23), Inches(4.7),
            [Inches(2.0), Inches(3.85), Inches(4.05), Inches(2.33)], hi_col=3, sz=12.5, header_sz=13)

# ════════ SLIDE 4 — MyAI画面 ════════
s = sl(); ft(s)
hdr(s, "PRODUCT", '"使う場面"が見える ── 実際のMyAI画面', "院長・スタッフ・患者、それぞれの画面で業務が変わる（画面はデモ・数値は例示）")
# 左：書類アシスト（実写真・大）
bx(s, Inches(0.55), Inches(1.85), Inches(5.7), Inches(4.7), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(0.55), Inches(1.85), Inches(5.7), Inches(0.06), RED)
s.shapes.add_picture(IMG1, Inches(0.8), Inches(2.2), width=Inches(5.2))
t(s, "① 書類アシスト｜カルテ→紹介状を数秒で下書き（院長の画面）", Inches(0.75), Inches(5.95), Inches(5.4), Inches(0.4), sz=11, bold=True, col=RED)
# 右上：ダッシュボード（実写真）
bx(s, Inches(6.45), Inches(1.85), Inches(6.33), Inches(2.75), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(6.45), Inches(1.85), Inches(6.33), Inches(0.06), RED)
s.shapes.add_picture(IMG2, Inches(8.05), Inches(2.12), width=Inches(3.95))
t(s, "② 成果ダッシュボード｜削減効果を可視化", Inches(6.45), Inches(4.66), Inches(6.33), Inches(0.3), sz=11, bold=True, col=RED)
# 右下：LINE問診（実写真・縦長）
bx(s, Inches(6.45), Inches(5.1), Inches(6.33), Inches(1.45), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(6.45), Inches(5.1), Inches(6.33), Inches(0.06), RED)
s.shapes.add_picture(IMG3, Inches(6.75), Inches(5.22), height=Inches(1.2))
t(s, "③ LINE問診・予約｜Lステップ連携（患者スマホ）", Inches(7.65), Inches(5.1), Inches(5.0), Inches(1.45), sz=11, bold=True, col=RED, anchor=MSO_ANCHOR.MIDDLE)

# ════════ SLIDE 5 — 実証データ ════════
s = sl(); ft(s)
hdr(s, "EVIDENCE", "「本当に自院でも変わるのか？」を証明する実証データ", "※当社実績ではなく、公開されている医療機関の導入実績です（出典明記）")
cases = [
    ("問診・予約のDX", "内科クリニック", "電話の嵐＋紙問診の手入力で受付過多", "Lステップで来院前にスマホ問診完了。電話を一次自動対応",
     ["初月にLINE予約 350件", "受付問診 70%削減", "月商 約1,500万円規模を押上げ"], "出典: Lステップ公開事例「そのだ内科」ほか"),
    ("書類作成のDX", "病院・総合クリニック", "紹介状・退院サマリ作成に医師が忙殺", "自院データを学習した生成AIが下書き自動生成",
     ["診断書 月400件を 50%削減", "退院サマリ 28分→8分(7割減)", "医師業務 月30h以上削減"], "出典: 新古賀/名古屋医療C/戸畑共立"),
    ("電話業務のDX", "地域密着クリニック", "インフル予約・時間外問合せで電話過多", "AI音声が一次対応、営業電話を自動ブロック",
     ["月1,200件の電話を自動化", "うち約400件(40%)を削減", "対応時間 全体70%削減"], "出典: IVRy 公開導入事例"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.85)
for i, (cat, who, bf, af, results, src) in enumerate(cases):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(4.75), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.06), RED)
    t(s, cat, cx, y0+Inches(0.18), cw, Inches(0.4), sz=13, bold=True, col=RED, align=PP_ALIGN.CENTER)
    t(s, who, cx, y0+Inches(0.62), cw, Inches(0.3), sz=11, bold=True, col=GRY, align=PP_ALIGN.CENTER)
    bx(s, cx+Inches(0.2), y0+Inches(1.0), cw-Inches(0.4), Inches(0.75), GRYBG)
    t(s, "Before", cx+Inches(0.32), y0+Inches(1.03), Inches(1.2), Inches(0.26), sz=9.5, bold=True, col=GRY)
    t(s, bf, cx+Inches(0.32), y0+Inches(1.26), cw-Inches(0.64), Inches(0.46), sz=10.5, col=INK, line_sp=1.05)
    bx(s, cx+Inches(0.2), y0+Inches(1.83), cw-Inches(0.4), Inches(0.75), REDBG)
    t(s, "After", cx+Inches(0.32), y0+Inches(1.86), Inches(1.2), Inches(0.26), sz=9.5, bold=True, col=RED)
    t(s, af, cx+Inches(0.32), y0+Inches(2.09), cw-Inches(0.64), Inches(0.46), sz=10.5, col=INK, line_sp=1.05)
    bx(s, cx+Inches(0.2), y0+Inches(2.68), cw-Inches(0.4), Inches(1.65), WHT, line=CARDLN, lw=0.75)
    bx(s, cx+Inches(0.2), y0+Inches(2.68), Inches(0.08), Inches(1.65), RED)
    t(s, "成果", cx+Inches(0.38), y0+Inches(2.72), Inches(1.2), Inches(0.3), sz=10, bold=True, col=RED)
    for j, rr in enumerate(results):
        t(s, "● " + rr, cx+Inches(0.38), y0+Inches(3.04)+Emu(int(Inches(0.4))*j), cw-Inches(0.6), Inches(0.4), sz=11, bold=True, col=INK, line_sp=1.0)
    t(s, src, cx, y0+Inches(4.4), cw, Inches(0.35), sz=8.5, col=GRY, align=PP_ALIGN.CENTER)

# ════════ SLIDE 6 — ROI早見表 ════════
s = sl(); ft(s)
hdr(s, "SIMULATION", "御院の規模なら、いくら残る？ ── 削減効果の早見表", "1モデルではなく、自院を当てはめられる規模別シミュレーション")
rows = [
    ("院の規模", "月の削減額", "成果報酬(30%)", "院に残る効果（年額）"),
    ("医師1・事務2", "約16万円", "4.8万円", "約134万円"),
    ("医師1・事務3（標準）", "約24万円", "7.2万円", "約202万円"),
    ("医師2・事務5", "約40万円", "12万円", "約336万円"),
]
light_table(s, rows, Inches(0.8), Inches(1.85), Inches(11.73), Inches(2.5),
            [Inches(3.4), Inches(2.5), Inches(2.6), Inches(3.23)], hi_col=3, sz=15, header_sz=13)
bx(s, Inches(0.8), Inches(4.7), Inches(11.73), Inches(2.0), CARD, line=CARDLN, lw=1.0)
t(s, "標準モデルの内訳（医師1・事務3 ／ 事務人件費 月83万円・総480時間・時間単価 約1,729円）", Inches(1.0), Inches(4.88), Inches(11.3), Inches(0.4), sz=12.5, bold=True, col=RED)
t(s, "問診 −42h ＋ レセプト −45h ＋ 電話 −40h ＋ 書類 −12h ＝ 月139時間削減（フルタイム約1名分の余力創出）", Inches(1.0), Inches(5.36), Inches(11.3), Inches(0.4), sz=13, col=INK)
t(s, "→ 月139h × 1,729円 ＝ 月 約24万円削減 → 年 約288万円。成果報酬を引いて、院に残る効果は 約202万円/年。", Inches(1.0), Inches(5.84), Inches(11.3), Inches(0.4), sz=13, bold=True, col=REDD)
t(s, "※「うちは事務3人だから…年200万か」と、その場で自院に当てはめられます。", Inches(1.0), Inches(6.32), Inches(11.3), Inches(0.4), sz=11, col=GRY)

# ════════ SLIDE 7 — 競合比較表（宮崎修正：当社を左に強調／競合2列を統合／0円特大）════════
s = sl(); ft(s)
hdr(s, "WHY RISK-ZERO", "当サービスだけ、唯一のノーリスク ── 業界価格の真実", "他社は効果が出なくても固定費が発生。御院の持ち出しは“ゼロ”から始められる")
rows = [
    ("比較項目", "当サービス（中立伴走）", "一般IT・DX業者", "Lステップ構築代行"),
    ("提供内容", "Lステップ×マイAI 全体最適化", "自社アプリ販売／高額導入支援", "LINE・Lステップ構築"),
    ("初期費用", "0円", "20万〜300万", "50〜100万超"),
    ("月額費用", "0円（基本料なし）", "1.5〜300万（固定）", "5〜20万（固定）"),
    ("報酬体系", "削減人件費の30%のみ", "稼働問わず固定課金", "固定保守料"),
    ("運用定着", "3ヶ月伴走で定着まで", "自力運用 or 別料金で高額", "現場運用に不慣れ"),
    ("院側のリスク", "ゼロ（効果0なら0円）", "高〜極めて高", "中"),
]
tb7 = light_table(s, rows, Inches(0.55), Inches(1.8), Inches(12.23), Inches(4.95),
                  [Inches(2.2), Inches(4.6), Inches(3.0), Inches(2.43)], hi_col=1, sz=12.5, header_sz=13)
# 「0円」を特大に（宮崎指摘）
def _big(cell, sz):
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(sz); r.font.bold = True; r.font.color.rgb = RED; r.font.name = FONT
_big(tb7.cell(2, 1), 30)   # 初期費用 0円
_big(tb7.cell(3, 1), 20)   # 月額費用 0円（基本料なし）

# ════════ SLIDE 8 — 3ステップ ════════
s = sl(); ft(s)
hdr(s, "HOW IT WORKS", "「うちでも本当に測定できる？」── リスクゼロの3ステップ", "成果報酬への最大の懸念=「測定の手間」は、コンサル側が巻き取ります")
steps = [
    ("STEP 1", "事前の無料コンサル", "貴院の負担ゼロ", "当方が競合分析と類似事例を持ち込み、業務フローとボトルネックを洗い出し。\n\n※「2週間時間を計測」等の煩雑な作業は一切要求しません。", None),
    ("STEP 2", "フェアな測定ルール", "双方納得の透明基準", "導入前後で同等の患者数帯に基づき削減時間を算出。人員変更月は除外。\n\n※最低契約期間以降は30日前通知で解約可。ロックインなし。", None),
    ("STEP 3", "手元に利益が残る", "リスクは私たちが取る", "【例：月24万円分を削減】\n削減額24万 × 30% ＝ お支払い 7.2万円", ("貴院に残る効果", "月 16.8万円", "（年 約200万円）／効果0なら支払い0円")),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.75)
CARDH = Inches(3.5)
for i, (st, ti, sub, body, hi) in enumerate(steps):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, CARDH, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.8), RED)
    t(s, st, cx, y0+Inches(0.08), cw, Inches(0.42), sz=19, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, sub, cx, y0+Inches(0.5), cw, Inches(0.28), sz=10.5, col=RGBColor(0xF2,0xD8,0xD6), align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.25), y0+Inches(0.92), cw-Inches(0.5), Inches(0.5), sz=16, bold=True, col=INK, align=PP_ALIGN.CENTER)
    if hi:
        t(s, body, cx+Inches(0.3), y0+Inches(1.42), cw-Inches(0.6), Inches(0.6), sz=11, col=GRY, line_sp=1.1)
        bx(s, cx+Inches(0.25), y0+Inches(2.0), cw-Inches(0.5), Inches(1.35), REDBG, line=RED, lw=1.5)
        t(s, hi[0], cx+Inches(0.25), y0+Inches(2.08), cw-Inches(0.5), Inches(0.3), sz=11, bold=True, col=RED, align=PP_ALIGN.CENTER)
        t(s, hi[1], cx+Inches(0.25), y0+Inches(2.36), cw-Inches(0.5), Inches(0.6), sz=25, bold=True, col=RED, align=PP_ALIGN.CENTER)
        t(s, hi[2], cx+Inches(0.3), y0+Inches(2.98), cw-Inches(0.6), Inches(0.4), sz=9.5, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.0)
    else:
        t(s, body, cx+Inches(0.3), y0+Inches(1.42), cw-Inches(0.6), Inches(1.95), sz=11.5, col=GRY, line_sp=1.18)
# 下の余白に追記（宮崎指摘）：結論バンド
band_y = y0 + CARDH + Inches(0.3)
bx(s, Inches(0.55), band_y, Inches(12.23), Inches(0.95), REDBG)
bx(s, Inches(0.55), band_y, Inches(0.1), Inches(0.95), RED)
t(s, "つまり ── 御院の持ち出しはゼロ。効果が出た“分だけ”、私たちはいただく。", Inches(0.85), band_y+Inches(0.1), Inches(11.6), Inches(0.4), sz=15, bold=True, col=REDD)
t(s, "✓ IT導入補助金サポート　　✓ ロックインなし（30日前通知で解約可）　　✓ 使いこなすまで3ヶ月伴走", Inches(0.85), band_y+Inches(0.54), Inches(11.6), Inches(0.35), sz=11.5, bold=True, col=INK)

# ════════ SLIDE 9 — 選ばれる理由 ════════
s = sl(); ft(s)
hdr(s, "OUR VALUE", "なぜ私たちは「完全成果報酬」で伴走できるのか？", "「なぜ他社は固定費で、御社は成果報酬でやれるのか？」への答え")
reasons = [
    ("中立選定", "特定メーカーに縛られない", "自社ツールの売り込みではない。Lステップ/マイAI/IVRから貴院の診療科・規模に最適な組合せを厳選。"),
    ("実践ノウハウ", "自社で泥臭く運用した知見", "ツールは「設計」より「運用」でつまずく。自社実践で現場の抵抗感・エラー回避まで熟知。"),
    ("一気通貫", "補助金申請＋定着まで", "IT導入補助金・医療情報化支援基金の活用を支援し初期コスト最小化。使いこなすまで3ヶ月伴走。"),
    ("リスクゼロ", "自信があるからこその設計", "事前分析で「本当に削減できる」確証と成功事例があるから、院にリスクを負わせないモデルが成立。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.4), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.85)
for i, (ti, sub, body) in enumerate(reasons):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, cw, Inches(0.06), RED)
    t(s, str(i+1), cx+Inches(0.35), cy+Inches(0.3), Inches(0.9), Inches(0.9), sz=34, bold=True, col=RED)
    t(s, ti, cx+Inches(1.3), cy+Inches(0.3), cw-Inches(1.5), Inches(0.5), sz=19, bold=True, col=INK)
    t(s, sub, cx+Inches(1.3), cy+Inches(0.85), cw-Inches(1.5), Inches(0.4), sz=12, bold=True, col=RED)
    t(s, body, cx+Inches(0.4), cy+Inches(1.4), cw-Inches(0.7), Inches(0.95), sz=12, col=GRY, line_sp=1.15)

# ════════ SLIDE 9.5 — 安心・セキュリティ（李牧追加：最大の穴を先回り）════════
s = sl(); ft(s)
hdr(s, "SECURITY & TRUST", "ご安心ください ── 患者情報・セキュリティ・運用の不安に先回り", "医療現場の不安に、導入前から“運用ルール”で答えを用意します")
trust = [
    ("患者情報の扱い", "目的外利用なし", "院内データはクローズドに扱い、目的外利用なし。匿名化・最小限の原則を前提に設計し、運用ルールを最初に書面で取り決めます。"),
    ("アクセス管理", "誰が・いつ・何を", "役割別の権限と利用ログを整備。最終的な運用ルールは院長が決定。「ブラックボックスにしない」が原則。"),
    ("スタッフ負担を最小化", "段階導入＋伴走", "一気に変えない。優先業務から段階導入し、使いこなすまで3ヶ月伴走。マニュアル＋チャットで現場を支えます。"),
    ("やめたい時はやめられる", "ロックインなし", "最低契約期間以降は30日前通知で解約可。成果ゼロなら費用ゼロ。囲い込みはしません。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.4), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.85)
for i, (ti, tag, body) in enumerate(trust):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, cw, Inches(0.06), RED)
    t(s, "✓", cx+Inches(0.35), cy+Inches(0.28), Inches(0.7), Inches(0.7), sz=30, bold=True, col=RED)
    t(s, ti, cx+Inches(1.15), cy+Inches(0.3), cw-Inches(1.35), Inches(0.5), sz=18, bold=True, col=INK)
    t(s, tag, cx+Inches(1.15), cy+Inches(0.85), cw-Inches(1.35), Inches(0.4), sz=12, bold=True, col=RED)
    t(s, body, cx+Inches(0.4), cy+Inches(1.4), cw-Inches(0.7), Inches(0.95), sz=11.5, col=GRY, line_sp=1.15)

# ════════ SLIDE 10 — CTA / クロージング ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "ご相談はお気軽に", Inches(0.9), Inches(0.7), Inches(11), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "リスクゼロで、次世代のクリニック運営を始めませんか？", Inches(0.9), Inches(1.15), Inches(11.7), Inches(0.8), sz=28, bold=True, col=INK)
t(s, "押し売りは一切なし。単なる問い合わせでなく、1時間の無料コンサルセッションとしてご利用ください。",
  Inches(0.92), Inches(1.95), Inches(11.5), Inches(0.5), sz=14, col=GRY)
items = [
    ("ネット上の競合分析", "貴院周辺クリニックのデジタル対応状況を分析"),
    ("現状フローの整理", "電話の多さ・書類の手間などボトルネックをその場で特定"),
]
for i, (ti, ds) in enumerate(items):
    cy = Inches(2.8) + Inches(1.05) * i
    bx(s, Inches(0.9), cy, Inches(6.7), Inches(0.92), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.9), cy, Inches(0.06), Inches(0.92), RED)
    t(s, "✓ " + ti, Inches(1.15), cy+Inches(0.12), Inches(6.3), Inches(0.4), sz=15, bold=True, col=RED)
    t(s, ds, Inches(1.15), cy+Inches(0.48), Inches(6.3), Inches(0.4), sz=11.5, col=GRY)
bx(s, Inches(0.9), Inches(5.0), Inches(6.7), Inches(1.75), CARD, line=CARDLN, lw=1.0)
t(s, "お申込み（その場でQR／用紙）", Inches(1.15), Inches(5.15), Inches(6.3), Inches(0.4), sz=14, bold=True, col=RED)
for i, it in enumerate(["クリニック名 ／ ご担当者名", "電話番号 ／ メールアドレス", "ご希望の面談日時（オンライン対応可）"]):
    t(s, "▢  " + it, Inches(1.25), Inches(5.58)+Inches(0.38)*i, Inches(6.2), Inches(0.38), sz=12.5, col=INK)
# 特典（赤地白文字）
bx(s, Inches(7.9), Inches(2.8), Inches(4.55), Inches(3.95), RED)
t(s, "★ 参加特典", Inches(7.9), Inches(3.05), Inches(4.55), Inches(0.5), sz=16, bold=True, col=WHT, align=PP_ALIGN.CENTER)
t(s, "貴院専用\n業務削減シミュレーション\n報告書（診断書）", Inches(7.9), Inches(3.65), Inches(4.55), Inches(1.6), sz=22, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.1)
t(s, "を無料で作成・プレゼント", Inches(7.9), Inches(5.35), Inches(4.55), Inches(0.5), sz=14, col=RGBColor(0xF2,0xD8,0xD6), align=PP_ALIGN.CENTER)
bx(s, Inches(8.2), Inches(6.0), Inches(3.95), Inches(0.5), WHT)
t(s, "今回は丁寧に伴走するため【限定 ◯ 名】様まで", Inches(8.2), Inches(6.0), Inches(3.95), Inches(0.5), sz=11, bold=True, col=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save("clinic_dx_v8.pptx")
print("saved clinic_dx_v8.pptx  /  slides:", len(prs.slides._sldIdLst))

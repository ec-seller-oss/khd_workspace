"""
KHD 経営操縦席 操作マニュアル ── 菊池が「見て操作できる」1冊。
デザインは house style（クリーム白×レンガ赤）に統一。
出力: cockpit_manual.pptx
SSoT: 操縦席本体 1ofLJOFuW5175OHDSy0J0561LzrfAU-M015o7ByGH5wc / memory project_cockpit_rebuild
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
BLU    = RGBColor(0x2E, 0x5A, 0x8C)   # 自動=青系
BLUBG  = RGBColor(0xE4, 0xEC, 0xF4)
GRN    = RGBColor(0x3B, 0x6E, 0x3B)   # 入力=緑系
GRNBG  = RGBColor(0xE4, 0xF0, 0xE4)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Hiragino Sans"
W = Inches(13.33); H = Inches(7.5)

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp: p.line_spacing = line_sp
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb

def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def rbx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.34), sz=11.5, col=GRY)

def ft(slide, n):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "KHD 経営操縦席  ｜  操作マニュアル", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)
    t(slide, str(n), Inches(12.4), H-Inches(0.42), Inches(0.5), Inches(0.32), sz=9, col=GRY, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, body, accent=RED, accbg=REDBG, tsz=13, bsz=11):
    bx(slide, x, y, w, h, WHT, line=CARDLN, lw=0.75)
    bx(slide, x, y, Inches(0.09), h, accent)
    t(slide, title, x+Inches(0.25), y+Inches(0.14), w-Inches(0.4), Inches(0.4), sz=tsz, bold=True, col=accent)
    t(slide, body, x+Inches(0.25), y+Inches(0.6), w-Inches(0.45), h-Inches(0.7), sz=bsz, col=INK, line_sp=1.15)

# ════════════════════════════════ 1. 表紙 ════════════════════════════════
s = sl()
bx(s, 0, 0, W, H, BG)
bx(s, Inches(0.0), Inches(2.55), W, Inches(0.06), RED)
t(s, "KHD ｜ 経営の操縦席", Inches(0.9), Inches(1.6), Inches(11), Inches(0.5), sz=15, bold=True, col=RED)
t(s, "経営操縦席 操作マニュアル", Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.0), sz=40, bold=True, col=INK)
t(s, "あなたは「入力」と「報告」だけ。あとは全部つながって、数字が動く。",
  Inches(0.92), Inches(3.9), Inches(11), Inches(0.5), sz=15, col=GRY)
t(s, "1枚のスプレッドシート ／ 入力3枚 + 見るだけ3枚 ／ 古田土会計 × マーケKPI",
  Inches(0.92), Inches(4.5), Inches(11), Inches(0.4), sz=12, col=RED, bold=True)
t(s, "2026-06-06 版", Inches(0.92), Inches(6.4), Inches(5), Inches(0.4), sz=11, col=GRY)

# ════════════════════════════════ 2. 全体像 ════════════════════════════════
s = sl()
hdr(s, "OVERVIEW", "操縦席の全体像 — 6つのタブ", "左の3枚に入れる → 右の3枚が自動で出る。これだけ覚えればいい。")
# 入力3枚
t(s, "◀ 入力する 3枚（あなたが触る）", Inches(0.6), Inches(1.95), Inches(6), Inches(0.4), sz=14, bold=True, col=GRN)
ins = [("02_作業DB", "毎日の行動・報告を1行ずつ"),
       ("顧客マスター", "人（先生・客）を番号で蓄積"),
       ("03_売上見込み", "案件＋今日の追客本数")]
for i,(ti,bo) in enumerate(ins):
    y = Inches(2.4 + i*1.35)
    card(s, Inches(0.6), y, Inches(5.5), Inches(1.15), ti, bo, accent=GRN, accbg=GRNBG, tsz=14, bsz=12)
# 矢印
t(s, "→", Inches(6.25), Inches(3.6), Inches(0.8), Inches(0.8), sz=40, bold=True, col=RED, align=PP_ALIGN.CENTER)
# 見る3枚
t(s, "▶ 見るだけ 3枚（自動で動く）", Inches(7.2), Inches(1.95), Inches(6), Inches(0.4), sz=14, bold=True, col=BLU)
outs = [("01_統合司令塔", "過去/現在/未来＋今日の一手"),
        ("04_PL ＋ 05_資金繰り", "利益の形＋現金の谷（ランウェイ）"),
        ("06_資産負債BS", "資産・負債・純資産")]
for i,(ti,bo) in enumerate(outs):
    y = Inches(2.4 + i*1.35)
    card(s, Inches(7.2), y, Inches(5.5), Inches(1.15), ti, bo, accent=BLU, accbg=BLUBG, tsz=14, bsz=12)
ft(s, 2)

# ════════════════════════════════ 3. 毎日のリズム ════════════════════════════════
s = sl()
hdr(s, "DAILY", "毎日のリズム — 朝・日中・夜の3拍子", "考えなくても回る。朝は見る、日中は報告、夜はポチッ。")
steps = [("朝", "☀️", "01司令塔を見る", "今日の一手・追客目標・ランウェイを確認。\n秘書（Claude）が朝ブリーフで\n『今日やる15人/案件』を提案する。", RED),
         ("日中", "🏃", "動いたら報告する", "電話した・会った・返信きた を\nClaudeに一言。02作業DBに\n自動で1行記録される。", REDD),
         ("夜", "🌙", "『ポチッ』と言う", "その日のチャットをClaudeが遡って\n顧客マスター/作業DBへ振り分け。\n足りない欄は選択肢で質問が来る。", GRY)]
for i,(lab,emo,ti,bo,ac) in enumerate(steps):
    x = Inches(0.6 + i*4.25)
    bx(s, x, Inches(2.2), Inches(3.95), Inches(4.3), WHT, line=CARDLN, lw=0.9)
    bx(s, x, Inches(2.2), Inches(3.95), Inches(0.85), ac)
    t(s, emo+"  "+lab, x, Inches(2.32), Inches(3.95), Inches(0.6), sz=20, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, x+Inches(0.2), Inches(3.25), Inches(3.6), Inches(0.5), sz=15, bold=True, col=ac, align=PP_ALIGN.CENTER)
    t(s, bo, x+Inches(0.3), Inches(3.95), Inches(3.4), Inches(2.3), sz=12, col=INK, align=PP_ALIGN.CENTER, line_sp=1.25)
ft(s, 3)

# ════════════════════════════════ 4. 入力① 02作業DB ════════════════════════════════
s = sl()
hdr(s, "INPUT 1 / 3", "02_作業DB — 行動と報告を1行ずつ", "『何を・誰に・どう動かしたか』の記録。ここが全部の出発点。")
t(s, "入れること（Claudeに言えば自動で1行になる）", Inches(0.6), Inches(1.95), Inches(8), Inches(0.4), sz=13, bold=True, col=GRN)
rows = [["列", "中身", "例"],
        ["日付/区分", "予定 / 実績 / 突発", "実績"],
        ["本部・案件・相手", "誰の何の件か", "04 / 山崎先生"],
        ["温度", "客の熱さ", "ホット"],
        ["内容・報告値", "やったこと・結果の数字", "電話→来週面談OK"],
        ["営業直結", "営業か内務か（時間KPI）", "営業○"],
        ["次アクション・期限", "次に動くこと", "6/15 訪問"]]
cw = [Inches(2.4), Inches(4.4), Inches(5.2)]
tb = s.shapes.add_table(len(rows),3, Inches(0.6), Inches(2.4), Inches(12.0), Inches(3.2)).table
for ci,wv in enumerate(cw): tb.columns[ci].width = wv
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        c = tb.cell(ri,ci); c.text=str(val); c.vertical_anchor=MSO_ANCHOR.MIDDLE
        c.margin_left=Inches(0.12); c.fill.solid()
        c.fill.fore_color.rgb = GRN if ri==0 else (CARD if ri%2 else BG)
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name=FONT; r.font.size=Pt(13 if ri==0 else 12)
                r.font.bold=(ri==0 or ci==0); r.font.color.rgb=WHT if ri==0 else INK
bx(s, Inches(0.6), Inches(5.85), Inches(12.0), Inches(0.8), GRNBG, line=CARDLN, lw=0.75)
t(s, "💡 コツ：『○○先生に電話して来週面談OK、温度ホット』とだけ言えばOK。列の名前は覚えなくていい。",
  Inches(0.8), Inches(5.95), Inches(11.6), Inches(0.6), sz=12.5, bold=True, col=GRN, anchor=MSO_ANCHOR.MIDDLE)
ft(s, 4)

# ════════════════════════════════ 5. 入力② 顧客マスター ════════════════════════════════
s = sl()
hdr(s, "INPUT 2 / 3", "顧客マスター — 人（WHO）を番号で貯める", "物件がまだ無い『人だけ』はここ。H001〜の番号で一生残す。")
card(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(2.0), "ここに入れる人",
     "・先生（医療テナント・承継の見込み）\n・買取再販のテレアポ先\n・まだ案件化してない知り合い\n\n→ H001, H002… と番号がつく", accent=GRN, accbg=GRNBG, tsz=14, bsz=12.5)
card(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.0), "各人につける情報",
     "・温度（ホット/ぬるい）\n・最終接触・次アクション\n・紐づく案件（あれば）\n・関連フォルダ（資料のDriveリンク）", accent=BLU, accbg=BLUBG, tsz=14, bsz=12.5)
bx(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.9), WHT, line=RED, lw=1.2)
t(s, "★ 顧客マスター と 02作業DB の使い分け（迷ったらこれ）", Inches(0.8), Inches(4.55), Inches(11.5), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "人だけ（物件まだ無い）          →  顧客マスター に貯める\n"
     "物件が絡む話になった          →  02作業DB で行動記録＋03で案件化\n"
     "顧客に物件が紐づいた          →  作業DBへ昇格（H番号を引き継ぐ）",
  Inches(0.95), Inches(5.05), Inches(11.4), Inches(1.2), sz=13.5, col=INK, line_sp=1.3)
ft(s, 5)

# ════════════════════════════════ 6. 入力③ 03売上見込み（ドライバー） ════════════════════════════════
s = sl()
hdr(s, "INPUT 3 / 3", "03_売上見込み — 案件と『今日の追客本数』", "案件を確度で積む → 家族と黒字から逆算 → 今日やる本数が出る。")
# 案件
card(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(1.5), "① 案件を積む（5〜14行）",
     "本部・科目・案件名・金額を入れる。\nClaudeに『○○の案件、満額110万、確度6割』\nと言えば入る。", accent=GRN, accbg=GRNBG, tsz=14, bsz=12)
# ドライバー逆算
bx(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(4.4), WHT, line=BLU, lw=1.1)
t(s, "② 営業ドライバー（自動で逆算）", Inches(7.0), Inches(2.12), Inches(5.5), Inches(0.4), sz=14, bold=True, col=BLU)
chain = ["家族が潰れないライン  659,285",
         "−  経常確定粗利（EC+賃料）  142,000",
         "=  🔴 毎月の穴  517,285",
         "÷ 単価 → 必要 成約数  1件",
         "÷ 転換率 → 必要アポ  5件",
         "÷ 転換率 → 必要CV  10件",
         "÷ 稼働日 → 💡 今日の追客  1件/日"]
for i,c in enumerate(chain):
    yy = Inches(2.6 + i*0.52)
    hi = (i==2 or i==6)
    bx(s, Inches(7.0), yy, Inches(5.5), Inches(0.44), REDBG if hi else BG, line=CARDLN, lw=0.5)
    t(s, c, Inches(7.2), yy, Inches(5.2), Inches(0.44), sz=12, bold=hi, col=RED if hi else INK, anchor=MSO_ANCHOR.MIDDLE)
card(s, Inches(0.6), Inches(3.7), Inches(5.9), Inches(2.7), "③ 古田土の未来会計図表",
     "売上(PQ)→粗利(MQ)→利益(G)の形を自動表示。\n\n・損益分岐点比率 F÷MQ（理想80%）\n・格付け SS〜D（今＝SS 超優良）\n・労働分配率・経営安全率\n\nマーケ用語（CV/IMP/CVR）との対応表も同じタブに。",
     accent=RED, accbg=REDBG, tsz=14, bsz=11.5)
ft(s, 6)

# ════════════════════════════════ 7. 見る① 01司令塔 ════════════════════════════════
s = sl()
hdr(s, "VIEW 1 / 3", "01_統合司令塔 — まずここを見る", "過去（資産）→ 現在（現金）→ 未来（谷）→ 今日の一手。1画面で経営が分かる。")
blocks = [("◆ 過去", "純資産・総資産・総負債\n（BS連動）", BLU),
          ("◆ 現在", "現預金とランウェイ\n（あと何ヶ月もつか）", RED),
          ("◆ 未来", "通期で一番苦しい月の現金\n（谷がプラスか）", REDD),
          ("◆ 今日の行動", "緊急→生命線→追客→仕込み\n＋営業直結比率（目標60%）", GRY)]
for i,(ti,bo,ac) in enumerate(blocks):
    x = Inches(0.6 + i*3.1)
    bx(s, x, Inches(2.3), Inches(2.9), Inches(3.9), WHT, line=CARDLN, lw=0.9)
    bx(s, x, Inches(2.3), Inches(2.9), Inches(0.7), ac)
    t(s, ti, x, Inches(2.4), Inches(2.9), Inches(0.5), sz=15, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, bo, x+Inches(0.25), Inches(3.25), Inches(2.4), Inches(2.7), sz=12, col=INK, align=PP_ALIGN.CENTER, line_sp=1.3)
bx(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.55), GRYBG)
t(s, "見方：上から下に読むだけ。赤い信号（🔴守り）が出たら投資を止める、🟢なら攻めてOK。",
  Inches(0.8), Inches(6.37), Inches(11.8), Inches(0.5), sz=12, bold=True, col=INK, anchor=MSO_ANCHOR.MIDDLE)
ft(s, 7)

# ════════════════════════════════ 8. 見る② 04PL+05資金繰り ════════════════════════════════
s = sl()
hdr(s, "VIEW 2 / 3", "04_PL ＋ 05_資金繰り — 利益と現金", "上＝利益の形（古田土）、下＝月ごとの現金の動き。1タブに上下で同居。")
card(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.2), "04_損益PL（上）",
     "・売上 → 変動費 → 粗利(MQ) → 固定費 → 経常\n・法人 / 個人事業 / 家計 / 全社の4本立て\n・第4期決算の実科目に接地（業務委託料/接待交際費 等）\n\n見るポイント：\n  損益分岐点比率（F÷MQ）が低いほど安全\n  → 今のKHDは原価ほぼ0のサービス業で SS級",
     accent=RED, accbg=REDBG, tsz=15, bsz=12.5)
card(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(4.2), "05_資金繰り（下）",
     "・月ごとの入金 − 出金 = 月末現金\n・主格別（法人/個人/家計）\n・ランウェイ＝現金があと何ヶ月もつか\n\n見るポイント：\n  最悪月（谷）の現金がプラスか\n  → ここがマイナスなら『毎月の穴』を\n     営業で埋めるのが生命線",
     accent=BLU, accbg=BLUBG, tsz=15, bsz=12.5)
ft(s, 8)

# ════════════════════════════════ 9. 色とルール ════════════════════════════════
s = sl()
hdr(s, "RULES", "色の意味 と やってはいけないこと", "色を見れば『触っていい所』が分かる。地雷だけ踏まないように。")
t(s, "■ セルの色", Inches(0.6), Inches(2.0), Inches(6), Inches(0.4), sz=14, bold=True, col=RED)
legend = [("🟡 黄（薄）", "あなたが入力する所", GRNBG, GRN),
          ("🟦 青（薄）", "自動計算（触らない）", BLUBG, BLU),
          ("🟧 金/橙", "結論・KPIの答え", REDBG, RED),
          ("⬜ 灰", "見出し・参考", GRYBG, GRY)]
for i,(lab,desc,bgc,ac) in enumerate(legend):
    y = Inches(2.5 + i*0.95)
    bx(s, Inches(0.6), y, Inches(1.9), Inches(0.7), bgc, line=ac, lw=1.0)
    t(s, lab, Inches(0.6), y, Inches(1.9), Inches(0.7), sz=12, bold=True, col=ac, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t(s, desc, Inches(2.7), y, Inches(3.6), Inches(0.7), sz=13, col=INK, anchor=MSO_ANCHOR.MIDDLE)
# 禁止
bx(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(4.4), RGBColor(0xFB,0xEE,0xEC), line=RED, lw=1.4)
t(s, "🚫 やってはいけない", Inches(7.0), Inches(2.15), Inches(5.5), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "① タブ名を変えない\n   → 05資金繰りの計算式 約40本が全部\n     #REF! で壊れる（過去の事故）\n\n"
     "② 青・灰のセルを書き換えない\n   → 自動計算が消える\n\n"
     "③ 行を勝手に削除/挿入しない\n   → 参照がズレて #REF!\n\n"
     "④ 触って壊れたら → Claudeに『直して』",
  Inches(7.0), Inches(2.7), Inches(5.5), Inches(3.6), sz=12.5, col=INK, line_sp=1.25)
ft(s, 9)

# ════════════════════════════════ 10. 困ったら ════════════════════════════════
s = sl()
hdr(s, "TROUBLE", "困ったとき — #REF! や #DIV/0! が出たら", "あわてない。だいたい2手で直る。")
tr = [("1", "まず Cmd + R でリロード", "画面が古いキャッシュのことが多い。\n再計算されてエラーが消えれば、それで終わり。", RED),
      ("2", "消えなければ Claude に言う", "『#REFが出てる』『○○タブの△行』とだけ伝える。\nClaudeがDriveから直接シートを読んで、\n壊れた参照を1本のスクリプトで直す。", BLU),
      ("3", "数字がおかしい時", "『○○の数字が合わない』と言えば、\n計算式をさかのぼって原因を特定する。", GRY)]
for i,(n,ti,bo,ac) in enumerate(tr):
    y = Inches(2.2 + i*1.45)
    bx(s, Inches(0.6), y, Inches(12.1), Inches(1.25), WHT, line=CARDLN, lw=0.9)
    bx(s, Inches(0.6), y, Inches(1.1), Inches(1.25), ac)
    t(s, n, Inches(0.6), y, Inches(1.1), Inches(1.25), sz=30, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t(s, ti, Inches(1.9), y+Inches(0.13), Inches(10.5), Inches(0.45), sz=15, bold=True, col=ac)
    t(s, bo, Inches(1.9), y+Inches(0.58), Inches(10.5), Inches(0.6), sz=12, col=INK, line_sp=1.15)
ft(s, 10)

# ════════════════════════════════ 11. まとめ ════════════════════════════════
s = sl()
bx(s, 0, 0, W, H, BG)
bx(s, Inches(0.0), Inches(2.4), W, Inches(0.06), RED)
t(s, "まとめ", Inches(0.9), Inches(1.5), Inches(11), Inches(0.5), sz=15, bold=True, col=RED)
t(s, "あなたは『入力』と『報告』だけ。", Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.8), sz=32, bold=True, col=INK)
t(s, "数字・利益・現金の谷は、つながって勝手に動く。", Inches(0.92), Inches(3.6), Inches(11), Inches(0.5), sz=17, col=GRY)
pts = ["朝＝01司令塔を見る   ／   日中＝動いたら報告   ／   夜＝『ポチッ』で振り分け",
       "迷ったら：人だけ→顧客マスター、物件絡む→作業DB",
       "壊れたら：リロード → ダメなら Claude に『直して』",
       "タブ名は変えない（40本の式が飛ぶ）"]
for i,p in enumerate(pts):
    y = Inches(4.4 + i*0.6)
    t(s, "✓  "+p, Inches(0.95), y, Inches(11.5), Inches(0.5), sz=13.5, bold=(i==0), col=RED if i==0 else INK)
ft(s, 11)

OUT = "/Users/kikuchikenta/01_honbu_docs_automation/cockpit_manual.pptx"
prs.save(OUT)
print("SAVED:", OUT, "/ slides:", len(prs.slides._sldIdLst))

"""
AI医療コンサル 受注ロードマップ（福井起点）スライド
v8 と同一デザインシステム（クリーム白×レンガ赤）を踏襲。
出力: fukui_medical_flow.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Hiragino Sans"
W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "AI医療コンサル 受注ロードマップ  ｜  04コンサル本部", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)


def light_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = REDD if is_hi else RED
            else:
                cell.fill.fore_color.rgb = REDBG if is_hi else (CARD if ri % 2 == 1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = WHT
                    elif is_hi:
                        r.font.color.rgb = RED
                    elif ci == 0:
                        r.font.color.rgb = INK
                    else:
                        r.font.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    return tb


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "SALES ROADMAP", Inches(0.9), Inches(1.9), Inches(11), Inches(0.45), sz=15, bold=True, col=RED)
t(s, "AI医療コンサル", Inches(0.88), Inches(2.45), Inches(11.5), Inches(0.9), sz=42, bold=True, col=INK)
t(s, "受注までの設計図", Inches(0.88), Inches(3.35), Inches(11.5), Inches(0.9), sz=42, bold=True, col=RED)
t(s, "福井面談を起点に、最初の1院を「初期0円・成果報酬・リスクゼロ」で獲得するまでの全工程。\n弾＝提案デッキ clinic_dx_v8（全11枚）。本書はその「撃ち方」。",
  Inches(0.9), Inches(4.45), Inches(11.4), Inches(1.0), sz=14, col=GRY, line_sp=1.3)
bx(s, Inches(0.9), Inches(6.05), Inches(11.5), Pt(1.2), LINE)
t(s, "チームてっかん  ｜  菊池 研太  ｜  04コンサル本部（李牧）", Inches(0.9), Inches(6.2), Inches(11), Inches(0.4), sz=13, bold=True, col=INK)

# ════════ SLIDE 2 — 全体フロー（6段階） ════════
s = sl(); ft(s)
hdr(s, "THE PROCESS", "福井起点 ── 受注までの6ステップ", "「福井に動いてもらう」のがゴールではない。最初の1院の「実名」を取り、自分で繋がるまでが設計。")
stages = [
    ("STAGE 0", "福井へ打診", "改まって30分の\n壁打ち枠を取る"),
    ("STAGE 1", "福井面談", "商品理解＋紹介\n＋座組み合意"),
    ("STAGE 2", "医師アポ", "「無料の削減試算」\nで30分を獲得"),
    ("STAGE 3", "提案(v8)", "Slide6で自院ROIを\nその場で自分ごと化"),
    ("STAGE 4", "無料コンサル", "競合分析＋報告書\nでGIVE・信頼構築"),
    ("STAGE 5", "成果報酬契約", "測定合意→伴走→\n承継/テナントへ"),
]
cw, gx, x0, y0, ch = Inches(1.95), Inches(0.13), Inches(0.45), Inches(2.15), Inches(3.7)
for i, (no, ti, goal) in enumerate(stages):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.7), RED)
    t(s, no, cx, y0+Inches(0.14), cw, Inches(0.42), sz=14, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, cx, y0+Inches(0.95), cw, Inches(0.5), sz=16, bold=True, col=INK, align=PP_ALIGN.CENTER)
    bx(s, cx+Inches(0.6), y0+Inches(1.55), cw-Inches(1.2), Pt(2), RED)
    t(s, goal, cx+Inches(0.12), y0+Inches(1.75), cw-Inches(0.24), Inches(1.8), sz=11, col=GRY,
      align=PP_ALIGN.CENTER, line_sp=1.2)
    if i < len(stages) - 1:
        t(s, "▶", cx+cw-Inches(0.02), y0+Inches(1.5), Inches(0.2), Inches(0.4), sz=13, bold=True, col=RED, align=PP_ALIGN.CENTER)
bx(s, Inches(0.45), Inches(6.15), Inches(12.4), Inches(0.7), REDBG)
bx(s, Inches(0.45), Inches(6.15), Inches(0.1), Inches(0.7), RED)
t(s, "李牧の鉄則 ── 福井は触媒。主語は常に菊池。並行で「温かいリスト」の医師を1名、依存の保険に取る。",
  Inches(0.75), Inches(6.27), Inches(12.0), Inches(0.45), sz=13, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)

# ════════ SLIDE 3 — 各段階のゴールと一手 ════════
s = sl(); ft(s)
hdr(s, "WHAT TO SAY", "各段階：ゴールと「最初の一言」", "迷ったらこの一言から入る。売り込まず、相談・無料診断（GIVE）で入口を開ける。")
rows = [
    ("段階", "この段階のゴール", "キーになる一言（入口）"),
    ("0 福井へ打診", "30分の壁打ち枠を取る", "「資料が完成形まで固まった。刺さりそうな先生いないか30分相談させて」"),
    ("1 福井面談", "具体名を1〜2件引き出す", "「最初に当てるなら誰がいい？ ◯◯先生あたりどうでしょう」"),
    ("2 医師アポ", "無料試算で初回30分", "「初期ゼロ・成果報酬です。まず無料で御院の削減効果を試算させてください」"),
    ("3 提案(v8)", "Slide6で自院ROIを体験", "「事務さん何名ですか？ …3名なら年で約200万、院に残ります」"),
    ("4 無料コンサル", "信頼の対価を積む", "「2週間計測みたいな面倒は一切不要。測定はこちらの仕事です」"),
    ("5 契約", "測定合意→伴走開始", "「効果が出た分の3割だけ。出なければ0円。リスクは私が取ります」"),
]
light_table(s, rows, Inches(0.5), Inches(1.85), Inches(12.33), Inches(4.6),
            [Inches(2.0), Inches(3.2), Inches(7.13)], hi_col=None, sz=12, header_sz=13)

# ════════ SLIDE 4 — 座組み（福井とのWIN-WIN） ════════
s = sl(); ft(s)
hdr(s, "THE DEAL", "座組み ── 福井（TAW）にもリターンが乗る形", "紹介してもらう以上、叔父にも必ず果実を。利益分配は口頭で終わらせず「書面化」する。")
deal = [
    ("菊池（KHD）", "AI導入の「中身」", "提案・無料コンサル・導入・3ヶ月伴走まで実行。リスク（成果報酬）はこちらが負う。"),
    ("福井（TAW）", "医師ネットワーク", "最初の1院への紹介・信頼の貸与。これが「温かい入口」になる。"),
    ("共に取る果実", "承継・テナント", "AIを入口に承継/テナント相談(110万/件)へ。ここは一緒に取りにいく。"),
]
cw, gx, x0, y0, ch = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.95), Inches(3.1)
for i, (ti, tag, body) in enumerate(deal):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.06), RED)
    t(s, ti, cx, y0+Inches(0.3), cw, Inches(0.5), sz=18, bold=True, col=INK, align=PP_ALIGN.CENTER)
    t(s, tag, cx, y0+Inches(0.9), cw, Inches(0.5), sz=15, bold=True, col=RED, align=PP_ALIGN.CENTER)
    bx(s, cx+Inches(0.5), y0+Inches(1.5), cw-Inches(1.0), Pt(2), LINE)
    t(s, body, cx+Inches(0.3), y0+Inches(1.7), cw-Inches(0.6), Inches(1.3), sz=12, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.25)
bx(s, Inches(0.55), Inches(5.4), Inches(12.23), Inches(1.0), REDBG)
bx(s, Inches(0.55), Inches(5.4), Inches(0.1), Inches(1.0), RED)
t(s, "★ 必ず書面化 ── 紹介料 or 承継・テナント案件化時の分配ルールを、契約前に1枚で取り決める。",
  Inches(0.85), Inches(5.6), Inches(11.8), Inches(0.6), sz=14, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)

# ════════ SLIDE 5 — 次アクション ════════
s = sl(); ft(s)
hdr(s, "NEXT ACTIONS", "次の一手 ── これだけは先に潰す", "資料完成はゴールじゃない。アポ1件がゴール。下の4つを動かす。")
acts = [
    ("1", "候補医師を用意", "福井に出す「カード」＝最初に当てる先生1〜2名の実名を、面談前に決めておく。"),
    ("2", "木曜30分を確保", "火曜に会う前にLINEで打診。改めて木曜に30分の壁打ち枠をもらう。"),
    ("3", "並行リスト1名", "福井依存の保険。リベ大の医師など「温かい接点」を1名、同じトークで当てる。"),
    ("4", "利益分配を書面化", "福井との座組み（紹介料 or 案件化時の分配）を1枚で取り決める。口頭で流さない。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.05), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.95)
for i, (no, ti, body) in enumerate(acts):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    t(s, no, cx+Inches(0.3), cy+Inches(0.25), Inches(0.8), Inches(0.9), sz=32, bold=True, col=RED)
    t(s, ti, cx+Inches(1.15), cy+Inches(0.3), cw-Inches(1.35), Inches(0.5), sz=18, bold=True, col=INK)
    t(s, body, cx+Inches(1.15), cy+Inches(0.9), cw-Inches(1.35), Inches(1.0), sz=12, col=GRY, line_sp=1.2)

prs.save("fukui_medical_flow.pptx")
print("saved fukui_medical_flow.pptx  /  slides:", len(prs.slides._sldIdLst))

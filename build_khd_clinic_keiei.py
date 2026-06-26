"""
クリニック経営の全体像と、KHDが効く所（現状→未来）。
前提揃え＋他院営業の両用。どこでも何でもできるAIの強み。クリーム白×レンガ赤。
出力: khd_clinic_keiei.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.55),Inches(0.35),Inches(12),Inches(0.4),sz=12.5,bold=True,col=RED)
    bx(slide,Inches(0.57),Inches(0.72),Inches(1.6),Pt(3),RED)
    t(slide,main,Inches(0.55),Inches(0.82),Inches(12.2),Inches(0.55),sz=20,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.57),Inches(1.36),Inches(12.2),Inches(0.3),sz=11,col=GRY)
def ft(slide,n):
    bx(slide,Inches(0.5),H-Inches(0.45),Inches(12.33),Pt(1.2),LN)
    t(slide,"KHD ｜ AI医療コンサル ── クリニック経営の全体像と、効く所",Inches(0.5),H-Inches(0.4),Inches(10),Inches(0.3),sz=8.5,col=GRY)
    t(slide,n,Inches(12.45),H-Inches(0.4),Inches(0.4),Inches(0.3),sz=8.5,col=GRY)
def tbl(slide,rows,x,y,w,h,col_w,sz=11,hsz=11.5):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            c=tb.cell(ri,ci); c.text=str(val); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.09); c.margin_right=Inches(0.05); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            c.fill.solid()
            if ri==0: c.fill.fore_color.rgb=RED
            else: c.fill.fore_color.rgb=(CARD if ri%2==1 else BG)
            for p in c.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(hsz if ri==0 else sz); r.font.bold=(ri==0 or ci==0)
                    r.font.color.rgb=(WHT if ri==0 else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
    return tb

# 1 表紙
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"KHD ｜ AI医療コンサル",Inches(0.9),Inches(1.1),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"クリニック経営の全体像と、\nKHDが効く所（現状 → 未来）",Inches(0.88),Inches(1.65),Inches(11.7),Inches(1.7),sz=27,bold=True,col=INK,line_sp=1.12)
t(s,"経営のどこで、AIが効くか。いま「できていないこと」が、こう変わる。",Inches(0.92),Inches(3.75),Inches(11.4),Inches(0.5),sz=13.5,col=GRY)
bx(s,Inches(0.9),Inches(4.55),Inches(11.5),Inches(1.2),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(4.55),Inches(0.1),Inches(1.2),RED)
t(s,"特定のツールを売る会社ではありません。経営の「どこでも」入って、AIで「何でも」小さく作って試せる——\nそれが、KHDのAI医療コンサルです。",Inches(1.2),Inches(4.55),Inches(11),Inches(1.2),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)

# 2 経営の全体像ポンチ図
s=sl(); ft(s,"2")
hdr(s,"まず全体像","クリニック経営は「5つの箱」で回っている","患者さんは ①→②→③→④ を回り、再診で①に戻る（継続）。⑤が土台。")
boxes=[("① 集患","見つけてもらう\n(HP・地図・紹介)"),("② 受付・予約","来てもらう\n(電話・受付・順番)"),
       ("③ 診療","診る・記録する\n(問診・カルテ)"),("④ 再診・継続","また来てもらう\n(次回予約・定期)")]
cw,gx,x0,yb=Inches(2.78),Inches(0.42),Inches(0.55),Inches(2.15)
for i,(ti,ds) in enumerate(boxes):
    cx=x0+(cw+gx)*i
    bx(s,cx,yb,cw,Inches(1.55),CARD,line=CARDLN,lw=1.0); bx(s,cx,yb,cw,Inches(0.06),RED)
    t(s,ti,cx+Inches(0.15),yb+Inches(0.2),cw-Inches(0.3),Inches(0.5),sz=15,bold=True,col=REDD,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.15),yb+Inches(0.75),cw-Inches(0.3),Inches(0.7),sz=11,col=INK,align=PP_ALIGN.CENTER,line_sp=1.2)
    if i<3: t(s,"→",cx+cw+Inches(0.02),yb+Inches(0.55),Inches(0.38),Inches(0.5),sz=20,bold=True,col=RED,align=PP_ALIGN.CENTER)
t(s,"↩ 再診・継続で ① に戻る（このループが太いほど、経営は安定）",Inches(0.55),Inches(3.95),Inches(12.2),Inches(0.4),sz=11.5,bold=True,col=REDD)
bx(s,Inches(0.55),Inches(4.5),Inches(12.23),Inches(1.4),GRYBG); bx(s,Inches(0.55),Inches(4.5),Inches(0.1),Inches(1.4),RED)
t(s,"⑤ バックオフィス（土台）",Inches(0.82),Inches(4.62),Inches(11.6),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"記帳・税務／IT導入補助金／採用・定着／承継・M&A・診療圏（開業・分院）。表に出ないが、ここが崩れると全部止まる。",Inches(0.85),Inches(5.08),Inches(11.6),Inches(0.7),sz=12,col=INK,line_sp=1.25)

# 3 現状→未来の表
s=sl(); ft(s,"3")
hdr(s,"現状 → 未来","いま「できていないこと」が、AIでこう変わる","ここを全員で揃える。推測でなく、現状を見える化してから。")
rows=[
 ("経営の箱","いま できていないこと（現状）","KHDのAIで こう変わる（未来）"),
 ("① 集患","HP/MEOが弱い。「何で知ったか」も紙で死蔵＝流入が不明","HP刷新・地図最適化＋問診で流入を自動で数値化"),
 ("② 受付・予約","電話・紙・アイコールのみ。来院の前後に手が届かない","LINEで問診・予約・リマインド・お知らせを足す"),
 ("③ 診療","入力・転記・紹介状・カルテ化の手間が残業に","AI入力・音声・Stream Deck・定型文で時短"),
 ("④ 再診・継続","人力の架電で取りこぼし→月末残業","次回予約・CPAP自動・離脱の掘り起こしで継続"),
 ("⑤ バックオフィス","補助金・記帳・採用・承継が手つかず","補助金申請支援・診療圏調査・承継コンサル"),
]
tbl(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(4.3),[Inches(2.3),Inches(5.1),Inches(4.83)],sz=11,hsz=11.5)

# 4 KHDが効く所マップ
s=sl(); ft(s,"4")
hdr(s,"どこでも入れる","KHDが効く所マップ ── 経営の箱 × 打ち手","特定の1機能でなく、経営のどこにでも。だから「刺さる1つ」から選べる。")
mp=[("① 集患","HP / MEO / オウンドメディア × AI / 流入の見える化"),
    ("② 受付・予約","公式LINE（問診・予約・リマインド・FAQ・お知らせ）"),
    ("③ 診療","AI入力 / 音声カルテ / Stream Deck / 紹介状AI"),
    ("④ 再診・継続","次回予約・CPAP自動 / 離脱の掘り起こし / 追客"),
    ("⑤ バックオフィス","IT補助金支援 / 診療圏調査 / 承継・M&A / 記帳効率化")]
y=Inches(2.0)
for ti,ds in mp:
    bx(s,Inches(0.55),y,Inches(12.23),Inches(0.82),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),y,Inches(0.1),Inches(0.82),RED)
    t(s,ti,Inches(0.8),y,Inches(2.5),Inches(0.82),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ds,Inches(3.4),y,Inches(9.2),Inches(0.82),sz=12.5,col=INK,anchor=MSO_ANCHOR.MIDDLE)
    y=y+Inches(0.92)
t(s,"→ 全部を一度に売りません。御院の課題に合う「1つ」から、無料で小さく。",Inches(0.55),Inches(6.65),Inches(12),Inches(0.0),sz=11,bold=True,col=REDD) if False else None

# 5 AIの強み
s=sl(); ft(s,"5")
hdr(s,"なぜ KHD か","「どこでも・何でもできるAI」の強み","製品を売る会社にはできない立ち位置")
whys=[("中立","特定ツールを売らない。課題から最適な打ち手を選ぶ"),
      ("実物を0円で作る","資料で終わらせない。AIで動く実物を作って見せる"),
      ("速い・柔軟","従来20〜250万のものを0円〜。その場で直せる"),
      ("経営×現場の両目線","数字（経営）と、スタッフの手間（現場）の両方を見る")]
cw,ch,gx,gy=Inches(6.0),Inches(1.5),Inches(0.23),Inches(0.25); x0,y0=Inches(0.55),Inches(2.0)
for i,(ti,ds) in enumerate(whys):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.25),cy+Inches(0.16),cw-Inches(0.45),Inches(0.45),sz=15,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.27),cy+Inches(0.65),cw-Inches(0.5),Inches(0.7),sz=11.5,col=INK,line_sp=1.2)
bx(s,Inches(0.55),Inches(5.5),Inches(12.23),Inches(0.95),REDBG); bx(s,Inches(0.55),Inches(5.5),Inches(0.1),Inches(0.95),RED)
t(s,"他社は「自社製品を売る」。KHDは「クリニックがAIを使いこなせるようにする」＝中立で、経営のどこでも伴走する。",Inches(0.82),Inches(5.5),Inches(11.8),Inches(0.95),sz=13,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# 6 はじめ方
s=sl(); ft(s,"6")
hdr(s,"はじめ方","まず「無料診断」で、現状を見える化","推測でなく数字で。そこから未来図を一緒に描きます。")
steps=[("① 無料診断","経営の5つの箱を一緒に棚卸し。今の数字を見える化"),
       ("② 刺さる1つ","効果が大きい所を1つ選ぶ（LINEとは限らない）"),
       ("③ 0円で小さく","AIで実物を作って試す。負担なく"),
       ("④ 測って広げる","前後を数字で比較。効いた所だけ拡大")]
cw,gx,x0=Inches(2.95),Inches(0.18),Inches(0.55)
for i,(ti,ds) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.2),cw,Inches(2.3),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.2),cw,Inches(0.06),RED)
    t(s,ti,cx+Inches(0.15),Inches(2.4),cw-Inches(0.3),Inches(0.6),sz=14,bold=True,col=REDD,align=PP_ALIGN.CENTER,line_sp=1.1)
    t(s,ds,cx+Inches(0.2),Inches(3.15),cw-Inches(0.4),Inches(1.2),sz=11.5,col=INK,align=PP_ALIGN.CENTER,line_sp=1.3)
    if i<3: t(s,"→",cx+cw-Inches(0.02),Inches(3.1),Inches(0.3),Inches(0.5),sz=16,bold=True,col=RED,align=PP_ALIGN.CENTER)
bx(s,Inches(0.55),Inches(4.85),Inches(12.23),Inches(1.0),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(4.85),Inches(0.1),Inches(1.0),TEALD)
t(s,"御院に合うかどうか、一緒に確かめながら。まず小さく、いつでも見直せます。",Inches(0.82),Inches(4.85),Inches(11.8),Inches(1.0),sz=14,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("khd_clinic_keiei.pptx")
print("saved khd_clinic_keiei.pptx slides:",len(prs.slides._sldIdLst))

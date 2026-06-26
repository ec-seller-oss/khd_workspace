# -*- coding: utf-8 -*-
"""
京橋クリニック AI医療コンサル｜6/15週 会食 当日決めることA4一枚
KHD配色: クリーム白#F9F6EF × レンガ赤#AA2E26 / Hiragino Sans
A4縦 1枚（pptx → 手動 or LibreOfficeでPDF化）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM = RGBColor(0xF9,0xF6,0xEF); RED = RGBColor(0xAA,0x2E,0x26)
INK = RGBColor(0x1A,0x1A,0x1A); RULE = RGBColor(0xDA,0xD6,0xCF)
BEIGE = RGBColor(0xF1,0xEC,0xE1); GREY = RGBColor(0x6B,0x6B,0x6B)
WHITE = RGBColor(0xFF,0xFF,0xFF); GREEN = RGBColor(0x2E,0x7D,0x32)
FONT = "Hiragino Sans"

# A4縦: 210mm x 297mm
prs = Presentation()
prs.slide_width = Emu(int(210/25.4*914400))
prs.slide_height = Emu(int(297/25.4*914400))
SW, SH = prs.slide_width, prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(x,y,w,h,fill=None,line=None,lw=1.0):
    sp = slide.shapes.add_shape(1, Inches(x),Inches(y),Inches(w),Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp

def txt(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.0,sb=0):
    tb = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0], tuple): runs=[runs]
    for i,para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        if sb: p.space_after = Pt(sb)
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.name=FONT; r.font.size=Pt(sz)
            r.font.bold=b; r.font.color.rgb=c
    return tb

AW = 210/25.4  # page width in inches
# 背景
rect(0,0,AW,297/25.4,fill=CREAM)
# 左赤縦バー
rect(0,0,0.10,297/25.4,fill=RED)

M = 0.55  # margin
CW = AW - M*2  # content width

# 社内秘バナー（対外提示禁止）
rect(M,0.18,CW,0.22,fill=RED)
txt(M,0.18,CW,0.22,[[("【社内秘・対外提示禁止】 菊池＋福井の作戦メモ — 山崎先生には見せない",9,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# ヘッダー
txt(M,0.45,CW,0.25,[[("MEETING DECISION SHEET ｜ 6/15週 会食",10.5,True,RED)]])
txt(M,0.72,CW,0.5,[[("京橋クリニック ｜ AI医療コンサル",19,True,INK)]])
txt(M,1.22,CW,0.25,[[("当日 決めること・残務・今後（A4一枚）",10.5,False,GREY)]])
rect(M,1.55,CW,0.03,fill=RED)
# メタ
txt(M,1.66,CW,0.25,[[("対象：山崎先生（内科）＋福井（事務長） ／ ",9,False,GREY),
                     ("中核信条：売り込まない・GIVE先行・信頼の対価",9,True,RED)]])

y = 2.05
# ===== 🎯 決めること =====
def band(yy, label, color):
    rect(M,yy,CW,0.34,fill=color)
    txt(M+0.12,yy,CW-0.24,0.34,[[(label,12,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
    return yy+0.34

y = band(y, "🎯 決めること", RED)
decisions = [
    ("1","将来ゴールの合意","京橋を最初の成功事例に → 他院へ横展開する絵を共有"),
    ("2","最初に着手する2課題","電話IVR ／ LINE問診（京橋の最大ペイン）の優先順を決める"),
    ("3","ベースライン数値の確定","福井提供の実数をROIタブに入力 → 「年◯◯万残る」をその場提示"),
    ("4","無料診断の許可＋データ開示","削減シミュ報告書の作成許可と既存データ閲覧の許可"),
    ("5","導入の段取り","誰が・いつ ／ 医療事務アンケートを全員分回収"),
    ("6","福井との利益分配","覚書で配分％を確定（口頭で流さない）"),
]
yy = y+0.06
for n,t,d in decisions:
    rect(M,yy,0.34,0.34,fill=BEIGE)
    txt(M,yy,0.34,0.34,[[(n,12,True,RED)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(M+0.45,yy-0.01,CW-0.45,0.36,[[(t+"　",10.5,True,INK),(d,9.5,False,GREY)]],anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.40
y = yy+0.05

# ===== ✅ 残務 =====
y = band(y, "✅ 残務（当日まで or 当日割り振る）", RGBColor(0x8A,0x8A,0x8A))
zanmu = [
    "未回収の事務アンケートを全員分回収（福井経由）＝現場の声を拾うGIVE",
    "問診票のデジタル化の可否確認（紙→LINE事前問診・集患経路の集計）",
    "既存ベンダー（アイコール予約）の契約・解約条件の確認",
    "決裁者の確認（山崎先生か理事長か）",
]
yy=y+0.08
for t in zanmu:
    txt(M+0.12,yy,0.2,0.28,[[("□",11,True,RED)]])
    txt(M+0.40,yy,CW-0.40,0.28,[[(t,9.8,False,INK)]])
    yy+=0.30
y=yy+0.05

# ===== ▶ 今後 =====
y = band(y, "▶ 今後の流れ", RED)
txt(M+0.12,y+0.10,CW-0.24,0.3,
    [[("翌週 医療事務ヒアリング ▶ 導入（IVR＋LINE問診） ▶ 90日 効果測定 ▶ 成果報酬 ▶ 成功事例化 ▶ 承継・テナント（110万/件）の本丸へ",10,True,INK)]])
y += 0.55

# プラス材料 カード
rect(M,y,CW,0.78,fill=BEIGE,line=RED,lw=1.0)
txt(M+0.15,y+0.08,CW-0.3,0.25,[[("＋ 提案のプラス材料",10,True,RED)]])
txt(M+0.15,y+0.33,CW-0.3,0.45,
    [[("・事務アンケート全員回収＝離職リスクの根拠＆現場の味方化　",9.3,False,INK)],
      [("・問診票活用＝①紙→LINE問診で手入力ゼロ ②「何で知ったか」で集患経路→再来院導線　※今ある紙を生かす提案",9.3,False,INK)]])
y += 0.92

# フッター
rect(M,y,CW,0.012,fill=RULE)
txt(M,y+0.06,CW,0.25,[[("ROI試算・WBS・料金体系の詳細 → 受注設計スプシ（8タブ）／ 福井利益分配は覚書ドラフト参照　｜　04コンサル本部（李牧） 2026-05-30",8,False,GREY)]])

out="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/京橋_当日決めることA4.pptx"
prs.save(out)
print("SAVED:", out)

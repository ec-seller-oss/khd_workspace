# -*- coding: utf-8 -*-
"""
【外部・先生提示用】京橋クリニック様 へのご提案サマリ A4一枚
KHD配色: クリーム白#F9F6EF × レンガ赤#AA2E26 / Hiragino Sans
社内用語(売り込まない/GIVE/利益分配/横展開/本丸/抜け漏れ)は一切載せない。院のメリットのみ。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM = RGBColor(0xF9,0xF6,0xEF); RED = RGBColor(0xAA,0x2E,0x26)
INK = RGBColor(0x1A,0x1A,0x1A); RULE = RGBColor(0xDA,0xD6,0xCF)
BEIGE = RGBColor(0xF1,0xEC,0xE1); GREY = RGBColor(0x6B,0x6B,0x6B)
WHITE = RGBColor(0xFF,0xFF,0xFF)
FONT = "Hiragino Sans"

prs = Presentation()
prs.slide_width = Emu(int(210/25.4*914400))
prs.slide_height = Emu(int(297/25.4*914400))
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(x,y,w,h,fill=None,line=None,lw=1.0):
    sp = slide.shapes.add_shape(1, Inches(x),Inches(y),Inches(w),Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp

def txt(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.0):
    tb = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0], tuple): runs=[runs]
    for i,para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.name=FONT; r.font.size=Pt(sz)
            r.font.bold=b; r.font.color.rgb=c
    return tb

AW = 210/25.4
rect(0,0,AW,297/25.4,fill=CREAM)
rect(0,0,0.10,297/25.4,fill=RED)
M = 0.55; CW = AW - M*2

# ヘッダー
txt(M,0.45,CW,0.25,[[("京橋クリニック様 へのご提案",10.5,True,RED)]])
txt(M,0.72,CW,0.5,[[("AI・LINEで、現場の負担を軽くする",19,True,INK)]])
txt(M,1.22,CW,0.25,[[("初期費用0円・成果報酬型のクリニックDX",10.5,False,GREY)]])
rect(M,1.55,CW,0.03,fill=RED)
txt(M,1.65,CW,0.22,[[("テナントアシスト・ウイン株式会社　菊池 研太",9,False,GREY)]])

def band(yy, label):
    rect(M,yy,CW,0.32,fill=RED)
    txt(M+0.12,yy,CW-0.24,0.32,[[(label,11.5,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
    return yy+0.32

y = 2.00
# 1. 御院で起きていること
y = band(y, "御院で、こんなことは起きていませんか？")
items = [
    "電話が鳴り止まず、受付が目の前の患者対応とぶつかる",
    "紙の問診票を、毎日手で入力し直している",
    "月初のレセプトに、まとまった時間がかかる",
    "ご予約の行き違いで、受付や診察の順番が混み合う",
]
yy=y+0.08
for t in items:
    txt(M+0.14,yy,0.2,0.28,[[("●",9,True,RED)]])
    txt(M+0.42,yy,CW-0.42,0.28,[[(t,10.5,False,INK)]])
    yy+=0.31
y=yy+0.06

# 2. AIでこう変わる
y = band(y, "AI・LINEで、こう変わります")
sol = [
    ("電話","AI音声が一次対応。営業電話も自動でブロック"),
    ("予約・問診","LINEで来院前にスマホ問診・順番のご案内"),
    ("書類・レセプト","マイAIが下書きと点検を補助"),
]
yy=y+0.08
for a,b in sol:
    rect(M+0.10,yy,1.45,0.32,fill=BEIGE)
    txt(M+0.10,yy,1.45,0.32,[[(a,10.5,True,RED)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(M+1.70,yy,CW-1.70,0.32,[[(b,10.5,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.38
y=yy+0.06

# 3. リスクゼロの仕組み
y = band(y, "御院の持ち出しは「ゼロ」から始められます")
rect(M,y+0.06,CW,0.72,fill=WHITE,line=RULE,lw=1.0)
txt(M+0.15,y+0.13,CW-0.3,0.62,
    [[("初期費用 0円　／　月額 0円　／　成果報酬 40%（最初の6ヶ月だけ）",11,True,RED)],
     [("7ヶ月目以降は 月5万円の保守（月1回MTG込み）。効果0なら費用0円・いつでも解約可・3ヶ月伴走",9.6,False,INK)]],
    space=1.15)
y += 0.88

# 4. 御院に残る試算
y = band(y, "御院に、いくら残るか（標準的な目安）")
rect(M,y+0.06,CW,0.72,fill=BEIGE)
txt(M+0.15,y+0.12,CW-0.3,0.62,
    [[("月 約24万円の削減　→　",11,True,INK),("1〜6ヶ月は院に 14.4万／月、7ヶ月目〜は月5万保守だけで院に 19万／月",10.5,True,RED)],
     [("12ヶ月で院に約200万円が残る計算　／　※当日、御院の実数で計算し直してお出しします",9.5,False,GREY)]],
    space=1.2)
y += 0.88

# 5. 進め方
y = band(y, "進め方（リスクのない3ステップ）")
steps = [
    ("1","無料で御院専用の「業務削減シミュレーション報告書」を作成します"),
    ("2","お困りの大きいところから、段階的に導入（3ヶ月伴走）"),
    ("3","効果を測ったうえで、出た成果に応じてのみご精算"),
]
yy=y+0.08
for n,t in steps:
    rect(M+0.10,yy,0.30,0.30,fill=RED)
    txt(M+0.10,yy,0.30,0.30,[[(n,11,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(M+0.52,yy,CW-0.52,0.30,[[(t,10.3,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.36
y=yy+0.05

rect(M,y,CW,0.012,fill=RULE)
txt(M,y+0.06,CW,0.3,[[("本日は、御院で一番お困りの業務を、一緒に整理させてください。無理におすすめすることはいたしません。",9.5,True,RED)]])

out="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/【先生提示用】京橋クリニック様_ご提案サマリA4.pptx"
prs.save(out)
print("SAVED:", out)

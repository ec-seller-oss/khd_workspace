"""
先生＋医療事務 共通プレゼン台本スライド。クリーム白×レンガ赤。
出力: kyobashi_dabon.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック ｜ 先生＋医療事務 共通台本 ｜ KHD AI医療コンサル",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=(TEALBG if is_hi else (CARD if ri%2==1 else BG))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or (ci==0)
                    r.font.color.rgb=(WHT if ri==0 else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
    return tb

# ════ SLIDE 1 — 表紙 ════
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"TALK SCRIPT",Inches(0.9),Inches(1.6),Inches(11),Inches(0.4),sz=15,bold=True,col=RED)
t(s,"先生＋医療事務に、",Inches(0.88),Inches(2.15),Inches(11.5),Inches(0.9),sz=38,bold=True,col=INK)
t(s,"同時に見せる台本",Inches(0.88),Inches(2.95),Inches(11.5),Inches(0.9),sz=38,bold=True,col=RED)
t(s,"クリニック全体が良くなるイメージを、先生にも事務にも持ってもらう。",Inches(0.9),Inches(3.95),Inches(11.5),Inches(0.5),sz=14,col=GRY)
cards=[("売り込まない","効果が出た分だけ・やめられる"),("一緒に測る","ベースライン→比較で公平・透明"),("小さく始める","まず1つ無料で")]
ox,ow,og=Inches(0.9),Inches(3.7),Inches(0.2)
for i,(ti,ds) in enumerate(cards):
    cx=ox+(ow+og)*i
    bx(s,cx,Inches(4.7),ow,Inches(1.4),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(4.7),ow,Inches(0.06),RED)
    t(s,ti,cx,Inches(4.9),ow,Inches(0.5),sz=17,bold=True,col=RED,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.2),Inches(5.45),ow-Inches(0.4),Inches(0.6),sz=11.5,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.1)
bx(s,Inches(0.9),Inches(6.55),Inches(11.5),Pt(1.2),LN)
t(s,"中核：相手目線でGIVE → 信頼の対価として収益（押し売りはしない）",Inches(0.9),Inches(6.65),Inches(11),Inches(0.4),sz=12.5,bold=True,col=INK)

# ════ SLIDE 2 — それぞれの見たい点 ════
s=sl(); ft(s)
hdr(s,"WHO SEES WHAT","それぞれの『見たい点』に、名指しで触る","医療事務＝日々の負担／先生＝経営と自分。どちらか一方だけの話にしない")
rows=[
 ("テーマ","医療事務が見たい点（日々の負担）","先生が見たい点（経営・自分）"),
 ("電話","鳴り止まない電話・クレーム対応が減る","人件費・離職リスクの低減"),
 ("受付","聞き取り・紙→電カル転記の二度手間が消える","患者満足（待ち時間・連絡）"),
 ("待ち時間","順番クレームが減る","院の回転・評判"),
 ("月末","まとめ会計の残業が減る","残業代・労務リスク"),
 ("自分","「私たちが楽になる」実感","先生ご自身の書類・連絡の手間も減る"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(4.0),[Inches(1.8),Inches(5.6),Inches(4.83)],hi_col=None,sz=12.5,header_sz=13)
t(s,"→ 両方に「あなたの困りごとが、ここで減ります」と直接触れるのがコツ。",Inches(0.55),Inches(6.1),Inches(12.2),Inches(0.4),sz=12.5,bold=True,col=REDD)

# ════ SLIDE 3 — 当日の流れ（6ステップ） ════
s=sl(); ft(s)
hdr(s,"FLOW","当日の流れ ── 6ステップの台本","つかみ → 実物 → 事務 → 先生 → 一緒に測る → 小さく始める")
steps=[("1","冒頭（共通）","「患者・先生・スタッフの三方が楽に。アンケートを読み込んで作りました」"),
       ("2","動く実物デモ","スマホでメニュー・問診・リマインドを実演。「掲示でなくスマホに届く」"),
       ("3","医療事務パート","電話・転記・待ち時間クレーム・月末残業が「ここで」減る（名指しで）"),
       ("4","先生パート","残業/人件費/離職リスク・患者満足・先生の書類手間・数字で見える"),
       ("5","一緒に測る（KPI）","導入前ベースライン→1〜2ヶ月で比較。指標は皆さんと相談して決める"),
       ("6","小さく始める","まず1つ無料で（再診リマインド or FAQ）。効果が出た分だけ・やめられる")]
cw,gx=Inches(3.95),Inches(0.24); x0=Inches(0.55)
for i,(no,ti,ds) in enumerate(steps):
    col=i%3; row=i//3
    cx=x0+(cw+gx)*col; cy=Inches(2.1)+(Inches(2.15))*row
    bx(s,cx,cy,cw,Inches(1.9),CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,cw,Inches(0.06),RED)
    bx(s,cx+Inches(0.22),cy+Inches(0.2),Inches(0.46),Inches(0.46),RED,shape=MSO_SHAPE.OVAL)
    t(s,no,cx+Inches(0.22),cy+Inches(0.2),Inches(0.46),Inches(0.46),sz=15,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ti,cx+Inches(0.8),cy+Inches(0.24),cw-Inches(0.9),Inches(0.5),sz=14,bold=True,col=INK)
    t(s,ds,cx+Inches(0.26),cy+Inches(0.86),cw-Inches(0.5),Inches(0.95),sz=10.5,col=GRY,line_sp=1.12)

# ════ SLIDE 4 — 楽になる中身（事務／先生 2カラム） ════
s=sl(); ft(s)
hdr(s,"THE VALUE","具体的に、ここが楽になる","事務＝日々／先生＝経営。最後は「まず1つ無料で」")
# 左：医療事務
bx(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.9),TEALBG,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(0.06),TEALD)
t(s,"医療事務（あなたたちが楽に）",Inches(0.8),Inches(2.12),Inches(5.5),Inches(0.4),sz=15,bold=True,col=TEALD)
for i,ln in enumerate(["① 鳴り止まない電話 → FAQ自動応答が一次対応","② 聞き取り・電カル転記 → Web問診で来院前に取得","③ 待ち時間クレーム → 順番をLINEで自動通知","④ 月末まとめ会計の残業 → 定期受診リマインド"]):
    t(s,ln,Inches(0.85),Inches(2.7)+Inches(0.72)*i,Inches(5.5),Inches(0.65),sz=12,col=INK,line_sp=1.1)
# 右：先生
bx(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.9),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(0.06),RED)
t(s,"先生（経営とご自身）",Inches(7.03),Inches(2.12),Inches(5.5),Inches(0.4),sz=15,bold=True,col=REDD)
for i,ln in enumerate(["① 残業・人件費・離職リスクの低減","② 患者満足（待ち時間・連絡のスムーズさ）","③ 先生ご自身の書類・連絡の手間も減る","④ 効果が「数字」で見える（KPI測定）"]):
    t(s,ln,Inches(7.08),Inches(2.7)+Inches(0.72)*i,Inches(5.5),Inches(0.65),sz=12,col=INK,line_sp=1.1)
bx(s,Inches(0.55),Inches(6.1),Inches(12.23),Inches(0.75),REDBG); bx(s,Inches(0.55),Inches(6.1),Inches(0.1),Inches(0.75),RED)
t(s,"クロージング：まず1つだけ無料で（再診リマインド or FAQ）→ 導入前のベースライン測定からスタート。",Inches(0.82),Inches(6.1),Inches(11.8),Inches(0.75),sz=13,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("kyobashi_dabon.pptx")
print("saved kyobashi_dabon.pptx slides:",len(prs.slides._sldIdLst))

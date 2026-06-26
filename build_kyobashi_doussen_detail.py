"""
京橋 集患→来院→再診 動線 詳細設計（前後比較）。
BEFORE/AFTERフロー・緻密比較表・流入ファネル・HP/MEOワイヤーフレーム(さいとう内科TTP)・問診/配信の細部。
クリーム白×レンガ赤。出力: kyobashi_doussen_detail.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
BLUEBG=RGBColor(0xE6,0xF1,0xFB); BLUED=RGBColor(0x18,0x5F,0xA5)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.55),Inches(0.35),Inches(12),Inches(0.4),sz=12.5,bold=True,col=RED)
    bx(slide,Inches(0.57),Inches(0.72),Inches(1.6),Pt(3),RED)
    t(slide,main,Inches(0.55),Inches(0.82),Inches(12.2),Inches(0.55),sz=20,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.57),Inches(1.36),Inches(12.2),Inches(0.3),sz=11,col=GRY)
def ft(slide,n):
    bx(slide,Inches(0.5),H-Inches(0.45),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック 集患→来院→再診 動線 詳細設計（前後比較）",Inches(0.5),H-Inches(0.4),Inches(10),Inches(0.3),sz=8.5,col=GRY)
    t(slide,n,Inches(12.45),H-Inches(0.4),Inches(0.4),Inches(0.3),sz=8.5,col=GRY)
def flowrow(slide,y,steps,accent,tagcol,tagbg):
    cw,gx,x0=Inches(2.18),Inches(0.22),Inches(0.55)
    for i,(ti,ds,tag) in enumerate(steps):
        cx=x0+(cw+gx)*i
        bx(slide,cx,y,cw,Inches(1.0),CARD,line=CARDLN,lw=1.0); bx(slide,cx,y,cw,Inches(0.06),accent)
        t(slide,ti,cx+Inches(0.12),y+Inches(0.12),cw-Inches(0.24),Inches(0.45),sz=12,bold=True,col=INK,align=PP_ALIGN.CENTER,line_sp=1.0)
        t(slide,ds,cx+Inches(0.12),y+Inches(0.55),cw-Inches(0.24),Inches(0.4),sz=9,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.05)
        if i<4: t(slide,"→",cx+cw-Inches(0.04),y+Inches(0.3),Inches(0.3),Inches(0.4),sz=15,bold=True,col=accent,align=PP_ALIGN.CENTER)
        bx(slide,cx,y+Inches(1.06),cw,Inches(0.5),tagbg)
        t(slide,tag,cx+Inches(0.1),y+Inches(1.06),cw-Inches(0.2),Inches(0.5),sz=9.5,bold=True,col=tagcol,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.0)
def tbl(slide,rows,x,y,w,h,col_w,sz=10.5,hsz=11):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            c=tb.cell(ri,ci); c.text=str(val); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.05); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            c.fill.solid()
            if ri==0: c.fill.fore_color.rgb=RED
            else: c.fill.fore_color.rgb=(CARD if ri%2==1 else BG)
            for p in c.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(hsz if ri==0 else sz); r.font.bold=(ri==0 or ci==0)
                    r.font.color.rgb=(WHT if ri==0 else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
    return tb

# 1 表紙
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"京橋クリニック 御中",Inches(0.9),Inches(1.15),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"集患 → 来院 → 再診\n動線 詳細設計（前後比較）",Inches(0.88),Inches(1.7),Inches(11.7),Inches(1.7),sz=27,bold=True,col=INK,line_sp=1.12)
t(s,"斎藤内科さんの構造をTTPし、人の流入が「前後」でどう変わるかを、緻密に設計しました。",Inches(0.92),Inches(3.8),Inches(11.4),Inches(0.5),sz=13.5,col=GRY)
bx(s,Inches(0.9),Inches(4.6),Inches(11.5),Inches(1.15),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(4.6),Inches(0.1),Inches(1.15),RED)
t(s,"見つける（HP/MEO）→ つながる（LINE）→ 来院前に問診 → 受付（アイコール）→ 再診（自動）。\n各段の「取りこぼし」と「手間」を、1つずつ潰します。",Inches(1.2),Inches(4.6),Inches(11),Inches(1.15),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)

# 2 BEFORE
s=sl(); ft(s,"2")
hdr(s,"いまの動線（BEFORE）","どこで「取りこぼす」か","各段に詰まりがある。だから新患が増えず、再診が抜け、残業になる。")
flowrow(s,Inches(2.05),[
 ("見つけにくい","HPが弱い・MEO未整備","✗ 新患が来ない"),
 ("電話が殺到","問い合わせが電話に集中","✗ 受付パンク"),
 ("紙の問診→転記","同じことを二度書く","✗ 手間・紙の山"),
 ("順番取っても来ない","遅刻・番号ゴネ","✗ クレーム対応"),
 ("再診が来ない","架電が大変・不在多い","✗ 月末残業"),
],GRY,REDD,REDBG)
bx(s,Inches(0.55),Inches(4.35),Inches(12.23),Inches(1.55),REDBG,line=RED,lw=1.0); bx(s,Inches(0.55),Inches(4.35),Inches(0.1),Inches(1.55),RED)
t(s,"BEFORE の正体",Inches(0.82),Inches(4.45),Inches(11.6),Inches(0.35),sz=13,bold=True,col=REDD)
t(s,"・入口（認知）が弱く、流入そのものが細い。 ・つながる手段が「来院時だけ」で、来院前後に手が届かない。\n・問診も再診も人海戦術＝面談シートの「外来19時／CPAP未来院で月末残業」がこれ。",Inches(0.85),Inches(4.82),Inches(11.6),Inches(1.0),sz=11.5,col=INK,line_sp=1.3)

# 3 AFTER
s=sl(); ft(s,"3")
hdr(s,"新しい動線（AFTER）","流入が増え、「ループ」が回る","入口を太くし、来院前後をLINEでつなぐ。順番(アイコール)は変えない。")
flowrow(s,Inches(2.05),[
 ("見つかる","HP・MEO（斎藤内科型）","○ 新患の流入↑"),
 ("LINEでつながる","HP/院内QR/問診から友だち","○ 来院前から接点"),
 ("来院前にWeb問診","記入済みで来院","○ 転記ゼロ・紙減"),
 ("受付はアイコール","順番はそのまま＋問診届く","○ 受付が軽い"),
 ("次回予約→自動再診","予約LINE＋リマインド/CPAP","○ 再診がループ"),
],TEALD,TEALD,TEALBG)
bx(s,Inches(0.55),Inches(4.35),Inches(12.23),Inches(1.55),TEALBG,line=TEALD,lw=1.0); bx(s,Inches(0.55),Inches(4.35),Inches(0.1),Inches(1.55),TEALD)
t(s,"AFTER の勘所",Inches(0.82),Inches(4.45),Inches(11.6),Inches(0.35),sz=13,bold=True,col=TEALD)
t(s,"・「見つける→つながる→来院→再診」が1本の線に。 ・離脱しかけの人もLINEで掘り起こし＝再診ループ。\n・アイコールは止めない。LINEは順番に触らず、来院「前」と「後」だけを担う。",Inches(0.85),Inches(4.82),Inches(11.6),Inches(1.0),sz=11.5,col=INK,line_sp=1.3)

# 4 緻密 前後比較表
s=sl(); ft(s,"4")
hdr(s,"何をすると、どう変わるか（前後比較）","場面ごとに「いま → やること → こう変わる」","ここが今回のキモ。1行ずつ潰す。")
rows=[
 ("場面","いま（BEFORE）","やること","こう変わる（AFTER）"),
 ("見つけてもらう","HP弱い・MEO未整備","HP刷新＋Googleビジネス最適化","指名検索・地図流入で新患増"),
 ("問い合わせ","電話に集中・受付パンク","LINE FAQ自動応答＋導線","一次対応をLINEが肩代わり・電話減"),
 ("友だち化","接点が来院時だけ","HP/院内QR/問診リンクで誘導","来院前から接点・後も追える"),
 ("問診","紙→転記（二度手間）","Web問診(一般/ESS/健診/化学物質)","来院前記入・転記ゼロ・紙減"),
 ("本人確認","マイナの紙コピー","事前にデジタル取得","紙と手間が減る（事務の声）"),
 ("当日受付","アイコール(来ない/遅刻)","アイコール継続＋順番通知は将来","受付は変えず、案内だけ届く"),
 ("再診(定期/CPAP)","来ない→月末残業","次回予約LINE＋CPAP自動送付","取りこぼし減・残業を平準化"),
 ("離脱者","事務がザオラル架電(不在多)","LINEで次回予約を掘り起こし","架電減・不在でも予約が取れる"),
]
tbl(s,rows,Inches(0.55),Inches(1.78),Inches(12.23),Inches(4.5),[Inches(1.95),Inches(3.2),Inches(3.35),Inches(3.73)],sz=10,hsz=11)

# 5 流入ファネル前後
s=sl(); ft(s,"5")
hdr(s,"流入の「数」が、どう変わるか","認知 → 接点 → 来院(新患) → 再診/継続","各段の「漏れ」を施策で塞ぐ＝同じ広告費でも残る人が増える")
stages=[("認知","HP/MEOで見つかる"),("接点(友だち)","LINE登録で離脱を捕捉"),("来院(新患)","問診で来院ハードル↓"),("再診/継続","予約・リマインドで継続")]
cw,gx,x0=Inches(2.95),Inches(0.18),Inches(0.55)
for i,(ti,ds) in enumerate(stages):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.1),cw,Inches(1.5),BLUEBG,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.1),cw,Inches(0.06),BLUED)
    t(s,ti,cx+Inches(0.12),Inches(2.25),cw-Inches(0.24),Inches(0.5),sz=14,bold=True,col=BLUED,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.15),Inches(2.78),cw-Inches(0.3),Inches(0.7),sz=10.5,col=INK,align=PP_ALIGN.CENTER,line_sp=1.2)
    if i<3: t(s,"▶",cx+cw-Inches(0.02),Inches(2.55),Inches(0.3),Inches(0.5),sz=14,bold=True,col=BLUED,align=PP_ALIGN.CENTER)
bx(s,Inches(0.55),Inches(3.95),Inches(6.0),Inches(2.0),REDBG); bx(s,Inches(0.55),Inches(3.95),Inches(0.1),Inches(2.0),RED)
t(s,"BEFORE（漏れが多い）",Inches(0.8),Inches(4.05),Inches(5.5),Inches(0.35),sz=12.5,bold=True,col=REDD)
t(s,"・認知が細い（指名検索/地図に出ない）\n・接点が来院時だけ→一見客が消える\n・紙問診で来院ハードル\n・再診は人力架電→不在で抜ける",Inches(0.82),Inches(4.45),Inches(5.5),Inches(1.45),sz=11,col=INK,line_sp=1.35)
bx(s,Inches(6.78),Inches(3.95),Inches(6.0),Inches(2.0),TEALBG); bx(s,Inches(6.78),Inches(3.95),Inches(0.1),Inches(2.0),TEALD)
t(s,"AFTER（漏れを塞ぐ）",Inches(7.03),Inches(4.05),Inches(5.5),Inches(0.35),sz=12.5,bold=True,col=TEALD)
t(s,"・HP/MEOで認知の入口を太く\n・LINE登録で一見客を捕捉→追える\n・Web問診で来院ハードル↓\n・予約/リマインド自動で再診が残る",Inches(7.03),Inches(4.45),Inches(5.5),Inches(1.45),sz=11,col=INK,line_sp=1.35)

# 6 HP/MEO ワイヤーフレーム
s=sl(); ft(s,"6")
hdr(s,"HP・MEO レイアウト設計","斎藤内科の構造をTTP → 京橋版ワイヤーフレーム","左＝ページの並び（上から）／右＝なぜ効くか（TTPの勝ち所）")
px,pw=Inches(0.55),Inches(5.3)
bx(s,px,Inches(1.85),pw,Inches(5.0),WHT,line=CARDLN,lw=1.2)
secs=[("ヘッダー：ロゴ ｜ ☎予約専用 ｜ [WEB予約][問診][LINE]",RED,WHT),
      ("ヒーロー：1メッセージ ＋ [予約する]大ボタン",CARD,INK),
      ("診療時間・アクセス（住所/電話/最寄駅 徒歩◯分）",GRYBG,INK),
      ("SNS / LINE 連携バー",CARD,INK),
      ("お知らせ（ブログ・通信＝月次更新）",GRYBG,INK),
      ("当院の特徴 ×4 ／ 診療科 カード ×9",CARD,INK),
      ("院長あいさつ・経歴／提携病院",GRYBG,INK),
      ("地図（NAP再）／フッター（NAP再）",CARD,INK)]
yy=Inches(1.98)
for txt,fill,tc in secs:
    bx(s,px+Inches(0.12),yy,pw-Inches(0.24),Inches(0.54),fill,line=CARDLN,lw=0.75)
    t(s,txt,px+Inches(0.24),yy,pw-Inches(0.45),Inches(0.54),sz=9.5,bold=(fill==RED),col=tc,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.0)
    yy=yy+Inches(0.6)
bx(s,px+pw-Inches(1.2),Inches(6.05),Inches(1.0),Inches(0.5),RED,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
t(s,"予約",px+pw-Inches(1.2),Inches(6.05),Inches(1.0),Inches(0.5),sz=10,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
t(s,"↑ 右下に常時フローティングCTA",px+Inches(0.12),Inches(6.55),pw,Inches(0.3),sz=8.5,col=GRY)
ax=Inches(6.2)
wins=[("予約専用ダイヤルをヘッダー独立","迷わせない。電話したい人を取りこぼさない"),
      ("WEB予約/問診/LINE＋右下フローティング","どの画面でも予約・問診に1タップ"),
      ("NAP（住所/電話/最寄駅）を3回以上","Googleマップ上位＝MEOの基本。地図流入を取る"),
      ("診療科を9つカード化","検索キーワードの入口を増やす（流入↑）"),
      ("院長経歴・提携病院・ブログ月次","信頼＋新鮮さ＝指名検索とクチコミに効く")]
yy=Inches(1.95)
for ti,ds in wins:
    bx(s,ax,yy,Inches(6.6),Inches(0.92),CARD,line=CARDLN,lw=1.0); bx(s,ax,yy,Inches(0.09),Inches(0.92),BLUED)
    t(s,ti,ax+Inches(0.2),yy+Inches(0.1),Inches(6.3),Inches(0.4),sz=11.5,bold=True,col=BLUED)
    t(s,ds,ax+Inches(0.22),yy+Inches(0.5),Inches(6.2),Inches(0.38),sz=10,col=INK,line_sp=1.1)
    yy=yy+Inches(1.0)

# 7 問診の細部
s=sl(); ft(s,"7")
hdr(s,"問診の細部","種類で分岐 → 必須/任意 → 自動採点 → 送信後の処理","ここまで決めておけば、GO後すぐ作れる")
rows=[
 ("問診の種類","主な項目","仕掛け"),
 ("一般","症状/発症/体温/既往・家族歴/アレルギー/喫煙飲酒/妊娠/ジェネリック","「当院を何で知ったか」を集患タグに自動分類"),
 ("睡眠時無呼吸(ESS)","身長体重＋8状況の眠気(0〜3)","合計を自動採点→11点以上は受付に印"),
 ("健診","健診メニュー別の事前確認","予約フォームと連結"),
 ("化学物質過敏症","専用問診（症状・経過）","送信→先生に電話フォロー依頼が飛ぶ"),
]
tbl(s,rows,Inches(0.55),Inches(1.8),Inches(12.23),Inches(2.95),[Inches(2.7),Inches(5.6),Inches(3.93)],sz=10.5,hsz=11)
bx(s,Inches(0.55),Inches(4.95),Inches(12.23),Inches(1.5),GRYBG); bx(s,Inches(0.55),Inches(4.95),Inches(0.1),Inches(1.5),RED)
t(s,"送信後の自動処理（共通）",Inches(0.82),Inches(5.05),Inches(11.6),Inches(0.35),sz=12.5,bold=True,col=REDD)
t(s,"① 患者さんへ自動お礼　② タグ付け（問診記入済／CPAP候補／集患経路 等）　③ 受付の一覧に届く（氏名・要点）\n④ 将来：電子カルテ(Medicom)へ取り込み＝二度入力を消す（連携の深さは要確認）。",Inches(0.85),Inches(5.42),Inches(11.6),Inches(1.0),sz=11,col=INK,line_sp=1.3)

# 8 配信の細部
s=sl(); ft(s,"8")
hdr(s,"お知らせ・リマインドの細部","いつ・誰に・何を・どんな文面で","押し売らない。お知らせ止まりで、判断は人。")
rows=[
 ("配信","対象（タグ）","タイミング","文面の方針"),
 ("再診リマインド(ON)","再診・定期","次回目安の数日前","来院の目安をそっと。お知らせ止まり"),
 ("CPAP 次回予約","CPAP候補/装着中","受診後＋定期","次回予約をご案内（自動）"),
 ("再診促し","最終来院から一定経過","期間後に1回","次回予約の導線。催促でなく案内"),
 ("健診・化学物質","該当・申込者","申込/案内時","専用問診・予約フォームへ誘導"),
 ("お知らせ/休診","友だち全員","随時","掲示の代わりに。読まれる形で"),
]
tbl(s,rows,Inches(0.55),Inches(1.8),Inches(12.23),Inches(3.4),[Inches(2.6),Inches(3.0),Inches(3.0),Inches(3.63)],sz=10.5,hsz=11)
bx(s,Inches(0.55),Inches(5.4),Inches(12.23),Inches(1.0),REDBG); bx(s,Inches(0.55),Inches(5.4),Inches(0.1),Inches(1.0),RED)
t(s,"※ どれも「お知らせ」。医療・予約可否の判断は、これまでどおり人が行います（中核：売り込まない／患者目線）。",Inches(0.82),Inches(5.4),Inches(11.8),Inches(1.0),sz=12,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# 9 まとめ＋GO
s=sl(); ft(s,"9")
hdr(s,"まとめ ＆ GO判断","この詳細設計で、作り込みに入ってよいか","GOで、まず問診から実物を作ります")
bx(s,Inches(0.55),Inches(1.85),Inches(12.23),Inches(1.5),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(1.85),Inches(0.12),Inches(1.5),TEALD)
t(s,"設計のまとめ",Inches(0.82),Inches(1.95),Inches(11.6),Inches(0.35),sz=13,bold=True,col=TEALD)
t(s,"入口を太く（HP/MEO）→ LINEでつなぐ → 来院前に問診 → 受付はアイコールのまま → 再診を自動ループ。\n各段の「取りこぼし」と「手間」を、前後比較表のとおり1行ずつ潰す。",Inches(0.85),Inches(2.32),Inches(11.6),Inches(1.0),sz=12,bold=True,col=INK,line_sp=1.3)
bx(s,Inches(0.55),Inches(3.6),Inches(6.0),Inches(2.4),REDBG,line=RED,lw=1.0); bx(s,Inches(0.55),Inches(3.6),Inches(0.1),Inches(2.4),RED)
t(s,"先に教えてほしい2点",Inches(0.82),Inches(3.7),Inches(5.5),Inches(0.4),sz=13,bold=True,col=REDD)
t(s,"①「最後の1週間」に架電、の運用の意味\n②サインエコーのLINE導入、とは？",Inches(0.82),Inches(4.15),Inches(5.5),Inches(1.6),sz=12,col=INK,line_sp=1.4)
bx(s,Inches(6.78),Inches(3.6),Inches(6.0),Inches(2.4),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(3.6),Inches(0.1),Inches(2.4),RED)
t(s,"GOで着手する順",Inches(7.03),Inches(3.7),Inches(5.5),Inches(0.4),sz=13,bold=True,col=REDD)
t(s,"P1：問診（一般・ESS）の実物→受付に届く\nP2：リマインド/CPAP/再診促し\nP3：HP・MEO（斎藤内科TTP）",Inches(7.03),Inches(4.15),Inches(5.5),Inches(1.6),sz=12,col=INK,line_sp=1.4)
t(s,"この詳細設計でGO、でよろしいですか？（直す所があれば、その行だけ言ってください）",Inches(0.55),Inches(6.2),Inches(12.2),Inches(0.4),sz=12.5,bold=True,col=REDD)

prs.save("kyobashi_doussen_detail.pptx")
print("saved kyobashi_doussen_detail.pptx slides:",len(prs.slides._sldIdLst))

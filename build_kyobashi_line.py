"""
京橋クリニック LINEハーネス 実装案デッキ（7/1当日持参用）
アンケート(260602_その他.pdf 事務3+看護師3)直結。クリーム白×レンガ赤。
売り込まずGIVE・アンダーコミット（アイコール連携は要確認）。
出力: kyobashi_line_jissouan.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
LNG    = RGBColor(0x06, 0xC7, 0x55)   # LINEグリーン
LNGD   = RGBColor(0x06, 0x7A, 0x35)

FONT = "Hiragino Sans"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "京橋クリニック  ｜  LINE活用 実装案  ｜  AI医療コンサル", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)


def light_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = REDD if is_hi else RED
            else:
                cell.fill.fore_color.rgb = REDBG if is_hi else (CARD if ri % 2 == 1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = WHT
                    elif is_hi:
                        r.font.color.rgb = RED
                    elif ci == 0:
                        r.font.color.rgb = INK
                    else:
                        r.font.color.rgb = RGBColor(0x3A,0x3A,0x3A)
    return tb


def line_phone(slide, x, y, w, h, title, bubbles):
    """LINE風スマホモック。bubbles=[(text, is_clinic_right_green)]"""
    bx(slide, x, y, w, h, LNG)
    bx(slide, x+Inches(0.12), y+Inches(0.12), w-Inches(0.24), h-Inches(0.24), WHT)
    bx(slide, x+Inches(0.12), y+Inches(0.12), w-Inches(0.24), Inches(0.5), LNGD)
    t(slide, title, x+Inches(0.28), y+Inches(0.12), w-Inches(0.5), Inches(0.5),
      sz=11, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
    by = y+Inches(0.78)
    for txt, right in bubbles:
        nlines = max(1, (len(txt)) // 14 + 1)
        bh = Inches(0.34) + Inches(0.2)*(nlines-1)
        bw = w-Inches(1.1)
        bxx = x+(w-bw-Inches(0.3)) if right else x+Inches(0.3)
        bx(slide, bxx, by, bw, bh, RGBColor(0x9A,0xE5,0xB0) if right else GRYBG)
        t(slide, txt, bxx+Inches(0.12), by, bw-Inches(0.24), bh, sz=9, col=INK, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.0)
        by = by + bh + Inches(0.14)


def arrow(slide, x, y, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xC9,0x6A,0x62)
    s.line.fill.background(); s.shadow.inherit = False
    return s


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "LINE × CLINIC", Inches(0.9), Inches(1.55), Inches(7.2), Inches(0.45), sz=15, bold=True, col=RED)
t(s, "アンケートの声を、", Inches(0.88), Inches(2.1), Inches(7.4), Inches(0.9), sz=40, bold=True, col=INK)
t(s, "LINEで減らす", Inches(0.88), Inches(2.9), Inches(7.4), Inches(0.9), sz=40, bold=True, col=RED)
t(s, "医療事務・看護師アンケート（事務3名・看護師3名）が示した\n残業の主因を、LINE導線で物理的に減らす実装案。初期構築は当社が巻き取り。",
  Inches(0.9), Inches(3.9), Inches(7.3), Inches(0.9), sz=13.5, col=GRY, line_sp=1.25)
ox, oy, ow, og = Inches(0.9), Inches(5.0), Inches(2.3), Inches(0.16)
offers = [("初期構築", "当社負担"), ("まず試す", "再診\nリマインド"), ("院の手間", "ゼロ")]
for i, (lab, val) in enumerate(offers):
    cx = ox + (ow + og) * i
    bx(s, cx, oy, ow, Inches(1.35), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, oy, ow, Inches(0.06), RED)
    t(s, lab, cx, oy+Inches(0.16), ow, Inches(0.35), sz=12, col=GRY, align=PP_ALIGN.CENTER)
    vsz = 20 if "\n" in val else 26
    t(s, val, cx, oy+Inches(0.5), ow, Inches(0.8), sz=vsz, bold=True, col=RED, align=PP_ALIGN.CENTER, line_sp=1.0)
bx(s, Inches(0.9), Inches(6.7), Inches(11.5), Pt(1.2), LINE)
t(s, "チームてっかん  ｜  菊池 研太  ｜  2026.06", Inches(0.9), Inches(6.82), Inches(11), Inches(0.4), sz=13, bold=True, col=INK)
# 右：LINEモック（再診リマインド）
line_phone(s, Inches(9.0), Inches(1.55), Inches(3.6), Inches(4.6), "京橋クリニック",
           [("そろそろCPAPの定期診察の時期です", True),
            ("ご予約はこちらから ▶", True),
            ("6/28(土) 10:00 で予約しました", False),
            ("当日の順番が近づくとお知らせします", True)])

# ════════ SLIDE 2 — アンケートが示す残業の主因 ════════
s = sl(); ft(s)
hdr(s, "THE VOICE", "アンケートが示す、残業の3つの主因", "事務3名・看護師3名の生の声（2026年6月実施／対応責任者 福井）")
issues = [
    ("01", "順番受付（アイコール）が機能不全", "「順番を取っても受付時間内に来ない」「注意書きを読まない」\n→ 受付が毎回クレーム対応に追われる。"),
    ("02", "CPAPの通院離脱で月末が爆発", "「毎月来る同意のはずが来ない」\n→ 月末最終日にまとめて会計 → 残業が大幅に増える。"),
    ("03", "説明しても、読まれない・伝わらない", "「説明も掲示もしているのに読まない患者が多すぎる」\n→ 紙・掲示では届かない。"),
]
cw, ch, gx = Inches(3.95), Inches(4.3), Inches(0.24)
x0, y0 = Inches(0.55), Inches(1.95)
for i, (no, ti, desc) in enumerate(issues):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.12), RED)
    t(s, no, cx+Inches(0.3), y0+Inches(0.28), Inches(1.5), Inches(0.8), sz=34, bold=True, col=RED)
    t(s, ti, cx+Inches(0.3), y0+Inches(1.15), cw-Inches(0.6), Inches(1.1), sz=16, bold=True, col=INK, line_sp=1.1)
    bx(s, cx+Inches(0.3), y0+Inches(2.35), cw-Inches(0.6), Pt(1.2), LINE)
    t(s, desc, cx+Inches(0.3), y0+Inches(2.55), cw-Inches(0.6), Inches(1.6), sz=12, col=GRY, line_sp=1.25)
band_y = y0 + ch + Inches(0.18)
bx(s, Inches(0.55), band_y, Inches(11.83), Inches(0.5), REDBG)
bx(s, Inches(0.55), band_y, Inches(0.1), Inches(0.5), RED)
t(s, "＋ HP導線：外来受付・化学物質過敏症の問い合わせが「電話して」導線になっている → 事務の電話対応をさらに増やしている。",
  Inches(0.8), band_y, Inches(11.5), Inches(0.5), sz=11.5, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)

# ════════ SLIDE 3 — LINEでこう減らす（before→after） ════════
s = sl(); ft(s)
hdr(s, "THE SOLUTION", "その主因、LINEでこう減らせます", "アンケートの声（Before）→ LINE導線（After）。重い見積りでなく、まず1つ試す前提")
rows = [
    ("機能", "Before（アンケートの声）", "After（LINE導線で）"),
    ("① 再診リマインド ★まず1つ", "CPAPが毎月来ない→月末まとめ会計で残業爆発", "予約日前にLINE自動リマインド→来院率UP→月末残業を平準化（効果が数字で出る）"),
    ("② 順番・待ち時間の通知", "順番を取っても来ない／注意を読まない→受付がクレーム", "順番接近で「あと◯人／受付期限◯分前」を自動プッシュ→受付の口頭対応を削減"),
    ("③ 周知・HP導線の整理", "説明・掲示しても読まない／HP問い合わせが電話導線", "リッチメニューに受付ルール・CPAP説明を集約。HPは電話→LINE/Web完結へ改修"),
]
light_table(s, rows, Inches(0.55), Inches(1.95), Inches(12.23), Inches(4.0),
            [Inches(3.0), Inches(4.4), Inches(4.83)], hi_col=2, sz=12, header_sz=13)
t(s, "★ 最初の一手は①再診リマインド。来院率と月末残業の変化を1〜2ヶ月測れば、効果が「数字」で見える＝成果報酬の入口。",
  Inches(0.55), Inches(6.1), Inches(12.2), Inches(0.5), sz=12, bold=True, col=REDD)

# ════════ SLIDE 4 — 実装の構成 ════════
s = sl(); ft(s)
hdr(s, "ARCHITECTURE", "実装の構成 ── 何を、どう繋ぐか", "市販ツールの組合せで実現。特別な開発は不要。初期構築・設定は当社が巻き取り")
comps = [
    ("患者のスマホ", "LINE", "順番通知・再診リマインド・\n受付ルールを受け取る", LNG, WHT),
    ("LINE公式アカウント", "無料開設", "クリニックの公式窓口。\nリッチメニューで常設案内", RED, WHT),
    ("配信ツール", "エルメ / Lステップ", "CPAP等をタグ分け→\n予約日リマインドを自動配信", CARD, INK),
    ("流入導線", "HP・院内QR", "HPと受付QRからLINE登録。\n将来HPリプレイス", CARD, INK),
]
cw, gx, x0, y0 = Inches(2.78), Inches(0.34), Inches(0.55), Inches(2.15)
for i, (lab, val, desc, fill, fg) in enumerate(comps):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(2.7), CARD if fill in (CARD,) else BG, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.85), fill)
    t(s, lab, cx, y0+Inches(0.12), cw, Inches(0.35), sz=12, bold=True, col=fg, align=PP_ALIGN.CENTER)
    t(s, val, cx, y0+Inches(0.45), cw, Inches(0.4), sz=14, bold=True, col=fg, align=PP_ALIGN.CENTER)
    t(s, desc, cx+Inches(0.2), y0+Inches(1.05), cw-Inches(0.4), Inches(1.5), sz=11, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.2)
    if i < 3:
        arrow(s, cx+cw+Inches(0.02), y0+Inches(1.05), Inches(0.3), Inches(0.5))
by = y0 + Inches(2.95)
bx(s, Inches(0.55), by, Inches(11.83), Inches(0.55), GRYBG)
bx(s, Inches(0.55), by, Inches(0.1), Inches(0.55), RED)
t(s, "⚠ 順番通知（②）は、既存「アイコール」のLINE連携可否を要確認。連携できない場合は受付ボタンで対象者へ通知する簡易運用で代替します。",
  Inches(0.8), by, Inches(11.5), Inches(0.55), sz=11.5, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)

# ════════ SLIDE 5 — 段階導入 ════════
s = sl(); ft(s)
hdr(s, "ROADMAP", "小さく速く始める ── 4段階", "いきなり全部やらない。効果が出る順に、無料で着手できる所から")
rows = [
    ("フェーズ", "やること", "ねらい・効果"),
    ("Phase 0（即・ほぼ0円）", "LINE公式開設＋リッチメニュー＋HP/受付にLINE登録QR", "「読まない」をプッシュで届ける土台をつくる"),
    ("Phase 1（再診リマインド）★", "CPAP・再診をLINEで分類→月次リマインドを自動配信", "来院率UP→月末残業を削減。効果測定を開始（成果報酬の入口）"),
    ("Phase 2（順番通知）", "アイコール連携を確認→順番・待ち時間のLINE通知", "受付のクレーム・口頭対応を削減"),
    ("Phase 3（HPリプレイス）", "HP全面改修・Web問診/予約・マーケ巻き取り", "電話導線を解消し、集患まで一気通貫"),
]
light_table(s, rows, Inches(0.55), Inches(1.95), Inches(12.23), Inches(4.3),
            [Inches(3.3), Inches(5.3), Inches(3.63)], hi_col=2, sz=12, header_sz=13)

# ════════ SLIDE 6 — 始め方 / クロージング ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "HOW TO START", Inches(0.9), Inches(0.7), Inches(11), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "まず、再診リマインド「1つ」から。", Inches(0.9), Inches(1.15), Inches(11.7), Inches(0.8), sz=30, bold=True, col=INK)
t(s, "全部を一度にやりません。効果が数字で出る①だけを無料で試して、続きはそれから決めます。",
  Inches(0.92), Inches(2.0), Inches(11.5), Inches(0.5), sz=14, col=GRY)
# 左：始め方カード
cards = [
    ("初期構築は当社が巻き取り", "LINE公式の開設・配信設定・文面づくりまで当社が実施。先生・事務の手間はゼロ。"),
    ("効果を1〜2ヶ月測る", "来院率と月末残業の変化を可視化。「数字」で続ける/やめるを判断できる。"),
    ("着手はほぼ0円", "LINE公式（無料枠）＋配信ツール無料プランで開始。本格運用でも月5千〜2万円程度。"),
]
for i, (ti, ds) in enumerate(cards):
    cy = Inches(2.7) + Inches(1.18) * i
    bx(s, Inches(0.9), cy, Inches(7.1), Inches(1.02), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.9), cy, Inches(0.06), Inches(1.02), RED)
    t(s, "✓ " + ti, Inches(1.15), cy+Inches(0.12), Inches(6.7), Inches(0.4), sz=15, bold=True, col=RED)
    t(s, ds, Inches(1.15), cy+Inches(0.5), Inches(6.7), Inches(0.45), sz=11.5, col=GRY, line_sp=1.05)
# 右：役割分担
bx(s, Inches(8.3), Inches(2.7), Inches(4.15), Inches(3.56), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(8.3), Inches(2.7), Inches(4.15), Inches(0.06), RED)
t(s, "役割分担", Inches(8.3), Inches(2.9), Inches(4.15), Inches(0.4), sz=15, bold=True, col=RED, align=PP_ALIGN.CENTER)
bx(s, Inches(8.55), Inches(3.5), Inches(3.65), Inches(1.2), GRYBG)
t(s, "宮崎", Inches(8.7), Inches(3.6), Inches(3.4), Inches(0.35), sz=13, bold=True, col=INK)
t(s, "技術：Stream Deck／\n電子カルテからの自動抽出・入力連携", Inches(8.7), Inches(3.95), Inches(3.4), Inches(0.7), sz=11, col=GRY, line_sp=1.15)
bx(s, Inches(8.55), Inches(4.85), Inches(3.65), Inches(1.25), REDBG)
t(s, "菊池", Inches(8.7), Inches(4.95), Inches(3.4), Inches(0.35), sz=13, bold=True, col=RED)
t(s, "患者導線：LINE／HP／マーケ\n＋アンケートを「残業削減」に翻訳", Inches(8.7), Inches(5.3), Inches(3.4), Inches(0.7), sz=11, col=INK, line_sp=1.15)
bx(s, Inches(0.9), Inches(6.5), Inches(11.55), Pt(1.2), LINE)
t(s, "押し売りはしません。アンケートの声に、まず1つで応える形でご提案します。", Inches(0.9), Inches(6.62), Inches(11.5), Inches(0.4), sz=12.5, bold=True, col=INK)

prs.save("kyobashi_line_jissouan.pptx")
print("saved kyobashi_line_jissouan.pptx  /  slides:", len(prs.slides._sldIdLst))

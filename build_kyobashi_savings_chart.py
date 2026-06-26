# -*- coding: utf-8 -*-
"""
京橋クリニック AI医療コンサル｜削減効果グラフ＋新料金（先生提示用）
新モデル: 初期0 / 月額0 / 成果報酬40%×6ヶ月 → 7ヶ月目〜月5万保守(月1MTG込み)
KHD配色: クリーム白#F9F6EF × レンガ赤#AA2E26
"""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
import numpy as np

plt.rcParams["font.family"] = ["Hiragino Sans"]
plt.rcParams["axes.unicode_minus"] = False

CREAM="#F9F6EF"; RED="#AA2E26"; INK="#1A1A1A"; BEIGE="#F1ECE1"
GREEN="#2E7D32"; RULE="#DAD6CF"; GREY="#6B6B6B"; LIGHTRED="#E8C7C4"; STAYGREEN="#7BA77C"

# ---- 前提（事務3名・標準モデル）----
SAVE = 240000          # 月削減額(円)
RATE = 0.40            # 成果報酬率(1-6ヶ月)
MAINT = 50000          # 7ヶ月目以降の保守(月)
months = np.arange(1, 13)
pay = np.where(months <= 6, SAVE*RATE, MAINT)       # 院のお支払い
stay = SAVE - pay                                   # 院に残る
cum_stay = np.cumsum(stay)                          # 院に残る累計

man = 10000.0  # 万円換算
fig, ax = plt.subplots(figsize=(8.6, 4.7), dpi=200)
fig.patch.set_facecolor(CREAM); ax.set_facecolor(CREAM)

# stacked bars: 院に残る(下) + お支払い(上) = 削減額24万
b1 = ax.bar(months, stay/man, width=0.62, color=STAYGREEN, label="院に残る効果", zorder=3)
b2 = ax.bar(months, pay/man, bottom=stay/man, width=0.62, color=RED, label="お支払い", zorder=3)

# 6→7ヶ月の境界
ax.axvline(6.5, color=GREY, lw=1.0, ls=(0,(4,3)), zorder=2)
ax.text(6.5, 26.6, "成果報酬6ヶ月  →  7ヶ月目〜 保守 月5万", color=GREY,
        fontsize=9, ha="center", va="bottom")

# 院に残るの値ラベル（代表月）
for m in [1, 6, 7, 12]:
    ax.text(m, (stay[m-1]/man)/2, f"{stay[m-1]/man:.1f}万", color="white",
            fontsize=8.5, ha="center", va="center", fontweight="bold", zorder=4)

ax.set_ylim(0, 28)
ax.set_xticks(months); ax.set_xticklabels([f"{m}" for m in months], fontsize=9)
ax.set_xlabel("導入からの月数", fontsize=9, color=INK)
ax.set_ylabel("月額（万円）", fontsize=9, color=INK)
ax.tick_params(axis="y", labelsize=8)
for s in ["top","right"]: ax.spines[s].set_visible(False)
for s in ["left","bottom"]: ax.spines[s].set_color(RULE)
ax.grid(axis="y", color=RULE, lw=0.6, zorder=0)

# 累計線（第2軸）
ax2 = ax.twinx()
ax2.plot(months, cum_stay/man, color=INK, lw=2.2, marker="o", ms=4.5,
         label="院に残る（累計）", zorder=5)
ax2.scatter([12],[cum_stay[-1]/man], s=60, color=RED, zorder=6)
ax2.annotate(f"12ヶ月で\n院に約{cum_stay[-1]/man:.0f}万円",
             xy=(12, cum_stay[-1]/man), xytext=(9.0, 120),
             fontsize=10.5, color=RED, fontweight="bold", ha="center",
             arrowprops=dict(arrowstyle="->", color=RED, lw=1.4))
ax2.set_ylim(0, 230)
ax2.set_ylabel("院に残る効果（累計・万円）", fontsize=9, color=INK)
ax2.tick_params(axis="y", labelsize=8)
for s in ["top"]: ax2.spines[s].set_visible(False)
ax2.spines["right"].set_color(RULE)

# 凡例（まとめ）
h1,l1 = ax.get_legend_handles_labels()
h2,l2 = ax2.get_legend_handles_labels()
ax.legend(h1+h2, l1+l2, loc="upper left", fontsize=8.5, frameon=False, ncol=1)

plt.tight_layout()
chart_png = "/Users/kikuchikenta/01_honbu_docs_automation/_savings_chart.png"
plt.savefig(chart_png, facecolor=CREAM, bbox_inches="tight")
print("chart saved:", chart_png)

# ============ スライド化（16:9・KHD配色） ============
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

cRED=RGBColor(0xAA,0x2E,0x26); cINK=RGBColor(0x1A,0x1A,0x1A); cCREAM=RGBColor(0xF9,0xF6,0xEF)
cBEIGE=RGBColor(0xF1,0xEC,0xE1); cWHITE=RGBColor(0xFF,0xFF,0xFF); cGREY=RGBColor(0x6B,0x6B,0x6B)
cGREEN=RGBColor(0x2E,0x7D,0x32); cRULE=RGBColor(0xDA,0xD6,0xCF)
FONT="Hiragino Sans"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
slide=prs.slides.add_slide(prs.slide_layouts[6])

def rect(x,y,w,h,fill=None,line=None,lw=1.0):
    sp=slide.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp
def txt(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.0):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=space
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.name=FONT; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
    return tb

rect(0,0,13.333,7.5,fill=cCREAM)
rect(0,0,0.14,7.5,fill=cRED)
txt(0.5,0.34,12.4,0.3,[[("SIMULATION ｜ 御院に残る削減効果",11,True,cRED)]])
txt(0.5,0.66,12.4,0.55,[[("使うほど、院にお金が残る ── 削減効果と料金",22,True,cINK)]])
txt(0.5,1.24,12.4,0.3,[[("標準モデル：医師1・事務3名（月の削減額 約24万円）／ ※当日、御院の実数で計算し直してお出しします",10.5,False,cGREY)]])
rect(0.5,1.58,12.33,0.03,fill=cRED)

# チャート画像
slide.shapes.add_picture(chart_png, Inches(0.45), Inches(1.8), width=Inches(8.5))

# 右：料金カード
RX=9.25; RW=3.6
rect(RX,1.85,RW,2.55,fill=cWHITE,line=cRULE,lw=1.0)
txt(RX+0.2,1.98,RW-0.4,0.3,[[("料金（リスクゼロ）",12,True,cRED)]])
rows=[("初期費用","0円"),("月額基本料","0円"),("成果報酬","削減額の40%"),("　成果報酬の期間","6ヶ月"),
      ("7ヶ月目以降","月5万円 保守"),("　保守に含む","月1回MTG・運用相談")]
yy=2.34
for k,v in rows:
    bold = not k.startswith("　")
    txt(RX+0.2,yy,1.75,0.28,[[(k,10 if bold else 9, bold, cINK if bold else cGREY)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(RX+1.95,yy,RW-2.15,0.28,[[(v,11 if bold else 9.5, True if bold else False, cRED if bold else cINK)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.32
txt(RX+0.2,yy+0.02,RW-0.4,0.3,[[("効果0なら費用0円・解約自由",9,False,cGREY)]])

# 右下：効果サマリ カード
rect(RX,4.55,RW,2.4,fill=cBEIGE)
txt(RX+0.2,4.68,RW-0.4,0.3,[[("御院に残る効果",12,True,cRED)]])
summ=[("1〜6ヶ月","月 14.4万円（60%）"),("7ヶ月目以降","月 19万円（24万−保守5万）"),("12ヶ月 累計","約 200万円")]
yy=5.06
for k,v in summ:
    txt(RX+0.2,yy,1.4,0.34,[[(k,9.5,True,cINK)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(RX+1.5,yy,RW-1.7,0.34,[[(v,11,True,cGREEN if "12" not in k else cRED)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.40
txt(RX+0.2,yy+0.04,RW-0.4,0.5,[[("成果報酬は最初の6ヶ月だけ。",9.5,True,cINK)],[("7ヶ月目からは月5万で、効果はまるごと院に。",9.5,False,cINK)]],space=1.1)

# 下部メッセージ
rect(0.45,6.95,8.5,0.012,fill=cRULE)
txt(0.45,7.02,8.5,0.3,[[("ポイント：成果報酬は6ヶ月で卒業。以降は月5万の保守だけで、削減効果はそのまま御院の手元に残り続けます。",10,True,cRED)]])

out="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/【先生提示用】京橋_削減効果と料金グラフ.pptx"
prs.save(out)
print("SAVED:", out)

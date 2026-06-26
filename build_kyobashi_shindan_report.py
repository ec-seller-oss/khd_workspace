# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); WHT=RGBColor(0xFF,0xFF,0xFF)
GRN=RGBColor(0x2E,0x7D,0x32); REDBG=RGBColor(0xF4,0xE4,0xE2)
FONT="Hiragino Sans"
W=Inches(13.333); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
SAV="/Users/kikuchikenta/01_honbu_docs_automation/_savings_chart.png"

def sl():
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    return s
def bx(s,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.shadow.inherit=False
    if col is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=col
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp
def t(s,x,y,w,h,text,sz=11,b=False,c=INK,al=PP_ALIGN.LEFT,an=MSO_ANCHOR.TOP,ls=None):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=an
    tf.margin_left=Emu(0); tf.margin_right=Emu(0); tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=al
        if ls: p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.name=FONT; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
    return tb
def hdr(s,eye,main,sub=""):
    t(s,Inches(0.6),Inches(0.4),Inches(12),Inches(0.35),eye,13,True,RED)
    bx(s,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(s,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),main,22,True,INK)
    if sub: t(s,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sub,11.5,False,GRY)
def ft(s,p):
    bx(s,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(s,Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),"テナントアシスト・ウイン株式会社 ｜ 菊池 研太",9,False,GRY)
    t(s,Inches(12.0),H-Inches(0.42),Inches(0.8),Inches(0.32),f"{p} / 5",9,False,GRY,al=PP_ALIGN.RIGHT)

s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,Inches(0.9),Inches(1.5),Inches(11),Inches(0.4),"京橋クリニック様",14,True,RED)
t(s,Inches(0.88),Inches(2.05),Inches(11.5),Inches(0.9),"業務削減シミュレーション報告書",34,True,INK)
t(s,Inches(0.9),Inches(2.95),Inches(11.5),Inches(0.5),"＋ オンライン診療（初診から）導入のご提案",20,True,RED)
t(s,Inches(0.92),Inches(3.7),Inches(11),Inches(0.6),"御院の現状を拝見し、削減できるポイントと進め方をまとめました。\n初期費用0円・成果報酬型。御院の持ち出しはゼロから始められます。",13.5,False,GRY,ls=1.3)
for i,(lab,val) in enumerate([("初期費用","0円"),("月額基本料","0円"),("成果報酬","40%")]):
    cx=Inches(0.9)+ (Inches(2.45)+Inches(0.18))*i
    bx(s,cx,Inches(4.9),Inches(2.45),Inches(1.2),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(4.9),Inches(2.45),Inches(0.06),RED)
    t(s,cx,Inches(5.05),Inches(2.45),Inches(0.3),lab,12,False,GRY,al=PP_ALIGN.CENTER)
    t(s,cx,Inches(5.32),Inches(2.45),Inches(0.7),val,32,True,RED,al=PP_ALIGN.CENTER)
t(s,Inches(0.9),Inches(6.3),Inches(11.5),Inches(0.3),"※成果報酬は最初の6ヶ月のみ。7ヶ月目以降は月5万円の保守（月1回の運用相談込み）",11,True,REDD)
t(s,Inches(0.9),Inches(6.85),Inches(11),Inches(0.3),"テナントアシスト・ウイン株式会社 ｜ 菊池 研太　2026.06",10,False,GRY)

s=sl(); ft(s,2)
hdr(s,"YOUR CLINIC","御院の現状 ── お困りごとを拝見しました","現場のお声と日々の業務から、負担が大きい5つのポイント")
items=[("電話","回線が受付人数より多く、問い合わせ対応が長くなりがち"),("予約","ご予約（アイコール）の行き違いで、診察が長引き・順番も前後"),("問診","紙の問診票を、毎日手で入力し直している"),("書類","紹介状の作成に二度手間が生じている"),("レセプト","月初の点検に、まとまった時間がかかる")]
y0=Inches(1.95)
for i,(tag,d) in enumerate(items):
    cy=y0+Inches(0.92)*i
    bx(s,Inches(0.6),cy,Inches(12.1),Inches(0.78),CARD if i%2 else WHT,line=CARDLN,lw=1.0)
    bx(s,Inches(0.6),cy,Inches(0.1),Inches(0.78),RED)
    bx(s,Inches(0.85),cy+Inches(0.19),Inches(1.5),Inches(0.42),RED)
    t(s,Inches(0.85),cy+Inches(0.19),Inches(1.5),Inches(0.42),tag,13,True,WHT,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    t(s,Inches(2.6),cy,Inches(10.0),Inches(0.78),d,13,False,INK,an=MSO_ANCHOR.MIDDLE)
t(s,Inches(0.6),Inches(6.65),Inches(12),Inches(0.3),"→ 事務スタッフの負担を軽くすることが、定着にもつながります。",11.5,True,GRN)

s=sl(); ft(s,3)
hdr(s,"SIMULATION","削減シミュレーション ── 御院に、いくら残るか","標準モデル（事務3名・月24万円削減）。当日、御院の実数で計算し直します")
s.shapes.add_picture(SAV,Inches(0.5),Inches(1.85),width=Inches(8.2))
RX,RW=Inches(9.0),Inches(3.85)
bx(s,RX,Inches(1.9),RW,Inches(2.2),CARD,line=CARDLN,lw=1.0)
t(s,RX+Inches(0.2),Inches(2.02),RW-Inches(0.4),Inches(0.3),"御院に残る効果",13,True,RED)
for i,(k,v) in enumerate([("1〜6ヶ月","月 14.4万円"),("7ヶ月目〜","月 19万円"),("12ヶ月 累計","約 200万円")]):
    yy=Inches(2.4)+Inches(0.45)*i
    t(s,RX+Inches(0.25),yy,Inches(1.6),Inches(0.4),k,12,True,INK)
    t(s,RX+Inches(1.6),yy,RW-Inches(1.85),Inches(0.4),v,15,True,REDD,al=PP_ALIGN.RIGHT)
bx(s,RX,Inches(4.35),RW,Inches(2.3),WHT,line=CARDLN,lw=1.0); bx(s,RX,Inches(4.35),RW,Inches(0.06),RED)
t(s,RX+Inches(0.2),Inches(4.48),RW-Inches(0.4),Inches(0.3),"安心の設計",13,True,RED)
for i,ln in enumerate(["効果が出なければ 費用0円","測定の手間は全部こちらが巻き取る","患者情報の扱いは最初に書面ルール化","ロックインなし・30日前通知で解約可"]):
    t(s,RX+Inches(0.22),Inches(4.85)+Inches(0.43)*i,RW-Inches(0.44),Inches(0.4),"✓ "+ln,11.5,True,INK)

s=sl(); ft(s,4)
hdr(s,"ONLINE CONSULTATION","オンライン診療（初診から）の進め方","先生が一番やりたいこと。今のカルテを変えずに、追加で始められます")
bx(s,Inches(0.6),Inches(1.95),Inches(6.0),Inches(4.6),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.6),Inches(1.95),Inches(6.0),Inches(0.06),RED)
t(s,Inches(0.8),Inches(2.08),Inches(5.6),Inches(0.3),"進め方（低リスク）",13,True,RED)
for i,(n,d) in enumerate([("1","今の電子カルテ（Medicom）はそのまま継続"),("2","オンライン診療を追加で導入（予約→問診→ビデオ→決済）"),("3","iPadで運用。現場の操作はこちらで設定・伴走"),("4","IT導入補助金（今期7/21締切）に間に合えば、コストを圧縮")]):
    yy=Inches(2.5)+Inches(0.78)*i
    bx(s,Inches(0.8),yy,Inches(0.42),Inches(0.42),RED)
    t(s,Inches(0.8),yy,Inches(0.42),Inches(0.42),n,13,True,WHT,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
    t(s,Inches(1.4),yy-Inches(0.02),Inches(5.0),Inches(0.7),d,12,False,INK,an=MSO_ANCHOR.MIDDLE,ls=1.1)
bx(s,Inches(6.85),Inches(1.95),Inches(5.85),Inches(4.6),WHT,line=CARDLN,lw=1.0); bx(s,Inches(6.85),Inches(1.95),Inches(5.85),Inches(0.06),RED)
t(s,Inches(7.05),Inches(2.08),Inches(5.5),Inches(0.3),"「初診から」やるための前提",13,True,RED)
t(s,Inches(7.07),Inches(2.5),Inches(5.5),Inches(3.9),"・2026年4月施行の改訂指針に沿って運用します。\n\n・初診のオンライン診療も、既往歴・服薬・アレルギー等を把握でき、医師が可能と判断した場合に実施できます。\n\n・初診は原則、顔写真付きの本人確認（マイナンバーカード・運転免許等）を行います。\n\n・対面とオンラインを適切に組み合わせる前提で、御院の方針に合わせて設計します。\n\n※最新の施行版指針に沿って、運用ルールを書面で取り決めます。",12,False,INK,ls=1.25)

s=sl(); ft(s,5)
hdr(s,"NEXT STEP","進め方 ── 効果を確認してから、導入へ","費用が発生するのは、効果が数字で確認できた後だけです")
steps=[("STEP 1","無料の業務診断","本報告書＋現場の業務を拝見し、削減ポイントを特定"),("STEP 2","段階導入","オンライン診療・入力時短など、効果の大きい所から少しずつ"),("STEP 3","効果測定・精算","測定は当方が実施。確認できた削減分の40%のみご精算")]
cw,gx,x0,y0=Inches(3.95),Inches(0.24),Inches(0.55),Inches(2.0)
for i,(st,ti,bd) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,y0,cw,Inches(2.9),CARD,line=CARDLN,lw=1.0); bx(s,cx,y0,cw,Inches(0.7),RED)
    t(s,cx,y0+Inches(0.12),cw,Inches(0.45),st,17,True,WHT,al=PP_ALIGN.CENTER)
    t(s,cx+Inches(0.2),y0+Inches(0.9),cw-Inches(0.4),Inches(0.6),ti,15,True,INK,al=PP_ALIGN.CENTER)
    t(s,cx+Inches(0.25),y0+Inches(1.6),cw-Inches(0.5),Inches(1.1),bd,11.5,False,GRY,al=PP_ALIGN.CENTER,ls=1.2)
by=y0+Inches(3.15)
bx(s,Inches(0.55),by,Inches(12.23),Inches(0.95),REDBG); bx(s,Inches(0.55),by,Inches(0.1),Inches(0.95),RED)
t(s,Inches(0.85),by+Inches(0.1),Inches(11.6),Inches(0.4),"御院の持ち出しはゼロから。効果が出た分だけ、私たちはいただきます。",15,True,REDD)
t(s,Inches(0.85),by+Inches(0.54),Inches(11.6),Inches(0.35),"✓ 成果報酬は最初の6ヶ月だけ（以降は月5万の保守のみ）　✓ IT導入補助金サポート　✓ 3ヶ月伴走",11.5,True,INK)

import shutil
prs.save("kyobashi_shindan_report.pptx")
d="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/京橋_業務削減シミュレーション報告書.pptx"
shutil.copy("kyobashi_shindan_report.pptx",d)
print("SAVED:",d,"/ slides:",len(prs.slides._sldIdLst))

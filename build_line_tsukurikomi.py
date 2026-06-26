"""
京橋 公式LINE 作り込み案（次回YESを取る用）。まず問診票から。クリーム白×レンガ赤。
出力: line_tsukurikomi.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=22,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.46),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック 公式LINE 作り込み案",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)

# ════ 1 表紙 ════
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"京橋クリニック 公式LINE",Inches(0.9),Inches(1.15),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"作り込み案（次回ご確認用）\nまずは「問診票」から。",Inches(0.88),Inches(1.7),Inches(11.7),Inches(1.6),sz=29,bold=True,col=INK,line_sp=1.12)
t(s,"前回（6/25）いただいた皆さまの声を、具体的な中身に落とし込みました。",Inches(0.92),Inches(3.75),Inches(11.4),Inches(0.5),sz=14,col=GRY)
bx(s,Inches(0.9),Inches(4.6),Inches(11.5),Inches(1.15),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(4.6),Inches(0.1),Inches(1.15),RED)
t(s,"全部を一度にやりません。まず「問診票」から、無料で小さく。効いた所から広げます。",Inches(1.2),Inches(4.6),Inches(11),Inches(1.15),sz=15,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 2 前回決まったこと ════
s=sl(); ft(s)
hdr(s,"前回（6/25）の確認","ここまで決まりました","この前提のうえで、中身を詰めます")
dec=[("① LINE導入＝方向OK","「やる方向で」とご判断いただきました"),
     ("② スタートは来週〜","具体的な日程は、福井さんと調整します"),
     ("③ まずは「問診票」から","入口は問診。小さく始めます"),
     ("④ リマインドはデフォルトON","看護師長のご意見を反映し、まずON設定")]
cw,ch,gx,gy=Inches(6.0),Inches(1.7),Inches(0.23),Inches(0.3); x0,y0=Inches(0.55),Inches(2.05)
for i,(ti,ds) in enumerate(dec):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.25),cy+Inches(0.2),cw-Inches(0.45),Inches(0.6),sz=15,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.27),cy+Inches(0.82),cw-Inches(0.5),Inches(0.7),sz=12,col=INK,line_sp=1.2)

# ════ 3 全体像（フェーズ） ════
s=sl(); ft(s)
hdr(s,"全体像","問診を入口に、段階で広げる","全部一度にやらない。問診から小さく、効いた所だけ次へ。")
bx(s,Inches(0.55),Inches(2.0),Inches(12.23),Inches(1.4),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(2.0),Inches(0.12),Inches(1.4),TEALD)
t(s,"まず（P1）",Inches(0.8),Inches(2.12),Inches(2.2),Inches(0.4),sz=14,bold=True,col=TEALD)
t(s,"① Web問診（一般／健診／化学物質過敏症）　＋　② 再診リマインド（デフォルトON）",Inches(0.85),Inches(2.55),Inches(11.6),Inches(0.7),sz=14,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
bx(s,Inches(0.55),Inches(3.55),Inches(12.23),Inches(1.4),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(3.55),Inches(0.12),Inches(1.4),RED)
t(s,"次に（P2）",Inches(0.8),Inches(3.67),Inches(2.2),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"③ CPAPの次回予約 自動送付　／　④ 再診促し（ザオラル→LINE）　／　⑤ 化学物質過敏症・健診の問診",Inches(0.85),Inches(4.1),Inches(11.6),Inches(0.7),sz=13.5,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
bx(s,Inches(0.55),Inches(5.1),Inches(12.23),Inches(1.0),REDBG); bx(s,Inches(0.55),Inches(5.1),Inches(0.1),Inches(1.0),RED)
t(s,"横串の効能：次回予約をきちんと取る＝取りこぼし防止＋「必ず取ってますよ」が国へのアピールにも。",Inches(0.82),Inches(5.1),Inches(11.8),Inches(1.0),sz=13.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 4 ①Web問診 ════
s=sl(); ft(s)
hdr(s,"① Web問診（ここからスタート）","来院前にスマホで記入 → 受付の聞き取り・転記が減る","先生からいただいたテンプレを反映。斉藤内科さんの導線を参考に。")
forms=[("一般の問診","ふだんの受診向け。来院前にスマホで記入"),
       ("健診の問診","健診の問診はLINEでOK（ご意見どおり）"),
       ("化学物質過敏症","LINEが特に響く層。問診後に先生がお電話フォロー")]
cw,gx,x0=Inches(3.96),Inches(0.18),Inches(0.55)
for i,(ti,ds) in enumerate(forms):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.1),cw,Inches(2.0),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.1),cw,Inches(0.6),RED)
    t(s,ti,cx+Inches(0.15),Inches(2.18),cw-Inches(0.3),Inches(0.45),sz=14,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ds,cx+Inches(0.25),Inches(2.85),cw-Inches(0.5),Inches(1.1),sz=12,col=INK,align=PP_ALIGN.CENTER,line_sp=1.3)
bx(s,Inches(0.55),Inches(4.35),Inches(12.23),Inches(1.55),GRYBG); bx(s,Inches(0.55),Inches(4.35),Inches(0.1),Inches(1.55),RED)
t(s,"設計のしかた",Inches(0.82),Inches(4.45),Inches(11),Inches(0.35),sz=13,bold=True,col=REDD)
t(s,"・ベンチマーク：斉藤内科さん（デジスマ診療）の問診導線を参考に、迷わない流れに。\n・先生からいただいたテンプレをベースに、京橋仕様へ。\n・送信されたら自動でお礼＋タグ付け。受付は一覧で確認できます。",
  Inches(0.85),Inches(4.82),Inches(11.6),Inches(1.0),sz=12,col=INK,line_sp=1.3)

# ════ 5 ②リマインド＋CPAP ════
s=sl(); ft(s)
hdr(s,"② 再診リマインド（デフォルトON）＋ CPAP","取りこぼしを、仕組みで防ぐ","看護師長のご意見を反映")
block=[("再診リマインド（デフォルトON）","看護師長のご意見どおり、まずはデフォルトでON。再診・定期の方に、来院の目安をLINEでそっとお知らせします。"),
       ("CPAP定期 → 次回予約を自動送付","これまで毎回は聞けていなかった次回予約を、事前設定で自動的にご案内。通院の途切れを防ぎます。")]
y=Inches(2.1)
for ti,ds in block:
    bx(s,Inches(0.55),y,Inches(12.23),Inches(1.75),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),y,Inches(0.1),Inches(1.75),RED)
    t(s,ti,Inches(0.85),y+Inches(0.2),Inches(11.6),Inches(0.45),sz=15,bold=True,col=REDD)
    t(s,ds,Inches(0.87),y+Inches(0.78),Inches(11.5),Inches(0.8),sz=12.5,col=INK,line_sp=1.3)
    y=y+Inches(1.95)
t(s,"※ どちらも「お知らせ」止まり。患者さんの判断・医療の判断は、これまでどおり人が行います。",Inches(0.55),Inches(6.05),Inches(12),Inches(0.4),sz=11.5,bold=True,col=REDD)

# ════ 6 ③再診促し ════
s=sl(); ft(s)
hdr(s,"③ 再診促し ── 事務の架電を、LINEで軽く","いまの大変さを、仕組みで肩代わり","先生のお電話は「本当に必要な人」だけに")
bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(2.9),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(0.55),GRY)
t(s,"いま",Inches(0.8),Inches(2.06),Inches(5.5),Inches(0.4),sz=14,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"・最近来ていない方へ、事務が一人ずつ架電（大変）。\n・その後、先生もお電話。\n・不在・ご出張で、結局つながらないことも多い。",Inches(0.82),Inches(2.7),Inches(5.5),Inches(2.0),sz=12.5,col=INK,line_sp=1.4)
bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(2.9),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(0.55),TEALD)
t(s,"これから",Inches(7.03),Inches(2.06),Inches(5.5),Inches(0.4),sz=14,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"・離れかけの方へ、LINEで次回予約の導線を。\n・事務の架電を減らし、不在でも予約が取れる。\n・先生のお電話は「本当に必要な人」だけに集中。",Inches(7.0),Inches(2.7),Inches(5.5),Inches(2.0),sz=12.5,col=INK,line_sp=1.4)
bx(s,Inches(0.55),Inches(5.1),Inches(12.23),Inches(0.95),REDBG); bx(s,Inches(0.55),Inches(5.1),Inches(0.1),Inches(0.95),RED)
t(s,"おまけの効能：次回予約をきちんと取る＝「必ず取ってますよ」が、国へのアピールにもなります。",Inches(0.82),Inches(5.1),Inches(11.8),Inches(0.95),sz=13.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 7 ご確認＆次の一手 ════
s=sl(); ft(s)
hdr(s,"ご確認 ＆ 次の一手","ここだけ、教えてください","あとは、こちらで作り込みます")
bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(2.6),REDBG,line=RED,lw=1.0); bx(s,Inches(0.55),Inches(2.0),Inches(0.1),Inches(2.6),RED)
t(s,"確認したい2点",Inches(0.82),Inches(2.12),Inches(5.5),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"①「最後の1週間」に架電、の運用の意味\n（再診促しの設計に効きます）\n\n② サインエコーのLINE導入、とは？\n（どんなご要望か、もう少し）",Inches(0.82),Inches(2.65),Inches(5.5),Inches(1.85),sz=12.5,col=INK,line_sp=1.35)
bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(2.6),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(2.0),Inches(0.1),Inches(2.6),RED)
t(s,"いただきたいもの・次の流れ",Inches(7.03),Inches(2.12),Inches(5.5),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"・問診テンプレの現物（反映します）\n\n・次回（納品アポ）で、この作り込みでGO\n　→ 設定 → 運用開始\n\n・問診の「実物」もお見せします",Inches(7.03),Inches(2.65),Inches(5.5),Inches(1.85),sz=12.5,col=INK,line_sp=1.35)
bx(s,Inches(0.55),Inches(4.95),Inches(12.23),Inches(1.0),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(4.95),Inches(0.1),Inches(1.0),TEALD)
t(s,"まず「問診」から、無料で小さく。合わなければ、いつでも見直せます。",Inches(0.82),Inches(4.95),Inches(11.8),Inches(1.0),sz=15,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("line_tsukurikomi.pptx")
print("saved line_tsukurikomi.pptx slides:",len(prs.slides._sldIdLst))

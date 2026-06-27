"""新Mac3/Mac2 セットアップ手順書 ── KHD標準デザイン(クリーム白×レンガ赤)。
ゆーし=PC初心者向けに具体化(初回は菊池代行/毎日はクリックだけ・リモートデスクトップ推奨)。
出力: KHD_新Mac3Mac2_セットアップ手順書.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=col; r.font.name=FONT
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)

def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"KHD 3台セットアップ手順書  ｜  Mac1 / Mac2 / Mac3  ｜  2026-06-26",Inches(0.5),H-Inches(0.42),Inches(11),Inches(0.32),sz=9,col=GRY)

def panel(s,x,y,w,h,title,lines,tcol=RED,lh=0.6,tsz=15):
    bx(s,x,y,w,h,CARD,line=CARDLN,lw=1.0); bx(s,x,y,w,Inches(0.06),RED)
    t(s,title,x+Inches(0.28),y+Inches(0.18),w-Inches(0.5),Inches(0.4),sz=tsz,bold=True,col=tcol)
    yy=y+Inches(0.8)
    for ln,col in lines:
        t(s,ln,x+Inches(0.3),yy,w-Inches(0.6),Inches(lh),sz=12.5,col=col,line_sp=1.12); yy=yy+Inches(lh)

def band(s,y,text,sub=""):
    hh=Inches(0.95) if sub else Inches(0.62)
    bx(s,Inches(0.55),y,Inches(12.23),hh,REDBG); bx(s,Inches(0.55),y,Inches(0.1),hh,RED)
    t(s,text,Inches(0.85),y+Inches(0.1),Inches(11.6),Inches(0.4),sz=14,bold=True,col=REDD)
    if sub: t(s,sub,Inches(0.85),y+Inches(0.54),Inches(11.6),Inches(0.35),sz=11.5,col=INK)

def bigsteps(s,x,y,w,steps,gap=1.02,nsz=20,tsz=15.5):
    for i,txt in enumerate(steps):
        yy=y+Inches(gap)*i
        c=s.shapes.add_shape(MSO_SHAPE.OVAL,x,yy,Inches(0.62),Inches(0.62))
        c.fill.solid(); c.fill.fore_color.rgb=RED; c.line.fill.background(); c.shadow.inherit=False
        t(s,str(i+1),x,yy,Inches(0.62),Inches(0.62),sz=nsz,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        t(s,txt,x+Inches(0.85),yy,w-Inches(0.95),Inches(0.62),sz=tsz,col=INK,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.05)

# ── S1 表紙 ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"MAC SETUP ｜ 3 DEVICES · 1 ACCOUNT",Inches(0.9),Inches(1.55),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"新Mac(Mac3) ＆ Mac2",Inches(0.88),Inches(2.15),Inches(11.7),Inches(0.9),sz=40,bold=True,col=INK)
t(s,"セットアップ手順書",Inches(0.88),Inches(2.95),Inches(11.7),Inches(0.9),sz=40,bold=True,col=RED)
t(s,"3台を1アカウントで動かし、データもClaudeも“Mac1と同じ状態”で起動するまで。",Inches(0.9),Inches(3.95),Inches(11.6),Inches(0.6),sz=14,col=GRY,line_sp=1.25)
roles=[("Mac1","現メイン（今のMac）"),("Mac2","ゆーしMac（PC初心者OK）"),("Mac3","新ハイスペック（母艦）")]
ox,ow,gap=Inches(0.9),Inches(3.83),Inches(0.2)
for i,(lab,role) in enumerate(roles):
    cx=ox+(ow+gap)*i
    bx(s,cx,Inches(4.95),ow,Inches(1.25),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(4.95),ow,Inches(0.06),RED)
    t(s,lab,cx,Inches(5.16),ow,Inches(0.45),sz=22,bold=True,col=RED,align=PP_ALIGN.CENTER)
    t(s,role,cx,Inches(5.66),ow,Inches(0.45),sz=12.5,col=INK,align=PP_ALIGN.CENTER)
bx(s,Inches(0.9),Inches(6.55),Inches(11.5),Pt(1.2),LINE)
t(s,"KHD  ｜  菊池 研太  ｜  2026-06-26",Inches(0.9),Inches(6.66),Inches(11),Inches(0.4),sz=13,bold=True,col=INK)

# ── S2 Apple連携の答え ──
s=sl(); ft(s)
hdr(s,"THE ANSWER","「ボタンひとつ同期」できてる？","結論：初期コピー＝Apple／継続同期＝git（両方準備済）")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(4.3),"✅ Apple でできる",[
 ("iCloud・同一Apple ID（菊池研太）設定済",INK),
 ("→ 移行アシスタントでMac1→Mac3を",GRY),
 ("　丸ごとコピー（ほぼボタン一発）",GRY),
 ("Handoff／コピペ共有／AirDrop も",INK),
 ("3台で使える",INK),
],lh=0.58)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(4.3),"⚠️ Apple では無理 → git で解決",[
 ("“ずっとボタン一つで同期”はApple非対応",INK),
 ("（iCloudはApple純正データだけ）",GRY),
 ("作業フォルダ＆Claudeの記憶は対象外",INK),
 ("→ git（Drive内・設定済）で同期",REDD),
 ("　コマンドはClaude代行＝実質ひと言",GRY),
],tcol=REDD,lh=0.58)

# ── S3 全体像 ──
s=sl(); ft(s)
hdr(s,"ARCHITECTURE","全体像 ── 1アカウント・3台","正本はGoogle Drive内の2リポジトリ（両方push済）")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(2.5),"khd_workspace.git",[
 ("作業フォルダ全体（資料・スプシ・秘書室）",INK),
 ("~/01_honbu_docs_automation",GRY),
])
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(2.5),"khd_memory.git",[
 ("Claudeの記憶（あなたの方針・過去の判断）",INK),
 ("~/.claude/.../memory",GRY),
])
band(s,Inches(4.75),"置き場＝マイドライブ/KHD_git_remote/（ec-seller アカウント内）",
 "iCloud＝Apple純正データ　｜　各Macで同じClaude(Anthropic)アカウントにログイン")

# ── S4 Mac3 Phase1 ──
s=sl(); ft(s)
hdr(s,"MAC3 ｜ PHASE 1","移行アシスタントで丸ごとコピー","「Mac1と同じ」を最速で作る（菊池さんが実施）")
bigsteps(s,Inches(0.7),Inches(1.95),Inches(12.0),[
 "電源を入れる → 言語・国・Wi-Fi を選ぶ",
 "「情報を転送」画面で “Mac／Time Machine／起動ディスクから” を選ぶ",
 "Mac1 を選び、転送する項目は全部チェック →「続ける」",
 "終わるまで待つ（数十分）。アプリ・データ・設定がそのまま入る",
],gap=0.92)
band(s,Inches(5.95),"Thunderboltケーブル直結が最速。Homebrew未導入でも移行で必要物（Node・claude等）は乗る")

# ── S5 Mac3 Phase2 ──
s=sl(); ft(s)
hdr(s,"MAC3 ｜ PHASE 2","サインイン ＆ Google Drive","同期の土台を通電する（菊池さんが実施）")
bigsteps(s,Inches(0.7),Inches(1.9),Inches(12.0),[
 "Apple ID（kemkemsp@yahoo.co.jp）でサインイン・iCloud を ON",
 "画面上のGoogle Drive（雲アイコン）→ ec-seller@kikuchi-hd.net でログイン",
 "Finder → サイドバー「Google Drive」→「マイドライブ」→「KHD_git_remote」",
 "そのフォルダを右クリック →「オフラインで使用可能にする」にチェック",
],gap=0.9)
band(s,Inches(5.85),"⚠️ ④の“オフラインで使用可能”を忘れると同期が失敗。ここだけ要注意")

# ── S6 Mac3 Phase3-4 ──
s=sl(); ft(s)
hdr(s,"MAC3 ｜ PHASE 3-4","同期 → Claude起動（ゴール）","ここまで来たらMac1と同じ状態")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.5),"③ 同期（かんたん）",[
 ("Claudeに「同期して」と言うだけ",INK),
 ("（git pull は Claude が実行）",GRY),
 ("→ 作業フォルダと記憶が最新になる",INK),
],lh=0.58)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.5),"④ Claude起動",[
 ("1. Launchpad →「その他」→「ターミナル」",INK),
 ("2. claude と打って Enter → ログイン",INK),
 ("3. 出るリンクで各連携(MCP)を許可",INK),
 ("4. /company と打つ → 秘書が起動",INK),
],lh=0.58)
band(s,Inches(5.75),"✓ 秘書が前回の続きを覚えていれば成功 ＝ Mac1と同じ状態")

# ── S7 Mac2 ゆーし：考え方 ──
s=sl(); ft(s)
hdr(s,"MAC2 ｜ ゆーしMac（PC初心者OK）","むずかしい作業は本人にさせない","初回は菊池が代行／毎日はクリックだけ。推奨＝母艦に“画面で入る”")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.5),"初回だけ（菊池が代行）",[
 ("母艦Mac3に「Chromeリモート",INK),
 ("デスクトップ」を設定（無料）",INK),
 ("ゆーしのMacには何も入れない",REDD),
 ("→ 難しい設定はゆーし不要",GRY),
],lh=0.6)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.5),"毎日（ゆーしがやる事）",[
 ("Chromeを開いてアイコンを押すだけ",INK),
 ("母艦Mac3の画面がそのまま出る",INK),
 ("→ その中でClaudeを使う",INK),
 ("（実際の処理は母艦が動く）",GRY),
],lh=0.6)
band(s,Inches(5.75),"なぜ？ ゆーしのMacにgitもClaudeも入れない＝壊れない・競合しない・実体は母艦1つ")

# ── S8 Mac2 ゆーし：毎日の使い方 ──
s=sl(); ft(s)
hdr(s,"MAC2 ｜ 毎日の使い方","この4ステップだけ覚えればOK","むずかしい言葉ゼロ。最初の数回だけ一緒にやれば慣れます")
bigsteps(s,Inches(0.8),Inches(1.95),Inches(11.8),[
 "Chrome（カラフルな丸いアイコン）を開く",
 "画面の「母艦Mac3」を押す →（最初だけ数字パスワード）",
 "母艦の画面が出る → いつものClaudeに話しかける",
 "終わったら、その窓を閉じるだけ（保存はClaudeが自動）",
],gap=0.96,tsz=16)
band(s,Inches(6.0),"困ったら菊池さんにLINE。※本人PCで直接動かしたい時は、初回導入を菊池が全部代行→以後はアイコン1つ")

# ── S9 日々の運用（3台）──
s=sl(); ft(s)
hdr(s,"DAILY OPS","日々の運用 ── 鉄則は2つだけ","菊池さんもゆーしさんもgitは覚えなくていい")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(2.3),"① 開始pull / 終了push",[
 ("作業開始＝git pull（最新を取込）",INK),
 ("作業終了＝git push（正本を更新）",INK),
 ("どちらもClaudeが代行",GRY),
])
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(2.3),"② 書き手は1台",[
 ("母艦Mac3をマスターに固定",INK),
 ("同時に2人で同じ物を書き換えない",GRY),
])
band(s,Inches(4.55),"📱 iPhone・Mac2は「母艦にリモートで入る」が安全",
 "母艦を遠隔操作するだけ＝書き手が増えず競合ゼロ。複数台・同時ログインもOK")

# ── S10 チェックリスト / 締め ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"CHECKLIST",Inches(0.9),Inches(0.65),Inches(11),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"この順でやれば、3台が立ち上がる",Inches(0.9),Inches(1.1),Inches(11.7),Inches(0.7),sz=26,bold=True,col=INK)
panel(s,Inches(0.55),Inches(2.1),Inches(6.0),Inches(3.6),"Mac3（母艦・菊池）",[
 ("□ 移行アシスタントで丸ごとコピー",INK),
 ("□ Apple ID／Drive(ec-seller)サインイン",INK),
 ("□ KHD_git_remote をオフライン化",INK),
 ("□ 「同期して」でgit pull（Claude）",INK),
 ("□ ターミナルで claude → /company 確認",INK),
],lh=0.58)
panel(s,Inches(6.78),Inches(2.1),Inches(6.0),Inches(3.6),"Mac2（ゆーし・初心者）",[
 ("□【初回・菊池】母艦にChromeリモート",INK),
 ("　 デスクトップを設定",GRY),
 ("□【毎日・ゆーし】Chrome→母艦を押す",INK),
 ("□ 母艦の画面でClaudeに話しかける",INK),
 ("□ 終わったら窓を閉じる",INK),
],lh=0.58)
bx(s,Inches(0.55),Inches(5.95),Inches(12.23),Pt(1.2),LINE)
t(s,"KHD  ｜  3台・1アカウント・記憶まで同期  ｜  菊池 研太",Inches(0.55),Inches(6.08),Inches(12),Inches(0.4),sz=12,bold=True,col=INK)

out="/Users/kikuchikenta/01_honbu_docs_automation/KHD_新Mac3Mac2_セットアップ手順書.pptx"
prs.save(out); print("saved:",out,"/ slides:",len(prs.slides._sldIdLst))

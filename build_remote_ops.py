# -*- coding: utf-8 -*-
"""01_オペレーティングモデル × Mac3台・リモート活用 デッキ（KHD標準クリーム白×レンガ赤）
既存「人は打席だけ」モデルに、旧Mac自宅常設×リモートでの3台使い分け＋スマホ→自宅Macリモート起動手順を統合。
出力: 01_オペモデル_Mac3台リモート活用_260628.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=col; r.font.name=FONT
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)

def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"KHD オペレーティングモデル × Mac3台・リモート活用  ｜  菊池 研太  ｜  2026-06-28",Inches(0.5),H-Inches(0.42),Inches(11),Inches(0.32),sz=9,col=GRY)

def panel(s,x,y,w,h,title,lines,tcol=RED,lh=0.6,tsz=15,bsz=12.5):
    bx(s,x,y,w,h,CARD,line=CARDLN,lw=1.0); bx(s,x,y,w,Inches(0.06),RED)
    t(s,title,x+Inches(0.28),y+Inches(0.16),w-Inches(0.5),Inches(0.45),sz=tsz,bold=True,col=tcol)
    yy=y+Inches(0.78)
    for ln,col in lines:
        t(s,ln,x+Inches(0.3),yy,w-Inches(0.6),Inches(lh),sz=bsz,col=col,line_sp=1.12); yy=yy+Inches(lh)

def band(s,y,text,sub=""):
    hh=Inches(0.95) if sub else Inches(0.62)
    bx(s,Inches(0.55),y,Inches(12.23),hh,REDBG); bx(s,Inches(0.55),y,Inches(0.1),hh,RED)
    t(s,text,Inches(0.85),y+Inches(0.1),Inches(11.6),Inches(0.4),sz=14,bold=True,col=REDD)
    if sub: t(s,sub,Inches(0.85),y+Inches(0.54),Inches(11.6),Inches(0.35),sz=11.5,col=INK)

def bigsteps(s,x,y,w,steps,gap=0.92,nsz=19,tsz=14):
    for i,(head,desc) in enumerate(steps):
        yy=y+Inches(gap)*i
        c=s.shapes.add_shape(MSO_SHAPE.OVAL,x,yy,Inches(0.56),Inches(0.56))
        c.fill.solid(); c.fill.fore_color.rgb=RED; c.line.fill.background(); c.shadow.inherit=False
        t(s,str(i+1),x,yy,Inches(0.56),Inches(0.56),sz=nsz,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        t(s,head,x+Inches(0.78),yy-Inches(0.02),w-Inches(0.85),Inches(0.34),sz=tsz,bold=True,col=INK,line_sp=1.0)
        t(s,desc,x+Inches(0.78),yy+Inches(0.3),w-Inches(0.85),Inches(0.5),sz=11,col=GRY,line_sp=1.05)

# ── S1 表紙 ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"OPERATING MODEL ｜ 1 HOME MAC RUNS THE REST",Inches(0.9),Inches(1.5),Inches(11.5),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"人は「打席」だけ。",Inches(0.88),Inches(2.1),Inches(11.7),Inches(0.9),sz=40,bold=True,col=INK)
t(s,"残りは、自宅Macが24時間回す。",Inches(0.88),Inches(2.95),Inches(11.7),Inches(0.9),sz=40,bold=True,col=RED)
t(s,"会って・聞いて・詰めて・信頼を作る＝人。記録/集計/資料/分析/可視化＝Claude。\nその“Claude側”を、自宅に常設した旧Macが、スマホ1本でいつでも回す。",
  Inches(0.9),Inches(3.95),Inches(11.6),Inches(0.9),sz=13.5,col=GRY,line_sp=1.25)
roles=[("旧Mac","自宅常設・リモート箱","Claude代行を24h・git母艦"),
       ("新Mac","持ち出し・打席機","商談/内見の相棒"),
       ("ゆーしMac","ゆーし専用","EC折半・本体DBは触らない")]
ox,ow,gap=Inches(0.9),Inches(3.83),Inches(0.2)
for i,(lab,role,sub) in enumerate(roles):
    cx=ox+(ow+gap)*i
    bx(s,cx,Inches(5.15),ow,Inches(1.32),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(5.15),ow,Inches(0.06),RED)
    t(s,lab,cx,Inches(5.32),ow,Inches(0.45),sz=21,bold=True,col=RED,align=PP_ALIGN.CENTER)
    t(s,role,cx,Inches(5.8),ow,Inches(0.34),sz=12,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,sub,cx,Inches(6.12),ow,Inches(0.34),sz=10.5,col=GRY,align=PP_ALIGN.CENTER)
t(s,"KHD  ｜  菊池 研太  ｜  2026-06-28",Inches(0.9),Inches(6.66),Inches(11),Inches(0.4),sz=12,bold=True,col=INK)

# ── S2 大原則（再掲＋今回の進化）──
s=sl(); ft(s)
hdr(s,"THE PRINCIPLE","経営の最小単位は1つ ── 人は打席に全集中","空いた時間は調査士と家族へ。“それ以外全部”を旧Macが肩代わりする")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(2.15),"◯ 人＝打席・判断・信頼",[
 ("会う／電話／ヒアリング／詰める",INK),
 ("クロージング・GO／見送り・指値・採用・GIVE",INK),
 ("→ 関係・判断・存在は自動化できない",GRY),
],lh=0.5)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(2.15),"AI Claude＝それ以外 全部",[
 ("記録(02転記)／集計(週次KPI)／資料作成",INK),
 ("分析(決め手)／リマインド／可視化",INK),
 ("→ 旧Macに常駐させ、24時間まわす",REDD),
],tcol=REDD,lh=0.5)
band(s,Inches(4.35),"今回の進化 ── 旧Macを「自宅常設のClaude代行サーバー」に格上げ",
     "持ち歩かない1台を母艦にする＝書き手1台ルールに最適。打席の合間にスマホから投げれば、帰宅前に“残り全部”が片付いている。")
panel(s,Inches(0.55),Inches(5.5),Inches(12.23),Inches(1.3),"⛔ 逃げの罠（ここだけ注意）",[
 ("「自動化の作り込み」自体が新しい謎の業務になりがち。物件の自動提案を作り込むより、レインズで自分で精査→ヒアリング→即提案が速い（6/27 羽鳥MEETの学び）。空いた時間は打席・調査士・家族へ。",INK),
],tcol=REDD,lh=0.5,bsz=12)

# ── S3 Mac3台の使い分け（核）──
s=sl(); ft(s)
hdr(s,"3 MACHINES","Mac3台の使い分け ── どれで何をやるか","“1案件1台”が原則。同じ顧客を2台で触らない（分岐衝突の実害あり）")
cols=[
 ("旧Mac｜自宅常設・リモート箱",RED,[
   ("● Claude代行サーバー（24時間）",INK),
   ("● git母艦＝“書き手1台”はここ",INK),
   ("● 常時2セッション：",REDD),
   ("　①03賃貸追客 ②01財務/KPI",GRY),
   ("● 重作業：査定・収支・謄本",INK),
   ("　報告書・診療圏・ダッシュボード",GRY),
   ("● スマホから叩く（次頁の手順）",INK),
 ]),
 ("新Mac｜持ち出し・打席機",REDD,[
   ("● 商談・内見・現地の相棒",INK),
   ("● 打席直後に“その場で1行”入力",INK),
   ("● 外出時はスマホ→旧Macに",GRY),
   ("　投げてもOK（PC開かなくていい）",GRY),
   ("● ハイスペック＝重い処理も可",INK),
   ("● ※書き手は原則旧Mac。新Macは",GRY),
   ("　“読む・投げる・その場メモ”中心",GRY),
 ]),
 ("ゆーしMac｜ゆーし専用",GRY,[
   ("● EC折半パートナー ゆーし用",INK),
   ("● ECや限定作業のみ",INK),
   ("● KHD本体DB・顧客マスターは",REDD),
   ("　触らない（1案件1台を厳守）",GRY),
   ("● clone前にKHD_git_remoteを",GRY),
   ("　“オフラインで使用可能”に",GRY),
   ("● 本体と同時pushしない",INK),
 ]),
]
cx0,cw,cg=Inches(0.55),Inches(3.94),Inches(0.21)
for i,(title,tc,lines) in enumerate(cols):
    x=cx0+(cw+cg)*i
    bx(s,x,Inches(1.95),cw,Inches(4.7),CARD,line=CARDLN,lw=1.0); bx(s,x,Inches(1.95),cw,Inches(0.06),tc)
    t(s,title,x+Inches(0.24),Inches(2.12),cw-Inches(0.4),Inches(0.7),sz=13.5,bold=True,col=tc,line_sp=1.0)
    yy=Inches(2.95)
    for ln,col in lines:
        t(s,ln,x+Inches(0.26),yy,cw-Inches(0.5),Inches(0.5),sz=11.5,col=col,line_sp=1.08); yy=yy+Inches(0.48)

# ── S4 リモートで投げる⬜代行カタログ ──
s=sl(); ft(s)
hdr(s,"WHAT TO THROW","スマホから投げる ── 打席直後に1行、裏で完了","「報告回数に応じて積み上がる」。朝晩に限らず“別れた直後”が理想")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(4.1),"✅ リモートで回す（⬜Claude代行）",[
 ("「○○さんと予算13万/なぜ＝駅近で握った」",INK),
 ("　→ 作業DB入力・KPI/日報・顧客マスター更新",GRY),
 ("追客ドラフト作成（送信は自分の最終チェック）",INK),
 ("査定・収支試算・謄本/契約精読",INK),
 ("診療圏調査・報告書・議事録→タスク化",INK),
 ("朝ブリーフ／週次KPI／ダッシュボード",INK),
 ("→ 脳に溜めず即ハンドオフ＝打席に集中",REDD),
],lh=0.5)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(4.1),"⛔ リモートでやらせない",[
 ("物件の自動ピック・自動提案",REDD),
 ("　→ 6/27羽鳥MEETの結論：ヒアリング→提案は",GRY),
 ("　　菊池が打席でやる方が速くて成約に近い",GRY),
 ("クロージング・YES取り（人の核）",INK),
 ("GO/見送り・指値・採用・撤退の最終判断",INK),
 ("",INK),
 ("→ “誰のYesに繋がるか”言えない作業はやめる",REDD),
],tcol=REDD,lh=0.5)
band(s,Inches(6.2),"同時セッションは“箱”でなく“投げ込み頻度”を増やす",
     "打席時間は1人分。常時2セッション（03追客／01財務）で十分。04医療・05物件は案件が立った時だけ起こす。")

# ── S5 リモート起動 手順書（#1）──
s=sl(); ft(s)
hdr(s,"HOW TO CONNECT","スマホ→自宅Mac リモート起動 手順","本命＝Claude Code 公式 Remote Control（SSH不要・コードはMac上のまま・無料）")
bigsteps(s,Inches(0.7),Inches(2.0),Inches(7.4),[
 ("旧Macを更新","ターミナルで `claude --version`（v2.1.51以上）。古ければ `claude upgrade`"),
 ("Remote Controlを開始","対象フォルダで `claude` 起動 → Remote Control を有効化（画面にQR/URLが出る）"),
 ("iPhoneに Claude アプリ","App Storeから入れ、同じアカウントでログイン"),
 ("QRを読んで接続","アプリのカメラ/「Code」タブで MacのQRを読む → セッションに接続"),
 ("Macを寝かせない","別タブで `caffeinate -di &`（終了は `killall caffeinate`）＋設定でスリープOFF"),
])
panel(s,Inches(8.45),Inches(2.0),Inches(4.35),Inches(2.55),"常時2セッション運用",[
 ("①「03追客」②「01財務」の2本を",INK),
 ("　名前付きで起動",GRY),
 ("iPhoneの「Code」タブで",INK),
 ("　タップ切替",GRY),
 ("緑点＝オンライン中の目印",GRY),
],lh=0.46,tsz=13.5)
panel(s,Inches(8.45),Inches(4.72),Inches(4.35),Inches(1.95),"⚠️ 電源OFFからは無理",[
 ("Wake-on-LANは不安定。",INK),
 ("旧Mac＝常時起動＋スリープなしが",REDD),
 ("一番手堅い（電気代だけ）",GRY),
],tcol=REDD,lh=0.46,tsz=13.5)
t(s,"※コマンドの正確な表記はバージョンで変わることがあります。`claude --help` と公式ドキュメント（code.claude.com/docs/remote-control）で最新を確認。",
  Inches(0.7),Inches(6.72),Inches(12),Inches(0.3),sz=9.5,col=GRY)

# ── S6 方式比較＆セキュリティ ──
s=sl(); ft(s)
hdr(s,"COMPARISON","3つの方式 ── 迷ったら本命1本でいい","結論：菊池さんは Remote Control 一択。他は知識として")
rows=[
 ("方式","セットアップ","操作感／用途","推奨"),
 ("① Remote Control（本命）","低・5分","アプリで快適。打席の合間に確認・指示","★★★★★"),
 ("② SSH+Tailscale+Blink","中・30分","ターミナル常駐。画面は小さい","★★★★"),
 ("③ リモートデスク(Jump等)","低","Mac画面ごと操作。遅く見づらい","★★★"),
]
ry,rh=Inches(2.0),Inches(0.62)
cwidths=[Inches(4.1),Inches(2.4),Inches(4.0),Inches(1.73)]
for ri,row in enumerate(rows):
    cx=Inches(0.55); head=(ri==0)
    for ci,cell in enumerate(row):
        bg=RED if head else (REDBG if ci==3 else (CARD if ri%2 else GRYBG))
        bx(s,cx,ry+rh*ri,cwidths[ci],rh,bg,line=CARDLN,lw=0.8)
        tc=WHT if head else (REDD if ci==3 else INK)
        t(s,cell,cx+Inches(0.18),ry+rh*ri,cwidths[ci]-Inches(0.3),rh,sz=12 if head else 11.5,
          bold=head or ci==0 or ci==3,col=tc,align=PP_ALIGN.CENTER if ci==3 else PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
        cx=cx+cwidths[ci]
panel(s,Inches(0.55),Inches(4.95),Inches(12.23),Inches(1.65),"🔐 セキュリティ（人に勧める時の心得）",[
 ("Remote Control＝コードは自宅Mac上で実行・通信はTLS・ポート開放不要（受け身で安全）。iPhoneを失くしたらアカウント側でデバイス登録を取消せば即遮断。",INK),
 ("SSH方式を使うなら必ず Tailscale等の閉域網＋鍵認証。SSHを直接インターネットに開けない。",GRY),
],lh=0.55,bsz=12)

# ── S7 1日の流れ（リモート前提）──
s=sl(); ft(s)
hdr(s,"DAILY FLOW","1日の流れ ── 朝の聖域は守り、打席に集中","旧Macは裏で回り続ける。第一声は「今日の打席は誰？」")
flow=[
 ("朝","調査士(5-7聖域)→当日の打席を1つ決める","旧Mac：朝ブリーフ自動（カネ直結2-3件＋今日の打席は誰か）"),
 ("日中","打席に立つ（会う・聞く・詰める）","スマホ→旧Mac：対話を02へ転記／追客ドラフト／資料作成"),
 ("夜","結果と「決め手」を一言／明日の方針","旧Mac：実績化・KPI更新・明日の打席キュー準備"),
 ("週次(月)","先週レビュー→今週の打席を選ぶ","旧Mac：週次KPIを自動再構築（本部別h・成約率・繰越）"),
]
yy=Inches(2.0)
for tag,human,ai in flow:
    bx(s,Inches(0.55),yy,Inches(1.5),Inches(1.0),REDBG); bx(s,Inches(0.55),yy,Inches(0.1),Inches(1.0),RED)
    t(s,tag,Inches(0.55),yy,Inches(1.5),Inches(1.0),sz=17,bold=True,col=REDD,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bx(s,Inches(2.2),yy,Inches(5.15),Inches(1.0),CARD,line=CARDLN,lw=0.8)
    t(s,"人＝打席",Inches(2.4),yy+Inches(0.1),Inches(4.8),Inches(0.3),sz=9.5,bold=True,col=RED)
    t(s,human,Inches(2.4),yy+Inches(0.4),Inches(4.8),Inches(0.55),sz=12,col=INK,line_sp=1.05)
    bx(s,Inches(7.5),yy,Inches(5.28),Inches(1.0),GRYBG,line=CARDLN,lw=0.8)
    t(s,"旧Mac＝Claude代行",Inches(7.7),yy+Inches(0.1),Inches(4.9),Inches(0.3),sz=9.5,bold=True,col=GRY)
    t(s,ai,Inches(7.7),yy+Inches(0.4),Inches(4.9),Inches(0.55),sz=12,col=INK,line_sp=1.05)
    yy=yy+Inches(1.12)

# ── S8 締め ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"ONE LINE",Inches(0.9),Inches(1.7),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"あなたは、打席に立つだけでいい。",Inches(0.88),Inches(2.5),Inches(11.7),Inches(0.9),sz=34,bold=True,col=INK)
t(s,"自宅Macが、残り全部を24時間回す。",Inches(0.88),Inches(3.35),Inches(11.7),Inches(0.9),sz=34,bold=True,col=RED)
t(s,"それ（会って・聞いて・詰める）が一番難しく、一番価値があり、反復で当たり前になる。\n空いた時間は、調査士と家族へ。",
  Inches(0.9),Inches(4.5),Inches(11.6),Inches(0.9),sz=15,col=INK,line_sp=1.3)
bx(s,Inches(0.9),Inches(5.85),Inches(11.5),Pt(1.2),LINE)
t(s,"KHD オペレーティングモデル × Mac3台・リモート活用  ｜  2026-06-28",Inches(0.9),Inches(5.98),Inches(11),Inches(0.4),sz=12,bold=True,col=INK)

out="01_オペモデル_Mac3台リモート活用_260628.pptx"
prs.save(out)
print("saved:",out, "slides:", len(prs.slides._sldIdLst))

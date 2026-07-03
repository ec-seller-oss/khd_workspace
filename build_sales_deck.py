# -*- coding: utf-8 -*-
"""
KHD AI医療コンサル 営業資料（医師・紹介先に見せるピッチデッキ）
フック＝無料・診療圏ミニ診断。デザイン=クリーム白×レンガ赤・ヒラギノ
出力: KHD_iryo_consul_sales.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE="/Users/kikuchikenta/01_honbu_docs_automation"
CREAM=(249,246,239); BRICK=(170,46,38); DARK=(43,43,43); MUTED=(138,129,122)
WHITE=(255,255,255); LINE=(224,216,204); SOFT=(245,224,220)

def rgb(t): return RGBColor(*t)
def add_bg(s,prs):
    b=s.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb=rgb(CREAM); b.line.fill.background(); b.shadow.inherit=False
    a=s.shapes.add_shape(1,0,0,Inches(0.14),prs.slide_height)
    a.fill.solid(); a.fill.fore_color.rgb=rgb(BRICK); a.line.fill.background(); a.shadow.inherit=False
def tb(s,x,y,w,h,text,size,color,bold=True,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    box=s.shapes.add_textbox(x,y,w,h); tf=box.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.name="ヒラギノ角ゴシック W6" if bold else "ヒラギノ角ゴシック W3"; r.font.color.rgb=rgb(color)
    return box
def bullets(s,x,y,w,h,items,size=18,color=DARK,gap=8):
    box=s.shapes.add_textbox(x,y,w,h); tf=box.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        r=p.add_run(); r.text="・"+it; r.font.size=Pt(size)
        r.font.name="ヒラギノ角ゴシック W3"; r.font.color.rgb=rgb(color)
    return box
def eyebrow(s,txt):
    tb(s,Inches(0.5),Inches(0.4),Inches(12),Inches(0.5),txt,14,BRICK)

def build(path):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]

    # 1 表紙
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.9),Inches(2.2),Inches(11.5),Inches(1.8),"開業・経営を、\n「数字」で勝たせる。",44,DARK)
    tb(s,Inches(0.95),Inches(4.4),Inches(11.5),Inches(0.8),"医療 × 不動産 × AI ─ 家業型の医療コンサル（KHD）",22,BRICK,bold=False)
    tb(s,Inches(0.95),Inches(6.5),Inches(11.5),Inches(0.6),"まずは候補地の「無料・診療圏ミニ診断」から",16,MUTED,bold=False)

    # 2 課題
    s=prs.slides.add_slide(blank); add_bg(s,prs); eyebrow(s,"先生が直面する現実")
    tb(s,Inches(0.5),Inches(1.0),Inches(12),Inches(1),"開業の成否は、腕より先に「立地」で決まる。",30,DARK)
    bullets(s,Inches(0.7),Inches(2.3),Inches(11.8),Inches(4),[
        "どれだけ良い医療でも、診療圏（人口×受療率×競合）を外すと患者は来ない。",
        "従来の診療圏調査は 業者に数十万円・数週間、しかも中身はブラックボックス。",
        "物件は「借りられるか・再建築・境界・科目縛り」まで見ないと後で詰まる。",
        "制度・法規・資金計画…開業は『バラバラの専門家』に振り回されがち。",
    ],size=20,gap=14)

    # 3 だから何が要るか
    s=prs.slides.add_slide(blank); add_bg(s,prs); eyebrow(s,"本当に必要なこと")
    tb(s,Inches(0.5),Inches(1.0),Inches(12),Inches(1),"「立地の数字」と「不動産」を、一人で通しで見る人。",28,DARK)
    bullets(s,Inches(0.7),Inches(2.3),Inches(11.8),Inches(4),[
        "診療圏（数字）→ 物件・契約（不動産）→ 事業計画・制度 を分断せず一気通貫。",
        "「数字上は勝てる」立地でも、物件が出ない・条件が合わなければ机上の空論。",
        "工場型（量産・担当がコロコロ替わる）ではなく、開業から承継まで同じ顔が伴走。",
    ],size=20,gap=14)

    # 4 KHDの強み＝尖り5軸
    s=prs.slides.add_slide(blank); add_bg(s,prs); eyebrow(s,"なぜKHDか（他にない5つ）")
    tb(s,Inches(0.5),Inches(1.0),Inches(12),Inches(0.9),"医療×不動産×AI を、一人称で回す家業型コンサル。",26,DARK)
    bullets(s,Inches(0.7),Inches(2.1),Inches(11.8),Inches(4.5),[
        "① 医療×AI ─ 診療圏をAIで10分・可視化（同業でここまでやる人はまずいない）",
        "② 不動産の目利き ─ 土地家屋調査士×宅建。境界・再建築・土地値・科目縛りまで",
        "③ 家業の伴走 ─ 26年の医療テナント実績。1人の医師の生涯に同じ顔ぶれで",
        "④ 厚利少本 ─ 数を追わず、売り込まず、信頼の対価としてお付き合いする",
        "⑤ 婦人科・小児科の濃い実績",
    ],size=19,gap=12)

    # 5 提供メニュー
    s=prs.slides.add_slide(blank); add_bg(s,prs); eyebrow(s,"ご提供できること")
    tb(s,Inches(0.5),Inches(1.0),Inches(12),Inches(0.9),"無料診断を入口に、必要な所だけ。",28,DARK)
    bullets(s,Inches(0.7),Inches(2.1),Inches(11.8),Inches(4.5),[
        "① 【無料】診療圏ミニ診断 ─ 候補地の「1日に何人来うるか」を概算でお返し",
        "② 開業立地・事業計画 ─ 診療圏の本調査、収支・資金計画、金融機関目線の作り込み",
        "③ テナント・物件・契約 ─ 物件探索、条件交渉、契約前の不動産リスク精査",
        "④ 医業承継 ─ 立地実績を引き継ぐ承継の目利き・値づけ",
    ],size=19,gap=12)

    # 6 AIで何が変わる（デモ図）
    s=prs.slides.add_slide(blank); add_bg(s,prs); eyebrow(s,"AIで、ここまで見える")
    tb(s,Inches(0.5),Inches(0.95),Inches(12),Inches(0.8),"数十万円・数週間 → 10分・可視化。",24,DARK)
    tb(s,Inches(0.7),Inches(2.2),Inches(11.8),Inches(3.5),
       "人口を積み上げて「約21人/日」。\n→ 競合1軒の按分で「4〜6人/日（月100〜160人）」まで圧縮。\n\n人口が多い＝勝てる立地、ではない。競合まで数値化して初めて判断できる。\n（図は面談時に実際のデモ画像を提示：積み上げ図／競合按分図）",22,DARK,bold=False)

    # 7 進め方
    s=prs.slides.add_slide(blank); add_bg(s,prs); eyebrow(s,"進め方（押し売りしません）")
    tb(s,Inches(0.5),Inches(1.0),Inches(12),Inches(0.9),"まず『無料の数字』から。それを見て決めてください。",26,DARK)
    bullets(s,Inches(0.7),Inches(2.2),Inches(11.8),Inches(4),[
        "STEP1 候補地と診療科目を教えていただく（フォーム or 会話で5分）",
        "STEP2 48時間で「診療圏ミニ診断」レポートを無料でお返し",
        "STEP3 見て「もっと詳しく」となれば、本調査・物件・計画へ（必要な所だけ）",
        "料金の目安：スポット調査／伴走コンサルはご相談。厚利少本で数は絞ります。",
    ],size=20,gap=14)

    # 8 CTA
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.9),Inches(2.4),Inches(11.5),Inches(1.6),"まずは、候補地の\n「無料・診療圏ミニ診断」から。",38,DARK)
    tb(s,Inches(0.95),Inches(4.6),Inches(11.5),Inches(1.2),
       "エリアと診療科目を教えていただくだけ。売り込みはしません。\nKHD（菊池）／ 医療×不動産×AI の家業型コンサル",20,BRICK,bold=False)

    prs.save(path); print("saved",path)

if __name__=="__main__":
    build(os.path.join(BASE,"KHD_iryo_consul_sales.pptx"))

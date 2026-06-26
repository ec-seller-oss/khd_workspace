from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── カラーパレット ──
BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # 濃紺
BG_MID    = RGBColor(0x1A, 0x2E, 0x44)   # 中紺
ACCENT    = RGBColor(0xE8, 0xA8, 0x00)   # ゴールド
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xCC, 0xD6, 0xE0)
ACCENT2   = RGBColor(0x4A, 0x9E, 0xCB)   # ライトブルー

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def label(slide, text, x, y):
    box(slide, x, y, Inches(0.08), Inches(0.35), ACCENT)
    txt(slide, text, x + Inches(0.18), y, Inches(10), Inches(0.45),
        size=11, color=LIGHT_GRAY)

def section_num(slide, num, x, y):
    txt(slide, f"{num:02d}", x, y, Inches(2), Inches(1.2),
        size=72, bold=True, color=ACCENT)

def divider_line(slide, y):
    ln = slide.shapes.add_shape(1,
        Inches(0.8), y, Inches(11.73), Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()

def footer(slide, text="医療専門不動産  |  ビジネスモデル"):
    box(slide, 0, H - Inches(0.38), W, Inches(0.38), BG_MID)
    txt(slide, text, Inches(0.5), H - Inches(0.35), Inches(10), Inches(0.32),
        size=9, color=LIGHT_GRAY)

def bullet_block(slide, items, x, y, w, line_h=Inches(0.38), size=16):
    for i, item in enumerate(items):
        txt(slide, item, x, y + i * line_h, w, line_h, size=size, color=WHITE)

def highlight_box(slide, text, x, y, w, h, fsize=20):
    box(slide, x, y, w, h, BG_MID)
    b = slide.shapes.add_shape(1, x, y, Inches(0.06), h, )
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT
    b.line.fill.background()
    txt(slide, text, x + Inches(0.2), y + Inches(0.1),
        w - Inches(0.3), h - Inches(0.2), size=fsize, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 01  カバー
# ════════════════════════════════════════════
sl = add_slide(); bg(sl)
box(sl, 0, 0, W, H, BG_DARK)
box(sl, 0, 0, Inches(0.5), H, ACCENT)
box(sl, 0, H - Inches(1.1), W, Inches(1.1), BG_MID)

txt(sl, "「物件を追うな、情報の非対称性を狩れ」",
    Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.7),
    size=18, color=ACCENT, italic=True)

txt(sl, "医療専門不動産の\nビジネスモデル",
    Inches(0.9), Inches(1.7), Inches(11), Inches(2.4),
    size=44, bold=True, color=WHITE)

txt(sl, "一般不動産屋との事業モデルの比較、そして「30年単位で勝つ」ための設計思想。",
    Inches(0.9), Inches(4.1), Inches(10), Inches(0.6),
    size=16, color=LIGHT_GRAY)

divider_line(sl, Inches(4.85))

txt(sl, "チームてっかん  |  菊池",
    Inches(0.9), Inches(5.1), Inches(6), Inches(0.45),
    size=18, bold=True, color=WHITE)
txt(sl, "宅建業免許：東京都知事 (4)　　2026年",
    Inches(0.9), Inches(5.55), Inches(7), Inches(0.4),
    size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 02  AGENDA
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.2), BG_MID)
txt(sl, "AGENDA  /  本日の構成", Inches(0.6), Inches(0.25),
    Inches(10), Inches(0.7), size=28, bold=True, color=WHITE)

sections = [
    ("01", "事業モデル", "なぜ「医療×不動産」は長期で勝てるのか。\n／情報の非対称性という最強の堀"),
    ("02", "営業ノウハウ", "再現できる「型」になるまでの実装。\n／ハンター式ヒアリングトーク含む"),
    ("03", "成約事例", "1件の入口が、どう収益として重なるか。\n／6,000万粗利事例の全解剖"),
    ("04", "事業展開", "チームとSNSで「ダムを作る」次のステージ。"),
]
col_w = Inches(3.0)
for i, (num, title, desc) in enumerate(sections):
    cx = Inches(0.5) + i * (col_w + Inches(0.18))
    cy = Inches(1.5)
    box(sl, cx, cy, col_w, Inches(4.8), BG_MID)
    box(sl, cx, cy, col_w, Inches(0.06), ACCENT)
    txt(sl, num, cx + Inches(0.2), cy + Inches(0.15),
        col_w, Inches(0.7), size=32, bold=True, color=ACCENT)
    txt(sl, title, cx + Inches(0.2), cy + Inches(0.85),
        col_w - Inches(0.3), Inches(0.55), size=18, bold=True, color=WHITE)
    txt(sl, desc, cx + Inches(0.2), cy + Inches(1.5),
        col_w - Inches(0.3), Inches(3.1), size=13, color=LIGHT_GRAY, wrap=True)

txt(sl, "業界構造 → 営業ノウハウ → 成約事例 → 事業展開の4部構成。最後に「明日からの動き方」を。",
    Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
    size=11, color=LIGHT_GRAY, italic=True)

# ════════════════════════════════════════════
# SLIDE 03  SPEAKER
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.1), BG_MID)
txt(sl, "SPEAKER  /  自己紹介", Inches(0.6), Inches(0.2),
    Inches(10), Inches(0.7), size=26, bold=True, color=WHITE)

txt(sl, "「カバン持ちから始まった、医療専門不動産」",
    Inches(0.7), Inches(1.25), Inches(8), Inches(0.5),
    size=19, bold=True, color=ACCENT)

episode = (
    "大叔父の医業コンサル法人を承継中。宅建業免許：東京都知事 (4)\n\n"
    "【原点エピソード】\n"
    "大叔父は医業コンサルの第一人者。その後ろを歩きながら医師の診察室に入り続けた。\n"
    "気づいたのは「医師は不動産の話を、信頼できる身内にしか相談しない」こと。\n"
    "情報は看板ではなく、信頼の深さに従って流れる──これが業態の本質。"
)
txt(sl, episode, Inches(0.7), Inches(1.85), Inches(7.5), Inches(2.8),
    size=14, color=WHITE, wrap=True)

takeaways = [
    ("01", "医療業界の構造的な「強さ」", "なぜ値下げ競争に巻き込まれないのか。"),
    ("02", "30年単位の顧客LTV設計", "開業→拡張→機材→承継→売却まで関わる仕組み。"),
    ("03", "横のつながり＝チームビルディング", "士業・医療機器・保険・金融を束ねる動き方。"),
    ("04", "明日から動ける、シンプルな原則", "派手な近道はない。1件ずつ、目の前を大切にする。"),
]
for i, (n, t, d) in enumerate(takeaways):
    rx = Inches(0.7) + (i % 2) * Inches(5.9)
    ry = Inches(4.8) + (i // 2) * Inches(0.75)
    txt(sl, f"{n}  {t}", rx, ry, Inches(5.6), Inches(0.38),
        size=13, bold=True, color=ACCENT2)
    txt(sl, d, rx + Inches(0.3), ry + Inches(0.37), Inches(5.3), Inches(0.33),
        size=11, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 04  この場の約束（NEW）
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.1), BG_MID)
txt(sl, "今日お伝えするのは「不動産の売り方」ではない。",
    Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
    size=24, bold=True, color=WHITE)

not_items = ["✕  物件の見つけ方", "✕  SUUMO・レインズの活用法", "✕  仲介手数料の交渉術"]
yes_items = ["✓  情報の源泉（門番）にどう近づくか", "✓  「悩み」から数千万の収益を生む設計", "✓  値下げ競争に入らず30年勝ち続ける構造"]

box(sl, Inches(0.7), Inches(1.35), Inches(5.5), Inches(2.4), BG_MID)
txt(sl, "NOT TEACHING", Inches(0.9), Inches(1.45), Inches(5), Inches(0.4),
    size=12, bold=True, color=LIGHT_GRAY)
for i, it in enumerate(not_items):
    txt(sl, it, Inches(0.9), Inches(1.95) + i * Inches(0.6),
        Inches(5.2), Inches(0.5), size=15, color=RGBColor(0xFF,0x66,0x66))

box(sl, Inches(6.9), Inches(1.35), Inches(6.0), Inches(2.4), BG_MID)
txt(sl, "TEACHING", Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.4),
    size=12, bold=True, color=ACCENT)
for i, it in enumerate(yes_items):
    txt(sl, it, Inches(7.1), Inches(1.95) + i * Inches(0.6),
        Inches(5.7), Inches(0.5), size=14, color=WHITE)

quote = "「物件を追っているうちは、仲介手数料（3%）の奴隷から抜け出せない。\n 10〜30%の粗利は、ネットの向こうではなく『人の悩みの深さ』の中にある。」"
box(sl, Inches(0.7), Inches(4.05), Inches(12.1), Inches(1.6), BG_MID)
b = sl.shapes.add_shape(1, Inches(0.7), Inches(4.05), Inches(0.1), Inches(1.6))
b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background()
txt(sl, quote, Inches(1.05), Inches(4.2), Inches(11.6), Inches(1.3),
    size=17, color=ACCENT, italic=True, wrap=True)

txt(sl, "SUUMOやレインズで同じ情報を見ている同業者との差は、『情報の流れ方の違い』だけ。",
    Inches(0.7), Inches(5.85), Inches(12), Inches(0.4), size=12, color=LIGHT_GRAY, italic=True)

# ════════════════════════════════════════════
# SLIDE 05  SECTION 01 ディバイダー
# ════════════════════════════════════════════
sl = add_slide(); bg(sl)
box(sl, 0, 0, Inches(0.5), H, ACCENT)
box(sl, 0, H - Inches(0.8), W, Inches(0.8), BG_MID)
txt(sl, "SECTION", Inches(0.9), Inches(1.0), Inches(10), Inches(0.8),
    size=22, color=LIGHT_GRAY, bold=False)
txt(sl, "01", Inches(0.9), Inches(1.7), Inches(5), Inches(2.0),
    size=96, bold=True, color=ACCENT)
txt(sl, "事業モデル", Inches(0.9), Inches(3.7), Inches(10), Inches(0.8),
    size=34, bold=True, color=WHITE)
txt(sl, "なぜ「医療×不動産」は、長期で勝てるのか。\n情報の非対称性・業界構造・長期事業の設計、の4枚で。",
    Inches(0.9), Inches(4.6), Inches(11), Inches(1.0),
    size=16, color=LIGHT_GRAY, wrap=True)
txt(sl, "医療専門不動産  |  ビジネスモデル",
    Inches(0.9), H - Inches(0.65), Inches(9), Inches(0.4), size=9, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 06  会社の成り立ち
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "01  /  会社の成り立ちと実績",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "「医療コンサル不動産」は、銀行の現場発注で生まれた。後付けではなく、現場で組み上げた業態。",
    Inches(0.6), Inches(0.75), Inches(11.5), Inches(0.3), size=12, color=LIGHT_GRAY)

steps = [
    ("01", "証券会社時代", "金融の原則を身につける"),
    ("02", "NECメディカル", "医療業界の構造に気付く"),
    ("03", "住友銀行との連携", "「貸せるがコンサルがいない」"),
    ("04", "事業化・会社設立", "型を文書化、事業として確立"),
]
sw = Inches(2.8)
for i, (n, t, d) in enumerate(steps):
    sx = Inches(0.55) + i * (sw + Inches(0.2))
    sy = Inches(1.3)
    box(sl, sx, sy, sw, Inches(1.7), BG_MID)
    box(sl, sx, sy, sw, Inches(0.06), ACCENT)
    txt(sl, n, sx+Inches(0.15), sy+Inches(0.12), sw, Inches(0.45), size=22, bold=True, color=ACCENT)
    txt(sl, t, sx+Inches(0.15), sy+Inches(0.6), sw-Inches(0.2), Inches(0.45), size=15, bold=True, color=WHITE)
    txt(sl, d, sx+Inches(0.15), sy+Inches(1.08), sw-Inches(0.2), Inches(0.5), size=12, color=LIGHT_GRAY)
    if i < 3:
        txt(sl, "→", sx+sw+Inches(0.02), sy+Inches(0.6), Inches(0.22), Inches(0.45), size=18, color=ACCENT)

# エピソード
ep_text = ("【エピソード：住友銀行の担当者の一言】\n"
           "「菊池さん、うちは医師に貸したい。でも事業計画書を正しく読めるコンサルがいないから融資の判断ができない」\n"
           "銀行はお金を持つが医療の目利きができない。私は人脈と知見を持つが融資力が弱い。\n"
           "「タッグを組めば最強だ」──翌日から連携が始まった。")
box(sl, Inches(0.55), Inches(3.2), Inches(12.2), Inches(1.6), BG_MID)
b2 = sl.shapes.add_shape(1, Inches(0.55), Inches(3.2), Inches(0.08), Inches(1.6))
b2.fill.solid(); b2.fill.fore_color.rgb = ACCENT; b2.line.fill.background()
txt(sl, ep_text, Inches(0.8), Inches(3.3), Inches(11.8), Inches(1.4),
    size=12, color=WHITE, wrap=True)

stats = [("140件超", "開業支援"), ("20件超", "事業承継"), ("17件超", "経営支援・救済")]
sw2 = Inches(3.7)
for i, (val, lab) in enumerate(stats):
    sx = Inches(0.55) + i * (sw2 + Inches(0.35))
    box(sl, sx, Inches(5.05), sw2, Inches(1.1), BG_MID)
    txt(sl, val, sx+Inches(0.2), Inches(5.12), sw2, Inches(0.55), size=30, bold=True, color=ACCENT)
    txt(sl, lab, sx+Inches(0.2), Inches(5.65), sw2, Inches(0.38), size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 07  業界の違いと構造的優位
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "02  /  業界の違いと構造的優位",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "医療不動産は「売買」ではなく「コンサルティング」。真の強みは業界横断の「チームビルディング」にある。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

# 比較ボックス
box(sl, Inches(0.55), Inches(1.25), Inches(5.8), Inches(2.5), RGBColor(0x1A,0x1A,0x2E))
txt(sl, "一般不動産　単発の「売買仲介」",
    Inches(0.75), Inches(1.35), Inches(5.4), Inches(0.45), size=14, bold=True, color=LIGHT_GRAY)
for i, t in enumerate(["情報はオープン、利幅は薄い","毎回ゼロから集客","1件決まれば、次は振り出し"]):
    txt(sl, f"・{t}", Inches(0.75), Inches(1.9)+i*Inches(0.55), Inches(5.4), Inches(0.48), size=13, color=LIGHT_GRAY)

box(sl, Inches(7.0), Inches(1.25), Inches(5.8), Inches(2.5), BG_MID)
box(sl, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.06), ACCENT)
txt(sl, "医療不動産　「コンサル」＋「チームビルディング」",
    Inches(7.2), Inches(1.35), Inches(5.4), Inches(0.45), size=14, bold=True, color=WHITE)
for i, t in enumerate(["開業→移転→拡張→機材→承継→売却と連鎖","信頼は次の紹介を呼ぶ","利益が長期で複利的に積み上がる"]):
    txt(sl, f"・{t}", Inches(7.2), Inches(1.9)+i*Inches(0.55), Inches(5.4), Inches(0.48), size=13, color=WHITE)

txt(sl, "VS", Inches(6.1), Inches(2.1), Inches(0.8), Inches(0.7), size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

reasons = [
    ("①  情報の非対称性", "医師は本業で忙しく不動産情報に疎い。\n比較サイトの安値競争に陥らず、信頼で選ばれる市場。"),
    ("②  ライフサイクルの長さ", "1人の医師に30年関わり続けられる。\n開業→拡張→機材→承継→引退後の自宅売却まで。"),
    ("③  チームビルディング", "士業・医療機器・薬品卸・保険・設計建築・金融──\n業界横断のチームを束ねる位置に立つ。"),
]
rw = Inches(3.9)
for i, (t, d) in enumerate(reasons):
    rx = Inches(0.55) + i * (rw + Inches(0.12))
    box(sl, rx, Inches(4.0), rw, Inches(2.1), BG_MID)
    box(sl, rx, Inches(4.0), rw, Inches(0.06), ACCENT)
    txt(sl, t, rx+Inches(0.15), Inches(4.1), rw-Inches(0.2), Inches(0.45), size=14, bold=True, color=ACCENT)
    txt(sl, d, rx+Inches(0.15), Inches(4.6), rw-Inches(0.2), Inches(1.4), size=12, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 08  情報の非対称性（NEW）
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "情報の非対称性  ──  核心スライド",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=ACCENT)
txt(sl, "良い情報は「不動産屋」には来ない。「信頼できる身内」にしか落ちない。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=13, color=LIGHT_GRAY)

# フロー図
flow = ["医師の本音", "門番\n（税理士・弁護士・銀行員）", "私たち\n（コンサルとして認知）", "不動産を\n「処方箋」として活用"]
colors = [BG_MID, BG_MID, RGBColor(0x1A,0x40,0x5E), RGBColor(0x1A,0x50,0x2E)]
fw = Inches(2.7)
for i, (f, c) in enumerate(zip(flow, colors)):
    fx = Inches(0.55) + i * (fw + Inches(0.45))
    box(sl, fx, Inches(1.45), fw, Inches(1.4), c)
    box(sl, fx, Inches(1.45), fw, Inches(0.06), ACCENT)
    txt(sl, f, fx+Inches(0.15), Inches(1.6), fw-Inches(0.2), Inches(1.1),
        size=14, bold=True, color=WHITE, wrap=True, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", fx+fw+Inches(0.05), Inches(1.85), Inches(0.4), Inches(0.5), size=20, color=ACCENT)
    txt(sl, "（信頼があって初めて流れる）" if i==0 else ("門番との関係構築で届く" if i==1 else ""),
        fx+Inches(0.1), Inches(2.95), fw-Inches(0.1), Inches(0.38), size=10, color=LIGHT_GRAY, italic=True)

# 2列比較
box(sl, Inches(0.55), Inches(3.55), Inches(5.8), Inches(1.5), BG_MID)
txt(sl, "一般流通物件", Inches(0.75), Inches(3.65), Inches(5.4), Inches(0.4), size=15, bold=True, color=LIGHT_GRAY)
txt(sl, "仲介手数料3%　全員が同じ情報を持つ", Inches(0.75), Inches(4.1), Inches(5.4), Inches(0.8), size=13, color=LIGHT_GRAY)

box(sl, Inches(7.0), Inches(3.55), Inches(5.8), Inches(1.5), BG_MID)
box(sl, Inches(7.0), Inches(3.55), Inches(5.8), Inches(0.06), ACCENT)
txt(sl, "非公開・難件", Inches(7.2), Inches(3.65), Inches(5.4), Inches(0.4), size=15, bold=True, color=ACCENT)
txt(sl, "粗利10〜30%　情報を持つのは1人だけ", Inches(7.2), Inches(4.1), Inches(5.4), Inches(0.8), size=13, color=WHITE)

ep = ("【エピソード】大叔父のカバン持ち中、医師が診察室でこう言った。\n"
      "「来年クリニックを移転しようと思っているが、どこに相談すればいいか分からなくて……」\n"
      "不動産屋として行ったわけではない。信頼ゼロの場に「連れ」として入ったからこそ聞けた話。\n"
      "情報は看板ではなく、信頼の深さに従って流れる。")
box(sl, Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.5), BG_MID)
b3 = sl.shapes.add_shape(1, Inches(0.55), Inches(5.2), Inches(0.08), Inches(1.5))
b3.fill.solid(); b3.fill.fore_color.rgb = ACCENT; b3.line.fill.background()
txt(sl, ep, Inches(0.8), Inches(5.3), Inches(11.8), Inches(1.3), size=12, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 09  長期事業モデル
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "03  /  長期事業モデルの作り方",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "1件の相談を、生涯にわたる収益ストリームに育てる。派手な急成長を狙わない。複利で積み上がる事業設計。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

streams = [
    ("0年目", "相談・受付", "紹介・無料"),
    ("0–1年目", "診療圏調査", "コンサルfee"),
    ("1年目", "事業計画・物件", "仲介手数料"),
    ("1年〜継続", "機材・備品", "卸売・継続収益"),
    ("5–30年目", "拡張・承継・売却", "再仲介・買取再販"),
]
sw3 = Inches(2.35)
for i, (yr, title, rev) in enumerate(streams):
    sx = Inches(0.45) + i * (sw3 + Inches(0.12))
    alpha = 0.5 + i * 0.12
    c = RGBColor(int(0x1A*alpha), int(0x40*alpha), int(0x5E*(0.5+i*0.1)))
    box(sl, sx, Inches(1.3), sw3, Inches(2.2), BG_MID)
    box(sl, sx, Inches(1.3), sw3, Inches(0.06), ACCENT)
    txt(sl, yr, sx+Inches(0.1), Inches(1.38), sw3-Inches(0.15), Inches(0.38), size=11, color=LIGHT_GRAY)
    txt(sl, title, sx+Inches(0.1), Inches(1.8), sw3-Inches(0.15), Inches(0.6), size=14, bold=True, color=WHITE, wrap=True)
    txt(sl, rev, sx+Inches(0.1), Inches(2.5), sw3-Inches(0.15), Inches(0.8), size=11, color=ACCENT, wrap=True)

txt(sl, "1人の医師から、生涯で  数千万〜数億円規模  の累積取引が積み上がる構造。",
    Inches(0.55), Inches(3.68), Inches(12), Inches(0.45), size=15, bold=True, color=ACCENT)

principles = [
    ("01 利他のマインド", "目先の手数料を追わず、医師のキャリア全体に伴走"),
    ("02 ハブとして繋ぐ", "全部やらない。最初に電話される存在になればいい"),
    ("03 双方向ネットワーク", "「もらう」関係を脱し、「紹介する側」にも回る"),
    ("04 30年の領域選び", "ライフサイクルが長い業界に身を置けばLTVが伸びる"),
    ("05 ノウハウの「型」化", "属人性を抜き、再現性を残す。1件ごとに型を更新"),
]
pw = Inches(2.35)
for i, (t, d) in enumerate(principles):
    px = Inches(0.45) + i * (pw + Inches(0.12))
    box(sl, px, Inches(4.25), pw, Inches(2.0), BG_MID)
    txt(sl, t, px+Inches(0.1), Inches(4.33), pw-Inches(0.15), Inches(0.45), size=12, bold=True, color=ACCENT2)
    txt(sl, d, px+Inches(0.1), Inches(4.85), pw-Inches(0.15), Inches(1.25), size=11, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 10  SECTION 02 ディバイダー
# ════════════════════════════════════════════
sl = add_slide(); bg(sl)
box(sl, 0, 0, Inches(0.5), H, ACCENT)
box(sl, 0, H - Inches(0.8), W, Inches(0.8), BG_MID)
txt(sl, "SECTION", Inches(0.9), Inches(1.0), Inches(10), Inches(0.8), size=22, color=LIGHT_GRAY)
txt(sl, "02", Inches(0.9), Inches(1.7), Inches(5), Inches(2.0), size=96, bold=True, color=ACCENT)
txt(sl, "営業ノウハウ", Inches(0.9), Inches(3.7), Inches(10), Inches(0.8), size=34, bold=True, color=WHITE)
txt(sl, "再現できる「型」になるまで、現場で何をやってきたか。\n5ステップ──士業との関係・相談ハブ化・診療圏調査・事業計画書・ハンター式ヒアリング。",
    Inches(0.9), Inches(4.6), Inches(11), Inches(1.1), size=16, color=LIGHT_GRAY, wrap=True)
txt(sl, "医療専門不動産  |  ビジネスモデル", Inches(0.9), H-Inches(0.65), Inches(9), Inches(0.4), size=9, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 11  STEP01 士業との関係構築
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "STEP 01  /  士業との関係構築",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "「お願いします」をやめた瞬間、紹介が来始めた。5年通った、泥臭い5年。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=12, color=LIGHT_GRAY)

before_items = ["士業事務所を回り案件を「もらう」立場で挨拶", "他業者と並列、紹介は来ない", "条件の悪い案件だけ回ってくる"]
after_items  = ["医師に「税理士いますか？いなければ◯◯先生を」と紹介", "士業に頼られる位置に、案件は双方向に", "コミュニティ内で「あの会社」と固有名詞で呼ばれる"]

box(sl, Inches(0.55), Inches(1.25), Inches(5.8), Inches(2.3), RGBColor(0x1A,0x1A,0x2E))
txt(sl, "BEFORE  「お願いします」の営業", Inches(0.75), Inches(1.33), Inches(5.4), Inches(0.4), size=13, bold=True, color=LIGHT_GRAY)
for i, t in enumerate(before_items):
    txt(sl, f"・{t}", Inches(0.75), Inches(1.83)+i*Inches(0.58), Inches(5.4), Inches(0.5), size=12, color=LIGHT_GRAY)

box(sl, Inches(7.0), Inches(1.25), Inches(5.8), Inches(2.3), BG_MID)
box(sl, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.06), ACCENT)
txt(sl, "AFTER  「ご紹介します」の営業", Inches(7.2), Inches(1.33), Inches(5.4), Inches(0.4), size=13, bold=True, color=ACCENT)
for i, t in enumerate(after_items):
    txt(sl, f"・{t}", Inches(7.2), Inches(1.83)+i*Inches(0.58), Inches(5.4), Inches(0.5), size=12, color=WHITE)

txt(sl, "VS", Inches(6.1), Inches(2.0), Inches(0.8), Inches(0.6), size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

ep = ("【転換点エピソード】税理士の先生への5回目の挨拶。「先週ある医師から信頼できる税理士を探していると聞きまして、ぜひ先生をご紹介したいのですが」\n"
      "その瞬間、先生の表情が変わった。1ヶ月後に電話が来た。「菊池さん、クリニックの移転で相談したい先生がいるんだけど──」\n"
      "最初に投げることで、初めて返ってくる。")
box(sl, Inches(0.55), Inches(3.7), Inches(12.2), Inches(1.3), BG_MID)
b4 = sl.shapes.add_shape(1, Inches(0.55), Inches(3.7), Inches(0.08), Inches(1.3))
b4.fill.solid(); b4.fill.fore_color.rgb = ACCENT; b4.line.fill.background()
txt(sl, ep, Inches(0.8), Inches(3.8), Inches(11.8), Inches(1.1), size=12, color=WHITE, wrap=True)

reality = [("断られても顔を出し続ける", "1回の挨拶では何も起きない。5回、10回通ってようやく覚えられる。"),
           ("「便利屋」上等で動く", "雑用でも引き受ける。「あの会社いつでも対応してくれる」が入口。"),
           ("勉強会・葬式まで出る", "業務外で会う回数が、信頼の上限を決める。"),
           ("「先に紹介」を実装する", "もらう前に、まず投げる。義理で紹介が返ってくる。")]
rw2 = Inches(2.95)
for i, (t, d) in enumerate(reality):
    rx = Inches(0.55) + i*(rw2+Inches(0.12))
    box(sl, rx, Inches(5.15), rw2, Inches(1.1), BG_MID)
    txt(sl, t, rx+Inches(0.12), Inches(5.23), rw2-Inches(0.2), Inches(0.38), size=11, bold=True, color=ACCENT2)
    txt(sl, d, rx+Inches(0.12), Inches(5.63), rw2-Inches(0.2), Inches(0.55), size=10, color=LIGHT_GRAY, wrap=True)

# ════════════════════════════════════════════
# SLIDE 12  STEP02 相談プラットフォーム化
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "STEP 02  /  相談プラットフォーム化",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "横のつながりが、そのまま「集客の入口」になる。「最初に電話される存在」になる。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=12, color=LIGHT_GRAY)

# HUBビジュアル - 中心
box(sl, Inches(5.4), Inches(2.6), Inches(2.5), Inches(1.1), ACCENT)
txt(sl, "医療チームのHUB", Inches(5.4), Inches(2.65), Inches(2.5), Inches(1.0),
    size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

partners = [
    ("医師・クリニック", "起点", Inches(5.55), Inches(1.3)),
    ("会計士・税理士", "提携", Inches(0.6), Inches(2.0)),
    ("弁護士・司法書士", "提携", Inches(0.6), Inches(3.2)),
    ("医療機器メーカー", "自社対応", Inches(0.6), Inches(4.4)),
    ("薬品卸・販売企業", "提携", Inches(0.6), Inches(5.5)),
    ("保険代理店", "提携", Inches(9.5), Inches(2.0)),
    ("設計建築・内装", "提携", Inches(9.5), Inches(3.2)),
    ("ハウスメーカー", "提携", Inches(9.5), Inches(4.4)),
    ("金融機関", "提携", Inches(9.5), Inches(5.5)),
    ("人材紹介", "提携", Inches(5.55), Inches(5.7)),
]
for name, role, px, py in partners:
    box(sl, px, py, Inches(2.8), Inches(0.65), BG_MID)
    txt(sl, name, px+Inches(0.1), py+Inches(0.05), Inches(2.0), Inches(0.35), size=12, bold=True, color=WHITE)
    txt(sl, role, px+Inches(0.1), py+Inches(0.38), Inches(1.5), Inches(0.25), size=9, color=ACCENT)

txt(sl, "自分で全部やる必要はない。「最初に電話される存在」になれば、案件は向こうから来る。",
    Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.4), size=13, bold=True, color=ACCENT, italic=True)

# ════════════════════════════════════════════
# SLIDE 13  STEP03 診療圏調査
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "STEP 03  /  診療圏調査",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "数字で「ここなら勝てる」を示す実物の資料。サンプル：T整形外科リウマチ科（川崎市宮前区／鷺沼駅圏）",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

process_steps = [
    ("01", "診療圏マップ確定", "地形・道路・公共交通で\nエリアを視覚化"),
    ("02", "競合施設の調査", "競合7施設をリスト化\n診療時間・推定外来数"),
    ("03", "診療圏人口集計", "10町丁目を集計\n総人口48,201人"),
    ("04", "受療率の適用", "厚労省データで年齢×\n性別を当て込み最適化"),
    ("05", "自院推定外来算定", "(人口×受療率)÷\n(競合数+1)で算出"),
]
pw2 = Inches(2.35)
for i, (n, t, d) in enumerate(process_steps):
    px = Inches(0.45) + i*(pw2+Inches(0.12))
    box(sl, px, Inches(1.28), pw2, Inches(1.9), BG_MID)
    box(sl, px, Inches(1.28), pw2, Inches(0.06), ACCENT)
    txt(sl, n, px+Inches(0.12), Inches(1.35), pw2, Inches(0.38), size=18, bold=True, color=ACCENT)
    txt(sl, t, px+Inches(0.12), Inches(1.75), pw2-Inches(0.2), Inches(0.45), size=12, bold=True, color=WHITE, wrap=True)
    txt(sl, d, px+Inches(0.12), Inches(2.25), pw2-Inches(0.2), Inches(0.8), size=10, color=LIGHT_GRAY, wrap=True)

output_items = [
    ("対象エリア", "鷺沼駅 / 半径2km"),
    ("診療圏人口", "48,201 人"),
    ("65歳以上比率", "23.8 %"),
    ("競合施設", "7 施設"),
    ("競合の総外来", "523 人/日"),
    ("自院推定（昼）", "133 人/日"),
    ("自院推定（夜）", "123 人/日"),
]
ow = Inches(1.65)
for i, (lab, val) in enumerate(output_items):
    ox = Inches(0.45) + i*(ow+Inches(0.12))
    box(sl, ox, Inches(3.38), ow, Inches(1.3), BG_MID)
    txt(sl, val, ox+Inches(0.1), Inches(3.45), ow-Inches(0.15), Inches(0.65), size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(sl, lab, ox+Inches(0.1), Inches(4.1), ow-Inches(0.15), Inches(0.45), size=9, color=LIGHT_GRAY, wrap=True, align=PP_ALIGN.CENTER)

txt(sl, "テンプレに見えるが、過去140件超の現場データに裏打ちされた一品物。AI時代でも、ここは経験で差が出る。",
    Inches(0.55), Inches(4.88), Inches(12.2), Inches(0.38), size=12, color=LIGHT_GRAY, italic=True)

txt(sl, "【現場の声】「こんな資料を出してくれた業者は初めてです。先生は本当に私のことを考えてくれているんですね」──数字は信頼を作る。感情の前に根拠を示せる専門家になれ。",
    Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.75), size=12, color=WHITE, italic=True, wrap=True)

# ════════════════════════════════════════════
# SLIDE 14  STEP04 事業計画書
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "STEP 04  /  事業計画書",
    Inches(0.6), Inches(0.18), Inches(10), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "銀行が通す事業計画書──融資が通れば、すべてが動く。必要資金2.3億円 / 借入2億円（利率1.3%・15年元利均等）",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

format_items = [
    ("①資金調達", "自己資金/借入/利率/返済"),
    ("②概算総事業費", "建設・不動産・諸費用・運転"),
    ("③施設概要", "法人・住所・診療科・床面積"),
    ("④減価償却", "建物・設備・医療機器"),
    ("⑤人件費", "常勤・パート・福利・賞与"),
    ("⑥収支算定", "外来×単価×日数/経費率"),
    ("⑦年度別損益", "5期分・税引後・返済財源"),
    ("⑧借入返済明細", "元金均等の年次CF"),
]
fw2 = Inches(2.8)
for i, (t, d) in enumerate(format_items[:4]):
    fx = Inches(0.55) + i*(fw2+Inches(0.1))
    box(sl, fx, Inches(1.28), fw2, Inches(1.0), BG_MID)
    txt(sl, t, fx+Inches(0.12), Inches(1.33), fw2-Inches(0.2), Inches(0.38), size=13, bold=True, color=ACCENT2)
    txt(sl, d, fx+Inches(0.12), Inches(1.73), fw2-Inches(0.2), Inches(0.42), size=11, color=LIGHT_GRAY)
for i, (t, d) in enumerate(format_items[4:]):
    fx = Inches(0.55) + i*(fw2+Inches(0.1))
    box(sl, fx, Inches(2.38), fw2, Inches(1.0), BG_MID)
    txt(sl, t, fx+Inches(0.12), Inches(2.43), fw2-Inches(0.2), Inches(0.38), size=13, bold=True, color=ACCENT2)
    txt(sl, d, fx+Inches(0.12), Inches(2.83), fw2-Inches(0.2), Inches(0.42), size=11, color=LIGHT_GRAY)

years = [("R8（初年度）","+2,970"),("R9","+7,962"),("R10","+12,966"),("R11","+18,618"),("R12","+25,822")]
yw = Inches(2.35)
for i, (yr, val) in enumerate(years):
    yx = Inches(0.55)+i*(yw+Inches(0.12))
    box(sl, yx, Inches(3.55), yw, Inches(1.2), BG_MID)
    box(sl, yx, Inches(3.55), yw, Inches(0.06), ACCENT)
    txt(sl, yr, yx+Inches(0.1), Inches(3.62), yw-Inches(0.15), Inches(0.38), size=11, color=LIGHT_GRAY)
    txt(sl, f"{val} 千円", yx+Inches(0.1), Inches(4.02), yw-Inches(0.15), Inches(0.58), size=18, bold=True, color=ACCENT)

txt(sl, "※ 5年で借入残高1.0億円・返済原資が剰余に転じる",
    Inches(0.55), Inches(4.9), Inches(12), Inches(0.35), size=11, color=LIGHT_GRAY)
txt(sl, "テンプレのように見えるが、銀行を通す形にできるのは経験の積み重ねだけ。\n「この計画書を持ってくる業者の案件なら、うちは前向きに検討します」──融資担当者を動かす資料が作れる。これが最大の差。",
    Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.85), size=12, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 15  STEP05 ハンター式ヒアリング（NEW）
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "STEP 05  /  ハンター式ヒアリングトーク  【NEW】",
    Inches(0.6), Inches(0.18), Inches(11), Inches(0.65), size=22, bold=True, color=ACCENT)
txt(sl, "「不動産屋に話さない本音」を引き出す3ステップ。入り口は絶対に「不動産」の話をしない。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=12, color=LIGHT_GRAY)

steps_h = [
    ("STEP①\n業界特有の「苦しみ」への共感",
     "「先生、最近は診療報酬の改定で、現場を回すだけでも手出しが増えてるんじゃないですか？」",
     "→ 狙い：「この人は自分の業界を知っている」というインサイダー認定を得る。"),
    ("STEP②\n最悪のシナリオの言語化",
     "「もし今月、メインバンクからの融資が止まったら、どの資産から手をつけますか？」\n「後継者の方とは、万が一の際の出口について具体的に話したことはありますか？」",
     "→ 狙い：相手が「言えない悩み」を持っているかを確認する。"),
    ("STEP③\n処方箋の提示（出口）",
     "「売ることが目的ではなく、先生の法人の『再起動』が目的です。\n 私が泥を被って、1ヶ月以内に現金を用意するスキームを組みます。」",
     "→ 「守るための現金化」として提案する。「業者」ではなく「救済者」として認識される。"),
]
sh = Inches(1.85)
for i, (t, q, note) in enumerate(steps_h):
    box(sl, Inches(0.55), Inches(1.28)+i*(sh+Inches(0.1)), Inches(12.2), sh, BG_MID)
    box(sl, Inches(0.55), Inches(1.28)+i*(sh+Inches(0.1)), Inches(0.08), sh, ACCENT)
    txt(sl, t, Inches(0.8), Inches(1.35)+i*(sh+Inches(0.1)),
        Inches(2.8), sh-Inches(0.2), size=13, bold=True, color=ACCENT, wrap=True)
    txt(sl, q, Inches(3.7), Inches(1.35)+i*(sh+Inches(0.1)),
        Inches(5.5), sh-Inches(0.25), size=12, color=WHITE, wrap=True, italic=True)
    txt(sl, note, Inches(9.3), Inches(1.35)+i*(sh+Inches(0.1)),
        Inches(3.3), sh-Inches(0.25), size=10, color=LIGHT_GRAY, wrap=True)

txt(sl, "相手が「この人は同じ側の人間だ」と感じた瞬間から、本当の営業が始まる。",
    Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35), size=12, color=ACCENT, italic=True)

# ════════════════════════════════════════════
# SLIDE 16  SECTION 03 ディバイダー
# ════════════════════════════════════════════
sl = add_slide(); bg(sl)
box(sl, 0, 0, Inches(0.5), H, ACCENT)
box(sl, 0, H-Inches(0.8), W, Inches(0.8), BG_MID)
txt(sl, "SECTION", Inches(0.9), Inches(1.0), Inches(10), Inches(0.8), size=22, color=LIGHT_GRAY)
txt(sl, "03", Inches(0.9), Inches(1.7), Inches(5), Inches(2.0), size=96, bold=True, color=ACCENT)
txt(sl, "成約事例", Inches(0.9), Inches(3.7), Inches(10), Inches(0.8), size=34, bold=True, color=WHITE)
txt(sl, "1件の入口が、どう収益として重なっていくか。\n6,000万粗利の全解剖と、対照的な2件──「単発で終わらない」事業の実態を。",
    Inches(0.9), Inches(4.6), Inches(11), Inches(1.0), size=16, color=LIGHT_GRAY, wrap=True)
txt(sl, "医療専門不動産  |  ビジネスモデル", Inches(0.9), H-Inches(0.65), Inches(9), Inches(0.4), size=9, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 17  CASE01 弁護士ルート
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "CASE 01  /  弁護士ルート  ──  横のつながりが直接効いた事例",
    Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), size=22, bold=True, color=WHITE)
txt(sl, "「医療を分かる人いる？」──弁護士からの1本の電話。一般不動産屋が断った居抜き案件。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

case_items = [
    ("01 入口", "弁護士から居抜きの紹介", "電源・水回り・フロア用途を評価できず一般屋で断られた案件が、医療を分かる窓口として回ってきた。"),
    ("02 なぜ我々に", "「医療だけ」と思われていない", "事務所内で「人柄含めて信頼できる窓口」と認知。不動産が絡む案件はカジュアルに投げてもらえる位置。"),
    ("03 アプローチ", "自社買取＋再販", "売り急ぎで市場に出すと足元を見られる。買取で時間を作り、医療向けにバリューアップして次の使い手へ。"),
    ("04 メッセージ", "専門特化＝逆に間口が広がる", "医療を切り口にすると、一般不動産の窓口にもなれる。コンサルできる軸を1本持つことが効く。"),
]
cw = Inches(2.95)
for i, (n, t, d) in enumerate(case_items):
    cx = Inches(0.55)+i*(cw+Inches(0.12))
    box(sl, cx, Inches(1.25), cw, Inches(2.8), BG_MID)
    box(sl, cx, Inches(1.25), cw, Inches(0.06), ACCENT)
    txt(sl, n, cx+Inches(0.12), Inches(1.32), cw-Inches(0.2), Inches(0.38), size=11, color=LIGHT_GRAY)
    txt(sl, t, cx+Inches(0.12), Inches(1.73), cw-Inches(0.2), Inches(0.5), size=13, bold=True, color=WHITE, wrap=True)
    txt(sl, d, cx+Inches(0.12), Inches(2.28), cw-Inches(0.2), Inches(1.65), size=11, color=LIGHT_GRAY, wrap=True)

cp_items = [("診療圏調査fee","未"), ("事業計画書fee","未"), ("物件仲介手数料","✓"), ("買取・再販差益","✓"), ("機材・リース","未"), ("士業・税理士紹介","✓"), ("顧問料","未"), ("承継・M&A","未")]
box(sl, Inches(0.55), Inches(4.25), Inches(12.2), Inches(0.4), BG_MID)
txt(sl, "CASHPOINT MAP", Inches(0.7), Inches(4.3), Inches(4), Inches(0.3), size=11, bold=True, color=LIGHT_GRAY)
cpw = Inches(1.45)
for i, (lab, st) in enumerate(cp_items):
    cx = Inches(0.55)+i*(cpw+Inches(0.06))
    c = ACCENT if st=="✓" else RGBColor(0x33,0x44,0x55)
    box(sl, cx, Inches(4.75), cpw, Inches(0.75), c)
    txt(sl, lab, cx+Inches(0.05), Inches(4.82), cpw-Inches(0.08), Inches(0.38), size=9, color=WHITE, wrap=True)
    txt(sl, st, cx+Inches(0.05), Inches(5.2), cpw-Inches(0.08), Inches(0.25), size=10, bold=True, color=BG_DARK if st=="✓" else LIGHT_GRAY)

ep2 = ("【エピソード：弁護士からの電話】「医療を分かる人いる？顧問先のクリニックが資金繰りで詰まっていて、居抜きで売りたいが不動産屋が全員断ってくる」\n"
       "断られた理由がすぐにわかった。医療機器の残存価値、内装の用途変更費用、電気容量──一般の業者には読めない。「明日8時に現地へ」翌朝確認、当日夕方に買取価格を提示した。")
box(sl, Inches(0.55), Inches(5.65), Inches(12.2), Inches(1.05), BG_MID)
b5 = sl.shapes.add_shape(1, Inches(0.55), Inches(5.65), Inches(0.08), Inches(1.05))
b5.fill.solid(); b5.fill.fore_color.rgb = ACCENT; b5.line.fill.background()
txt(sl, ep2, Inches(0.8), Inches(5.73), Inches(11.8), Inches(0.9), size=11, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 18  CASE02 危機対応
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "CASE 02  /  危機対応  ──  長期顧客化の典型例",
    Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), size=22, bold=True, color=WHITE)
txt(sl, "「閉院 vs 救済」の岐路で、医師の人生に伴走した数年。事務長の使い込みで資金破綻寸前。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

case2 = [
    ("01 入口", "事務長の使い込みで資金破綻寸前", "「人」の問題で経営が傾いた典型例。チームを設計してきたから相談が来る。"),
    ("02 普通なら", "閉院→物件売却で「終わり」", "医師は廃業、設備は二束三文、患者は離散。一般不動産屋ならここまで。"),
    ("03 アプローチ", "自宅売却＋チームから人材派遣", "資金は不動産売却で確保。チームから事務長を送り込み、診療を止めずに再建。"),
    ("04 結果", "売却益→派遣→顧問料の連結", "1件の不動産案件が、派遣収入と顧問料の継続収益に変わった。"),
]
cw2 = Inches(2.95)
for i, (n, t, d) in enumerate(case2):
    cx = Inches(0.55)+i*(cw2+Inches(0.12))
    box(sl, cx, Inches(1.25), cw2, Inches(2.5), BG_MID)
    box(sl, cx, Inches(1.25), cw2, Inches(0.06), ACCENT)
    txt(sl, n, cx+Inches(0.12), Inches(1.32), cw2-Inches(0.2), Inches(0.38), size=11, color=LIGHT_GRAY)
    txt(sl, t, cx+Inches(0.12), Inches(1.73), cw2-Inches(0.2), Inches(0.5), size=13, bold=True, color=WHITE, wrap=True)
    txt(sl, d, cx+Inches(0.12), Inches(2.28), cw2-Inches(0.2), Inches(1.35), size=11, color=LIGHT_GRAY, wrap=True)

rev_items = [
    ("Y1", "自宅売却仲介＋利益", "資金繰り解消、当社の取引利益"),
    ("Y1", "事務長の派遣収入", "チームから人材を送り売上化"),
    ("継続", "顧問料（毎月）", "経営支援の継続収入"),
    ("Y2+", "クリニック移転・機材", "立て直し後の物件・設備案件"),
    ("Y5+", "事業承継の伴走", "後継候補接続・承継後も継続"),
]
rw3 = Inches(2.35)
for i, (yr, t, d) in enumerate(rev_items):
    rx = Inches(0.55)+i*(rw3+Inches(0.12))
    box(sl, rx, Inches(3.92), rw3, Inches(1.5), BG_MID)
    box(sl, rx, Inches(3.92), rw3, Inches(0.06), ACCENT)
    txt(sl, yr, rx+Inches(0.12), Inches(3.98), rw3-Inches(0.2), Inches(0.35), size=13, bold=True, color=ACCENT)
    txt(sl, t, rx+Inches(0.12), Inches(4.35), rw3-Inches(0.2), Inches(0.45), size=12, bold=True, color=WHITE, wrap=True)
    txt(sl, d, rx+Inches(0.12), Inches(4.82), rw3-Inches(0.2), Inches(0.5), size=10, color=LIGHT_GRAY, wrap=True)

txt(sl, "学び：「人」の問題まで踏み込めると、不動産1件＋派遣＋顧問料で収益が重なる。チームがあるからこそ動ける。",
    Inches(0.55), Inches(5.58), Inches(12.2), Inches(0.38), size=12, color=LIGHT_GRAY, italic=True)
txt(sl, "→ 次のスライド：6,000万が生まれた瞬間──この構造を徹底解剖",
    Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.35), size=12, color=ACCENT, bold=True)

# ════════════════════════════════════════════
# SLIDE 19  6,000万が生まれた瞬間（NEW）
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), RGBColor(0x1A,0x2E,0x1A))
txt(sl, "6,000万が生まれた瞬間  ──  破産寸前→資産現金化→再起動",
    Inches(0.6), Inches(0.15), Inches(12), Inches(0.65), size=22, bold=True, color=ACCENT)
txt(sl, "医療法人救済の全記録。半年で完結。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=12, color=LIGHT_GRAY)

timeline = [
    ("0ヶ月目\n接触", "大叔父（医業コンサル第一人者）のカバン持ちとして医療法人の経営相談に同席。不動産業者としてではなく「信頼できる連れ」として入室。"),
    ("3〜5ヶ月目\n発掘", "「事務長が資金を横領していた」という極秘情報が門番（大叔父）を通じて届く。資金ショートまであと2週間。"),
    ("6ヶ月目\n信頼の借用", "「大叔父が信頼している人」というお墨付きで院長に面談。警戒心はゼロ。「任せます」の一言をその場でもらった。"),
    ("決済期\nスピードが命", "ノンバンク×プロジェクト融資の知見を駆使。自宅＋医院不動産を48時間で査定・提示。「閉院」を「再起動」に変えた。"),
]
tw = Inches(2.95)
for i, (t, d) in enumerate(timeline):
    tx = Inches(0.55)+i*(tw+Inches(0.12))
    box(sl, tx, Inches(1.28), tw, Inches(3.0), BG_MID)
    box(sl, tx, Inches(1.28), tw, Inches(0.06), ACCENT)
    txt(sl, t, tx+Inches(0.12), Inches(1.35), tw-Inches(0.2), Inches(0.7), size=13, bold=True, color=ACCENT, wrap=True)
    txt(sl, d, tx+Inches(0.12), Inches(2.1), tw-Inches(0.2), Inches(2.0), size=12, color=WHITE, wrap=True)

box(sl, Inches(0.55), Inches(4.45), Inches(12.2), Inches(0.9), RGBColor(0x1A,0x3A,0x1A))
txt(sl, "結果：自社粗利  約6,000万円",
    Inches(0.75), Inches(4.52), Inches(6), Inches(0.45), size=22, bold=True, color=ACCENT)
txt(sl, "院長の感想：「あなたがいなかったら終わっていた」",
    Inches(0.75), Inches(4.97), Inches(12), Inches(0.3), size=14, color=WHITE, italic=True)

ep3 = ("【院長の一言】決済が終わった夜、院長から電話があった。「最初は半信半疑でした。でもあなたは逃げなかった。48時間で現金を用意してくれた。うちの病院は終わりかけていたのに、今ここで診療を続けられている。本当にありがとうございます」\n"
       "不動産は手段に過ぎない。誰かの人生の危機に、最後まで伴走できるかどうか。それだけです。")
box(sl, Inches(0.55), Inches(5.52), Inches(12.2), Inches(1.25), BG_MID)
b6 = sl.shapes.add_shape(1, Inches(0.55), Inches(5.52), Inches(0.08), Inches(1.25))
b6.fill.solid(); b6.fill.fore_color.rgb = ACCENT; b6.line.fill.background()
txt(sl, ep3, Inches(0.8), Inches(5.6), Inches(11.8), Inches(1.1), size=12, color=WHITE, wrap=True)

# ════════════════════════════════════════════
# SLIDE 20  PATTERN CATALOG
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "PATTERN CATALOG  /  案件のバラエティ",
    Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), size=24, bold=True, color=WHITE)
txt(sl, "チームを束ねている結果、これらが日常的に発生する──「単発仲介」では捉えられない世界。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

patterns = [
    ("P01", "新規開業の物件選定", "診療圏調査→候補抽出→賃貸契約まで士業ハブで完結。"),
    ("P02", "クリニック移転・拡張", "現物件の出口設計と新物件選定を同時に。"),
    ("P03", "医療法人化に伴う本部移転", "登記要件と運営動線を両立。会計士・行政書士と並走。"),
    ("P04", "機材入替・リース紹介", "提携機材商社へ。設備更新タイミングで再接触。"),
    ("P05", "内装リニューアル", "提携内装業者へ繋ぎ、紹介マージン。動線改善も。"),
    ("P06", "事業承継・M&A", "後継候補との接続、銀行調整、設備譲渡。1案件で複数収益。"),
    ("P07", "廃業・救済（債権者対応）", "閉院ではなく事業継続の道を提示。弁護士・会計士と組み数年並走。"),
    ("P08", "多店舗展開支援", "2院目・3院目。組織設計と物件取得を同時伴走。"),
]
pw3 = Inches(2.95)
for i, (n, t, d) in enumerate(patterns):
    px = Inches(0.45) + (i%4)*(pw3+Inches(0.12))
    py = Inches(1.28) + (i//4)*Inches(2.35)
    box(sl, px, py, pw3, Inches(2.15), BG_MID)
    box(sl, px, py, pw3, Inches(0.06), ACCENT)
    txt(sl, n, px+Inches(0.12), py+Inches(0.12), pw3, Inches(0.38), size=16, bold=True, color=ACCENT)
    txt(sl, t, px+Inches(0.12), py+Inches(0.55), pw3-Inches(0.2), Inches(0.5), size=12, bold=True, color=WHITE, wrap=True)
    txt(sl, d, px+Inches(0.12), py+Inches(1.1), pw3-Inches(0.2), Inches(0.9), size=11, color=LIGHT_GRAY, wrap=True)

txt(sl, "年140件超の中で、「単発で終わる案件」はほとんど無い。どのパターンも、次のパターンの入口になっている。",
    Inches(0.45), Inches(6.9), Inches(12.2), Inches(0.38), size=12, color=LIGHT_GRAY, italic=True)

# ════════════════════════════════════════════
# SLIDE 21  SECTION 04 ディバイダー（NEW）
# ════════════════════════════════════════════
sl = add_slide(); bg(sl)
box(sl, 0, 0, Inches(0.5), H, ACCENT)
box(sl, 0, H-Inches(0.8), W, Inches(0.8), BG_MID)
txt(sl, "SECTION", Inches(0.9), Inches(1.0), Inches(10), Inches(0.8), size=22, color=LIGHT_GRAY)
txt(sl, "04", Inches(0.9), Inches(1.7), Inches(5), Inches(2.0), size=96, bold=True, color=ACCENT)
txt(sl, "事業展開", Inches(0.9), Inches(3.7), Inches(10), Inches(0.8), size=34, bold=True, color=WHITE)
txt(sl, "チームとSNSで「信頼のダムを作る」次のステージ。\n1人でやる限界を超えるための設計。",
    Inches(0.9), Inches(4.6), Inches(11), Inches(1.0), size=16, color=LIGHT_GRAY, wrap=True)
txt(sl, "医療専門不動産  |  ビジネスモデル", Inches(0.9), H-Inches(0.65), Inches(9), Inches(0.4), size=9, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 22  次のステージ（NEW）
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "次のステージ  /  「信頼のダムを作る」  【NEW】",
    Inches(0.6), Inches(0.18), Inches(11), Inches(0.65), size=22, bold=True, color=ACCENT)
txt(sl, "あなたが構築すべきは「不動産を売る仕組み」ではなく、「誰にも言えない悩みが自動的に集まってくるダム」。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=11, color=LIGHT_GRAY)

txt(sl, "入り口はコンサル、出口はハンター。",
    Inches(0.6), Inches(1.2), Inches(12), Inches(0.45), size=18, bold=True, color=WHITE)

sns_items = [
    ("① 恐怖と救済\n（エピソード発信）",
     "「今月ショートする医療法人。事務長は逃亡。残されたのは動かない不動産。そこで私がやったこと──」\n→ 特定層が震えるリアルな実話を投稿。"),
    ("② 情報の非対称性\n（専門家発信）",
     "「不動産業者に相談すると叩かれるだけの物件も、医療コンサルの視点を加えると価値が3倍になる理由」\n→ 裏技・専門知識を公開し権威を確立。"),
    ("③ 逃げない覚悟\n（人間性発信）",
     "「カバン持ちから学んだ、最後の一人まで向き合う泥臭さ」\n→ 顧客が最も恐れる「業者の逃げ」を否定する。"),
]
sw4 = Inches(3.9)
for i, (t, d) in enumerate(sns_items):
    sx = Inches(0.55)+i*(sw4+Inches(0.12))
    box(sl, sx, Inches(1.8), sw4, Inches(2.8), BG_MID)
    box(sl, sx, Inches(1.8), sw4, Inches(0.06), ACCENT)
    txt(sl, t, sx+Inches(0.15), Inches(1.88), sw4-Inches(0.25), Inches(0.75), size=13, bold=True, color=ACCENT, wrap=True)
    txt(sl, d, sx+Inches(0.15), Inches(2.7), sw4-Inches(0.25), Inches(1.75), size=12, color=WHITE, wrap=True)

team_steps = ["各業界の「太客の隣」に潜り込む\n（医療・福祉・士業など）",
              "難件を拾ってくる\n（入り口コンサル）",
              "出口の刈り取りを\n指揮・バックアップ",
              "粗利を分配し\n実績としてコンテンツ化"]
tw2 = Inches(2.95)
for i, t in enumerate(team_steps):
    tx = Inches(0.55)+i*(tw2+Inches(0.12))
    box(sl, tx, Inches(4.8), tw2, Inches(1.1), RGBColor(0x1A,0x38,0x50))
    txt(sl, f"Step {i+1}", tx+Inches(0.12), Inches(4.87), tw2-Inches(0.2), Inches(0.3), size=10, color=ACCENT)
    txt(sl, t, tx+Inches(0.12), Inches(5.2), tw2-Inches(0.2), Inches(0.6), size=11, color=WHITE, wrap=True)

phrases = ["「不動産屋を名乗っているうちは、仲介手数料の枠から出られない。」",
           "「一生営業したくないなら、一生モノの『信頼のダム』を今すぐ掘れ。」"]
for i, p in enumerate(phrases):
    txt(sl, p, Inches(0.6), Inches(6.08)+i*Inches(0.38), Inches(12.2), Inches(0.35),
        size=12, color=ACCENT, italic=True, bold=True)

# ════════════════════════════════════════════
# SLIDE 23  CLOSING
# ════════════════════════════════════════════
sl = add_slide(); bg(sl); footer(sl)
box(sl, 0, 0, W, Inches(1.05), BG_MID)
txt(sl, "CLOSING  /  明日から動ける、ノウハウの結論",
    Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), size=22, bold=True, color=WHITE)
txt(sl, "派手な近道はない。1件ずつ、目の前を大切にする。",
    Inches(0.6), Inches(0.75), Inches(12), Inches(0.3), size=13, color=LIGHT_GRAY)

closes = [
    ("フットワーク軽く", "連絡は即返信、スピードで信頼を作る。挨拶は出向く。"),
    ("顔を出し続けろ", "1回では覚えられない。3回、5回、10回。義理が信頼に化ける。"),
    ("足を使え", "現地・士業事務所・医師の診察室。机では分からないことが現場にある。"),
    ("頭を使え", "30年スパンで考える。1件の手数料ではなく、生涯LTVから逆算する。"),
    ("金を使え", "情報・人脈・道具に投資する。ケチると、機会のほうが先に去る。"),
    ("まず1件取れ", "完璧を待たずに動く。1件目を「事故なく」完遂すれば、次の100件の入口になる。"),
]
cw3 = Inches(3.9)
for i, (t, d) in enumerate(closes):
    cx = Inches(0.45)+(i%3)*(cw3+Inches(0.12))
    cy = Inches(1.28)+(i//3)*Inches(1.35)
    box(sl, cx, cy, cw3, Inches(1.2), BG_MID)
    box(sl, cx, cy, cw3, Inches(0.06), ACCENT)
    txt(sl, t, cx+Inches(0.15), cy+Inches(0.12), cw3-Inches(0.25), Inches(0.38), size=14, bold=True, color=ACCENT)
    txt(sl, d, cx+Inches(0.15), cy+Inches(0.55), cw3-Inches(0.25), Inches(0.55), size=11, color=WHITE, wrap=True)

ep4 = ("【エピソード：最初の1件】正直、何もわからなかった。事業計画書の書き方も銀行への提案の仕方も完璧じゃなかった。\n"
       "でも逃げなかった。深夜まで資料を直した。融資担当者に電話をかけ続けた。\n"
       "その1件が完遂できた日、医師からこう言われた。「次の先生も紹介しますよ」──それが2件目の入口だった。完璧を待っていたら、その1件も2件目もなかった。")
box(sl, Inches(0.45), Inches(4.1), Inches(12.2), Inches(1.4), BG_MID)
b7 = sl.shapes.add_shape(1, Inches(0.45), Inches(4.1), Inches(0.08), Inches(1.4))
b7.fill.solid(); b7.fill.fore_color.rgb = ACCENT; b7.line.fill.background()
txt(sl, ep4, Inches(0.7), Inches(4.2), Inches(11.8), Inches(1.2), size=12, color=WHITE, wrap=True)

txt(sl, "そして──「医師を分かる業者」になれ。士業も医師も、結局そこを見ている。",
    Inches(0.45), Inches(5.65), Inches(12.2), Inches(0.42), size=15, bold=True, color=ACCENT)

# ════════════════════════════════════════════
# SLIDE 24  エンドスライド
# ════════════════════════════════════════════
sl = add_slide(); bg(sl)
box(sl, 0, 0, Inches(0.5), H, ACCENT)
box(sl, 0, H-Inches(1.1), W, Inches(1.1), BG_MID)
divider_line(sl, Inches(3.5))

txt(sl, "ご清聴ありがとうございました。",
    Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.6), size=20, color=LIGHT_GRAY)
txt(sl, "1人の医師に、\n30年寄り添う。",
    Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.8), size=46, bold=True, color=WHITE)
txt(sl, "これが、医療専門不動産のビジネスモデル。",
    Inches(0.9), Inches(3.65), Inches(11.5), Inches(0.55), size=20, color=LIGHT_GRAY)
txt(sl, "チームてっかん  |  菊池",
    Inches(0.9), Inches(5.1), Inches(7), Inches(0.5), size=20, bold=True, color=WHITE)
txt(sl, "ご相談・面談、お気軽にどうぞ。",
    Inches(0.9), Inches(5.65), Inches(7), Inches(0.42), size=14, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════
out = "/Users/kikuchikenta/01_honbu_docs_automation/medical_realestate_seminar_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")

"""
追客システムの使い方ガイド（KHD BtoC共通）
デザイン=クリーム白×レンガ赤（reference_slide_pipeline準拠）。
画面イメージはモックアップ（実スクショ代替）で再現：スプシ／Googleタスク／KPIダッシュボード。
出力: tsuikyaku_guide.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x4F)
FONT = "Hiragino Sans"
W, H = Inches(13.33), Inches(7.5)

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK, align=PP_ALIGN.LEFT,
      italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp: p.line_spacing = line_sp
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb

def bx(slide, x, y, w, h, col, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)

def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "追客システムの使い方  ｜  KHD 05物件調達・03事業運営  ｜  BtoC共通", Inches(0.5), H-Inches(0.42), Inches(11), Inches(0.32), sz=9, col=GRY)

def chip(slide, x, y, w, txt, fill=RED, fg=WHT, sz=11):
    bx(slide, x, y, w, Inches(0.34), fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t(slide, txt, x, y, w, Inches(0.34), sz=sz, bold=True, col=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def arrow(slide, x, y, w, h):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = RED; a.line.fill.background(); a.shadow.inherit = False
    return a

# ── 画面モック：スプレッドシート ──
def sheet_mock(slide, x, y, w, h):
    bx(slide, x, y, w, h, WHT, line=CARDLN, lw=1.0)
    bx(slide, x, y, w, Inches(0.06), RED)
    t(slide, "📋 ①追客リスト（Googleスプレッドシート）", x+Inches(0.22), y+Inches(0.16), w-Inches(0.4), Inches(0.32), sz=12, bold=True, col=RED)
    cols = ["顧客名","ステージ","温度","次アクション","期限","状態"]
    cw   = [Inches(1.55),Inches(1.55),Inches(0.95),Inches(3.0),Inches(0.95),Inches(1.05)]
    rows = [
        ["持倉様","⑦クロージング","HOT","小池先生に登記スケ確認","6/9","対応中"],
        ["曾我先生","④物件提案","WARM","福井下書きの送信判断","6/5","対応中"],
        ["内山先生","④物件提案","WARM","西新宿物件を紹介送付","6/5","対応中"],
    ]
    gx = x+Inches(0.22); gy = y+Inches(0.62); rh = Inches(0.48)
    # header
    cxx = gx
    for ci,c in enumerate(cols):
        hi = ci in (3,4)
        bx(slide, cxx, gy, cw[ci], rh, REDD if hi else RED)
        t(slide, c, cxx, gy, cw[ci], rh, sz=10.5, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cxx += cw[ci]
    # rows
    for ri,row in enumerate(rows):
        cxx = gx; ry = gy+rh*(ri+1)
        for ci,val in enumerate(row):
            hi = ci in (3,4)
            fill = REDBG if hi else (CARD if ri%2==0 else BG)
            bx(slide, cxx, ry, cw[ci], rh, fill, line=LINE, lw=0.5)
            col = RED if (ci==2 and val=="HOT") else INK
            t(slide, val, cxx+Inches(0.04), ry, cw[ci]-Inches(0.06), rh, sz=9.5,
              bold=(ci in(0,2)), col=col, align=(PP_ALIGN.LEFT if ci==3 else PP_ALIGN.CENTER), anchor=MSO_ANCHOR.MIDDLE)
            cxx += cw[ci]
    # callout
    cy2 = gy+rh*(len(rows)+1)+Inches(0.12)
    bx(slide, gx, cy2, w-Inches(0.44), Inches(0.5), GRYBG, line=RED, lw=1.2)
    t(slide, "▲ この「次アクション＋期限」を書くと、10分以内にタスクが自動で湧く（＝強制力）",
      gx+Inches(0.15), cy2, w-Inches(0.7), Inches(0.5), sz=10.5, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)

# ── 画面モック：Googleタスク（スマホ風）──
def task_mock(slide, x, y, w, h):
    bx(slide, x, y, w, h, WHT, line=CARDLN, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bx(slide, x, y, w, Inches(0.5), RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t(slide, "✅ Googleタスク（自動生成）", x, y+Inches(0.05), w, Inches(0.4), sz=12, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tasks = [
        ("【追客】持倉様 ⑦クロージング｜小池先生に登記スケ確認", "6/9"),
        ("【追客】曾我先生 ④物件提案｜福井下書きの送信判断", "6/5"),
        ("【追客】内山先生 ④物件提案｜西新宿物件を紹介送付", "6/5"),
    ]
    ty = y+Inches(0.7)
    for txt, due in tasks:
        bx(slide, x+Inches(0.2), ty, Inches(0.26), Inches(0.26), None, line=GRY, lw=1.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        t(slide, txt, x+Inches(0.58), ty-Inches(0.04), w-Inches(1.6), Inches(0.5), sz=9.5, col=INK, line_sp=0.95)
        chip(slide, x+w-Inches(0.95), ty-Inches(0.02), Inches(0.75), due, fill=REDBG, fg=REDD, sz=9)
        bx(slide, x+Inches(0.2), ty+Inches(0.52), w-Inches(0.4), Pt(0.8), LINE)
        ty += Inches(0.66)

# ── 画面モック：KPIダッシュボード ──
def kpi_mock(slide, x, y, w, h):
    bx(slide, x, y, w, h, WHT, line=CARDLN, lw=1.0)
    bx(slide, x, y, w, Inches(0.06), RED)
    t(slide, "📊 KPIダッシュボード（毎時 自動更新）", x+Inches(0.25), y+Inches(0.16), w-Inches(0.4), Inches(0.35), sz=12, bold=True, col=RED)
    tiles = [("今週の追客数","8","件"),("アクティブ追客","5","件"),("SLA期限切れ","1","件"),("成約率","50","%")]
    tw = (w-Inches(0.5)-Inches(0.45))/4; tx = x+Inches(0.25); tyy = y+Inches(0.65)
    for i,(lab,val,unit) in enumerate(tiles):
        cx = tx+(tw+Inches(0.15))*i
        red_tile = (i==2)
        bx(slide, cx, tyy, tw, Inches(1.15), REDBG if red_tile else CARD, line=(RED if red_tile else CARDLN), lw=(1.4 if red_tile else 1.0))
        t(slide, lab, cx, tyy+Inches(0.12), tw, Inches(0.3), sz=9.5, col=(REDD if red_tile else GRY), align=PP_ALIGN.CENTER)
        t(slide, val, cx, tyy+Inches(0.38), tw, Inches(0.6), sz=30, bold=True, col=RED, align=PP_ALIGN.CENTER)
        t(slide, unit, cx, tyy+Inches(0.92), tw, Inches(0.22), sz=9, col=GRY, align=PP_ALIGN.CENTER)
    # ステージ別バー
    t(slide, "ステージ別 件数", x+Inches(0.25), tyy+Inches(1.4), w-Inches(0.5), Inches(0.3), sz=10.5, bold=True, col=INK)
    bars = [("④提案",3),("⑤アポ",1),("⑥内見",1),("⑦クロ",1),("⑧成約",0)]
    bxx0 = x+Inches(0.4); byy = tyy+Inches(2.95); bw = Inches(0.85); maxh = Inches(1.05); gap=Inches(0.55)
    mx = max(b[1] for b in bars) or 1
    for i,(lab,v) in enumerate(bars):
        bh = Emu(int(maxh*(v/mx))) if v>0 else Emu(int(Inches(0.03)))
        cx = bxx0+(bw+gap)*i
        bx(slide, cx, byy-bh, bw, bh, RED if i==3 else RGBColor(0xC9,0x6A,0x62))
        t(slide, str(v), cx-Inches(0.1), byy-bh-Inches(0.26), bw+Inches(0.2), Inches(0.24), sz=10, bold=True, col=INK, align=PP_ALIGN.CENTER)
        t(slide, lab, cx-Inches(0.15), byy+Inches(0.04), bw+Inches(0.3), Inches(0.24), sz=9, col=GRY, align=PP_ALIGN.CENTER)

# ════════ S1 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "FOLLOW-UP ENGINE ｜ 追客システム", Inches(0.9), Inches(1.5), Inches(7.5), Inches(0.45), sz=15, bold=True, col=RED)
t(s, "追客システムの使い方", Inches(0.88), Inches(2.05), Inches(7.6), Inches(0.9), sz=40, bold=True, col=INK)
t(s, "スプシに書くだけ。あとは全自動。", Inches(0.88), Inches(2.95), Inches(7.6), Inches(0.8), sz=30, bold=True, col=RED)
t(s, "仲介トップセールス（羽鳥師匠・野仲先輩）のトークをTTPした、\n事業横断で使える追客の型。タスク化とKPI集計はGASが裏で回す。",
  Inches(0.9), Inches(3.95), Inches(7.4), Inches(0.9), sz=13.5, col=GRY, line_sp=1.25)
# 右：3チップの流れ
ox, oy = Inches(8.7), Inches(2.0)
flow = [("①","スプシに1行書く"),("②","タスクが自動で湧く"),("③","完了でKPIが立つ")]
for i,(no,txt) in enumerate(flow):
    cy = oy+Inches(1.35)*i
    bx(s, ox, cy, Inches(3.9), Inches(1.05), CARD, line=CARDLN, lw=1.0)
    bx(s, ox, cy, Inches(0.06), Inches(1.05), RED)
    t(s, no, ox+Inches(0.25), cy+Inches(0.18), Inches(0.7), Inches(0.7), sz=30, bold=True, col=RED)
    t(s, txt, ox+Inches(1.0), cy+Inches(0.3), Inches(2.8), Inches(0.5), sz=15, bold=True, col=INK, anchor=MSO_ANCHOR.MIDDLE)
    if i<2: bx(s, ox+Inches(1.7), cy+Inches(1.05), Inches(0.5), Inches(0.3), REDBG, shape=MSO_SHAPE.DOWN_ARROW)
bx(s, Inches(0.9), Inches(6.78), Inches(11.5), Pt(1.2), LINE)
t(s, "KHD ｜ 菊池 研太 ｜ 05物件調達・03事業運営", Inches(0.9), Inches(6.9), Inches(11), Inches(0.4), sz=13, bold=True, col=INK)

# ════════ S2 全体像 ════════
s = sl(); ft(s)
hdr(s, "THE FLOW", "全体像｜あなたは「書くだけ」、あとはGASが回す", "PC不要・Google側で自動稼働（gas_tsuikyaku.gs：10分ごと＋毎時）")
cols = [
    ("あなた", "スプシに1行＋次アクション", ["新しい客を1行追加","次アクションと期限を書く","追客したらタスクを完了"], CARD),
    ("GAS（自動）", "Googleが裏で処理", ["行→タスク自動生成","完了タスクをKPI集計","成約/失注でタスク自動クローズ"], REDBG),
    ("結果（自動）", "見える化される", ["タスクが湧いて忘れない","KPIダッシュボード更新","履歴・最終接触日が残る"], CARD),
]
cw, gx, x0, y0 = Inches(3.85), Inches(0.55), Inches(0.55), Inches(2.1)
for i,(ti,sub,items,fill) in enumerate(cols):
    cx = x0+(cw+gx)*i
    bx(s, cx, y0, cw, Inches(3.9), fill, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.7), RED)
    t(s, ti, cx, y0+Inches(0.1), cw, Inches(0.5), sz=17, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t(s, sub, cx, y0+Inches(0.82), cw, Inches(0.4), sz=12, bold=True, col=RED, align=PP_ALIGN.CENTER)
    for j,it in enumerate(items):
        t(s, "● "+it, cx+Inches(0.3), y0+Inches(1.45)+Inches(0.62)*j, cw-Inches(0.5), Inches(0.55), sz=12, col=INK, line_sp=1.0)
    if i<2: arrow(s, cx+cw+Inches(0.06), y0+Inches(1.6), Inches(0.42), Inches(0.6))

# ════════ S3 日々の3ステップ ════════
s = sl(); ft(s)
hdr(s, "DAILY", "毎日やることは、たった3つ", "この3つだけ。残りは全部自動で回る")
steps = [
    ("STEP 1","客を1行足す","反響が来たら","顧客名・ステージ・温度を入れる。\nこれがKPIの母数になる。", "→ DBに登録"),
    ("STEP 2","次アクション＋期限を書く","動かす一手を決める","「誰に・何を・いつまで」を書く。\n10分以内にタスクが自動で湧く。", "→ タスク自動生成"),
    ("STEP 3","タスクを完了にする","追客したら","LINEでも電話でもメールでもOK。\nやったらタスクを閉じるだけ。", "→ KPI＋1・履歴自動"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(2.0)
for i,(st,ti,sub,body,res) in enumerate(steps):
    cx = x0+(cw+gx)*i
    bx(s, cx, y0, cw, Inches(3.7), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.7), RED)
    t(s, st, cx, y0+Inches(0.1), cw, Inches(0.5), sz=18, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t(s, ti, cx+Inches(0.2), y0+Inches(0.85), cw-Inches(0.4), Inches(0.6), sz=16, bold=True, col=INK, align=PP_ALIGN.CENTER)
    t(s, sub, cx, y0+Inches(1.45), cw, Inches(0.3), sz=11, bold=True, col=RED, align=PP_ALIGN.CENTER)
    t(s, body, cx+Inches(0.3), y0+Inches(1.85), cw-Inches(0.6), Inches(1.2), sz=11.5, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.15)
    bx(s, cx+Inches(0.3), y0+Inches(3.05), cw-Inches(0.6), Inches(0.45), REDBG)
    t(s, res, cx+Inches(0.3), y0+Inches(3.05), cw-Inches(0.6), Inches(0.45), sz=11, bold=True, col=REDD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ════════ S4 画面① スプシ ════════
s = sl(); ft(s)
hdr(s, "SCREEN 1", "画面イメージ①｜追客リスト（スプレッドシート）", "ここが司令塔。客を入れて、次アクションを書く。※画面はイメージ")
sheet_mock(s, Inches(0.55), Inches(1.95), Inches(12.23), Inches(4.6))

# ════════ S5 画面② タスク ════════
s = sl(); ft(s)
hdr(s, "SCREEN 2", "画面イメージ②｜自動で湧くタスク → 完了でKPI", "スプシに書いた次アクションが、10分以内にタスク化される。※画面はイメージ")
task_mock(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(4.4))
# 右：完了したら何が起きるか
rx = Inches(7.0)
bx(s, rx, Inches(2.0), Inches(5.78), Inches(4.4), CARD, line=CARDLN, lw=1.0)
bx(s, rx, Inches(2.0), Inches(5.78), Inches(0.06), RED)
t(s, "タスクを「完了」にすると、自動で…", rx+Inches(0.3), Inches(2.25), Inches(5.2), Inches(0.4), sz=14, bold=True, col=RED)
effs = [
    ("KPI ＋1","今週の追客数（完了タスク）がカウントされる"),
    ("履歴が残る","メモ欄に「[6/9済]登記スケ確認」と自動追記"),
    ("最終接触日 更新","いつ動いたかが自動で入る"),
    ("行が解放","次アクションを書けば、また次のタスクが湧く"),
]
for i,(ti,ds) in enumerate(effs):
    cy = Inches(2.85)+Inches(0.82)*i
    bx(s, rx+Inches(0.3), cy, Inches(5.18), Inches(0.68), WHT, line=CARDLN, lw=0.75)
    bx(s, rx+Inches(0.3), cy, Inches(0.08), Inches(0.68), RED)
    t(s, "✓ "+ti, rx+Inches(0.5), cy+Inches(0.06), Inches(2.2), Inches(0.55), sz=12.5, bold=True, col=REDD, anchor=MSO_ANCHOR.MIDDLE)
    t(s, ds, rx+Inches(2.5), cy+Inches(0.06), Inches(2.9), Inches(0.55), sz=10, col=GRY, anchor=MSO_ANCHOR.MIDDLE, line_sp=0.95)

# ════════ S6 画面③ ダッシュボード ════════
s = sl(); ft(s)
hdr(s, "SCREEN 3", "画面イメージ③｜KPIダッシュボード（毎時自動更新）", "今週の追客数＝完了タスク件数（LINE/電話/メール どれでもOK）。※数値は例示")
kpi_mock(s, Inches(2.4), Inches(1.95), Inches(8.5), Inches(4.7))

# ════════ S7 全自動 vs 人 ════════
s = sl(); ft(s)
hdr(s, "WHO DOES WHAT", "何が全自動で、何があなたの仕事か", "境界をはっきりさせる。あなたの手は最小限")
# 左：全自動
bx(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(4.3), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(0.6), RED)
t(s, "🤖 全自動（もう触らない）", Inches(0.55), Inches(2.05), Inches(6.0), Inches(0.5), sz=15, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i,it in enumerate(["タスク生成（行→タスク）","KPIダッシュボード更新","完了タスクの集計＝追客数","成約/失注でタスククローズ","履歴・最終接触日の記録"]):
    t(s, "● "+it, Inches(0.9), Inches(2.85)+Inches(0.62)*i, Inches(5.4), Inches(0.55), sz=13, col=INK)
# 右：人
bx(s, Inches(6.78), Inches(2.0), Inches(6.0), Inches(4.3), GRYBG, line=CARDLN, lw=1.0)
bx(s, Inches(6.78), Inches(2.0), Inches(6.0), Inches(0.6), REDD)
t(s, "🙋 あなた（3つだけ）", Inches(6.78), Inches(2.05), Inches(6.0), Inches(0.5), sz=15, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i,it in enumerate(["新しい客を1行足す","次アクション＋期限を書く","追客したらタスクを完了にする"]):
    cy = Inches(2.95)+Inches(0.9)*i
    bx(s, Inches(7.1), cy, Inches(5.35), Inches(0.7), WHT, line=CARDLN, lw=0.75)
    bx(s, Inches(7.1), cy, Inches(0.08), Inches(0.7), RED)
    t(s, str(i+1)+". "+it, Inches(7.35), cy, Inches(5.0), Inches(0.7), sz=14, bold=True, col=INK, anchor=MSO_ANCHOR.MIDDLE)
t(s, "※ 文面まで残したくなったら「Gmail軸」を後付け可能", Inches(6.78), Inches(5.95), Inches(6.0), Inches(0.35), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)

# ════════ S8 はじめ方＋締め ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "GET STARTED ｜ はじめ方", Inches(0.9), Inches(0.7), Inches(11), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "5分・1回だけのセットアップで、ずっと自動", Inches(0.9), Inches(1.15), Inches(11.7), Inches(0.7), sz=26, bold=True, col=INK)
setup = [
    ("1","スプシ → 拡張機能 → Apps Script を開く"),
    ("2","Drive「05_追客GAS_貼付用.gs.txt」を全文コピペ"),
    ("3","サービス＋から「Tasks API」を追加"),
    ("4","関数 setupTriggers を実行して承認"),
]
for i,(no,txt) in enumerate(setup):
    cy = Inches(2.1)+Inches(0.85)*i
    bx(s, Inches(0.9), cy, Inches(7.2), Inches(0.7), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.9), cy, Inches(0.7), Inches(0.7), RED)
    t(s, no, Inches(0.9), cy, Inches(0.7), Inches(0.7), sz=22, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t(s, txt, Inches(1.8), cy, Inches(6.2), Inches(0.7), sz=13.5, bold=True, col=INK, anchor=MSO_ANCHOR.MIDDLE)
# 右：締めメッセージ
bx(s, Inches(8.5), Inches(2.1), Inches(3.95), Inches(3.6), RED)
t(s, "あなたの仕事は", Inches(8.5), Inches(2.6), Inches(3.95), Inches(0.5), sz=15, bold=True, col=WHT, align=PP_ALIGN.CENTER)
t(s, "書いて、閉じる", Inches(8.5), Inches(3.2), Inches(3.95), Inches(0.8), sz=30, bold=True, col=WHT, align=PP_ALIGN.CENTER)
t(s, "それだけ。", Inches(8.5), Inches(4.05), Inches(3.95), Inches(0.6), sz=22, bold=True, col=RGBColor(0xF2,0xD8,0xD6), align=PP_ALIGN.CENTER)
t(s, "追客の数が、毎週\n自動でKPIに積み上がる", Inches(8.5), Inches(4.75), Inches(3.95), Inches(0.8), sz=12.5, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.2)
bx(s, Inches(0.9), Inches(6.5), Inches(11.5), Pt(1.2), LINE)
t(s, "詰まったらエラーを貼ってください。伴走します。", Inches(0.9), Inches(6.62), Inches(11), Inches(0.4), sz=12, bold=True, col=GRY)

prs.save("tsuikyaku_guide.pptx")
print("saved tsuikyaku_guide.pptx  /  slides:", len(prs.slides._sldIdLst))

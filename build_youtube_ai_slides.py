"""
YouTube「KHD × Claude Code 活用ガイド」 画面サポート用スライド
build_slides_minimal.py のテーマ（ダーク×ゴールド）を流用。
出力: youtube_ai_guide_2026.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG  = RGBColor(0x0D, 0x1B, 0x2A)
ACC = RGBColor(0xE8, 0xA8, 0x00)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
LGR = RGBColor(0xCC, 0xD6, 0xE0)
MID = RGBColor(0x1A, 0x2E, 0x44)
BLU = RGBColor(0x4A, 0x9E, 0xCB)
GRN = RGBColor(0x4C, 0xC2, 0x8C)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=20, bold=False, col=WHT,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
        r.font.italic = italic
    return tb


def bx(slide, x, y, w, h, col, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = col
    if lw:
        s.line.color.rgb = col
    else:
        s.line.fill.background()
    return s


def hdr(slide, main, sub=""):
    bx(slide, 0, 0, W, Inches(1.1), MID)
    t(slide, main, Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), sz=24, bold=True)
    if sub:
        t(slide, sub, Inches(0.6), Inches(0.78), Inches(12), Inches(0.3), sz=11, col=LGR)


def ft(slide):
    bx(slide, 0, H-Inches(0.38), W, Inches(0.38), MID)
    t(slide, "KHD × Claude Code  活用ガイド",
      Inches(0.5), H-Inches(0.35), Inches(10), Inches(0.32), sz=9, col=LGR)


def cmd_box(slide, cmd, x, y, w):
    """コマンドを目立つチップで表示"""
    bx(slide, x, y, w, Inches(0.6), RGBColor(0x10, 0x24, 0x38))
    bx(slide, x, y, Inches(0.08), Inches(0.6), ACC)
    t(slide, cmd, x+Inches(0.2), y+Inches(0.1), w-Inches(0.3), Inches(0.4),
      sz=18, bold=True, col=ACC)


# ── 01 カバー ──────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "AIアシスタントでKHD業務を丸ごと自動化",
  Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.55), sz=17, col=ACC, italic=True)
t(s, "KHD × Claude Code\n完全活用ガイド",
  Inches(0.9), Inches(1.6), Inches(11.5), Inches(2.0), sz=44, bold=True)
t(s, "物件評価・業者メール・朝のタスク整理を、コマンド1つで。特別なITスキルは不要。",
  Inches(0.9), Inches(3.7), Inches(11), Inches(0.55), sz=15, col=LGR)
bx(s, Inches(0.9), Inches(4.5), Inches(10), Inches(0.04), ACC)
t(s, "菊池ホールディングス（KHD）  |  2026年",
  Inches(0.9), Inches(4.65), Inches(10), Inches(0.45), sz=14, col=LGR)

# ── 02 この動画でできること（Before→After）──
s = sl(); ft(s)
hdr(s, "この動画でできるようになること")
bx(s, Inches(0.6), Inches(1.45), Inches(5.7), Inches(4.4), MID)
t(s, "BEFORE  /  今", Inches(0.85), Inches(1.6), Inches(5.2), Inches(0.4), sz=13, col=LGR)
before = ["・PDFが届くたびに手で評価", "・KPI計算を電卓で", "・メールを毎回ゼロから作文",
          "・朝、何からやるか毎回迷う", "・銀行選びは勘"]
for i, x in enumerate(before):
    t(s, x, Inches(0.85), Inches(2.2)+i*Inches(0.62), Inches(5.2), Inches(0.55),
      sz=14, col=RGBColor(0xFF, 0x88, 0x88))
bx(s, Inches(6.9), Inches(1.45), Inches(6.0), Inches(4.4), MID)
bx(s, Inches(6.9), Inches(1.45), Inches(6.0), Inches(0.06), ACC)
t(s, "AFTER  /  Claude Code導入後", Inches(7.15), Inches(1.6), Inches(5.5), Inches(0.4),
  sz=13, bold=True, col=ACC)
after = ["・物件PDF → 30秒で評価レポート", "・CF率/CCR/土地値を自動算出",
         "・業者メールはコマンド1つで下書き", "・毎朝 /wbs-morning でタスク一覧",
         "・融資先も判断軸つきで提案"]
for i, x in enumerate(after):
    t(s, x, Inches(7.15), Inches(2.2)+i*Inches(0.62), Inches(5.5), Inches(0.55),
      sz=14, col=WHT)

# ── 03 起動方法 ────────────────────────────
s = sl(); ft(s)
hdr(s, "STEP 0  /  起動方法", "やることは1つだけ")
t(s, "デスクトップの『Claude』を開くだけ。",
  Inches(0.7), Inches(1.5), Inches(12), Inches(0.6), sz=26, bold=True, col=ACC)
pts = ["ターミナルも、難しいコマンドも一切不要",
       "開いた瞬間に『01_honbu_docs_automation』が自動ロード",
       "KHDの判断基準・不動産KPI・業者情報がすべてAIに読み込み済み",
       "特別なセットアップは不要。開くだけ。"]
for i, x in enumerate(pts):
    bx(s, Inches(0.7), Inches(2.55)+i*Inches(0.85), Inches(0.14), Inches(0.6), ACC)
    t(s, x, Inches(1.0), Inches(2.6)+i*Inches(0.85), Inches(11.5), Inches(0.6), sz=16)

# ── 04 物件評価 3ステップ ──────────────────
s = sl(); ft(s)
hdr(s, "1  /  物件評価の使い方", "物件PDFが届いたら、この3ステップ")
steps = [("STEP 1", "PDFのパスをコピー", "Finderで右クリック → 『パスのコピー』"),
         ("STEP 2", "チャットに入力", "/property-intake  と打って\nコピーしたパスを貼り付け"),
         ("STEP 3", "Enterで待つ", "約30秒で評価レポートが出力される")]
cw = Inches(4.0)
for i, (no, ti, di) in enumerate(steps):
    cx = Inches(0.55)+i*(cw+Inches(0.14))
    bx(s, cx, Inches(1.55), cw, Inches(3.9), MID)
    bx(s, cx, Inches(1.55), cw, Inches(0.06), ACC)
    t(s, no, cx+Inches(0.2), Inches(1.75), cw-Inches(0.3), Inches(0.5), sz=15, bold=True, col=ACC)
    t(s, ti, cx+Inches(0.2), Inches(2.35), cw-Inches(0.3), Inches(0.9), sz=18, bold=True)
    t(s, di, cx+Inches(0.2), Inches(3.4), cw-Inches(0.3), Inches(1.8), sz=13, col=LGR)
cmd_box(s, "/property-intake  <PDFのパス>", Inches(0.55), Inches(5.7), Inches(12.2))

# ── 05 KPIの見方 ───────────────────────────
s = sl(); ft(s)
hdr(s, "1  /  評価レポートの見方", "3つの数字を見れば、買いか見送りか分かる")
kpis = [("CF率", "1.5%以上が買い", "毎月の手残り率", BLU),
        ("CCR", "15%以上が合格", "自己資金に対するリターン", GRN),
        ("土地値割合", "0.4以上が安全圏", "担保価値の指標", ACC)]
cw = Inches(4.0)
for i, (k, base, desc, col) in enumerate(kpis):
    cx = Inches(0.55)+i*(cw+Inches(0.14))
    bx(s, cx, Inches(1.55), cw, Inches(3.0), MID)
    bx(s, cx, Inches(1.55), cw, Inches(0.06), col)
    t(s, k, cx+Inches(0.2), Inches(1.8), cw-Inches(0.3), Inches(0.6), sz=22, bold=True, col=col)
    t(s, base, cx+Inches(0.2), Inches(2.6), cw-Inches(0.3), Inches(0.6), sz=16, bold=True)
    t(s, desc, cx+Inches(0.2), Inches(3.3), cw-Inches(0.3), Inches(1.0), sz=13, col=LGR)
t(s, "3つとも未達なら自動で『見送り』判定 → 『DBに追加？』にYESでスプレッドシートに自動記録。",
  Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.9), sz=15, col=ACC)

# ── 06 毎朝 ────────────────────────────────
s = sl(); ft(s)
hdr(s, "2  /  毎朝の使い方")
cmd_box(s, "/wbs-morning", Inches(0.7), Inches(1.5), Inches(7.0))
pts = ["全社WBSから未完了タスクを自動で引っ張ってくる",
       "今日やるべきことをリストアップ",
       "毎朝5:57に昨日の学びをAIが自動でメモリ記録（何もしなくてOK）"]
for i, x in enumerate(pts):
    bx(s, Inches(0.7), Inches(2.6)+i*Inches(0.85), Inches(0.14), Inches(0.6), ACC)
    t(s, x, Inches(1.0), Inches(2.65)+i*Inches(0.85), Inches(11.5), Inches(0.6), sz=16)

# ── 07 業者メール ──────────────────────────
s = sl(); ft(s)
hdr(s, "3  /  業者・銀行へのメール下書き")
cmd_box(s, "/vendor-mail  相手先と用件を伝えるだけ", Inches(0.7), Inches(1.5), Inches(11.9))
t(s, "例）/vendor-mail  〇〇社に△△案件をクローズしたい旨を伝えたい",
  Inches(0.7), Inches(2.45), Inches(12), Inches(0.5), sz=15, col=LGR, italic=True)
pts = ["KHDのスタイルに合ったメール文がすぐ出てくる",
       "あとはコピーしてGmailに貼るだけ"]
for i, x in enumerate(pts):
    bx(s, Inches(0.7), Inches(3.3)+i*Inches(0.85), Inches(0.14), Inches(0.6), ACC)
    t(s, x, Inches(1.0), Inches(3.35)+i*Inches(0.85), Inches(11.5), Inches(0.6), sz=16)

# ── 08 融資先 ──────────────────────────────
s = sl(); ft(s)
hdr(s, "4  /  融資先の選び方")
cmd_box(s, "/bank-approach  物件の場所と金額を伝える", Inches(0.7), Inches(1.5), Inches(11.9))
pts = ["どの銀行に打診すべきか、優先順位つきで提案",
       "判断軸はすでに組み込み済み"]
for i, x in enumerate(pts):
    bx(s, Inches(0.7), Inches(2.6)+i*Inches(0.85), Inches(0.14), Inches(0.6), ACC)
    t(s, x, Inches(1.0), Inches(2.65)+i*Inches(0.85), Inches(11.5), Inches(0.6), sz=16)
t(s, "（地域別の融資先判断軸が内蔵されています）",
  Inches(0.7), Inches(4.5), Inches(12), Inches(0.5), sz=13, col=LGR, italic=True)

# ── 09 まとめ + CTA ────────────────────────
s = sl(); ft(s)
hdr(s, "まとめ  /  難しいことは何もない")
items = ["1.  アプリを開くだけでKHDの情報が自動ロード",
         "2.  物件PDF → /property-intake → 30秒で評価",
         "3.  毎朝 /wbs-morning でタスク整理",
         "4.  メール下書きは /vendor-mail",
         "5.  融資先選定は /bank-approach"]
for i, x in enumerate(items):
    t(s, x, Inches(0.8), Inches(1.5)+i*Inches(0.7), Inches(11.5), Inches(0.6), sz=18, bold=True)
bx(s, Inches(0.7), Inches(5.2), Inches(12.0), Inches(1.1), MID)
bx(s, Inches(0.7), Inches(5.2), Inches(0.08), Inches(1.1), ACC)
t(s, "コマンドを打てば、AIが全部やります。",
  Inches(0.95), Inches(5.45), Inches(11.6), Inches(0.7), sz=20, bold=True, col=ACC)

prs.save("youtube_ai_guide_2026.pptx")
print("✅ youtube_ai_guide_2026.pptx を生成しました（全9枚）")

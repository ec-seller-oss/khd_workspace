"""ゆーしMac セットアップ・マニュアル(本人がフル設定) KHDクリーム白×レンガ赤。
出力: KHD_ゆーしMac_セットアップマニュアル.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); DARK=RGBColor(0x2C,0x2C,0x2A); MONOC=RGBColor(0xF1,0xEC,0xE1)
WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"; MONO="Menlo"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None,fontname=FONT):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold
        r.font.color.rgb=col; r.font.name=fontname
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)

def ft(slide,p):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"ゆーしMac セットアップ・マニュアル  ｜  KHD",Inches(0.5),H-Inches(0.42),Inches(9),Inches(0.32),sz=9,col=GRY)
    t(slide,p,Inches(11.3),H-Inches(0.42),Inches(1.5),Inches(0.32),sz=9,col=GRY,align=PP_ALIGN.RIGHT)

def code(slide,x,y,w,lines,sz=14):
    h=Inches(0.28)+Inches(0.32)*len(lines)
    bx(slide,x,y,w,h,DARK)
    yy=y+Inches(0.13)
    for ln in lines:
        t(slide,ln,x+Inches(0.2),yy,w-Inches(0.35),Inches(0.32),sz=sz,col=MONOC,fontname=MONO); yy=yy+Inches(0.32)
    return h

def num(slide,x,y,n,r=0.5):
    c=slide.shapes.add_shape(MSO_SHAPE.OVAL,x,y,Inches(r),Inches(r))
    c.fill.solid(); c.fill.fore_color.rgb=RED; c.line.fill.background(); c.shadow.inherit=False
    t(slide,str(n),x,y,Inches(r),Inches(r),sz=int(r*44),bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def star(slide,x,y,w,text):
    bx(slide,x,y,w,Inches(0.5),REDBG); bx(slide,x,y,Inches(0.08),Inches(0.5),RED)
    t(slide,"★ "+text,x+Inches(0.2),y,w-Inches(0.3),Inches(0.5),sz=12,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ── S1 表紙 ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"YUUSHI MAC SETUP MANUAL",Inches(0.9),Inches(1.5),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"ゆーしMac セットアップ・マニュアル",Inches(0.88),Inches(2.1),Inches(11.7),Inches(0.9),sz=38,bold=True,col=INK)
t(s,"自分のMacで Claude を動かすまで（インストール完全版）",Inches(0.9),Inches(3.0),Inches(11.6),Inches(0.5),sz=18,bold=True,col=RED)
t(s,"アプリ名・メニュー名・コマンドまで、この通りに進めればOK。\n所要 約60分。つまずいたら、その画面をスクショして菊池さんにLINE。",Inches(0.9),Inches(3.7),Inches(11.6),Inches(0.9),sz=14,col=GRY,line_sp=1.3)
bx(s,Inches(0.9),Inches(5.0),Inches(11.5),Inches(1.2),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(5.0),Inches(11.5),Inches(0.06),RED)
t(s,"用意するもの",Inches(1.1),Inches(5.15),Inches(11),Inches(0.35),sz=13,bold=True,col=RED)
t(s,"・Mac（ゆーし）　・ネット環境　・菊池さんから受け取る“ログイン情報・フォルダのパス”（★の所で使う）",Inches(1.1),Inches(5.5),Inches(11),Inches(0.6),sz=13,col=INK,line_sp=1.2)
bx(s,Inches(0.9),Inches(6.5),Inches(11.5),Pt(1.2),LINE)
t(s,"KHD  ｜  菊池 研太  ｜  2026-06-26",Inches(0.9),Inches(6.6),Inches(11),Inches(0.4),sz=12,bold=True,col=INK)

# ── S2 全体の流れ ──
s=sl(); ft(s,"2 / 9")
hdr(s,"OVERVIEW","全体の流れ ── 7ステップ","上から順にやればOK。★は菊池さんと連携する所")
flow=[("1","ターミナルを開く"),("2","開発ツール(git)を入れる"),("3","Node.js を入れる"),
 ("4","Claude Code を入れて ログイン ★"),("5","Google Drive を入れる ★"),
 ("6","データを取り込む(clone) ★"),("7","起動して /company で確認")]
x0,y0=Inches(0.7),Inches(2.0); cw=Inches(5.9); gx=Inches(0.6); rh=Inches(0.92)
for i,(n,lab) in enumerate(flow):
    cx=x0+(cw+gx)*(i//4); cy=y0+rh*(i%4)
    num(s,cx,cy,n,r=0.56)
    t(s,lab,cx+Inches(0.75),cy,cw-Inches(0.8),Inches(0.56),sz=16,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
star(s,Inches(6.6),Inches(5.7),Inches(6.2),"★3箇所は、菊池さんから情報をもらってから進む")

# ── S3 事前準備 ──
s=sl(); ft(s,"3 / 9")
hdr(s,"BEFORE YOU START","事前に 菊池さんから受け取るもの（★）","この4つが揃ってから始めるとスムーズ")
items=[
 ("① Claudeのログイン情報","claudeに入る時に使う（菊池さんの案内に従う）"),
 ("② Google Driveの入り方","どのアカウントで入るか／フォルダの共有"),
 ("③ 取り込むフォルダの正確なパス","STEP6のgit cloneでそのまま貼る1行"),
 ("④ つなぐ連携(MCP)はどれか","必要なものだけ。菊池さんが指定"),
]
y=2.05
for i,(ti,de) in enumerate(items):
    cy=Inches(y)+Inches(1.08)*i
    bx(s,Inches(0.6),cy,Inches(12.1),Inches(0.92),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.6),cy,Inches(0.1),Inches(0.92),RED)
    t(s,ti,Inches(0.85),cy+Inches(0.12),Inches(6.0),Inches(0.6),sz=16,bold=True,col=INK)
    t(s,de,Inches(6.7),cy+Inches(0.12),Inches(5.8),Inches(0.7),sz=12.5,col=GRY,anchor=MSO_ANCHOR.MIDDLE)

# ── S4 STEP1-2 ──
s=sl(); ft(s,"4 / 9")
hdr(s,"STEP 1-2","ターミナルを開く → 開発ツール(git)","“ターミナル”は文字でMacに命令する黒い画面。こわくない")
num(s,Inches(0.65),Inches(1.95),1,r=0.5); t(s,"ターミナルを開く",Inches(1.3),Inches(1.95),Inches(11),Inches(0.5),sz=16,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
t(s,"Finder →「アプリケーション」→「ユーティリティ」→「ターミナル」をダブルクリック\n（早い方法：⌘+スペース →「ターミナル」と打ってEnter）",Inches(1.3),Inches(2.5),Inches(11.3),Inches(0.8),sz=13,col=GRY,line_sp=1.25)
num(s,Inches(0.65),Inches(3.6),2,r=0.5); t(s,"開発ツール(git)を入れる",Inches(1.3),Inches(3.6),Inches(11),Inches(0.5),sz=16,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
t(s,"ターミナルに下を貼ってEnter → 出た窓で「インストール」をクリック → 終わるまで待つ",Inches(1.3),Inches(4.15),Inches(11.3),Inches(0.4),sz=13,col=GRY)
code(s,Inches(1.3),Inches(4.6),Inches(8.5),["xcode-select --install"],sz=15)
star(s,Inches(0.6),Inches(6.0),Inches(12.1),"うまくいかない時は、ターミナルの画面をスクショして菊池さんにLINE")

# ── S5 STEP3 ──
s=sl(); ft(s,"5 / 9")
hdr(s,"STEP 3","Node.js を入れる","Claudeを動かすための土台ソフト")
steps=[
 "ブラウザで  nodejs.org  を開く",
 "「LTS」と書かれた macOS Installer (.pkg) をダウンロード",
 "ダウンロードした .pkg をダブルクリック →「続ける」→ 同意 →「インストール」",
 "Macのパスワードを入力（入れたら完了）",
]
y=2.0
for i,stp in enumerate(steps):
    cy=Inches(y)+Inches(0.62)*i
    num(s,Inches(0.65),cy,i+1,r=0.46)
    t(s,stp,Inches(1.25),cy,Inches(11.4),Inches(0.46),sz=14.5,col=INK,anchor=MSO_ANCHOR.MIDDLE)
t(s,"入ったか確認：ターミナルで下を打つ → 「v22…」のように数字が出ればOK",Inches(0.65),Inches(4.7),Inches(11.5),Inches(0.4),sz=13,col=GRY)
code(s,Inches(0.65),Inches(5.15),Inches(6.0),["node -v"],sz=15)

# ── S6 STEP4 ──
s=sl(); ft(s,"6 / 9")
hdr(s,"STEP 4","Claude Code を入れて ログイン ★","ここが本体。コマンドは1行ずつコピペでOK")
t(s,"① インストール（ターミナルに貼ってEnter・数分待つ）",Inches(0.65),Inches(1.95),Inches(11.5),Inches(0.4),sz=14,bold=True,col=INK)
code(s,Inches(0.65),Inches(2.4),Inches(11.0),["npm install -g @anthropic-ai/claude-code"],sz=15)
t(s,"② 入ったか確認（バージョン番号が出ればOK）",Inches(0.65),Inches(3.2),Inches(11.5),Inches(0.4),sz=14,bold=True,col=INK)
code(s,Inches(0.65),Inches(3.65),Inches(6.0),["claude --version"],sz=15)
t(s,"③ ログイン（下を打つ→ブラウザが開く→案内のアカウントでログイン）",Inches(0.65),Inches(4.45),Inches(11.5),Inches(0.4),sz=14,bold=True,col=INK)
code(s,Inches(0.65),Inches(4.9),Inches(6.0),["claude"],sz=15)
star(s,Inches(0.6),Inches(6.0),Inches(12.1),"③のログインは“菊池さんから受け取ったClaudeのログイン情報”を使う")

# ── S7 STEP5-6 ──
s=sl(); ft(s,"7 / 9")
hdr(s,"STEP 5-6","Google Drive → データ取り込み(clone) ★","会社のデータと“Claudeの記憶”を取り込む")
num(s,Inches(0.65),Inches(1.92),5,r=0.46); t(s,"Google Drive を入れる",Inches(1.25),Inches(1.92),Inches(11),Inches(0.46),sz=15,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
t(s,"google.com/drive/download からインストール → ログイン（★どのアカウントかは菊池さんと確認）\nFinder →「Google Drive」→「マイドライブ」→「KHD_git_remote」を右クリック →「オフラインで使用可能にする」",Inches(1.25),Inches(2.42),Inches(11.4),Inches(0.9),sz=13,col=GRY,line_sp=1.25)
num(s,Inches(0.65),Inches(3.75),6,r=0.46); t(s,"データを取り込む(clone)",Inches(1.25),Inches(3.75),Inches(11),Inches(0.46),sz=15,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
t(s,"ターミナルに下を貼ってEnter（パスは★菊池さんが正確な1行を渡す）",Inches(1.25),Inches(4.25),Inches(11.4),Inches(0.4),sz=13,col=GRY)
code(s,Inches(1.25),Inches(4.7),Inches(11.4),['git clone "<KHD_git_remoteのパス>/khd_workspace.git" ~/01_honbu_docs_automation'],sz=13)
star(s,Inches(0.6),Inches(6.05),Inches(12.1),"記憶(memory)は起動後にClaudeへ「記憶を同期して」と頼めば自動で正しい場所へ入る")

# ── S8 STEP7 ──
s=sl(); ft(s,"8 / 9")
hdr(s,"STEP 7","起動して /company で確認（ゴール）","ここまで来たら、いつものClaudeが使える")
t(s,"ターミナルに、上から順に1行ずつ：",Inches(0.65),Inches(1.95),Inches(11.5),Inches(0.4),sz=14,bold=True,col=INK)
code(s,Inches(0.65),Inches(2.4),Inches(8.5),["cd ~/01_honbu_docs_automation","claude"],sz=15)
steps=[
 "Claudeが立ち上がったら  /company  と入力 → 秘書が起動",
 "「記憶を同期して」と頼む → 過去の方針・記録が入る ★",
 "画面に出るリンクで、必要な連携(MCP)を許可 ★（菊池さん指定の分だけ）",
]
y=3.7
for i,stp in enumerate(steps):
    cy=Inches(y)+Inches(0.6)*i
    num(s,Inches(0.65),cy,i+1,r=0.44)
    t(s,stp,Inches(1.25),cy,Inches(11.4),Inches(0.44),sz=14,col=INK,anchor=MSO_ANCHOR.MIDDLE)
star(s,Inches(0.6),Inches(6.0),Inches(12.1),"秘書が前回の続きを覚えていれば成功＝セットアップ完了")

# ── S9 困った時＆毎日のルール／チェックリスト ──
s=sl(); ft(s,"9 / 9")
hdr(s,"DAILY & HELP","毎日の使い方 と 困った時","むずかしく考えない。困ったら菊池さんにLINE")
bx(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(2.5),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(0.06),RED)
t(s,"毎日の使い方",Inches(0.8),Inches(2.12),Inches(5.5),Inches(0.4),sz=15,bold=True,col=RED)
t(s,"始める時：「同期して」と言う（最新を取込）\n終わる時：「保存して」と言う（みんなに反映）\n※pull/pushはClaudeが代行。gitは覚えなくてOK",Inches(0.85),Inches(2.6),Inches(5.5),Inches(1.7),sz=13,col=INK,line_sp=1.3)
bx(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(2.5),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(0.06),RED)
t(s,"困った時",Inches(7.03),Inches(2.12),Inches(5.5),Inches(0.4),sz=15,bold=True,col=RED)
t(s,"・赤い文字やエラーが出たら、その画面を\n　スクショして菊池さんにLINE（文ごと）\n・何度やり直してもMacは壊れない\n・同じ物を母艦と同時に書き換えない",Inches(7.03),Inches(2.6),Inches(5.6),Inches(1.7),sz=13,col=INK,line_sp=1.3)
bx(s,Inches(0.55),Inches(4.65),Inches(12.23),Inches(2.0),REDBG); bx(s,Inches(0.55),Inches(4.65),Inches(0.1),Inches(2.0),RED)
t(s,"チェックリスト",Inches(0.8),Inches(4.78),Inches(11),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"□ ターミナルを開いた　□ git(xcode-select)　□ Node.js　□ Claude Code＋ログイン★",Inches(0.8),Inches(5.25),Inches(11.7),Inches(0.4),sz=13.5,col=INK)
t(s,"□ Google Drive＋KHD_git_remoteをオフライン化★　□ データをclone★　□ /company で確認",Inches(0.8),Inches(5.7),Inches(11.7),Inches(0.4),sz=13.5,col=INK)
t(s,"□「記憶を同期して」で記憶が入った　□ 必要な連携(MCP)を許可した",Inches(0.8),Inches(6.15),Inches(11.7),Inches(0.4),sz=13.5,col=INK)

out="/Users/kikuchikenta/01_honbu_docs_automation/KHD_ゆーしMac_セットアップマニュアル.pptx"
prs.save(out); print("saved:",out,"/ slides:",len(prs.slides._sldIdLst))

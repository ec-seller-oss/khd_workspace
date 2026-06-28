# -*- coding: utf-8 -*-
"""03_ゆーしMac かんたん運用ガイド（PC初心者ゆーし向け・KHDクリーム白×レンガ赤）
方針：ローカルで動かす旧方式(データ無しで使えなかった)をやめ、Chromeリモートで自宅の旧Mac(常時起動)に
入って、そこのクロードを使うだけ、に作り替え。大きい字・最小手順・安心トーン。
出力: 03_ゆーしMac_かんたん運用ガイド_260628.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); GRN=RGBColor(0x2E,0x7D,0x46)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold
        r.font.color.rgb=col; r.font.name=FONT
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.6),sz=25,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.5),Inches(12.1),Inches(0.35),sz=13,col=GRY)

def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"ゆーしMac かんたん運用ガイド  ｜  困ったら 菊池さんにLINE  ｜  2026-06-28",Inches(0.5),H-Inches(0.42),Inches(11),Inches(0.32),sz=9,col=GRY)

def panel(s,x,y,w,h,title,lines,tcol=RED,lh=0.62,tsz=17,bsz=14):
    bx(s,x,y,w,h,CARD,line=CARDLN,lw=1.0); bx(s,x,y,w,Inches(0.07),tcol)
    t(s,title,x+Inches(0.3),y+Inches(0.2),w-Inches(0.5),Inches(0.5),sz=tsz,bold=True,col=tcol)
    yy=y+Inches(0.9)
    for ln,col in lines:
        t(s,ln,x+Inches(0.32),yy,w-Inches(0.6),Inches(lh),sz=bsz,col=col,line_sp=1.15); yy=yy+Inches(lh)

def band(s,y,text,col=REDBG,tc=REDD,h=0.7):
    bx(s,Inches(0.55),y,Inches(12.23),Inches(h),col); bx(s,Inches(0.55),y,Inches(0.1),Inches(h),RED)
    t(s,text,Inches(0.9),y,Inches(11.6),Inches(h),sz=15,bold=True,col=tc,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.1)

# ── S1 表紙 ──
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"YUUSHI ｜ かんたん運用ガイド",Inches(0.9),Inches(1.55),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"パソコンで クロードを使う",Inches(0.88),Inches(2.2),Inches(11.7),Inches(0.9),sz=44,bold=True,col=INK)
t(s,"― 自宅のMacに「入って」使うだけ ―",Inches(0.9),Inches(3.15),Inches(11.7),Inches(0.7),sz=24,bold=True,col=RED)
t(s,"むずかしいことは、考えなくてOK 🙆\nインストールも、同期も、もう要りません。Chromeを開いて入るだけ。",
  Inches(0.9),Inches(4.1),Inches(11.6),Inches(0.9),sz=15,col=GRY,line_sp=1.3)
bx(s,Inches(0.9),Inches(5.5),Inches(11.5),Inches(0.95),REDBG); bx(s,Inches(0.9),Inches(5.5),Inches(0.1),Inches(0.95),RED)
t(s,"今までのやり方（自分のMacで動かす）は、もうやめてOK。",Inches(1.2),Inches(5.62),Inches(11),Inches(0.4),sz=15,bold=True,col=REDD)
t(s,"これからは「菊池さんの自宅Mac（ずっと起動中）」に入って、いつものクロードを使います。",Inches(1.2),Inches(6.02),Inches(11),Inches(0.35),sz=13,col=INK)
t(s,"KHD  ｜  菊池 研太  ｜  2026-06-28",Inches(0.9),Inches(6.75),Inches(11),Inches(0.35),sz=11,col=GRY)

# ── S2 しくみ：なぜ前はダメ→今は大丈夫 ──
s=sl(); ft(s)
hdr(s,"WHY IT WORKS NOW","なんで前は使えなかった？ ── 今は大丈夫","“クロードと会社のデータがある場所”に入るかどうか、の違いだけ")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.0),"😣 前（自分のMacで動かす）",[
 ("自分のMacでクロードを起動した",INK),
 ("でも 会社のデータ・記憶が",INK),
 ("　そのMacに入っていなかった",GRY),
 ("→「何も知らないクロード」になり",REDD),
 ("　使えなかった",REDD),
],tcol=GRY,lh=0.56)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.0),"🙆 今（自宅Macに入って使う）",[
 ("菊池さんの自宅Mac（24時間 起動）に",INK),
 ("　Chromeで“入って”使う",INK),
 ("そこに 全データ・記憶が そろってる",INK),
 ("→ 最初から「いつものクロード」",GRN),
 ("　インストールも同期も要らない",GRN),
],tcol=GRN,lh=0.56)
band(s,Inches(5.25),"ゆーしのMac＝ただの「窓口」。本物のクロードと会社のデータは、菊池さんの自宅Macの中にあります。",h=1.0)

# ── S3 初回だけ（菊池さんがやる）──
s=sl(); ft(s)
hdr(s,"FIRST TIME ONLY","最初の1回だけ ── ほぼ菊池さんがやります","ゆーしは「準備できたよ」の連絡を待つだけでOK")
panel(s,Inches(0.55),Inches(1.95),Inches(7.6),Inches(4.5),"菊池さんがやること",[
 ("① 自宅の旧Mac（ずっと起動中）に",INK),
 ("　Chromeリモートデスクトップを設定",GRY),
 ("② ゆーしのChromeに「母艦Mac」が",INK),
 ("　出るように共有",GRY),
 ("③ 数字のパスワード（PIN）をLINEで渡す",INK),
 ("④「準備できたよ」とゆーしに連絡",INK),
],lh=0.6,bsz=15)
panel(s,Inches(8.35),Inches(1.95),Inches(4.43),Inches(4.5),"ゆーしがやること",[
 ("Chromeが入っているか だけ確認",INK),
 ("（無ければ google.com/chrome",GRY),
 ("　からインストール）",GRY),
 ("",INK),
 ("あとは菊池さんの",INK),
 ("「準備できたよ」を 待つ ☕",GRN),
],tcol=GRN,lh=0.58,bsz=14)

# ── S4 毎日の使い方（4ステップ・大きく）──
s=sl(); ft(s)
hdr(s,"EVERY DAY","毎日の使い方 ── たった4ステップ","上から順にやるだけ。覚えなくて大丈夫")
steps=[
 ("Chrome（カラフルな丸）を 開く",""),
 ("「母艦Mac」を 押す → 数字パスワードを入れる","菊池さんからもらったPIN"),
 ("いつものクロードに「○○して」と 話しかける","ふつうに日本語で書くだけ"),
 ("終わったら、まどを 閉じる","保存はクロードがやる。何もしなくてOK"),
]
yy=Inches(2.05); gap=Inches(1.18)
for i,(main,sub) in enumerate(steps):
    y=yy+gap*i
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.7),y,Inches(0.82),Inches(0.82))
    c.fill.solid(); c.fill.fore_color.rgb=RED; c.line.fill.background(); c.shadow.inherit=False
    t(s,str(i+1),Inches(0.7),y,Inches(0.82),Inches(0.82),sz=28,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bx(s,Inches(1.75),y,Inches(11.0),Inches(0.82),CARD,line=CARDLN,lw=0.8); bx(s,Inches(1.75),y,Inches(0.08),Inches(0.82),RED)
    if sub:
        t(s,main,Inches(2.05),y+Inches(0.1),Inches(10.5),Inches(0.4),sz=18,bold=True,col=INK)
        t(s,sub,Inches(2.05),y+Inches(0.5),Inches(10.5),Inches(0.3),sz=12.5,col=GRY)
    else:
        t(s,main,Inches(2.05),y,Inches(10.5),Inches(0.82),sz=18,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)

# ── S5 入った後の確認 ＆ 困った時 ──
s=sl(); ft(s)
hdr(s,"CHECK & HELP","入れたか確認 ＆ 困った時","この2つだけ覚えておけばOK")
panel(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.1),"✅ 入れたか の確認",[
 ("クロードの画面で /company と打つ",INK),
 ("→ 秘書が出てくれば 成功 🎉",GRN),
 ("「同期して」と言うと、最新の状態に",INK),
 ("　なります（最初に1回 言えばOK）",GRY),
],tcol=GRN,lh=0.6,bsz=14)
panel(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.1),"🆘 困った時",[
 ("赤い字・エラーが出たら、何も触らず",INK),
 ("　画面をスクショ → 菊池さんにLINE",REDD),
 ("何回やり直しても Macは壊れません",INK),
 ("わからない時は、まず菊池さんに聞く",GRY),
],tcol=REDD,lh=0.6,bsz=14)
band(s,Inches(5.35),"⚠️ ひとつだけ：菊池さんが作業中の案件は、同時に書かないで（先にLINEで一声）。",h=1.0)

# ── S6 チェックリスト＆一言 ──
s=sl(); ft(s)
hdr(s,"CHECKLIST","これだけ覚えればOK","")
panel(s,Inches(0.55),Inches(1.8),Inches(12.23),Inches(2.7),"毎日の流れ（コピーして手元に）",[
 ("1. Chromeを開く　→　2.「母艦Mac」を押す（数字PIN）",INK),
 ("3. クロードに「○○して」と話しかける　→　4. 終わったら窓を閉じる",INK),
 ("※ 最初の1回だけ「同期して」。保存はクロードが自動。困ったらスクショ→LINE。",GRY),
],lh=0.7,bsz=16)
bx(s,Inches(0.55),Inches(4.8),Inches(12.23),Inches(1.5),REDBG); bx(s,Inches(0.55),Inches(4.8),Inches(0.1),Inches(1.5),RED)
t(s,"あなたは「話しかける」だけ。",Inches(0.95),Inches(5.0),Inches(11.4),Inches(0.5),sz=22,bold=True,col=REDD)
t(s,"インストールも同期も、もう気にしなくて大丈夫。Chromeで入って、いつも通り頼むだけです。",
  Inches(0.95),Inches(5.6),Inches(11.4),Inches(0.5),sz=14,col=INK)

out="03_ゆーしMac_かんたん運用ガイド_260628.pptx"
prs.save(out)
print("saved:",out,"slides:",len(prs.slides._sldIdLst))

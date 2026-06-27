# -*- coding: utf-8 -*-
"""新Mac3/Mac2 セットアップ手順書 9枚 pptx生成（KHDダークテーマ・テキスト軽量）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG  = RGBColor(0x0D,0x1B,0x2A)   # navy
ACC = RGBColor(0xE8,0xA8,0x00)   # amber
WHT = RGBColor(0xFF,0xFF,0xFF)
LGR = RGBColor(0xCC,0xD6,0xE0)
BLU = RGBColor(0x4A,0x9E,0xCB)
GRN = RGBColor(0x4C,0xC2,0x8A)
ORG = RGBColor(0xF0,0x9F,0x27)

W, H = Inches(13.33), Inches(7.5)
prs = Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK = prs.slide_layouts[6]

def slide(title, kicker=""):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    # accent bar
    bar = s.shapes.add_shape(1, Inches(0.55), Inches(0.55), Inches(0.12), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.85), Inches(0.45), Inches(12), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; r = p.add_run(); r.text=title
    r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=WHT
    if kicker:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text=kicker
        r2.font.size=Pt(14); r2.font.color.rgb=ACC
    return s

def body(s, lines, top=1.85, size=17, gap=0.07):
    """lines: list of (text, color, indent_level)"""
    tb = s.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.7), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap=True
    first=True
    for text,col,ind in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.space_after=Pt(8); p.level=ind
        r=p.add_run(); r.text=text
        r.font.size=Pt(size if ind==0 else size-2)
        r.font.color.rgb=col
        if ind==0 and (text.startswith("【") or text.endswith("：")):
            r.font.bold=True
    return s

# S1 タイトル
s=slide("新Mac(Mac3) ＆ Mac2 セットアップ手順書","3台を1アカウントで動かし、Claudeが起動するまで ／ 2026-06-26")
body(s,[
 ("Mac1 ＝ 現メイン（今のMac）", WHT,0),
 ("Mac2 ＝ ゆーしMac（既存・サブ）", WHT,0),
 ("Mac3 ＝ 新ハイスペック（母艦に昇格・Mac1の代替）", ACC,0),
 ("", LGR,0),
 ("Mac3はまず「Mac1と同じように使える」が目標。次ページから手順。", LGR,0),
], top=2.4, size=20)

# S2 Apple連携
s=slide("「ボタンひとつ同期」できてる？","結論：初期コピーはApple／継続同期はgit（両方準備済）")
body(s,[
 ("✅ iCloud・同一Apple ID（菊池研太）設定済", GRN,0),
 ("→ 移行アシスタントでMac1→Mac3を丸ごとクローン（ほぼボタン一発）", LGR,1),
 ("✅ Handoff／ユニバーサルクリップボード／AirDropも使える", GRN,0),
 ("⚠️ ただし「継続的にボタンひとつ同期」はAppleに無い", ORG,0),
 ("iCloudはApple純正データだけ（書類/写真/Safari/鍵）", LGR,1),
 ("→ 作業フォルダ＆Claudeの脳はgit（Drive内・設定済）で同期。コマンドはClaude代行", BLU,0),
])

# S3 全体像
s=slide("全体像（1アカウント・3台）","正本はGoogle Drive内の2リポジトリ（両方push済）")
body(s,[
 ("khd_workspace.git ＝ 作業フォルダ  ~/01_honbu_docs_automation", WHT,0),
 ("khd_memory.git ＝ Claudeの脳（記憶）  ~/.claude/.../memory", WHT,0),
 ("置き場 ＝ マイドライブ/KHD_git_remote/（ec-sellerアカウント内）", LGR,1),
 ("iCloud ＝ Apple純正データ（書類/写真/Safari/鍵）", LGR,0),
 ("各Macで同じAnthropicアカウントにログイン", LGR,0),
])

# S4 Mac3 P1
s=slide("【Mac3】Phase 1 ｜ 移行アシスタント","「Mac1と同じ」を最速で作る一手")
body(s,[
 ("1. 初期設定ウィザード →「Macから情報を移行」を選ぶ", WHT,0),
 ("2. Mac1を選び、丸ごとコピー（アプリ＋データ＋設定）", WHT,0),
 ("3. Claude・Cursor・Chrome・Drive・Office・LINE＋Node＋claude CLI が全部乗る", WHT,0),
 ("Thunderboltケーブル直結が最速（Wi-Fiより圧倒的）", BLU,0),
 ("Homebrewは未導入だが、移行で必要物は乗るので原則不要", LGR,0),
])

# S5 Mac3 P2
s=slide("【Mac3】Phase 2 ｜ サインイン＆Drive","同期の土台を通電する")
body(s,[
 ("Apple ID：kemkemsp@yahoo.co.jp でログイン・iCloud Drive ON", WHT,0),
 ("Google Drive Desktop：ec-seller@kikuchi-hd.net でログイン→同期開始", WHT,0),
 ("⚠️ 最重要：Finderで マイドライブ/KHD_git_remote を右クリック", ORG,0),
 ("→「オフラインで使用可能にする」（忘れると同期が失敗）", ORG,1),
])

# S6 Mac3 P3
s=slide("【Mac3】Phase 3 ｜ リポジトリ同期","移行で来たローカルを、Drive正本に合わせるだけ")
body(s,[
 ("作業：~/01_honbu_docs_automation（origin設定済→pullで最新化）", WHT,0),
 ("脳：~/.claude/.../memory（同上）", WHT,0),
 ("✅ ユーザー名・パスがMac1と同一なら、メモリの保存先まで自動一致", GRN,0),
 ("実コマンド（git pull）はClaudeが実行。菊池さんは「同期して」と言うだけ", LGR,0),
])

# S7 Mac3 P4
s=slide("【Mac3】Phase 4 ｜ Claude起動（ゴール）","ここまで来たらMac3は母艦として稼働")
body(s,[
 ("1. ターミナルで claude → 同アカウントでログイン", WHT,0),
 ("2. MCP再認証：Drive／Calendar／Gmail／Tasks／Notion／LINE／TradingView", WHT,0),
 ("3. プロジェクトで /company → 秘書起動＋メモリ反映を確認", WHT,0),
 ("✅ 秘書が前回の続きを覚えていれば成功＝Mac1と同じ状態", GRN,0),
])

# S8 Mac2
s=slide("【Mac2 ゆーしMac】既存Mac＝クローン経路","移行アシスタントは使わない（既に使用中のため）")
body(s,[
 ("1. Google Drive Desktop＝ec-sellerでログイン→KHD_git_remoteをオフライン化", WHT,0),
 ("2. Node・claude CLI が無ければ入れる（移行が無いため）", WHT,0),
 ("3. 2リポジトリを同じパスへ git clone（作業＋脳）", WHT,0),
 ("4. claudeログイン＋MCP再認証 → /company", WHT,0),
 ("書き手は母艦Mac3。Mac2は読み／軽編集中心（競合回避）", LGR,0),
])

# S9 運用
s=slide("日々の運用 ＆ チェックリスト","鉄則は2つだけ")
body(s,[
 ("開始＝git pull ／ 終了＝git push（Claudeが代行）", ACC,0),
 ("書き手は1台（母艦Mac3）。push後はDrive同期完了を待つ／他Macはpull前に待つ", WHT,0),
 ("", LGR,0),
 ("Mac3：移行→サインイン→KHD_git_remoteオフライン→pull→claude→/company", BLU,0),
 ("Mac2：Driveログイン→オフライン→clone→node/claude→/company", BLU,0),
])

out="/Users/kikuchikenta/01_honbu_docs_automation/KHD_新Mac3Mac2_セットアップ手順書.pptx"
prs.save(out)
print("saved:", out)

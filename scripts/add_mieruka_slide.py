"""
clinic_dx_v8.pptx に「経営の見える化」GIVEカードスライドを追加する
挿入位置: 最終スライド（ご相談はお気軽に）の直前
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import copy
import shutil

SRC  = "/Users/kikuchikenta/01_honbu_docs_automation/clinic_dx_v8.pptx"
DEST = "/Users/kikuchikenta/01_honbu_docs_automation/clinic_dx_v8.pptx"

# ===== カラー定義（clinic_dx 準拠）=====
NAVY    = RGBColor(0x1C, 0x35, 0x57)
RED     = RGBColor(0xAA, 0x2E, 0x26)
CREAM   = RGBColor(0xF9, 0xF6, 0xEF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW  = RGBColor(0xFF, 0xF9, 0xC4)
BLUE    = RGBColor(0xE3, 0xF2, 0xFD)
GRAY_LT = RGBColor(0xEC, 0xEF, 0xF1)
GREEN   = RGBColor(0xE8, 0xF5, 0xE9)
DARK_GR = RGBColor(0x2E, 0x5F, 0x8A)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=Pt(11), bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    txbox = slide.shapes.add_textbox(l, t, w, h)
    tf = txbox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or NAVY
    if line_spacing:
        from pptx.util import Pt as pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txbox

def add_rect_text(slide, text, l, t, w, h,
                  fill_rgb=CREAM, line_rgb=None,
                  font_size=Pt(11), bold=False, text_color=None,
                  align=PP_ALIGN.LEFT, v_anchor=None):
    shape = add_rect(slide, l, t, w, h, fill_rgb, line_rgb, Pt(1.5))
    tf = shape.text_frame
    tf.word_wrap = True
    if v_anchor:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = text_color or NAVY
    return shape

# ===== メイン処理 =====
prs = Presentation(SRC)
W = prs.slide_width   # 12192000 EMU = 13.33in
H = prs.slide_height  # 6858000  EMU = 7.5in

# ブランクレイアウト（index=6 が通常ブランク）
layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout)

inch = Inches(1)

# --------- 背景: クリーム ---------
bg = add_rect(slide, 0, 0, W, H, CREAM)

# --------- ヘッダーバー（レンガ赤） ---------
hdr_h = Inches(0.85)
add_rect(slide, 0, 0, W, hdr_h, RED)

# ヘッダー左: タイトル
add_text_box(slide,
    "GIVE①   経営の見える化シート　|　先生専用の「1枚」を無料でお作りします",
    Inches(0.3), Inches(0.1), Inches(9.5), Inches(0.65),
    font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ヘッダー右: ラベル
add_text_box(slide,
    "クリニックDX「My AI」  ｜  AI医療コンサル",
    Inches(9.8), Inches(0.25), Inches(3.3), Inches(0.4),
    font_size=Pt(9), color=WHITE, align=PP_ALIGN.RIGHT)

# --------- サブタイトル ---------
add_text_box(slide,
    "クリニックも店舗も「見える化」すれば、次の一手が変わる",
    Inches(0.4), Inches(0.95), W - Inches(0.8), Inches(0.38),
    font_size=Pt(13), bold=False, color=DARK_GR, align=PP_ALIGN.LEFT)

# --------- 3カラム構成（左: 悩み / 中: ダッシュボード / 右: 解決）---------
col_top  = Inches(1.4)
col_h    = Inches(3.5)
col_w    = Inches(3.8)
gap      = Inches(0.25)

# ---- 左カラム: 先生が抱えるお悩み ----
lx = Inches(0.25)
add_rect(slide, lx, col_top, col_w, Inches(0.38), RED)
add_text_box(slide, "😰 先生が抱えるお悩み",
    lx + Inches(0.05), col_top, col_w, Inches(0.36),
    font_size=Pt(12), bold=True, color=WHITE)

pain_items = [
    "💸  来月の資金繰りが不安…",
    "📊  売上はあるのに手元現金が少ない",
    "🔢  損益分岐点（何人来れば黒字か）が",
    "　　わからない",
    "📋  決算書を見ても経営判断に使えない",
    "⏰  月次確認を税理士に任せていて",
    "　　自分では把握できていない",
]
pain_y = col_top + Inches(0.42)
for item in pain_items:
    add_text_box(slide, item,
        lx + Inches(0.08), pain_y, col_w - Inches(0.1), Inches(0.4),
        font_size=Pt(10.5), color=NAVY)
    pain_y += Inches(0.42)

# ---- 中カラム: ダッシュボードイメージ ----
cx = lx + col_w + gap
add_rect(slide, cx, col_top, col_w, Inches(0.38), NAVY)
add_text_box(slide, "📐 見える化ダッシュボード（サンプル）",
    cx + Inches(0.05), col_top, col_w, Inches(0.36),
    font_size=Pt(11), bold=True, color=WHITE)

# ミニダッシュボード表示
db_y = col_top + Inches(0.44)
db_items = [
    (YELLOW,  "🟡  月初現金（入力）",    "¥ 2,850,000"),
    (GRAY_LT, "⬜  来月予測売上（自動）", "¥ 4,200,000"),
    (BLUE,    "🟦  BEP来客数（自動）",   "38 人/日"),
    (GRAY_LT, "⬜  固定費合計（自動）",  "¥ 1,680,000"),
    (GREEN,   "🟢  月次損益（自動）",    "+ ¥ 520,000"),
    (GRAY_LT, "⬜  FL比率（自動）",      "57.3%（適正）"),
    (BLUE,    "🟦  資金繰り着地（自動）","¥ 3,370,000 →安全"),
]
for fill, label, val in db_items:
    row_h = Inches(0.38)
    add_rect_text(slide, label,
        cx + Inches(0.05), db_y, col_w * 0.58, row_h,
        fill_rgb=fill, line_rgb=GRAY_LT,
        font_size=Pt(9.5), text_color=NAVY)
    add_rect_text(slide, val,
        cx + Inches(0.05) + col_w * 0.59, db_y, col_w * 0.36, row_h,
        fill_rgb=fill, line_rgb=GRAY_LT,
        font_size=Pt(9.5), bold=True, text_color=NAVY, align=PP_ALIGN.RIGHT)
    db_y += row_h + Inches(0.04)

# ← 矢印（左から中）
add_text_box(slide,
    "→",
    cx - Inches(0.22), col_top + Inches(1.5), Inches(0.22), Inches(0.5),
    font_size=Pt(20), bold=True, color=RED, align=PP_ALIGN.CENTER)

# ---- 右カラム: 解決・GIVE内容 ----
rx = cx + col_w + gap
add_rect(slide, rx, col_top, col_w, Inches(0.38), DARK_GR)
add_text_box(slide, "✅ GIVEで先生が手に入るもの",
    rx + Inches(0.05), col_top, col_w, Inches(0.36),
    font_size=Pt(12), bold=True, color=WHITE)

solve_items = [
    "📍  月初現金を入れるだけで",
    "    売上見込・資金繰り・損益が全自動",
    "",
    "🏥  クリニック専用にカスタマイズ",
    "    （保険診療単価・外来数をベース）",
    "",
    "📈  来月の資金ショートを",
    "    2〜3ヶ月前から予知できる",
    "",
    "💡  銀行に見せられる",
    "    経営計画書の素材にもなる",
]
sol_y = col_top + Inches(0.44)
for item in solve_items:
    add_text_box(slide, item,
        rx + Inches(0.08), sol_y, col_w - Inches(0.1), Inches(0.32),
        font_size=Pt(10), color=NAVY)
    sol_y += Inches(0.32)

# ← 矢印（中から右）
add_text_box(slide,
    "→",
    rx - Inches(0.22), col_top + Inches(1.5), Inches(0.22), Inches(0.5),
    font_size=Pt(20), bold=True, color=RED, align=PP_ALIGN.CENTER)

# --------- メニューバー（4段）---------
menu_top = Inches(5.05)
menu_h   = Inches(0.55)
menu_w   = (W - Inches(0.5)) / 4
menus = [
    (RED,     "M0  無料GIVE",      "先生専用シートを1回無料で作成・お見せします",           WHITE, WHITE),
    (DARK_GR, "M1  構築サポート",  "20〜40万（買い切り）導入・移行をフルサポート",           WHITE, WHITE),
    (NAVY,    "M2  月次伴走",       "3〜5万/月（定額）月初に一緒に数字を読む",               WHITE, YELLOW),
    (GRAY_LT, "M3  フルカスタム",  "50万〜  銀行提出書類・AI連携まで",                       NAVY,  NAVY),
]
mx = Inches(0.25)
for fill, title, desc, tc, dc in menus:
    add_rect_text(slide, f"{title}\n{desc}",
        mx, menu_top, menu_w - Inches(0.08), menu_h,
        fill_rgb=fill, line_rgb=None,
        font_size=Pt(9), bold=False, text_color=tc,
        align=PP_ALIGN.CENTER)
    # タイトル部分を太字にしたいが add_rect_text は単一runなので
    # 別のtextboxを重ねてタイトルだけ強調
    add_text_box(slide, title,
        mx + Inches(0.05), menu_top + Inches(0.05),
        menu_w - Inches(0.18), Inches(0.28),
        font_size=Pt(11), bold=True, color=tc, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc,
        mx + Inches(0.05), menu_top + Inches(0.28),
        menu_w - Inches(0.18), Inches(0.26),
        font_size=Pt(8.5), bold=False, color=dc, align=PP_ALIGN.CENTER)
    mx += menu_w

# --------- フッター ---------
ftr_top = H - Inches(0.32)
add_rect(slide, 0, ftr_top, W, Inches(0.32), NAVY)
add_text_box(slide,
    "★ 本商品は医療コンサルの「おまけ」ではなく、先生の経営と資金繰りを守る「本命GIVE」です。  ｜  KHD 04_コンサル（テナントアシスト・ウイン）",
    Inches(0.2), ftr_top + Inches(0.04), W - Inches(0.4), Inches(0.26),
    font_size=Pt(8), color=WHITE, align=PP_ALIGN.CENTER)

# --------- スライドを最終スライドの直前に移動 ---------
# 現在は末尾に追加されているので、最後から2番目に移動
xml_slides = prs.slides._sldIdLst
# 末尾スライドのelement
new_slide_el = xml_slides[-1]
# 最後の「ご相談」スライドの直前（-2番目）へ移動
last_el = xml_slides[-2]
xml_slides.remove(new_slide_el)
# last_elの前に挿入
last_el.addprevious(new_slide_el)

prs.save(DEST)
print(f"✅ 保存完了: {DEST}")
print(f"   スライド数: {len(prs.slides)}")

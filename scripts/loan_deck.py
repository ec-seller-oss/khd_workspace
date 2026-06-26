#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
融資資料デッキ生成エンジン（汎用・データ駆動）
================================================================
screen_property の判定データ → 銀行に出せる融資資料デッキ(pptx)。
build(data, out) を呼ぶだけ。data に無い項目はスライドを省略/“—”でフォールバック。
North Star成果物。KHDデザイン（クリーム白×レンガ赤）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa
from pathlib import Path

CREAM = RGBColor(0xF9, 0xF6, 0xEF); BRICK = RGBColor(0xAA, 0x2E, 0x26)
INK = RGBColor(0x33, 0x2A, 0x28); GRAY = RGBColor(0x77, 0x70, 0x6E)
GREEN = RGBColor(0x1F, 0x7A, 0x1F); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINEC = RGBColor(0xE3, 0xD8, 0xD1)
FONT = "Meiryo"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _bg(s):
    r = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = CREAM; r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); return s


def box(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, ls=1.1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = color
    return tb


def bar(s, x, y, w, h, color=BRICK):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background(); r.shadow.inherit = False
    return r


def header(s, no, title):
    bar(s, 0, 0.55, 0.22, 0.65, BRICK)
    box(s, 0.5, 0.5, 11.8, 0.8, title, size=26, color=BRICK, bold=True)
    box(s, 12.4, 0.6, 0.7, 0.5, f"{no:02d}", size=16, color=GRAY, bold=True, align=PP_ALIGN.RIGHT)
    bar(s, 0.5, 1.32, 12.3, 0.02, LINEC)


def kpi_card(s, x, y, w, label, value, vcolor=BRICK, sub=""):
    c = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.5))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.color.rgb = LINEC; c.line.width = Pt(1); c.shadow.inherit = False
    box(s, x, y + 0.12, w, 0.4, label, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    box(s, x, y + 0.45, w, 0.7, value, size=24, color=vcolor, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        box(s, x, y + 1.12, w, 0.35, sub, size=10, color=GRAY, align=PP_ALIGN.CENTER)


def table(s, x, y, w, rows, col1=0.25, rh=0.5, fs=13):
    cy = y
    for i, (k, v) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else RGBColor(0xF2, 0xEC, 0xE6)
        r = s.shapes.add_shape(1, Inches(x), Inches(cy), Inches(w), Inches(rh))
        r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
        box(s, x + 0.12, cy + 0.03, w * col1 - 0.2, rh, k, size=fs, color=BRICK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        box(s, x + w * col1, cy + 0.03, w * (1 - col1) - 0.15, rh, str(v), size=fs, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        cy += rh


def footer(s):
    box(s, 0.5, 7.05, 9, 0.35,
        "菊池ホールディングス（KHD）／ 本資料は融資ご相談用の概算であり最終条件は協議のうえ確定します",
        size=9, color=GRAY)


def g(d, k, default="—"):
    v = d.get(k)
    return v if v not in (None, "") else default


def build(data, out):
    """data(dict) → 融資資料デッキpptx。最小=表紙/案件概要/物件/路線価/収支/融資/結び。"""
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    props = data.get("props", [])
    is_income = data.get("is_income", False)
    no = 0

    # 表紙
    s = add_slide(prs)
    bar(s, 0, 2.55, 13.333, 0.10, BRICK); bar(s, 0, 4.55, 13.333, 0.04, LINEC)
    box(s, 0.95, 1.4, 11, 0.5, "CONFIDENTIAL ／ 融資ご相談資料", size=12, color=GRAY)
    box(s, 0.9, 2.75, 11.5, 1.0, g(data, "title", "融資ご相談資料"), size=44, color=BRICK, bold=True)
    box(s, 0.95, 3.78, 11.5, 0.7, g(data, "sub", ""), size=19, color=INK)
    if data.get("bank"):
        box(s, 0.95, 4.75, 11.5, 0.5, data["bank"], size=18, color=INK, bold=True)
    box(s, 0.95, 6.1, 11.5, 0.9, f"{g(data,'company','菊池ホールディングス（KHD）')}\n{g(data,'rep','')}",
        size=15, color=GRAY)

    # 案件概要
    no += 1; s = add_slide(prs); header(s, no, "案件概要")
    kc = data.get("kpi", {})
    cards = [(k, v) for k, v in kc.items()][:4]
    cw = 2.9; x = 0.5
    for label, (val, col, sub) in [(k, kc[k]) for k in list(kc)[:4]]:
        kpi_card(s, x, 2.0, cw, label, val, col, sub); x += cw + 0.2
    table(s, 0.5, 3.9, 12.3, data.get("overview", [("所在地", "—")]), col1=0.20, rh=0.5)
    footer(s)

    # 物件個別
    for i, pr in enumerate(props):
        no += 1; s = add_slide(prs); header(s, no, f"物件{i+1} ― {g(pr,'name','')}")
        cx = 0.5
        for lab, key, col in [("取得価格", "price", BRICK), ("年間賃料", "rent_y", BRICK),
                              ("利回り", "yield", GREEN), ("現況", "occ", GREEN)]:
            kpi_card(s, cx, 1.55, 3.0, lab, g(pr, key), col); cx += 3.2
        rows = [(k, g(pr, v)) for k, v in [("所在地", "addr"), ("構造／築年", "struct_year"),
                ("土地面積", "land"), ("建物面積", "bldg"), ("月額賃料", "rent_m"), ("特記", "note")]]
        has_img = pr.get("img") and Path(pr["img"]).exists()
        table(s, 0.5, 3.35, (7.1 if has_img else 12.3), rows, col1=(0.28 if has_img else 0.20), rh=0.5)
        if has_img:
            fr = s.shapes.add_shape(1, Inches(7.95), Inches(3.30), Inches(4.85), Inches(2.95))
            fr.fill.solid(); fr.fill.fore_color.rgb = WHITE; fr.line.color.rgb = LINEC; fr.shadow.inherit = False
            s.shapes.add_picture(pr["img"], Inches(8.05), Inches(3.40), width=Inches(4.65))
            box(s, 7.95, 6.30, 4.85, 0.3, "現況外観", size=10, color=GRAY, align=PP_ALIGN.CENTER)
        footer(s)

    # 路線価・担保
    rk = data.get("rosenka", {})
    if rk:
        no += 1; s = add_slide(prs); header(s, no, "担保・積算評価（融資の根拠）")
        cards = [("相続税路線価", "souzoku", INK), ("土地値割合", "tochine", GREEN),
                 ("積算割合", "sekisan", BRICK)]
        cx = 0.5
        for lab, key, col in cards:
            kpi_card(s, cx, 1.6, 3.9, lab, g(rk, key), col); cx += 4.1
        box(s, 0.5, 3.5, 12.3, 3.0,
            "■ 土地値割合（玉川式・投資の下値）＝路線価×面積÷価格。0.4以上で安全圏。\n"
            f"　本件 {g(rk,'tochine')}。\n"
            "■ 積算割合（銀行の融資根拠）＝積算評価÷価格。1.0以上で積算が出て融資が付きやすい。\n"
            f"　本件 {g(rk,'sekisan')}。積算＝土地（路線価×面積）＋建物（再調達単価×延床×残存÷耐用）。\n"
            f"　参考：{g(rk,'coverB','—')}\n"
            "■ 路線価はreinfolib周辺取引価格より推定。最終資料は国税庁路線価図で実値確認。",
            size=13, color=INK, ls=1.25)
        footer(s)

    # 収支・収益性
    no += 1; s = add_slide(prs); header(s, no, "収支・収益性")
    inc = data.get("income", {})
    if is_income and inc:
        cx = 0.5
        for lab, key, col in [("年間賃料", "rent_y", BRICK), ("表面利回り", "omote", GREEN),
                              ("実質利回り", "real", GREEN), ("税引前CF", "cf", GREEN)]:
            kpi_card(s, cx, 1.6, 3.0, lab, g(inc, key), col); cx += 3.2
        table(s, 0.5, 3.5, 12.3, inc.get("rows", []), col1=0.30, rh=0.5)
    else:
        cap = data.get("capital", {})
        cx = 0.5
        for lab, key, col in [("取得価格", "price", INK), ("想定売却額", "uridashi", BRICK), ("粗利率", "arari", GREEN)]:
            kpi_card(s, cx, 1.6, 3.9, lab, g(cap, key), col); cx += 4.1
        table(s, 0.5, 3.5, 12.3, cap.get("rows", [("メモ", "—")]), col1=0.30, rh=0.5)
    footer(s)

    # 融資相談
    no += 1; s = add_slide(prs); header(s, no, "融資のご相談")
    kpi_card(s, 0.5, 1.6, 4.0, "融資希望額", g(data, "loan_req"), BRICK)
    kpi_card(s, 4.7, 1.6, 4.0, "返済原資", g(data, "gensen", "—"), INK)
    kpi_card(s, 8.9, 1.6, 3.9, "スキーム", g(data, "scheme", "—"), INK)
    box(s, 0.5, 3.5, 12.3, 2.6,
        "■ ご相談事項\n・融資可否・融資可能額／対象費目\n・金利・期間・自己資金割合のご条件\n"
        "・抵当権設定・必要書類のご確認\n\n※ 詳細書類（レントロール・確定測量等）は順次ご提出いたします。",
        size=15, color=INK, ls=1.3)
    footer(s)

    # 結び
    s = add_slide(prs); bar(s, 0, 2.4, 13.333, 0.08, BRICK)
    box(s, 0.9, 2.7, 11.5, 1.4, g(data, "closing",
        "本案件へのご支援を賜りますようお願い申し上げます。"), size=20, color=INK, ls=1.3)
    box(s, 0.95, 5.4, 11.5, 0.9, f"{g(data,'company','菊池ホールディングス（KHD）')}\n{g(data,'rep','')}",
        size=15, color=GRAY)

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out

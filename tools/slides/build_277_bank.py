"""
277 高松二丁目 収益物件取得 ── 岩手銀行(西原様)提出用 査定/融資資料 v3
菊池FB全反映(2回目):
・敬称徹底(トラステン合同会社様/岩手大学生活協同組合様/岩手銀行様)・人柄が伝わるトーン
・私道=三浦哲様 単独100%(課税明細)。「共有」表記の経緯/手前道路の権利は謄本取得・調査中リスクとして明記＋位置指定申請図(三浦設計士作成)貼付
・消費税=売買代金内に収める前提を銀行と事前共有(手出し増を許さない)
・土地値=合算＋2区画分割(アパート用地/戸建用地)それぞれ路線価相当・評価額
・DSCR=計算式/根拠を可視化(NOI→返済→DSCR)
・銀行に出さない【社内用】収支全結果・消費税精査・末永様依頼・抜け漏れをスライドで追加(タイトルに削除明記)
・側溝/下水・私設管(ガス/下水)・解体(3棟/アスベスト)・ガス貸与・耐震・告知・高低差/地中/越境・必要書面チェックリスト等を別紙網羅
・銀行新情報: 保証協会上乗せ金利なし/担保調査手数料11万/印紙4万/保証料一括・借入内包
・KIKUCHIホールディングス株式会社 代表取締役 菊池研太 で統一
出力: 277_bank.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); GRN=RGBColor(0x2E,0x7D,0x32)
NAVY=RGBColor(0x2B,0x3A,0x55)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
ENT="KIKUCHIホールディングス株式会社　代表取締役　菊池 研太"
APT="/tmp/277_apt.png"; HOUSE="/tmp/277_house.png"
MYSO_A="/tmp/277_myso_apt.png"; MYSO_H="/tmp/277_myso_house.png"; KOZ="/tmp/277_koz.png"
ROAD="/tmp/277_dc_p2.png"; KAITAI="/tmp/277_kaitai_p2.png"

prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def sl(internal=False):
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    if internal:
        bx(s,Inches(0),Inches(0),W,Inches(0.32),NAVY)
        t(s,"【社内管理用 ／ 銀行提出時は削除】",Inches(0.3),Inches(0.02),Inches(12.7),Inches(0.28),sz=11,bold=True,col=WHT,align=PP_ALIGN.CENTER)
    return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub="",top=0.38):
    t(slide,eyebrow,Inches(0.6),Inches(top),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(top+0.38),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(top+0.5),Inches(12.1),Inches(0.55),sz=22,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(top+1.04),Inches(12.1),Inches(0.34),sz=11.5,col=GRY)
def ft(slide,n):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"岩手県盛岡市高松二丁目 収益物件取得プロジェクト ｜ KIKUCHIホールディングス株式会社",Inches(0.5),H-Inches(0.42),Inches(11),Inches(0.32),sz=9,col=GRY)
    t(slide,str(n),W-Inches(0.9),H-Inches(0.42),Inches(0.4),Inches(0.32),sz=9,col=GRY,align=PP_ALIGN.RIGHT)
def table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12,hi_rows=None,sub_rows=None):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    hi_rows=hi_rows or []; sub_rows=sub_rows or []
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.05)
            cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            elif ri in hi_rows: cell.fill.fore_color.rgb=REDBG
            elif ri in sub_rows: cell.fill.fore_color.rgb=GRYBG
            else: cell.fill.fore_color.rgb=REDBG if is_hi else (CARD if ri%2==1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or is_hi or (ci==0) or (ri in hi_rows) or (ri in sub_rows)
                    r.font.color.rgb=WHT if ri==0 else (RED if (is_hi or ri in hi_rows) else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
    return tb
def framed_pic(slide,path,x,y,w,h,cap=None):
    bx(slide,x-Inches(0.06),y-Inches(0.06),w+Inches(0.12),h+Inches(0.12),CARD,CARDLN,1.0)
    bx(slide,x-Inches(0.06),y-Inches(0.06),w+Inches(0.12),Inches(0.06),RED)
    slide.shapes.add_picture(path,x,y,w,h)
    if cap: t(slide,cap,x,y+h+Inches(0.08),w,Inches(0.28),sz=9.5,col=GRY,align=PP_ALIGN.CENTER)
def card(slide,x,y,w,h,title,body,big=None,accent=RED):
    bx(slide,x,y,w,h,CARD,CARDLN,1.0); bx(slide,x,y,w,Inches(0.06),accent)
    t(slide,title,x+Inches(0.2),y+Inches(0.16),w-Inches(0.4),Inches(0.4),sz=12.5,bold=True,col=accent)
    if big: t(slide,big,x+Inches(0.2),y+Inches(0.5),w-Inches(0.4),Inches(0.6),sz=24,bold=True,col=INK)
    if body: t(slide,body,x+Inches(0.2),y+(Inches(1.12) if big else Inches(0.58)),w-Inches(0.4),h-Inches(1.2),sz=11,col=INK,line_sp=1.15)

# ════════ 銀行提出スライド ════════
# S1 表紙
s=sl()
bx(s,Inches(0.6),Inches(2.15),Pt(5),Inches(2.55),RED)
t(s,"FINANCING PROPOSAL ｜ 物件査定・融資ご相談資料",Inches(0.95),Inches(2.2),Inches(11.6),Inches(0.5),sz=15,bold=True,col=RED)
t(s,"岩手県盛岡市高松二丁目",Inches(0.92),Inches(2.8),Inches(11.6),Inches(0.9),sz=38,bold=True,col=INK)
t(s,"収益アパート＋戸建　取得プロジェクト",Inches(0.92),Inches(3.7),Inches(11.6),Inches(0.8),sz=28,bold=True,col=RED)
bx(s,Inches(0.95),Inches(4.7),Inches(8.0),Pt(1.2),LINE)
t(s,"木造AP8戸（アーバンキューブ）＋木造戸建3棟（花みずき貸家）／ 取得3,700万円・土地値割合 約90.5%の堅い担保",Inches(0.95),Inches(4.85),Inches(11.7),Inches(0.5),sz=12.5,col=GRY)
t(s,ENT,Inches(0.95),Inches(6.05),Inches(11.6),Inches(0.4),sz=15,bold=True,col=INK)
t(s,"2026年6月19日　／　株式会社岩手銀行様",Inches(0.95),Inches(6.5),Inches(9),Inches(0.35),sz=11,col=GRY)

# S2 目次
s=sl(); hdr(s,"CONTENTS","目次")
items=[("01","案件概要（物件・マイソク・地番/所有者・道路・担保価値）"),
       ("02","事業計画（賃料・市場/競合・出口戦略・実績）"),
       ("03","収支計画・融資条件のご相談（DSCR根拠つき）"),
       ("04","スケジュール／売却を見据えた必要書面チェックリスト"),
       ("05","ごあいさつ／別紙：現在調査中の事項")]
y=1.95
for no,tt in items:
    bx(s,Inches(0.9),Inches(y),Inches(0.7),Inches(0.62),RED)
    t(s,no,Inches(0.9),Inches(y+0.1),Inches(0.7),Inches(0.45),sz=19,bold=True,col=WHT,align=PP_ALIGN.CENTER)
    t(s,tt,Inches(1.8),Inches(y+0.12),Inches(11),Inches(0.45),sz=15,bold=True,col=INK)
    y+=0.92
ft(s,2)

# S3 合算サマリ
s=sl(); hdr(s,"01 案件概要","物件概要（合算サマリ）","アパート・戸建を一体取得。土地が取得価格の約9割を占める担保性の高い案件です。")
rows=[["項目","内　容"],
      ["所在地","岩手県盛岡市高松二丁目 34-5（アパート）／ 34番30号（戸建3棟）"],
      ["交通","いわて銀河鉄道「青山」駅 徒歩22分"],
      ["構成","木造AP 1K×8戸（アーバンキューブ）＋ 木造戸建 3K×3棟（花みずき貸家）"],
      ["取得価格","37,000千円（アパート25,000千円＋戸建12,000千円・消費税込）"],
      ["土地面積","合計 693.82㎡（宅地646.72㎡＋私道47.10㎡）"],
      ["用途地域","第二種中高層住居専用地域／建ぺい60%・容積200%・市街化区域"],
      ["接道","東側6.0m 私道（位置指定道路）／私道は売主様の単独所有・持分100%"],
      ["管理","アパート＝岩手大学生活協同組合様／戸建＝トラステン合同会社様"],
      ["現況／利回り","賃貸中（OC・入居10/11）／表面13.0%・実質11.1%（満室想定）"]]
table(s,rows,Inches(0.6),Inches(1.78),Inches(12.1),Inches(4.9),[Inches(2.3),Inches(9.8)],sz=12,header_sz=13,hi_rows=[4,9])
ft(s,3)

# S4 アーバン明細
s=sl(); hdr(s,"01 案件概要","① アーバンキューブ（木造AP 1K×8戸）","1997年築・2023年外壁改修済。岩手大学至近で学生需要が安定し、満室稼働中です。")
rows=[["項目","内容"],["価格","25,000千円"],["土地面積","公簿335.66㎡（うち私道持分47.10㎡）"],
      ["建物","木造2階建 238.48㎡／1K×8戸"],["築年月","1997年3月（外壁改修2023年）"],
      ["現況","賃貸中・入居率8/8（満室）"],["年間賃料","4,032千円（満室）／表面16.1%"],
      ["管理／設備","岩手大学生活協同組合様／B・T別・プロパン・FF暖房・ネット対応"]]
table(s,rows,Inches(0.6),Inches(1.78),Inches(7.0),Inches(4.3),[Inches(1.9),Inches(5.1)],sz=11.5,header_sz=12.5,hi_rows=[6])
framed_pic(s,APT,Inches(8.0),Inches(1.95),Inches(4.4),Inches(3.3),"アーバンキューブ 外観")
t(s,"※ 設備・現況の出所：不動産管理トラステン合同会社様「info-sheet（マイソク）」。次頁に現物を添付します。",Inches(0.6),Inches(6.25),Inches(12),Inches(0.3),sz=10,col=GRY)
ft(s,4)

# S5 アーバン マイソク
s=sl(); hdr(s,"01 案件概要","① アーバンキューブ 物件資料（マイソク・添付）","出所＝トラステン合同会社様 info-sheet。B・T別/プロパン/FF暖房/ネット対応の記載元です。")
framed_pic(s,MYSO_A,Inches(1.5),Inches(1.95),Inches(10.3),Inches(4.4))
ft(s,5)

# S6 花みずき明細
s=sl(); hdr(s,"01 案件概要","② 花みずき貸家3棟（木造戸建 3K×3棟）","1968年築。退去時は更地化し戸建用地として分割売却できる“出口の保険”を持ちます。")
rows=[["項目","内容"],["価格","12,000千円"],["土地面積","公簿358.16㎡（108.3坪）／私道持分なし"],
      ["建物","木造2階建 計163.92㎡／3K×3棟"],["築年月","1968年7月"],
      ["現況","賃貸中・入居率2/3"],["年間賃料","1,300千円／表面10.8%（満室15.7%）"],
      ["管理／設備","トラステン合同会社様／プロパン・上下水道"]]
table(s,rows,Inches(0.6),Inches(1.78),Inches(7.0),Inches(4.3),[Inches(1.9),Inches(5.1)],sz=11.5,header_sz=12.5,hi_rows=[6])
framed_pic(s,HOUSE,Inches(8.0),Inches(1.95),Inches(4.4),Inches(2.75),"花みずき貸家 外観")
t(s,"※ 設備・現況の出所：トラステン合同会社様「info-sheet（マイソク）」。次頁に現物を添付します。",Inches(0.6),Inches(6.25),Inches(12),Inches(0.3),sz=10,col=GRY)
ft(s,6)

# S7 花みずき マイソク
s=sl(); hdr(s,"01 案件概要","② 花みずき貸家3棟 物件資料（マイソク・添付）","出所＝トラステン合同会社様 info-sheet。")
framed_pic(s,MYSO_H,Inches(1.5),Inches(1.95),Inches(10.3),Inches(4.4))
ft(s,7)

# S8 地番×所有者×持分
s=sl(); hdr(s,"01 案件概要 ★最重要","地番・所有者・持分の一覧（公図／固定資産税課税明細より）","売主は実質3者様。私道は売主様の単独所有（共有ではありません）。")
rows=[["地番","地目／用途","面積㎡","所有者","持分"],
      ["高松2-10-4","宅地","181.18","三浦 哲 様","単独"],
      ["高松2-10-5","宅地","94.39","三浦 哲 様","単独"],
      ["高松2-10-2","宅地","12.42","三浦 哲 様","単独"],
      ["高松2-10-10","宅地","0.57","三浦 哲 様","単独"],
      ["― アパート宅地 小計","","288.56","",""],
      ["高松2-10-1","公衆用道路(位置指定)","7.54","三浦 哲 様","100%"],
      ["高松2-10-11","公衆用道路(位置指定)","15.37","三浦 哲 様","100%"],
      ["高松2-10-12","公衆用道路(位置指定)","22.16","三浦 哲 様","100%"],
      ["高松2-11-8","公衆用道路(位置指定)","2.03","三浦 哲 様","100%"],
      ["― 私道 小計","","47.10","",""],
      ["高松2-10-6","宅地","312.59","末永 ちゑ子 様","単独"],
      ["高松2-11-4","宅地","45.57","末永 ちゑ子 様","単独"],
      ["― 戸建宅地 小計","","358.16","",""],
      ["合計（土地）","","693.82","",""]]
table(s,rows,Inches(0.55),Inches(1.78),Inches(8.2),Inches(4.75),
      [Inches(1.9),Inches(2.7),Inches(1.2),Inches(1.7),Inches(0.7)],
      sz=10.5,header_sz=11,sub_rows=[5,11,14],hi_rows=[6,7,8,9])
framed_pic(s,KOZ,Inches(9.05),Inches(1.85),Inches(3.5),Inches(4.4),"公図（高松2丁目）")
t(s,"※建物：アーバンキューブ＝有限会社三浦設計士様（三浦哲様と同一住所・同グループ）／花みずき3棟＝末永ちゑ子様。",
  Inches(0.55),Inches(6.6),Inches(8.3),Inches(0.5),sz=9,col=GRY,line_sp=1.05)
ft(s,8)

# S9 道路・接道＋申請図
s=sl(); hdr(s,"01 案件概要","道路・接道の状況（位置指定道路）","私道は売主様の単独所有。右図は位置指定道路の申請図（三浦設計士様 作成）です。")
card(s,Inches(0.6),Inches(1.8),Inches(5.0),Inches(2.0),"現状（課税明細で確認）",
     "・東側6.0m 私道（位置指定道路）に接道\n・私道4筆(計47.10㎡)は三浦哲様の単独所有・持分100%（公衆用道路）\n・複数地権者の共有ではなく、通行・掘削の同意取得リスクは小さい",accent=GRN)
card(s,Inches(0.6),Inches(3.95),Inches(5.0),Inches(2.5),"調査中（謄本取得・並行ヒアリング）",
     "・マイソク様式の「共有持分」表記と単独所有の経緯（いつ・どの登記で現状になったか）\n・私道に接続する“手前の道路区間”の地番・所有者\n・将来共有化していた場合の通行・掘削権、覚書・使用料の有無\n→ 私道全筆の登記簿謄本を取得し確定させます",accent=RED)
framed_pic(s,ROAD,Inches(6.0),Inches(1.8),Inches(6.6),Inches(4.4),"位置指定道路 申請図（三浦設計士様 作成）")
ft(s,9)

# S10 評価額・路線価・土地値（合算＋2区画）
s=sl(); hdr(s,"01 案件概要","評価額・路線価と土地値割合（合算／2区画別）","売却時はアパート用地・戸建用地の2区画に分けて評価。いずれも担保性は良好です。")
rows=[["区画","面積","固定資産税評価額","路線価相当(相続税)","摘要"],
      ["アパート用地（三浦哲様）","288.56㎡","9,651,754円","約 38,200円/㎡","建付地・現況OC想定"],
      ["戸建用地（末永様）","358.16㎡","13,004,788円","約 41,500円/㎡","退去時2区画分譲想定"],
      ["合計（私道は非課税）","646.72㎡","22,656,542円","―","実勢 約3,240〜3,560万"]]
table(s,rows,Inches(0.55),Inches(1.8),Inches(12.2),Inches(2.3),
      [Inches(3.4),Inches(1.7),Inches(2.8),Inches(2.5),Inches(1.8)],sz=10.5,header_sz=11,hi_rows=[3])
card(s,Inches(0.6),Inches(4.35),Inches(3.9),Inches(1.95),"土地値割合（合算）",None,big="約 90.5%")
t(s,"実勢土地値(中央約3,350万)÷取得3,700万。🟢安全圏。",Inches(0.8),Inches(5.55),Inches(3.5),Inches(0.6),sz=10.5,col=INK)
card(s,Inches(4.7),Inches(4.35),Inches(3.9),Inches(1.95),"計算根拠",None)
t(s,"固定資産税評価額 ÷ 0.7 ＝ 公示価格相当\n公示相当 ×1.0〜1.1（地方）＝ 実勢価格\n出所：令和8年度 固定資産税課税明細書",Inches(4.9),Inches(4.95),Inches(3.5),Inches(1.2),sz=10.5,col=INK,line_sp=1.2)
card(s,Inches(8.8),Inches(4.35),Inches(3.85),Inches(1.95),"分譲時の出口（戸建用地）",None)
t(s,"戸建用地358㎡を2区画(各約54坪)に分割。\n実需 約1,000万円×2＝約2,000万円を想定。",Inches(9.0),Inches(4.95),Inches(3.5),Inches(1.0),sz=10.5,col=INK,line_sp=1.2)
ft(s,10)

# S11 レントロール
s=sl(); hdr(s,"02 事業計画","レントロール（現況・出所：トラステン合同会社様 2026/3/30）","アパートは満室稼働。戸建は1室空室＝満室化でNOIが伸びます。")
rows=[["物件","戸数","入居","現況 月額","満室想定 月額","表面利回り"],
      ["アーバンキューブ","8戸","8/8","336千円","336千円","16.1%"],
      ["花みずき貸家","3棟","2/3","108千円","153千円","10.8 → 15.7%"],
      ["合計","11","10/11","444千円","489千円","13.0%（合算）"]]
table(s,rows,Inches(0.6),Inches(1.95),Inches(12.1),Inches(2.5),
      [Inches(3.0),Inches(1.3),Inches(1.3),Inches(2.0),Inches(2.3),Inches(2.2)],sz=12,header_sz=12,hi_rows=[3])
t(s,"・アパートは岩手大学生活協同組合様の管理で安定運営（管理費5%）。\n・戸建3棟は1室空室。満室化で年間賃料 約130万円→約184万円に増加し、NOIがさらに改善します。",
  Inches(0.7),Inches(4.8),Inches(12),Inches(1.0),sz=12.5,col=INK,line_sp=1.25)
ft(s,11)

# S12 市場
s=sl(); hdr(s,"02 事業計画","周辺の賃貸需給と立地","岩手大学・周辺施設による底堅い賃貸需要。学生〜単身者の実需が厚いエリアです。")
card(s,Inches(0.7),Inches(1.9),Inches(3.85),Inches(3.5),"需要の源泉","・いわて銀河鉄道「青山」駅圏\n・岩手大学に近く学生の単身需要が安定\n・大学生協様の入居付けルートを確保\n・1K中心＝景気変動に強い実需層")
card(s,Inches(4.75),Inches(1.9),Inches(3.85),Inches(3.5),"供給・競合","・周辺は築古木造1K/戸建が中心\n・新築供給は限定的\n・同エリアの中古木造は表面10〜14%で流通（次頁の成約事例）")
card(s,Inches(8.8),Inches(1.9),Inches(3.85),Inches(3.5),"本件の位置づけ","・取得利回り 表面13〜16%は\n　エリア上位水準\n・土地値が堅く賃料＋担保の二重の裏付け\n・空室1室は伸びしろ")
ft(s,12)

# S13 成約事例
s=sl(); hdr(s,"02 事業計画","盛岡市の収益物件 成約事例（市場の換価性）","同エリア・同型の木造OCアパートが現に成約。担保の換価性が高いことを示します。")
rows=[["成約事例","構造／戸数","築年","成約価格","表面利回り","成約時期"],
      ["★サンコーポ高松2（高松2丁目）","木造4戸","1983","1,380万円","14%（満室）","R2/12"],
      ["リバーサイドメイプル（高松4丁目）","木造6戸","1992","760万円","─","R7/2"],
      ["レジデンス深沢（松尾町・OC）","軽量鉄骨4戸","1982","2,000万円","10.5%","R3/4"],
      ["メゾン悠（三本柳）","木造1K6室","1998","1,450万円","─","R2/2"],
      ["パステルパレス（上堂）","RC","2003","5,250万円","7.3%","R5/4"]]
table(s,rows,Inches(0.55),Inches(1.95),Inches(12.2),Inches(3.0),
      [Inches(4.2),Inches(2.1),Inches(1.2),Inches(1.9),Inches(1.6),Inches(1.2)],sz=11,header_sz=11.5,hi_rows=[1])
t(s,"★同じ高松2丁目「サンコーポ高松2」（木造・OC・満室表面14%）が1,380万で成約済＝本件アーバンキューブとほぼ同型・同エリア・同利回り。\n　現況オーナーチェンジのまま売却できる市場が実在し、担保処分時の換価性も裏付けられます。",
  Inches(0.6),Inches(5.15),Inches(12.1),Inches(1.0),sz=11,col=REDD,bold=True,line_sp=1.2)
t(s,"出所：レインズ／盛岡市 収益物件 成約事例（2026/4 取得）。",Inches(0.6),Inches(6.25),Inches(10),Inches(0.3),sz=9,col=GRY)
ft(s,13)

# S14 事業方針
s=sl(); hdr(s,"02 事業計画","事業方針 ― 返済原資は長期の安定賃料","返済原資は賃料収入です。売却は将来の選択肢で、堅い土地値が下支えします。")
card(s,Inches(0.7),Inches(1.9),Inches(3.85),Inches(3.6),"主軸：長期保有・賃料返済","・本件は収益物件の長期保有が主軸\n・返済原資は安定賃料（DSCR1.40）\n・大学需要で高稼働、管理会社様へ委託\n・空室1室の解消でNOI改善",accent=GRN)
card(s,Inches(4.75),Inches(1.9),Inches(3.85),Inches(3.6),"将来の選択肢①：現況売却","・市況・資産入替の判断時に、入居中(OC)のまま投資家へ売却も可能\n・同エリアでOC物件が現に流通（前頁）\n・残債を返済し売却益を確定")
card(s,Inches(8.8),Inches(1.9),Inches(3.85),Inches(3.6),"将来の選択肢②：更地分割","・戸建の退去時は解体・更地化し戸建用地として実需へ分割売却\n・土地値が出口を下支え\n・出口を複線化しリスクを限定")
t(s,"※②③は出口の安全性（担保処分価値）を示すもので、短期転売を前提とした計画ではありません。返済はあくまで賃料収入で行います。",
  Inches(0.7),Inches(5.75),Inches(12),Inches(0.5),sz=10.5,col=GRY,line_sp=1.1)
ft(s,14)

# S15 DSCR計算根拠
s=sl(); hdr(s,"03 収支計画と融資条件","返済余力（DSCR）の計算根拠","初めてご覧の方にも追えるよう、NOI→返済→DSCRを段階で示します。")
rows=[["① 純収益（NOI）の算出","金額"],
      ["満室想定 年間賃料（AP4,032＋戸建満室1,884）","5,916千円"],
      ["▲ 空室・滞納損（約5%）","▲296千円"],
      ["▲ 運営費（管理5%・固都税・損保・修繕/原状回復 等）","▲1,620千円"],
      ["＝ 純収益 NOI","約4,000千円"]]
table(s,rows,Inches(0.6),Inches(1.85),Inches(6.6),Inches(2.6),[Inches(4.7),Inches(1.9)],sz=11,header_sz=11.5,hi_rows=[4])
rows2=[["② 返済とDSCR","値"],
       ["借入額","45,000千円"],
       ["金利／期間","2.5%／20年（元利均等）"],
       ["年間返済額","2,861千円"],
       ["DSCR ＝ NOI ÷ 年間返済","4,000 ÷ 2,861 ＝ 1.40"],
       ["年間CF ＝ NOI − 返済","＋1,139千円"]]
table(s,rows2,Inches(7.4),Inches(1.85),Inches(5.3),Inches(3.0),[Inches(2.5),Inches(2.8)],sz=11,header_sz=11.5,hi_rows=[4,5])
t(s,"◆DSCR（Debt Service Coverage Ratio）＝銀行が返済余力をみる代表指標。1.0で返済と純収益が均衡、1.2以上で安全圏とされます。本件は1.40。\n　本融資は長期保有・賃料返済型のため、短期転売の粗利率・IRRではなく、DSCR・実質利回り・土地値割合を主指標としています。",
  Inches(0.6),Inches(5.0),Inches(12.1),Inches(1.2),sz=11,col=INK,line_sp=1.25)
ft(s,15)

# S16 保有シナリオ
s=sl(); hdr(s,"03 収支計画と融資条件","保有シナリオ別の見込み（担保処分価値の裏付け）","賃料でCFを積みつつ、必要時は現況売却で残債を完済できます。")
rows=[["保有","累計CF","売却時残債","保証料返戻","売却後の手取り見込み"],
      ["5年保有 → 現況OC売却（3,700万）","+5,686千円","35,762千円","+2,453千円","約 +8,200千円"],
      ["10年保有 → 現況OC売却（3,700万）","+11,372千円","25,295千円","+1,338千円","約 +23,200千円"]]
table(s,rows,Inches(0.6),Inches(2.0),Inches(12.1),Inches(2.0),
      [Inches(4.4),Inches(2.0),Inches(2.2),Inches(1.9),Inches(1.6)],sz=11.5,header_sz=11.5,hi_col=4)
t(s,"※売却価格は現況（築年・賃料）横ばいの保守前提。繰上完済による信用保証料の未経過分返戻を加算。\n※あくまで担保処分時の換価可能性を示すものです。詳細は別紙 収支計算シート（玉川式・長期保有版）をご参照ください。",
  Inches(0.7),Inches(4.35),Inches(12),Inches(1.0),sz=11,col=GRY,line_sp=1.2)
ft(s,16)

# S17 総事業費・融資条件（新情報反映）
s=sl(); hdr(s,"03 収支計画と融資条件","総事業費と融資条件（貴行ご回答反映）","保証料は一括で借入金に内包、上乗せ金利なし。手出しはほぼ生じない設計です。")
rows=[["費目","金額"],
      ["物件取得費（消費税込・総額）","37,000千円"],
      ["諸費用（仲介手数料・登記・取得税・保険）","約2,327千円"],
      ["印紙税","40千円"],
      ["不動産担保調査手数料（55千×2件）","110千円"],
      ["信用保証料（①297＋②149・一括）","4,460千円"],
      ["総事業費 計","約43,937千円"]]
table(s,rows,Inches(0.6),Inches(1.8),Inches(7.0),Inches(3.3),[Inches(4.9),Inches(2.1)],sz=11,header_sz=12,hi_rows=[6])
rows2=[["融資条件","内容"],
       ["① アパート","30,000千円／20年／本物件担保"],
       ["② 貸家","15,000千円／20年／本物件担保"],
       ["借入合計","45,000千円（保証料を内包）"],
       ["金利","2.5%（保証協会の上乗せ金利なし）"]]
table(s,rows2,Inches(7.9),Inches(1.8),Inches(4.8),Inches(2.6),[Inches(1.8),Inches(3.0)],sz=11,header_sz=12,hi_rows=[3])
t(s,"◆前提のお願い（消費税）：売買代金37,000千円は消費税を含む総額とし、これを超える消費税等の別途負担は行わない前提で売買契約を締結します。建物（法人所有）分の消費税は税込価格に含まれることを契約書で明記します。\n◆借入45,000千円 ＞ 総事業費 約43,937千円 ＝ 自己持ち出しはほぼ発生しません。",
  Inches(0.6),Inches(5.2),Inches(12.1),Inches(1.2),sz=11,col=REDD,bold=True,line_sp=1.25)
ft(s,17)

# S18 スケジュール（右=詳細）
s=sl(); hdr(s,"04 スケジュール","全体スケジュール（予定）","左＝大枠の時期、右＝具体的な実務内容。")
steps=[("2026年6月","融資申込・条件確定","本資料＋収支計算シート・試算表(2〜4月)・前年同月比を提出。融資可否・金利・保証料・別途費用を確定。"),
       ("2026年6〜7月","売買契約・重説","トラステン合同会社様と売買契約締結。重説、私道通行・掘削権の承継特約、告知事項・消費税(税込)の確認。"),
       ("2026年7月","融資実行・決済","金消契約 → 融資実行 → 売買代金決済・所有権移転登記・抵当権設定。"),
       ("決済後〜","管理開始・長期運営","大学生協様/トラステン様で運営継続。空室1室の客付け、賃料でのCF積上げ・約定返済。"),
       ("5〜10年後","出口（選択）","資産入替判断時に現況OC売却、または戸建退去時に更地化・2区画分譲で残債完済。")]
y=2.0
for d,short,detail in steps:
    bx(s,Inches(0.7),Inches(y),Inches(2.3),Inches(0.78),CARD,CARDLN,1.0)
    t(s,d,Inches(0.72),Inches(y+0.07),Inches(2.26),Inches(0.32),sz=12,bold=True,col=RED,align=PP_ALIGN.CENTER)
    t(s,short,Inches(0.72),Inches(y+0.4),Inches(2.26),Inches(0.32),sz=11,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,detail,Inches(3.25),Inches(y+0.05),Inches(9.4),Inches(0.78),sz=11,col=INK,line_sp=1.1,anchor=MSO_ANCHOR.MIDDLE)
    y+=0.88
ft(s,18)

# S19 売却を見据えた必要書面チェックリスト（運用）
s=sl(); hdr(s,"04 資料管理","売却を見据えた“必要書面”チェックリスト（運用）","いつ退去が出ても売れる状態を今から整え、遠隔地でも安心して運営・売却できます。")
rows=[["売却・運営に必要な書面","入手状況","売却時に買主へ交付"],
      ["登記簿謄本・公図・地積測量図","公図○／測量図 未（自社で確定測量）","✓"],
      ["建築確認済証・検査済証（特にアパート）","交渉中（確認申請図とセットで依頼）","✓"],
      ["賃貸借契約書（全戸）・レントロール","レントロール○／契約書 未","✓"],
      ["告知書（事件事故・心理的瑕疵）","未取得","✓"],
      ["設備一覧・修繕履歴（外壁2023済）","一部のみ","✓"],
      ["私道の通行・掘削承諾書／覚書","未（謄本取得後に手当て）","✓"],
      ["越境の覚書・境界確認書","未","✓"],
      ["公課証明書・固定資産税評価証明","まもなく入手予定","✓"]]
table(s,rows,Inches(0.6),Inches(1.85),Inches(12.1),Inches(4.4),
      [Inches(5.3),Inches(4.6),Inches(2.2)],sz=10.5,header_sz=11.5,hi_col=2,hi_rows=[2,4])
t(s,"※本表は残務と連動し、入手のつど更新します。売却時交付物を先回りで揃える＝資産価値・換価性の維持につながります。",
  Inches(0.6),Inches(6.3),Inches(12),Inches(0.3),sz=9.5,col=GRY)
ft(s,19)

# S20 代表挨拶
s=sl(); hdr(s,"05 ごあいさつ","経営方針 ― フローとストックの両立")
body=("このたびは弊社の事業計画をご高覧賜り、誠にありがとうございます。\n\n"
      "弊社は、独自の物件検索・AI査定により割安な不動産を取得し、加工・再販する「フロー収益」を進めてまいりました。\n\n"
      "本件は、土地値の堅い収益物件を取得し、安定賃料による「ストック収益」を積み上げる中長期投資です。市況に左右されにくい賃料収入で財務体質を強化し、フローとストックの両立による持続的成長の基盤といたします。\n\n"
      "御行の良きパートナーとして、誠実に・末永くお取引させていただきたく、引き続きのご支援を賜りますようお願い申し上げます。")
t(s,body,Inches(0.9),Inches(1.95),Inches(11.5),Inches(3.8),sz=14,col=INK,line_sp=1.4)
t(s,ENT,Inches(0.9),Inches(6.05),Inches(11),Inches(0.4),sz=15,bold=True,col=REDD)
ft(s,20)

# S21 別紙：調査中事項
s=sl(); hdr(s,"05 別紙","現在調査中の事項・残論点（透明性のためのご開示）","査定精度を上げるため、以下を仲介・市・売主・銀行へ確認中です。")
rows=[["区分","調査中の事項","収支／担保への影響"],
      ["耐震","AP=1997(新耐震)／花みずき=1968(旧耐震)","旧耐震は融資・出口に影響。要評価"],
      ["排水","花みずき脇の側溝の用途(雨水/汚水・公共/私設)","垂れ流しでないか・更地売りの是正費"],
      ["私設管","位置指定道路内：水道のほかガス・下水の有無、使用料・覚書・徴収方法","インフラ負担・近隣関係"],
      ["ガス","プロパン無償貸与契約・解約精算金（確認中）","解約時の精算リスク"],
      ["瑕疵","告知事項(事件事故・心理的瑕疵)","出口価格に直結。未取得"],
      ["建物","確認済証・検査済証／確認申請図","再建築・売却性。AP分は交渉中"],
      ["土地(最重要)","高低差・地中埋設物・越境・分合筆要否","“見えない”最大リスク。要現地・測量"],
      ["書面","確定測量図・通行掘削承諾書・越境覚書・建物プラン作図可否・建築見積(外構/インフラ込)","出口の実現性"],
      ["設備","エアコン・プロパン・給湯器の所有/状態・修繕履歴","引継ぎ負担・修繕引当"],
      ["賃貸","賃貸借契約書・入居者属性","立退き/承継の確度"],
      ["消費税","売買代金が税込総額か(建物=法人所有)","手出し増の回避（要契約明記）"]]
table(s,rows,Inches(0.55),Inches(1.78),Inches(12.2),Inches(4.7),
      [Inches(1.5),Inches(6.4),Inches(4.3)],sz=9.5,header_sz=11,hi_rows=[1,7,12])
ft(s,21)

# ════════ 社内管理用スライド（銀行提出時は削除）════════
# S22 解体見積（添付）
s=sl(internal=True); hdr(s,"社内管理用","解体見積（花みずき3棟）と確認事項",
    "東日産業様 御見積 4,246,000円（税込）。3棟分(延床163.92㎡)。アスベスト除去は明示行なし＝要確認。",top=0.55)
framed_pic(s,KAITAI,Inches(0.8),Inches(2.0),Inches(5.2),Inches(4.2),"東日産業様 見積内訳")
card(s,Inches(6.4),Inches(2.0),Inches(6.2),Inches(2.0),"確認事項（更地化コスト）",
     "・本見積は花みずき3棟(10-6他・延床163.92㎡)の解体＝3棟分で正。\n・内外の残置物・地中埋設物は『別途』と明記＝見積外（追加費用）。\n・アスベストの記載なし＝別途/要事前調査。\n・相見積（2〜3社）を取得し総額を確定。")
card(s,Inches(6.4),Inches(4.15),Inches(6.2),Inches(2.05),"更地分譲の試算（社内）",
     "戸建用地358㎡→2区画分譲：売上 約2,000万(実需1,000万×2)。\n−解体600 −測量50 −立退100 −売却費60 ＝ 粗利ほぼゼロ。\n→ 更地転売は出口の保険。本線はあくまで長期保有。")
ft(s,22)

# S23 収支 全結果（社内）
s=sl(internal=True); hdr(s,"社内管理用","収支計算 全結果（玉川式・長期保有版）",
    "別紙Excel『277_長期シミュレーション_合わせ込み』の主要アウトプット。金利感応度つき。",top=0.55)
rows=[["金利","年間返済","年間CF(税前)","DSCR","CCR/CF率","5年後残債"],
      ["2.5%","2,861千円","+1,137千円","1.40","フルローン","35,762千円"],
      ["3.0%","2,981千円","+1,004千円","1.34","─","─"],
      ["3.5%","3,113千円","+867千円","1.28","─","─"]]
table(s,rows,Inches(0.6),Inches(2.0),Inches(12.1),Inches(2.0),
      [Inches(1.6),Inches(2.4),Inches(2.6),Inches(1.6),Inches(2.1),Inches(2.4)],sz=11,header_sz=11,hi_rows=[1])
rows2=[["短期(参考・銀行非提出)","結果"],
       ["現況一括転売(〜1年)","残債>売価で ▲約708万（不可）"],
       ["更地転売・エンド実需2,000万","粗利 ほぼゼロ"],
       ["更地転売・ビルダー1,600万","粗利 ▲約398万（赤字）"]]
table(s,rows2,Inches(0.6),Inches(4.3),Inches(7.2),Inches(2.0),[Inches(4.2),Inches(3.0)],sz=11,header_sz=11,hi_rows=[])
card(s,Inches(8.1),Inches(4.3),Inches(4.6),Inches(2.0),"結論（社内）",
     "短期は全て不利＝長期保有CF＋現況OC売りが最適解。\nフルローンで手出しゼロ＝ダウンサイド限定。\n最大変数は金利上昇→5年前後で現況OC売り前倒しも検討。")
ft(s,23)

# S24 末永様への依頼（社内）
s=sl(internal=True); hdr(s,"社内管理用","末永ちゑ子様への依頼（立退き・承継の戦略）",
    "戸建退去をスムーズに進めるための具体的な依頼内容（売主様の協力を引き出す）。",top=0.55)
card(s,Inches(0.7),Inches(2.0),Inches(5.9),Inches(4.0),"依頼すること（具体化）",
     "①現入居者2世帯との関係・連絡先・人柄の引継ぎ（紹介の一筆）\n②これまでの賃貸条件・口約束・滞納/トラブル有無の開示\n③立退き交渉時の“顔つなぎ”同席 or 紹介状\n④定期借家への切替 or 退去合意のタイミング相談\n⑤過去の修繕・設備更新の記憶（履歴の代替）\n⑥告知事項（事件事故等）の確認")
card(s,Inches(6.8),Inches(2.0),Inches(5.85),Inches(4.0),"狙い・進め方",
     "・売主様は入居者と長い関係＝信頼の“のれん”を承継すれば立退きが円滑。\n・いきなり立退き条件でなく、まず関係承継＋情報開示をGIVEとして依頼。\n・トラステン合同会社様を窓口に、売主様の負担が軽い形で協力を引き出す。\n・退去合意が取れた住戸から更地化・分譲の段取りへ。",accent=GRN)
ft(s,24)

# S25 補足・抜け漏れメモ（社内）
s=sl(internal=True); hdr(s,"社内管理用","補足・抜け漏れメモ（随時追記）",
    "気づいた論点をスライド単位で蓄積。確認が取れ次第、本編へ反映。",top=0.55)
rows=[["#","論点","状態／次アクション"],
      ["1","ガス無償貸与の解約精算金（プロパン）","メールで確認中"],
      ["2","公課証明書・評価証明の入手","まもなく入手予定"],
      ["3","賃貸借契約書（全戸）・敷金・特約","トラステン様へ依頼"],
      ["4","入居者属性（学生/社会人・滞納履歴）","トラステン様へ依頼"],
      ["5","エアコン・給湯器・設備の所有/リース","売主様へ確認（難航可能性）"],
      ["6","建物プラン作図可否（戸建2戸 or 新築AP）","設計へ打診"],
      ["7","建築見積（外構・インフラ工事込）","出口確度を上げるため取得"],
      ["8","金利上昇対策（固定/一部固定の可否）","西原様へ相談"]]
table(s,rows,Inches(0.6),Inches(2.0),Inches(12.1),Inches(4.2),
      [Inches(0.7),Inches(6.6),Inches(4.8)],sz=11,header_sz=11.5)
ft(s,25)

prs.save("277_bank.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")

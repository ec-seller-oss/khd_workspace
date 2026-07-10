"""
クリニックDXセミナー資料 "勝てる版" v2 を生成する。
台本: 【改訂台本v2】クリニックDXセミナー資料（Google Doc / 2026-05-26）に準拠。
設計: ①数字を自分ごと化 ②Before→Afterを情景で ③第三者実証データを出典付きで誠実に引用
MyAIのUI画面は図形でモック化。実写が要る箇所は画像差込枠を配置。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── カラーパレット（医療×テックの清潔感）─────────────
NAVY   = RGBColor(0x0A, 0x1F, 0x33)   # 表紙・章扉の背景
NAVY2  = RGBColor(0x10, 0x2A, 0x43)   # ヘッダーバー
LIGHT  = RGBColor(0xF5, 0xF8, 0xFB)   # コンテンツページ背景
CARD   = RGBColor(0xFF, 0xFF, 0xFF)   # カード白
CARDLN = RGBColor(0xDD, 0xE6, 0xEE)   # カード枠線
TEAL   = RGBColor(0x12, 0xA5, 0x9B)   # メインアクセント
TEALD  = RGBColor(0x0C, 0x7A, 0x73)   # 濃ティール
GREEN  = RGBColor(0x1E, 0xB8, 0x6B)   # 成果・院に残る効果
GREEND = RGBColor(0x14, 0x8A, 0x50)
RED    = RGBColor(0xE0, 0x4F, 0x5F)   # 課題・リスク
AMBER  = RGBColor(0xF0, 0xA8, 0x20)
INK    = RGBColor(0x1A, 0x2B, 0x3A)   # 本文濃
GRAY   = RGBColor(0x5B, 0x6B, 0x7B)   # サブ
LGRAY  = RGBColor(0xEC, 0xF1, 0xF5)   # 薄背景
BLUEBG = RGBColor(0xE7, 0xF1, 0xFA)
GREENBG= RGBColor(0xE6, 0xF6, 0xEC)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Hiragino Sans"
_MK = "/Users/kikuchikenta/01_honbu_docs_automation/myai_mockups"
IMG1 = _MK + "/screen1_doc.png"     # 書類自動生成
IMG2 = _MK + "/screen2_dash.png"    # ダッシュボード
IMG3 = _MK + "/screen3_line.png"    # スマホLINE問診
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = col
        r.font.name = FONT
    return tb


def rect(slide, x, y, w, h, fill, line=None, lw=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def card(slide, x, y, w, h, fill=CARD, line=CARDLN, lw=1.0):
    return rect(slide, x, y, w, h, fill, line, lw, rounded=True)


def hdr(slide, main, sub=""):
    rect(slide, 0, 0, W, Inches(1.15), NAVY2)
    rect(slide, 0, Inches(1.15), W, Pt(3), TEAL)
    t(slide, main, Inches(0.6), Inches(0.20), Inches(12.1), Inches(0.6),
      sz=25, bold=True, col=WHT)
    if sub:
        t(slide, sub, Inches(0.62), Inches(0.78), Inches(12.1), Inches(0.32),
          sz=12, col=RGBColor(0xB9, 0xCB, 0xD9))


def chip(slide, text, x, y, w, h, fill, col=WHT, sz=12, bold=True):
    c = rect(slide, x, y, w, h, fill, rounded=True)
    t(slide, text, x, y, w, h, sz=sz, bold=bold, col=col,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def style_table(tbl, header_fill=NAVY2, header_col=WHT, body_fill=WHT,
                alt_fill=LGRAY, sz=12, header_sz=12):
    n_rows = len(tbl.rows)
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.fore_color.rgb = body_fill if ri % 2 == 1 else alt_fill
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0)
                    r.font.color.rgb = header_col if ri == 0 else INK


def set_cell(cell, text, col=INK, bold=False, sz=None, fill=None):
    cell.text = text
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = FONT
            r.font.color.rgb = col
            r.font.bold = bold
            if sz:
                r.font.size = Pt(sz)


def myai_mock(slide, x, y, w, h):
    """MyAIのUIモック：紹介状をAIが自動下書きする画面イメージ。"""
    rect(slide, x, y, w, h, RGBColor(0x16, 0x33, 0x4D), line=TEAL, lw=1.5, rounded=True)
    # 上部バー
    rect(slide, x+Inches(0.12), y+Inches(0.12), w-Inches(0.24), Inches(0.42),
         TEAL, rounded=True)
    t(slide, "マイAI ｜ 紹介状ジェネレーター", x+Inches(0.28), y+Inches(0.13),
      w-Inches(0.5), Inches(0.4), sz=11, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
    # 入力カード
    rect(slide, x+Inches(0.22), y+Inches(0.68), w-Inches(0.44), Inches(0.55),
         RGBColor(0x21, 0x44, 0x60), rounded=True)
    t(slide, "カルテから抽出中… 患者ID / 主訴 / 既往歴 / 処方",
      x+Inches(0.36), y+Inches(0.70), w-Inches(0.6), Inches(0.5),
      sz=9.5, col=RGBColor(0xBE, 0xD6, 0xE6), anchor=MSO_ANCHOR.MIDDLE)
    # 出力プレビュー
    rect(slide, x+Inches(0.22), y+Inches(1.34), w-Inches(0.44), h-Inches(2.05),
         WHT, rounded=True)
    t(slide, "AI生成プレビュー", x+Inches(0.38), y+Inches(1.42),
      w-Inches(0.6), Inches(0.3), sz=9, bold=True, col=TEALD)
    for i, ln in enumerate(["─────────────────────",
                            "拝啓  時下ますますご清栄のことと…",
                            "下記の患者をご紹介申し上げます。",
                            "診断: ＿＿＿  /  所見: ＿＿＿＿＿",
                            "─────────────────────"]):
        t(slide, ln, x+Inches(0.38), y+Inches(1.74)+Emu(int(Inches(0.30))*i),
          w-Inches(0.7), Inches(0.3), sz=9, col=GRAY)
    # 生成ボタン
    rect(slide, x+Inches(0.22), y+h-Inches(0.62), w-Inches(0.44), Inches(0.44),
         GREEN, rounded=True)
    t(slide, "⚡ 8分 → 数秒で下書き完了", x+Inches(0.22), y+h-Inches(0.62),
      w-Inches(0.44), Inches(0.44), sz=11, bold=True, col=WHT,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def photo_ph(slide, x, y, w, h, label):
    """実写画像の差込枠（朝、写真をはめる用）。"""
    rect(slide, x, y, w, h, LGRAY, line=TEAL, lw=1.25, rounded=True)
    t(slide, "🖼  画像差込\n" + label, x, y, w, h, sz=12, bold=True, col=TEALD,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.2)


# ════════════════════════════════════════════════════
# SLIDE 1 — 表紙
# ════════════════════════════════════════════════════
s = sl(NAVY)
rect(s, 0, 0, Inches(0.28), H, TEAL)
chip(s, "次世代クリニック DX", Inches(0.9), Inches(0.7), Inches(3.0), Inches(0.45),
     TEALD, sz=13)
t(s, '"人を増やさず"\n診療の質を上げる。', Inches(0.9), Inches(1.35),
  Inches(7.4), Inches(1.9), sz=44, bold=True, col=WHT, line_sp=1.05)
t(s, 'クリニックDX「My AI」  ※マイAIは仮称', Inches(0.92), Inches(3.35),
  Inches(7.4), Inches(0.4), sz=15, col=TEAL, bold=True)
t(s, "Lステップ ×「マイAI」× AI音声(IVR) で実現する、リスクゼロのDX",
  Inches(0.92), Inches(3.85), Inches(7.4), Inches(0.5), sz=14, col=RGBColor(0xC4,0xD3,0xDF))
# オファー3カード
ox, oy, ow, og = Inches(0.9), Inches(4.7), Inches(2.35), Inches(0.18)
offers = [("初期費用", "0円"), ("月額基本料", "0円"), ("成果報酬", "30%")]
for i, (lab, val) in enumerate(offers):
    cx = ox + (ow + og) * i
    card(s, cx, oy, ow, Inches(1.35), CARD, CARDLN)
    t(s, lab, cx, oy+Inches(0.16), ow, Inches(0.35), sz=12, col=GRAY,
      align=PP_ALIGN.CENTER)
    t(s, val, cx, oy+Inches(0.48), ow, Inches(0.7), sz=34, bold=True, col=TEALD,
      align=PP_ALIGN.CENTER)
t(s, "成果報酬は「実際に人件費相当額を削減できた分」だけ／医療現場の実証データに基づく設計／IT導入補助金サポート対応",
  Inches(0.9), Inches(6.35), Inches(7.4), Inches(0.7), sz=10.5, col=RGBColor(0x9F,0xB4,0xC4))
# 右：MyAI実画面（書類生成）をデバイス風カードに
card(s, Inches(8.45), Inches(2.3), Inches(4.5), Inches(3.1), WHT, CARDLN)
s.shapes.add_picture(IMG1, Inches(8.65), Inches(2.72), width=Inches(4.1))
t(s, "実際のMyAI画面：カルテ→紹介状を数秒で自動下書き",
  Inches(8.45), Inches(5.48), Inches(4.5), Inches(0.5), sz=10.5,
  col=RGBColor(0x9F, 0xB4, 0xC4), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 2 — 4大課題
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "現場はもう限界では？ ── これ、御院でも起きていませんか？",
    "ネット上の生の声が示す、クリニックのリアルな4大課題")
issues = [
    ("①", "採用しても、定着しない", "70万〜150万円", "医療事務の採用単価。なのに保険制度の複雑さで即戦力にならず、早期離職の悪循環。"),
    ("②", "カルテ入力で、患者と目が合わない", "診察時間の半分", "が電子カルテ入力。診療後も紹介状・診断書作成で2時間超の残業。"),
    ("③", "月初10日のレセプトが、毎月重い", "月10〜20時間", "の残業が常態化。目視チェックによる算定漏れの不安も重なる。"),
    ("④", "鳴り止まない電話で、受付がパンク", "月1,200件超", "の電話対応。「今空いてますか？」で目の前の患者対応が滞る。"),
]
cw, ch = Inches(5.95), Inches(2.45)
gx, gy = Inches(0.55), Inches(0.5)
x0, y0 = Inches(0.55), Inches(1.55)
for i, (no, ti, num, desc) in enumerate(issues):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    card(s, cx, cy, cw, ch)
    rect(s, cx, cy, Inches(0.12), ch, RED, rounded=False)
    t(s, no, cx+Inches(0.3), cy+Inches(0.2), Inches(0.8), Inches(0.6),
      sz=30, bold=True, col=RED)
    t(s, ti, cx+Inches(1.1), cy+Inches(0.28), cw-Inches(1.3), Inches(0.6),
      sz=17, bold=True, col=INK)
    t(s, num, cx+Inches(0.35), cy+Inches(1.05), cw-Inches(0.6), Inches(0.7),
      sz=30, bold=True, col=RED)
    t(s, desc, cx+Inches(0.35), cy+Inches(1.75), cw-Inches(0.6), Inches(0.6),
      sz=12, col=GRAY, line_sp=1.1)

# ════════════════════════════════════════════════════
# SLIDE 3 — 自動化マップ
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "既存の業務をどう変えるか？ ── 完全自動化マップ",
    "個別ツールではなく、クリニック全体を1つのエコシステムで最適化する")
rows = [
    ("業務領域", "導入前（アナログ）", "導入後", "削減効果"),
    ("問診・予約", "電話受付＋紙問診をカルテへ手入力", "Lステップで来院前にスマホ上で完結", "問診業務 40〜70%削減"),
    ("電話対応", "全件スタッフが受話", "AI音声(IVR)が一次対応・LINE誘導", "電話件数 最大40〜80%削減"),
    ("書類作成", "過去カルテを探しつつ手作業", "自院学習「マイAI」が下書き自動生成", "書類時間 月30時間以上削減"),
    ("レセプト点検", "月末月初に目視で残業", "AIが点検ルールで算定漏れ自動チェック", "点検時間 最大1/20に短縮"),
]
tb = s.shapes.add_table(len(rows), 4, Inches(0.55), Inches(1.5),
                        Inches(12.23), Inches(5.2)).table
tb.columns[0].width = Inches(2.1)
tb.columns[1].width = Inches(3.9)
tb.columns[2].width = Inches(3.9)
tb.columns[3].width = Inches(2.33)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        tb.cell(ri, ci).text = val
style_table(tb, sz=13, header_sz=13)
# 削減効果列を緑強調
for ri in range(1, len(rows)):
    set_cell(tb.cell(ri, 3), rows[ri][3], col=GREEND, bold=True, sz=12.5,
             fill=GREENBG)
    set_cell(tb.cell(ri, 0), rows[ri][0], col=NAVY2, bold=True, sz=13)

# ════════════════════════════════════════════════════
# SLIDE 3.5 — 実際のMyAI画面（ソリューションの直後）
# ════════════════════════════════════════════════════
s = sl()
hdr(s, '"使う場面"が見える ── 実際のMyAI画面',
    "院長・スタッフ・患者、それぞれの画面で業務が変わる（画面はデモ・数値は例示）")
# 左：書類生成（大）
card(s, Inches(0.55), Inches(1.5), Inches(7.15), Inches(4.85))
s.shapes.add_picture(IMG1, Inches(0.8), Inches(1.85), width=Inches(6.5))
t(s, "① 書類アシスト｜カルテ→紹介状を数秒で下書き（院長の画面）",
  Inches(0.8), Inches(6.0), Inches(6.7), Inches(0.35), sz=12, bold=True, col=NAVY2)
# 右上：ダッシュボード
card(s, Inches(7.95), Inches(1.5), Inches(4.83), Inches(2.86))
s.shapes.add_picture(IMG2, Inches(8.42), Inches(1.62), width=Inches(3.9))
t(s, "② 成果ダッシュボード｜削減効果を可視化", Inches(7.95), Inches(4.05),
  Inches(4.83), Inches(0.3), sz=11, bold=True, col=NAVY2, align=PP_ALIGN.CENTER)
# 右下：スマホLINE
card(s, Inches(7.95), Inches(4.5), Inches(4.83), Inches(1.85))
s.shapes.add_picture(IMG3, Inches(8.3), Inches(4.6), height=Inches(1.65))
t(s, "③ LINE問診・予約\nLステップ連携（患者スマホ）", Inches(9.35), Inches(4.95),
  Inches(3.3), Inches(1.0), sz=11.5, bold=True, col=NAVY2, line_sp=1.2)

# ════════════════════════════════════════════════════
# SLIDE 4 — 実証データ3事例
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "「本当に自院でも変わるのか？」を証明する実証データ",
    "※以下は当社実績ではなく、公開されている医療機関の導入実績です（出典明記）")
cases = [
    ("問診・予約のDX", "内科クリニック",
     "電話の嵐＋紙問診の手入力で受付過多",
     "Lステップで来院前にスマホ問診完了。電話を一次自動対応",
     ["初月にLINE予約 350件", "受付問診 70%削減", "月商 約1,500万円規模を押上げ"],
     "出典: Lステップ公開事例「そのだ内科」ほか"),
    ("書類作成のDX", "病院・総合クリニック",
     "紹介状・退院サマリ作成に医師が忙殺",
     "自院データを学習した生成AIが下書き自動生成",
     ["診断書 月400件を 50%削減", "退院サマリ 28分→8分(7割減)", "医師業務 月30h以上削減"],
     "出典: 新古賀/名古屋医療センター/戸畑共立"),
    ("電話業務のDX", "地域密着クリニック",
     "インフル予約・時間外問合せで電話過多",
     "AI音声が一次対応、営業電話を自動ブロック",
     ["月1,200件の電話を自動化", "うち約400件(40%)を削減", "対応時間 全体70%削減"],
     "出典: IVRy 公開導入事例"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.5)
for i, (cat, who, bf, af, results, src) in enumerate(cases):
    cx = x0 + (cw + gx) * i
    card(s, cx, y0, cw, Inches(5.25))
    chip(s, cat, cx+Inches(0.2), y0+Inches(0.2), cw-Inches(0.4), Inches(0.42),
         TEALD, sz=12.5)
    t(s, who, cx+Inches(0.2), y0+Inches(0.7), cw-Inches(0.4), Inches(0.3),
      sz=11, bold=True, col=GRAY, align=PP_ALIGN.CENTER)
    # Before
    rect(s, cx+Inches(0.2), y0+Inches(1.08), cw-Inches(0.4), Inches(0.85),
         LGRAY, rounded=True)
    t(s, "Before", cx+Inches(0.32), y0+Inches(1.12), Inches(1.2), Inches(0.28),
      sz=9.5, bold=True, col=GRAY)
    t(s, bf, cx+Inches(0.32), y0+Inches(1.36), cw-Inches(0.64), Inches(0.55),
      sz=10.5, col=INK, line_sp=1.05)
    # After
    rect(s, cx+Inches(0.2), y0+Inches(2.0), cw-Inches(0.4), Inches(0.85),
         BLUEBG, rounded=True)
    t(s, "After", cx+Inches(0.32), y0+Inches(2.04), Inches(1.2), Inches(0.28),
      sz=9.5, bold=True, col=RGBColor(0x2A,0x6E,0xA8))
    t(s, af, cx+Inches(0.32), y0+Inches(2.28), cw-Inches(0.64), Inches(0.55),
      sz=10.5, col=INK, line_sp=1.05)
    # 成果
    rect(s, cx+Inches(0.2), y0+Inches(2.92), cw-Inches(0.4), Inches(1.75),
         GREENBG, rounded=True)
    t(s, "成果", cx+Inches(0.32), y0+Inches(2.96), Inches(1.2), Inches(0.3),
      sz=10, bold=True, col=GREEND)
    for j, rr in enumerate(results):
        t(s, "● " + rr, cx+Inches(0.32), y0+Inches(3.28)+Emu(int(Inches(0.43))*j),
          cw-Inches(0.6), Inches(0.42), sz=12, bold=True, col=GREEND, line_sp=1.0)
    t(s, src, cx+Inches(0.2), y0+Inches(4.8), cw-Inches(0.4), Inches(0.4),
      sz=8.5, col=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 5 — ROI早見表
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "御院の規模なら、いくら残る？ ── 削減効果の早見表",
    "1モデルではなく、自院を当てはめられる規模別シミュレーション")
rows = [
    ("院の規模", "月の削減額", "成果報酬(30%)", "院に残る効果（年額）"),
    ("医師1・事務2", "約16万円", "4.8万円", "約134万円"),
    ("医師1・事務3（標準）", "約24万円", "7.2万円", "約202万円"),
    ("医師2・事務5", "約40万円", "12万円", "約336万円"),
]
tb = s.shapes.add_table(len(rows), 4, Inches(0.8), Inches(1.65),
                        Inches(11.73), Inches(2.9)).table
tb.columns[0].width = Inches(3.4)
tb.columns[1].width = Inches(2.5)
tb.columns[2].width = Inches(2.6)
tb.columns[3].width = Inches(3.23)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        tb.cell(ri, ci).text = val
style_table(tb, sz=16, header_sz=14)
for ri in range(1, len(rows)):
    set_cell(tb.cell(ri, 3), rows[ri][3], col=GREEND, bold=True, sz=21,
             fill=GREENBG)
set_cell(tb.cell(2, 0), rows[2][0], col=NAVY2, bold=True, sz=16, fill=RGBColor(0xE2,0xF0,0xEE))
# 内訳の補足
card(s, Inches(0.8), Inches(4.85), Inches(11.73), Inches(2.05), RGBColor(0xF0,0xF5,0xF9), CARDLN)
t(s, "標準モデルの内訳（医師1・事務3 ／ 事務人件費 月83万円・総480時間・時間単価 約1,729円）",
  Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.4), sz=12.5, bold=True, col=NAVY2)
t(s, "問診 −42h ＋ レセプト −45h ＋ 電話 −40h ＋ 書類 −12h  ＝  月139時間削減（フルタイム約1名分の余力創出）",
  Inches(1.0), Inches(5.45), Inches(11.3), Inches(0.4), sz=13, col=INK)
t(s, "→ 月139h × 1,729円 ＝ 月 約24万円削減 → 年 約288万円。うち成果報酬を引いて、院に残る効果は 約202万円/年。",
  Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.4), sz=13, bold=True, col=GREEND)
t(s, "※「うちは事務3人だから…年200万か」と、その場で自院に当てはめられます。",
  Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.4), sz=11, col=GRAY)

# ════════════════════════════════════════════════════
# SLIDE 6 — 競合比較表
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "他社に頼むと「先行投資リスク」を抱える ── 業界価格の真実",
    "他社は全部、効果が出なくても固定費が発生。当サービスだけ唯一のノーリスク")
rows = [
    ("比較項目", "一般ITベンダー", "大手DXコンサル", "Lステップ構築代行", "当サービス（中立型伴走）"),
    ("提供内容", "自社アプリ単体販売", "業務分析・高額導入支援", "LINE・Lステップ構築", "Lステップ×マイAI 全体最適化"),
    ("初期費用", "20万〜数十万", "100〜300万", "50〜100万超", "0円"),
    ("月額費用", "1.5〜5万(固定)", "50〜300万(固定)", "5〜20万(固定)", "0円（基本料なし）"),
    ("報酬体系", "稼働問わず固定", "稼働ベース固定", "固定保守料", "削減できた人件費額の30%のみ"),
    ("運用定着", "自力運用", "別料金で高額", "現場運用に不慣れ", "使いこなすまで3ヶ月伴走"),
    ("院側のリスク", "高", "極めて高", "中", "ゼロ（効果0なら支払い0）"),
]
tb = s.shapes.add_table(len(rows), 5, Inches(0.45), Inches(1.45),
                        Inches(12.43), Inches(5.35)).table
tb.columns[0].width = Inches(1.95)
tb.columns[1].width = Inches(2.35)
tb.columns[2].width = Inches(2.35)
tb.columns[3].width = Inches(2.35)
tb.columns[4].width = Inches(3.43)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        tb.cell(ri, ci).text = val
style_table(tb, sz=11.5, header_sz=11.5)
# 当サービス列(4)を強調
for ri in range(len(rows)):
    if ri == 0:
        set_cell(tb.cell(ri, 4), rows[ri][4], col=WHT, bold=True, sz=12, fill=GREEND)
    else:
        set_cell(tb.cell(ri, 4), rows[ri][4], col=GREEND, bold=True, sz=12, fill=GREENBG)
    set_cell(tb.cell(ri, 0), rows[ri][0], col=(WHT if ri == 0 else NAVY2), bold=True, sz=11.5,
             fill=(NAVY2 if ri == 0 else RGBColor(0xEC,0xF1,0xF5)))

# ════════════════════════════════════════════════════
# SLIDE 7 — 成果報酬の3ステップ
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "「うちでも本当に測定できる？」にお答えする、リスクゼロの3ステップ",
    "成果報酬への最大の懸念=「測定の手間」は、コンサル側が巻き取ります")
steps = [
    ("STEP 1", "事前の無料コンサル", "貴院の負担ゼロ",
     "当方が競合分析と類似事例を持ち込み、業務フローとボトルネックを洗い出し。\n\n※「2週間時間を計測する」等の煩雑な作業は一切要求しません。", NAVY2, None),
    ("STEP 2", "フェアな測定ルール", "双方納得の透明基準",
     "導入前後で同等の患者数帯に基づき削減時間を算出。人員変更月は除外。\n\n※最低契約期間以降は30日前通知で解約可。ロックインなし。", TEALD, None),
    ("STEP 3", "手元に利益が残る", "リスクは私たちが取る",
     "【例：月24万円分を削減】\n削減額24万 × 30% ＝ お支払い 7.2万円", GREEND,
     ("貴院に残る効果", "月 16.8万円", "（年 約200万円）／効果0なら支払い0円")),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.55)
for i, (st, ti, sub, body, accent, hi) in enumerate(steps):
    cx = x0 + (cw + gx) * i
    card(s, cx, y0, cw, Inches(5.2))
    rect(s, cx, y0, cw, Inches(0.95), accent, rounded=True)
    rect(s, cx, y0+Inches(0.55), cw, Inches(0.4), accent)  # 角丸下を四角で隠す
    t(s, st, cx, y0+Inches(0.1), cw, Inches(0.45), sz=20, bold=True, col=WHT,
      align=PP_ALIGN.CENTER)
    t(s, sub, cx, y0+Inches(0.56), cw, Inches(0.32), sz=11, col=RGBColor(0xE8,0xF3,0xF1),
      align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.25), y0+Inches(1.15), cw-Inches(0.5), Inches(0.55),
      sz=17, bold=True, col=INK, align=PP_ALIGN.CENTER)
    t(s, body, cx+Inches(0.3), y0+Inches(1.8), cw-Inches(0.6), Inches(2.0),
      sz=11.5, col=INK, line_sp=1.15)
    if hi:
        rect(s, cx+Inches(0.25), y0+Inches(3.45), cw-Inches(0.5), Inches(1.5),
             GREENBG, line=GREEN, lw=1.5, rounded=True)
        t(s, hi[0], cx+Inches(0.25), y0+Inches(3.55), cw-Inches(0.5), Inches(0.3),
          sz=11, bold=True, col=GREEND, align=PP_ALIGN.CENTER)
        t(s, hi[1], cx+Inches(0.25), y0+Inches(3.85), cw-Inches(0.5), Inches(0.6),
          sz=27, bold=True, col=GREEND, align=PP_ALIGN.CENTER)
        t(s, hi[2], cx+Inches(0.3), y0+Inches(4.5), cw-Inches(0.6), Inches(0.4),
          sz=10, col=GRAY, align=PP_ALIGN.CENTER, line_sp=1.0)

# ════════════════════════════════════════════════════
# SLIDE 8 — 選ばれる理由
# ════════════════════════════════════════════════════
s = sl()
hdr(s, "なぜ私たちは「完全成果報酬」で伴走できるのか？",
    "「なぜ他社は固定費で、御社は成果報酬でやれるのか？」への答え")
reasons = [
    ("中立選定", "特定メーカーに縛られない",
     "自社ツールの売り込みではない。Lステップ/マイAI/IVRから貴院の診療科・規模に最適な組合せを厳選。", TEALD),
    ("実践ノウハウ", "自社で泥臭く運用した知見",
     "ツールは「設計」より「運用」でつまずく。自社実践で現場の抵抗感・エラー回避まで熟知。", TEALD),
    ("一気通貫", "補助金申請＋定着まで",
     "IT導入補助金・医療情報化支援基金の活用を支援し初期コスト最小化。使いこなすまで3ヶ月伴走。", AMBER),
    ("リスクゼロ", "自信があるからこその設計",
     "事前分析で「本当に削減できる」確証と成功事例があるから、院にリスクを負わせないモデルが成立。", GREEND),
]
cw, ch, gx, gy = Inches(5.95), Inches(2.45), Inches(0.55), Inches(0.5)
x0, y0 = Inches(0.55), Inches(1.55)
for i, (ti, sub, body, ac) in enumerate(reasons):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    card(s, cx, cy, cw, ch)
    rect(s, cx, cy, Inches(0.12), ch, ac)
    chip(s, str(i+1), cx+Inches(0.35), cy+Inches(0.32), Inches(0.7), Inches(0.7), ac, sz=22)
    t(s, ti, cx+Inches(1.25), cy+Inches(0.3), cw-Inches(1.4), Inches(0.5),
      sz=19, bold=True, col=INK)
    t(s, sub, cx+Inches(1.25), cy+Inches(0.85), cw-Inches(1.4), Inches(0.4),
      sz=12, bold=True, col=ac)
    t(s, body, cx+Inches(0.4), cy+Inches(1.4), cw-Inches(0.7), Inches(0.95),
      sz=12.5, col=GRAY, line_sp=1.15)

# ════════════════════════════════════════════════════
# SLIDE 9 — CTA（無料コンサルの価値）
# ════════════════════════════════════════════════════
s = sl(NAVY)
rect(s, 0, 0, Inches(0.28), H, TEAL)
t(s, "まずは「無料 AI・LINE活用コンサル」から", Inches(0.9), Inches(0.7),
  Inches(11.5), Inches(0.7), sz=30, bold=True, col=WHT)
t(s, "押し売り一切なし。単なる問い合わせでなく、1時間の無料セッションとしてご利用ください。",
  Inches(0.92), Inches(1.5), Inches(11.5), Inches(0.5), sz=14, col=RGBColor(0xC4,0xD3,0xDF))
items = [
    ("ネット上の競合分析", "貴院周辺クリニックのデジタル対応状況を分析"),
    ("現状フローの整理", "電話の多さ・書類の手間などボトルネックをその場で特定"),
]
for i, (ti, ds) in enumerate(items):
    cy = Inches(2.35) + Inches(1.15) * i
    card(s, Inches(0.9), cy, Inches(6.7), Inches(0.95), RGBColor(0x16,0x33,0x4D), TEALD)
    t(s, "✓ " + ti, Inches(1.15), cy+Inches(0.13), Inches(6.3), Inches(0.4),
      sz=15, bold=True, col=TEAL)
    t(s, ds, Inches(1.15), cy+Inches(0.5), Inches(6.3), Inches(0.4),
      sz=11.5, col=RGBColor(0xC4,0xD3,0xDF))
# 特典カード（強調）
card(s, Inches(7.9), Inches(2.35), Inches(4.55), Inches(3.4), GREEN, GREEND, lw=2)
t(s, "★ 参加特典", Inches(7.9), Inches(2.6), Inches(4.55), Inches(0.5),
  sz=16, bold=True, col=WHT, align=PP_ALIGN.CENTER)
t(s, "貴院専用\n業務削減シミュレーション\n報告書（診断書）", Inches(7.9), Inches(3.2),
  Inches(4.55), Inches(1.6), sz=22, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.1)
t(s, "を無料で作成・プレゼント", Inches(7.9), Inches(4.95), Inches(4.55), Inches(0.5),
  sz=14, col=WHT, align=PP_ALIGN.CENTER)
chip(s, "今回は丁寧に伴走するため【限定 ◯ 名】様まで", Inches(0.9), Inches(6.3),
     Inches(11.55), Inches(0.6), AMBER, col=NAVY, sz=15)

# ════════════════════════════════════════════════════
# SLIDE 10 — クロージング
# ════════════════════════════════════════════════════
s = sl(NAVY)
rect(s, 0, 0, Inches(0.28), H, TEAL)
t(s, "リスクゼロで、\n次世代のクリニック運営を始めませんか？", Inches(0.9), Inches(0.85),
  Inches(11.5), Inches(1.6), sz=34, bold=True, col=WHT, line_sp=1.05)
t(s, "「とりあえず自院ならどこを自動化できるのか知りたい」── その動機で構いません。押し売りは一切いたしません。",
  Inches(0.92), Inches(2.75), Inches(11.5), Inches(0.5), sz=14, col=RGBColor(0xC4,0xD3,0xDF))
# 申込項目カード
card(s, Inches(0.9), Inches(3.6), Inches(7.1), Inches(3.1), RGBColor(0x16,0x33,0x4D), TEALD)
t(s, "お申込み（その場でQR／用紙）", Inches(1.2), Inches(3.85), Inches(6.5), Inches(0.5),
  sz=16, bold=True, col=TEAL)
for i, it in enumerate(["クリニック名 ／ ご担当者名",
                        "電話番号 ／ メールアドレス",
                        "ご希望の面談日時（オンライン対応可）"]):
    t(s, "▢  " + it, Inches(1.3), Inches(4.5)+Inches(0.6)*i, Inches(6.4), Inches(0.5),
      sz=14, col=WHT)
# QRプレースホルダ
photo_ph(s, Inches(8.55), Inches(3.6), Inches(3.9), Inches(3.1), "申込QRコード")

prs.save("clinic_dx_v2_img.pptx")
print("saved clinic_dx_v2_img.pptx  /  slides:", len(prs.slides._sldIdLst))

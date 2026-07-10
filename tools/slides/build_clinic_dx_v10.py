# -*- coding: utf-8 -*-
"""
クリニックDX「My AI」 v10 ── 忙しい医師向けショート版（全5枚）
設計思想: 先生の意思決定を1つに絞る＝「医療事務への30分ヒアリングのOK」だけ。
契約判断は無料の削減シミュレーション報告書を見てから（=今日は決めなくていい）。
料金: 初期0/月額0/成果報酬40%×6ヶ月→7ヶ月目〜月5万保守(月1MTG込)。
配色: クリーム白#F9F6EF × レンガ赤#AA2E26（KHD標準）
出力: clinic_dx_v10.pptx（詳細はv9を別添）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
GRN    = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Hiragino Sans"
_MK = "/Users/kikuchikenta/01_honbu_docs_automation/myai_mockups"
IMG1 = _MK + "/screen1_doc.png"
CHART_SAVINGS = "/Users/kikuchikenta/01_honbu_docs_automation/_savings_chart.png"
W = Inches(13.33); H = Inches(7.5)

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp: p.line_spacing = line_sp
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = FONT
    return tb

def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)

def ft(slide, page):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "クリニックDX「My AI」  ｜  テナントアシスト・ウイン株式会社  菊池 研太", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)
    t(slide, f"{page} / 5", Inches(12.0), H-Inches(0.42), Inches(0.8), Inches(0.32), sz=9, col=GRY, align=PP_ALIGN.RIGHT)

# ════════ 1. 表紙＝結論 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "AI × CLINIC DX ｜ 京橋クリニック様", Inches(0.9), Inches(0.95), Inches(8), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "現場の負担を減らして、", Inches(0.88), Inches(1.5), Inches(8), Inches(0.8), sz=36, bold=True, col=INK)
t(s, "院に 年約200万円 残す。", Inches(0.88), Inches(2.25), Inches(8), Inches(0.8), sz=36, bold=True, col=RED)
t(s, "電話・問診・書類・レセプトをAIとLINEで自動化。\n御院の持ち出しはゼロから始められます。",
  Inches(0.9), Inches(3.2), Inches(7.4), Inches(0.8), sz=14, col=GRY, line_sp=1.3)
# オファー3カード
ox, oy, ow, og = Inches(0.9), Inches(4.25), Inches(2.45), Inches(0.18)
for i, (lab, val) in enumerate([("初期費用", "0円"), ("月額基本料", "0円"), ("成果報酬", "40%")]):
    cx = ox + (ow + og) * i
    bx(s, cx, oy, ow, Inches(1.25), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, oy, ow, Inches(0.06), RED)
    t(s, lab, cx, oy+Inches(0.16), ow, Inches(0.35), sz=12, col=GRY, align=PP_ALIGN.CENTER)
    t(s, val, cx, oy+Inches(0.44), ow, Inches(0.7), sz=33, bold=True, col=RED, align=PP_ALIGN.CENTER)
t(s, "成果報酬は最初の6ヶ月だけ。7ヶ月目以降は 月5万円の保守のみ（月1回の運用相談込み）。\n効果が出なければ費用は0円。",
  Inches(0.9), Inches(5.7), Inches(7.5), Inches(0.7), sz=12.5, bold=True, col=REDD, line_sp=1.3)
# 右：MyAI画面
bx(s, Inches(8.7), Inches(1.4), Inches(4.2), Inches(3.3), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(8.7), Inches(1.4), Inches(4.2), Inches(0.06), RED)
s.shapes.add_picture(IMG1, Inches(8.95), Inches(1.75), width=Inches(3.7))
t(s, "実際の画面：カルテ→紹介状を数秒で下書き", Inches(8.7), Inches(4.8), Inches(4.2), Inches(0.4), sz=10, col=GRY, align=PP_ALIGN.CENTER)
# 導入の流れ（静かに置く）
bx(s, Inches(8.7), Inches(5.35), Inches(4.2), Inches(1.3), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(8.7), Inches(5.35), Inches(4.2), Inches(0.06), RED)
t(s, "導入の流れ", Inches(8.9), Inches(5.48), Inches(3.8), Inches(0.3), sz=11.5, bold=True, col=RED)
t(s, "無料の業務診断（30分） → 御院専用レポート\n→ 数字をご確認のうえ、導入のご判断", Inches(8.9), Inches(5.82), Inches(3.85), Inches(0.7), sz=10.5, col=INK, line_sp=1.25)
t(s, "テナントアシスト・ウイン株式会社 ｜ 菊池 研太", Inches(0.9), Inches(6.95), Inches(8), Inches(0.35), sz=11, col=GRY)

# ════════ 2. 御院で起きていること ════════
s = sl(); ft(s, 2)
hdr(s, "THE PROBLEM", "御院でも、起きていませんか？", "現場の4大負担 ── どれか1つでも当てはまれば、削減余地があります")
issues = [
    ("電話", "鳴り止まない電話で受付がパンク", "月1,200件超", "AI音声が一次対応・営業電話ブロック"),
    ("問診", "紙の問診票を毎日手入力", "受付業務の山", "LINEで来院前にスマホ問診"),
    ("書類", "紹介状・診断書で診療後も残業", "診察時間の半分", "AIがカルテから数秒で下書き"),
    ("レセプト", "月初の点検が毎月重い", "月10〜20時間", "AIが算定漏れを自動チェック"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.3), Inches(0.45), Inches(0.35)
x0, y0 = Inches(0.55), Inches(1.95)
for i, (tag, ti, num, sol) in enumerate(issues):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    bx(s, cx+Inches(0.32), cy+Inches(0.25), Inches(1.15), Inches(0.42), RED)
    t(s, tag, cx+Inches(0.32), cy+Inches(0.25), Inches(1.15), Inches(0.42), sz=13, bold=True, col=WHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t(s, ti, cx+Inches(1.65), cy+Inches(0.28), cw-Inches(1.85), Inches(0.6), sz=14.5, bold=True, col=INK)
    t(s, num, cx+Inches(0.34), cy+Inches(0.92), cw-Inches(0.6), Inches(0.6), sz=24, bold=True, col=RED)
    t(s, "→ " + sol, cx+Inches(0.36), cy+Inches(1.62), cw-Inches(0.65), Inches(0.5), sz=12, bold=True, col=GRN)

# ════════ 3. いくら残るか（グラフ1枚） ════════
s = sl(); ft(s, 3)
hdr(s, "SIMULATION", "使うほど、院にお金が残る", "標準モデル（事務3名・月24万円削減）。当日、御院の実数で計算し直してお出しします")
s.shapes.add_picture(CHART_SAVINGS, Inches(0.45), Inches(1.8), width=Inches(8.4))
RX, RW = Inches(9.1), Inches(3.75)
bx(s, RX, Inches(1.85), RW, Inches(2.2), CARD, line=CARDLN, lw=1.0)
t(s, "御院に残る効果", RX+Inches(0.2), Inches(1.98), RW-Inches(0.4), Inches(0.3), sz=13, bold=True, col=RED)
for _i, (_k, _v) in enumerate([("1〜6ヶ月", "月 14.4万円"), ("7ヶ月目〜", "月 19万円"), ("12ヶ月 累計", "約 200万円")]):
    _yy = Inches(2.38) + Inches(0.45) * _i
    t(s, _k, RX+Inches(0.25), _yy, Inches(1.6), Inches(0.4), sz=12, bold=True, col=INK)
    t(s, _v, RX+Inches(1.6), _yy, RW-Inches(1.85), Inches(0.4), sz=15, bold=True, col=REDD, align=PP_ALIGN.RIGHT)
bx(s, RX, Inches(4.3), RW, Inches(2.35), WHT, line=CARDLN, lw=1.0)
bx(s, RX, Inches(4.3), RW, Inches(0.06), RED)
t(s, "安心の設計", RX+Inches(0.2), Inches(4.45), RW-Inches(0.4), Inches(0.3), sz=13, bold=True, col=RED)
for _i, _line in enumerate(["効果が出なければ 費用0円", "測定の手間は全部こちらが巻き取る", "患者情報の扱いは最初に書面ルール化", "ロックインなし・30日前通知で解約可"]):
    t(s, "✓ " + _line, RX+Inches(0.22), Inches(4.85)+Inches(0.44)*_i, RW-Inches(0.44), Inches(0.4), sz=11.5, bold=True, col=INK)

# ════════ 4. 進め方（無料診断→効果確認→導入） ════════
s = sl(); ft(s, 4)
hdr(s, "HOW IT WORKS", "進め方 ── 効果を確認してから、導入へ", "費用が発生するのは効果が数字で確認できた後。判断材料はすべて先にお渡しします")
steps = [
    ("STEP 1", "無料の業務診断", "受付・事務の業務の流れを30分ほど伺います。診療の妨げにならない形で行います。", True),
    ("STEP 2", "御院専用レポート", "削減できる業務と金額を試算した報告書を無料でお渡しします。", False),
    ("STEP 3", "段階導入", "効果の大きい業務から少しずつ。スタッフ説明会・3ヶ月伴走付き。", False),
    ("STEP 4", "効果測定・精算", "測定は当方が実施（御院の手間なし）。確認できた削減分の40%のみ精算。", False),
]
cw, gx, x0, y0 = Inches(2.95), Inches(0.22), Inches(0.55), Inches(1.95)
for i, (when, ti, body, hi) in enumerate(steps):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(3.3), REDBG if hi else CARD, line=RED if hi else CARDLN, lw=2.0 if hi else 1.0)
    bx(s, cx, y0, cw, Inches(0.7), RED if hi else GRYBG)
    t(s, when, cx, y0+Inches(0.12), cw, Inches(0.45), sz=17, bold=True, col=WHT if hi else INK, align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.2), y0+Inches(0.9), cw-Inches(0.4), Inches(0.75), sz=15.5, bold=True, col=REDD if hi else INK, align=PP_ALIGN.CENTER, line_sp=1.1)
    t(s, body, cx+Inches(0.25), y0+Inches(1.75), cw-Inches(0.5), Inches(1.4), sz=11.5, col=GRY, line_sp=1.25)
    if i < 3:
        t(s, "›", cx+cw+Inches(0.0), y0+Inches(1.3), Inches(0.25), Inches(0.6), sz=22, bold=True, col=RED, align=PP_ALIGN.CENTER)
# 下バンド
band_y = y0 + Inches(3.6)
bx(s, Inches(0.55), band_y, Inches(12.23), Inches(0.95), REDBG)
bx(s, Inches(0.55), band_y, Inches(0.1), Inches(0.95), RED)
t(s, "費用が発生するのは、導入して、効果が数字で確認できた後だけ。", Inches(0.85), band_y+Inches(0.1), Inches(11.6), Inches(0.4), sz=15, bold=True, col=REDD)
t(s, "✓ 成果報酬40%は最初の6ヶ月だけ（以降は月5万の保守のみ）　✓ IT導入補助金サポート　✓ 3ヶ月伴走", Inches(0.85), band_y+Inches(0.54), Inches(11.6), Inches(0.35), sz=11.5, bold=True, col=INK)

# ════════ 5. 料金とサポート設計（すべて込み） ════════
s = sl(); ft(s, 5)
hdr(s, "PRICING & CARE", "料金とサポート ── 不安になりがちな点は、先に設計に入れています",
    "別料金・追加請求はありません。下記すべて込みの料金です")
# 左：料金サマリ
LX, LW = Inches(0.55), Inches(4.6)
bx(s, LX, Inches(1.95), LW, Inches(2.65), WHT, line=CARDLN, lw=1.0)
bx(s, LX, Inches(1.95), LW, Inches(0.06), RED)
t(s, "料金", LX+Inches(0.25), Inches(2.1), LW-Inches(0.5), Inches(0.3), sz=14, bold=True, col=RED)
for _i, (_k, _v) in enumerate([
    ("初期費用・月額基本料", "0円"),
    ("成果報酬", "削減額の40%"),
    ("成果報酬の期間", "最初の6ヶ月のみ"),
    ("7ヶ月目以降", "月5万円（保守のみ）"),
    ("効果が出なかった場合", "費用0円"),
]):
    _yy = Inches(2.5) + Inches(0.41) * _i
    t(s, _k, LX+Inches(0.25), _yy, Inches(2.5), Inches(0.38), sz=11.5, col=INK)
    t(s, _v, LX+Inches(2.55), _yy, LW-Inches(2.8), Inches(0.38), sz=12.5, bold=True, col=REDD, align=PP_ALIGN.RIGHT)
# 左下：保守に含むもの
bx(s, LX, Inches(4.8), LW, Inches(1.85), CARD, line=CARDLN, lw=1.0)
t(s, "保守（月5万円）に含むもの", LX+Inches(0.25), Inches(4.93), LW-Inches(0.5), Inches(0.3), sz=12.5, bold=True, col=RED)
for _i, _line in enumerate(["月1回の運用ミーティング・改善相談", "AIツールのアップデート対応・軽微な改修", "スタッフからの問い合わせ窓口（チャット）"]):
    t(s, "・" + _line, LX+Inches(0.27), Inches(5.3)+Inches(0.4)*_i, LW-Inches(0.5), Inches(0.38), sz=11, col=INK)
# 右：先回りのケア（チェックリスト2列風）
RX2, RW2 = Inches(5.45), Inches(7.35)
bx(s, RX2, Inches(1.95), RW2, Inches(4.7), CARD, line=CARDLN, lw=1.0)
bx(s, RX2, Inches(1.95), RW2, Inches(0.06), RED)
t(s, "料金に含まれるサポート（先回りの設計）", RX2+Inches(0.25), Inches(2.1), RW2-Inches(0.5), Inches(0.3), sz=14, bold=True, col=RED)
cares = [
    ("無料の業務診断＋御院専用レポート", "導入判断の材料は、費用が発生する前にすべてお渡し"),
    ("効果測定は当方が実施", "「2週間時間を計って」等のお願いは一切しません"),
    ("測定ルールの事前書面合意", "同等の患者数帯で比較・人員変更月は除外。後から揉めない設計"),
    ("患者情報の取り扱いルール書面化", "目的外利用なし・3省2ガイドライン準拠。導入前に取り決め"),
    ("スタッフ説明会＋3ヶ月伴走", "現場が使いこなせるまで。マニュアルも当方で用意"),
    ("IT導入補助金の申請支援", "使える補助金があれば、申請までお手伝い"),
    ("ロックインなし", "最低期間以降は30日前通知で解約可。データはお返しします"),
]
for _i, (_ti, _ds) in enumerate(cares):
    _yy = Inches(2.52) + Inches(0.59) * _i
    t(s, "✓", RX2+Inches(0.25), _yy, Inches(0.35), Inches(0.5), sz=15, bold=True, col=RED)
    t(s, _ti, RX2+Inches(0.62), _yy, Inches(3.6), Inches(0.5), sz=11.5, bold=True, col=INK)
    t(s, _ds, RX2+Inches(4.25), _yy+Inches(0.02), RW2-Inches(4.5), Inches(0.55), sz=9.5, col=GRY, line_sp=1.1)

import shutil
prs.save("clinic_dx_v10.pptx")
_dst = "/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/clinic_dx_v10_short.pptx"
shutil.copy("clinic_dx_v10.pptx", _dst)
print("saved clinic_dx_v10.pptx  /  slides:", len(prs.slides._sldIdLst))
print("copied to:", _dst)

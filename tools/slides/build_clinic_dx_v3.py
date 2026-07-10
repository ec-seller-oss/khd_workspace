"""
クリニックDX「My AI」 v3 "カッコいい版" 完成品（全10ページ）
SSoT: 提案スライド完全設計書（Google Doc 1QfTxY6oHsj3r4LCdM7N50KzZXwUvqPWbOCdnGVgPojM）。
核: 医者の4大悩み → AI商品(マイAI/Lステップ/IVR)がどう解決するか を分かりやすく。
デザイン: ディープ深夜ネイビー × エレクトリックティール、大胆タイポ、ヒーロー数字、帯グラデ。
出力: clinic_dx_v3.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── パレット ─────────────
INK0   = RGBColor(0x05, 0x10, 0x1E)
NAVY   = RGBColor(0x08, 0x18, 0x2B)
NAVY2  = RGBColor(0x0E, 0x24, 0x3C)
TEAL   = RGBColor(0x1A, 0xD6, 0xCB)
TEALD  = RGBColor(0x10, 0x8E, 0x88)
MINT   = RGBColor(0x8E, 0xF7, 0xEF)
GOLD   = RGBColor(0xF5, 0xC2, 0x42)
GREEN  = RGBColor(0x2D, 0xE0, 0x8E)
GREEND = RGBColor(0x16, 0x9E, 0x60)
RED    = RGBColor(0xFF, 0x6B, 0x7A)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
GRY    = RGBColor(0x9C, 0xB3, 0xC4)
CARD   = RGBColor(0x0F, 0x29, 0x42)
CARD2  = RGBColor(0x13, 0x31, 0x4F)
CARDLN = RGBColor(0x21, 0x48, 0x66)

FONT = "Hiragino Sans"
_MK = "/Users/kikuchikenta/01_honbu_docs_automation/myai_mockups"
IMG1 = _MK + "/screen1_doc.png"
IMG2 = _MK + "/screen2_dash.png"
IMG3 = _MK + "/screen3_line.png"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl(bg=INK0):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=WHT,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def rect(slide, x, y, w, h, fill, line=None, lw=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def grad_band(slide, x, y, w, h, c1, c2, angle=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background(); shp.shadow.inherit = False
    shp.fill.gradient()
    st = shp.fill.gradient_stops
    st[0].position = 0.0; st[0].color.rgb = c1
    st[1].position = 1.0; st[1].color.rgb = c2
    try:
        shp.fill.gradient_angle = angle
    except Exception:
        pass
    return shp


def card(slide, x, y, w, h, fill=CARD, line=CARDLN, lw=1.0):
    return rect(slide, x, y, w, h, fill, line, lw, rounded=True)


def chip(slide, text, x, y, w, h, fill, col, sz=12, bold=True):
    rect(slide, x, y, w, h, fill, rounded=True)
    t(slide, text, x, y, w, h, sz=sz, bold=bold, col=col,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def hdr(slide, eyebrow, main, sub=""):
    rect(slide, 0, 0, W, Inches(1.3), INK0)
    rect(slide, 0, Inches(1.3), W, Pt(2.5), TEAL)
    chip(slide, eyebrow, Inches(0.6), Inches(0.24), Inches(2.7), Inches(0.4), TEAL, INK0, sz=11)
    t(slide, main, Inches(0.6), Inches(0.66), Inches(12.1), Inches(0.55),
      sz=23, bold=True, col=WHT)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.0), Inches(12.1), Inches(0.3),
          sz=11.5, col=GRY)


def dark_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    """ダーク基調のテーブル。hi_col列をティール強調。"""
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = TEALD if is_hi else NAVY2
            else:
                cell.fill.fore_color.rgb = (RGBColor(0x10, 0x33, 0x33) if is_hi
                                            else (CARD if ri % 2 == 1 else CARD2))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = WHT
                    elif is_hi:
                        r.font.color.rgb = MINT
                    elif ci == 0:
                        r.font.color.rgb = WHT
                    else:
                        r.font.color.rgb = GRY
    return tb


# ════════ SLIDE 1 — 表紙 ════════
s = sl(INK0)
grad_band(s, Inches(7.4), 0, Inches(5.93), H, NAVY2, TEALD, angle=120)
rect(s, Inches(7.4), 0, Inches(0.05), H, TEAL)
rect(s, 0, 0, Inches(0.16), H, TEAL)
chip(s, "NEXT-GEN  CLINIC  DX", Inches(0.95), Inches(0.8), Inches(3.55), Inches(0.46), TEAL, INK0, sz=12.5)
t(s, '"人を増やさず"', Inches(0.9), Inches(1.45), Inches(6.6), Inches(1.0), sz=49, bold=True, col=WHT, line_sp=1.0)
t(s, '診療の質を、上げる。', Inches(0.9), Inches(2.4), Inches(6.6), Inches(1.0), sz=49, bold=True, col=TEAL, line_sp=1.0)
t(s, "クリニックのAI・LINE導入 まるごとおまかせ", Inches(0.92), Inches(3.62), Inches(6.4), Inches(0.5), sz=18, bold=True, col=WHT)
t(s, "独自カスタマイズの「マイAI」×「Lステップ」×AI音声(IVR)で実現する、\n初期費用0円・完全成果報酬型のクリニックDX。", Inches(0.92), Inches(4.12),
  Inches(6.5), Inches(0.9), sz=13, col=GRY, line_sp=1.25)
ox, oy, ow, og = Inches(0.92), Inches(5.45), Inches(2.05), Inches(0.22)
offers = [("初期費用", "0", "円"), ("月額基本料", "0", "円"), ("成果報酬", "30", "%")]
for i, (lab, val, unit) in enumerate(offers):
    cx = ox + (ow + og) * i
    t(s, lab, cx, oy, ow, Inches(0.3), sz=11, col=GRY)
    t(s, val, cx, oy+Inches(0.28), Inches(1.4), Inches(0.85), sz=50, bold=True, col=GOLD)
    t(s, unit, cx+Inches(1.0), oy+Inches(0.72), Inches(0.8), Inches(0.4), sz=18, bold=True, col=GOLD)
    if i < 2:
        rect(s, cx+ow+Inches(0.02), oy+Inches(0.32), Pt(1.2), Inches(0.8), CARDLN)
t(s, "成果報酬は「実際に人件費相当額を削減できた分」だけ", Inches(0.92), Inches(6.6), Inches(6.4), Inches(0.4), sz=10, col=GRY)
# 右：MyAI画面
rect(s, Inches(8.35), Inches(2.0), Inches(4.45), Inches(3.3), CARD, line=TEAL, lw=1.5, rounded=True)
rect(s, Inches(8.35), Inches(2.0), Inches(4.45), Inches(0.5), TEALD, rounded=True)
t(s, "マイAI ｜ 紹介状ジェネレーター", Inches(8.55), Inches(2.02), Inches(4.0), Inches(0.46), sz=11, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
s.shapes.add_picture(IMG1, Inches(8.6), Inches(2.66), width=Inches(3.95))
t(s, "実際のMyAI画面：カルテ → 紹介状を数秒で自動下書き", Inches(8.35), Inches(5.42), Inches(4.45), Inches(0.6), sz=10.5, col=MINT, align=PP_ALIGN.CENTER)
t(s, "菊池ホールディングス（KHD）  |  @khd_medical01  |  ※マイAIは仮称", Inches(0.92), Inches(7.0), Inches(11), Inches(0.4), sz=9.5, col=GRY)

# ════════ SLIDE 2 — 4大課題 ════════
s = sl(NAVY)
hdr(s, "THE  PROBLEM", "現場はもう限界では？ ── 御院でも起きていませんか？", "ネット上の生の声が示す、クリニックのリアルな4大お悩み")
issues = [
    ("01", "採用しても、定着しない", "70万〜150万", "円", "医療事務の採用単価。保険制度の複雑さで即戦力にならず早期離職の悪循環。"),
    ("02", "カルテ入力で患者と目が合わない", "診察の半分", "", "が電子カルテ入力。診療後も紹介状・診断書作成で2時間超の残業。"),
    ("03", "月初10日のレセプトが毎月重い", "月10〜20", "h残業", "が常態化。目視チェックによる算定漏れの不安も重なる。"),
    ("04", "鳴り止まない電話で受付がパンク", "月1,200", "件超", "の電話対応。「今空いてますか？」で目の前の患者対応が滞る。"),
]
cw, ch, gx, gy = Inches(6.05), Inches(2.45), Inches(0.5), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.62)
for i, (no, ti, num, unit, desc) in enumerate(issues):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    card(s, cx, cy, cw, ch)
    rect(s, cx, cy, Inches(0.14), ch, RED)
    t(s, no, cx+Inches(0.32), cy+Inches(0.2), Inches(1.1), Inches(0.6), sz=15, bold=True, col=RED)
    t(s, ti, cx+Inches(1.0), cy+Inches(0.24), cw-Inches(1.2), Inches(0.7), sz=16, bold=True, col=WHT)
    t(s, num, cx+Inches(0.34), cy+Inches(1.0), cw-Inches(2.0), Inches(0.8), sz=37, bold=True, col=RED)
    if unit:
        t(s, unit, cx+Inches(0.34), cy+Inches(1.18), cw-Inches(0.6), Inches(0.5), sz=18, bold=True, col=RED, align=PP_ALIGN.RIGHT)
    t(s, desc, cx+Inches(0.34), cy+Inches(1.8), cw-Inches(0.6), Inches(0.6), sz=11.5, col=GRY, line_sp=1.1)

# ════════ SLIDE 3 — 自動化マップ（核：ニーズ→AI解決）════════
s = sl(NAVY)
hdr(s, "THE  SOLUTION", "既存の業務をどう変えるか？ ── 完全自動化マップ", "医者の悩み（左）を、AI商品（右）がどう解決するか。クリニック全体を1つのエコシステムで最適化")
rows = [
    ("業務領域", "現状のアナログ業務", "マイAI・Lステップ・IVR導入後", "削減効果"),
    ("問診・予約", "電話受付＋紙問診をカルテへ手入力", "Lステップで来院前にスマホ上で完結", "問診業務 40〜70%減"),
    ("電話対応", "全件スタッフが受話", "AI音声(IVR)が一次対応・LINE誘導", "電話件数 最大40〜80%減"),
    ("書類作成", "過去カルテを探しつつ手作業", "自院学習「マイAI」が下書き自動生成", "書類時間 月30h以上減"),
    ("レセプト点検", "月末月初に目視で残業", "AIが4,000ルールで算定漏れ自動チェック", "点検時間 最大1/20"),
]
dark_table(s, rows, Inches(0.55), Inches(1.55), Inches(12.23), Inches(5.1),
           [Inches(2.0), Inches(3.85), Inches(4.05), Inches(2.33)], hi_col=3, sz=12.5, header_sz=13)

# ════════ SLIDE 4 — 実際のMyAI画面 ════════
s = sl(NAVY)
hdr(s, "PRODUCT", '"使う場面"が見える ── 実際のMyAI画面', "院長・スタッフ・患者、それぞれの画面で業務が変わる（画面はデモ・数値は例示）")
card(s, Inches(0.55), Inches(1.55), Inches(7.15), Inches(4.95))
s.shapes.add_picture(IMG1, Inches(0.8), Inches(1.9), width=Inches(6.5))
t(s, "① 書類アシスト｜カルテ→紹介状を数秒で下書き（院長の画面）", Inches(0.8), Inches(6.05), Inches(6.7), Inches(0.35), sz=12, bold=True, col=MINT)
card(s, Inches(7.95), Inches(1.55), Inches(4.83), Inches(2.9))
s.shapes.add_picture(IMG2, Inches(8.42), Inches(1.68), width=Inches(3.9))
t(s, "② 成果ダッシュボード｜削減効果を可視化", Inches(7.95), Inches(4.12), Inches(4.83), Inches(0.3), sz=11, bold=True, col=MINT, align=PP_ALIGN.CENTER)
card(s, Inches(7.95), Inches(4.6), Inches(4.83), Inches(1.9))
s.shapes.add_picture(IMG3, Inches(8.3), Inches(4.72), height=Inches(1.66))
t(s, "③ LINE問診・予約\nLステップ連携（患者スマホ）", Inches(9.35), Inches(5.05), Inches(3.3), Inches(1.0), sz=11.5, bold=True, col=MINT, line_sp=1.2)

# ════════ SLIDE 5 — 実証データ3事例 ════════
s = sl(NAVY)
hdr(s, "EVIDENCE", "「本当に自院でも変わるのか？」を証明する実証データ", "※当社実績ではなく、公開されている医療機関の導入実績です（出典明記）")
cases = [
    ("問診・予約のDX", "内科クリニック", "電話の嵐＋紙問診の手入力で受付過多", "Lステップで来院前にスマホ問診完了。電話を一次自動対応",
     ["初月にLINE予約 350件", "受付問診 70%削減", "月商 約1,500万円規模を押上げ"], "出典: Lステップ公開事例「そのだ内科」ほか"),
    ("書類作成のDX", "病院・総合クリニック", "紹介状・退院サマリ作成に医師が忙殺", "自院データを学習した生成AIが下書き自動生成",
     ["診断書 月400件を 50%削減", "退院サマリ 28分→8分(7割減)", "医師業務 月30h以上削減"], "出典: 新古賀/名古屋医療C/戸畑共立"),
    ("電話業務のDX", "地域密着クリニック", "インフル予約・時間外問合せで電話過多", "AI音声が一次対応、営業電話を自動ブロック",
     ["月1,200件の電話を自動化", "うち約400件(40%)を削減", "対応時間 全体70%削減"], "出典: IVRy 公開導入事例"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.55)
for i, (cat, who, bf, af, results, src) in enumerate(cases):
    cx = x0 + (cw + gx) * i
    card(s, cx, y0, cw, Inches(5.05))
    chip(s, cat, cx+Inches(0.2), y0+Inches(0.18), cw-Inches(0.4), Inches(0.42), TEALD, WHT, sz=12.5)
    t(s, who, cx+Inches(0.2), y0+Inches(0.66), cw-Inches(0.4), Inches(0.3), sz=11, bold=True, col=GRY, align=PP_ALIGN.CENTER)
    rect(s, cx+Inches(0.2), y0+Inches(1.0), cw-Inches(0.4), Inches(0.8), CARD2, rounded=True)
    t(s, "Before", cx+Inches(0.32), y0+Inches(1.04), Inches(1.2), Inches(0.28), sz=9.5, bold=True, col=RED)
    t(s, bf, cx+Inches(0.32), y0+Inches(1.28), cw-Inches(0.64), Inches(0.5), sz=10.5, col=WHT, line_sp=1.05)
    rect(s, cx+Inches(0.2), y0+Inches(1.88), cw-Inches(0.4), Inches(0.8), RGBColor(0x0E,0x33,0x44), rounded=True)
    t(s, "After", cx+Inches(0.32), y0+Inches(1.92), Inches(1.2), Inches(0.28), sz=9.5, bold=True, col=TEAL)
    t(s, af, cx+Inches(0.32), y0+Inches(2.16), cw-Inches(0.64), Inches(0.5), sz=10.5, col=WHT, line_sp=1.05)
    rect(s, cx+Inches(0.2), y0+Inches(2.78), cw-Inches(0.4), Inches(1.75), RGBColor(0x0C,0x33,0x28), line=GREEND, lw=1.0, rounded=True)
    t(s, "成果", cx+Inches(0.32), y0+Inches(2.82), Inches(1.2), Inches(0.3), sz=10, bold=True, col=GREEN)
    for j, rr in enumerate(results):
        t(s, "● " + rr, cx+Inches(0.32), y0+Inches(3.14)+Emu(int(Inches(0.43))*j), cw-Inches(0.6), Inches(0.42), sz=11.5, bold=True, col=GREEN, line_sp=1.0)
    t(s, src, cx+Inches(0.2), y0+Inches(4.6), cw-Inches(0.4), Inches(0.4), sz=8.5, col=GRY, align=PP_ALIGN.CENTER)

# ════════ SLIDE 6 — ROI早見表 ════════
s = sl(NAVY)
hdr(s, "SIMULATION", "御院の規模なら、いくら残る？ ── 削減効果の早見表", "1モデルではなく、自院を当てはめられる規模別シミュレーション")
rows = [
    ("院の規模", "月の削減額", "成果報酬(30%)", "院に残る効果（年額）"),
    ("医師1・事務2", "約16万円", "4.8万円", "約134万円"),
    ("医師1・事務3（標準）", "約24万円", "7.2万円", "約202万円"),
    ("医師2・事務5", "約40万円", "12万円", "約336万円"),
]
dark_table(s, rows, Inches(0.8), Inches(1.6), Inches(11.73), Inches(2.7),
           [Inches(3.4), Inches(2.5), Inches(2.6), Inches(3.23)], hi_col=3, sz=15, header_sz=13)
card(s, Inches(0.8), Inches(4.6), Inches(11.73), Inches(2.2), CARD)
t(s, "標準モデルの内訳（医師1・事務3 ／ 事務人件費 月83万円・総480時間・時間単価 約1,729円）", Inches(1.0), Inches(4.78), Inches(11.3), Inches(0.4), sz=12.5, bold=True, col=TEAL)
t(s, "問診 −42h ＋ レセプト −45h ＋ 電話 −40h ＋ 書類 −12h ＝ 月139時間削減（フルタイム約1名分の余力創出）", Inches(1.0), Inches(5.26), Inches(11.3), Inches(0.4), sz=13, col=WHT)
t(s, "→ 月139h × 1,729円 ＝ 月 約24万円削減 → 年 約288万円。成果報酬を引いて、院に残る効果は 約202万円/年。", Inches(1.0), Inches(5.74), Inches(11.3), Inches(0.4), sz=13, bold=True, col=GREEN)
t(s, "※「うちは事務3人だから…年200万か」と、その場で自院に当てはめられます。", Inches(1.0), Inches(6.28), Inches(11.3), Inches(0.4), sz=11, col=GRY)

# ════════ SLIDE 7 — 競合比較表 ════════
s = sl(NAVY)
hdr(s, "WHY  RISK-ZERO", "他社に頼むと「先行投資リスク」を抱える ── 業界価格の真実", "他社は効果が出なくても固定費が発生。当サービスだけ唯一のノーリスク")
rows = [
    ("比較項目", "一般ITベンダー", "大手DXコンサル", "Lステップ構築代行", "当サービス（中立伴走）"),
    ("提供内容", "自社アプリ単体販売", "業務分析・高額導入支援", "LINE・Lステップ構築", "Lステップ×マイAI 全体最適化"),
    ("初期費用", "20万〜数十万", "100〜300万", "50〜100万超", "0円"),
    ("月額費用", "1.5〜5万(固定)", "50〜300万(固定)", "5〜20万(固定)", "0円（基本料なし）"),
    ("報酬体系", "稼働問わず固定", "稼働ベース固定", "固定保守料", "削減人件費の30%のみ"),
    ("運用定着", "自力運用", "別料金で高額", "現場運用に不慣れ", "3ヶ月伴走"),
    ("院側のリスク", "高", "極めて高", "中", "ゼロ（効果0なら0円）"),
]
dark_table(s, rows, Inches(0.45), Inches(1.5), Inches(12.43), Inches(5.25),
           [Inches(1.95), Inches(2.35), Inches(2.35), Inches(2.35), Inches(3.43)], hi_col=4, sz=11, header_sz=11.5)

# ════════ SLIDE 8 — 成果報酬3ステップ ════════
s = sl(NAVY)
hdr(s, "HOW  IT  WORKS", "「うちでも本当に測定できる？」── リスクゼロの3ステップ", "成果報酬への最大の懸念=「測定の手間」は、コンサル側が巻き取ります")
steps = [
    ("STEP 1", "事前の無料コンサル", "貴院の負担ゼロ", "当方が競合分析と類似事例を持ち込み、業務フローとボトルネックを洗い出し。\n\n※「2週間時間を計測」等の煩雑な作業は一切要求しません。", TEALD, None),
    ("STEP 2", "フェアな測定ルール", "双方納得の透明基準", "導入前後で同等の患者数帯に基づき削減時間を算出。人員変更月は除外。\n\n※最低契約期間以降は30日前通知で解約可。ロックインなし。", TEALD, None),
    ("STEP 3", "手元に利益が残る", "リスクは私たちが取る", "【例：月24万円分を削減】\n削減額24万 × 30% ＝ お支払い 7.2万円", GREEND,
     ("貴院に残る効果", "月 16.8万円", "（年 約200万円）／効果0なら支払い0円")),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.6)
for i, (st, ti, sub, body, accent, hi) in enumerate(steps):
    cx = x0 + (cw + gx) * i
    card(s, cx, y0, cw, Inches(5.1))
    rect(s, cx, y0, cw, Inches(0.95), accent, rounded=True)
    rect(s, cx, y0+Inches(0.55), cw, Inches(0.4), accent)
    t(s, st, cx, y0+Inches(0.1), cw, Inches(0.45), sz=20, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, sub, cx, y0+Inches(0.57), cw, Inches(0.32), sz=11, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.25), y0+Inches(1.12), cw-Inches(0.5), Inches(0.55), sz=17, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, body, cx+Inches(0.3), y0+Inches(1.75), cw-Inches(0.6), Inches(2.0), sz=11.5, col=GRY, line_sp=1.15)
    if hi:
        rect(s, cx+Inches(0.25), y0+Inches(3.4), cw-Inches(0.5), Inches(1.5), RGBColor(0x0C,0x33,0x28), line=GREEN, lw=1.5, rounded=True)
        t(s, hi[0], cx+Inches(0.25), y0+Inches(3.5), cw-Inches(0.5), Inches(0.3), sz=11, bold=True, col=GREEN, align=PP_ALIGN.CENTER)
        t(s, hi[1], cx+Inches(0.25), y0+Inches(3.8), cw-Inches(0.5), Inches(0.6), sz=26, bold=True, col=GREEN, align=PP_ALIGN.CENTER)
        t(s, hi[2], cx+Inches(0.3), y0+Inches(4.45), cw-Inches(0.6), Inches(0.4), sz=10, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.0)

# ════════ SLIDE 9 — 選ばれる理由 ════════
s = sl(NAVY)
hdr(s, "OUR  VALUE", "なぜ私たちは「完全成果報酬」で伴走できるのか？", "「なぜ他社は固定費で、御社は成果報酬でやれるのか？」への答え")
reasons = [
    ("中立選定", "特定メーカーに縛られない", "自社ツールの売り込みではない。Lステップ/マイAI/IVRから貴院の診療科・規模に最適な組合せを厳選。", TEAL),
    ("実践ノウハウ", "自社で泥臭く運用した知見", "ツールは「設計」より「運用」でつまずく。自社実践で現場の抵抗感・エラー回避まで熟知。", TEAL),
    ("一気通貫", "補助金申請＋定着まで", "IT導入補助金・医療情報化支援基金の活用を支援し初期コスト最小化。使いこなすまで3ヶ月伴走。", GOLD),
    ("リスクゼロ", "自信があるからこその設計", "事前分析で「本当に削減できる」確証と成功事例があるから、院にリスクを負わせないモデルが成立。", GREEN),
]
cw, ch, gx, gy = Inches(6.05), Inches(2.45), Inches(0.5), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.62)
for i, (ti, sub, body, ac) in enumerate(reasons):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    card(s, cx, cy, cw, ch)
    rect(s, cx, cy, Inches(0.14), ch, ac)
    chip(s, str(i+1), cx+Inches(0.35), cy+Inches(0.3), Inches(0.7), Inches(0.7), ac, INK0, sz=22)
    t(s, ti, cx+Inches(1.25), cy+Inches(0.28), cw-Inches(1.4), Inches(0.5), sz=19, bold=True, col=WHT)
    t(s, sub, cx+Inches(1.25), cy+Inches(0.83), cw-Inches(1.4), Inches(0.4), sz=12, bold=True, col=ac)
    t(s, body, cx+Inches(0.4), cy+Inches(1.4), cw-Inches(0.7), Inches(0.95), sz=12, col=GRY, line_sp=1.15)

# ════════ SLIDE 10 — CTA / クロージング ════════
s = sl(INK0)
grad_band(s, 0, 0, W, Inches(2.2), NAVY2, TEALD, angle=45)
rect(s, 0, 0, Inches(0.16), H, TEAL)
t(s, "リスクゼロで、次世代のクリニック運営を始めませんか？", Inches(0.9), Inches(0.55), Inches(11.8), Inches(0.8), sz=29, bold=True, col=WHT)
t(s, "押し売りは一切なし。単なる問い合わせでなく、1時間の無料コンサルセッションとしてご利用ください。", Inches(0.92), Inches(1.4), Inches(11.5), Inches(0.5), sz=14, col=MINT)
items = [
    ("ネット上の競合分析", "貴院周辺クリニックのデジタル対応状況を分析"),
    ("現状フローの整理", "電話の多さ・書類の手間などボトルネックをその場で特定"),
]
for i, (ti, ds) in enumerate(items):
    cy = Inches(2.7) + Inches(1.1) * i
    card(s, Inches(0.9), cy, Inches(6.7), Inches(0.95))
    t(s, "✓ " + ti, Inches(1.15), cy+Inches(0.13), Inches(6.3), Inches(0.4), sz=15, bold=True, col=TEAL)
    t(s, ds, Inches(1.15), cy+Inches(0.5), Inches(6.3), Inches(0.4), sz=11.5, col=GRY)
card(s, Inches(0.9), Inches(4.9), Inches(6.7), Inches(1.85), CARD)
t(s, "お申込み（その場でQR／用紙）", Inches(1.15), Inches(5.05), Inches(6.3), Inches(0.4), sz=14, bold=True, col=TEAL)
for i, it in enumerate(["クリニック名 ／ ご担当者名", "電話番号 ／ メールアドレス", "ご希望の面談日時（オンライン対応可）"]):
    t(s, "▢  " + it, Inches(1.25), Inches(5.5)+Inches(0.4)*i, Inches(6.2), Inches(0.4), sz=12.5, col=WHT)
# 特典カード
card(s, Inches(7.9), Inches(2.7), Inches(4.55), Inches(4.05), RGBColor(0x0C,0x33,0x28), GREEN, lw=2)
t(s, "★ 参加特典", Inches(7.9), Inches(2.95), Inches(4.55), Inches(0.5), sz=16, bold=True, col=GREEN, align=PP_ALIGN.CENTER)
t(s, "貴院専用\n業務削減シミュレーション\n報告書（診断書）", Inches(7.9), Inches(3.6), Inches(4.55), Inches(1.6), sz=22, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.1)
t(s, "を無料で作成・プレゼント", Inches(7.9), Inches(5.35), Inches(4.55), Inches(0.5), sz=14, col=MINT, align=PP_ALIGN.CENTER)
chip(s, "今回は丁寧に伴走するため【限定 ◯ 名】様まで", Inches(7.9), Inches(6.05), Inches(4.55), Inches(0.55), GOLD, INK0, sz=13)

prs.save("clinic_dx_v3.pptx")
print("saved clinic_dx_v3.pptx  /  slides:", len(prs.slides._sldIdLst))

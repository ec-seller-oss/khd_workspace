"""
クリニックDX「My AI」 v3 "カッコいい版" デザインサンプル（表紙＋課題1枚）
方向性確認用。OKなら全11枚へ展開。
デザイン: ディープ深夜ネイビー × エレクトリックティール、大胆タイポ、ヒーロー数字、
帯グラデ風レイヤー、余白多め。
出力: clinic_dx_v3_sample.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── パレット（攻めた医療×テック）─────────────
INK0   = RGBColor(0x05, 0x10, 0x1E)   # 最暗ネイビー（背景ベース）
NAVY   = RGBColor(0x08, 0x18, 0x2B)
NAVY2  = RGBColor(0x0E, 0x24, 0x3C)
TEAL   = RGBColor(0x1A, 0xD6, 0xCB)   # エレクトリックティール（主役）
TEALD  = RGBColor(0x10, 0x8E, 0x88)
MINT   = RGBColor(0x8E, 0xF7, 0xEF)
GOLD   = RGBColor(0xF5, 0xC2, 0x42)   # 0円/30% の一撃だけ
GREEN  = RGBColor(0x2D, 0xE0, 0x8E)
RED    = RGBColor(0xFF, 0x6B, 0x7A)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
GRY    = RGBColor(0x9C, 0xB3, 0xC4)
CARD   = RGBColor(0x0F, 0x29, 0x42)   # ダークカード
CARDLN = RGBColor(0x1E, 0x44, 0x60)

FONT = "Hiragino Sans"
_MK = "/Users/kikuchikenta/01_honbu_docs_automation/myai_mockups"
IMG1 = _MK + "/screen1_doc.png"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl(bg=INK0):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=WHT,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = col
        r.font.name = FONT
    return tb


def rect(slide, x, y, w, h, fill, line=None, lw=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def grad_band(slide, x, y, w, h, c1, c2, angle=0):
    """2色グラデーションの帯。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background(); shp.shadow.inherit = False
    shp.fill.gradient()
    stops = shp.fill.gradient_stops
    stops[0].position = 0.0; stops[0].color.rgb = c1
    stops[1].position = 1.0; stops[1].color.rgb = c2
    try:
        shp.fill.gradient_angle = angle
    except Exception:
        pass
    return shp


def chip(slide, text, x, y, w, h, fill, col, sz=12, bold=True):
    rect(slide, x, y, w, h, fill, rounded=True)
    t(slide, text, x, y, w, h, sz=sz, bold=bold, col=col,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 1 — 表紙（カッコいい版）
# ════════════════════════════════════════════════════
s = sl(INK0)
# 右側に斜めティール→ネイビーのグラデ大面（奥行き）
grad_band(s, Inches(7.4), 0, Inches(5.93), H, NAVY2, TEALD, angle=120)
rect(s, Inches(7.4), 0, Inches(0.05), H, TEAL)
# 左の極細アクセント
rect(s, 0, 0, Inches(0.16), H, TEAL)
# eyebrow
chip(s, "NEXT-GEN  CLINIC  DX", Inches(0.95), Inches(0.85), Inches(3.55), Inches(0.46),
     TEAL, INK0, sz=12.5)
# ヒーロー見出し
t(s, '"人を増やさず"', Inches(0.9), Inches(1.5), Inches(6.6), Inches(1.0),
  sz=50, bold=True, col=WHT, line_sp=1.0)
t(s, '診療の質を、上げる。', Inches(0.9), Inches(2.45), Inches(6.6), Inches(1.0),
  sz=50, bold=True, col=TEAL, line_sp=1.0)
# サブ
t(s, "クリニックDX「My AI」", Inches(0.92), Inches(3.7), Inches(6.4), Inches(0.5),
  sz=20, bold=True, col=WHT)
t(s, "Lステップ ×「マイAI」× AI音声(IVR)\nで実現する、リスクゼロのDX。", Inches(0.92), Inches(4.25),
  Inches(6.4), Inches(0.9), sz=14, col=GRY, line_sp=1.25)
# ヒーロー数字（offer）
ox, oy, ow, og = Inches(0.92), Inches(5.55), Inches(2.05), Inches(0.22)
offers = [("初期費用", "0", "円"), ("月額基本料", "0", "円"), ("成果報酬", "30", "%")]
for i, (lab, val, unit) in enumerate(offers):
    cx = ox + (ow + og) * i
    t(s, lab, cx, oy, ow, Inches(0.3), sz=11, col=GRY)
    t(s, val, cx, oy+Inches(0.28), Inches(1.4), Inches(0.85), sz=52, bold=True, col=GOLD)
    t(s, unit, cx+Inches(1.05), oy+Inches(0.72), Inches(0.8), Inches(0.4), sz=18, bold=True, col=GOLD)
    if i < 2:
        rect(s, cx+ow+Inches(0.02), oy+Inches(0.3), Pt(1.2), Inches(0.8), CARDLN)
# 右：MyAI画面をフローティングデバイス風
rect(s, Inches(8.35), Inches(2.05), Inches(4.45), Inches(3.25), CARD, line=TEAL, lw=1.5, rounded=True)
rect(s, Inches(8.35), Inches(2.05), Inches(4.45), Inches(0.5), TEALD, rounded=True)
t(s, "マイAI ｜ 紹介状ジェネレーター", Inches(8.55), Inches(2.07), Inches(4.0), Inches(0.46),
  sz=11, bold=True, col=WHT, anchor=MSO_ANCHOR.MIDDLE)
s.shapes.add_picture(IMG1, Inches(8.6), Inches(2.7), width=Inches(3.95))
t(s, "実際のMyAI画面：カルテ → 紹介状を数秒で自動下書き",
  Inches(8.35), Inches(5.45), Inches(4.45), Inches(0.6), sz=10.5, col=MINT, align=PP_ALIGN.CENTER)
# フッター
t(s, "菊池ホールディングス（KHD）   |   @khd_medical01   |   ※マイAIは仮称",
  Inches(0.92), Inches(6.95), Inches(11), Inches(0.4), sz=10, col=GRY)

# ════════════════════════════════════════════════════
# SLIDE 2 — 4大課題（カッコいい版・ヒーロー数字強調）
# ════════════════════════════════════════════════════
s = sl(NAVY)
# ヘッダー
rect(s, 0, 0, W, Inches(1.25), INK0)
rect(s, 0, Inches(1.25), W, Pt(3), TEAL)
chip(s, "THE  PROBLEM", Inches(0.6), Inches(0.22), Inches(2.4), Inches(0.4), TEAL, INK0, sz=11)
t(s, "現場はもう限界では？ ── 御院でも起きていませんか？", Inches(0.6), Inches(0.62), Inches(12), Inches(0.55),
  sz=24, bold=True, col=WHT)
issues = [
    ("01", "採用しても、定着しない", "70万〜150万", "円", "医療事務の採用単価。保険制度の複雑さで即戦力にならず早期離職の悪循環。"),
    ("02", "カルテ入力で患者と目が合わない", "診察の半分", "", "が電子カルテ入力。診療後も紹介状・診断書作成で2時間超の残業。"),
    ("03", "月初10日のレセプトが毎月重い", "月10〜20", "h残業", "が常態化。目視チェックによる算定漏れの不安も重なる。"),
    ("04", "鳴り止まない電話で受付がパンク", "月1,200", "件超", "の電話対応。「今空いてますか？」で目の前の患者対応が滞る。"),
]
cw, ch, gx, gy = Inches(6.05), Inches(2.5), Inches(0.5), Inches(0.42)
x0, y0 = Inches(0.55), Inches(1.6)
for i, (no, ti, num, unit, desc) in enumerate(issues):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    rect(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0, rounded=True)
    rect(s, cx, cy, Inches(0.14), ch, RED, rounded=False)
    t(s, no, cx+Inches(0.32), cy+Inches(0.22), Inches(1.1), Inches(0.6), sz=15, bold=True, col=RED)
    t(s, ti, cx+Inches(1.0), cy+Inches(0.26), cw-Inches(1.2), Inches(0.7), sz=16, bold=True, col=WHT)
    t(s, num, cx+Inches(0.34), cy+Inches(1.05), cw-Inches(0.6), Inches(0.8), sz=38, bold=True, col=RED)
    if unit:
        t(s, unit, cx+Inches(0.34), cy+Inches(1.05), cw-Inches(0.5), Inches(0.8), sz=38, bold=True, col=RED, align=PP_ALIGN.RIGHT)
    t(s, desc, cx+Inches(0.34), cy+Inches(1.85), cw-Inches(0.6), Inches(0.6), sz=11.5, col=GRY, line_sp=1.1)

prs.save("clinic_dx_v3_sample.pptx")
print("saved clinic_dx_v3_sample.pptx  /  slides:", len(prs.slides._sldIdLst))

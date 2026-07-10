"""
クリニックDX「My AI」 v4 ── 「260526_医療専門不動産のビジネスモデル」デッキに
デザイン・色を統一した版（紺×ゴールド／build_slides_minimal.py 準拠）。
SSoT: 提案スライド完全設計書（Google Doc 1QfTxY6oHsj3r4LCdM7N50KzZXwUvqPWbOCdnGVgPojM）。
保存先: Drive「260526_AI医療コンサル」(固定)。
出力: clinic_dx_v4.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 参照デッキと同一パレット（build_slides_minimal.py）──
BG   = RGBColor(0x0D, 0x1B, 0x2A)   # 紺・背景
ACC  = RGBColor(0xE8, 0xA8, 0x00)   # ゴールド・主アクセント
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
LGR  = RGBColor(0xCC, 0xD6, 0xE0)   # ライトグレー本文
MID  = RGBColor(0x1A, 0x2E, 0x44)   # ミッド紺・バー/カード
MID2 = RGBColor(0x12, 0x24, 0x3A)   # やや暗いカード（テーブル交互）
BLU  = RGBColor(0x4A, 0x9E, 0xCB)   # サブアクセント青
RED  = RGBColor(0xFF, 0x66, 0x66)   # 課題・リスク
GRN  = RGBColor(0x6F, 0xC6, 0x8C)   # 成果（控えめ緑）
GOLDBG = RGBColor(0x2E, 0x28, 0x12)  # ゴールド強調セル背景

FONT = "Hiragino Sans"
_MK = "/Users/kikuchikenta/01_honbu_docs_automation/myai_mockups"
IMG1 = _MK + "/screen1_doc.png"
IMG2 = _MK + "/screen2_dash.png"
IMG3 = _MK + "/screen3_line.png"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=WHT,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, main, sub=""):
    bx(slide, 0, 0, W, Inches(1.1), MID)
    t(slide, main, Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.6), sz=24, bold=True, col=WHT)
    if sub:
        t(slide, sub, Inches(0.62), Inches(0.78), Inches(12.1), Inches(0.3), sz=11.5, col=LGR)


def ft(slide):
    bx(slide, 0, H-Inches(0.38), W, Inches(0.38), MID)
    t(slide, "クリニックDX「My AI」  |  AI医療コンサル", Inches(0.5), H-Inches(0.35), Inches(10), Inches(0.32), sz=9, col=LGR)


def navy_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = ACC if is_hi else MID
            else:
                cell.fill.fore_color.rgb = GOLDBG if is_hi else (MID if ri % 2 == 1 else MID2)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = BG if is_hi else WHT
                    elif is_hi:
                        r.font.color.rgb = ACC
                    elif ci == 0:
                        r.font.color.rgb = WHT
                    else:
                        r.font.color.rgb = LGR
    return tb


def myai_mock(slide, x, y, w, h):
    """MyAI 紹介状ジェネレーター画面（図形描画・画像レス）。"""
    bx(slide, x, y, w, h, MID)
    bx(slide, x, y, w, Inches(0.06), ACC)
    bx(slide, x+Inches(0.18), y+Inches(0.2), w-Inches(0.36), Inches(0.42), RGBColor(0x24,0x3A,0x52))
    t(slide, "マイAI ｜ 紹介状ジェネレーター", x+Inches(0.32), y+Inches(0.21), w-Inches(0.5), Inches(0.4),
      sz=11, bold=True, col=ACC, anchor=MSO_ANCHOR.MIDDLE)
    bx(slide, x+Inches(0.18), y+Inches(0.74), w-Inches(0.36), Inches(0.5), RGBColor(0x10,0x22,0x36))
    t(slide, "カルテから抽出中… 患者ID / 主訴 / 既往歴 / 処方", x+Inches(0.3), y+Inches(0.76),
      w-Inches(0.5), Inches(0.46), sz=9.5, col=LGR, anchor=MSO_ANCHOR.MIDDLE)
    bx(slide, x+Inches(0.18), y+Inches(1.36), w-Inches(0.36), h-Inches(2.1), WHT)
    t(slide, "AI生成プレビュー", x+Inches(0.34), y+Inches(1.44), w-Inches(0.6), Inches(0.3), sz=9, bold=True, col=RGBColor(0x0C,0x7A,0x73))
    for i, ln in enumerate(["──────────────────",
                            "拝啓  時下ますますご清栄のことと…",
                            "下記の患者をご紹介申し上げます。",
                            "診断: ＿＿＿  /  所見: ＿＿＿＿",
                            "──────────────────"]):
        t(slide, ln, x+Inches(0.34), y+Inches(1.74)+Emu(int(Inches(0.28))*i), w-Inches(0.66), Inches(0.28),
          sz=9, col=RGBColor(0x5B,0x6B,0x7B))
    bx(slide, x+Inches(0.18), y+h-Inches(0.6), w-Inches(0.36), Inches(0.42), ACC)
    t(slide, "⚡ 8分 → 数秒で下書き完了", x+Inches(0.18), y+h-Inches(0.6), w-Inches(0.36), Inches(0.42),
      sz=11, bold=True, col=BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def dash_mock(slide, x, y, w, h):
    """成果ダッシュボード画面（棒グラフ風）。"""
    bx(slide, x, y, w, h, MID)
    bx(slide, x, y, w, Inches(0.06), ACC)
    t(slide, "📊 削減効果ダッシュボード", x+Inches(0.25), y+Inches(0.18), w-Inches(0.4), Inches(0.35), sz=11, bold=True, col=ACC)
    bars = [("問診", 0.70), ("電話", 0.55), ("書類", 0.62), ("レセプト", 0.90)]
    base_x = x+Inches(0.3); base_y = y+Inches(0.7); bw = Inches(0.55); maxh = Inches(1.5); gap = Inches(0.45)
    for i, (lab, v) in enumerate(bars):
        bxx = base_x + (bw+gap)*i
        bh = Emu(int(maxh*v))
        bx(slide, bxx, base_y+maxh-bh, bw, bh, ACC if i == 3 else BLU)
        t(slide, str(int(v*100))+"%", bxx-Inches(0.1), base_y+maxh-bh-Inches(0.28), bw+Inches(0.2), Inches(0.26), sz=9, bold=True, col=WHT, align=PP_ALIGN.CENTER)
        t(slide, lab, bxx-Inches(0.15), base_y+maxh+Inches(0.04), bw+Inches(0.3), Inches(0.26), sz=9, col=LGR, align=PP_ALIGN.CENTER)


def line_mock(slide, x, y, w, h):
    """LINE問診・予約画面（チャット風）。"""
    bx(slide, x, y, w, h, RGBColor(0x06,0xC7,0x55))
    bx(slide, x+Inches(0.12), y+Inches(0.12), w-Inches(0.24), h-Inches(0.24), RGBColor(0xE9,0xF4,0xEC))
    t(slide, "LINE 問診・予約", x+Inches(0.28), y+Inches(0.2), w-Inches(0.5), Inches(0.3), sz=10, bold=True, col=RGBColor(0x06,0x7A,0x35))
    bubbles = [("来院前にスマホで問診✓", True), ("予約日をお選びください", False), ("○/○ 10:00 で予約完了", True)]
    by = y+Inches(0.6)
    for txt, right in bubbles:
        bw = w-Inches(1.0)
        bxx = x+(w-bw-Inches(0.28)) if right else x+Inches(0.28)
        bx(slide, bxx, by, bw, Inches(0.4), RGBColor(0x9A,0xE5,0xB0) if right else WHT)
        t(slide, txt, bxx+Inches(0.1), by+Inches(0.02), bw-Inches(0.2), Inches(0.36), sz=8.5, col=RGBColor(0x1A,0x2B,0x3A), anchor=MSO_ANCHOR.MIDDLE)
        by = by + Inches(0.48)


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "AIで、人を増やさず診療の質を上げる。", Inches(0.9), Inches(0.8), Inches(7.2), Inches(0.5),
  sz=17, col=ACC, italic=True)
t(s, "クリニックのAI・LINE導入\nまるごとおまかせ", Inches(0.9), Inches(1.45), Inches(7.3), Inches(1.9),
  sz=42, bold=True, col=WHT, line_sp=1.05)
t(s, "独自カスタマイズの「マイAI」×「Lステップ」×AI音声(IVR)で実現する、\n初期費用0円・完全成果報酬型のクリニックDX。",
  Inches(0.92), Inches(3.5), Inches(7.3), Inches(0.9), sz=13.5, col=LGR, line_sp=1.25)
# オファー3カード
ox, oy, ow, og = Inches(0.9), Inches(4.7), Inches(2.3), Inches(0.16)
offers = [("初期費用", "0円"), ("月額基本料", "0円"), ("成果報酬", "30%")]
for i, (lab, val) in enumerate(offers):
    cx = ox + (ow + og) * i
    bx(s, cx, oy, ow, Inches(1.35), MID)
    bx(s, cx, oy, ow, Inches(0.06), ACC)
    t(s, lab, cx, oy+Inches(0.18), ow, Inches(0.35), sz=12, col=LGR, align=PP_ALIGN.CENTER)
    t(s, val, cx, oy+Inches(0.5), ow, Inches(0.7), sz=34, bold=True, col=ACC, align=PP_ALIGN.CENTER)
t(s, "成果報酬は「実際に人件費相当額を削減できた分」だけ／医療現場の実証データに基づく設計／IT導入補助金サポート対応",
  Inches(0.9), Inches(6.35), Inches(7.3), Inches(0.7), sz=10.5, col=RGBColor(0x9F,0xB4,0xC4))
# 右：MyAI画面（図形モック）
myai_mock(s, Inches(8.45), Inches(2.2), Inches(4.45), Inches(3.2))
t(s, "実際のMyAI画面：カルテ → 紹介状を数秒で自動下書き",
  Inches(8.45), Inches(5.5), Inches(4.45), Inches(0.5), sz=10.5, col=RGBColor(0x9F,0xB4,0xC4), align=PP_ALIGN.CENTER)
t(s, "菊池ホールディングス（KHD）  |  @khd_medical01  |  ※マイAIは仮称",
  Inches(0.9), Inches(7.0), Inches(11), Inches(0.4), sz=9.5, col=RGBColor(0x9F,0xB4,0xC4))

# ════════ SLIDE 2 — 4大課題 ════════
s = sl(); ft(s)
hdr(s, "現場はもう限界では？ ── 御院でも起きていませんか？", "ネット上の生の声が示す、クリニックのリアルな4大お悩み")
issues = [
    ("01", "採用しても、定着しない", "70万〜150万円", "医療事務の採用単価。保険制度の複雑さで即戦力にならず早期離職の悪循環。"),
    ("02", "カルテ入力で患者と目が合わない", "診察時間の半分", "が電子カルテ入力。診療後も紹介状・診断書作成で2時間超の残業。"),
    ("03", "月初10日のレセプトが毎月重い", "月10〜20時間", "の残業が常態化。目視チェックによる算定漏れの不安も重なる。"),
    ("04", "鳴り止まない電話で受付がパンク", "月1,200件超", "の電話対応。「今空いてますか？」で目の前の患者対応が滞る。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.42), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.5)
for i, (no, ti, num, desc) in enumerate(issues):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, MID)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    t(s, no, cx+Inches(0.3), cy+Inches(0.2), Inches(1.0), Inches(0.6), sz=28, bold=True, col=RED)
    t(s, ti, cx+Inches(1.05), cy+Inches(0.28), cw-Inches(1.25), Inches(0.6), sz=16, bold=True, col=WHT)
    t(s, num, cx+Inches(0.32), cy+Inches(1.02), cw-Inches(0.6), Inches(0.7), sz=30, bold=True, col=ACC)
    t(s, desc, cx+Inches(0.34), cy+Inches(1.72), cw-Inches(0.6), Inches(0.6), sz=11.5, col=LGR, line_sp=1.1)

# ════════ SLIDE 3 — 自動化マップ（ニーズ→AI解決の核）════════
s = sl(); ft(s)
hdr(s, "既存の業務をどう変えるか？ ── 完全自動化マップ", "医者の悩み（左）を、AI商品（右）がどう解決するか。クリニック全体を1つのエコシステムで最適化")
rows = [
    ("業務領域", "現状のアナログ業務", "マイAI・Lステップ・IVR導入後", "削減効果"),
    ("問診・予約", "電話受付＋紙問診をカルテへ手入力", "Lステップで来院前にスマホ上で完結", "問診業務 40〜70%減"),
    ("電話対応", "全件スタッフが受話", "AI音声(IVR)が一次対応・LINE誘導", "電話件数 最大40〜80%減"),
    ("書類作成", "過去カルテを探しつつ手作業", "自院学習「マイAI」が下書き自動生成", "書類時間 月30h以上減"),
    ("レセプト点検", "月末月初に目視で残業", "AIが4,000ルールで算定漏れ自動チェック", "点検時間 最大1/20"),
]
navy_table(s, rows, Inches(0.55), Inches(1.5), Inches(12.23), Inches(5.1),
           [Inches(2.0), Inches(3.85), Inches(4.05), Inches(2.33)], hi_col=3, sz=12.5, header_sz=13)

# ════════ SLIDE 4 — 実際のMyAI画面 ════════
s = sl(); ft(s)
hdr(s, '"使う場面"が見える ── 実際のMyAI画面', "院長・スタッフ・患者、それぞれの画面で業務が変わる（画面はデモ・数値は例示）")
myai_mock(s, Inches(0.55), Inches(1.5), Inches(5.6), Inches(4.85))
t(s, "① 書類アシスト｜カルテ→紹介状を数秒で下書き（院長の画面）", Inches(0.55), Inches(6.0), Inches(5.6), Inches(0.35), sz=11, bold=True, col=ACC)
dash_mock(s, Inches(6.45), Inches(1.5), Inches(6.33), Inches(2.95))
t(s, "② 成果ダッシュボード｜削減効果を可視化", Inches(6.45), Inches(4.5), Inches(6.33), Inches(0.3), sz=11, bold=True, col=ACC)
line_mock(s, Inches(6.45), Inches(4.95), Inches(6.33), Inches(1.4))
t(s, "③ LINE問診・予約｜Lステップ連携（患者スマホ）", Inches(6.45), Inches(6.0), Inches(6.33), Inches(0.3), sz=11, bold=True, col=ACC)

# ════════ SLIDE 5 — 実証データ3事例 ════════
s = sl(); ft(s)
hdr(s, "「本当に自院でも変わるのか？」を証明する実証データ", "※当社実績ではなく、公開されている医療機関の導入実績です（出典明記）")
cases = [
    ("問診・予約のDX", "内科クリニック", "電話の嵐＋紙問診の手入力で受付過多", "Lステップで来院前にスマホ問診完了。電話を一次自動対応",
     ["初月にLINE予約 350件", "受付問診 70%削減", "月商 約1,500万円規模を押上げ"], "出典: Lステップ公開事例「そのだ内科」ほか"),
    ("書類作成のDX", "病院・総合クリニック", "紹介状・退院サマリ作成に医師が忙殺", "自院データを学習した生成AIが下書き自動生成",
     ["診断書 月400件を 50%削減", "退院サマリ 28分→8分(7割減)", "医師業務 月30h以上削減"], "出典: 新古賀/名古屋医療C/戸畑共立"),
    ("電話業務のDX", "地域密着クリニック", "インフル予約・時間外問合せで電話過多", "AI音声が一次対応、営業電話を自動ブロック",
     ["月1,200件の電話を自動化", "うち約400件(40%)を削減", "対応時間 全体70%削減"], "出典: IVRy 公開導入事例"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.5)
for i, (cat, who, bf, af, results, src) in enumerate(cases):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(5.05), MID)
    bx(s, cx, y0, cw, Inches(0.06), ACC)
    t(s, cat, cx, y0+Inches(0.18), cw, Inches(0.4), sz=13, bold=True, col=ACC, align=PP_ALIGN.CENTER)
    t(s, who, cx, y0+Inches(0.62), cw, Inches(0.3), sz=11, bold=True, col=LGR, align=PP_ALIGN.CENTER)
    bx(s, cx+Inches(0.2), y0+Inches(1.0), cw-Inches(0.4), Inches(0.8), MID2)
    t(s, "Before", cx+Inches(0.32), y0+Inches(1.04), Inches(1.2), Inches(0.28), sz=9.5, bold=True, col=RED)
    t(s, bf, cx+Inches(0.32), y0+Inches(1.28), cw-Inches(0.64), Inches(0.5), sz=10.5, col=WHT, line_sp=1.05)
    bx(s, cx+Inches(0.2), y0+Inches(1.88), cw-Inches(0.4), Inches(0.8), MID2)
    t(s, "After", cx+Inches(0.32), y0+Inches(1.92), Inches(1.2), Inches(0.28), sz=9.5, bold=True, col=BLU)
    t(s, af, cx+Inches(0.32), y0+Inches(2.16), cw-Inches(0.64), Inches(0.5), sz=10.5, col=WHT, line_sp=1.05)
    bx(s, cx+Inches(0.2), y0+Inches(2.78), cw-Inches(0.4), Inches(1.75), MID2)
    bx(s, cx+Inches(0.2), y0+Inches(2.78), Inches(0.08), Inches(1.75), ACC)
    t(s, "成果", cx+Inches(0.36), y0+Inches(2.82), Inches(1.2), Inches(0.3), sz=10, bold=True, col=ACC)
    for j, rr in enumerate(results):
        t(s, "● " + rr, cx+Inches(0.36), y0+Inches(3.14)+Emu(int(Inches(0.43))*j), cw-Inches(0.6), Inches(0.42), sz=11.5, bold=True, col=WHT, line_sp=1.0)
    t(s, src, cx, y0+Inches(4.6), cw, Inches(0.4), sz=8.5, col=LGR, align=PP_ALIGN.CENTER)

# ════════ SLIDE 6 — ROI早見表 ════════
s = sl(); ft(s)
hdr(s, "御院の規模なら、いくら残る？ ── 削減効果の早見表", "1モデルではなく、自院を当てはめられる規模別シミュレーション")
rows = [
    ("院の規模", "月の削減額", "成果報酬(30%)", "院に残る効果（年額）"),
    ("医師1・事務2", "約16万円", "4.8万円", "約134万円"),
    ("医師1・事務3（標準）", "約24万円", "7.2万円", "約202万円"),
    ("医師2・事務5", "約40万円", "12万円", "約336万円"),
]
navy_table(s, rows, Inches(0.8), Inches(1.5), Inches(11.73), Inches(2.6),
           [Inches(3.4), Inches(2.5), Inches(2.6), Inches(3.23)], hi_col=3, sz=15, header_sz=13)
bx(s, Inches(0.8), Inches(4.45), Inches(11.73), Inches(2.2), MID)
t(s, "標準モデルの内訳（医師1・事務3 ／ 事務人件費 月83万円・総480時間・時間単価 約1,729円）", Inches(1.0), Inches(4.62), Inches(11.3), Inches(0.4), sz=12.5, bold=True, col=ACC)
t(s, "問診 −42h ＋ レセプト −45h ＋ 電話 −40h ＋ 書類 −12h ＝ 月139時間削減（フルタイム約1名分の余力創出）", Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.4), sz=13, col=WHT)
t(s, "→ 月139h × 1,729円 ＝ 月 約24万円削減 → 年 約288万円。成果報酬を引いて、院に残る効果は 約202万円/年。", Inches(1.0), Inches(5.58), Inches(11.3), Inches(0.4), sz=13, bold=True, col=GRN)
t(s, "※「うちは事務3人だから…年200万か」と、その場で自院に当てはめられます。", Inches(1.0), Inches(6.12), Inches(11.3), Inches(0.4), sz=11, col=LGR)

# ════════ SLIDE 7 — 競合比較表 ════════
s = sl(); ft(s)
hdr(s, "他社に頼むと「先行投資リスク」を抱える ── 業界価格の真実", "他社は効果が出なくても固定費が発生。当サービスだけ唯一のノーリスク")
rows = [
    ("比較項目", "一般ITベンダー", "大手DXコンサル", "Lステップ構築代行", "当サービス（中立伴走）"),
    ("提供内容", "自社アプリ単体販売", "業務分析・高額導入支援", "LINE・Lステップ構築", "Lステップ×マイAI 全体最適化"),
    ("初期費用", "20万〜数十万", "100〜300万", "50〜100万超", "0円"),
    ("月額費用", "1.5〜5万(固定)", "50〜300万(固定)", "5〜20万(固定)", "0円（基本料なし）"),
    ("報酬体系", "稼働問わず固定", "稼働ベース固定", "固定保守料", "削減人件費の30%のみ"),
    ("運用定着", "自力運用", "別料金で高額", "現場運用に不慣れ", "3ヶ月伴走"),
    ("院側のリスク", "高", "極めて高", "中", "ゼロ（効果0なら0円）"),
]
navy_table(s, rows, Inches(0.45), Inches(1.45), Inches(12.43), Inches(5.25),
           [Inches(1.95), Inches(2.35), Inches(2.35), Inches(2.35), Inches(3.43)], hi_col=4, sz=11, header_sz=11.5)

# ════════ SLIDE 8 — 成果報酬3ステップ ════════
s = sl(); ft(s)
hdr(s, "「うちでも本当に測定できる？」── リスクゼロの3ステップ", "成果報酬への最大の懸念=「測定の手間」は、コンサル側が巻き取ります")
steps = [
    ("STEP 1", "事前の無料コンサル", "貴院の負担ゼロ", "当方が競合分析と類似事例を持ち込み、業務フローとボトルネックを洗い出し。\n\n※「2週間時間を計測」等の煩雑な作業は一切要求しません。", None),
    ("STEP 2", "フェアな測定ルール", "双方納得の透明基準", "導入前後で同等の患者数帯に基づき削減時間を算出。人員変更月は除外。\n\n※最低契約期間以降は30日前通知で解約可。ロックインなし。", None),
    ("STEP 3", "手元に利益が残る", "リスクは私たちが取る", "【例：月24万円分を削減】\n削減額24万 × 30% ＝ お支払い 7.2万円", ("貴院に残る効果", "月 16.8万円", "（年 約200万円）／効果0なら支払い0円")),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.5)
for i, (st, ti, sub, body, hi) in enumerate(steps):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, Inches(5.1), MID)
    bx(s, cx, y0, cw, Inches(0.9), MID2)
    bx(s, cx, y0, cw, Inches(0.06), ACC)
    t(s, st, cx, y0+Inches(0.12), cw, Inches(0.45), sz=20, bold=True, col=ACC, align=PP_ALIGN.CENTER)
    t(s, sub, cx, y0+Inches(0.56), cw, Inches(0.32), sz=11, col=LGR, align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.25), y0+Inches(1.1), cw-Inches(0.5), Inches(0.55), sz=17, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, body, cx+Inches(0.3), y0+Inches(1.72), cw-Inches(0.6), Inches(2.0), sz=11.5, col=LGR, line_sp=1.15)
    if hi:
        bx(s, cx+Inches(0.25), y0+Inches(3.4), cw-Inches(0.5), Inches(1.5), MID2, line=ACC, lw=1.5)
        t(s, hi[0], cx+Inches(0.25), y0+Inches(3.5), cw-Inches(0.5), Inches(0.3), sz=11, bold=True, col=ACC, align=PP_ALIGN.CENTER)
        t(s, hi[1], cx+Inches(0.25), y0+Inches(3.8), cw-Inches(0.5), Inches(0.6), sz=26, bold=True, col=ACC, align=PP_ALIGN.CENTER)
        t(s, hi[2], cx+Inches(0.3), y0+Inches(4.45), cw-Inches(0.6), Inches(0.4), sz=10, col=LGR, align=PP_ALIGN.CENTER, line_sp=1.0)

# ════════ SLIDE 9 — 選ばれる理由 ════════
s = sl(); ft(s)
hdr(s, "なぜ私たちは「完全成果報酬」で伴走できるのか？", "「なぜ他社は固定費で、御社は成果報酬でやれるのか？」への答え")
reasons = [
    ("中立選定", "特定メーカーに縛られない", "自社ツールの売り込みではない。Lステップ/マイAI/IVRから貴院の診療科・規模に最適な組合せを厳選。"),
    ("実践ノウハウ", "自社で泥臭く運用した知見", "ツールは「設計」より「運用」でつまずく。自社実践で現場の抵抗感・エラー回避まで熟知。"),
    ("一気通貫", "補助金申請＋定着まで", "IT導入補助金・医療情報化支援基金の活用を支援し初期コスト最小化。使いこなすまで3ヶ月伴走。"),
    ("リスクゼロ", "自信があるからこその設計", "事前分析で「本当に削減できる」確証と成功事例があるから、院にリスクを負わせないモデルが成立。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.42), Inches(0.45), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.5)
for i, (ti, sub, body) in enumerate(reasons):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, MID)
    bx(s, cx, cy, cw, Inches(0.06), ACC)
    t(s, str(i+1), cx+Inches(0.35), cy+Inches(0.3), Inches(0.9), Inches(0.9), sz=34, bold=True, col=ACC)
    t(s, ti, cx+Inches(1.3), cy+Inches(0.3), cw-Inches(1.5), Inches(0.5), sz=19, bold=True, col=WHT)
    t(s, sub, cx+Inches(1.3), cy+Inches(0.85), cw-Inches(1.5), Inches(0.4), sz=12, bold=True, col=ACC)
    t(s, body, cx+Inches(0.4), cy+Inches(1.4), cw-Inches(0.7), Inches(0.95), sz=12, col=LGR, line_sp=1.15)

# ════════ SLIDE 10 — CTA / クロージング ════════
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "リスクゼロで、次世代のクリニック運営を始めませんか？", Inches(0.9), Inches(0.7), Inches(11.7), Inches(0.8), sz=29, bold=True, col=WHT)
t(s, "押し売りは一切なし。単なる問い合わせでなく、1時間の無料コンサルセッションとしてご利用ください。",
  Inches(0.92), Inches(1.55), Inches(11.5), Inches(0.5), sz=14, col=LGR)
items = [
    ("ネット上の競合分析", "貴院周辺クリニックのデジタル対応状況を分析"),
    ("現状フローの整理", "電話の多さ・書類の手間などボトルネックをその場で特定"),
]
for i, (ti, ds) in enumerate(items):
    cy = Inches(2.45) + Inches(1.1) * i
    bx(s, Inches(0.9), cy, Inches(6.7), Inches(0.95), MID)
    bx(s, Inches(0.9), cy, Inches(0.06), Inches(0.95), ACC)
    t(s, "✓ " + ti, Inches(1.15), cy+Inches(0.13), Inches(6.3), Inches(0.4), sz=15, bold=True, col=ACC)
    t(s, ds, Inches(1.15), cy+Inches(0.5), Inches(6.3), Inches(0.4), sz=11.5, col=LGR)
bx(s, Inches(0.9), Inches(4.75), Inches(6.7), Inches(1.85), MID)
t(s, "お申込み（その場でQR／用紙）", Inches(1.15), Inches(4.9), Inches(6.3), Inches(0.4), sz=14, bold=True, col=ACC)
for i, it in enumerate(["クリニック名 ／ ご担当者名", "電話番号 ／ メールアドレス", "ご希望の面談日時（オンライン対応可）"]):
    t(s, "▢  " + it, Inches(1.25), Inches(5.35)+Inches(0.4)*i, Inches(6.2), Inches(0.4), sz=12.5, col=WHT)
# 特典
bx(s, Inches(7.9), Inches(2.45), Inches(4.55), Inches(4.15), MID, line=ACC, lw=2)
t(s, "★ 参加特典", Inches(7.9), Inches(2.7), Inches(4.55), Inches(0.5), sz=16, bold=True, col=ACC, align=PP_ALIGN.CENTER)
t(s, "貴院専用\n業務削減シミュレーション\n報告書（診断書）", Inches(7.9), Inches(3.35), Inches(4.55), Inches(1.6), sz=22, bold=True, col=WHT, align=PP_ALIGN.CENTER, line_sp=1.1)
t(s, "を無料で作成・プレゼント", Inches(7.9), Inches(5.1), Inches(4.55), Inches(0.5), sz=14, col=LGR, align=PP_ALIGN.CENTER)
bx(s, Inches(7.9), Inches(5.85), Inches(4.55), Inches(0.55), ACC)
t(s, "今回は丁寧に伴走するため【限定 ◯ 名】様まで", Inches(7.9), Inches(5.85), Inches(4.55), Inches(0.55), sz=13, bold=True, col=BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save("clinic_dx_v4.pptx")
print("saved clinic_dx_v4.pptx  /  slides:", len(prs.slides._sldIdLst))

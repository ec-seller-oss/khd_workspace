"""
クリニック公式LINE 立ち上げパイプライン（再利用テンプレ／京橋実例）
クリーム白×レンガ赤。出力: clinic_line_pipeline.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"クリニック公式LINE 立ち上げパイプライン ｜ KHD AI医療コンサル",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def arrow(slide,x,y,w,h):
    s=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,h); s.fill.solid()
    s.fill.fore_color.rgb=RGBColor(0xC9,0x6A,0x62); s.line.fill.background(); s.shadow.inherit=False; return s
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=(TEALBG if is_hi else (CARD if ri%2==1 else BG))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or (ci==0)
                    r.font.color.rgb=(WHT if ri==0 else (TEALD if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A))))
    return tb

# ════ SLIDE 1 — パイプライン全体像 ════
s=sl(); ft(s)
hdr(s,"PIPELINE","クリニック公式LINE 立ち上げパイプライン（8ステップ）","1院で型を確立 → 2院目以降はテンプレ流用で最短化。京橋クリニックで実証済み")
steps=[("1","アカウント開設","LINE公式＋Messaging API\n＋Login＋LIFF"),
       ("2","Harness デプロイ","Cloudflare無料枠\n（Claude Codeで構築）"),
       ("3","リッチメニュー","6タイルを設置\n（PNGテンプレ）"),
       ("4","FAQ自動応答","16ルール＋医療\nエスカレーション"),
       ("5","挨拶・イントロ","友だち追加時の\n自動メッセージ"),
       ("6","フォーム","Web問診／予約\n（LIFF）"),
       ("7","タグ＋リマインド","CPAP等の対象に\n定期配信"),
       ("8","院の実態反映","固有値・受付方式\nで仕上げ")]
cw,gx=Inches(2.78),Inches(0.34); x0=Inches(0.55)
for i,(no,ti,ds) in enumerate(steps):
    col=i%4; row=i//4
    cx=x0+(cw+gx)*col; cy=Inches(2.15)+(Inches(2.15))*row
    fill=REDBG if i==7 else CARD
    bx(s,cx,cy,cw,Inches(1.85),fill,line=CARDLN,lw=1.0); bx(s,cx,cy,cw,Inches(0.06),RED)
    bx(s,cx+Inches(0.25),cy+Inches(0.22),Inches(0.5),Inches(0.5),RED,shape=MSO_SHAPE.OVAL)
    t(s,no,cx+Inches(0.25),cy+Inches(0.22),Inches(0.5),Inches(0.5),sz=16,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ti,cx+Inches(0.85),cy+Inches(0.26),cw-Inches(0.95),Inches(0.5),sz=13.5,bold=True,col=INK)
    t(s,ds,cx+Inches(0.28),cy+Inches(0.95),cw-Inches(0.5),Inches(0.8),sz=10.5,col=GRY,line_sp=1.1)
    if col<3: arrow(s,cx+cw+Inches(0.02),cy+Inches(0.7),Inches(0.28),Inches(0.5))
t(s,"赤い⑧だけが院ごとに変わる中心。①〜⑦は型として次の院へそのまま流用できる。",Inches(0.55),Inches(6.45),Inches(12.2),Inches(0.4),sz=12.5,bold=True,col=REDD)

# ════ SLIDE 2 — 汎用 vs 院別カスタム ════
s=sl(); ft(s)
hdr(s,"REUSABLE vs CUSTOM","次の院でそのまま使う所 ／ 院ごとに変える所","比較表で「再利用ポイント」を固定 → 横展開のたびに作り直さない")
rows=[
 ("工程","汎用（テンプレ流用・次もそのまま）","院ごとにカスタマイズ"),
 ("基盤","Harnessデプロイ手順・MCP/API投入の型","アカウント名・各チャネルID"),
 ("リッチメニュー","6タイルの構成・PNG雛形","タイル文言（予約制／当日受付制で変わる）"),
 ("FAQ自動応答","16ルールの構造・医療エスカレーション設計","診療時間／電話／住所／休診／診療科目"),
 ("フォーム","Web問診・予約フォームの型","予約制の有無で転用（例：予約→健診申込）"),
 ("リマインド","タグ＋セグメント配信の仕組み","対象疾患・頻度（例：CPAP月次）"),
 ("挨拶・イントロ","文面テンプレ（{formUrl}差込）","院名・受付方式の表現"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(4.5),[Inches(2.2),Inches(5.6),Inches(4.43)],hi_col=1,sz=12,header_sz=13)
t(s,"緑＝汎用テンプレ（次の院でも流用）／右＝院別の差し替えだけ。これが横展開の生産性。",Inches(0.55),Inches(6.6),Inches(12.2),Inches(0.4),sz=12,bold=True,col=TEALD)

# ════ SLIDE 3 — 京橋の実カスタム ════
s=sl(); ft(s)
hdr(s,"CASE: 京橋クリニック","実例 ── 京橋クリニックでの院別カスタマイズ","「院の実態に合わせる」⑧の中身。次の院でもこの粒度でヒアリング→反映する")
cards=[
 ("受付方式","予約制ではなく『当日の順番受付制』。→ 予約タイルを廃止せず「健診・各種予約」に転用（健診・化学物質・禁煙外来は予約制）"),
 ("固有値（FAQへ反映）","月火木金 9:30-12:30/14:30-18:00（水は午前 院長不在）・土日祝休／☎03-3563-5011／中央区京橋2-5-22 キムラヤビル2F・京橋駅徒歩1分"),
 ("対象疾患・リマインド","睡眠時無呼吸（CPAP）の定期受診 → 「CPAP定期」タグ＋リマインド配信"),
 ("院長・診療科","院長 山崎 明男／一般内科・呼吸器・循環器・消化器・アレルギー・禁煙外来・産業医 ほか"),
]
y=Inches(2.05)
for i,(ti,ds) in enumerate(cards):
    cy=y+(Inches(1.18))*i
    bx(s,Inches(0.55),cy,Inches(12.23),Inches(1.04),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),cy,Inches(0.1),Inches(1.04),RED)
    t(s,ti,Inches(0.85),cy+Inches(0.12),Inches(3.0),Inches(0.8),sz=13.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ds,Inches(3.9),cy+Inches(0.1),Inches(8.6),Inches(0.85),sz=11.5,col=INK,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.12)
t(s,"※ 院別カスタムは「公式サイト＋先生ヒアリング」で収集 → ClaudeがAPIで即反映（京橋は所要数分）。",Inches(0.55),Inches(6.95),Inches(12.2),Inches(0.4),sz=11,col=GRY)

# ════ SLIDE 4 — 効果測定KPI（医療事務面談で確定） ════
s=sl(); ft(s)
hdr(s,"KPI ── 時短効果の測定","医療事務の面談で確定するKPI（想定）","導入前にベースラインを測り、導入後と比較する。指標・対象は現地ヒアリングで調整（変わる前提）")
rows=[
 ("KPI（候補）","測り方（導入前ベースライン）","LINEで効く所"),
 ("受付の電話本数 / 日","数日サンプルで着信件数を記録","FAQ自動応答が一次対応→件数減"),
 ("Web問診の事前記入率","来院前の記入数 ÷ 来院数","受付の聞き取り・電カル転記を削減"),
 ("受付の聞き取り・転記 時間","数件をストップウォッチ計測","Web問診で来院前に取得→短縮"),
 ("待ち時間・順番クレーム / 日","受付が件数をメモ","順番をLINEで自動通知→減"),
 ("定期受診の来院率（CPAP等）","対象者の受診 ÷ 対象者","リマインド配信で取りこぼし防止"),
 ("月末の残業時間","勤怠記録","通院の平準化で月末の山を削減"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(4.3),[Inches(3.5),Inches(4.6),Inches(4.13)],hi_col=None,sz=11.5,header_sz=12.5)
bx(s,Inches(0.55),Inches(6.4),Inches(12.23),Inches(0.62),REDBG); bx(s,Inches(0.55),Inches(6.4),Inches(0.1),Inches(0.62),RED)
t(s,"※ 数値・対象・優先順位は医療事務の面談で確定（現場が無理なく測れる粒度に）。押し付けず「一緒に測る」＝公平・透明に。",Inches(0.82),Inches(6.4),Inches(11.8),Inches(0.62),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("clinic_line_pipeline.pptx")
print("saved clinic_line_pipeline.pptx slides:",len(prs.slides._sldIdLst))

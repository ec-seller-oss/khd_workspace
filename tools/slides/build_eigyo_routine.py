"""
KHD 不動産業 営業ルーティン ＆ 物件フォルダDB（まず見る1枚デッキ）
クリーム白×レンガ赤。営業マニュアルDocs(1cbEBwPr)＋物件フォルダ箱(01-08)を構造化。
出力: eigyo_routine.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"; C=PP_ALIGN.CENTER; L=PP_ALIGN.LEFT
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=L,italic=False,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=22,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.46),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"KHD 不動産業 ｜ 営業ルーティン＆物件フォルダDB ｜ まず見る1枚（都度更新）",Inches(0.5),H-Inches(0.42),Inches(11),Inches(0.32),sz=9,col=GRY)
def table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=REDBG if is_hi else (CARD if ri%2==1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment=L if ci==0 else L
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or is_hi or (ci==0)
                    r.font.color.rgb=WHT if ri==0 else (RED if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
    return tb
def fcard(s,x,y,w,num,name,desc,hi=False):
    bx(s,x,y,w,Inches(0.06),RED)
    bx(s,x,y+Inches(0.06),w,Inches(1.18),REDBG if hi else CARD,CARDLN)
    t(s,num,x,y+Inches(0.13),w,Inches(0.3),sz=15,bold=True,col=RED,align=C)
    t(s,name,x,y+Inches(0.5),w,Inches(0.3),sz=11.5,bold=True,col=INK,align=C)
    t(s,desc,x+Inches(0.05),y+Inches(0.8),w-Inches(0.1),Inches(0.4),sz=8.5,col=GRY,align=C)

# ── S1 表紙 ──
s=sl()
bx(s,Inches(0.6),Inches(2.2),Inches(0.09),Inches(2.5),RED)
t(s,"KHD REAL ESTATE — SALES ROUTINE",Inches(0.95),Inches(2.2),Inches(11),Inches(0.5),sz=15,bold=True,col=RED)
t(s,"営業ルーティン ＆ 物件フォルダDB",Inches(0.92),Inches(2.85),Inches(11.5),Inches(1.0),sz=40,bold=True,col=INK)
t(s,"新規獲得 → 追客 → 査定 → 銀行上申。",Inches(0.95),Inches(4.05),Inches(11.5),Inches(0.5),sz=20,bold=True,col=RED)
t(s,"次に何をやるか、人に聞かず見える化する1枚。",Inches(0.95),Inches(4.55),Inches(11.5),Inches(0.5),sz=16,col=GRY)
bx(s,Inches(0.62),Inches(5.7),Inches(8),Pt(1.5),LINE)
t(s,"KHD 不動産業（買取再販・バイセル式TTP）  ｜  2026-06-04  ｜  現状版・都度更新して育てる",Inches(0.62),Inches(5.85),Inches(12),Inches(0.4),sz=12,col=GRY)

# ── S2 全体動線 01-08 ──
s=sl()
hdr(s,"WHOLE FLOW — 全体動線","物件1件のライフサイクル ＝ フォルダ番号が、そのまま作業順序",
    "物件ごとにDB番号を振り、中に 01〜08 の標準箱を置く。番号順に進めれば「今どこ・次に何」が分かる。")
folders=[("01_概要","物件情報・マイソク",False),("02_査定","収支入力・ざっと査定 ★",True),
("03_銀行","上申書・送付状で打診 ★",True),("04_仲介","客付け・媒介"[:],False),
("05_工事","解体・リフォーム",False),("06_領収書","経費・精算",False),
("07_保有中","賃貸管理",False),("08_売却","決済・出口",False)]
cw=Inches(2.92); gap=Inches(0.13)
for i,(num,desc,hi) in enumerate(folders):
    col=i%4; rowi=i//4
    x=Inches(0.6)+ (cw+gap)*col
    y=Inches(2.25)+ Inches(1.55)*rowi
    fcard(s,x,y,cw,num,num.split("_")[1],desc,hi)
bx(s,Inches(0.6),Inches(5.6),Inches(12.13),Inches(0.7),REDBG,CARDLN)
t(s,"★ 02 査定 → 03 銀行上申 までが「最速営業のやること」。まずここを最短で回す。細かい落とし穴は後工程で潰す。",
  Inches(0.8),Inches(5.72),Inches(11.8),Inches(0.5),sz=13,bold=True,col=RED)
ft(s)

# ── S3 フォルダDBの作り方 ──
s=sl()
hdr(s,"FOLDER DB — 物件フォルダの作り方","DB番号を振って、過去案件も“人に聞かず”見える化する")
steps=[("①","物件に DB番号 を振る","例：277_高松2丁目。番号＝物件の背番号。"),
("②","標準箱 01〜08 をコピー","空の標準フォルダを物件名の中に置く。"),
("③","番号順に作業を進める","今どこ・次に何をやるかが一目で分かる。"),
("④","過去案件も同じ型に揃える","資料から自分で遡れる＝属人化を防ぐ。")]
for i,(n,ti,ds) in enumerate(steps):
    y=Inches(2.15)+Inches(1.12)*i
    bx(s,Inches(0.6),y,Inches(7.4),Inches(0.96),CARD,CARDLN); bx(s,Inches(0.6),y,Inches(0.07),Inches(0.96),RED)
    t(s,n,Inches(0.78),y+Inches(0.22),Inches(0.6),Inches(0.5),sz=22,bold=True,col=RED)
    t(s,ti,Inches(1.5),y+Inches(0.13),Inches(6.3),Inches(0.4),sz=15,bold=True,col=INK)
    t(s,ds,Inches(1.5),y+Inches(0.52),Inches(6.4),Inches(0.4),sz=11,col=GRY)
table(s,[["命名ルール","例"],["物件フォルダ","DB番号_物件名"],["　例","277_高松2丁目"],
["中の標準箱","01_概要〜08_売却"],["銀行雛形","送付状.docx を流用"]],
Inches(8.3),Inches(2.15),Inches(4.4),Inches(2.6),[Inches(1.8),Inches(2.6)],sz=11,header_sz=12)
t(s,"※ この型に揃えるほど、誰でも・自分の過去案件からでも次の一手が見える。",Inches(8.3),Inches(4.95),Inches(4.5),Inches(0.6),sz=10,col=GRY)
ft(s)

# ── S4 ①新規獲得（架電） ──
s=sl()
hdr(s,"STEP ① 拾う（仕入）","業者架電で物件を引く ── 数のゲーム","格納先＝01_概要。物件が出たら情報・マイソクをここへ。")
table(s,[["項目","やること"],
["やること","架電リスト上から電話 → 訳アリ・築古・相続・空き家を聞く → 住所/価格/概要をメモ → 査定キューへ渡す"],
["格納先","01_概要 に物件情報・マイソクを入れる（ここが起点）"],
["KPI","1日のコール数 ／ 紹介獲得数（数を最優先）"],
["狙う物件","戸建・相続・空き家（土地のみは架電に不向き）"],
["禁止","押し売り・長文・嘘の緊急性。『また情報を回したい』と思われる関係づくりが最優先"]],
Inches(0.6),Inches(2.05),Inches(12.1),Inches(3.6),[Inches(2.0),Inches(10.1)],sz=12.5,header_sz=13)
ft(s)

# ── S5 ②査定（最速営業のコア） ──
s=sl()
hdr(s,"STEP ② 査る（査定）── 最速営業のコア","収支入力 → ざっと査定 → 銀行上申 を最短で回す")
flow=[("02_査定","収支シートに\n住所・面積・想定価格を入力"),
("ざっと査定","土地値割合0.4以上 ＆\n粗利率20%以上か即判定"),
("03_銀行","上申書＋送付状で\n銀行/ノンバンクへ打診")]
fw=Inches(3.75)
for i,(ti,ds) in enumerate(flow):
    x=Inches(0.7)+Inches(4.05)*i
    bx(s,x,Inches(2.2),fw,Inches(0.07),RED); bx(s,x,Inches(2.27),fw,Inches(1.5),CARD,CARDLN)
    t(s,ti,x,Inches(2.4),fw,Inches(0.4),sz=15,bold=True,col=RED,align=C)
    t(s,ds,x,Inches(2.85),fw,Inches(0.8),sz=12,col=INK,align=C,line_sp=1.1)
    if i<2: t(s,"→",x+fw-Inches(0.05),Inches(2.6),Inches(0.5),Inches(0.6),sz=24,bold=True,col=RED,align=C)
table(s,[["判断基準","ライン／メモ"],
["土地値割合","0.4以上（路線価×面積÷購入価格）＝担保の安全圏"],
["粗利率","20%以上（出口価格から逆算）"],
["再建築可否","必ず確認。不可なら無尽で買い→出口は三井トラストL&F"],
["融資先(自社)","実績ゼロは地銀・信金PJ当面不可。物件担保系ノンバンク＝日本住宅無尽/全宅ファイナンス/アサックス/SBIエステート"]],
Inches(0.6),Inches(4.1),Inches(12.1),Inches(2.5),[Inches(2.3),Inches(9.8)],sz=12,header_sz=12.5)
ft(s)

# ── S6 ③追客 ──
s=sl()
hdr(s,"STEP ③ 出す（客付け・追客）","実需=仲介 / 投資家=ツクビト。追客は“型”で回す")
table(s,[["観点","やること"],
["勝ち筋3つ","① 速度（反響60分以内接触）② LINE登録（追客の生命線）③ 悩みの魔法質問『前の部屋で困ったことは？』"],
["使うシート","追客管理シート(BtoC)で温度管理 ＋ トップセールス分析(羽鳥/野仲)の型"],
["羽鳥の型","商品を先に出さない／お願い感を出さない／質問で話させる（営業＝恋愛）"],
["中核信条","売り込まない。相手目線でGIVE、信頼の対価で収益化（信頼が先・収益は結果）"]],
Inches(0.6),Inches(2.05),Inches(12.1),Inches(3.2),[Inches(2.2),Inches(9.9)],sz=12.5,header_sz=13)
ft(s)

# ── S7 その先＋落とし穴 ──
s=sl()
hdr(s,"STEP ④+ 買う〜売る ／ 落とし穴","04仲介 → 05工事 → 06領収書 → 07保有 → 08売却")
chain=["04_仲介","05_工事","06_領収書","07_保有中","08_売却"]
for i,c in enumerate(chain):
    x=Inches(0.6)+Inches(1.18)*i
    bx(s,x,Inches(2.15),Inches(1.05),Inches(0.62),CARD,CARDLN)
    t(s,c,x,Inches(2.27),Inches(1.05),Inches(0.4),sz=9.5,bold=True,col=INK,align=C)
    if i<4: t(s,"→",x+Inches(1.0),Inches(2.18),Inches(0.25),Inches(0.5),sz=15,bold=True,col=RED,align=C)
t(s,"⚠ 落とし穴（先に潰す）",Inches(0.6),Inches(3.1),Inches(11),Inches(0.4),sz=14,bold=True,col=RED)
table(s,[["項目","確認ポイント"],
["再建築可否","不可は出口が限定。無尽で買い→三井トラストL&F"],
["接道・越境","接道2m以上。越境・残置・地中埋設物を現地で確認"],
["契約条文","05特約条文を契約に必ず入れる（瑕疵・残置・解体）"],
["出口の二択","実需=一般仲介（高値）／投資家=ツクビト（速い）"]],
Inches(0.6),Inches(3.55),Inches(12.1),Inches(2.5),[Inches(2.6),Inches(9.5)],sz=12,header_sz=12.5)
ft(s)

# ── S8 締め：都度更新 ──
s=sl()
bx(s,Inches(0.6),Inches(2.2),Inches(0.09),Inches(2.4),RED)
t(s,"UPDATE RULE — 育てるマニュアル",Inches(0.95),Inches(2.2),Inches(11),Inches(0.5),sz=15,bold=True,col=RED)
t(s,"このマニュアルは“未完成”。",Inches(0.92),Inches(2.85),Inches(11.5),Inches(0.7),sz=30,bold=True,col=INK)
t(s,"落とし穴・勝ち筋は、箱とスプシに追記して育てる。",Inches(0.92),Inches(3.6),Inches(11.5),Inches(0.7),sz=24,bold=True,col=RED)
t(s,"気づいた落とし穴・型は、物件フォルダ(01-08)と営業ハブスプシの『自社ルール欄』に都度追記。\n人に聞かずとも次にやることが見える状態を、全員で更新し続ける。",
  Inches(0.95),Inches(4.5),Inches(11.5),Inches(1.0),sz=14,col=GRY,line_sp=1.3)
bx(s,Inches(0.95),Inches(5.85),Inches(8),Pt(1.5),LINE)
t(s,"KHD 不動産業 ｜ 営業ルーティン＆物件フォルダDB ｜ 2026-06-04（現状版）",Inches(0.95),Inches(6.0),Inches(12),Inches(0.4),sz=12,col=GRY)

prs.save("eigyo_routine.pptx")
print("saved eigyo_routine.pptx  slides:",len(prs.slides._sldIdLst))

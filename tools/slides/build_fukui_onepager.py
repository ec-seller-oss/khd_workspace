"""
福井さん（事務長）向け A4一枚：公式LINE 運用で大切にしたいこと。
やさしい言葉・丁寧・シンプル。クリーム白×レンガ赤。出力: fukui_onepager.pptx（A4縦）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x5A,0x5A,0x5A); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"
W=Inches(8.27); H=Inches(11.69)   # A4縦
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
s=prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb=BG

def t(text,x,y,w,h,sz=14,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=1.25):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def box(x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sh=s.shapes.add_shape(shape,x,y,w,h)
    if col is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=col
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False; return sh

ML=Inches(0.62)  # 左マージン
CW=Inches(7.03)  # コンテンツ幅
# ヘッダー
box(Inches(0.5),Inches(0.5),Pt(4),Inches(0.95),RED)
t("公式LINE ― 運用で大切にしたいこと",Inches(0.72),Inches(0.5),Inches(7.2),Inches(0.6),sz=21,bold=True,col=INK)
t("京橋クリニックの皆さまへ",Inches(0.72),Inches(1.12),Inches(7.2),Inches(0.35),sz=13,col=RED,bold=True)
# はじめに
t("このたび、患者さんからのお問い合わせや受付のご負担を少しでも軽くするために、京橋クリニックさんの公式LINEを準備いたしました。スタッフの皆さまに安心してお使いいただけるよう、運用で大切にしたい点を一枚にまとめました。",
  ML,Inches(1.62),CW,Inches(1.0),sz=13,col=GRY,line_sp=1.35)

points=[
 ("スタッフの仕事を置きかえるものではありません",
  "LINEは「受付の代わり」ではなく、面倒な一次対応を引き受ける“助手”です。皆さまが患者さんに向き合う時間を増やすための道具としてお使いください。"),
 ("お薬や症状のご相談は、必ず人がお答えします",
  "むずかしいご質問には自動で答えないようにしています。「お電話・受付でご案内します」とお伝えする設計です。医療の判断は、これまでどおり先生とスタッフが行います。"),
 ("患者さんの個人情報は、これまでどおり大切に",
  "ご予約やWeb問診でいただいた情報は、院内のルールにそって、いつもと同じように扱います。"),
 ("まずは「1つだけ」、無料で小さく始めます",
  "すべてを一度に始めることはありません。効果が見えたものから、少しずつ広げていきます。ご負担をかけない進め方にします。"),
 ("効果は、皆さまと一緒に確かめます",
  "「お電話が減ったか」「楽になったか」を一緒に見ていきます。スタッフを評価するためのものではありません。"),
 ("合わなければ、いつでも見直せます",
  "無理に続けるものではありません。当院は「当日の順番受付制」のままです（予約制には変えません）。困ったことがあれば、すぐにご相談ください。"),
]
y=Inches(2.7)
for i,(ti,ds) in enumerate(points):
    box(ML,y,CW,Inches(1.15),CARD,line=CARDLN,lw=1.0); box(ML,y,Inches(0.1),Inches(1.15),RED)
    box(ML+Inches(0.28),y+Inches(0.2),Inches(0.5),Inches(0.5),RED,shape=MSO_SHAPE.OVAL)
    t(str(i+1),ML+Inches(0.28),y+Inches(0.2),Inches(0.5),Inches(0.5),sz=15,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(ti,ML+Inches(0.95),y+Inches(0.14),CW-Inches(1.1),Inches(0.4),sz=14,bold=True,col=REDD)
    t(ds,ML+Inches(0.97),y+Inches(0.52),CW-Inches(1.2),Inches(0.6),sz=11.5,col=INK,line_sp=1.22)
    y=y+Inches(1.27)

# むすび
box(ML,y+Inches(0.05),CW,Inches(1.0),RGBColor(0xF4,0xE4,0xE2)); box(ML,y+Inches(0.05),Inches(0.1),Inches(1.0),RED)
t("本日は、皆さまの「これは要らない」「これがあると助かる」というお声を、ぜひお聞かせください。\n一緒に、より良いかたちにしていきたいと思っております。どうぞよろしくお願いいたします。",
  ML+Inches(0.3),y+Inches(0.05),CW-Inches(0.5),Inches(1.0),sz=12.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.3)
t("KHD（AI医療コンサル）  ｜  京橋クリニック 公式LINE のご案内",ML,H-Inches(0.55),CW,Inches(0.3),sz=9.5,col=GRY)

prs.save("fukui_onepager.pptx")
print("saved fukui_onepager.pptx (A4縦)")

"""
gitって何？／GitHub移行とは？ 社内向けやさしい解説スライド
KHD標準デザイン（クリーム白×レンガ赤）。2026-07-02 の「main 2」トラブルを題材に、
菊池さん向けに①gitの仕組み ②今の運用のメリデメ ③GitHub移行の意味 ④今後の方針 を解説。
出力: git_explainer_260702.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Hiragino Sans"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK,
      align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=23, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "gitって何？／GitHub移行とは　｜　KHD秘書室 Tech Briefing", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)


def light_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w):
        tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0:
                cell.fill.fore_color.rgb = REDD if is_hi else RED
            else:
                cell.fill.fore_color.rgb = REDBG if is_hi else (CARD if ri % 2 == 1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0:
                        r.font.color.rgb = WHT
                    elif is_hi:
                        r.font.color.rgb = RED
                    elif ci == 0:
                        r.font.color.rgb = INK
                    else:
                        r.font.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    return tb


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "TECH BRIEFING ｜ 秘書室より", Inches(0.9), Inches(0.85), Inches(9), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "gitって何？", Inches(0.88), Inches(1.5), Inches(10), Inches(0.9), sz=44, bold=True, col=INK)
t(s, "そしてGitHub移行って何？", Inches(0.88), Inches(2.35), Inches(10.5), Inches(0.9), sz=36, bold=True, col=RED)
t(s, "2026-07-02、2台Mac同期の金庫で実際に起きた「main 2」トラブルを題材に、\n仕組み・今のメリデメ・今後の方針をやさしくまとめました。",
  Inches(0.92), Inches(3.35), Inches(10.8), Inches(0.9), sz=14.5, col=GRY, line_sp=1.3)
# 結論3カード
ox, oy, ow, og = Inches(0.9), Inches(4.6), Inches(3.55), Inches(0.25)
offers = [("今すぐ移行", "不要"), ("壊れた時", "AIが自動復旧"), ("次に検討", "事件が増えたら")]
for i, (lab, val) in enumerate(offers):
    cx = ox + (ow + og) * i
    bx(s, cx, oy, ow, Inches(1.35), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, oy, ow, Inches(0.06), RED)
    t(s, lab, cx, oy+Inches(0.2), ow, Inches(0.35), sz=12.5, col=GRY, align=PP_ALIGN.CENTER)
    t(s, val, cx, oy+Inches(0.52), ow, Inches(0.7), sz=24, bold=True, col=RED, align=PP_ALIGN.CENTER)
bx(s, Inches(0.9), Inches(6.55), Inches(11.5), Pt(1.2), LINE)
t(s, "先読みでの結論 ── 詳しくは次ページから", Inches(0.9), Inches(6.68), Inches(11), Inches(0.4), sz=12.5, col=GRY)

# ════════ SLIDE 2 — 今日、何が起きたか（3ステップ） ════════
s = sl(); ft(s)
hdr(s, "WHAT HAPPENED TODAY", "今日、実際に何が起きたのか？", "「同期が失敗しました」の裏側で起きていたこと（結果：実害なし・復旧済み）")
steps = [
    ("STEP 1", "2台のMacが同時に書き込み", "菊池さんのMacと、ゆーしさんが使うMacが\n同じ「金庫（gitリポジトリ）」の同じ台帳(main)に\nほぼ同時に書き込みをした。"),
    ("STEP 2", "Google Driveが「複製」を作成", "Driveは元々“ファイルを丸ごと同期する”道具。\n同じファイルへの同時書き込みを検知すると、\n安全策として複製ファイル「main 2」を作った。"),
    ("STEP 3", "gitが読めず、同期が全面停止", "gitは「main」という名前の台帳しか読めない。\n「main 2」という半端な名前のファイルが残っていると\n以後の同期作業ぜんぶが止まる。"),
]
cw, gx, x0, y0 = Inches(3.95), Inches(0.24), Inches(0.55), Inches(1.9)
CARDH = Inches(3.3)
for i, (st, ti, body) in enumerate(steps):
    cx = x0 + (cw + gx) * i
    bx(s, cx, y0, cw, CARDH, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, y0, cw, Inches(0.7), RED)
    t(s, st, cx, y0+Inches(0.1), cw, Inches(0.42), sz=17, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, cx+Inches(0.25), y0+Inches(0.86), cw-Inches(0.5), Inches(0.65), sz=15.5, bold=True, col=INK, align=PP_ALIGN.CENTER, line_sp=1.1)
    t(s, body, cx+Inches(0.3), y0+Inches(1.62), cw-Inches(0.6), Inches(1.6), sz=11.5, col=GRY, line_sp=1.25)
band_y = y0 + CARDH + Inches(0.3)
bx(s, Inches(0.55), band_y, Inches(12.23), Inches(0.85), REDBG)
bx(s, Inches(0.55), band_y, Inches(0.1), Inches(0.85), RED)
t(s, "結果：中身は既に本流に含まれていた古いコミットへの参照だけ ＝ データ損失ゼロ。AIが原因特定〜削除〜復旧まで完了。",
  Inches(0.85), band_y+Inches(0.24), Inches(11.6), Inches(0.4), sz=13, bold=True, col=REDD)

# ════════ SLIDE 3 — gitって何？ ════════
s = sl(); ft(s)
hdr(s, "WHAT IS GIT", "そもそも「git（ギット）」って何？", "一言で言うと──「変更履歴を全部記録してくれる台帳」")
cards = [
    ("① Wordの変更履歴の、フォルダ版", "Wordの「変更履歴を記録」機能を、フォルダ全体・全ファイルに対してやってくれる仕組み。\n誰が・いつ・何を変えたかを全部覚えていて、いつでも過去の状態に戻せる。"),
    ("② 「金庫」＝ 中央の台帳置き場", "複数の人・複数のパソコンで同じ台帳を共有するために、台帳の“正本”を1ヶ所（金庫）に置く。\n今のKHDでは、この金庫をGoogle Drive上に置いている。"),
    ("③ pull ＝ 金庫から取り出す", "自分のMacに、金庫にある最新の台帳（他の人の作業内容）を反映させること。\n作業を始める前に必ず行う「入口」の作業。"),
    ("④ push ＝ 金庫に預ける", "自分のMacでの作業内容を、金庫に反映させて他の人も見れるようにすること。\n作業が終わったら行う「出口」の作業。"),
]
cw, ch, gx, gy = Inches(6.0), Inches(2.35), Inches(0.45), Inches(0.35)
x0, y0 = Inches(0.55), Inches(1.85)
for i, (ti, body) in enumerate(cards):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    bx(s, cx, cy, cw, ch, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch, RED)
    t(s, ti, cx+Inches(0.32), cy+Inches(0.22), cw-Inches(0.6), Inches(0.55), sz=15.5, bold=True, col=RED)
    t(s, body, cx+Inches(0.32), cy+Inches(0.86), cw-Inches(0.6), Inches(1.35), sz=12, col=GRY, line_sp=1.3)

# ════════ SLIDE 4 — 今の運用の仕組み図解 ════════
s = sl(); ft(s)
hdr(s, "HOW WE USE IT TODAY", "今の運用 ── 「Google Drive上に置いたgit金庫」", "2台のMacが、Google Driveというフォルダの中にある金庫を経由してやり取りしている")
by, bh = Inches(2.4), Inches(2.0)
box_w = Inches(3.0)
mac1_x, vault_x, mac2_x = Inches(0.7), Inches(5.15), Inches(9.6)
bx(s, mac1_x, by, box_w, bh, CARD, line=CARDLN, lw=1.0)
bx(s, mac1_x, by, box_w, Inches(0.06), RED)
t(s, "💻 菊池さんのMac", mac1_x, by+Inches(0.3), box_w, Inches(0.4), sz=14, bold=True, col=INK, align=PP_ALIGN.CENTER)
t(s, "作業フォルダ＋記憶(memory)", mac1_x, by+Inches(0.85), box_w, Inches(0.4), sz=11, col=GRY, align=PP_ALIGN.CENTER)
t(s, "自分だけの手元コピー", mac1_x, by+Inches(1.3), box_w, Inches(0.4), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)

bx(s, vault_x, by-Inches(0.1), box_w, bh+Inches(0.2), REDBG, line=RED, lw=1.5)
bx(s, vault_x, by-Inches(0.1), box_w, Inches(0.06), RED)
t(s, "🔒 git金庫（.git）", vault_x, by+Inches(0.22), box_w, Inches(0.4), sz=14, bold=True, col=RED, align=PP_ALIGN.CENTER)
t(s, "Google Drive上に設置", vault_x, by+Inches(0.68), box_w, Inches(0.35), sz=11, col=REDD, align=PP_ALIGN.CENTER)
t(s, "台帳の“正本”を保管\n(workspace金庫／memory金庫)", vault_x, by+Inches(1.1), box_w, Inches(0.9), sz=10.5, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.2)

bx(s, mac2_x, by, box_w, bh, CARD, line=CARDLN, lw=1.0)
bx(s, mac2_x, by, box_w, Inches(0.06), RED)
t(s, "💻 ゆーしさんのMac", mac2_x, by+Inches(0.3), box_w, Inches(0.4), sz=14, bold=True, col=INK, align=PP_ALIGN.CENTER)
t(s, "同じアカウントで利用", mac2_x, by+Inches(0.85), box_w, Inches(0.4), sz=11, col=GRY, align=PP_ALIGN.CENTER)
t(s, "自分だけの手元コピー", mac2_x, by+Inches(1.3), box_w, Inches(0.4), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)

t(s, "⇄", Inches(3.75), by+Inches(0.55), Inches(1.35), Inches(0.9), sz=32, bold=True, col=RED, align=PP_ALIGN.CENTER)
t(s, "pull／push", Inches(3.6), by+Inches(1.35), Inches(1.6), Inches(0.35), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)
t(s, "⇄", Inches(8.2), by+Inches(0.55), Inches(1.35), Inches(0.9), sz=32, bold=True, col=RED, align=PP_ALIGN.CENTER)
t(s, "pull／push", Inches(8.05), by+Inches(1.35), Inches(1.6), Inches(0.35), sz=10.5, col=GRY, align=PP_ALIGN.CENTER)

note_y = by + bh + Inches(0.5)
bx(s, Inches(0.7), note_y, Inches(11.9), Inches(1.5), GRYBG)
t(s, "💡 誤解しやすいポイント", Inches(0.95), note_y+Inches(0.18), Inches(11.4), Inches(0.35), sz=13, bold=True, col=RED)
t(s, "「Google Driveに保存＝全部自動で同期される」と思われがちですが、それは普通のファイルの話。\nこの金庫(.git)の中身は、Driveの自動同期とは別に「pull→作業→push」という手続きを踏んで初めて反映されます。\n2台が同時にこの手続きを行うと、今回のような競合（main 2事件）が起きる構造的なリスクがあります。",
  Inches(0.95), note_y+Inches(0.55), Inches(11.4), Inches(0.9), sz=11.5, col=INK, line_sp=1.25)

# ════════ SLIDE 5 — 今の運用のメリット・デメリット ════════
s = sl(); ft(s)
hdr(s, "PROS & CONS", "今の運用（Drive上のgit金庫）のメリット・デメリット", "コストゼロで動いているが、破損リスクを抱えたまま使っている状態")
colw, colh, gx = Inches(5.85), Inches(4.7), Inches(0.5)
x0, y0 = Inches(0.55), Inches(1.85)
bx(s, x0, y0, colw, colh, CARD, line=CARDLN, lw=1.0)
bx(s, x0, y0, colw, Inches(0.55), RED)
t(s, "◎ メリット", x0, y0+Inches(0.1), colw, Inches(0.4), sz=16, bold=True, col=WHT, align=PP_ALIGN.CENTER)
merits = ["追加費用ゼロ（Google Drive契約の中で完結、他サービス契約不要）",
          "既にセットアップ済みで、日常的には問題なく動いている",
          "全部KHD管理下のGoogleアカウント内＝外部に情報が出ない"]
for i, mtext in enumerate(merits):
    my = y0+Inches(0.85)+Inches(1.2)*i
    t(s, "✓", x0+Inches(0.3), my, Inches(0.4), Inches(0.5), sz=18, bold=True, col=RED)
    t(s, mtext, x0+Inches(0.7), my, colw-Inches(1.0), Inches(1.1), sz=12.5, col=INK, line_sp=1.25)

x1 = x0 + colw + gx
bx(s, x1, y0, colw, colh, CARD, line=CARDLN, lw=1.0)
bx(s, x1, y0, colw, Inches(0.55), REDD)
t(s, "△ デメリット・リスク", x1, y0+Inches(0.1), colw, Inches(0.4), sz=16, bold=True, col=WHT, align=PP_ALIGN.CENTER)
demerits = ["2台が同時に書き込むと壊れやすい（今回のような事件が今後も起きうる）",
            "壊れた時の直し方に専門知識が要る（放置すると同期停止に気づかない恐れ）",
            "Google Driveは本来「gitの精密な同時書き込み」向けの道具ではなく、相性が良くない"]
for i, dtext in enumerate(demerits):
    dy = y0+Inches(0.85)+Inches(1.2)*i
    t(s, "!", x1+Inches(0.32), dy, Inches(0.4), Inches(0.5), sz=18, bold=True, col=REDD)
    t(s, dtext, x1+Inches(0.7), dy, colw-Inches(1.0), Inches(1.1), sz=12.5, col=INK, line_sp=1.25)

# ════════ SLIDE 6 — GitHub移行とは？ ════════
s = sl(); ft(s)
hdr(s, "WHAT IS GITHUB MIGRATION", "「GitHubへの移行」って何をすること？", "一言で言うと──「金庫番を、プロのサービスに外注する」こと")
cards2 = [
    ("GitHubとは", "gitの金庫（さっき説明した“正本置き場”）を専門にホスティングしてくれるサービス。世界中の会社・開発者が使う実績あるインフラ。個人の非公開リポジトリは無料。"),
    ("何が変わるか", "Google Drive経由をやめて、2台のMacが直接GitHubに対してpull/pushするようになる。日々の使い方（作業する→保存する）は今と同じで、裏側の“金庫の置き場所”だけが変わる。"),
    ("なぜ安定するか", "GitHubは同時書き込みの競合を正式な手順（マージ・ロック）で処理する専用インフラ。Driveのファイルコピー機能に頼らないため、今回のような「main 2」型の破損がそもそも起きない。"),
    ("必要な準備", "最初の一度だけ、GitHubアカウント作成と認証設定が必要（内務作業）。それ以降は特別な操作は増えず、菊池さんが覚えることは実質増えない。"),
]
cw2, ch2, gx2, gy2 = Inches(6.0), Inches(2.35), Inches(0.45), Inches(0.35)
x0b, y0b = Inches(0.55), Inches(1.85)
for i, (ti, body) in enumerate(cards2):
    cx = x0b + (cw2 + gx2) * (i % 2)
    cy = y0b + (ch2 + gy2) * (i // 2)
    bx(s, cx, cy, cw2, ch2, CARD, line=CARDLN, lw=1.0)
    bx(s, cx, cy, Inches(0.12), ch2, RED)
    t(s, ti, cx+Inches(0.32), cy+Inches(0.22), cw2-Inches(0.6), Inches(0.45), sz=15.5, bold=True, col=RED)
    t(s, body, cx+Inches(0.32), cy+Inches(0.78), cw2-Inches(0.6), Inches(1.4), sz=11.5, col=GRY, line_sp=1.28)

# ════════ SLIDE 7 — 比較表 ════════
s = sl(); ft(s)
hdr(s, "COMPARISON", "今の方式 vs GitHub移行後 ── 比較表", "同じ「2台Mac同期」でも、裏側のインフラが変わるだけ。日々の使い方はほぼ同じ")
rows = [
    ("比較項目", "今：Drive上のgit金庫", "移行後：GitHub（非公開リポ）"),
    ("費用", "0円（Drive契約内）", "0円（個人非公開リポは無料）"),
    ("同時書き込みの安全性", "低い（今回のような破損が起きうる）", "高い（専用インフラが競合を正式処理）"),
    ("壊れた時の直しやすさ", "内部ファイルを手作業で調査・修復が必要", "GitHub側が状態を保証、基本的に壊れない"),
    ("日々の使い方", "pull→作業→push（今と同じ）", "pull→作業→push（変化なし）"),
    ("導入の手間", "なし（対応済み）", "初回のみアカウント作成・認証設定が必要"),
]
light_table(s, rows, Inches(0.55), Inches(1.85), Inches(12.23), Inches(4.4),
            [Inches(3.4), Inches(4.4), Inches(4.43)], hi_col=2, sz=12.5, header_sz=13)

# ════════ SLIDE 8 — 今後の方針 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "CONCLUSION", Inches(0.9), Inches(0.7), Inches(11), Inches(0.4), sz=14, bold=True, col=RED)
t(s, "今後の方針", Inches(0.9), Inches(1.15), Inches(11.7), Inches(0.8), sz=32, bold=True, col=INK)
t(s, "利益直結モード中は、営業直結タスクを最優先。GitHub移行は「内務」なので今は着手しない。",
  Inches(0.92), Inches(2.0), Inches(11.5), Inches(0.5), sz=14, col=GRY)
items = [
    ("① 今すぐは移行しない", "実害は毎回AIが自動で直せる範囲。営業直結タスク優先の方針を崩さない。"),
    ("② push確認は毎回聞かない", "fast-forwardで安全に完了するpushは自動実行。データが本当にぶつかる時だけ確認する。"),
    ("③ 「main 2」的な破損は自動復旧", "祖先コミットへの参照だけと確認できれば、AIが単独で削除・復旧してよいルールを秘書室に明記済み。"),
    ("④ 移行の判断トリガー", "同種の破損が繰り返し発生する／2台稼働の頻度が増えるなら、その時点で菊池さんの明示指示で着手する。"),
]
for i, (ti, ds) in enumerate(items):
    cy = Inches(2.75) + Inches(1.02) * i
    bx(s, Inches(0.9), cy, Inches(11.5), Inches(0.88), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.9), cy, Inches(0.06), Inches(0.88), RED)
    t(s, ti, Inches(1.15), cy+Inches(0.1), Inches(3.4), Inches(0.68), sz=13.5, bold=True, col=RED, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.1)
    t(s, ds, Inches(4.65), cy+Inches(0.1), Inches(7.55), Inches(0.68), sz=12, col=INK, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.2)
bx(s, Inches(0.9), Inches(6.85), Inches(11.5), Pt(1.2), LINE)
t(s, "詳細ルール：.company/secretary/CLAUDE.md「💻マルチMac運用ルール」に反映済み（2026-07-02）", Inches(0.9), Inches(6.95), Inches(11), Inches(0.4), sz=10.5, col=GRY)

prs.save("git_explainer_260702.pptx")
print("saved git_explainer_260702.pptx  /  slides:", len(prs.slides._sldIdLst))

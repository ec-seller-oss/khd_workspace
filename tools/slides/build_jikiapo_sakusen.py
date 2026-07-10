"""
次アポ作戦・チーム共有（YES最速→納品／役割／タスク／アイコール道筋）。
クリーム白×レンガ赤。出力: jikiapo_sakusen.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック 次アポ作戦・チーム共有（内部）",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.09); cell.margin_right=Inches(0.06); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=(TEALBG if is_hi else (CARD if ri%2==1 else BG))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or (ci==0) or is_hi
                    r.font.color.rgb=(WHT if ri==0 else (TEALD if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A))))
    return tb

# ════ 1 矯正 ════
s=sl(); ft(s)
hdr(s,"軌道修正","次アポの本質 ── ゴールは「先生のYES」。実機デモは主役じゃない","宮崎の指摘で本質に戻す。実機を見せてもYESは取れない。")
bx(s,Inches(0.55),Inches(2.05),Inches(12.23),Inches(1.5),REDBG,line=RED,lw=1.5); bx(s,Inches(0.55),Inches(2.05),Inches(0.12),Inches(1.5),RED)
t(s,"実機（LINE・Stream Deck）を見せても、YESは貰えない。",Inches(0.85),Inches(2.2),Inches(11.6),Inches(0.5),sz=18,bold=True,col=REDD)
t(s,"次アポの目的＝提案に「YES」を最速でもらう → そのまま納品（設定・運用開始）へ繋げる。",Inches(0.85),Inches(2.78),Inches(11.6),Inches(0.5),sz=15,bold=True,col=INK)
bx(s,Inches(0.55),Inches(3.85),Inches(12.23),Inches(2.0),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(3.85),Inches(12.23),Inches(0.06),RED)
t(s,"⚠ 自分（菊池）のズレ ＝ サボり場所が間違っていた",Inches(0.82),Inches(4.0),Inches(11.7),Inches(0.4),sz=15,bold=True,col=REDD)
t(s,"・作り込み・実機準備（楽しい所）に時間をかけ、本丸の「YES取り（クロージング）」と「日程の主導」を宮崎任せにしていた。\n・宮崎に握らせて自分が怠っていた＝一番大事な所をサボっていた。\n・矯正：次アポは「作って見せる」でなく「YESを取って納品に進める」。日程もクロージングも、菊池が自分で握る。",
  Inches(0.85),Inches(4.5),Inches(11.6),Inches(1.3),sz=12.5,col=INK,line_sp=1.3)

# ════ 2 役割 ════
s=sl(); ft(s)
hdr(s,"当日の役割","当日、誰が何をやるか","主役は「実機」でなく「YESを取る会話」。それぞれの持ち場を明確に。")
roles=[("先生（山崎先生）","YESを出す相手。「まず1つ、やってみよう」を引き出す。決裁者として握る。",REDBG,REDD),
       ("医療事務","ヒアリング相手。運用の声・ダメ出しをもらう。「使う人」の納得を作る。",TEALBG,TEALD),
       ("宮崎","技術担当。実機は当日無理→動画/口頭で補足。主役でなく裏方。",CARD,INK),
       ("菊池","主導・握り役。日程もクロージングも自分で取る（人任せにしない）。",CARD,INK)]
cw,ch,gx,gy=Inches(6.0),Inches(2.0),Inches(0.45),Inches(0.35); x0,y0=Inches(0.55),Inches(2.0)
for i,(ti,ds,fill,tc) in enumerate(roles):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,fill,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.12),ch,RED)
    t(s,ti,cx+Inches(0.3),cy+Inches(0.22),cw-Inches(0.5),Inches(0.5),sz=16,bold=True,col=tc)
    t(s,ds,cx+Inches(0.32),cy+Inches(0.9),cw-Inches(0.6),Inches(0.9),sz=12,col=INK,line_sp=1.2)

# ════ 3 YESと決めること ════
s=sl(); ft(s)
hdr(s,"GET THE YES","次アポで取りに行く「YES」と、その場で決めること","YESが出たら、間髪入れず納品（設定）の日程まで握る")
bx(s,Inches(0.55),Inches(2.0),Inches(12.23),Inches(1.0),TEALBG,line=TEALD,lw=1.5); bx(s,Inches(0.55),Inches(2.0),Inches(0.12),Inches(1.0),TEALD)
t(s,"取りに行くYES：「まず1つ（再診リマインド or FAQ自動応答）を無料で試す」を、その場で握る。",Inches(0.85),Inches(2.0),Inches(11.6),Inches(1.0),sz=15,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE)
deci=[("① どの1つから","再診リマインド／FAQ、どちらを最初に？"),
      ("② 効果の指標","医療事務と「何を測るか」を1つ決める"),
      ("③ 次回の日程","納品（設定・運用開始）の日をその場で押さえる"),
      ("④ 進める窓口","誰と進めるか（事務長・担当者）を確定")]
cw,ch,gx,gy=Inches(6.0),Inches(1.3),Inches(0.45),Inches(0.3); x0,y0=Inches(0.55),Inches(3.3)
for i,(ti,ds) in enumerate(deci):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.28),cy+Inches(0.14),cw-Inches(0.5),Inches(0.4),sz=14,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.3),cy+Inches(0.6),cw-Inches(0.55),Inches(0.6),sz=12,col=INK,line_sp=1.1)
t(s,"※「持ち帰って検討」で終わらせない。小さい1つだから、その場でYESを取りに行く。",Inches(0.55),Inches(6.55),Inches(12),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 4 日程の打ち手 ════
s=sl(); ft(s)
hdr(s,"日程：待たない","宮崎の予定で「待たない」── 菊池が行く","実機が主役でない＝宮崎不在でも菊池が単独で行ける")
bx(s,Inches(0.55),Inches(2.0),Inches(12.23),Inches(1.2),CARD,line=CARDLN,lw=1.0)
t(s,"状況：山崎先生の候補（6/25 昼・夜／6/26 昼／7/3 昼・夜）→ 宮崎は全部埋まり。宮崎「待つとどんどん遅れる」。",
  Inches(0.82),Inches(2.1),Inches(11.7),Inches(1.0),sz=13,col=INK,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)
bx(s,Inches(0.55),Inches(3.4),Inches(12.23),Inches(1.7),REDBG,line=RED,lw=1.5); bx(s,Inches(0.55),Inches(3.4),Inches(0.12),Inches(1.7),RED)
t(s,"打ち手：菊池が 6/25 か 6/26 に「単独で」行く。",Inches(0.85),Inches(3.55),Inches(11.6),Inches(0.5),sz=18,bold=True,col=REDD)
t(s,"・実機デモは主役でない → 宮崎不在でも成立。宮崎はStream Deckの動画/説明材料だけ用意。\n・宮崎の空き待ちで機を逃さない。YESを取りに行くのは菊池の仕事。\n・どうしても日程が合わなければ、その場で新候補を即もらう（持ち帰らない）。",
  Inches(0.85),Inches(4.05),Inches(11.6),Inches(1.0),sz=12.5,col=INK,line_sp=1.3)
t(s,"→ 今日中に山崎先生へ「6/25 or 6/26 で伺います」と返す。",Inches(0.55),Inches(5.4),Inches(12),Inches(0.4),sz=13,bold=True,col=REDD)

# ════ 5 タスク化 ════
s=sl(); ft(s)
hdr(s,"TASKS","タスク化 ── 誰が・いつまでに","アポごとに「決めること」を残す。曖昧にしない。")
rows=[
 ("タスク","担当","期限"),
 ("山崎先生へ日程確定（6/25 or 6/26で伺う意向＋可否）","菊池","今日中"),
 ("提案資料9枚・福井A4・実機LINE（問診自動URL）を当日セット","菊池","前日"),
 ("Stream Deckの動画/説明材料を用意（当日見せられない分）","宮崎","次アポまで"),
 ("当日：YESを取る → 納品（設定）日程まで握る","菊池","当日"),
 ("アイコール連携 3点確認（並行・YESには必須でない）","宮崎","随時"),
 ("固有値・個人情報同意文言・KPIを現地で確定","菊池","当日"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(4.2),[Inches(7.6),Inches(2.4),Inches(2.23)],hi_col=None,sz=12,header_sz=12.5)
t(s,"→ 次回アポ後も、この形（タスク＋担当＋期限＋決めたこと）でチーム共有を回す。",Inches(0.55),Inches(6.35),Inches(12),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 6 前提整理：i-CALLとは／一般クリニック比較 ════
s=sl(); ft(s)
hdr(s,"前提整理（チーム共有・必読）","i-CALL とは？ 一般クリニックの受付と、京橋は何が違うか","知らないと提案がブレる。まずチーム全員（福井さん含む）が同じ絵を持つ。")
bx(s,Inches(0.55),Inches(1.9),Inches(12.23),Inches(1.2),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(1.9),Inches(0.1),Inches(1.2),RED)
t(s,"i-CALL（アイコール）＝ 診療の「順番受付・予約」システム",Inches(0.82),Inches(2.0),Inches(11.7),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"患者さんがスマホ/PC/電話で順番を取り、混雑を確認、順番が近づくと通知（クラウド型・全国1,200施設）。\n→ 京橋の医療事務さんが「アイコールで順番取っても来ない」と課題に挙げている＝京橋は実際にこれで「当日の順番受付」をしている（事務の生の声・推測でない）。",
  Inches(0.84),Inches(2.42),Inches(11.6),Inches(0.7),sz=11.5,col=INK,line_sp=1.25)
rows=[("受付の観点","一般的なクリニック（予約制が多い）","京橋クリニック（当日順番受付制）"),
      ("受付の方法","時間を指定して事前に予約","来院順／WEBで順番を取得（i-CALL）"),
      ("使う仕組み","予約システム","順番受付システム（i-CALL）"),
      ("患者さんの困り","予約が取れない・キャンセル","順番取っても来ない・待ち時間が不明"),
      ("公式LINEで足すこと","予約確定・前日リマインド","順番が近づくと通知・待ち人数の表示")]
light_table(s,rows,Inches(0.55),Inches(3.28),Inches(12.23),Inches(2.35),[Inches(2.9),Inches(4.6),Inches(4.73)],hi_col=2,sz=11,header_sz=11.5)
bx(s,Inches(0.55),Inches(5.82),Inches(12.23),Inches(1.0),TEALBG,line=TEALD,lw=1.0); bx(s,Inches(0.55),Inches(5.82),Inches(0.1),Inches(1.0),TEALD)
t(s,"LINE・電子カルテとの相性：i-CALLは通知に「電話・メール・LINE」対応＋電子カルテ約30社連携（京橋のMedicom含む）。だからLINEとも電カルとも「相性よく」提案できる。※京橋のプランで実際にLINE連携できるか／どの公式LINEに出すかは要確認。",
  Inches(0.82),Inches(5.82),Inches(11.7),Inches(1.0),sz=11,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.2)

# ════ 7 役割分担：置き換える？分ける？ ════
s=sl(); ft(s)
hdr(s,"役割の明確化","i-CALL と 公式LINE ── 置き換える？ 役割を分ける？","結論：技術的には置き換えも可能。でも京橋は「役割を分けて共存」が正解。")
bx(s,Inches(0.55),Inches(1.95),Inches(12.23),Inches(0.72),REDBG,line=RED,lw=1.0); bx(s,Inches(0.55),Inches(1.95),Inches(0.1),Inches(0.72),RED)
t(s,"置き換えも「できなくはない」。が、京橋は既にi-CALL稼働中＝役割を分けて共存が正解。",Inches(0.82),Inches(1.95),Inches(11.7),Inches(0.72),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)
bx(s,Inches(0.55),Inches(2.85),Inches(6.0),Inches(2.0),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(2.85),Inches(6.0),Inches(0.5),RED)
t(s,"i-CALL ＝ 順番受付エンジン（受付の心臓部）",Inches(0.7),Inches(2.88),Inches(5.7),Inches(0.45),sz=13,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"役割：採番・順番管理／電話の自動受付／電カル(Medicom)連携／混雑のリアルタイム表示\n強み：医療現場で実績1,200施設。止められない受付の核。",Inches(0.78),Inches(3.5),Inches(5.55),Inches(1.25),sz=11,col=INK,line_sp=1.3)
bx(s,Inches(6.78),Inches(2.85),Inches(6.0),Inches(2.0),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(2.85),Inches(6.0),Inches(0.5),TEALD)
t(s,"公式LINE ＝ 患者との接点",Inches(6.93),Inches(2.88),Inches(5.7),Inches(0.45),sz=13,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"役割：Web問診・予約・FAQ自動応答・お知らせ・リマインド・追客\n強み：患者と繋がる・読まれる・無料で柔軟に直せる。",Inches(7.01),Inches(3.5),Inches(5.55),Inches(1.25),sz=11,col=INK,line_sp=1.3)
bx(s,Inches(0.55),Inches(5.0),Inches(12.23),Inches(1.5),GRYBG); bx(s,Inches(0.55),Inches(5.0),Inches(0.1),Inches(1.5),RED)
t(s,"なぜ「置き換え」でなく「役割分担」か",Inches(0.82),Inches(5.08),Inches(11.7),Inches(0.35),sz=13,bold=True,col=REDD)
t(s,"・京橋は既にi-CALL稼働中＝受付の心臓部。引っこ抜く＝受付を止めるリスク＋「まず1つ無料で小さく」に反する。\n・我々の今のLINEに「採番・整理券」機能は無い→置き換えるなら新規開発が要る（大工事）。電話受付の高齢者も切ってしまう。\n・だから〈i-CALL＝順番の中身／公式LINE＝患者接点〉と分け、「順番通知だけ」をLINEに寄せれば理想（統一）。",Inches(0.85),Inches(5.45),Inches(11.7),Inches(1.0),sz=10.5,col=INK,line_sp=1.25)
t(s,"※ 置き換えがアリなのは「i-CALLが無い新規クリニック」や「i-CALL費用を見直す」時。京橋では今やらない。",Inches(0.55),Inches(6.58),Inches(12.2),Inches(0.35),sz=10.5,bold=True,col=REDD)

# ════ 8 アイコール道筋（参考） ════
s=sl(); ft(s)
hdr(s,"参考：順番のlive化","アイコール(i-CALL)連携の道筋 ── 菊池の理解用","「順番をLINEで」は、我々が作るより i-CALL のLINE機能を活かすのが筋")
steps=[("①","アイコール＝順番受付","京橋が使う診療予約システム i-CALL。通知は電話/メール/LINEに対応"),
       ("②","作るでなく活かす","i-CALLは元々LINE通知できる。順番のlive化は「i-CALL機能を使う」が早い"),
       ("③","宮崎が3点確認","プランにLINE通知あるか／既存の公式LINEに紐付け可か／API有無"),
       ("④","統一 or 住み分け","統一できれば1つのLINEで順番も予約も問診も。無理なら住み分け")]
cw,gx,x0=Inches(2.95),Inches(0.18),Inches(0.55)
for i,(no,ti,ds) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.2),cw,Inches(2.6),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.2),cw,Inches(0.06),RED)
    bx(s,cx+Inches(0.2),Inches(2.4),Inches(0.5),Inches(0.5),RED,shape=MSO_SHAPE.OVAL)
    t(s,no,cx+Inches(0.2),Inches(2.4),Inches(0.5),Inches(0.5),sz=16,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ti,cx+Inches(0.15),Inches(3.05),cw-Inches(0.3),Inches(0.6),sz=13,bold=True,col=INK,align=PP_ALIGN.CENTER,line_sp=1.05)
    t(s,ds,cx+Inches(0.2),Inches(3.7),cw-Inches(0.4),Inches(1.0),sz=10.5,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.15)
t(s,"※ これは「将来像」。次アポのYESには必須でない＝先に出しすぎない（売り込みすぎ注意）。中央ビジコムMedicomもi-CALL連携対応で電カル側の整合は取りやすい。",
  Inches(0.55),Inches(5.1),Inches(12.2),Inches(0.6),sz=11.5,bold=True,col=REDD,line_sp=1.2)

prs.save("jikiapo_sakusen.pptx")
print("saved jikiapo_sakusen.pptx slides:",len(prs.slides._sldIdLst))

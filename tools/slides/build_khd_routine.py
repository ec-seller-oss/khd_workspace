# -*- coding: utf-8 -*-
"""KHD 経営ルーティン運用ガイド スライド生成 2026-06-04
毎朝/週次(事業部報告会)/月次(家族会議) を初見でも分かるよう網羅・構造化。
デザイン=クリーム白#F9F6EF × レンガ赤#AA2E26（KHD標準）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM = RGBColor(0xF9,0xF6,0xEF); BRICK = RGBColor(0xAA,0x2E,0x26)
INK = RGBColor(0x2B,0x2B,0x2B); SUB = RGBColor(0x6B,0x6B,0x6B)
GREEN = RGBColor(0x2E,0x7D,0x5B); GOLD = RGBColor(0xC9,0x9A,0x2E)
WHITE = RGBColor(0xFF,0xFF,0xFF); LINE = RGBColor(0xE3,0xDD,0xD2)

prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0,0, SW, SH); r.fill.solid(); r.fill.fore_color.rgb=bg
    r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s,x,y,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=4,line_sp=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[[(runs,18,INK,False)]]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.line_spacing=line_sp
        for (t,sz,col,bold) in para:
            r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name='Hiragino Kaku Gothic ProN'
    return tb

def titlebar(s,kicker,title):
    box(s,0,0,13.333,1.15,fill=BRICK)
    txt(s,0.6,0.12,12,0.4,[[(kicker,12,RGBColor(0xF0,0xD6,0xD2),True)]])
    txt(s,0.6,0.40,12.1,0.7,[[(title,26,WHITE,True)]])

def chip(s,x,y,w,text,col,tcol=WHITE,h=0.42,sz=12):
    box(s,x,y,w,h,fill=col); txt(s,x,y+0.02,w,h,[[(text,sz,tcol,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ===== 1 表紙 =====
s=slide(BRICK)
box(s,0,0,13.333,7.5,fill=BRICK)
box(s,0.9,2.0,11.5,3.4,fill=CREAM)
txt(s,1.3,2.35,10.7,0.5,[[("KHD 経営ルーティン 運用ガイド",30,BRICK,True)]])
txt(s,1.3,3.05,10.7,0.5,[[("毎朝 ・ 週次(事業部報告会) ・ 月次(家族会議) で家族と事業を回す",16,INK,False)]])
txt(s,1.3,3.7,10.7,1.4,[[("・初めて見る人（家族・本部長）でも、これ1冊で「いつ・どこを見て・何をするか」が分かる",13,SUB,False)],
    [("・正本は1つ＝Googleスプレッドシート「2026_KHD PJ一覧_v2」。迷ったらここへ戻る",13,SUB,False)],
    [("・合言葉：①統合司令塔だけ見る／月初は📸1クリック／時間は営業へ寄せる",13,SUB,False)]])
txt(s,1.3,5.0,10.7,0.4,[[("菊池ホールディングス（KHD）  2026-06　v1",12,SUB,False)]])

# ===== 2 なぜ要る =====
s=slide(); titlebar(s,"WHY ・ なぜこの仕組みが要るか","「現在地」を常に見て、時間を営業に寄せるため")
txt(s,0.6,1.5,12.1,0.5,[[("KHDの事業群が唯一の収入源（安定給与ゼロ）。だから――",15,INK,True)]])
items=[("①  毎日 現在地を知る","あと何ヶ月もつ(ランウェイ)・黒字か。不安でなく数字で経営する。",GREEN),
       ("②  内務を減らし営業へ","ツール磨きは0円。時間を商談・追客・仕込みに寄せる(目標60%)。",BRICK),
       ("③  家族で同じ数字を見る","妻も同じ画面で納得→『今これに集中』が家族の合意になる。",GOLD)]
for i,(h,d,c) in enumerate(items):
    y=2.2+i*1.45; box(s,0.6,y,12.1,1.25,fill=WHITE,line=LINE)
    box(s,0.6,y,0.16,1.25,fill=c)
    txt(s,1.0,y+0.18,11.4,0.5,[[(h,18,c,True)]])
    txt(s,1.0,y+0.68,11.4,0.5,[[(d,14,INK,False)]])

# ===== 3 全体像 =====
s=slide(); titlebar(s,"OVERVIEW ・ 全体像","3つのリズム ＋ 1つの正本(SSoT)")
# SSoT中央
box(s,4.9,1.5,3.5,1.0,fill=BRICK)
txt(s,4.9,1.62,3.5,0.8,[[("正本 SSoT",13,RGBColor(0xF0,0xD6,0xD2),True)],[("KHD PJ一覧_v2 (スプレッドシート)",12,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
cards=[("毎朝 (5分)","現在地を見る","①統合司令塔で\nランウェイ/損益/今日の営業1点",GREEN,0.6),
       ("週次","事業部報告会","6本部KPI・時間配分(②)\nパイプライン更新",GOLD,4.9),
       ("月次","家族会議＋月初締め","📸残高反映→純資産・ランウェイ\n妻と資源配分を決める",BRICK,9.2)]
for (k,t,d,c,x) in cards:
    box(s,x,3.1,3.5,2.8,fill=WHITE,line=LINE); box(s,x,3.1,3.5,0.7,fill=c)
    txt(s,x,3.18,3.5,0.6,[[(k,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.2,3.95,3.1,0.5,[[(t,16,c,True)]],align=PP_ALIGN.CENTER)
    txt(s,x+0.2,4.55,3.1,1.2,[[(ln,13,INK,False)] for ln in d.split("\n")],align=PP_ALIGN.CENTER,line_sp=1.1)
txt(s,0.6,6.2,12.1,0.5,[[("すべて同じ正本を見る。日々→週→月で粒度が上がるだけ。迷ったら正本に戻る。",13,SUB,False)]],align=PP_ALIGN.CENTER)

# ===== 4 ツール =====
s=slide(); titlebar(s,"TOOLS ・ 使う道具","どのツールを・何に使うか")
tools=[("📗 Googleスプレッドシート","正本(SSoT)。BS/資金繰り/損益/借入/司令塔。数字はここに集約。"),
       ("💰 「財務」メニュー(自作)","スプシ上部の専用ボタン。残高反映・スナップショット等を1クリック。"),
       ("📥 MF(マネーフォワード)","口座/カード残高の自動取得元。Claudeが読んで反映。"),
       ("🗓 Googleカレンダー / Tasks","予定と今日の締切タスク。実績時間も自動集計。"),
       ("🗒 Notion","成果物・議事録・ルール・作業ログの保管(検索で全部辿れる)。"),
       ("🤖 Claude(私) + Chrome","MF読取・明細DL・資料作成を代行。月初の残高更新を担当。")]
for i,(h,d) in enumerate(tools):
    col=i%2; row=i//2; x=0.6+col*6.15; y=1.55+row*1.75
    box(s,x,y,5.9,1.55,fill=WHITE,line=LINE)
    txt(s,x+0.25,y+0.18,5.4,0.5,[[(h,16,BRICK,True)]])
    txt(s,x+0.25,y+0.72,5.4,0.7,[[(d,13,INK,False)]],line_sp=1.05)

# ===== 5 大原則 =====
s=slide(); titlebar(s,"RULES ・ 大原則","迷ったらこの5つに従う")
rules=[("1","①統合司令塔だけ見る","普段はここだけ。現在地(ランウェイ/損益)が1枚で分かる。"),
       ("2","月初は『📸今日のスナップショット』を1クリック","MF残高反映＋その日の記録。手入力ほぼゼロ。"),
       ("3","信金(朝日/大東京/法人TB)は2ヶ月毎に手入力","MF非対応。窓口/電話で取得→🏦タブに打つだけ。"),
       ("4","営業の入金見込みは『売上見込み』に1行足す","案件が出たらここに。資金繰りへ自動反映。"),
       ("5","成果物・決定はNotionへ。推測の数字は書かない","証票(MF/明細)が根拠。ベタ打ちで埋めない。")]
for i,(n,h,d) in enumerate(rules):
    y=1.55+i*1.08; box(s,0.6,y,12.1,0.95,fill=WHITE,line=LINE)
    box(s,0.6,y,0.9,0.95,fill=BRICK); txt(s,0.6,y,0.9,0.95,[[(n,24,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.7,y+0.13,10.8,0.5,[[(h,16,INK,True)]])
    txt(s,1.7,y+0.55,10.8,0.4,[[(d,13,SUB,False)]])

# ===== 6 毎朝 =====
s=slide(); titlebar(s,"DAILY ・ 毎朝ルーティン(約5分)","現在地を見て、今日の営業を1点決める")
steps=[("STEP 1","スプシ『①統合司令塔』を開く","ランウェイ(あと何ヶ月)・全社損益・時間配分%を確認"),
       ("STEP 2","Google Tasksで今日の締切を見る","期限つきタスク。営業直結を最優先に並べ替え"),
       ("STEP 3","今日の『営業1点』を宣言","例：福井宛送信／栄町決済プッシュ／医療追客1件")]
for i,(k,h,d) in enumerate(steps):
    y=1.7+i*1.55; box(s,0.6,y,12.1,1.35,fill=WHITE,line=LINE); box(s,0.6,y,1.7,1.35,fill=GREEN)
    txt(s,0.6,y,1.7,1.35,[[(k,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,2.5,y+0.22,10,0.5,[[(h,17,INK,True)]])
    txt(s,2.5,y+0.78,10,0.4,[[(d,13,SUB,False)]])
txt(s,0.6,6.5,12.1,0.5,[[("ポイント：財務タブをいじらない。見るだけ→手は営業へ。",13,BRICK,True)]],align=PP_ALIGN.CENTER)

# ===== 7 週次 事業部報告会 =====
s=slide(); titlebar(s,"WEEKLY ・ 週次 事業部報告会","6本部のKPIと時間配分を点検")
txt(s,0.6,1.4,12.1,0.4,[[("アジェンダ（順に・15〜30分）",14,BRICK,True)]])
ag=[("① 現在地共有","①統合司令塔：ランウェイ・損益・営業直結%(目標60%)"),
    ("② 6本部KPI",  "01経営/02資金/03事業/04コンサル/05物件＋秘書室。各①数字②次の打ち手"),
    ("③ パイプライン更新","『売上見込み』で確度・着金月を更新→資金繰りに反映"),
    ("④ 時間配分レビュー","②本部マトリクス：実績h(カレンダー自動)。内務過多なら営業へ寄せる"),
    ("⑤ 来週の集中1つ","本部横断で『最も収益に効く1手』を決める")]
for i,(h,d) in enumerate(ag):
    y=1.95+i*1.0; box(s,0.6,y,12.1,0.88,fill=WHITE,line=LINE); box(s,0.6,y,0.16,0.88,fill=GOLD)
    txt(s,0.95,y+0.1,3.3,0.6,[[(h,15,GOLD,True)]])
    txt(s,4.3,y+0.13,8.2,0.6,[[(d,13,INK,False)]],line_sp=1.05)
txt(s,0.6,7.0,12.1,0.4,[[("使うタブ：①統合司令塔 / ②本部マトリクス / 売上見込み",12,SUB,False)]],align=PP_ALIGN.CENTER)

# ===== 8 月次 家族会議＋月初締め =====
s=slide(); titlebar(s,"MONTHLY ・ 月次 家族会議＋月初締め","数字を締めて、家族で資源配分を決める")
txt(s,0.6,1.4,12.1,0.4,[[("手順（毎月1日・約20分）",14,BRICK,True)]])
ms=[("1","残高を締める","『💰財務→📸今日のスナップショット』1クリック(MF反映＋当月記録)。信金は2ヶ月毎に🏦へ手入力"),
    ("2","純資産・ランウェイ確認","③BS(純資産)・④資金繰り(現金/ランウェイ)。先6ヶ月で谷(ショート)が無いか"),
    ("3","損益を確認","⑦損益：事業が黒字か・損益分岐に届くか。EC粗利は実績連動"),
    ("4","家族で資源配分","②マトリクス(時間×金×家族)を見て『来月どこに時間とお金を割くか』を妻と合意"),
    ("5","記録","決定事項をNotionへ。次月の家族会議で振り返る")]
for i,(n,h,d) in enumerate(ms):
    y=1.95+i*1.02; box(s,0.6,y,12.1,0.9,fill=WHITE,line=LINE)
    box(s,0.6,y,0.85,0.9,fill=BRICK); txt(s,0.6,y,0.85,0.9,[[(n,22,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.65,y+0.1,3.0,0.6,[[(h,15,INK,True)]])
    txt(s,4.7,y+0.13,7.9,0.65,[[(d,12.5,SUB,False)]],line_sp=1.03)

# ===== 9 数字の読み方 =====
s=slide(); titlebar(s,"LITERACY ・ 数字の読み方(初めての人へ)","この4つだけ分かればOK")
defs=[("ランウェイ","今の現金で、あと何ヶ月もつか。3未満=危険/6以上=安全。今は約11ヶ月。",GREEN),
      ("損益分岐","事業が黒字になる粗利ライン。粗利がこれを超えれば事業はプラス。",BRICK),
      ("純資産","資産−負債。家族の本当の財産。月次で増えてるかを見る。",GOLD),
      ("粗利","売上−原価。EC等が稼いだ実額。固定費を賄えるかの源泉。",INK)]
for i,(h,d,c) in enumerate(defs):
    col=i%2; row=i//2; x=0.6+col*6.15; y=1.6+row*2.35
    box(s,x,y,5.9,2.1,fill=WHITE,line=LINE); box(s,x,y,5.9,0.7,fill=c)
    txt(s,x,y+0.1,5.9,0.5,[[(h,18,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.3,y+0.95,5.3,1.0,[[(d,14,INK,False)]],line_sp=1.15)

# ===== 10 自動 vs 手動 =====
s=slide(); titlebar(s,"AUTO/MANUAL ・ 自動と手動","何が自動で、何を手で入れるか")
box(s,0.6,1.55,6.0,5.2,fill=WHITE,line=LINE); box(s,0.6,1.55,6.0,0.65,fill=GREEN)
txt(s,0.6,1.62,6.0,0.5,[[("🟢 自動（触らなくていい）",16,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.95,2.4,5.3,4.2,[[("・MF連携の銀行/証券/ポイント残高",13,INK,False)],
    [("・EC粗利（販売管理表→自動集計）",13,INK,False)],
    [("・借入返済・残高（⑤借入=証票連動）",13,INK,False)],
    [("・売上見込み→資金繰りの入金",13,INK,False)],
    [("・総資産・自己資本・ランウェイ(数式)",13,INK,False)],
    [("→ 月初に『📸』を押すだけで最新化",13,GREEN,True)]],line_sp=1.3)
box(s,6.75,1.55,5.95,5.2,fill=WHITE,line=LINE); box(s,6.75,1.55,5.95,0.65,fill=GOLD)
txt(s,6.75,1.62,5.95,0.5,[[("🟡 手動（人が入れる・理由あり）",16,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,7.1,2.4,5.3,4.2,[[("・信金(朝日/大東京/法人TB)＝MF非対応・2ヶ月毎",13,INK,False)],
    [("・MF連携不調(みずほ/PayPay等)＝要再連携",13,INK,False)],
    [("・証券の銘柄内訳/Amazonポイント＝MFは合算のみ",13,INK,False)],
    [("・現金・売掛金・現物＝実数を見て入力",13,INK,False)],
    [("・法人の案件売上/原価＝成約時に入力",13,INK,False)],
    [("→ BSのBL/BM列に『理由』を常時表示",13,GOLD,True)]],line_sp=1.3)

# ===== 11 タブの地図 =====
s=slide(); titlebar(s,"MAP ・ タブの歩き方","①統合司令塔を中心に")
mp=[("🎯 毎日見る","①統合司令塔",GREEN),
    ("📊 中身を見る","③BS / ④資金繰り / ⑦損益 / ②本部マトリクス",INK),
    ("🔧 月初に更新","📥MF残高 / 🏦残高クイック入力",BRICK),
    ("🗄 データ源(自動)","⑤借入 / 売上見込み / 役員報酬 / 諸経費 / 税金",SUB),
    ("🧹 整理候補","借入返済・借入条件(旧) / 🕛KPI / パス(機密)",GOLD)]
for i,(k,t,c) in enumerate(mp):
    y=1.6+i*1.04; box(s,0.6,y,12.1,0.9,fill=WHITE,line=LINE); box(s,0.6,y,3.4,0.9,fill=c)
    txt(s,0.6,y,3.4,0.9,[[(k,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.2,y+0.25,8.3,0.5,[[(t,15,INK,True)]])

# ===== 12 まとめ =====
s=slide(BRICK); box(s,0,0,13.333,7.5,fill=BRICK)
txt(s,0.9,0.9,11.5,0.7,[[("まとめ ― これだけ覚える",30,WHITE,True)]])
box(s,0.9,2.0,11.5,4.2,fill=CREAM)
txt(s,1.3,2.35,10.7,3.6,[
    [("毎朝　",18,BRICK,True),("①統合司令塔を見て、今日の営業を1点決める",16,INK,False)],
    [("週次　",18,GOLD,True),("事業部報告会＝6本部KPI＋時間配分(②)＋パイプライン更新",16,INK,False)],
    [("月初　",18,BRICK,True),("『📸』1クリックで残高反映→純資産/ランウェイ→妻と資源配分",16,INK,False)],
    [("信金　",18,SUB,True),("朝日/大東京/法人TB だけ2ヶ月毎に手入力",15,INK,False)],
    [("原則　",18,GREEN,True),("数字は証票が根拠・成果物はNotion・時間は営業へ",16,INK,False)],
    [("",8,INK,False)],
    [("迷ったら → 正本『KHD PJ一覧_v2』の ①統合司令塔 に戻る",15,BRICK,True)]],line_sp=1.25,sp_after=10)

prs.save('/Users/kikuchikenta/01_honbu_docs_automation/KHD_経営ルーティン運用ガイド.pptx')
print("OK saved")

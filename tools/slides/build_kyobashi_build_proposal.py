"""
京橋 公式LINE＋HP/MEO 作り込み提案（GO判断用）。
問診票(実票)動線・i-CALL↔LINE導線分離・他メニュー(メモ)・HP/MEO(さいとう内科TTP)。
クリーム白×レンガ赤。出力: kyobashi_build_proposal.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
BLUEBG=RGBColor(0xE6,0xF1,0xFB); BLUED=RGBColor(0x18,0x5F,0xA5)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=21,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.46),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide,n):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック 公式LINE＋HP/MEO 作り込み提案（GO判断用）",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
    t(slide,n,Inches(12.4),H-Inches(0.42),Inches(0.5),Inches(0.32),sz=9,col=GRY)

# ════ 1 表紙 ════
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"京橋クリニック 御中",Inches(0.9),Inches(1.1),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"公式LINE ＋ HP/MEO\n作り込み提案（GO判断用）",Inches(0.88),Inches(1.65),Inches(11.7),Inches(1.7),sz=28,bold=True,col=INK,line_sp=1.12)
t(s,"いただいた問診票・面談シート・前回メモをもとに、具体的な作り込みに落とし込みました。",Inches(0.92),Inches(3.75),Inches(11.4),Inches(0.5),sz=13.5,col=GRY)
bx(s,Inches(0.9),Inches(4.55),Inches(11.5),Inches(1.2),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(4.55),Inches(0.1),Inches(1.2),RED)
t(s,"このGOをいただけたら、まず「問診票」から作り込みに着手します。\nアイコール（順番受付）はそのまま活かし、LINEは役割を分けて重ねます。",Inches(1.2),Inches(4.55),Inches(11),Inches(1.2),sz=14,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)

# ════ 2 全体像（3レイヤー役割分担） ════
s=sl(); ft(s,"2")
hdr(s,"全体像","3つのレイヤーで、役割を分ける","それぞれの持ち場を分け、重ねる。置き換えない。")
layers=[("集患（入口）","HP ・ MEO（Googleマップ）","患者さんに見つけてもらう。さいとう内科さんの型をTTP",BLUEBG,BLUED),
        ("患者接点（つなぐ）","公式LINE","問診・次回予約・お知らせ・リマインド・FAQ・追客",REDBG,REDD),
        ("受付（院内）","アイコール（順番受付）","当日の順番・採番・電話受付・混雑表示。受付の心臓部＝そのまま",CARD,INK)]
y=Inches(2.05)
for ti,name,ds,fill,tc in layers:
    bx(s,Inches(0.55),y,Inches(12.23),Inches(1.45),fill,line=CARDLN,lw=1.0); bx(s,Inches(0.55),y,Inches(0.12),Inches(1.45),tc)
    t(s,ti,Inches(0.8),y+Inches(0.16),Inches(2.6),Inches(1.1),sz=13,bold=True,col=tc,anchor=MSO_ANCHOR.MIDDLE)
    t(s,name,Inches(3.5),y+Inches(0.2),Inches(3.6),Inches(1.05),sz=16,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
    t(s,ds,Inches(7.2),y+Inches(0.2),Inches(5.4),Inches(1.05),sz=12,col=INK,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.2)
    y=y+Inches(1.6)
t(s,"→ まず真ん中（公式LINE）の「問診」から。次にHP/MEO、と段階で。",Inches(0.55),Inches(6.95)-Inches(0.0),Inches(12),Inches(0.0),sz=11,bold=True,col=REDD) if False else None

# ════ 3 i-CALL↔LINE 導線の分け方 ════
s=sl(); ft(s,"3")
hdr(s,"導線の分け方（重要）","アイコール と LINE、どう分けるか","順番には触らない。LINEは「その手前と後」を担う。")
bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(3.1),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(0.55),GRY)
t(s,"アイコール（変えない）",Inches(0.8),Inches(2.06),Inches(5.5),Inches(0.4),sz=14,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"・当日の順番受付・採番\n・電話の自動受付（高齢の方）\n・混雑・待ち人数の表示\n・電子カルテ(Medicom)連携\n＝受付の心臓部。そのまま稼働。",Inches(0.82),Inches(2.7),Inches(5.5),Inches(2.3),sz=12.5,col=INK,line_sp=1.4)
bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(3.1),REDBG,line=RED,lw=1.0); bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(0.55),RED)
t(s,"公式LINE（足す）",Inches(7.03),Inches(2.06),Inches(5.5),Inches(0.4),sz=14,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"・来院「前」：Web問診\n・来院「後」：次回予約・リマインド\n・お知らせ・FAQ自動応答・追客\n・離脱しかけの方の掘り起こし\n＝順番には触らない。手前と後ろを担う。",Inches(7.03),Inches(2.7),Inches(5.5),Inches(2.3),sz=12.5,col=INK,line_sp=1.4)
bx(s,Inches(0.55),Inches(5.35),Inches(12.23),Inches(1.05),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(5.35),Inches(0.1),Inches(1.05),TEALD)
t(s,"重なるのは「順番のお知らせ」だけ。将来、アイコールのLINE通知機能で公式LINEに寄せられれば1本化（要確認）。それまでは住み分け。",Inches(0.82),Inches(5.35),Inches(11.8),Inches(1.05),sz=12.5,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)

# ════ 4 ①問診のLINE動線 ════
s=sl(); ft(s,"4")
hdr(s,"① 問診のLINE動線（ここから着手）","友だち追加 → 問診を選んで → スマホで記入 → 受付に届く","来院前に記入完了で、受付の聞き取り・転記・紙が減る")
steps=[("友だち追加","QR/リンクから"),("問診を選ぶ","一般／睡眠時無呼吸／健診／化学物質過敏症"),
       ("スマホで記入","来院前に自宅で"),("送信","自動でお礼＋タグ"),("受付に届く","一覧で確認・カルテへ")]
cw,gx,x0=Inches(2.3),Inches(0.18),Inches(0.55)
for i,(ti,ds) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.15),cw,Inches(1.7),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.15),cw,Inches(0.06),RED)
    t(s,str(i+1),cx+Inches(0.1),Inches(2.25),Inches(0.5),Inches(0.4),sz=13,bold=True,col=RED)
    t(s,ti,cx+Inches(0.1),Inches(2.62),cw-Inches(0.2),Inches(0.5),sz=13,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.12),Inches(3.12),cw-Inches(0.24),Inches(0.65),sz=9.5,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.1)
    if i<4: t(s,"→",cx+cw-Inches(0.02),Inches(2.6),Inches(0.3),Inches(0.5),sz=16,bold=True,col=RED,align=PP_ALIGN.CENTER)
bx(s,Inches(0.55),Inches(4.2),Inches(12.23),Inches(1.75),GRYBG); bx(s,Inches(0.55),Inches(4.2),Inches(0.1),Inches(1.75),RED)
t(s,"いただいた実問診票をベースに（2種類）",Inches(0.82),Inches(4.3),Inches(11.6),Inches(0.35),sz=13,bold=True,col=REDD)
t(s,"・①一般問診票（症状・既往・アレルギー・喫煙飲酒・いびき/無呼吸・ジェネリック希望・当院を何で知ったか 等）\n・②ESS／Epworth（睡眠時無呼吸・CPAP用。8つの状況の眠気を4段階→合計11点以上で疑い）\n→ 様式は京橋さんのものをそのまま電子化。斉藤内科さん（デジスマ）の「迷わない流れ」を参考に。",
  Inches(0.85),Inches(4.7),Inches(11.6),Inches(1.2),sz=11.5,col=INK,line_sp=1.3)

# ════ 5 問診の中身＋紙削減 ════
s=sl(); ft(s,"5")
hdr(s,"① 問診の中身（実票どおり）＋ おまけの効果","聞いている項目はそのまま。デジタルにするだけ","受付の手間と「紙」を、まとめて減らせる")
bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(3.0),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(0.55),RED)
t(s,"一般問診票（そのまま電子化）",Inches(0.8),Inches(2.06),Inches(5.5),Inches(0.4),sz=13,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"症状・発症時期／体温／既往歴・家族歴／アレルギー／喫煙・飲酒／いびき・無呼吸／妊娠（女性）／ジェネリック希望／★当院を何で知ったか（＝集患の分析データに）",Inches(0.82),Inches(2.7),Inches(5.5),Inches(2.2),sz=12,col=INK,line_sp=1.4)
bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(3.0),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(0.55),RED)
t(s,"ESS（睡眠時無呼吸・CPAP）",Inches(7.03),Inches(2.06),Inches(5.5),Inches(0.4),sz=13,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"身長・体重＋8つの状況の眠気を4段階（0〜3）で。合計が自動計算され、11点以上は要注意の目安に。\n→ CPAP対象の拾い上げ・定期フォローに直結。",Inches(7.03),Inches(2.7),Inches(5.5),Inches(2.2),sz=12,col=INK,line_sp=1.4)
bx(s,Inches(0.55),Inches(5.2),Inches(12.23),Inches(0.85),REDBG); bx(s,Inches(0.55),Inches(5.2),Inches(0.1),Inches(0.85),RED)
t(s,"おまけ：本人確認の「マイナンバー紙コピー」も、事前のデジタル取得で減らせます（事務の声に対応）。",Inches(0.82),Inches(5.2),Inches(11.8),Inches(0.85),sz=12.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 6 ②LINEの他メニュー（メモベース） ════
s=sl(); ft(s,"6")
hdr(s,"② LINEの他メニュー（前回メモから）","問診の次に、効く順で広げる","看護師長・事務の皆さまの声をそのまま反映")
rows=[("メニュー","中身","効果"),
      ("再診リマインド（デフォルトON）","再診・定期の方に来院の目安をお知らせ","取りこぼし防止（看護師長のご意見）"),
      ("CPAP 次回予約 自動送付","毎回聞けていなかった次回予約を事前設定で自動案内","月末の駆け込み・残業を減らす"),
      ("再診促し（ザオラル→LINE）","離脱しかけの方へLINEで次回予約導線","事務の架電を減らし、不在でも予約が取れる"),
      ("化学物質過敏症の問診","専用問診→先生がお電話フォロー","LINEが特に響く層。証明書運用と連動"),
      ("健診の問診","健診用フォーム","受付の手間を軽く"),
      ("次回予約をきちんと取る","LINEで確実に取得・記録","取りこぼし防止＋国へのアピールにも")]
n=len(rows); tb=s.shapes.add_table(n,3,Inches(0.55),Inches(2.0),Inches(12.23),Inches(4.0)).table
tb.first_row=False; tb.horz_banding=False; tb.columns[0].width=Inches(3.7); tb.columns[1].width=Inches(5.0); tb.columns[2].width=Inches(3.53)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Inches(0.1); cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
        cell.fill.solid(); cell.fill.fore_color.rgb=(RED if ri==0 else (CARD if ri%2==1 else BG))
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name=FONT; r.font.size=Pt(11 if ri else 12); r.font.bold=(ri==0 or ci==0)
                r.font.color.rgb=(WHT if ri==0 else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))

# ════ 7 ③HP・MEO（さいとう内科TTP） ════
s=sl(); ft(s,"7")
hdr(s,"③ HP・MEO（集患の入口・b）","さいとう内科さんの「勝ち型」をTTPして内製化","作るだけでなく、見つけてもらう。LINE/問診への入口にも。")
wins=[("多層の予約導線＋常時CTA","電話／WEB予約／WEB問診／LINEを、画面に常に出す（フローティング）"),
      ("予約専用ダイヤルを目立たせる","ヘッダーに独立配置。迷わせない"),
      ("MEO：NAPを反復","住所・電話・最寄駅(徒歩◯分)を3回以上。診療科をカード化"),
      ("信頼づくり","院長経歴・提携病院・院長ブログ/通信を月次更新。右下にAIチャットボット")]
cw,ch,gx,gy=Inches(6.0),Inches(1.6),Inches(0.23),Inches(0.3); x0,y0=Inches(0.55),Inches(2.05)
for i,(ti,ds) in enumerate(wins):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,BLUED)
    t(s,ti,cx+Inches(0.25),cy+Inches(0.18),cw-Inches(0.45),Inches(0.5),sz=13.5,bold=True,col=BLUED)
    t(s,ds,cx+Inches(0.27),cy+Inches(0.72),cw-Inches(0.5),Inches(0.8),sz=11.5,col=INK,line_sp=1.25)
t(s,"→ HP・MEOは「集患」レーン。LINE・問診と繋いで、見つける→つながる→来院、を一本に。",Inches(0.55),Inches(5.7),Inches(12.2),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 8 作り込みの段取り ════
s=sl(); ft(s,"8")
hdr(s,"作り込みの段取り","フェーズで、無理なく","まず問診。効果を見ながら次へ。")
ph=[("P1（まず）","Web問診（一般・ESS）＋ 再診リマインド（ON）","テンプレ反映→実物→受付に届くまで"),
    ("P2（次に）","CPAP自動・再診促し・化学物質過敏症・健診","効いた所から順に追加"),
    ("P3（その後）","HP・MEO（集患の入口）","さいとう内科TTPで内製化")]
cw,gx,x0=Inches(3.96),Inches(0.18),Inches(0.55)
for i,(p,sol,ds) in enumerate(ph):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.05),cw,Inches(2.5),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.05),cw,Inches(0.65),RED)
    t(s,p,cx+Inches(0.15),Inches(2.14),cw-Inches(0.3),Inches(0.48),sz=14,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    t(s,sol,cx+Inches(0.2),Inches(2.85),cw-Inches(0.4),Inches(0.9),sz=12.5,bold=True,col=REDD,align=PP_ALIGN.CENTER,line_sp=1.2)
    t(s,ds,cx+Inches(0.25),Inches(3.75),cw-Inches(0.5),Inches(0.7),sz=11,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.2)
bx(s,Inches(0.55),Inches(4.75),Inches(12.23),Inches(1.2),GRYBG); bx(s,Inches(0.55),Inches(4.75),Inches(0.1),Inches(1.2),RED)
t(s,"進め方：問診テンプレは新さん（事務窓口）とやり取り。開始日は福井さんと調整。\nスタートは来週〜。アイコールは止めず、LINEは順番に触らずに重ねます。",Inches(0.82),Inches(4.75),Inches(11.8),Inches(1.2),sz=12.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.3)

# ════ 9 効果測定 ════
s=sl(); ft(s,"9")
hdr(s,"効果の測り方","一緒に「数字」で確かめる","評価のためでなく、楽になったかを見るため")
rows=[("見る指標","どう良くなるか"),
      ("Web問診の事前記入率","来院前に取得→受付の聞き取り・転記・紙が減る"),
      ("受付の電話本数 / 日","リマインド・FAQで一次対応→件数が減る"),
      ("CPAP定期の来院率","次回予約の自動案内→取りこぼしが減る"),
      ("月末の残業時間","通院の平準化→月末の山・19時残業が減る"),
      ("新患数 / HP・MEO経由の流入","HP/MEOで見つけてもらう→新患が増える")]
n=len(rows); tb=s.shapes.add_table(n,2,Inches(0.55),Inches(2.0),Inches(12.23),Inches(3.6)).table
tb.first_row=False; tb.horz_banding=False; tb.columns[0].width=Inches(4.8); tb.columns[1].width=Inches(7.43)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Inches(0.12); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        cell.fill.solid(); cell.fill.fore_color.rgb=(RED if ri==0 else (CARD if ri%2==1 else BG))
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name=FONT; r.font.size=Pt(12.5); r.font.bold=(ri==0 or ci==0)
                r.font.color.rgb=(WHT if ri==0 else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
bx(s,Inches(0.55),Inches(5.85),Inches(12.23),Inches(0.6),GRYBG); bx(s,Inches(0.55),Inches(5.85),Inches(0.1),Inches(0.6),RED)
t(s,"導入前のいまの数字（電話本数・残業など）を、先に軽く取っておきます（ベースライン）。",Inches(0.82),Inches(5.85),Inches(11.8),Inches(0.6),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 10 GO判断 ════
s=sl(); ft(s,"10")
hdr(s,"GO判断","この方向で、作り込みに入ってよいか","GOをいただけたら、まず問診から着手します")
bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(2.7),REDBG,line=RED,lw=1.0); bx(s,Inches(0.55),Inches(2.0),Inches(0.1),Inches(2.7),RED)
t(s,"先に教えてほしい2点",Inches(0.82),Inches(2.12),Inches(5.5),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"①「最後の1週間」に架電、の運用の意味\n（再診促しの設計に効きます）\n\n② サインエコーのLINE導入、とは？\n（どんなご要望か、もう少し）",Inches(0.82),Inches(2.65),Inches(5.5),Inches(1.95),sz=12.5,col=INK,line_sp=1.35)
bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(2.7),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(2.0),Inches(0.1),Inches(2.7),RED)
t(s,"GOをもらったら",Inches(7.03),Inches(2.12),Inches(5.5),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"・問診テンプレを反映して、まず問診の実物を作る\n・受付に届くところまで通す\n・開始日は福井さんと調整\n・新さん（事務）と細部をやり取り",Inches(7.03),Inches(2.65),Inches(5.5),Inches(1.95),sz=12.5,col=INK,line_sp=1.35)
bx(s,Inches(0.55),Inches(4.95),Inches(12.23),Inches(1.0),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(4.95),Inches(0.1),Inches(1.0),TEALD)
t(s,"アイコールは活かす。LINEは順番に触らず重ねる。まず問診から、無料で小さく。これでGO、でよろしいですか？",Inches(0.82),Inches(4.95),Inches(11.8),Inches(1.0),sz=14,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("kyobashi_build_proposal.pptx")
print("saved kyobashi_build_proposal.pptx slides:",len(prs.slides._sldIdLst))

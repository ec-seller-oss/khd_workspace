"""
京橋 公式LINE ゴールイメージ（菊池が一目で分かる用）。
①リッチメニューのタップ先3画面（予約/Web問診/FAQ）②再診リマインド3本 の"完成形"を見せる。
クリーム白×レンガ赤。出力: kyobashi_line_goal.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF)
LNG=RGBColor(0x06,0xC7,0x55); LNGD=RGBColor(0x06,0x7A,0x35); BUB=RGBColor(0x8D,0xE0,0x55)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック ｜ 公式LINE ゴールイメージ ｜ AI医療コンサル",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def band(slide,y,text,col=REDBG,bar=RED,tc=REDD):
    bx(slide,Inches(0.55),y,Inches(12.23),Inches(0.62),col)
    bx(slide,Inches(0.55),y,Inches(0.1),Inches(0.62),bar)
    t(slide,text,Inches(0.82),y,Inches(11.8),Inches(0.62),sz=12.5,bold=True,col=tc,anchor=MSO_ANCHOR.MIDDLE)
def arrow(slide,x,y,w,h):
    s=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,h); s.fill.solid()
    s.fill.fore_color.rgb=RGBColor(0xC9,0x6A,0x62); s.line.fill.background(); s.shadow.inherit=False; return s
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=REDBG if is_hi else (CARD if ri%2==1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or is_hi or (ci==0)
                    if ri==0: r.font.color.rgb=WHT
                    elif ci==0: r.font.color.rgb=INK
                    else: r.font.color.rgb=RGBColor(0x3A,0x3A,0x3A)
    return tb

def phone(slide,x,y,w,h,title):
    bx(slide,x,y,w,h,WHT,line=CARDLN,lw=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bx(slide,x,y,w,Inches(0.5),LNG)
    t(slide,title,x+Inches(0.15),y,w-Inches(0.3),Inches(0.5),sz=11,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
    return y+Inches(0.62)

def rowbox(slide,x,y,w,label,val):
    bx(slide,x,y,w,Inches(0.5),GRYBG)
    t(slide,label,x+Inches(0.15),y,Inches(1.3),Inches(0.5),sz=10,col=GRY,anchor=MSO_ANCHOR.MIDDLE)
    t(slide,val,x+Inches(1.35),y,w-Inches(1.5),Inches(0.5),sz=11,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)

def bubble(slide,x,y,w,text,sz=10,fill=WHT,tc=INK):
    nlines=max(1,len(text)//16+1); h=Inches(0.36)+Inches(0.22)*(nlines-1)
    bx(slide,x,y,w,h,fill,line=CARDLN if fill==WHT else None,lw=0.75)
    t(slide,text,x+Inches(0.12),y,w-Inches(0.24),h,sz=sz,col=tc,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.05)
    return y+h+Inches(0.12)

# ════════ SLIDE 1 — 患者×クリニック マトリクス（宮崎指摘：視点を分ける）════════
s=sl(); ft(s)
hdr(s,"THE GOAL","完成したらこうなる ── 来院の流れに沿って、両方が変わる","左から「問い合わせ→予約→来院前→来院→次回へ」の順。各段で患者の体験とクリニックの効果を分けて整理（宮崎指摘を反映）")
rows=[
 ("視点","① 問い合わせ","② 予約","③ 来院前","④ 来院・受付","⑤ 次回へ"),
 ("患者は","質問に即・自動で回答","LINEで予約・変更","スマホで事前に問診","あと何人かをLINEで通知","再診をリマインド"),
 ("クリニックは","難しい質問だけ電話＝本数減","電話予約が減る","聞き取り・転記が消える","クレーム対応が減る","取りこぼし防止＝残業減"),
]
light_table(s,rows,Inches(0.45),Inches(2.05),Inches(12.43),Inches(2.6),[Inches(1.7),Inches(2.15),Inches(2.15),Inches(2.15),Inches(2.14),Inches(2.14)],hi_col=None,sz=11,header_sz=11.5)
t(s,"患者は「便利・待ち時間が減る」／クリニックは「電話・転記・取りこぼしが減って残業が減る」。同じ仕組みの「裏表」。",Inches(0.45),Inches(4.9),Inches(12.4),Inches(0.4),sz=12,bold=True,col=REDD)
band(s,Inches(5.45),"今できているのは、この表の中身ぜんぶの「実物」（スライド・リッチメニュー・動くデモ4本・リマインド3本・FAQ・問診・予約）。")
band(s,Inches(6.25),"残るはクリニック固有値の確定と、LIVE化の同意（先生の「作ってOK」＋LINE登録の電話認証＝菊池の1タップ）。",col=GRYBG,bar=RED,tc=INK)

# ════════ SLIDE 2 — なぜ電話が減るのか（自動応答チャットボット）※宮崎指摘 ════════
TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
s=sl(); ft(s)
hdr(s,"WHY FEWER CALLS","なぜ電話が減るのか ── 自動応答チャットボット","問い合わせの大半をLINEが即・自動で解決。難しい質問だけ電話へ。受付の一次対応をLINEが肩代わりする")
sy=Inches(2.05); sh=Inches(2.45)
# ステージ1：問い合わせ
x1=Inches(0.6); w1=Inches(2.95)
bx(s,x1,sy,w1,sh,CARD,line=CARDLN,lw=1.0); bx(s,x1,sy,w1,Inches(0.06),RED)
t(s,"患者の問い合わせ",x1,sy+Inches(0.18),w1,Inches(0.4),sz=14,bold=True,col=INK,align=PP_ALIGN.CENTER)
for i,q in enumerate(["「今日やってる?」","「予約したい」","「持ち物は?」","「本人確認は?」"]):
    t(s,q,x1+Inches(0.32),sy+Inches(0.72)+Inches(0.4)*i,w1-Inches(0.6),Inches(0.38),sz=11.5,col=GRY)
arrow(s,x1+w1+Inches(0.03),sy+Inches(0.95),Inches(0.32),Inches(0.55))
# ステージ2：自動応答
x2=Inches(4.0); w2=Inches(3.05)
bx(s,x2,sy,w2,sh,REDBG,line=RED,lw=1.2); bx(s,x2,sy,w2,Inches(0.06),RED)
t(s,"LINEが自動で返信",x2,sy+Inches(0.2),w2,Inches(0.4),sz=15,bold=True,col=REDD,align=PP_ALIGN.CENTER)
t(s,"キーワード／FAQで\n即・24時間・待たせず回答\n（本人確認・受付ルール・\n診療時間・予約方法 など）",x2+Inches(0.25),sy+Inches(0.78),w2-Inches(0.5),Inches(1.5),sz=11.5,col=INK,align=PP_ALIGN.CENTER,line_sp=1.2)
arrow(s,x2+w2+Inches(0.03),sy+Inches(0.95),Inches(0.32),Inches(0.55))
# ステージ3：分岐（解決／電話）
x3=Inches(7.55); w3=Inches(5.2)
bx(s,x3,sy,w3,Inches(1.15),TEALBG,line=CARDLN,lw=1.0); bx(s,x3,sy,Inches(0.1),Inches(1.15),TEALD)
t(s,"その場で解決 → 電話ゼロ",x3+Inches(0.3),sy+Inches(0.16),w3-Inches(0.5),Inches(0.4),sz=14,bold=True,col=TEALD)
t(s,"よくある質問の大半（例：7〜8割）はLINEで完結。患者は電話不要・待ち時間ゼロ。",x3+Inches(0.3),sy+Inches(0.58),w3-Inches(0.55),Inches(0.5),sz=11,col=INK,line_sp=1.05)
bx(s,x3,sy+Inches(1.3),w3,Inches(1.15),GRYBG,line=CARDLN,lw=1.0); bx(s,x3,sy+Inches(1.3),Inches(0.1),Inches(1.15),RED)
t(s,"難しい質問だけ → お電話をご案内",x3+Inches(0.3),sy+Inches(1.46),w3-Inches(0.5),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"自動で解決しない件のみ「お電話はこちら」。受付が出る電話が必要な分だけに絞られる。",x3+Inches(0.3),sy+Inches(1.88),w3-Inches(0.55),Inches(0.5),sz=11,col=INK,line_sp=1.05)
# 電話本数 before/after（イメージ）
t(s,"受付が対応する電話の本数（イメージ）",Inches(0.6),Inches(4.75),Inches(6),Inches(0.35),sz=12,bold=True,col=INK)
t(s,"現状",Inches(0.6),Inches(5.16),Inches(1.1),Inches(0.4),sz=11,col=GRY,anchor=MSO_ANCHOR.MIDDLE)
bx(s,Inches(1.7),Inches(5.16),Inches(9.5),Inches(0.4),RED); t(s,"全部 受付が出る",Inches(1.85),Inches(5.16),Inches(9),Inches(0.4),sz=10.5,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"導入後",Inches(0.6),Inches(5.66),Inches(1.1),Inches(0.4),sz=11,col=GRY,anchor=MSO_ANCHOR.MIDDLE)
bx(s,Inches(1.7),Inches(5.66),Inches(9.5),Inches(0.4),GRYBG)
bx(s,Inches(1.7),Inches(5.66),Inches(2.4),Inches(0.4),RED); t(s,"難しい分だけ",Inches(1.85),Inches(5.66),Inches(4),Inches(0.4),sz=10.5,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
t(s,"LINEが解決（電話なし）",Inches(4.3),Inches(5.66),Inches(6),Inches(0.4),sz=10.5,col=GRY,anchor=MSO_ANCHOR.MIDDLE)
band(s,Inches(6.4),"アンケート直結：事務の「鳴り止まない電話・クレーム対応」を、LINEの自動応答が一次対応で巻き取る → 受付が出る電話の本数が減る。")

# ════════ SLIDE 3 — ①タップ先3画面のゴール ════════
s=sl(); ft(s)
hdr(s,"患者が使う画面 ①","リッチメニューを押した先 ── この3画面の中身をつくる","【患者側の体験】完成後、患者はLINE内でここまで完結 → 【クリニック側】受付の聞き取り・紙→電カル転記が減る")
px=[Inches(0.95),Inches(5.0),Inches(9.05)]; pw=Inches(3.3); ph=Inches(3.95); py=Inches(2.0)
# 予約
cy=phone(s,px[0],py,pw,ph,"ご予約・再診")
rowbox(s,px[0]+Inches(0.2),cy,pw-Inches(0.4),"診療科","一般内科 ▾"); cy+=Inches(0.6)
rowbox(s,px[0]+Inches(0.2),cy,pw-Inches(0.4),"希望日","6/28(土) ▾"); cy+=Inches(0.6)
rowbox(s,px[0]+Inches(0.2),cy,pw-Inches(0.4),"時間帯","10:00 ▾"); cy+=Inches(0.68)
bx(s,px[0]+Inches(0.2),cy,pw-Inches(0.4),Inches(0.5),RED); t(s,"予約を確定する",px[0]+Inches(0.2),cy,pw-Inches(0.4),Inches(0.5),sz=12,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# Web問診
cy=phone(s,px[1],py,pw,ph,"Web問診（来院前）")
for lab in ["□ 発熱・咳・鼻水","□ 胸の症状（CPAP）","既往歴：＿＿＿＿＿","服用中の薬：＿＿＿"]:
    bx(s,px[1]+Inches(0.2),cy,pw-Inches(0.4),Inches(0.42),GRYBG); t(s,lab,px[1]+Inches(0.32),cy,pw-Inches(0.5),Inches(0.42),sz=10.5,col=INK,anchor=MSO_ANCHOR.MIDDLE); cy+=Inches(0.5)
bx(s,px[1]+Inches(0.2),cy,pw-Inches(0.4),Inches(0.5),RED); t(s,"送信する",px[1]+Inches(0.2),cy,pw-Inches(0.4),Inches(0.5),sz=12,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# FAQ
cy=phone(s,px[2],py,pw,ph,"よくある質問")
for lab in ["› 本人確認・マイナ保険証","› 受付・順番のルール","› 化学物質過敏症のご相談","› 休診日・診療時間"]:
    bx(s,px[2]+Inches(0.2),cy,pw-Inches(0.4),Inches(0.5),CARD,line=CARDLN,lw=0.75); t(s,lab,px[2]+Inches(0.32),cy,pw-Inches(0.55),Inches(0.5),sz=10.5,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE); cy+=Inches(0.58)
labs=["① 予約フォーム（予約URL）","② Web問診フォーム","③ FAQページ"]
for i,l in enumerate(labs):
    t(s,l,px[i],Inches(6.05),pw,Inches(0.35),sz=11,bold=True,col=RED,align=PP_ALIGN.CENTER)
band(s,Inches(6.5),"①でつくるのは、この3画面の「中身」。アンケート直結：予約=CPAP離脱対策／問診=待ち時間短縮／FAQ=本人確認・化学物質過敏症の電話を減らす。")

# ════════ SLIDE 3 — ②再診リマインド3本のゴール ════════
s=sl(); ft(s)
hdr(s,"患者に届くメッセージ ②","自動で届く ── 再診リマインド3本の中身をつくる","【患者側】患者のLINEにこう届く → 【クリニック側】CPAPの取りこぼしが減り、月末の残業が減る")
# 左：LINEスレッド
px0,py0,pw0,ph0=Inches(0.95),Inches(2.0),Inches(4.5),Inches(4.55)
bx(s,px0,py0,pw0,ph0,RGBColor(0x97,0xB3,0xCE),line=CARDLN,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
bx(s,px0,py0,pw0,Inches(0.5),LNG); t(s,"京橋クリニック",px0+Inches(0.15),py0,pw0-Inches(0.3),Inches(0.5),sz=11,bold=True,col=WHT,anchor=MSO_ANCHOR.MIDDLE)
cy=py0+Inches(0.66); bw=pw0-Inches(0.6)
cy=bubble(s,px0+Inches(0.25),cy,bw,"【京橋クリニック】CPAPをご利用の皆さまへ。今月の定期診察の時期です。データ確認と処方のため、ご予約をお願いします ▶ ご予約はこちら",sz=9.5)
cy=bubble(s,px0+Inches(0.25),cy,bw,"ご予約ありがとうございます。6/28(土)10:00 でお取りしました。変更は『予約変更』とお送りください。",sz=9.5,fill=BUB,tc=RGBColor(0x13,0x36,0x0A))
cy=bubble(s,px0+Inches(0.25),cy,bw,"明日6/28(土)10:00 のご予約です。保険証をお持ちください。来院前のWeb問診で待ち時間が短くなります ▶ 問診を始める",sz=9.5)
# 右：3本の狙い
rx=Inches(5.85); rw=Inches(6.9); ry=Inches(2.0)
mtx=[("① CPAP定期リマインド","タイミング：毎月初／対象=CPAP・定期タグ","狙い：毎月来ない人の取りこぼし防止＝月末まとめ会計の残業を削る（効果が数字で出る一手）"),
     ("② 予約確認（自動返信）","タイミング：予約された直後","狙い：予約の取り違え・無断キャンセルを防ぐ。変更導線もLINE内に集約"),
     ("③ 前日リマインド","タイミング：予約前日18時","狙い：来院率UP＋Web問診へ誘導で当日の受付・待ち時間を短縮")]
for i,(ti,tm,ai) in enumerate(mtx):
    cy2=ry+(Inches(1.42))*i
    bx(s,rx,cy2,rw,Inches(1.28),CARD,line=CARDLN,lw=1.0); bx(s,rx,cy2,Inches(0.1),Inches(1.28),RED)
    t(s,ti,rx+Inches(0.3),cy2+Inches(0.12),rw-Inches(0.5),Inches(0.4),sz=14,bold=True,col=INK)
    t(s,tm,rx+Inches(0.3),cy2+Inches(0.55),rw-Inches(0.5),Inches(0.3),sz=10.5,bold=True,col=RED)
    t(s,ai,rx+Inches(0.3),cy2+Inches(0.84),rw-Inches(0.5),Inches(0.4),sz=10.5,col=GRY,line_sp=1.05)
band(s,Inches(6.65),"②でつくるのは、この3本の「文面」。配信はエルメ無料で着手 → 来院率と月末残業を1〜2ヶ月測れば、成果報酬の根拠になる。")

prs.save("kyobashi_line_goal_v4.pptx")
print("saved kyobashi_line_goal_v4.pptx  slides:",len(prs.slides._sldIdLst))

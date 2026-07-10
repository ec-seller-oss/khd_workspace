"""
京橋 全体設計：システム×動線 体系比較（人の動きまで）。
システム構成の体系化＋動線(初診/無呼吸CPAP/再診)の現状→LINE後を人の動きで比較。
クリーム白×レンガ赤。出力: kyobashi_taikei.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
BLUED=RGBColor(0x18,0x5F,0xA5)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sh=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=col
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False; return sh
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.55),Inches(0.35),Inches(12),Inches(0.4),sz=12.5,bold=True,col=RED)
    bx(slide,Inches(0.57),Inches(0.72),Inches(1.6),Pt(3),RED)
    t(slide,main,Inches(0.55),Inches(0.82),Inches(12.2),Inches(0.55),sz=20,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.57),Inches(1.36),Inches(12.2),Inches(0.3),sz=11,col=GRY)
def ft(slide,n):
    bx(slide,Inches(0.5),H-Inches(0.45),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック 全体設計：システム×動線 体系比較（人の動きまで）",Inches(0.5),H-Inches(0.4),Inches(10),Inches(0.3),sz=8.5,col=GRY)
    t(slide,n,Inches(12.45),H-Inches(0.4),Inches(0.4),Inches(0.3),sz=8.5,col=GRY)
def tbl(slide,rows,x,y,w,h,col_w,sz=10,hsz=11,hi=None):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            c=tb.cell(ri,ci); c.text=str(val); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.05); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            c.fill.solid(); ishi=(hi is not None and ci==hi)
            if ri==0: c.fill.fore_color.rgb=(TEALD if ishi else RED)
            else: c.fill.fore_color.rgb=(TEALBG if ishi else (CARD if ri%2==1 else BG))
            for p in c.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(hsz if ri==0 else sz); r.font.bold=(ri==0 or ci==0)
                    r.font.color.rgb=(WHT if ri==0 else (TEALD if ishi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A))))
    return tb

# 1 表紙
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"京橋クリニック 御中",Inches(0.9),Inches(1.15),Inches(11),Inches(0.45),sz=15,bold=True,col=RED)
t(s,"全体設計：システム × 動線\n体系比較（人の動きまで）",Inches(0.88),Inches(1.7),Inches(11.7),Inches(1.7),sz=26,bold=True,col=INK,line_sp=1.12)
t(s,"どのシステムが何を担い、患者さん・スタッフの動きが、現状からどう変わるか。",Inches(0.92),Inches(3.8),Inches(11.4),Inches(0.5),sz=13.5,col=GRY)
bx(s,Inches(0.9),Inches(4.6),Inches(11.5),Inches(1.15),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(4.6),Inches(0.1),Inches(1.15),RED)
t(s,"システムの体系（誰が何を担う）→ 動線の比較（人の動き）。\n初診／睡眠時無呼吸(CPAP)／再診の3つの動きを、現状とLINE後で並べます。",Inches(1.2),Inches(4.6),Inches(11),Inches(1.15),sz=13.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)

# 2 システム構成（体系）
s=sl(); ft(s,"2")
hdr(s,"システムの体系","どのシステムが、何を担うか（5つの層）","役割を層で分ける。重ねるが、置き換えない。")
layers=[("集患層","HP・MEO（Googleマップ）","見つけてもらう（検索・地図・口コミ）",RED),
        ("接点層","公式LINE","問診・予約・お知らせ・リマインド・追客",RED),
        ("受付層","i-CALL（アイコール）","当日の順番受付・電話受付・混雑表示",GRY),
        ("記録層","電子カルテ Medicom","診療の記録・レセプト（将来：問診を取込）",GRY),
        ("管理層","CPAP機器・遠隔モニタリング","CPAPの使用データを確認（機器側）",GRY)]
y=Inches(1.95)
for nm,sysn,role,ac in layers:
    bx(s,Inches(0.55),y,Inches(12.23),Inches(0.92),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),y,Inches(0.12),Inches(0.92),ac)
    t(s,nm,Inches(0.8),y,Inches(2.0),Inches(0.92),sz=13,bold=True,col=(REDD if ac==RED else GRY),anchor=MSO_ANCHOR.MIDDLE)
    t(s,sysn,Inches(2.85),y,Inches(3.6),Inches(0.92),sz=13.5,bold=True,col=INK,anchor=MSO_ANCHOR.MIDDLE)
    t(s,role,Inches(6.6),y,Inches(6.0),Inches(0.92),sz=11.5,col=INK,anchor=MSO_ANCHOR.MIDDLE)
    y=y+Inches(1.0)
t(s,"赤＝今回うちが足す/磨く層（集患・接点）。グレー＝既存のまま（受付・記録・機器）。",Inches(0.55),Inches(6.95)-Inches(0.0),Inches(12),Inches(0.0),sz=11,bold=True,col=REDD) if False else t(s,"赤＝今回うちが足す/磨く層（集患・接点）。グレー＝既存のまま（受付・記録・機器）。",Inches(0.55),Inches(6.95),Inches(12.2),Inches(0.3),sz=10.5,bold=True,col=REDD)

# 3 システム比較表
s=sl(); ft(s,"3")
hdr(s,"システム比較（機能ごと）","機能 × いま × LINE導入後 × 担い手","どの機能を、どのシステムが担うか。重複と空白をなくす。")
rows=[
 ("機能","いま","LINE導入後","担い手"),
 ("集患","口コミ・看板頼み","HP/MEOで検索・地図流入","HP・MEO"),
 ("問い合わせ","電話に集中","LINE FAQ自動＋難問のみ電話","公式LINE"),
 ("問診","紙→カルテ転記","来院前にWeb問診","LINE→電カル(将来)"),
 ("予約・順番","電話・i-CALL","i-CALL(順番)＋LINE(健診/予約)","i-CALL／LINE"),
 ("お知らせ・リマインド","掲示・口頭(読まれない)","LINEで自動配信(届く)","公式LINE"),
 ("再診・CPAP管理","人力(記憶・架電)","タグ＋次回予約・自動リマインド","公式LINE"),
 ("記録","紙＋Medicomで二度入力","Web問診をMedicomへ(将来)","電カル"),
]
tbl(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(4.4),[Inches(2.4),Inches(3.4),Inches(3.9),Inches(2.53)],sz=10.5,hsz=11,hi=2)

# 4 動線① 初診
s=sl(); ft(s,"4")
hdr(s,"動線① 初診の患者","人の動き：現状 → LINE後（患=患者／事=受付）","同じ来院でも、誰がどこで動くかが変わる")
rows=[
 ("段階","現状（人の動き）","LINE後（人の動き）"),
 ("見つける","患:口コミ/たまたま→電話で確認","患:検索/地図でHP→LINE友だち追加"),
 ("問い合わせ","患:電話　事:1件ずつ口頭対応","患:LINEで自己解決　事:難問だけ対応"),
 ("来院前","（何もなし）","患:スマホでWeb問診を記入"),
 ("受付","患:紙問診を記入　事:カルテへ転記","患:記入済で受付　事:転記ほぼゼロ"),
 ("診察〜会計","事:本人確認の紙コピー等","事:情報が揃い、聞き直しが減る"),
]
tbl(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(3.7),[Inches(2.2),Inches(5.0),Inches(5.03)],sz=10.5,hsz=11,hi=2)
bx(s,Inches(0.55),Inches(5.75),Inches(12.23),Inches(0.7),GRYBG); bx(s,Inches(0.55),Inches(5.75),Inches(0.1),Inches(0.7),RED)
t(s,"変わる人の動き：患者は「電話・紙」→「スマホで完結」。受付は「転記・聞き取り」が消える。",Inches(0.82),Inches(5.75),Inches(11.8),Inches(0.7),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# 5 動線② 無呼吸 CPAP
s=sl(); ft(s,"5")
hdr(s,"動線② 睡眠時無呼吸（CPAP）の患者 ★","ここが京橋の最大の詰まり：CPAP定期に来ない→月末残業","標準のSASパスを参考（斉藤内科の詳細は非公開）")
rows=[
 ("段階","現状（人の動き）","LINE後（人の動き）"),
 ("気づき・受診","患:いびき/眠気で来院","患:同じ＋LINEでESS事前記入"),
 ("スクリーニング","事:紙のESSを記入・採点","患:ESSを事前回答→自動採点で候補に印"),
 ("簡易検査","事:口頭で案内・日程調整","患:LINEで日程案内・リマインド"),
 ("CPAP導入","医:導入・説明","医:同じ（人が行う）"),
 ("定期フォロー","患:次回を忘れ来ない　事:月末に架電・残業","患:次回予約をLINEで取得→自動リマインド"),
]
tbl(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(3.7),[Inches(2.2),Inches(5.0),Inches(5.03)],sz=10.5,hsz=11,hi=2)
bx(s,Inches(0.55),Inches(5.75),Inches(12.23),Inches(0.7),REDBG); bx(s,Inches(0.55),Inches(5.75),Inches(0.1),Inches(0.7),RED)
t(s,"効きどころ：CPAP定期の「来ない→架電→月末残業」を、次回予約LINE＋自動リマインドで断つ。",Inches(0.82),Inches(5.75),Inches(11.8),Inches(0.7),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# 6 動線③ 再診/離脱
s=sl(); ft(s,"6")
hdr(s,"動線③ 再診・離脱しかけの患者","人の動き：現状 → LINE後","事務の架電を、仕組みで肩代わり")
rows=[
 ("段階","現状（人の動き）","LINE後（人の動き）"),
 ("受診後","事:次回を口頭で（取りこぼし）","患:次回予約をLINEで取得・記録"),
 ("通院の谷間","（接点なし）","患:リマインドが届く（お知らせ止まり）"),
 ("離脱しかけ","事:最近来ない人へ一人ずつ架電","事:LINEで次回予約導線→自動で促す"),
 ("不在・出張","事:つながらず何度も架電","患:都合よい時にLINEで予約"),
 ("先生フォロー","医:多くの人に電話","医:本当に必要な人だけに集中"),
]
tbl(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(3.7),[Inches(2.2),Inches(5.0),Inches(5.03)],sz=10.5,hsz=11,hi=2)
bx(s,Inches(0.55),Inches(5.75),Inches(12.23),Inches(0.7),GRYBG); bx(s,Inches(0.55),Inches(5.75),Inches(0.1),Inches(0.7),RED)
t(s,"変わる人の動き：事務の「一人ずつ架電」が減り、不在でも患者が自分で予約できる。",Inches(0.82),Inches(5.75),Inches(11.8),Inches(0.7),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# 7 まとめ 差分
s=sl(); ft(s,"7")
hdr(s,"まとめ","人の動きの差分 ── 誰の・どの動作が、どう変わるか","減る動作（事務）と、増える接点（患者）を一望")
rows=[
 ("誰","減る動作（負担）","増える/良くなる動き"),
 ("患者さん","電話・紙の記入・順番の問い合わせ","スマホで問診・予約・お知らせ受信"),
 ("受付・事務","転記・聞き取り・架電・月末残業","本当に要る対応に集中できる"),
 ("看護師","紙ESSの採点・二度手間","事前データで準備が早い"),
 ("先生","多数への架電・情報の聞き直し","必要な人だけ・情報が揃った状態で診療"),
 ("経営","流入が見えない・取りこぼし","流入が数値化・再診が継続（売上の土台）"),
]
tbl(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(3.9),[Inches(2.3),Inches(4.7),Inches(5.23)],sz=10.5,hsz=11,hi=2)
bx(s,Inches(0.55),Inches(5.95),Inches(12.23),Inches(0.95),TEALBG,line=TEALD,lw=1.2); bx(s,Inches(0.55),Inches(5.95),Inches(0.1),Inches(0.95),TEALD)
t(s,"要点：システムは「足すだけ・置き換えない」。人の動きは「事務の手間が減り、患者の接点が増える」。",Inches(0.82),Inches(5.95),Inches(11.8),Inches(0.95),sz=13,bold=True,col=TEALD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("kyobashi_taikei.pptx")
print("saved kyobashi_taikei.pptx slides:",len(prs.slides._sldIdLst))

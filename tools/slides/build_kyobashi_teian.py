"""
京橋クリニック向け 公式LINE 導入のご提案（先生・医療事務に見せる用 v2）。
比較表追加／さりげなさ／感情面配慮。クリーム白×レンガ赤。出力: kyobashi_teian.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"京橋クリニック 公式LINE 導入のご提案",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def arrow(slide,x,y,w,h):
    s=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,h); s.fill.solid()
    s.fill.fore_color.rgb=RGBColor(0xC9,0x6A,0x62); s.line.fill.background(); s.shadow.inherit=False; return s
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.09); cell.margin_right=Inches(0.06); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=(TEALBG if is_hi else (CARD if ri%2==1 else BG))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or (ci==0) or is_hi
                    r.font.color.rgb=(WHT if ri==0 else (TEALD if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A))))
    return tb

# ════ 1 表紙 ════
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"京橋クリニック 御中",Inches(0.9),Inches(1.55),Inches(11),Inches(0.4),sz=15,bold=True,col=RED)
t(s,"公式LINE 導入のご提案",Inches(0.88),Inches(2.15),Inches(11.6),Inches(0.9),sz=38,bold=True,col=INK)
t(s,"患者さん・スタッフの皆さま・先生、",Inches(0.9),Inches(3.15),Inches(11.6),Inches(0.6),sz=21,bold=True,col=RED)
t(s,"みんなが少し楽になります。",Inches(0.9),Inches(3.72),Inches(11.6),Inches(0.6),sz=21,bold=True,col=RED)
t(s,"スタッフの皆さまのアンケート（事務・看護師の声）をもとに、御院に合わせて準備しました。\n面倒な一次対応を公式LINEが引き受け、皆さんが患者さんに向き合える時間を増やすご提案です。",
  Inches(0.92),Inches(4.6),Inches(11.4),Inches(0.9),sz=14,col=GRY,line_sp=1.3)
bx(s,Inches(0.9),Inches(6.4),Inches(11.5),Pt(1.2),LN)
t(s,"ご提案：KHD（AI医療コンサル）　｜　まずは小さく、いつでも見直せます",Inches(0.9),Inches(6.55),Inches(11),Inches(0.4),sz=12.5,bold=True,col=INK)

# ════ 2 いま起きていること ════
s=sl(); ft(s)
hdr(s,"いま起きていること","アンケートから見えた、現場の4つの負担","スタッフの皆さまの声をそのまま起点にしています（＝御院の実態に合わせています）")
issues=[("鳴り止まない電話","「番号取ったのに」等の問い合わせ対応で、つい手が止まる"),
        ("受付の聞き取り・転記","紙に書いて電子カルテに写す——同じことを二度",),
        ("待ち時間の問い合わせ","順番・待ち時間を何度も聞かれ、説明しても伝わりにくい"),
        ("月末のいそがしさ","定期受診の取りこぼし→月末にまとめて会計→残業に")]
cw,ch,gx,gy=Inches(6.0),Inches(2.0),Inches(0.45),Inches(0.35); x0,y0=Inches(0.55),Inches(2.0)
for i,(ti,ds) in enumerate(issues):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.12),ch,RED)
    t(s,ti,cx+Inches(0.3),cy+Inches(0.22),cw-Inches(0.5),Inches(0.5),sz=16,bold=True,col=INK)
    t(s,ds,cx+Inches(0.32),cy+Inches(0.95),cw-Inches(0.6),Inches(0.9),sz=12,col=GRY,line_sp=1.15)
t(s,"→ この「面倒な一次対応」を公式LINEが引き受け、皆さんは患者さんに向き合えます。",Inches(0.55),Inches(6.5),Inches(12),Inches(0.4),sz=13,bold=True,col=REDD)

# ════ 3 比較表（説明用カンペ） ════
s=sl(); ft(s)
hdr(s,"よくある運用との比較","一般的なクリニックの運用と、LINEで改善した院の違い","ひと目で「御院がどう変わるか」が分かるように（ご説明用）")
rows=[
 ("場面","よくある運用（一般的な院）","LINEを活かした院では","京橋クリニックでは"),
 ("問い合わせ","ほぼ全部 電話でスタッフが対応","自動応答が一次対応、難しい分だけ電話","当日受付制に合わせFAQ＋順番通知"),
 ("受付・記入","紙に記入 → 電子カルテへ転記","来院前にスマホでWeb問診","Web問診で受付の手間を軽く"),
 ("順番・待ち時間","口頭で都度ご説明","順番が近づくとLINEで自動通知","「あと◯人」をLINEでお知らせ"),
 ("お知らせ","掲示・口頭（読まれにくい）","LINEで届く・読まれる","定期受診（CPAP等）にリマインド"),
 ("効果の把握","測っていないことが多い","電話・残業の減り方を数字で確認","導入の前後を一緒に測る"),
]
light_table(s,rows,Inches(0.45),Inches(1.95),Inches(12.43),Inches(4.2),[Inches(1.85),Inches(3.4),Inches(3.55),Inches(3.63)],hi_col=3,sz=10.5,header_sz=11.5)
t(s,"※ 公開されている導入事例でも、電話対応や残業の削減が報告されています（院により差はあります）。",Inches(0.45),Inches(6.35),Inches(12.4),Inches(0.4),sz=11,col=GRY)

# ════ 4 患者の流れがこう変わる ════
s=sl(); ft(s)
hdr(s,"どう変わるか","公式LINEで、患者さんの流れがこう変わる","来院の流れに沿って、LINEが受付の一次対応を担います")
flow=[("①問い合わせ","「やってる?」等に\nLINEが自動で回答"),
      ("②順番・受付","当日の順番が\n近づくとLINE通知"),
      ("③来院前","スマホでWeb問診\n（待ち時間が短縮）"),
      ("④来院・受付","受付はスムーズ\n転記の手間も軽く"),
      ("⑤次回へ","定期受診の方へ\nリマインドでご案内")]
cw,gx,x0=Inches(2.32),Inches(0.18),Inches(0.55)
for i,(ti,ds) in enumerate(flow):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(2.3),cw,Inches(2.3),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(2.3),cw,Inches(0.06),RED)
    t(s,ti,cx+Inches(0.15),Inches(2.55),cw-Inches(0.3),Inches(0.5),sz=14,bold=True,col=RED,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.15),Inches(3.25),cw-Inches(0.3),Inches(1.2),sz=11.5,col=INK,align=PP_ALIGN.CENTER,line_sp=1.2)
    if i<4: arrow(s,cx+cw+Inches(0.0),Inches(3.2),Inches(0.2),Inches(0.45))
t(s,"患者さんは電話せずLINEで完結。スタッフは「問い合わせ対応の波」から解放されます。",Inches(0.55),Inches(5.1),Inches(12),Inches(0.4),sz=13,bold=True,col=REDD)
bx(s,Inches(0.55),Inches(5.75),Inches(12.23),Inches(0.7),GRYBG); bx(s,Inches(0.55),Inches(5.75),Inches(0.1),Inches(0.7),RED)
t(s,"※ 京橋クリニックは「当日の順番受付制」。健康診断・化学物質過敏症など予約制サービスは専用フォームでお申し込みいただけます。",Inches(0.82),Inches(5.75),Inches(11.8),Inches(0.7),sz=11.5,col=INK,anchor=MSO_ANCHOR.MIDDLE)

# ════ 5 皆さまのメリット ════
s=sl(); ft(s)
hdr(s,"皆さまのメリット","スタッフの皆さま と 先生、それぞれに効きます","どちらか一方でなく、クリニック全体が少し楽になります")
bx(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(3.55),TEALBG,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(1.95),Inches(6.0),Inches(0.06),TEALD)
t(s,"医療事務・看護師の皆さま",Inches(0.8),Inches(2.13),Inches(5.5),Inches(0.4),sz=16,bold=True,col=TEALD)
for i,ln in enumerate(["・鳴り止まない電話 → 自動応答が一次対応してくれます","・聞き取り・転記 → 来院前のWeb問診でぐっと楽に","・待ち時間の問い合わせ → 順番をLINEが自動でお知らせ","・面倒な一次対応が減り、患者さんに向き合う時間が増えます"]):
    t(s,ln,Inches(0.85),Inches(2.72)+Inches(0.66)*i,Inches(5.5),Inches(0.6),sz=12,col=INK,line_sp=1.1)
bx(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(3.55),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(1.95),Inches(6.0),Inches(0.06),RED)
t(s,"先生",Inches(7.03),Inches(2.13),Inches(5.5),Inches(0.4),sz=16,bold=True,col=REDD)
for i,ln in enumerate(["・残業が減り、スタッフが長く働きやすい環境に","・患者満足（待ち時間・連絡のスムーズさ）","・書類や連絡のご負担も軽くなります","・効果は「数字」で見える形でご報告します"]):
    t(s,ln,Inches(7.08),Inches(2.72)+Inches(0.66)*i,Inches(5.5),Inches(0.6),sz=12,col=INK,line_sp=1.1)
bx(s,Inches(0.55),Inches(5.7),Inches(12.23),Inches(0.78),REDBG); bx(s,Inches(0.55),Inches(5.7),Inches(0.1),Inches(0.78),RED)
t(s,"※ LINEは受付の「代わり」ではありません。面倒な一次対応を引き受け、皆さんが患者さんに向き合う時間を増やすための仕組みです。",Inches(0.82),Inches(5.7),Inches(11.8),Inches(0.78),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 6 実際の画面（実物スクショ） ════
s=sl(); ft(s)
hdr(s,"実際の画面","御院の公式LINEは、もう実際にこう動いています","下は本物のスクリーンショットです（当日はスマホで実機もお見せします）")
SC=["sc2.jpg","sc1.jpg","sc3.jpg"]
caps=["メニュー → 順番・健診予約・ごあいさつ","診療時間・アクセス・FAQに自動で回答","持ち物・Web問診の受付"]
fw,fh,gap=Inches(2.5),Inches(4.5),Inches(0.8); x0=Inches(2.1)
for i in range(3):
    fx=x0+(fw+gap)*i
    bx(s,fx,Inches(1.95),fw,fh,CARD,line=CARDLN,lw=1.0); bx(s,fx,Inches(1.95),fw,Inches(0.06),RED)
    s.shapes.add_picture(SC[i],fx+Inches(0.27),Inches(2.05),height=Inches(4.25))
    t(s,caps[i],fx-Inches(0.1),Inches(6.5),fw+Inches(0.2),Inches(0.45),sz=10.5,bold=True,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.0)

# ════ 7 効果を一緒に測る ════
s=sl(); ft(s)
hdr(s,"効果の測り方","効果は、一緒に「数字」で確かめます","いまの状態を測ってから、導入後と比べます（無理なく測れるものだけ）")
rows=[("見る指標（例）","どう良くなるか"),
      ("受付の電話本数 / 日","自動応答が一次対応 → 件数が減る"),
      ("Web問診の事前記入率","来院前に取得 → 聞き取り・転記が減る"),
      ("待ち時間・順番の問い合わせ","順番をLINE通知 → 減る"),
      ("定期受診の来院率","リマインドで取りこぼし防止 → 上がる"),
      ("月末の残業","通院の平準化 → 月末の山が減る")]
n=len(rows); tb=s.shapes.add_table(n,2,Inches(0.55),Inches(2.0),Inches(12.23),Inches(3.7)).table
tb.first_row=False; tb.horz_banding=False; tb.columns[0].width=Inches(5.0); tb.columns[1].width=Inches(7.23)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Inches(0.12); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        cell.fill.solid(); cell.fill.fore_color.rgb=(RED if ri==0 else (CARD if ri%2==1 else BG))
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name=FONT; r.font.size=Pt(12.5); r.font.bold=(ri==0 or ci==0)
                r.font.color.rgb=(WHT if ri==0 else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
bx(s,Inches(0.55),Inches(5.95),Inches(12.23),Inches(0.7),GRYBG); bx(s,Inches(0.55),Inches(5.95),Inches(0.1),Inches(0.7),RED)
t(s,"どの指標を見るかは、医療事務の皆さまと相談して決めます。スタッフを評価するためではなく、楽になったかを確かめるためです。",Inches(0.82),Inches(5.95),Inches(11.8),Inches(0.7),sz=11.5,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 8 一緒に良くする（ダメ出し歓迎） ════
s=sl(); ft(s)
hdr(s,"一緒に良くする","皆さんの目で、もっと良くしてください","現場の皆さんが、患者さんの動きを一番ご存知です。今日は「一緒に作る」時間にしたいです")
asks=[("他院で「これいいな」は？","他のクリニックで見た良い仕組み・工夫があれば、ぜひ教えてください"),
      ("この文面・順番、伝わりますか？","患者さんに伝わりにくい所・言い回しは、遠慮なくご指摘ください"),
      ("現場の「困りごと」、まだありますか？","ここでしか分からない手間・ヒヤリは、まだ拾いきれていないかも"),
      ("うまくいかなそうな所は？","「これは無理そう」というダメ出しこそ、改善のヒントになります")]
cw,ch,gx,gy=Inches(6.0),Inches(1.85),Inches(0.45),Inches(0.3); x0,y0=Inches(0.55),Inches(2.05)
for i,(ti,ds) in enumerate(asks):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.12),ch,RED)
    t(s,ti,cx+Inches(0.3),cy+Inches(0.2),cw-Inches(0.5),Inches(0.5),sz=15,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.32),cy+Inches(0.85),cw-Inches(0.6),Inches(0.85),sz=11.5,col=GRY,line_sp=1.15)
bx(s,Inches(0.55),Inches(6.5),Inches(12.23),Inches(0.6),REDBG); bx(s,Inches(0.55),Inches(6.5),Inches(0.1),Inches(0.6),RED)
t(s,"今日は「提案を聞く場」でなく、皆さんと「一緒に作る場」。一言いただくほど、患者さんもスタッフも楽になります。",Inches(0.82),Inches(6.5),Inches(11.8),Inches(0.6),sz=12,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 9 装置連動（将来の連動イメージ） ════
s=sl(); ft(s)
hdr(s,"全体の連動","患者さんの「ひとつのLINE」に、院内の仕組みがつながる","順番・予約・問診・受付が一つの窓口に集約。バラバラを「ひとつ」に。")
HUBx,HUBy,HUBw,HUBh=Inches(5.0),Inches(3.35),Inches(3.33),Inches(1.4)
bx(s,HUBx,HUBy,HUBw,HUBh,RED,line=REDD,lw=1.5)
t(s,"京橋クリニック 公式LINE",HUBx,HUBy+Inches(0.24),HUBw,Inches(0.5),sz=15,bold=True,col=WHT,align=PP_ALIGN.CENTER)
t(s,"患者さんの「ひとつの窓口」",HUBx,HUBy+Inches(0.78),HUBw,Inches(0.4),sz=11,bold=True,col=RGBColor(0xF7,0xDD,0xDA),align=PP_ALIGN.CENTER)
bx(s,Inches(5.6),Inches(1.98),Inches(2.13),Inches(0.85),TEALBG,line=TEALD,lw=1.0)
t(s,"患者さん（スマホ）",Inches(5.6),Inches(1.98),Inches(2.13),Inches(0.85),sz=12,bold=True,col=TEALD,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
bx(s,Inches(6.42),Inches(2.83),Inches(0.5),Inches(0.55),RED,shape=MSO_SHAPE.UP_DOWN_ARROW)
bx(s,Inches(0.7),Inches(3.5),Inches(3.05),Inches(1.1),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.7),Inches(3.5),Inches(0.1),Inches(1.1),RED)
t(s,"i-CALL（順番受付）",Inches(0.92),Inches(3.64),Inches(2.8),Inches(0.4),sz=13,bold=True,col=INK)
t(s,"順番が近づくとLINEで通知",Inches(0.94),Inches(4.08),Inches(2.75),Inches(0.5),sz=10.5,col=GRY,line_sp=1.1)
bx(s,Inches(3.85),Inches(3.83),Inches(1.05),Inches(0.5),RED,shape=MSO_SHAPE.LEFT_RIGHT_ARROW)
bx(s,Inches(9.58),Inches(3.5),Inches(3.05),Inches(1.1),CARD,line=CARDLN,lw=1.0); bx(s,Inches(12.53),Inches(3.5),Inches(0.1),Inches(1.1),RED)
t(s,"電子カルテ Medicom",Inches(9.8),Inches(3.64),Inches(2.8),Inches(0.4),sz=13,bold=True,col=INK)
t(s,"問診・予約の情報を連携",Inches(9.82),Inches(4.08),Inches(2.75),Inches(0.5),sz=10.5,col=GRY,line_sp=1.1)
bx(s,Inches(8.43),Inches(3.83),Inches(1.05),Inches(0.5),RED,shape=MSO_SHAPE.LEFT_RIGHT_ARROW)
bx(s,Inches(5.3),Inches(5.32),Inches(2.73),Inches(0.92),TEALBG,line=TEALD,lw=1.0)
t(s,"Stream Deck（受付操作）",Inches(5.3),Inches(5.38),Inches(2.73),Inches(0.4),sz=12,bold=True,col=TEALD,align=PP_ALIGN.CENTER)
t(s,"定型対応をワンタッチ",Inches(5.3),Inches(5.78),Inches(2.73),Inches(0.4),sz=10,col=TEALD,align=PP_ALIGN.CENTER)
bx(s,Inches(6.42),Inches(4.75),Inches(0.5),Inches(0.55),RED,shape=MSO_SHAPE.UP_DOWN_ARROW)
bx(s,Inches(0.55),Inches(6.55),Inches(12.23),Inches(0.62),GRYBG); bx(s,Inches(0.55),Inches(6.55),Inches(0.1),Inches(0.62),RED)
t(s,"中央ビジコム Medicom は i-CALL の連携対応リストに有り＝電カル側の整合がとりやすい。※連携の可否は確認中（将来の連動イメージ）。",Inches(0.82),Inches(6.55),Inches(11.8),Inches(0.62),sz=11,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 10 はじめ方 ════
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"はじめ方",Inches(0.9),Inches(0.85),Inches(11),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"まず1つだけ、無料で試しませんか？",Inches(0.88),Inches(1.35),Inches(11.6),Inches(0.8),sz=30,bold=True,col=INK)
t(s,"全部を一度に始めません。効果が出た所から、少しずつ広げます。",Inches(0.92),Inches(2.2),Inches(11.4),Inches(0.5),sz=14,col=GRY)
steps=[("STEP 1","まず1つ","再診リマインド か FAQ自動応答 を無料で開始"),
       ("STEP 2","一緒に測る","導入前→1〜2ヶ月で、楽になったかを数字でご報告"),
       ("STEP 3","広げる","効いた所だけ広げます。合わなければ見直せます")]
cw,gx,x0=Inches(3.95),Inches(0.24),Inches(0.9)
for i,(st,ti,ds) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,Inches(3.1),cw,Inches(2.2),CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(3.1),cw,Inches(0.7),RED)
    t(s,st,cx,Inches(3.22),cw,Inches(0.45),sz=18,bold=True,col=WHT,align=PP_ALIGN.CENTER)
    t(s,ti,cx,Inches(3.95),cw,Inches(0.5),sz=16,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,ds,cx+Inches(0.25),Inches(4.5),cw-Inches(0.5),Inches(0.8),sz=12,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.2)
bx(s,Inches(0.9),Inches(5.7),Inches(11.5),Inches(0.85),REDBG); bx(s,Inches(0.9),Inches(5.7),Inches(0.1),Inches(0.85),RED)
t(s,"御院に合うかどうか、一緒に確かめながら進めます。",Inches(1.2),Inches(5.7),Inches(11),Inches(0.85),sz=15,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

prs.save("kyobashi_teian.pptx")
print("saved kyobashi_teian.pptx slides:",len(prs.slides._sldIdLst))

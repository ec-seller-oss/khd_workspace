#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
KHD「土地から新築」精査フロー スライド生成
------------------------------------------------
菊池研太が自分のワークフローとして組んだ、マイソク→ボリューム→仕入判断の精査フロー。
根拠は建築基準法・自治体条例という公知の法令の一般的考え方であり、特定有料教材の複製ではない。
デザインシステム: クリーム白 #F9F6EF × レンガ赤 #AA2E26（KHD標準）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM = RGBColor(0xF9, 0xF6, 0xEF)
BRICK = RGBColor(0xAA, 0x2E, 0x26)
INK   = RGBColor(0x33, 0x2A, 0x28)
GRAY  = RGBColor(0x6B, 0x60, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=CREAM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, l, t, w, h, text, size=18, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Hiragino Kaku Gothic ProN'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def header(slide, kicker, title):
    box(slide, 0, 0, Inches(0.28), SH, fill=BRICK)
    txt(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4),
        kicker, size=14, color=BRICK, bold=True)
    txt(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.9),
        title, size=30, color=INK, bold=True)
    box(slide, Inches(0.72), Inches(1.7), Inches(2.2), Inches(0.06), fill=BRICK)


# ---------- 1. 表紙 ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 0, Inches(2.6), SW, Inches(0.08), fill=BRICK)
txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.5),
    'KHD 03 / 土地から新築', size=18, color=CREAM, bold=True)
txt(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1.8),
    '「土地から新築」精査フロー', size=46, color=WHITE, bold=True)
txt(s, Inches(1), Inches(4.3), Inches(11.3), Inches(1.2),
    'マイソク1枚から「何階・何造・何戸・粗利が出るか」を\n自分で叩くための再現パイプライン', size=22, color=CREAM)
txt(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
    '菊池研太 / 自前スキル化メモ（建築基準法・自治体条例の一般原則に基づく一次スクリーニング）',
    size=12, color=GRAY)

# ---------- 2. 全体像 ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, 'OVERVIEW', '全体像：マイソク → ボリューム → 仕入判断')
steps = [
    ('1', '敷地を読む', '用途/建蔽/容積\n防火指定/斜線/日影'),
    ('2', '駐車場を先に', '建物より先に\n台数と配置を決める'),
    ('3', '建築可能範囲', '斜線・離隔・\n敷地内通路で外形'),
    ('4', '構造3択', '木造200㎡→木三共\n→耐火木造→RC'),
    ('5', '戸数を絞る', '単純割り算を\n信じない（×0.75）'),
]
x = Inches(0.7); w = Inches(2.32); gap = Inches(0.12); top = Inches(2.4)
for i, (n, t1, t2) in enumerate(steps):
    bx = Emu(int(x) + i * (int(w) + int(gap)))
    box(s, bx, top, w, Inches(2.5), fill=WHITE, line=BRICK, line_w=1.5)
    box(s, bx, top, w, Inches(0.7), fill=BRICK)
    txt(s, bx, Inches(2.5), w, Inches(0.6), f'STEP {n}', size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, bx, Inches(3.25), w, Inches(0.6), t1, size=18, color=BRICK, bold=True, align=PP_ALIGN.CENTER)
    txt(s, bx, Inches(3.9), w, Inches(1.0), t2, size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.8),
    '▶ この5ステップを毎回同じ順で回す。順番（特に②駐車場先行）を崩さないのが要。',
    size=16, color=INK, bold=True)
txt(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.8),
    '▶ 出力＝「3階／木三共 or 耐火木造／約◯戸」の当たり → 加工費・出口価格で粗利試算へ接続。',
    size=16, color=INK)

# ---------- 3. Step1 ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, 'STEP 1', '敷地を正確に読む（法定上限を出す）')
items = [
    ('実効容積率', '指定容積と「道路幅員×係数(住居0.4/商業0.6)×100」の小さい方。\n前面道路が狭いと容積を使い切れない。'),
    ('延床・建築面積の上限', '延床上限＝敷地×実効容積。建築面積上限＝敷地×建蔽。\nここが戸数の天井になる。'),
    ('防火指定を必ず確認', '「なし」を勝手に建蔽+10%しない。法22条/準防火/防火で\n構造・コストが変わる。毎回チェック。'),
    ('規制フラグ', '北側斜線/高度地区・日影・ワンルーム条例・敷地内通路。\n自治体ごとに全く違う＝エリア独自ルールを調べる癖。'),
]
y = Inches(2.3)
for h, b in items:
    box(s, Inches(0.7), y, Inches(0.18), Inches(1.0), fill=BRICK)
    txt(s, Inches(1.05), y, Inches(3.4), Inches(1.0), h, size=18, color=BRICK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.6), y, Inches(8.0), Inches(1.0), b, size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y) + int(Inches(1.15)))

# ---------- 4. Step2-3 ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, 'STEP 2-3', '駐車場を先に置く → 斜線・離隔で外形を決める')
box(s, Inches(0.7), Inches(2.3), Inches(5.8), Inches(4.3), fill=WHITE, line=BRICK, line_w=1.5)
txt(s, Inches(1.0), Inches(2.5), Inches(5.3), Inches(0.6), 'STEP2 駐車場（建物より先）', size=20, color=BRICK, bold=True)
txt(s, Inches(1.0), Inches(3.2), Inches(5.3), Inches(3.2),
    '・建物から考えると駐車場が収まらない\n・まず「車路なし」で何台置けるか\n・足りなければ車路を引く判断\n・配置パターン（直角／一列）を比較\n・台数が出口（賃料・売却）に直結する\n  地域では台数を死守',
    size=15, color=INK)
box(s, Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.3), fill=WHITE, line=BRICK, line_w=1.5)
txt(s, Inches(7.1), Inches(2.5), Inches(5.3), Inches(0.6), 'STEP3 建築可能範囲', size=20, color=BRICK, bold=True)
txt(s, Inches(7.1), Inches(3.2), Inches(5.3), Inches(3.2),
    '・道路斜線：幅員×勾配(住居1.25/商業1.5)\n・北側斜線/高度地区 → 北側を削る\n・日影規制 → 高さ10m超(4階〜)は要注意\n  東西に長い敷地ほど不利\n・隣地離隔 約0.5〜0.6m\n・敷地内通路／窓先空地（自治体条例）',
    size=15, color=INK)

# ---------- 5. Step4 構造 ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, 'STEP 4', '構造3択を法規から逆算する')
cards = [
    ('木造（200㎡未満）', '延床200㎡未満×3階以下なら\n避難・採光の規制が一気に緩む。\n狭小地の最有利ゾーン。', '最有利'),
    ('木三共 / 耐火木造', '木三共=コスト最安級だが延床\n伸びにくい。伸ばしたい時は\n耐火木造で共用部を屋内化。', 'コスト軸'),
    ('RC造', '4階以上・容積を使い切る時。\nコスト高。日影/斜線をRCで\n詰める。王道とは限らない。', '容積軸'),
]
x = Inches(0.7); w = Inches(3.9); gap = Inches(0.25); top = Inches(2.4)
for i, (t1, b, tag) in enumerate(cards):
    bx = Emu(int(x) + i * (int(w) + int(gap)))
    box(s, bx, top, w, Inches(3.6), fill=WHITE, line=BRICK, line_w=1.5)
    box(s, bx, top, w, Inches(0.65), fill=BRICK)
    txt(s, bx, top, w, Inches(0.65), tag, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, bx, Emu(int(top)+int(Inches(0.9))), w, Inches(0.7), t1, size=19, color=BRICK, bold=True, align=PP_ALIGN.CENTER)
    txt(s, bx, Emu(int(top)+int(Inches(1.7))), w, Inches(1.8), b, size=14, color=INK, align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.8),
    '▶ 選ぶ順番：木造200㎡未満 → 木三共 → 耐火木造 → RC。日影・容積・コストから逆算する。',
    size=16, color=INK, bold=True)

# ---------- 6. Step5 戸数 ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, 'STEP 5', '戸数を現実に絞る（単純割り算を信じない）')
box(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(1.5), fill=BRICK)
txt(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.5),
    '「容積 ÷ 住戸面積 ＝ 戸数」は“上限の幻”。\n廊下・階段・EV・バルコニー・避難通路で、現実は約75%に落ちる。',
    size=22, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.6),
    '例：延床上限288㎡ ÷ 専有25㎡ ＝ 11.5戸（幻）  →  ×0.75 ＝ 現実 約8戸',
    size=20, color=INK, bold=True)
txt(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.6),
    '・業者の「容積いっぱいで◯戸」を鵜呑みにしない。即座に「廊下と避難通路で減る」と返せる。\n'
    '・重層長屋・法的採光の工夫など、戸数を伸ばす打ち手は別途（建築士と詰める領域）。\n'
    '・戸数×想定賃料 or 分譲単価 → 売上 → 粗利率20%以上か（KHD仕入基準）で足切り。',
    size=15, color=INK)

# ---------- 7. KPIゲート ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, 'GATE', 'KHD 仕入KPIゲート（玉川式・辛口基準）')
rows = [
    ('土地値割合', '0.4 以上', '路線価×面積 ÷ 仕入価格。担保価値の最低ライン'),
    ('粗利率（キャピタル）', '20% 以上', '粗利益 ÷ 売上予定額。1人デベの主軸'),
    ('CF率（保有時）', '1.5〜2.0%', '税引前CF ÷ 購入価格'),
    ('CCR（保有時）', '15〜20%', '税引前CF ÷ 自己資金'),
    ('再建築・接道', '可 / 要確認', '再建不可は出口が限定。越境・共有も現調'),
]
y = Inches(2.35)
box(s, Inches(0.7), y, Inches(3.3), Inches(0.6), fill=INK)
box(s, Inches(4.0), y, Inches(2.4), Inches(0.6), fill=INK)
box(s, Inches(6.4), y, Inches(6.2), Inches(0.6), fill=INK)
txt(s, Inches(0.7), y, Inches(3.3), Inches(0.6), '指標', size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(4.0), y, Inches(2.4), Inches(0.6), '基準', size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(6.4), y, Inches(6.2), Inches(0.6), '意味', size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
y = Emu(int(y) + int(Inches(0.6)))
for i, (a, b, c) in enumerate(rows):
    f = WHITE if i % 2 == 0 else CREAM
    box(s, Inches(0.7), y, Inches(3.3), Inches(0.7), fill=f, line=GRAY, line_w=0.5)
    box(s, Inches(4.0), y, Inches(2.4), Inches(0.7), fill=f, line=GRAY, line_w=0.5)
    box(s, Inches(6.4), y, Inches(6.2), Inches(0.7), fill=f, line=GRAY, line_w=0.5)
    txt(s, Inches(0.8), y, Inches(3.1), Inches(0.7), a, size=14, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.0), y, Inches(2.4), Inches(0.7), b, size=14, color=BRICK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.5), y, Inches(6.0), Inches(0.7), c, size=13, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y) + int(Inches(0.7)))

# ---------- 8. 運用 ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 0, Inches(1.5), SW, Inches(0.08), fill=BRICK)
txt(s, Inches(1), Inches(0.7), Inches(11), Inches(0.6), 'これからの運用', size=18, color=BRICK, bold=True)
txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.0),
    'メンバーシップに頼らず、自分の資産で回す', size=34, color=WHITE, bold=True)
txt(s, Inches(1), Inches(3.4), Inches(11.3), Inches(3.2),
    '① 自前ツール  land_screening_pipeline.py\n'
    '    → マイソク数値を入れて5ステップの当たりを即出力（仕入の足切り）\n\n'
    '② 要点ノート（Notion・自分の言葉）\n'
    '    → ボリューム検討マスター／実例メソッド。考え方の引き出し\n\n'
    '③ 正式判断は建築士・自治体窓口\n'
    '    → ツールは“当たり”まで。詰めは専門家と現調で。ここは外注思考',
    size=17, color=CREAM)

prs.save('/Users/kikuchikenta/01_honbu_docs_automation/KHD_土地から新築_精査フロー_v1.pptx')
print('SAVED KHD_土地から新築_精査フロー_v1.pptx')

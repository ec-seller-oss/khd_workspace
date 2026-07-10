#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
KHD「土地から新築」精査フロー スライド 大容量版 v2
------------------------------------------------------
全14回の実例を網羅した社内学習資産。図版は全てKHD自作の模式図
（公知の建築基準法・条例の考え方を図解したオリジナル）。
わん1級建築士塾の本文・図面・スクショは一切複製しない。
自社物件素材を差し込めるプレースホルダ付き。
デザインシステム: クリーム白 #F9F6EF × レンガ赤 #AA2E26
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CREAM=RGBColor(0xF9,0xF6,0xEF); BRICK=RGBColor(0xAA,0x2E,0x26)
INK=RGBColor(0x33,0x2A,0x28); GRAY=RGBColor(0x6B,0x60,0x5C)
WHITE=RGBColor(0xFF,0xFF,0xFF); LBRICK=RGBColor(0xE7,0xCF,0xCC)
SUN=RGBColor(0xE0,0x9A,0x2A); SKY=RGBColor(0xCF,0xDD,0xE3)
FONT='Hiragino Kaku Gothic ProN'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def bg(s,c=CREAM):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape,l,t,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def line_shape(s,l,t,w,h,color=INK,lw=2.0):
    sp=s.shapes.add_connector(2,l,t,Emu(int(l)+int(w)),Emu(int(t)+int(h)))
    sp.line.color.rgb=color; sp.line.width=Pt(lw); sp.shadow.inherit=False; return sp
def txt(s,l,t,w,h,text,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=FONT
    return tb
def header(s,kicker,title):
    box(s,0,0,Inches(0.28),SH,fill=BRICK)
    txt(s,Inches(0.7),Inches(0.4),Inches(11.8),Inches(0.4),kicker,size=13,color=BRICK,bold=True)
    txt(s,Inches(0.7),Inches(0.72),Inches(12.2),Inches(0.85),title,size=27,color=INK,bold=True)
    box(s,Inches(0.72),Inches(1.55),Inches(2.0),Inches(0.05),fill=BRICK)
def placeholder(s,l,t,w,h,label):
    box(s,l,t,w,h,fill=RGBColor(0xEE,0xEA,0xE2),line=GRAY,lw=1.0)
    txt(s,l,t,w,h,label,size=12,color=GRAY,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ============ 1. 表紙 ============
s=prs.slides.add_slide(BLANK); bg(s,INK)
box(s,0,Inches(2.5),SW,Inches(0.08),fill=BRICK)
txt(s,Inches(1),Inches(1.4),Inches(11.3),Inches(0.5),'KHD 03 / 土地から新築　社内学習資産',size=17,color=CREAM,bold=True)
txt(s,Inches(1),Inches(2.7),Inches(11.3),Inches(1.6),'「土地から新築」精査フロー　大容量版',size=42,color=WHITE,bold=True)
txt(s,Inches(1),Inches(4.2),Inches(11.3),Inches(1.4),'マイソク1枚から「何階・何造・何戸・粗利」を自分で叩く再現パイプライン\n＋ 実例14回から抽出した“数字の当て方”カタログ',size=21,color=CREAM)
txt(s,Inches(1),Inches(6.55),Inches(11.3),Inches(0.5),'菊池研太 / 図版は全てKHD自作の模式図（建築基準法・自治体条例の一般原則を図解）。第三者有料教材の複製は含まない。',size=11,color=GRAY)

# ============ 2. 3資産の関係 ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'HOW TO USE','3つの資産を使い分ける')
cards=[('① ルール','法規のタテ軸＋方法のヨコ軸','ボリューム検討マスター/実例メソッド要点ノート（Notion）'),
       ('② カタログ','実例14回の“当て方”','実例Vol.1〜14 要点集（Notion）。似た条件の回を引く'),
       ('③ ツール','マイソク→当たり自動算出','land_screening_pipeline.py（5ステップを即出力）')]
x=Inches(0.7);w=Inches(3.95);gap=Inches(0.2);top=Inches(2.1)
for i,(a,b,c) in enumerate(cards):
    bx=Emu(int(x)+i*(int(w)+int(gap)))
    box(s,bx,top,w,Inches(3.4),fill=WHITE,line=BRICK,lw=1.5)
    box(s,bx,top,w,Inches(0.75),fill=BRICK)
    txt(s,bx,top,w,Inches(0.75),a,size=19,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,bx,Emu(int(top)+int(Inches(1.0))),w,Inches(0.8),b,size=17,color=BRICK,bold=True,align=PP_ALIGN.CENTER)
    txt(s,bx,Emu(int(top)+int(Inches(1.9))),w,Inches(1.3),c,size=13,color=GRAY,align=PP_ALIGN.CENTER)
txt(s,Inches(0.7),Inches(5.8),Inches(12),Inches(1.2),'▶ ①②で「何が建つか」を当て、ツールで一次足切り → ③KPIゲートで「儲かるか」を判定 → 専門家へ。\n   収益判断（土地値・粗利・CF）は投資家＝KHD側の領域。建築士は建築精度の領域。',size=15,color=INK,bold=True)

# ============ 3. 全体像 5ステップ ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'OVERVIEW','全体像：マイソク → ボリューム → 仕入判断')
steps=[('1','敷地を読む','用途/建蔽/容積\n防火指定/斜線/日影'),('2','駐車場を先に','建物より先に\n台数と配置'),
       ('3','建築可能範囲','斜線・離隔・\n通路で外形'),('4','構造3択','木造200㎡→木三共\n→耐火木造→RC'),
       ('5','戸数を絞る','単純割り算を\n信じない(×0.75)')]
x=Inches(0.7);w=Inches(2.32);gap=Inches(0.12);top=Inches(2.2)
for i,(n,t1,t2) in enumerate(steps):
    bx=Emu(int(x)+i*(int(w)+int(gap)))
    box(s,bx,top,w,Inches(2.5),fill=WHITE,line=BRICK,lw=1.5); box(s,bx,top,w,Inches(0.65),fill=BRICK)
    txt(s,bx,top,w,Inches(0.65),f'STEP {n}',size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,bx,Emu(int(top)+int(Inches(0.85))),w,Inches(0.6),t1,size=17,color=BRICK,bold=True,align=PP_ALIGN.CENTER)
    txt(s,bx,Emu(int(top)+int(Inches(1.5))),w,Inches(0.9),t2,size=12,color=GRAY,align=PP_ALIGN.CENTER)
    if i<4: txt(s,Emu(int(bx)+int(w)-int(Inches(0.05))),Emu(int(top)+int(Inches(0.9))),Inches(0.25),Inches(0.5),'▶',size=16,color=BRICK,bold=True)
txt(s,Inches(0.7),Inches(5.1),Inches(12),Inches(1.6),'▶ 毎回同じ順で回す。特に②駐車場先行と⑤現実戸数を崩さない。\n▶ 出力＝「3階／木三共 or 耐火木造／約◯戸」→ 加工費・出口価格で粗利試算へ接続。',size=16,color=INK)

# ============ 4-5. 横断テーマ ============
themes1=['法定上限をそのまま信じない（狭小地は斜線・通路・採光で必ず削られる）',
 '200㎡未満は“魔法のライン”：木三共A/B制約が外れ、バルコニー不要＋屋内階段＋45分準耐火',
 '容積率は「指定 vs 道路幅員×係数(住居0.4/非住居0.6)」の小さい方',
 '旗竿地の容積は竿幅でなく“接道する建基法上の道路幅”で算定']
themes2=['日影で高さ確保しやすい順：①北側接道 ②東西に短い ③東西に長い＝多くNG',
 '駐車場は建物より先・道路に直角配置で「車路(シャロ)」を作らない',
 '竿4m以上の旗竿地は防火避難規定が外れ木三共が建つ＝安く拾える狙い目',
 '戸数増の定石：長屋×共同住宅併用／40㎡超はカウント外／3階有効範囲の精算']
for idx,(ttl,items) in enumerate([('①',themes1),('②',themes2)]):
    s=prs.slides.add_slide(BLANK); bg(s); header(s,'KEY PATTERNS',f'全14回で繰り返し出た再現パターン {ttl}')
    y=Inches(2.1)
    for it in items:
        box(s,Inches(0.7),y,Inches(0.5),Inches(1.0),fill=BRICK)
        txt(s,Inches(0.7),y,Inches(0.5),Inches(1.0),'✓',size=22,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        box(s,Inches(1.3),y,Inches(11.3),Inches(1.0),fill=WHITE,line=LBRICK,lw=1.0)
        txt(s,Inches(1.6),y,Inches(10.8),Inches(1.0),it,size=16,color=INK,anchor=MSO_ANCHOR.MIDDLE)
        y=Emu(int(y)+int(Inches(1.18)))

# ============ 6. [図] 建築可能範囲の4要素 ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'DIAGRAM','建築可能範囲の4要素（KHD自作模式図）')
# 敷地 box
sx,sy,sw,sh=Inches(1.2),Inches(2.4),Inches(5.6),Inches(4.0)
box(s,sx,sy,sw,sh,fill=RGBColor(0xEF,0xE9,0xDD),line=GRAY,lw=1.0)
# 道路 (下)
box(s,sx,Emu(int(sy)+int(sh)),sw,Inches(0.5),fill=RGBColor(0xCF,0xCF,0xCF))
txt(s,sx,Emu(int(sy)+int(sh)),sw,Inches(0.5),'前面道路',size=11,color=INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# 道路斜線 (斜め線 下から上へ)
line_shape(s,sx,Emu(int(sy)+int(sh)),Inches(2.2),Emu(-int(sh)),color=BRICK,lw=2.5)
txt(s,Emu(int(sx)+int(Inches(0.2))),Emu(int(sy)+int(Inches(2.4))),Inches(2.2),Inches(0.4),'道路斜線(幅員×1.25/1.5)',size=10,color=BRICK,bold=True)
# 北側斜線 (上辺から内側)
line_shape(s,Emu(int(sx)+int(sw)-int(Inches(1.6))),sy,Inches(1.6),Inches(1.3),color=RGBColor(0x2A,0x6E,0xA0),lw=2.5)
txt(s,Emu(int(sx)+int(sw)-int(Inches(2.2))),Emu(int(sy)-int(Inches(0.05))),Inches(2.2),Inches(0.4),'北側斜線/高度地区',size=10,color=RGBColor(0x2A,0x6E,0xA0),bold=True,align=PP_ALIGN.RIGHT)
# 0.6m離隔 (内側点線枠) — 内側の小box
box(s,Emu(int(sx)+int(Inches(0.4))),Emu(int(sy)+int(Inches(0.4))),Emu(int(sw)-int(Inches(0.8))),Emu(int(sh)-int(Inches(0.8))),line=INK,lw=1.0)
txt(s,Emu(int(sx)+int(Inches(0.55))),Emu(int(sy)+int(Inches(0.5))),Inches(3),Inches(0.4),'隣地離隔 0.6m',size=10,color=INK)
# 中央 建築可能範囲ラベル
txt(s,sx,Emu(int(sy)+int(Inches(1.7))),sw,Inches(0.6),'建築可能範囲',size=18,color=BRICK,bold=True,align=PP_ALIGN.CENTER)
# 右側 解説
rx=Inches(7.2)
txt(s,rx,Inches(2.3),Inches(5.5),Inches(0.5),'4要素を引き算して外形を出す',size=18,color=INK,bold=True)
for i,(t1) in enumerate(['① 道路斜線：幅員×1.25(住居)/1.5(非住居)','② 北側斜線/高度地区：北側を削る','③ 隣地離隔：約0.5〜0.6m','④ 敷地内通路2.1m/窓先空地（条例）']):
    txt(s,rx,Emu(int(Inches(2.9))+i*int(Inches(0.7))),Inches(5.5),Inches(0.6),t1,size=15,color=INK)
txt(s,rx,Inches(5.9),Inches(5.5),Inches(1.0),'※天空率を使うと道路斜線が緩和され、セットバックは概算0.7〜0.8掛けに短縮できる',size=13,color=GRAY)

# ============ 7. [図] 200㎡未満の魔法 ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'DIAGRAM','「延床200㎡未満」の効能（KHD自作模式図）')
# 左 200㎡以上
lx=Inches(0.9)
box(s,lx,Inches(2.2),Inches(5.4),Inches(0.7),fill=GRAY)
txt(s,lx,Inches(2.2),Inches(5.4),Inches(0.7),'延床 200㎡ 以上',size=18,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for i,t1 in enumerate(['× バルコニー2㎡が必須（専有が痩せる）','× 屋内階段にできない（セキュリティ弱）','× 60分準耐火（外壁コスト高）','× 木三共は条件A/Bの充足が必要']):
    txt(s,lx,Emu(int(Inches(3.1))+i*int(Inches(0.62))),Inches(5.4),Inches(0.6),t1,size=15,color=INK)
# 右 200㎡未満
rx=Inches(7.0)
box(s,rx,Inches(2.2),Inches(5.4),Inches(0.7),fill=BRICK)
txt(s,rx,Inches(2.2),Inches(5.4),Inches(0.7),'延床 200㎡ 未満',size=18,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for i,t1 in enumerate(['◎ バルコニー不要 → 専有面積に回せる','◎ 屋内階段OK（セキュリティUP）','◎ 45分準耐火（外壁仕様の選択肢増）','◎ 木三共A/B制約が外れ配置自由']):
    txt(s,rx,Emu(int(Inches(3.1))+i*int(Inches(0.62))),Inches(5.4),Inches(0.6),t1,size=15,color=BRICK,bold=True)
txt(s,Inches(0.9),Inches(6.0),Inches(11.5),Inches(1.0),'▶ 狭小地は「なるべく多く」でなく「200㎡の壁を超えないギリギリ」で容積参入面積を最大化する（板橋Vol.14の最適解）。',size=15,color=INK,bold=True)

# ============ 8. [図] 木三共 条件A/B ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'DIAGRAM','木三共：条件A/B マトリクス（200㎡超で必要）')
txt(s,Inches(0.7),Inches(1.8),Inches(12),Inches(0.5),'条件Aから1つ＋条件Bから1つを満たすと、200㎡超でも木造3階を「準耐火」で建てられる',size=15,color=INK)
# 2x2 grid
gx,gy,cw,ch=Inches(0.9),Inches(2.5),Inches(5.7),Inches(0.95)
data=[('条件A①','避難上有効なバルコニー2㎡以上'),('条件A②','共用廊下・階段を開放（屋外階段）'),
      ('条件B③','居室開口の外壁面から3m離隔（道路面除く）'),('条件B④','共用廊下・階段を開放（屋外階段）')]
for i,(a,b) in enumerate(data):
    ty=Emu(int(gy)+i*int(int(ch)+int(Inches(0.12))))
    box(s,gx,ty,Inches(1.4),ch,fill=BRICK)
    txt(s,gx,ty,Inches(1.4),ch,a,size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    box(s,Emu(int(gx)+int(Inches(1.5))),ty,Inches(4.6),ch,fill=WHITE,line=LBRICK,lw=1.0)
    txt(s,Emu(int(gx)+int(Inches(1.65))),ty,Inches(4.4),ch,b,size=13,color=INK,anchor=MSO_ANCHOR.MIDDLE)
# 右 結論
rx=Inches(7.1)
box(s,rx,Inches(2.5),Inches(5.5),Inches(4.0),fill=RGBColor(0xF1,0xE6,0xE4),line=BRICK,lw=1.5)
txt(s,Emu(int(rx)+int(Inches(0.3))),Inches(2.7),Inches(5.0),Inches(0.6),'屋内階段にできるのは',size=16,color=INK,bold=True)
txt(s,Emu(int(rx)+int(Inches(0.3))),Inches(3.25),Inches(5.0),Inches(0.7),'A① ＋ B③ の組合せのみ',size=22,color=BRICK,bold=True)
txt(s,Emu(int(rx)+int(Inches(0.3))),Inches(4.1),Inches(5.0),Inches(2.2),'＝ バルコニーあり＋居室窓は道路側のみ。\nそれ以外はA②/B④の開放階段（踊り場2㎡以上を屋外開放→外周配置）で住戸自由度が落ちる。\n\n→ 1フロア3戸かつ全室道路側は厳しい。\n　だから狭小地は200㎡未満で組む方が強い。',size=13,color=INK)

# ============ 9. [図] 日影 高さ確保 優先順位 ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'DIAGRAM','日影規制：高さを確保しやすい敷地の順（KHD自作模式図）')
patterns=[('① 北側接道','◎ 最有利','道路緩和で\n10m超でも可'),('② 東西に短い','○ 有利','日影が\n回りにくい'),('③ 東西に長い','× 多くNG','日影が\n伸びやすい')]
x=Inches(1.0);w=Inches(3.6);gap=Inches(0.45);top=Inches(2.3)
for i,(t1,verdict,note) in enumerate(patterns):
    bx=Emu(int(x)+i*(int(w)+int(gap)))
    box(s,bx,top,w,Inches(3.2),fill=WHITE,line=GRAY,lw=1.0)
    # 太陽
    box(s,Emu(int(bx)+int(Inches(1.45))),Emu(int(top)+int(Inches(0.25))),Inches(0.7),Inches(0.7),fill=SUN,shape=MSO_SHAPE.OVAL)
    # 敷地形状
    if i==2:
        box(s,Emu(int(bx)+int(Inches(0.5))),Emu(int(top)+int(Inches(1.3))),Inches(2.6),Inches(0.8),fill=SKY,line=INK,lw=1.0)
    else:
        box(s,Emu(int(bx)+int(Inches(1.05))),Emu(int(top)+int(Inches(1.2))),Inches(1.5),Inches(1.1),fill=SKY,line=INK,lw=1.0)
    if i==0:
        box(s,Emu(int(bx)+int(Inches(0.5))),Emu(int(top)+int(Inches(2.45))),Inches(2.6),Inches(0.25),fill=RGBColor(0xCF,0xCF,0xCF))
        txt(s,Emu(int(bx)+int(Inches(0.5))),Emu(int(top)+int(Inches(2.42))),Inches(2.6),Inches(0.3),'道路(北)',size=9,color=INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,bx,top,w,Inches(0.45),t1,size=16,color=BRICK,bold=True,align=PP_ALIGN.CENTER)
    txt(s,bx,Emu(int(top)+int(Inches(2.75))),w,Inches(0.45),verdict,size=17,color=INK,bold=True,align=PP_ALIGN.CENTER)
txt(s,Inches(1.0),Inches(5.9),Inches(11.5),Inches(1.0),'▶ 10m基準は実質“緩和”（9.9mでも日影NGが出る敷地は多い）。北側接道でないなら半地下/盛土で高さ10m未満に抑える。',size=15,color=INK,bold=True)

# ============ 10. [図] 駐車場 直角 vs 車路 ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'DIAGRAM','駐車場：直角配置で「車路」を作らない（KHD自作模式図）')
# NG 左
lx=Inches(0.9)
txt(s,lx,Inches(1.9),Inches(5.4),Inches(0.5),'× NG：縦列・斜め → 車路(シャロ)が建物を圧迫',size=15,color=GRAY,bold=True)
box(s,lx,Inches(2.5),Inches(5.4),Inches(3.4),fill=WHITE,line=GRAY,lw=1.0)
box(s,Emu(int(lx)+int(Inches(0.3))),Inches(2.7),Inches(2.0),Inches(3.0),fill=RGBColor(0xDD,0xCF,0xA0),line=GRAY,lw=0.5) # 車路
txt(s,Emu(int(lx)+int(Inches(0.3))),Inches(4.0),Inches(2.0),Inches(0.5),'車路',size=12,color=INK,align=PP_ALIGN.CENTER)
for i in range(3):
    box(s,Emu(int(lx)+int(Inches(2.5))),Emu(int(Inches(2.8))+i*int(Inches(0.95))),Inches(2.5),Inches(0.8),fill=SKY,line=INK,lw=0.8)
    txt(s,Emu(int(lx)+int(Inches(2.5))),Emu(int(Inches(2.8))+i*int(Inches(0.95))),Inches(2.5),Inches(0.8),'車',size=11,color=INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# OK 右
rx=Inches(7.0)
txt(s,rx,Inches(1.9),Inches(5.4),Inches(0.5),'◎ OK：道路に直角 → 道路を切り返しに使う',size=15,color=BRICK,bold=True)
box(s,rx,Inches(2.5),Inches(5.4),Inches(3.4),fill=WHITE,line=BRICK,lw=1.5)
for i in range(4):
    box(s,Emu(int(rx)+int(Inches(0.3))+i*int(Inches(1.25))),Inches(2.8),Inches(1.05),Inches(2.0),fill=SKY,line=INK,lw=0.8)
    txt(s,Emu(int(rx)+int(Inches(0.3))+i*int(Inches(1.25))),Inches(3.6),Inches(1.05),Inches(0.5),'車',size=11,color=INK,align=PP_ALIGN.CENTER)
box(s,rx,Inches(5.0),Inches(5.4),Inches(0.45),fill=RGBColor(0xCF,0xCF,0xCF))
txt(s,rx,Inches(5.0),Inches(5.4),Inches(0.45),'幅6m以上の道路（切り返しに使う）',size=11,color=INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.9),Inches(6.05),Inches(11.5),Inches(1.0),'3原則：①道路に直角 ②幅6m以上の道路に面して ③有効奥行5m×幅2.3m死守（寸法は自治体）。ピロティ車庫は構造体で実2.8〜3.0m要。',size=14,color=INK,bold=True)

# ============ 11. [図] 旗竿地 竿幅員ルール ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'DIAGRAM','旗竿地：竿幅員で建つものが決まる（KHD自作模式図）')
# 旗竿の絵
box(s,Inches(1.0),Inches(2.4),Inches(3.6),Inches(2.6),fill=RGBColor(0xEF,0xE9,0xDD),line=GRAY,lw=1.0) # 旗
box(s,Inches(2.4),Inches(5.0),Inches(0.8),Inches(1.3),fill=RGBColor(0xDD,0xCF,0xA0),line=GRAY,lw=1.0) # 竿
txt(s,Inches(2.0),Inches(5.4),Inches(1.6),Inches(0.5),'竿(路地状)',size=11,color=INK,align=PP_ALIGN.CENTER)
box(s,Inches(1.0),Emu(int(Inches(6.3))),Inches(3.6),Inches(0.4),fill=RGBColor(0xCF,0xCF,0xCF))
txt(s,Inches(1.0),Inches(6.3),Inches(3.6),Inches(0.4),'道路',size=10,color=INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# 右 ルール表
rx=Inches(5.2)
rows=[('竿幅員','建てられるもの（目安）'),('2m未満','建築不可'),('2〜3m未満','竿長15m以下（横浜例）'),
      ('3m以上','竿長25mまで（横浜例）'),('4m以上','防火避難規定が外れ木三共OK＝狙い目')]
ry=Inches(2.3)
for i,(a,b) in enumerate(rows):
    f=INK if i==0 else (RGBColor(0xF1,0xE6,0xE4) if i==4 else (WHITE if i%2 else CREAM))
    tc=WHITE if i==0 else INK
    box(s,rx,ry,Inches(2.2),Inches(0.7),fill=f,line=GRAY,lw=0.5)
    box(s,Emu(int(rx)+int(Inches(2.2))),ry,Inches(5.0),Inches(0.7),fill=f,line=GRAY,lw=0.5)
    txt(s,Emu(int(rx)+int(Inches(0.1))),ry,Inches(2.0),Inches(0.7),a,size=14,color=tc,bold=(i==0 or i==4),anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Emu(int(rx)+int(Inches(2.3))),ry,Inches(4.8),Inches(0.7),b,size=13,color=(BRICK if i==4 else tc),bold=(i==0 or i==4),anchor=MSO_ANCHOR.MIDDLE)
    ry=Emu(int(ry)+int(Inches(0.72)))
txt(s,Inches(1.0),Inches(6.85),Inches(11.5),Inches(0.5),'※容積率は竿幅でなく“接道する建基法上の道路幅”で算定。旗竿か否かは地域条例で変わる（要確認）。',size=12,color=GRAY)

# ============ 12〜25. 各Vol カード（14枚） ============
vols=[
 ('Vol.1','東京都足立区','住居系・日影測定面4m・東西に長い','耐火木造3階 12戸（延床258.99㎡）','単純割り算は裏切られる→日影で4階不可。RCはコスパ悪・木三共は延床伸びず→耐火木造3階が最適解'),
 ('Vol.2','東京都新宿区','第2種高度地区・ワンルーム条例10戸以上','木造3階 200㎡未満 9戸','200㎡未満の3効能＝バルコニー不要/屋内階段/45分準耐火。200㎡超は専有がむしろ痩せる'),
 ('Vol.3','東京都大田区','奥行深い東西長・日影測定面6.5m・指定容積400%優先','RC壁式 半地下4階 14戸','半地下で建物高さ9.9m<10mに抑え日影規制を回避し4階を成立'),
 ('Vol.4','東京都新宿区','接道4m未満・前面道路容積273%','木造耐火3階 8戸（容積算定200㎡未満）','接道長さの罠＝200㎡以下共同住宅は4m未満でも可だが3階以下が必須'),
 ('Vol.5','神奈川県横浜市(旗竿)','竿4m未満・共住100㎡未満/長屋3戸/2階以下','木造2階 8戸（共同5+長屋3）','共同住宅と長屋の併用で戸数を捻出。旗竿は参入障壁＝狙い目'),
 ('Vol.6','神奈川県横浜市','二方向避難(市条例6条)200㎡未満0.9m/以上2.0m','木造3階 9戸（延床175.68㎡）','わずか4.5㎡差で必要敷地が約1.5倍。200㎡未満で必要敷地が激減'),
 ('Vol.7','東京都中野区','北側接道でない・ワンルーム3階かつ12戸','木造3階 200㎡未満 8戸（→9戸も）','全斜線・通路・離隔後の「3階有効範囲」を精算し200㎡÷3階が収まれば1戸増'),
 ('Vol.8','東京都世田谷区','西側接道・第3種高度・複合制約','壁式RC4階 16戸（長屋2+共同14）','日影は北側接道>東西短い>東西長い。天空率でSB0.7〜0.8掛け。40㎡超はカウント外'),
 ('Vol.9','東京都豊島区(旗竿)','竿接道・容積は道路幅で算定160%','木造3階 重層長屋 9戸','容積は竿幅でなく道路幅で算定。重層長屋は3300mm開口幅が何個入るか。区で扱い差'),
 ('Vol.10','東京都足立区(旗竿)','竿4m以上・旗竿の正確な理解','木造3階 200㎡未満 6戸','竿4m以上は防火避難規定が外れ木三共OK＝他人がスルーする安く拾える狙い目'),
 ('Vol.11','東京都世田谷区(旗竿)','路地状4m未満は3階共同住宅不可','重層長屋 8戸','法的採光の工夫で建物幅3300→2700mmに縮め6戸→8戸'),
 ('Vol.12','北海道札幌市','駐車場義務・建蔽100%・2面道路12m','壁式RC4階 19戸','駐車場は道路に直角で車路なし。1階1区画を駐車場+通路に充て上階すっきり。北海道は盛土'),
 ('Vol.13','岐阜県岐阜市','大型645㎡・法22条建蔽60%','木三共 3層 15戸（延床502.8㎡）','大型でも手順同じ。駐車場先行で8→15台。道路斜線vs木三共8.2mでSB0.28m'),
 ('Vol.14','東京都板橋区(L字)','準工業・複合制約・建蔽70%・指定200%','木三200㎡未満 9戸','通路位置が階段→構造要件→延床に連鎖。3案比較は「容積参入面積＋コスト」で判断'),
]
for code,area,site,answer,core in vols:
    s=prs.slides.add_slide(BLANK); bg(s)
    header(s,f'実例カタログ　{code}',area)
    # 答えバッジ
    box(s,Inches(0.7),Inches(1.85),Inches(12.0),Inches(0.95),fill=BRICK)
    txt(s,Inches(1.0),Inches(1.85),Inches(11.4),Inches(0.95),f'答え：{answer}',size=22,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    # 敷地条件
    txt(s,Inches(0.7),Inches(3.1),Inches(2.2),Inches(0.5),'敷地条件',size=15,color=BRICK,bold=True)
    box(s,Inches(2.9),Inches(3.05),Inches(9.8),Inches(0.9),fill=WHITE,line=LBRICK,lw=1.0)
    txt(s,Inches(3.1),Inches(3.05),Inches(9.5),Inches(0.9),site,size=15,color=INK,anchor=MSO_ANCHOR.MIDDLE)
    # この回の核
    txt(s,Inches(0.7),Inches(4.2),Inches(2.2),Inches(0.5),'この回の核',size=15,color=BRICK,bold=True)
    box(s,Inches(2.9),Inches(4.15),Inches(9.8),Inches(1.5),fill=RGBColor(0xF1,0xE6,0xE4),line=BRICK,lw=1.2)
    txt(s,Inches(3.1),Inches(4.15),Inches(9.5),Inches(1.5),core,size=16,color=INK,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    # 自社素材プレースホルダ
    placeholder(s,Inches(0.7),Inches(5.95),Inches(12.0),Inches(0.95),'▢ ここに自社物件（栄町・船橋など）のマイソク/写真/自作ボリューム図を差し込む（任意）')

# ============ KPIゲート ============
s=prs.slides.add_slide(BLANK); bg(s); header(s,'GATE','KHD 仕入KPIゲート（玉川式・辛口基準）')
txt(s,Inches(0.7),Inches(1.75),Inches(12),Inches(0.5),'※わん流＝建築ボリューム精度／KHD＝収益判定。役割が違う。当たりが出たらこのゲートで足切り。',size=13,color=GRAY)
rows=[('指標','基準','意味'),('土地値割合','0.4以上','路線価×面積÷仕入。担保価値の最低ライン'),
 ('粗利率(キャピタル)','20%以上','粗利益÷売上予定額。1人デベの主軸'),
 ('CF率(保有時)','1.5〜2.0%','税引前CF÷購入価格'),('CCR(保有時)','15〜20%','税引前CF÷自己資金'),
 ('再建築・接道','可/要確認','再建不可は出口限定。越境・共有も現調')]
y=Inches(2.4)
for i,(a,b,c) in enumerate(rows):
    f=INK if i==0 else (WHITE if i%2 else CREAM); tc=WHITE if i==0 else INK
    box(s,Inches(0.7),y,Inches(3.3),Inches(0.66),fill=f,line=GRAY,lw=0.5)
    box(s,Inches(4.0),y,Inches(2.4),Inches(0.66),fill=f,line=GRAY,lw=0.5)
    box(s,Inches(6.4),y,Inches(6.2),Inches(0.66),fill=f,line=GRAY,lw=0.5)
    txt(s,Inches(0.85),y,Inches(3.1),Inches(0.66),a,size=14,color=tc,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(4.0),y,Inches(2.4),Inches(0.66),b,size=14,color=(tc if i==0 else BRICK),bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(6.55),y,Inches(6.0),Inches(0.66),c,size=12,color=(tc if i==0 else GRAY),anchor=MSO_ANCHOR.MIDDLE)
    y=Emu(int(y)+int(Inches(0.66)))

# ============ 運用 ============
s=prs.slides.add_slide(BLANK); bg(s,INK)
box(s,0,Inches(1.5),SW,Inches(0.08),fill=BRICK)
txt(s,Inches(1),Inches(0.7),Inches(11),Inches(0.6),'確認作業フロー',size=17,color=BRICK,bold=True)
txt(s,Inches(1),Inches(1.9),Inches(11.3),Inches(1.0),'実物件1件で回す',size=32,color=WHITE,bold=True)
txt(s,Inches(1),Inches(3.2),Inches(11.3),Inches(3.4),
 '① マイソクの数値を land_screening_pipeline.py に入れる\n    → 「3階／木三共 or 耐火木造／約◯戸」の当たりを出す\n\n'
 '② 実例Vol要点集で似た条件の回を引いて照合\n    → 旗竿なら竿幅、駐車場ありなら直角配置、狭小なら200㎡未満…\n\n'
 '③ KPIゲートで収益判定（土地値0.4／粗利20%）\n    → 通れば建築士・自治体で正式確認へ',size=17,color=CREAM)

# ============ 出典/ディスクレーマー ============
s=prs.slides.add_slide(BLANK); bg(s)
header(s,'NOTE','本デッキの位置づけ（出典・権利）')
txt(s,Inches(0.9),Inches(2.2),Inches(11.5),Inches(4.0),
 '・本デッキは菊池研太が「土地から新築」の仕入判断を自走するために、建築基準法・各自治体条例\n  という公知の法令の一般的な考え方を、自分の言葉とワークフローに落とし込んだKHDの社内資産です。\n\n'
 '・図版はすべてKHDが自作した模式図であり、特定の有料教材の本文・図面・スクリーンショットの\n  複製は一切含みません。\n\n'
 '・学びの土台として「わん1級建築士塾」(@taka5588) の有料記事を課金会員として参照しましたが、\n  本デッキはその要点を抽象化・再構成した学習成果であり、再配布・自社コンテンツ転用はしません。\n\n'
 '・「▢ プレースホルダ」には、菊池自身が権利を持つ自社物件（栄町・船橋等）のマイソク・写真・\n  自作ボリューム図のみを差し込んでください。',
 size=15,color=INK)

prs.save('/Users/kikuchikenta/01_honbu_docs_automation/KHD_土地から新築_精査フロー_v2_大容量版.pptx')
print('SAVED v2; slides=', len(prs.slides._sldIdLst))

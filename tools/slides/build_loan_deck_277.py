#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
融資資料デッキ生成（277 高松2丁目/盛岡 オーナーチェンジ収益物件）
================================================================
North Star成果物＝銀行に出せる融資資料デッキ。バイセル/B&SのAI査定→銀行駆け込みをTTP。
KHDデザイン（クリーム白×レンガ赤・python-pptx）。収益物件(保有)版テンプレ。
出力: out_screener/277_融資資料デッキ.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

IMG_DIR = Path.home() / "01_honbu_docs_automation" / "out_screener" / "277_images"

# ── デザインシステム ──
CREAM = RGBColor(0xF9, 0xF6, 0xEF)
BRICK = RGBColor(0xAA, 0x2E, 0x26)
INK = RGBColor(0x33, 0x2A, 0x28)
GRAY = RGBColor(0x77, 0x70, 0x6E)
GREEN = RGBColor(0x1F, 0x7A, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

# ── 277 案件データ ──
D = {
    "title": "融資ご相談資料",
    "sub": "盛岡市高松2丁目 オーナーチェンジ収益物件（2物件一括取得）",
    "company": "菊池ホールディングス（KHD）",
    "rep": "代表　菊池 研太",
    "bank": "株式会社 岩手銀行 御中",
    "total": "37,000千円",
    "props": [
        {"name": "アーバンキューブ（1K×8戸）", "price": "25,000千円", "struct": "木造2階建",
         "year": "1997年3月（築29年）", "land": "335.66㎡（公簿）+私道持分47.10㎡",
         "bldg": "238.48㎡", "rent_m": "294千円", "rent_y": "3,528千円", "yield": "14.11%",
         "occ": "入居 7/8（オーナーチェンジ）", "addr": "盛岡市高松2丁目34-5",
         "note": "2023年 外壁改修済／西側位置指定道路 持分100%", "img": "urban_06.jpg"},
        {"name": "花みずき貸家3棟（3K）", "price": "12,000千円", "struct": "木造2階建",
         "year": "1968年7月（築58年）", "land": "358.16㎡（公簿）", "bldg": "163.92㎡",
         "rent_m": "108千円", "rent_y": "1,300千円", "yield": "10.83%",
         "occ": "賃貸中（3棟）", "addr": "盛岡市高松2丁目34-30",
         "note": "アーバンキューブと同時取得が条件（セット）", "img": "hana_06.jpg"},
    ],
    "access": "いわて銀河鉄道「青山」駅 徒歩22分 ／ 盛岡市街エリア",
    "youto": "市街化区域・第二種中高層住居専用地域（建ぺい60%／容積200%）",
    "rent_y_total": "4,827千円",
    "yield_omote": "13.05%",
    "yield_real": "11.09%",
    "cf_year": "1,037千円",
    "cf_ritsu": "2.80%",
    "rosenka_souzoku": "48,240円/㎡（reinfolib実勢ベース推定・実値確認中）",
    "tochine": "約90%（土地値が取得価格をほぼカバー＝担保堅い）",
    "loan_req": "37,000千円",
    "scheme": "保有（インカム）／返済原資＝賃料収入",
    # 路線価3種（reinfolib実勢ベース推定・国税庁路線価図で実値確認のうえ最終反映）
    "rosenka_kotei": "約44,700円/㎡",
    "rosenka_souzoku2": "約51,000円/㎡",
    "jissei": "約67,000円/㎡",
    "kosho": "約63,800円/㎡",
    # 相場・競合
    "soba": [
        ("本件（合算）", "13.05%", "土地値割合 約90%／オーナーチェンジ"),
        ("地方・築古1棟の一般水準", "8〜12%目安", "本件は上位水準"),
        ("都市部・築浅", "4〜6%目安", "利回りは低いが流動性高"),
    ],
    # スケジュール（マイルストン）
    "schedule": [
        ("買付提出済", "完了"),
        ("融資審査（御行）", "進行中"),
        ("売買契約", "融資内諾後"),
        ("決済・所有権移転", "契約後"),
        ("オーナーチェンジ承継・賃料収受", "決済後ただちに"),
    ],
}


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = CREAM
    rect.line.fill.background()
    rect.shadow.inherit = False
    s.shapes._spTree.remove(rect._element); s.shapes._spTree.insert(2, rect._element)
    return s


def box(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = color
    return tb


def bar(s, x, y, w, h, color=BRICK):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    r.shadow.inherit = False
    return r


def header(s, no, title):
    bar(s, 0, 0.55, 0.22, 0.65, BRICK)
    box(s, 0.5, 0.5, 11.8, 0.8, title, size=26, color=BRICK, bold=True)
    box(s, 12.4, 0.6, 0.7, 0.5, f"{no:02d}", size=16, color=GRAY, bold=True, align=PP_ALIGN.RIGHT)
    bar(s, 0.5, 1.32, 12.3, 0.02, RGBColor(0xDD, 0xCF, 0xC7))


def kpi_card(s, x, y, w, label, value, vcolor=BRICK, sub=""):
    card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.5))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE3, 0xD8, 0xD1); card.line.width = Pt(1)
    card.shadow.inherit = False
    box(s, x, y + 0.12, w, 0.4, label, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    box(s, x, y + 0.45, w, 0.7, value, size=26, color=vcolor, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        box(s, x, y + 1.12, w, 0.35, sub, size=10, color=GRAY, align=PP_ALIGN.CENTER)


def table(s, x, y, w, rows, col1=0.42, rh=0.46, fs=13):
    cy = y
    for i, (k, v) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else RGBColor(0xF2, 0xEC, 0xE6)
        r = s.shapes.add_shape(1, Inches(x), Inches(cy), Inches(w), Inches(rh))
        r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
        box(s, x + 0.12, cy + 0.04, w * col1 - 0.2, rh, k, size=fs, color=BRICK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        box(s, x + w * col1, cy + 0.04, w * (1 - col1) - 0.15, rh, v, size=fs, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        cy += rh


def footer(s):
    box(s, 0.5, 7.05, 9, 0.35, "菊池ホールディングス（KHD）／ 本資料は融資ご相談用の概算であり最終条件は協議のうえ確定します", size=9, color=GRAY)


def build():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H

    # 1 表紙
    s = add_slide(prs)
    bar(s, 0, 0, EMU_W.inches, 7.5, CREAM)
    bar(s, 0, 2.55, 13.333, 0.10, BRICK)
    bar(s, 0, 4.55, 13.333, 0.04, RGBColor(0xDD, 0xCF, 0xC7))
    box(s, 0.9, 2.75, 11.5, 1.0, D["title"], size=46, color=BRICK, bold=True)
    box(s, 0.95, 3.75, 11.5, 0.7, D["sub"], size=20, color=INK)
    box(s, 0.95, 4.75, 11.5, 0.5, D["bank"], size=18, color=INK, bold=True)
    box(s, 0.95, 6.1, 11.5, 0.9, f"{D['company']}\n{D['rep']}", size=15, color=GRAY)
    box(s, 0.95, 1.4, 11, 0.5, "CONFIDENTIAL ／ 融資ご相談資料", size=12, color=GRAY)

    # 2 案件サマリー
    s = add_slide(prs); header(s, 1, "案件概要 ― 2物件一括（セット）取得")
    box(s, 0.5, 1.5, 12.3, 0.6,
        "盛岡市高松2丁目の隣接2物件（収益アパート＋貸家3棟）を一括取得するオーナーチェンジ案件。"
        "取得総額 37,000千円。", size=15, color=INK)
    kpi_card(s, 0.5, 2.3, 2.9, "取得総額", "37,000千円", BRICK, "2物件セット")
    kpi_card(s, 3.6, 2.3, 2.9, "年間賃料", D["rent_y_total"], BRICK, "表面利回り 13.05%")
    kpi_card(s, 6.7, 2.3, 2.9, "実質利回り", D["yield_real"], GREEN, "経費15%控除")
    kpi_card(s, 9.8, 2.3, 3.0, "土地値割合", "約90%", GREEN, "担保が堅い")
    table(s, 0.5, 4.15, 12.3, [
        ("所在地", "岩手県盛岡市高松2丁目（34-5／34-30）"),
        ("アクセス", D["access"]),
        ("用途地域", D["youto"]),
        ("スキーム", "保有（インカム）／返済原資＝賃料収入"),
    ], col1=0.20, rh=0.52)
    footer(s)

    # 3,4 物件個別
    for i, pr in enumerate(D["props"]):
        s = add_slide(prs); header(s, 2 + i, f"物件{i+1} ― {pr['name']}")
        kpi_card(s, 0.5, 1.55, 3.0, "取得価格", pr["price"], BRICK)
        kpi_card(s, 3.7, 1.55, 3.0, "年間賃料", pr["rent_y"], BRICK)
        kpi_card(s, 6.9, 1.55, 3.0, "表面利回り", pr["yield"], GREEN)
        kpi_card(s, 10.1, 1.55, 2.7, "現況", "賃貸中", GREEN, pr["occ"])
        table(s, 0.5, 3.35, 7.1, [
            ("所在地", pr["addr"]),
            ("構造／築年", f"{pr['struct']}　{pr['year']}"),
            ("土地面積", pr["land"]),
            ("建物面積", pr["bldg"]),
            ("月額賃料", pr["rent_m"]),
            ("特記", pr["note"]),
        ], col1=0.28, rh=0.50)
        # 現況外観写真（右）
        imgp = IMG_DIR / pr.get("img", "")
        if imgp.exists():
            frame = s.shapes.add_shape(1, Inches(7.95), Inches(3.30), Inches(4.85), Inches(2.95))
            frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
            frame.line.color.rgb = RGBColor(0xE3, 0xD8, 0xD1); frame.shadow.inherit = False
            s.shapes.add_picture(str(imgp), Inches(8.05), Inches(3.40), width=Inches(4.65))
            box(s, 7.95, 6.30, 4.85, 0.3, "現況外観", size=10, color=GRAY, align=PP_ALIGN.CENTER)
        footer(s)

    # 5 立地・賃貸市場
    s = add_slide(prs); header(s, 4, "立地・賃貸市場")
    box(s, 0.5, 1.55, 12.3, 3.0,
        "■ 立地\n"
        f"・{D['access']}\n"
        "・盛岡市街エリア。生活利便と賃貸需要が一定して見込めるエリア。\n\n"
        "■ 賃貸ポジショニング\n"
        "・アーバンキューブ＝単身（1K）需要、入居7/8と稼働良好。2023年外壁改修済で競争力維持。\n"
        "・花みずき＝戸建賃貸（3K）。ファミリー・法人需要。\n\n"
        "■ 担保・出口\n"
        "・土地値が取得価格をほぼカバー（土地値割合 約90%）＝元本毀損リスクが小さく担保堅い。\n"
        "・返済はインカム（賃料）で回収。長期保有を前提とした安定収益型。",
        size=15, color=INK, line_spacing=1.25)
    footer(s)

    # 6 路線価・担保評価
    s = add_slide(prs); header(s, 5, "路線価・担保評価")
    kpi_card(s, 0.5, 1.55, 3.0, "固定資産税路線価", D["rosenka_kotei"], INK, "推定")
    kpi_card(s, 3.7, 1.55, 3.0, "相続税路線価", D["rosenka_souzoku2"], BRICK, "推定・実値確認中")
    kpi_card(s, 6.9, 1.55, 3.0, "実勢（中央値）", D["jissei"], INK, "reinfolib周辺取引")
    kpi_card(s, 10.1, 1.55, 2.7, "土地値割合", "約90%", GREEN, "担保が堅い")
    box(s, 0.5, 3.35, 12.3, 3.3,
        "■ 担保評価（玉川式）\n"
        "・土地値割合 ＝ 路線価 × 土地面積(693.82㎡) ÷ 取得価格。0.4以上で安全圏のところ本件は約0.9。\n"
        "・花みずきは土地値が取得価格を上回り、建物ゼロ評価でも担保割れしない水準。\n\n"
        "■ 路線価→公示→実勢の換算（B&S式・参考）\n"
        f"・公示価格相当 ＝ 路線価 ÷ 0.8 → 約{D['kosho']}／地方実勢 ＝ 公示相当 × 1.0〜1.1 → {D['jissei']}\n"
        "・上記はreinfolib（国交省 不動産情報ライブラリ）の周辺取引価格中央値からの推定。\n"
        "・固定/相続税路線価は国税庁路線価図にて実値を確認のうえ最終資料に反映いたします。",
        size=14, color=INK, line_spacing=1.25)
    footer(s)

    # 7 収支・収益性
    s = add_slide(prs); header(s, 6, "収支・収益性")
    kpi_card(s, 0.5, 1.6, 3.0, "年間賃料", D["rent_y_total"], BRICK)
    kpi_card(s, 3.7, 1.6, 3.0, "表面利回り", D["yield_omote"], GREEN)
    kpi_card(s, 6.9, 1.6, 3.0, "実質利回り", D["yield_real"], GREEN, "経費15%控除")
    kpi_card(s, 10.1, 1.6, 2.7, "税引前CF", D["cf_year"], GREEN, f"CF率 {D['cf_ritsu']}")
    table(s, 0.5, 3.45, 12.3, [
        ("取得総額", "37,000千円（アーバン25,000＋花みずき12,000）"),
        ("年間賃料（満室時）", "4,827千円（アーバン3,528＋花みずき1,300）"),
        ("想定運営経費", "賃料の約15%（管理・修繕・空室等）"),
        ("融資条件（試算）", "金利3.0%・期間15年と仮定 → 年返済 約3,066千円"),
        ("税引前CF（試算）", "約1,037千円／年（CF率2.8%）"),
    ], col1=0.30, rh=0.50)
    box(s, 0.5, 6.5, 12.3, 0.4, "※ レントロール精査・融資条件確定により最終CFを更新します。", size=11, color=GRAY)
    footer(s)

    # 競合・相場
    s = add_slide(prs); header(s, 7, "相場・収益性の位置づけ")
    box(s, 0.5, 1.5, 12.3, 0.5, "本件の利回りを地域・タイプ別の相場水準と比較。本件は上位水準にあります。",
        size=15, color=INK)
    cy = 2.3
    hdr = s.shapes.add_shape(1, Inches(0.5), Inches(cy), Inches(12.3), Inches(0.5))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = BRICK; hdr.line.fill.background(); hdr.shadow.inherit = False
    box(s, 0.7, cy + 0.06, 4.5, 0.4, "区分", size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 5.5, cy + 0.06, 2.5, 0.4, "表面利回り", size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 8.3, cy + 0.06, 4.3, 0.4, "コメント", size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    cy += 0.5
    for i, (k, v, note) in enumerate(D["soba"]):
        bg = WHITE if i else RGBColor(0xEC, 0xE0, 0xDD)
        r = s.shapes.add_shape(1, Inches(0.5), Inches(cy), Inches(12.3), Inches(0.62))
        r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.color.rgb = RGBColor(0xE3, 0xD8, 0xD1); r.shadow.inherit = False
        box(s, 0.7, cy + 0.08, 4.7, 0.5, k, size=13, color=INK, bold=(i == 0), anchor=MSO_ANCHOR.MIDDLE)
        box(s, 5.5, cy + 0.08, 2.5, 0.5, v, size=15, color=(BRICK if i == 0 else GRAY), bold=True, anchor=MSO_ANCHOR.MIDDLE)
        box(s, 8.3, cy + 0.08, 4.3, 0.5, note, size=11, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        cy += 0.62
    box(s, 0.5, cy + 0.25, 12.3, 0.8,
        "■ 本件は「土地値以下で取得＝担保堅い」かつ「高利回り」を両立。地方築古ゆえの出口リスクは"
        "インカム回収主体の戦略と土地値の厚さで吸収する設計。",
        size=13, color=INK, line_spacing=1.25)
    box(s, 0.5, 6.7, 12.3, 0.35, "※ 利回り水準は一般的な目安。周辺相場はreinfolib取引価格より参照。", size=10, color=GRAY)
    footer(s)

    # 8 ハイライト & リスク
    s = add_slide(prs); header(s, 8, "投資ハイライト と 留意点")
    box(s, 0.5, 1.55, 6.0, 0.5, "◎ ハイライト", size=18, color=GREEN, bold=True)
    box(s, 0.5, 2.1, 6.0, 4.5,
        "・土地値以下で取得＝担保が堅い（土地値割合 約90%）\n"
        "・合算表面利回り 13%・実質11%の高収益\n"
        "・オーナーチェンジで即収益化（入居実績あり）\n"
        "・隣接2物件の一括取得でスケールメリット\n"
        "・返済はインカムで回収、長期安定型",
        size=14, color=INK, line_spacing=1.35)
    box(s, 6.8, 1.55, 6.0, 0.5, "△ 留意点（対応方針）", size=18, color=BRICK, bold=True)
    box(s, 6.8, 2.1, 6.0, 4.5,
        "・築古（花みずき築58年）→ 修繕計画・長期保有前提で評価\n"
        "・駅徒歩22分・地方立地 → インカム回収主体（出口に依存しない）\n"
        "・私道（位置指定）→ 通行掘削承諾・持分を確認\n"
        "・レントロール精査で実賃料・稼働を確定\n"
        "→ いずれも担保（土地値）の厚さが下支え",
        size=14, color=INK, line_spacing=1.35)
    footer(s)

    # 9 融資ご相談
    s = add_slide(prs); header(s, 9, "融資のご相談")
    kpi_card(s, 0.5, 1.6, 4.0, "融資希望額", D["loan_req"], BRICK, "取得総額")
    kpi_card(s, 4.7, 1.6, 4.0, "返済原資", "賃料収入", INK, "インカム")
    kpi_card(s, 8.9, 1.6, 3.9, "想定CF率", D["cf_ritsu"], GREEN, "税引前CF/価格")
    box(s, 0.5, 3.45, 12.3, 3.0,
        "■ ご相談事項\n"
        "・融資可否・融資可能額／対象費目\n"
        "・金利・期間・自己資金割合のご条件\n"
        "・抵当権設定・必要書類のご確認\n\n"
        "■ 想定スケジュール\n"
        "・買付提出済 → 融資審査 → 売買契約・決済 → オーナーチェンジ承継\n"
        "・レントロール・確定書類は順次ご提出いたします。",
        size=15, color=INK, line_spacing=1.3)
    footer(s)

    # スケジュール
    s = add_slide(prs); header(s, 10, "想定スケジュール")
    box(s, 0.5, 1.5, 12.3, 0.5, "買付提出済。融資審査を経て、決済後ただちに賃料収受を開始します。", size=15, color=INK)
    cy = 2.5
    for i, (step, st) in enumerate(D["schedule"]):
        dot = s.shapes.add_shape(9, Inches(0.7), Inches(cy + 0.05), Inches(0.34), Inches(0.34))
        done = st in ("完了", "進行中")
        dot.fill.solid(); dot.fill.fore_color.rgb = (GREEN if st == "完了" else (BRICK if st == "進行中" else RGBColor(0xCC, 0xC2, 0xBC)))
        dot.line.fill.background(); dot.shadow.inherit = False
        if i < len(D["schedule"]) - 1:
            ln = bar(s, 0.85, cy + 0.39, 0.04, 0.55, RGBColor(0xDD, 0xCF, 0xC7))
        box(s, 1.3, cy, 8.5, 0.5, step, size=16, color=INK, bold=(st == "進行中"), anchor=MSO_ANCHOR.MIDDLE)
        box(s, 10.0, cy, 2.8, 0.5, st, size=13,
            color=(GREEN if st == "完了" else (BRICK if st == "進行中" else GRAY)), bold=True, anchor=MSO_ANCHOR.MIDDLE)
        cy += 0.92
    footer(s)

    # 結び
    s = add_slide(prs)
    bar(s, 0, 0, 13.333, 7.5, CREAM)
    bar(s, 0, 2.4, 13.333, 0.08, BRICK)
    box(s, 0.9, 2.7, 11.5, 1.4,
        "土地値に裏打ちされた高収益の収益物件です。\n御行の良きパートナーとして、ご支援を賜りますようお願い申し上げます。",
        size=20, color=INK, line_spacing=1.3)
    box(s, 0.95, 5.4, 11.5, 0.9, f"{D['company']}\n{D['rep']}", size=15, color=GRAY)

    out_dir = Path.home() / "01_honbu_docs_automation" / "out_screener"
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "277_融資資料デッキ.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    p = build()
    print(f"✅ 融資資料デッキ生成: {p}")

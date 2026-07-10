# -*- coding: utf-8 -*-
"""
新MacBook Air M4 セットアップ ＆ Claude開発環境 移行マニュアル
  - スプレッドシート(.xlsx) : 購入品チェックリスト + 移行手順 全部入り
  - スライド(.pptx)        : 同内容のサマリー版（KHDデザインシステム）
実機調査(2026-06-22)に基づく。機密値(トークン)は一切含めない。
"""
import os, shutil
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ===== デザインシステム（KHD：クリーム白 × レンガ赤）=====
CREAM   = "F9F6EF"
BRICK   = "AA2E26"
DARK    = "2B2B2B"
GRAY    = "8A8A8A"
ALT     = "F2ECE0"
LINE    = "D8CFC0"
WHITE   = "FFFFFF"
JP_FONT = "Hiragino Sans"

OUT_DIR = "/Users/kikuchikenta/01_honbu_docs_automation"
XLSX = os.path.join(OUT_DIR, "KHD_新MacBook移行マニュアル_v1.xlsx")
PPTX = os.path.join(OUT_DIR, "KHD_新MacBook移行マニュアル_v1.pptx")

# =========================================================
#  1) スプレッドシート
# =========================================================
def fill(hex_): return PatternFill("solid", fgColor=hex_)
def f(sz=11, b=False, color=DARK): return Font(name=JP_FONT, size=sz, bold=b, color=color)
thin = Side(style="thin", color=LINE)
BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)

def build_xlsx():
    wb = Workbook()

    # ---- 00_凡例 ----
    ws = wb.active; ws.title = "00_凡例"
    ws.sheet_view.showGridLines = False
    ws["B2"] = "新MacBook Air M4  セットアップ ＆ Claude開発環境 移行マニュアル"
    ws["B2"].font = f(16, True, BRICK)
    ws["B3"] = "作成: 2026-06-22 / 対象: 13インチ MacBook Air (整備済) M4 10C-CPU/10C-GPU ミッドナイト"
    ws["B3"].font = f(10, False, GRAY)
    legend = [
        ("色／記号", "意味", "あなたの動き"),
        ("🟡 入力", "自分で買う・自分で実行する作業", "ポチる／手を動かす"),
        ("🟧 要再ログイン", "アカウント再認証・トークン再入力が要る", "ログインし直す"),
        ("🟦 確認", "正しく移行できたかチェックする項目", "目視チェック"),
        ("⬜ 自動", "移行アシスタント／同期で自動的に入る", "待つだけ"),
    ]
    r = 5
    for i,(a,b,c) in enumerate(legend):
        head = (i==0)
        for j,val in enumerate((a,b,c)):
            cell = ws.cell(row=r, column=2+j, value=val)
            cell.font = f(11, head, WHITE if head else DARK)
            cell.fill = fill(BRICK if head else (ALT if i%2 else CREAM))
            cell.border = BORDER
            cell.alignment = Alignment(vertical="center", wrap_text=True)
        r += 1
    note = ("使い方：①01タブで備品をポチる → ②03タブの順でソフトを入れる → "
            "③04タブの手順を上から実行 → ④05でMCP再接続 → ⑤06で動作確認。"
            "最短ルートは『04手順A＝Appleの移行アシスタント』。")
    ws.cell(row=r+1, column=2, value="📌 "+note).font = f(11, True, DARK)
    ws.merge_cells(start_row=r+1, start_column=2, end_row=r+2, end_column=8)
    ws.cell(row=r+1, column=2).alignment = Alignment(wrap_text=True, vertical="top")
    for col,w in {"A":2,"B":22,"C":40,"D":22,"E":14,"F":14,"G":14,"H":14}.items():
        ws.column_dimensions[col].width = w

    def make_sheet(title, headers, rows, widths, tagcol=0):
        ws = wb.create_sheet(title)
        ws.sheet_view.showGridLines = False
        ws.cell(row=1, column=1, value=title.replace("_"," ")).font = f(14, True, BRICK)
        for j,h in enumerate(headers):
            c = ws.cell(row=3, column=1+j, value=h)
            c.font = f(11, True, WHITE); c.fill = fill(BRICK); c.border = BORDER
            c.alignment = Alignment(vertical="center", horizontal="center", wrap_text=True)
        for i,row in enumerate(rows):
            for j,val in enumerate(row):
                c = ws.cell(row=4+i, column=1+j, value=val)
                c.font = f(10, False, DARK)
                c.fill = fill(ALT if i%2 else CREAM)
                c.border = BORDER
                c.alignment = Alignment(vertical="center", wrap_text=True,
                                        horizontal="center" if j==tagcol else "left")
        for j,w in enumerate(widths):
            ws.column_dimensions[get_column_letter(1+j)].width = w
        ws.freeze_panes = "A4"
        return ws

    # ---- 01_購入品チェックリスト ----
    make_sheet("01_購入品チェックリスト",
        ["区分","品目","おすすめ／メモ","優先度","目安価格","検索リンク"],
        [
         ["🟡","🛡 AppleCare+ ★期限30日","落下・水濡れ補償。整備済でも追加可。購入から30日以内のみ加入可","最優先","約2.5万前後","Apple公式 > マイサポート"],
         ["🟡","📺 アンチグレアフィルム(1位)","パワーサポート アンチグレア（M4/2025/13で選ぶ）","高","2,000-3,500","https://www.amazon.co.jp/s?k=パワーサポート+MacBook+Air+13+M4+2025+アンチグレア"],
         ["🟡","　〃 (2位 コスパ)","NIMASO アンチグレア 2枚入りガイド付き","高","1,500-2,800","https://www.amazon.co.jp/s?k=NIMASO+MacBook+Air+13+M4+2025+アンチグレア"],
         ["🟡","　〃 (3位)","ミヤビックス OverLay Plus（のぞき見版も有）","中","2,000-3,500","https://www.amazon.co.jp/s?k=ミヤビックス+OverLay+MacBook+Air+13+M4+2025"],
         ["🟡","🔌 USB-Cハブ 7in1","Airはポート2つ→HDMI/USB-A/SD/PD用。Anker/Satechi/Belkin","高","4,000-7,000","https://www.amazon.co.jp/s?k=Anker+USB-C+ハブ+7in1+MacBook"],
         ["🟡","🧹 クリーニングクロス","ミッドナイトは指紋が目立つ→常備で満足度UP","高","500-1,980","https://www.amazon.co.jp/s?k=MacBook+クリーニングクロス+マイクロファイバー"],
         ["🟡","💼 スリーブ／インナーケース","持ち歩き保護。tomtoc / Incase / Bellroy","中","2,500-6,000","https://www.amazon.co.jp/s?k=MacBook+Air+13+M4+スリーブケース"],
         ["🟡","🎨 天板スキン(任意)","wraplus/dbrand。指紋＋傷防止・放熱に有利・剥がせて下取り綺麗","任意","3,000-4,000","https://www.amazon.co.jp/s?k=wraplus+MacBook+Air+13+M4+スキン"],
         ["🟡","🔋 GaN充電器35W(任意)","小型予備。持ち歩き用。Anker等","任意","3,000-5,000","https://www.amazon.co.jp/s?k=Anker+GaN+35W+充電器"],
         ["⬜","⛔ キーボードカバー","非推奨：液晶にキー跡・放熱阻害・打鍵感低下。買わない","-","-","-"],
        ],
        [6,26,40,8,12,52])

    # ---- 02_移行ハード ----
    make_sheet("02_移行ハード",
        ["要否","ハード","用途","メモ"],
        [
         ["🟡","新MacBook Air M4(購入済)","母艦","arm64。旧Mac(Apple Silicon)と同アーキ→移行は素直"],
         ["🟡","旧Mac本体","移行元","電源ON・同一Wi-Fiに置く。移行アシスタントで吸い出す"],
         ["🟦","同一Wi-Fi 環境","移行アシスタント通信","有線(USB-C/Thunderbolt直結)ならさらに高速・安定"],
         ["任意","USB-C↔USB-Cケーブル","直結移行","Wi-Fiが遅い時。Thunderbolt対応だと最速"],
         ["任意","外付けSSD","バックアップ保険","移行前にTime Machineを取っておくと安心"],
         ["🟦","Apple ID / iCloud","各種同期・AppleCare+","新Macで同じApple IDにサインイン"],
        ],
        [6,28,24,52])

    # ---- 03_移行ソフト_インストール順 ----
    make_sheet("03_移行ソフト_インストール順",
        ["順","ソフト/環境","入れ方（コマンド/入手先）","備考"],
        [
         ["1","macOSアップデート / Apple ID","設定 > 一般 > ソフトウェアアップデート","最新化＋サインイン"],
         ["2","Xcode Command Line Tools","xcode-select --install","git・コンパイラ。最初に必須"],
         ["3","Node.js v24系","nodejs.org の .pkg、または旧Mac同様の方法","npm同梱。/usr/local/bin/node"],
         ["4","npm 設定(prefix/cache)","npm config set prefix ~/.npm-global ／ npm config set cache ~/.npm-cache-khd","旧Mac設定を再現"],
         ["5","PATH追加(zsh)","echo 'export PATH=~/.npm-global/bin:~/.local/bin:$PATH' >> ~/.zshrc","claude・自作CLIを通す"],
         ["6","Claude Code","npm install -g @anthropic-ai/claude-code","本体(現行 2.1.158)"],
         ["7","uv / uvx","curl -LsSf https://astral.sh/uv/install.sh | sh","tradingview MCPがuvx依存"],
         ["8","Python パッケージ","python3 -m pip install python-pptx openpyxl pandas google-api-python-client google-auth google-auth-oauthlib requests-oauthlib","スライド/スプシ/Google系生成用"],
         ["9","Google Drive デスクトップ","google.com/drive/download → ec-seller@… でサインイン","CloudStorage同期が自動復活"],
         ["10","Google Chrome","google.com/chrome","Claude in Chrome拡張用"],
         ["11","Cursor / VS Code(任意)","各公式サイト","旧Macに導入済。エディタ"],
        ],
        [4,26,58,28])

    # ---- 04_移行手順 ----
    make_sheet("04_移行手順",
        ["手順","ステップ","内容","区分"],
        [
         ["A 推奨","A-0 バックアップ","旧MacでTime Machine or 重要フォルダ控え","🟦"],
         ["A 推奨","A-1 移行アシスタント","新Mac初期設定で『この情報を転送』→旧Macから移行","⬜"],
         ["A 推奨","A-2 そのまま全部移行","~/.claude 一式・~/.claude.json・~/.local/bin・~/.zshrc・01_honbu_docs_automation・Node/アプリが丸ごと入る","⬜"],
         ["A 推奨","A-3 再ログイン3点","①Claudeにログイン ②Google Driveサインイン ③Apple ID","🟧"],
         ["A 推奨","A-4 動作確認","06タブのチェックを実施","🟦"],
         ["B 手動","B-1 ソフト導入","03タブの順で1→11をインストール","🟡"],
         ["B 手動","B-2 Claude設定移植","旧Mac ~/.claude を新Macへコピー(CLAUDE.md/settings/agents/skills/projects/memory/scheduled-tasks/plugins)","🟡"],
         ["B 手動","B-3 .claude.json移植 or 再設定","旧 ~/.claude.json をコピー(MCP+トークン込)。しない場合は05で手動再追加","🟧"],
         ["B 手動","B-4 自作CLI移植","~/.local/bin の khd-log/khd-scan/khd-yt-upload をコピー＋実行権限(chmod +x)","🟡"],
         ["B 手動","B-5 プロジェクト移植","~/01_honbu_docs_automation を丸ごとコピー(.company/ secretary・scripts含む)","🟡"],
         ["B 手動","B-6 プラグイン","company@cc-company マーケット再追加→/plugin で再インストール(または ~/.claude/plugins コピー)","🟧"],
         ["B 手動","B-7 MCP再接続","05タブ参照","🟧"],
         ["B 手動","B-8 動作確認","06タブ","🟦"],
        ],
        [8,20,60,8], tagcol=3)

    # ---- 05_MCP再接続_トークン ----
    make_sheet("05_MCP再接続_トークン",
        ["MCP","種別","再接続の仕方","必要な認証"],
        [
         ["notion","stdio(npx)","claude mcp add notion -- npx -y @notionhq/notion-mcp-server","環境変数 NOTION_TOKEN を再設定(機密)"],
         ["googletasks","http(Composio)","claude mcp add-json で Composio tool_router の URL を登録","Composio側URL(トークン埋込)を再取得"],
         ["tradingview","stdio(uvx)","プロジェクト .mcp.json に定義済→uvxが入っていれば自動","uvxの導入(03-7)が前提"],
         ["Google Drive / Gmail / Calendar","アカウント連携","Claudeに同一アカウントでログイン→コネクタが同期/再認可","各GoogleのOAuth再承認"],
         ["KHD_APPS_SCRIPT_URL","環境変数","khd-log用。export KHD_APPS_SCRIPT_URL=… を ~/.zshrc 等へ","Apps ScriptのWeb App URL"],
        ],
        [22,16,56,30])
    # 注意行
    ws = wb["05_MCP再接続_トークン"]
    rr = ws.max_row + 2
    ws.cell(row=rr, column=1, value="⚠ トークン/URLはこの表に書かない。1Password等 or 旧Mac ~/.claude.json から移すこと").font = f(10, True, BRICK)
    ws.merge_cells(start_row=rr, start_column=1, end_row=rr, end_column=4)

    # ---- 06_動作確認チェック ----
    make_sheet("06_動作確認チェック",
        ["No","確認コマンド/操作","期待結果","OK"],
        [
         ["1","claude --version","2.1.x が出る","□"],
         ["2","node --version / npm --version","v24系 / 11系","□"],
         ["3","which khd-log","~/.local/bin/khd-log","□"],
         ["4","claude 起動→/mcp","notion・googletasks・tradingview・Drive/Gmail/Calendarが緑","□"],
         ["5","Notion検索を1回試す","COMPANY KHDまとめ が引ける","□"],
         ["6","python3 -c 'import pptx,openpyxl,pandas'","エラーなし","□"],
         ["7","ls ~/.claude/projects/.../memory","Auto memoryが揃っている","□"],
         ["8","Google Drive同期フォルダを開く","CloudStorageが復活している","□"],
         ["9","プロジェクトで /company","秘書室(.company)が読める","□"],
         ["10","AppleCare+加入状況","設定>一般>情報 or マイサポートで有効","□"],
        ],
        [4,44,40,6], tagcol=3)

    wb.save(XLSX)
    print("XLSX saved:", XLSX)

# =========================================================
#  2) スライド
# =========================================================
def rgb(h): return RGBColor.from_string(h)

def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def bg(slide, color=CREAM):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(color)

    def box(slide, l, t, w, h, fill_=None, line_=None, line_w=1.0):
        from pptx.enum.shapes import MSO_SHAPE
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        sp.adjustments[0] = 0.06
        if fill_: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill_)
        else: sp.fill.background()
        if line_: sp.line.color.rgb = rgb(line_); sp.line.width = Pt(line_w)
        else: sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def txt(slide, l, t, w, h, text, size=18, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.05):
        tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line_sp
            r = p.add_run(); r.text = ln
            r.font.name = JP_FONT; r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = rgb(color)
        return tb

    def header(slide, kicker, title):
        box(slide, Inches(0), Inches(0), SW, Inches(1.15), fill_=BRICK)
        txt(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.35), kicker,
            size=12, bold=True, color="F4D9D5")
        txt(slide, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.62), title,
            size=24, bold=True, color=WHITE)

    def bullets(slide, l, t, w, items, size=15, gap=10):
        tb = slide.shapes.add_textbox(l, t, w, Inches(5)); tf = tb.text_frame
        tf.word_wrap = True
        for i,(mark, s, b) in enumerate(items):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.space_after = Pt(gap); p.line_spacing = 1.08
            r1 = p.add_run(); r1.text = mark+"  "
            r1.font.name=JP_FONT; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=rgb(BRICK)
            r2 = p.add_run(); r2.text = s
            r2.font.name=JP_FONT; r2.font.size=Pt(size); r2.font.bold=b; r2.font.color.rgb=rgb(DARK)
        return tb

    # ---- S1 表紙 ----
    s = prs.slides.add_slide(blank); bg(s)
    box(s, Inches(0), Inches(2.3), SW, Inches(2.9), fill_=BRICK)
    txt(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.5),
        "KHD｜社内マニュアル  2026-06-22", size=14, bold=True, color="F4D9D5")
    txt(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.3),
        "新MacBook Air M4 セットアップ\n＆ Claude開発環境 まるごと移行", size=34, bold=True, color=WHITE, line_sp=1.05)
    txt(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.6),
        "対象機：13インチ MacBook Air（整備済）M4 10コアCPU/10コアGPU・ミッドナイト",
        size=14, bold=False, color=GRAY)
    txt(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.6),
        "本書 = 買う物リスト＋移行ハード／ソフト／手順／MCP再接続／動作確認（詳細はスプシ連動）",
        size=12, bold=False, color=GRAY)

    # ---- S2 結論サマリー ----
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "CONCLUSION｜最短ルート", "結論：『移行アシスタント』＋『再ログイン3点』で9割終わる")
    box(s, Inches(0.6), Inches(1.45), Inches(5.9), Inches(5.4), fill_=WHITE, line_=LINE)
    txt(s, Inches(0.9), Inches(1.65), Inches(5.4), Inches(0.5), "✅ 最短(推奨)：手順A", size=18, bold=True, color=BRICK)
    bullets(s, Inches(0.9), Inches(2.25), Inches(5.4), [
        ("①","旧Macを横に置く（電源ON・同Wi-Fi）", False),
        ("②","新Mac初期設定で『移行アシスタント』", True),
        ("③","~/.claude 一式・.claude.json・自作CLI・プロジェクトが丸ごと移動", False),
        ("④","再ログイン3点：Claude / Google Drive / Apple ID", True),
        ("⑤","/mcp が全部緑＝完了", False),
    ], size=15, gap=12)
    box(s, Inches(6.8), Inches(1.45), Inches(5.9), Inches(5.4), fill_=ALT, line_=LINE)
    txt(s, Inches(7.1), Inches(1.65), Inches(5.4), Inches(0.5), "🧱 同アーキで安心", size=18, bold=True, color=BRICK)
    bullets(s, Inches(7.1), Inches(2.25), Inches(5.4), [
        ("●","旧Mac＝Apple Silicon、新Mac＝M4 → 同じarm64", True),
        ("●","Node・Claude・Pythonライブラリがそのまま動く", False),
        ("●","移行アシスタントが使えない時のみ『手順B＝手動再構築』", False),
        ("⚠","トークン/URLは表に書かず旧.claude.jsonから移す", True),
        ("📅","AppleCare+は購入30日以内が期限", True),
    ], size=15, gap=12)

    # ---- S3 購入したMac仕様 ----
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "DEVICE｜購入機", "購入した1台")
    specs = [("機種","13インチ MacBook Air（整備済製品）"),
             ("チップ","Apple M4（10コアCPU / 10コアGPU）"),
             ("カラー","ミッドナイト（※指紋が目立つ→クロス常備）"),
             ("アーキ","Apple Silicon arm64（旧Macと同一＝移行容易）"),
             ("ポート","Thunderbolt 4(USB-C)×2 / MagSafe / イヤホン端子"),
             ("外部出力","M4から外部ディスプレイ2台対応")]
    y = 1.6
    for k,v in specs:
        box(s, Inches(0.8), Inches(y), Inches(2.6), Inches(0.7), fill_=BRICK)
        txt(s, Inches(0.9), Inches(y+0.12), Inches(2.4), Inches(0.5), k, size=14, bold=True, color=WHITE)
        box(s, Inches(3.5), Inches(y), Inches(8.9), Inches(0.7), fill_=WHITE, line_=LINE)
        txt(s, Inches(3.7), Inches(y+0.12), Inches(8.6), Inches(0.5), v, size=14, bold=False, color=DARK)
        y += 0.82

    # ---- S4 買うべき備品 ----
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "ACCESSORIES｜備品", "買うべき備品（保護＝AppleCare+が最優先）")
    box(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.75), fill_=BRICK)
    txt(s, Inches(0.85), Inches(1.52), Inches(11.6), Inches(0.5),
        "🛡 最優先：AppleCare+（落下・水濡れ補償／購入30日以内のみ加入可・約2.5万前後）",
        size=16, bold=True, color=WHITE)
    txt(s, Inches(0.7), Inches(2.35), Inches(12), Inches(0.4), "📺 アンチグレアフィルム 3種ランキング", size=15, bold=True, color=BRICK)
    rank = [("🥇","パワーサポート アンチグレア","映り込み・指紋を自然に抑え見やすさNo.1。貼り精度◎"),
            ("🥈","NIMASO アンチグレア","コスパ最強・2枚入りガイド付きで失敗しにくい"),
            ("🥉","ミヤビックス OverLay Plus","専用カット精度高・反射防止強・のぞき見版も")]
    y=2.75
    for m,t_,d in rank:
        box(s, Inches(0.7), Inches(y), Inches(12), Inches(0.62), fill_=WHITE, line_=LINE)
        txt(s, Inches(0.85), Inches(y+0.1), Inches(0.6), Inches(0.45), m, size=18, bold=True)
        txt(s, Inches(1.5), Inches(y+0.12), Inches(4.0), Inches(0.45), t_, size=14, bold=True, color=DARK)
        txt(s, Inches(5.6), Inches(y+0.13), Inches(7.0), Inches(0.45), d, size=12, color=GRAY)
        y+=0.72
    txt(s, Inches(0.7), Inches(y+0.05), Inches(12), Inches(1.2),
        "🔌 USB-Cハブ7in1（ポート2つ補完・必須）／🧹 クロス／💼 スリーブ／🎨 天板スキン(任意)／🔋 GaN35W(任意)\n"
        "⛔ キーボードカバーは非推奨（液晶にキー跡・放熱阻害・打鍵感低下）",
        size=13, bold=False, color=DARK, line_sp=1.2)

    # ---- S5 移行ソフト(順) ----
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "SOFTWARE｜導入順", "新Macに入れるソフト（この順で）")
    steps = [
        ("1","macOS更新＋Apple IDサインイン"),
        ("2","Xcode Command Line Tools（xcode-select --install）"),
        ("3","Node.js v24系（nodejs.org .pkg）"),
        ("4","npm設定：prefix ~/.npm-global / cache ~/.npm-cache-khd"),
        ("5","PATH追加：~/.npm-global/bin と ~/.local/bin を ~/.zshrc へ"),
        ("6","Claude Code：npm i -g @anthropic-ai/claude-code"),
        ("7","uv / uvx（curl -LsSf https://astral.sh/uv/install.sh | sh）"),
        ("8","pip：python-pptx openpyxl pandas google-api-python-client 他"),
        ("9","Google Drive デスクトップ（ec-seller@… でサインイン）"),
        ("10","Chrome / Cursor / VS Code（任意）"),
    ]
    col_x=[Inches(0.7), Inches(6.9)]; y0=1.55
    for i,(n,t_) in enumerate(steps):
        col=0 if i<5 else 1; row=i%5
        x=col_x[col]; y=y0+row*1.0
        box(s, x, Inches(y), Inches(0.7), Inches(0.85), fill_=BRICK)
        txt(s, x, Inches(y+0.18), Inches(0.7), Inches(0.5), n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        box(s, x+Inches(0.85), Inches(y), Inches(5.1), Inches(0.85), fill_=WHITE, line_=LINE)
        txt(s, x+Inches(1.0), Inches(y+0.1), Inches(4.85), Inches(0.7), t_, size=12.5, bold=False, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    # ---- S6 手順A/B ----
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "STEPS｜移行手順", "手順A（推奨：自動）と 手順B（手動：クリーン再構築）")
    box(s, Inches(0.6), Inches(1.4), Inches(5.95), Inches(5.5), fill_=WHITE, line_=BRICK, line_w=1.5)
    txt(s, Inches(0.85), Inches(1.55), Inches(5.5), Inches(0.5), "A. 移行アシスタント（推奨）", size=17, bold=True, color=BRICK)
    bullets(s, Inches(0.85), Inches(2.15), Inches(5.5), [
        ("A-0","旧MacでTime Machineバックアップ", False),
        ("A-1","新Mac初期設定→『情報を転送』", True),
        ("A-2","~/.claude/.claude.json/.local/bin/プロジェクト等が丸ごと移動", False),
        ("A-3","再ログイン：Claude / Drive / Apple ID", True),
        ("A-4","/mcp と 06動作確認", False),
    ], size=14, gap=11)
    box(s, Inches(6.8), Inches(1.4), Inches(5.95), Inches(5.5), fill_=ALT, line_=LINE)
    txt(s, Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.5), "B. 手動再構築（クリーンに作る時）", size=17, bold=True, color=BRICK)
    bullets(s, Inches(7.05), Inches(2.15), Inches(5.5), [
        ("B-1","ソフトを03の順で導入", False),
        ("B-2","~/.claude をコピー(CLAUDE.md/agents/skills/projects/memory…)", True),
        ("B-3",".claude.json移植 or 05で手動再追加", False),
        ("B-4","~/.local/bin の自作CLIをコピー＋chmod +x", False),
        ("B-5","~/01_honbu_docs_automation を丸ごとコピー", True),
        ("B-6","plugin: company@cc-company 再追加", False),
        ("B-7","MCP再接続（05）", True),
    ], size=14, gap=9)

    # ---- S7 MCP再接続＆動作確認 ----
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "MCP & CHECK｜仕上げ", "MCP再接続 ＆ 動作確認")
    txt(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.4), "🔌 MCP再接続（要OAuth/トークン）", size=15, bold=True, color=BRICK)
    mcp=[("notion","npx @notionhq/notion-mcp-server＋NOTION_TOKEN"),
         ("googletasks","Composio tool_router URL を再登録"),
         ("tradingview","プロジェクト.mcp.json＋uvx（自動）"),
         ("Drive/Gmail/Calendar","Claude同一アカウントでログイン→同期/再認可")]
    y=1.8
    for k,v in mcp:
        box(s, Inches(0.7), Inches(y), Inches(3.0), Inches(0.6), fill_=BRICK)
        txt(s, Inches(0.8), Inches(y+0.08), Inches(2.8), Inches(0.45), k, size=12.5, bold=True, color=WHITE)
        box(s, Inches(3.8), Inches(y), Inches(8.8), Inches(0.6), fill_=WHITE, line_=LINE)
        txt(s, Inches(3.95), Inches(y+0.08), Inches(8.5), Inches(0.45), v, size=12.5, color=DARK)
        y+=0.7
    txt(s, Inches(0.7), Inches(y+0.1), Inches(12), Inches(0.4), "🟦 動作確認（抜粋）", size=15, bold=True, color=BRICK)
    txt(s, Inches(0.7), Inches(y+0.55), Inches(12), Inches(1.4),
        "□ claude --version が 2.1.x  □ /mcp が全部緑  □ Notionで『COMPANY KHDまとめ』が引ける\n"
        "□ python3 -c 'import pptx,openpyxl,pandas' でエラー無し  □ memory が揃う  □ /company で秘書室が読める  □ AppleCare+有効",
        size=13, color=DARK, line_sp=1.3)
    txt(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
        "⚠ トークン/URLはスライド・スプシに書かない。旧Mac ~/.claude.json か 1Password から移すこと。",
        size=12, bold=True, color=BRICK)

    prs.save(PPTX)
    print("PPTX saved:", PPTX)

if __name__ == "__main__":
    build_xlsx()
    build_pptx()
    print("DONE")

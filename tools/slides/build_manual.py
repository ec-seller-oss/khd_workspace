# -*- coding: utf-8 -*-
"""KHD 財務×営業 操縦マニュアル【完成版v2】。
番号付きタブ(01_〜11_)・税理士修正後の確定数字・実物レイアウト忠実再現。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CREAM=RGBColor(0xF9,0xF6,0xEF); BRICK=RGBColor(0xAA,0x2E,0x26)
DARK=RGBColor(0x33,0x33,0x33); WHITE=RGBColor(0xFF,0xFF,0xFF); GRAY=RGBColor(0x88,0x88,0x88)
Y=RGBColor(0xFF,0xF2,0xCC); O=RGBColor(0xFC,0xE5,0xCD); B=RGBColor(0xCF,0xE2,0xF3)
G=RGBColor(0xEC,0xEC,0xEC); GREEN=RGBColor(0xD9,0xEA,0xD3); GOLD=RGBColor(0xFF,0xE5,0x99)
FONT="Meiryo"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; blank=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(blank)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    bg.fill.solid(); bg.fill.fore_color.rgb=CREAM; bg.line.fill.background(); bg.shadow.inherit=False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2,bg._element); return s
def box(s,x,y,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp
def txt(s,x,y,w,h,t,size=14,color=DARK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(4);tf.margin_right=Pt(4);tf.margin_top=Pt(2);tf.margin_bottom=Pt(2)
    for i,ln in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=color
    return tb
def header(s,title,sub=None):
    box(s,0,0,13.333,1.0,fill=BRICK)
    txt(s,0.4,0.08,12.5,0.85,title,size=25,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    if sub: txt(s,0.4,0.74,12.5,0.3,sub,size=12,color=RGBColor(0xF0,0xD8,0xD5))
def chip(s,x,y,w,h,label,fill,tc=DARK,size=12,bold=True):
    box(s,x,y,w,h,fill=fill,line=RGBColor(0xDD,0xDD,0xDD),lw=0.75)
    txt(s,x,y,w,h,label,size=size,color=tc,bold=bold,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# 1 表紙
s=slide(); box(s,0,0,13.333,7.5,fill=BRICK)
txt(s,0,2.0,13.333,1.3,"KHD 財務 × 営業 操縦マニュアル",size=40,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
txt(s,0,3.35,13.333,0.6,"【完成版 v2】開けば分かる｜タブは使う順に 01〜11",size=20,color=GOLD,bold=True,align=PP_ALIGN.CENTER)
txt(s,0,4.1,13.333,0.7,"日々のタイムライン × どのタブ・どのセル × 判断の仕方",size=16,color=RGBColor(0xF0,0xD8,0xD5),align=PP_ALIGN.CENTER)
txt(s,0,5.9,13.333,0.5,"本体スプシ『01_KHD PJ一覧_v2_260601』  /  2026-06-05 税理士・損益分岐 修正反映済",size=12,color=RGBColor(0xE8,0xC8,0xC4),align=PP_ALIGN.CENTER)

# 2 タブバー（実物の並び）
s=slide(); header(s,"① タブは「使う順」01〜11 に整理済","見えてるのは操縦席だけ。データ系13枚は非表示・不要2枚は削除")
txt(s,0.4,1.2,12.5,0.4,"▼ 画面下のタブ（左から使う順）",size=13,color=BRICK,bold=True)
tabs=["pass","01_統合司令塔","02_追客リスト","03_売上見込み","04_損益PL","05_資金繰り","06_資産負債BS","07_本部マトリクス","08_残高クイック入力","09_MF残高","10_収益パイプライン","11_KPI"]
x=0.4; y=1.65
for i,t in enumerate(tabs):
    w=0.9 if t=="pass" else min(1.7,0.5+len(t)*0.135)
    fill=WHITE if t=="pass" else (O if i in(1,2) else (Y if i in(3,) else (B if i in(4,5,6) else CREAM)))
    if x+w>12.9: x=0.4; y+=0.62
    box(s,x,y,w,0.5,fill=fill,line=RGBColor(0xCC,0xCC,0xCC),lw=0.75)
    txt(s,x,y,w,0.5,t,size=9.5,color=DARK,bold=(t!="pass"),align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    x+=w+0.08
# 凡例＝意味
txt(s,0.4,3.0,12.5,0.4,"色＝役割（タブも中のセルも同じルール）",size=13,color=BRICK,bold=True)
leg=[("🟧 毎日見る：01司令塔・02追客",O),("🟡 入力：03売上見込み",Y),("🟦 見る(損益/現金)：04PL・05資金繰り・06BS",B),("⬜ 補助/データ：07〜11・pass",CREAM)]
lx=0.4
for lab,c in leg:
    box(s,lx,3.45,3.05,0.55,fill=c,line=RGBColor(0xDD,0xDD,0xDD),lw=0.6)
    txt(s,lx+0.05,3.45,2.95,0.55,lab,size=9.5,color=DARK,anchor=MSO_ANCHOR.MIDDLE); lx+=3.12
box(s,0.4,4.25,12.5,2.6,fill=WHITE,line=RGBColor(0xDD,0xDD,0xDD),lw=1)
txt(s,0.6,4.4,12.1,2.4,
"・非表示（データ源・スポット13枚）：⑤借入／借入返済／借入条件／税金／諸経費／クレカ用途／按分／役員報酬／経費削減・損切り／銀行提出前クリーンアップ／事業計画／未来会計図表／WBS\n   → 計算の元なので消さず非表示。表示メニュー（左下の三本線）からいつでも戻せる。\n・削除（不要2枚）：⑥使い方（このマニュアルで代替）／📋按分棚卸し_テンプレ（重複）\n・pass（ID/パスワード/サブスク管理）＝よく使うので表示のまま先頭に。\n\n★ ポイント：番号を変えても中の計算・参照は自動で追従。日常は 01〜06 を見れば回る。",
size=13,color=DARK,ls=1.25)

# 3 連動マップ
s=slide(); header(s,"② 連動マップ：数字はこう流れる","入口(黄)を正直に埋めるだけ。あとは自動で右へ")
def node(s,x,y,w,h,title,sub,fill):
    box(s,x,y,w,h,fill=fill,line=RGBColor(0xCC,0xCC,0xCC),lw=1)
    txt(s,x+0.1,y+0.1,w-0.2,0.5,title,size=13,color=DARK,bold=True)
    txt(s,x+0.1,y+0.6,w-0.2,h-0.7,sub,size=10,color=DARK,ls=1.0)
def arrow(s,x,y,w):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(0.32)); a.shadow.inherit=False
    a.fill.solid(); a.fill.fore_color.rgb=BRICK; a.line.fill.background()
node(s,0.4,1.5,3.0,1.6,"✍ 03_売上見込み","E満額/F確度/G着金月\n＝金額のSSoT",Y)
node(s,0.4,3.5,3.0,1.5,"✍ 04_損益PL ドライバー","B55/B56転換率\nB60稼働日",Y)
node(s,0.4,5.3,3.0,1.4,"🎯 02_追客リスト","誰を追う・確度\n次アクション",O)
arrow(s,3.5,2.2,1.0); arrow(s,3.5,4.15,1.0)
node(s,4.7,1.5,3.3,1.6,"05_資金繰り(自動)","入金=売上見込み連動\n→月末現金・ランウェイ",B)
node(s,4.7,3.4,3.3,1.7,"04_損益PL 逆算(自動)","家族ライン−粗利=穴\n÷単価=必要件数→追客数",B)
arrow(s,8.1,2.55,1.0); arrow(s,8.1,3.95,1.0)
node(s,9.3,1.9,3.6,2.9,"01_統合司令塔","毎日の操縦席\n・現預金/ランウェイ\n・通期の谷\n・毎月の穴\n・今日の一手/追客目標\n・営業直結比率",O)
arrow(s,3.5,5.9,5.7)
txt(s,0.4,6.95,12.5,0.4,"※触るのは黄(03売上見込みE/F/G・04PLのB55/B56/B60・02追客)だけ。青/データは見るだけ。",size=12,color=BRICK,bold=True)

# 4 毎日のタイムライン
s=slide(); header(s,"③ 毎日のタイムライン：いつ・どこを見て・何を打つ")
rows=[("🌅 朝(5分)","01_統合司令塔",O,"見る：今日の優先①②／今日の一手(B52)／追客目標(B62)／現預金・ランウェイ","→ 今日やる1〜2件を決める"),
 ("🌅 朝","02_追客リスト",O,"見る：優先スコア順／C3『今日の推奨追客』","→ 今日追う1人を確定"),
 ("☀ 日中","02_追客リスト",Y,"打つ：追客したら最終接触・次アクションを更新","→ 動いた案件はその場で1行"),
 ("☀ 日中","03_売上見込み",Y,"打つ：案件が動いたら 確度(F)・着金月(G)／新規はE/F/G追記","→ 資金繰りに自動反映"),
 ("🌙 夜(3分)","01_統合司令塔",O,"見る：営業直結比率(目標60%)／今日の一手は消化したか","→ 日報に結果○△×")]
y=1.25
for tag,tab,c,body,note in rows:
    box(s,0.4,y,12.5,1.02,fill=WHITE,line=RGBColor(0xDD,0xDD,0xDD),lw=0.75)
    chip(s,0.5,y+0.12,1.7,0.78,tag,c,size=11)
    txt(s,2.35,y+0.08,2.3,0.85,tab,size=12,color=BRICK,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.7,y+0.06,8.0,0.55,body,size=11.5,color=DARK)
    txt(s,4.7,y+0.6,8.0,0.38,note,size=10,color=GRAY); y+=1.1
txt(s,0.4,7.0,12.5,0.35,"毎日触るのは『01_司令塔』と『02_追客リスト』の2枚だけ。",size=12,color=BRICK,bold=True)

# 5 週次・月次
s=slide(); header(s,"④ 週次・月次ルーティン")
txt(s,0.4,1.15,8,0.4,"■ 週次（営業会議・15分）",size=16,color=BRICK,bold=True)
for i,(tab,body) in enumerate([("02_追客リスト","全件のステージ/確度を見直し、要再開を拾う"),
    ("03_売上見込み","期待値合計B21・採用額B20を確認（今月いくら見込めるか）"),
    ("10_収益パイプライン","取りこぼしチェック>0なら売上見込みに起票（例:高松）"),
    ("04_損益PL","B55/B56(転換率)を実績で更新→1日の必要追客数が現実化")]):
    y=1.6+i*0.68; chip(s,0.5,y,2.7,0.55,tab,B,size=10.5)
    txt(s,3.3,y,9.5,0.55,body,size=12,color=DARK,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.4,4.5,8,0.4,"■ 月初（経営の棚卸し・20分）",size=16,color=BRICK,bold=True)
for i,(tab,body) in enumerate([("03_売上見込み B2","シナリオ(弱気/現実/強気)を当月想定に切替"),
    ("08_残高クイック入力","各口座残高を実額入力→BS・資金繰りの月初現金が更新"),
    ("05_資金繰り","固定費(B14:B26)に変動あれば反映／月末現金・谷を確認"),
    ("04_損益PL B60","当月の稼働日を入力")]):
    y=4.95+i*0.6; chip(s,0.5,y,2.7,0.52,tab,Y,size=10.5)
    txt(s,3.3,y,9.5,0.52,body,size=12,color=DARK,anchor=MSO_ANCHOR.MIDDLE)

def core(title,sub,role,watch,inputs,judge,layout_rows):
    s=slide(); header(s,title,sub)
    box(s,0.4,1.2,6.1,5.9,fill=WHITE,line=RGBColor(0xDD,0xDD,0xDD),lw=1); yy=1.35
    txt(s,0.6,yy,5.7,0.4,"● 役割",size=13,color=BRICK,bold=True); yy+=0.42
    txt(s,0.6,yy,5.7,0.6,role,size=12,color=DARK,ls=1.05); yy+=0.78
    txt(s,0.6,yy,5.7,0.4,"👀 見る（頻度つき）",size=13,color=BRICK,bold=True); yy+=0.42
    txt(s,0.6,yy,5.7,1.1,watch,size=11,color=DARK,ls=1.1); yy+=1.18
    txt(s,0.6,yy,5.7,0.4,"✍ 打つ",size=13,color=BRICK,bold=True); yy+=0.42
    txt(s,0.6,yy,5.7,0.9,inputs,size=11,color=DARK,ls=1.1); yy+=0.98
    txt(s,0.6,yy,5.7,0.4,"⚖ 判断",size=13,color=BRICK,bold=True); yy+=0.42
    txt(s,0.6,yy,5.7,1.0,judge,size=11,color=DARK,ls=1.1)
    box(s,6.7,1.2,6.2,5.9,fill=WHITE,line=RGBColor(0xDD,0xDD,0xDD),lw=1)
    txt(s,6.85,1.3,5.9,0.35,"▼ 実物イメージ（色＝運用ルール）",size=11,color=GRAY,bold=True); ry=1.75
    for cells in layout_rows:
        rx=6.85
        for (label,w,fill) in cells:
            box(s,rx,ry,w,0.42,fill=fill,line=RGBColor(0xCC,0xCC,0xCC),lw=0.5)
            txt(s,rx+0.03,ry,w-0.05,0.42,label,size=9,color=DARK,anchor=MSO_ANCHOR.MIDDLE); rx+=w
        ry+=0.45
    return s

core("コア① 01_統合司令塔 — 毎日の操縦席","朝晩ここを開く。あなたが言った『大事なシート』",
 "過去(BS)→現在(現金/ランウェイ)→未来(谷/穴)→今日やること を1枚に。",
 "毎日：今日の優先①〜⑤／今日の一手／追客目標\n毎日：現預金・ランウェイ(🟢/🟡/🔴)\n週次：通期の谷・毎月の穴・営業直結比率",
 "なし（全自動表示）。『見て動く』専用。",
 "・ランウェイ🔴<3ヶ月→投資凍結\n・営業直結<60%→内務を切って追客へ\n・①緊急→②生命線→③追客 の順で処理",
 [[("01_統合司令塔 — 過去→現在→未来＋行動",6.0,BRICK)],
  [("現預金 7,520,809円",3.0,O),("ランウェイ 11.9ヶ月🟢",3.0,O)],
  [("通期の谷 4,602,420円 🟢",6.0,B)],
  [("毎月の穴 434,857円(税理士修正後)",6.0,B)],
  [("家族ライン 577,587 − 粗利 142,731",6.0,B)],
  [("今日の一手：医療 月0.7件で家族ライン突破",6.0,O)],
  [("③追客：持倉様(栄町)→6/20決済",6.0,O)],
  [("営業直結比率 26%(目標60%) 🔴",6.0,B)]])
core("コア② 02_追客リスト — 誰を追うか(全案件SSoT)","医療/物件/cloud mil/協業を優先スコア順に",
 "全本部の見込み客を1枚に。優先スコアで自動並べ替え、今日の1人を司令塔へ。",
 "毎日：上位の優先スコア／C3『今日の推奨追客』\n週次：要再開(COLD)を拾い直す",
 "随時：追客したら最終接触・次アクション・期限を更新\n新規客は1行追加",
 "・優先スコア高=今日触る\n・確度はステージで決める\n・賃貸新規は『05_追客管理シート』で受けHOT化で昇格",
 [[("優先",0.7,G),("対象名",2.0,G),("確度%",0.9,Y),("温度",0.9,G),("状態",1.5,Y)],
  [("105",0.7,B),("内山先生",2.0,WHITE),("65",0.9,Y),("WARM",0.9,WHITE),("進行",1.5,Y)],
  [("90",0.7,B),("山崎先生(京橋)",2.0,WHITE),("90",0.9,Y),("HOT",0.9,WHITE),("進行",1.5,Y)],
  [("140",0.7,B),("持倉様(栄町)",2.0,WHITE),("100",0.9,Y),("HOT",0.9,WHITE),("進行",1.5,Y)],
  [("70",0.7,B),("曾我先生",2.0,WHITE),("70",0.9,Y),("WARM",0.9,WHITE),("提案中",1.5,Y)]])
core("コア③ 03_売上見込み — 金額の真実源(SSoT)","ここのE/F/Gが資金繰り・PL・司令塔すべての土台",
 "案件ごとの満額・確度・着金月を持つ唯一の場所。確度加重(M)と採用額(K)を自動算出し05資金繰りへ。",
 "週次：期待値合計B21（確度加重の収入見込み）／採用額B20",
 "案件が動くたび：満額(E)・確度%(F)・着金月(G)\n月初：シナリオB2(弱気/現実/強気)／新規は空き行に1行",
 "・確度はステージで(提案前10/提案済30/前向き50/交渉中70/内諾90)\n・着金月は2026/06形式(資金繰りと一致)\n・科目は3つだけ(栄町売却/医療テナント/買取再販)",
 [[("科目",2.2,Y),("満額E",1.3,Y),("確度F",1.0,Y),("着金月G",1.5,Y)],
  [("栄町売却",2.2,WHITE),("1,900,000",1.3,WHITE),("100%",1.0,WHITE),("2026/06",1.5,WHITE)],
  [("医療テナント",2.2,WHITE),("1,100,000",1.3,WHITE),("60%",1.0,WHITE),("2026/09",1.5,WHITE)],
  [("(新規はここへ)",2.2,Y),("—",1.3,Y),("—",1.0,Y),("—",1.5,Y)],
  [("採用額合計 B20",4.5,B),("1,900,000",2.0,B)],
  [("期待値合計 B21",4.5,B),("2,560,000",2.0,B)]])
core("コア④ 04_損益PL 営業ドライバー — 逆算エンジン","家族ラインから『今日の追客数』まで自動逆算。経常利益は金色で強調",
 "家族が潰れないライン−確定粗利=毎月の穴。これを必要件数→必要アポ→1日の追客数に翻訳。",
 "毎日：今日の一手(B52)／追客目標(B62)／時間の一手(B66)\n月次：必要件数・1日の追客数・月次トラッキング",
 "月1回：B55追客→アポ／B56アポ→成約 転換率\n月初：B60稼働日",
 "・穴>0=営業せよのサイン\n・転換率は実績で更新→1日の追客目標が信用できる値に\n・経常利益(金色)＝結論。ここがプラスか見る",
 [[("家族ライン B33",3.3,G),("577,587円",2.5,G)],
  [("− 確定粗利 B32",3.3,G),("142,731円",2.5,G)],
  [("= 毎月の穴 B34",3.3,B),("434,857円",2.5,B)],
  [("÷ 医療確度込/件",3.3,G),("660,000円",2.5,G)],
  [("= 必要医療件数",3.3,B),("0.66件/月",2.5,B)],
  [("→ 1日の追客",3.3,O),("約1件/日",2.5,O)],
  [("経常利益(全社・金色)",3.3,GOLD),("結論ライン",2.5,GOLD)]])
core("コア⑤ 05_資金繰り — 会社の体力計","入金は売上見込みから自動。見るのが基本。税理士修正反映済",
 "月初現金→入金−出金→月末現金を8ヶ月先まで。通期の谷とランウェイで『いつ苦しいか』。",
 "週次/月次：月末現預金(B29:I29＝体力)／ランウェイ判定／通期の谷",
 "変わった時だけ：固定費・出金(B14:B26)／育休(B10)\n※入金(B5:B8)は触らない＝売上見込み連動",
 "・税理士を正値に修正(法人18,333/個人46,750)→穴が67千円縮小\n・月末現金がどこかで赤→入金前倒し or 出金カット\n・栄町190万(6月)が生命線",
 [[("月初現金(司令塔連動)",4.0,G),("7,520,809",2.0,G)],
  [("入金=売上見込み連動",4.0,G),("自動",2.0,G)],
  [("✍ 法人税理士(修正)",4.0,Y),("18,333",2.0,Y)],
  [("✍ 個人税理士(修正)",4.0,Y),("46,750",2.0,Y)],
  [("月末現金=体力計",4.0,B),("8,743,221",2.0,B)],
  [("ランウェイ判定",4.0,B),("🟢攻めOK",2.0,B)]])

# 最終 3原則
s=slide(); header(s,"⑤ これだけ守れば回る — 3原則")
for i,(t,c,b) in enumerate([("① 触るのは黄色だけ",Y,"03売上見込みのE/F/G・04PLのB55/B56/B60・02追客。\nあとは全部自動。青/データは見るだけ。"),
      ("② 数字は1か所だけ(SSoT)",B,"金額=03売上見込み／追客=02追客リスト。\n同じ情報を2か所に書かない(二重入力禁止)。"),
      ("③ 毎日は2枚だけ開く",O,"朝晩『01_統合司令塔』と『02_追客リスト』。\n動いた案件はその場で1行更新＝報連相がデータになる。")]):
    y=1.5+i*1.78; box(s,0.8,y,11.7,1.55,fill=WHITE,line=RGBColor(0xDD,0xDD,0xDD),lw=1); box(s,0.8,y,0.35,1.55,fill=c)
    txt(s,1.3,y+0.12,11.0,0.5,t,size=18,color=BRICK,bold=True)
    txt(s,1.3,y+0.68,11.0,0.8,b,size=13,color=DARK,ls=1.1)
txt(s,0.8,7.0,11.7,0.4,"色：🟡入力 / 🟧毎日見る / 🟦週次月次で見る / ⬜自動・データ（各タブ下部にも凡例あり）",size=12,color=GRAY,align=PP_ALIGN.CENTER)

prs.save("/Users/kikuchikenta/01_honbu_docs_automation/KHD_操縦マニュアル.pptx")
print("saved v2 /", len(prs.slides._sldIdLst), "slides")

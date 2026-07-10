#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""KHD AI査定エンジン 操作マニュアル＆全体設計デッキ（視覚化）"""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent / "scripts"))
import loan_deck as L
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM, BRICK, INK, GRAY, GREEN, WHITE = L.CREAM, L.BRICK, L.INK, L.GRAY, L.GREEN, L.WHITE
BLUE = RGBColor(0x2B, 0x57, 0x9A)


def chip(s, x, y, w, h, text, fill=WHITE, fg=INK, size=12, bold=True, line=BRICK):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.color.rgb = line; r.line.width = Pt(1.25); r.shadow.inherit = False
    L.box(s, x, y + 0.02, w, h, text, size=size, color=fg, bold=bold,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)


def arrow(s, x, y, txt="→"):
    L.box(s, x, y, 0.5, 0.4, txt, size=22, color=BRICK, bold=True, align=PP_ALIGN.CENTER)


def shot(s, path, x, y, maxw, maxh, caption=None):
    """実画面/実成果物のPNGを枠付きで配置（アスペクト比維持・中央寄せ）"""
    from PIL import Image
    iw, ih = Image.open(path).size
    r = iw / ih
    w, h = maxw, maxw / r
    if h > maxh:
        h, w = maxh, maxh * r
    px = x + (maxw - w) / 2
    py = y + (maxh - h) / 2
    fr = s.shapes.add_shape(1, Inches(px - 0.05), Inches(py - 0.05),
                            Inches(w + 0.10), Inches(h + 0.10))
    fr.fill.solid(); fr.fill.fore_color.rgb = WHITE
    fr.line.color.rgb = GRAY; fr.line.width = Pt(1.25); fr.shadow.inherit = False
    s.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(w), Inches(h))
    if caption:
        L.box(s, x, py + h + 0.10, maxw, 0.4, caption, size=12, color=GRAY,
              align=PP_ALIGN.CENTER)


SHOTS = Path(__file__).parent / "out_screener" / "shots"


def build():
    prs = L.Presentation(); prs.slide_width = L.EMU_W; prs.slide_height = L.EMU_H

    # 1 表紙
    s = L.add_slide(prs)
    L.bar(s, 0, 2.5, 13.333, 0.10, BRICK)
    L.box(s, 0.9, 1.3, 11, 0.5, "操作マニュアル ＆ 全体設計", size=14, color=GRAY)
    L.box(s, 0.9, 2.7, 11.5, 1.1, "🔍 KHD AI査定エンジン", size=44, color=BRICK, bold=True)
    L.box(s, 0.95, 3.85, 11.5, 0.7, "物件を入れる → 自動で査定 → 良物件は銀行資料まで全自動", size=19, color=INK)
    L.box(s, 0.95, 6.2, 11.5, 0.6, "菊池ホールディングス（KHD）／ 2026-06-03", size=13, color=GRAY)

    # 2 全体像（1枚）
    s = L.add_slide(prs); L.header(s, 1, "全体像 ― 物件が来てから銀行資料まで")
    y = 2.0
    chip(s, 0.5, y, 2.2, 1.0, "①入口\nPDF/住所/URL\n/メール/Drive", fill=RGBColor(0xEC,0xE0,0xDD))
    arrow(s, 2.75, y + 0.3)
    chip(s, 3.3, y, 2.0, 1.0, "②査定\n路線価自動\n玉川式KPI", fill=WHITE)
    arrow(s, 5.35, y + 0.3)
    chip(s, 5.9, y, 2.0, 1.0, "③判定\n🟢🟡🔴", fill=WHITE)
    arrow(s, 7.95, y + 0.3)
    chip(s, 8.5, y, 2.1, 1.0, "④蓄積\n物DB＋\n自動査定DB", fill=WHITE)
    arrow(s, 10.65, y + 0.3)
    chip(s, 11.2, y, 1.7, 1.0, "⑤通知\nNotion/\nGoogleタスク", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN)
    # 🟢分岐
    L.box(s, 3.3, 3.25, 9, 0.4, "▼ 判定が🟢買いの時だけ自動で続行", size=13, color=GREEN, bold=True)
    chip(s, 5.9, 3.75, 2.6, 0.9, "⑥融資資料デッキ\n自動生成(pptx)", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN)
    arrow(s, 8.55, 4.0)
    chip(s, 9.1, 3.75, 3.2, 0.9, "⑦銀行へ\n（速攻駆け込み）", fill=RGBColor(0xF7,0xE9,0xE7), line=BRICK)
    L.box(s, 0.5, 5.2, 12.3, 1.6,
          "■ 一言で：物件情報を渡せば、路線価まで自動で引いて「買い/見送り」を判定。\n"
          "　 『買い🟢』なら、銀行に出せる融資資料デッキまで勝手に作ってスマホに通知。\n"
          "■ 渡し方は4つ（次ページ）。良し悪しの基準は玉川式（土地値・利回り・CF）。",
          size=14, color=INK, ls=1.3)
    L.footer(s)

    # 3 入口4つ（誰がどこで何を）
    s = L.add_slide(prs); L.header(s, 2, "入口4つ ― どこで・誰が・何を")
    L.table(s, 0.5, 1.55, 12.3, [
        ("📄 マイソクPDF", "菊池がDrive『_査定受け箱』にPDFを入れる → 毎朝7時 自動査定"),
        ("🏠 住所だけ", "菊池が「○○市○○ 2500万 査定して」と言う → その場で査定"),
        ("🔗 ポータルURL", "健美家/楽待のURLを貼る → 詳細を自動取得して本査定"),
        ("📧 業者メール", "GASが業者メールを自動収集→URL抽出 → 毎朝 自動査定（手入力ゼロ）"),
    ], col1=0.24, rh=1.05, fs=14)
    L.box(s, 0.5, 6.2, 12.3, 0.7,
          "★ ③メール経路をONにすると、業者から来る物件は『何もしなくても』自動査定されます。",
          size=14, color=BRICK, bold=True)
    L.footer(s)

    # 4 査定の中身＋🟢定義
    s = L.add_slide(prs); L.header(s, 3, "査定の中身 ＝ 玉川式KPI（🟢の定義）")
    L.box(s, 0.5, 1.5, 12.3, 0.5, "路線価はreinfolib（国交省）から自動取得。下の基準を全部クリアで🟢買い。", size=14, color=INK)
    L.box(s, 0.5, 2.2, 6.0, 0.45, "◆ 収益物件（保有）", size=15, color=BRICK, bold=True)
    L.table(s, 0.5, 2.7, 6.0, [
        ("実質利回り", "8%以上 で🟢"),
        ("CF率", "1.5%以上 で🟢"),
        ("CCR", "15%以上 で🟢"),
        ("土地値割合(担保)", "0.4以上 で🟢"),
    ], col1=0.55, rh=0.6, fs=13)
    L.box(s, 6.8, 2.2, 6.0, 0.45, "◆ 再販（売却）", size=15, color=BRICK, bold=True)
    L.table(s, 6.8, 2.7, 6.0, [
        ("粗利率", "20%以上 で🟢"),
        ("土地値割合(担保)", "0.4以上 で🟢"),
        ("（参考）件名査定", "利回り7%↑=🟢"),
    ], col1=0.55, rh=0.6, fs=13)
    L.box(s, 0.5, 5.7, 12.3, 1.0,
          "■ 総合判定：🔴が1つでも→見送り／🟡や欠損が残る→要検討／全部🟢→『買い』。\n"
          "■ 『買い🟢』になった時だけ、融資資料デッキ生成＋物DB蓄積まで自動で進む。",
          size=14, color=INK, ls=1.3)
    L.footer(s)

    # 5 【実物】査定結果はこう出る（判定xlsx）
    s = L.add_slide(prs); L.header(s, 4, "【実物】査定結果はこう出る（自動生成）")
    L.box(s, 0.5, 1.45, 12.3, 0.5,
          "物件を渡すと、この判定シートが自動で出来る。総合判定🟢/数値は全部自動計算。",
          size=14, color=INK)
    shot(s, SHOTS / "判定_277_高松_盛岡(合算)_2026-06-02.xlsx.png", 0.7, 2.0, 7.6, 4.6,
         caption="▲ 実例：盛岡市高松2丁目（収益）／🟢買い・土地値90.5%・実質利回り11.1%")
    L.box(s, 8.7, 2.2, 4.1, 0.45, "◆ ここが自動で埋まる", size=15, color=BRICK, bold=True)
    L.box(s, 8.7, 2.75, 4.1, 3.6,
          "・総合判定（🟢🟡🔴）\n"
          "・価格／土地面積\n"
          "・土地値割合（路線価から自動）\n"
          "・モード（収益/再販を自動判別）\n"
          "・表面/実質利回り\n"
          "・CF率・CCR・粗利率\n\n"
          "→ 手入力ゼロ。路線価は\n　 reinfolib(国交省)から自動取得。",
          size=13, color=INK, ls=1.35)
    L.footer(s)

    # 6 結果はどこで見る
    s = L.add_slide(prs); L.header(s, 5, "結果はどこで見る？（スマホOK）")
    L.kpi_card(s, 0.5, 1.7, 3.9, "一覧で見る", "Notion", BRICK, "🤖自動査定結果DB")
    L.kpi_card(s, 4.6, 1.7, 3.9, "🟢が届く", "Googleタスク", GREEN, "買い候補がプッシュ")
    L.kpi_card(s, 8.7, 1.7, 4.1, "デッキを取る", "Drive", BLUE, "_査定結果デッキ/")
    L.table(s, 0.5, 3.6, 12.3, [
        ("Notion 🤖自動査定結果DB", "全件（物件/判定/利回り/土地値/価格/入口）。判定で🟢フィルタ"),
        ("Googleタスク", "🟢買い候補がスマホのタスクに自動で立つ（毎朝）"),
        ("Drive『_査定結果デッキ』", "🟢物件の融資資料pptxが自動で入る → 銀行へ"),
        ("物DB「物」タブ", "🟢はSSoT台帳にも自動蓄積"),
    ], col1=0.32, rh=0.62, fs=13)
    L.footer(s)

    # 7 【実物】自動生成された融資資料デッキ
    s = L.add_slide(prs); L.header(s, 6, "【実物】🟢なら銀行資料まで自動生成")
    L.box(s, 0.5, 1.45, 12.3, 0.5,
          "判定が🟢買いの時だけ、銀行へ出せる融資資料デッキ(pptx)が勝手に出来る。",
          size=14, color=INK)
    shot(s, SHOTS / "277_融資資料デッキ.pptx.png", 0.7, 2.0, 7.6, 4.6,
         caption="▲ 実物：岩手銀行 御中／盛岡市高松2丁目 融資ご相談資料（全自動生成）")
    L.box(s, 8.7, 2.2, 4.1, 0.45, "◆ デッキの中身（自動）", size=15, color=BRICK, bold=True)
    L.box(s, 8.7, 2.75, 4.1, 3.6,
          "・表紙（宛先銀行を自動差込）\n"
          "・物件概要／所在地図\n"
          "・収支・KPIサマリ\n"
          "・土地値/積算（融資根拠）\n"
          "・出口戦略\n\n"
          "→ Drive『_査定結果デッキ』に\n　 自動で格納。あとは持参するだけ。",
          size=13, color=INK, ls=1.35)
    L.footer(s)

    # 8 完全無人の構成
    s = L.add_slide(prs); L.header(s, 7, "完全無人の仕組み（裏側）")
    y = 1.9
    chip(s, 0.5, y, 2.6, 1.1, "GAS（Google側）\n業者メール収集\n10分ごと・無人", fill=RGBColor(0xEC,0xE0,0xDD))
    arrow(s, 3.15, y + 0.35)
    chip(s, 3.7, y, 2.4, 1.1, "スプシ\n物件URLが\n溜まる", fill=WHITE)
    arrow(s, 6.15, y + 0.35)
    chip(s, 6.7, y, 2.8, 1.1, "launchd（Mac）\n毎朝7時に起動\nauto_pipeline", fill=WHITE)
    arrow(s, 9.55, y + 0.35)
    chip(s, 10.1, y, 2.7, 1.1, "査定→DB→\nデッキ→通知\n（無人）", fill=RGBColor(0xE7,0xF0,0xE7), line=GREEN)
    L.box(s, 0.5, 3.4, 12.3, 1.0,
          "■ メール経路＝GAS（Googleの中で勝手に動く）＋ Mac毎朝7時のlaunchd。\n"
          "■ Drive経路＝受け箱にPDFを入れれば、同じlaunchdが拾って査定。", size=14, color=INK, ls=1.3)
    L.box(s, 0.5, 4.6, 12.3, 0.45, "■ ワンタイム設定（1回だけ）", size=15, color=BRICK, bold=True)
    L.table(s, 0.5, 5.1, 12.3, [
        ("① launchctl load", "毎朝7時の自動起動（✅済）"),
        ("② Tasks API有効化＋認証", "🟢をGoogleタスク通知（要：GCPで有効化）"),
        ("③ GAS貼付＋トリガー", "業者メールを自動収集（要：script.google.comで貼付）"),
    ], col1=0.32, rh=0.55, fs=13)
    L.footer(s)

    # 9 日々の使い方
    s = L.add_slide(prs); L.header(s, 8, "日々の使い方（菊池の動き）")
    L.box(s, 0.5, 1.6, 12.3, 0.45, "◆ 毎朝（30秒）", size=15, color=BRICK, bold=True)
    L.box(s, 0.7, 2.1, 12.0, 0.9,
          "スマホでNotion『🤖自動査定結果DB』を開く → 判定🟢を見る → 気になる物件のデッキ(Drive)を確認。\n"
          "Googleタスクに🟢通知が来ていれば、そこからジャンプ。", size=14, color=INK, ls=1.3)
    L.box(s, 0.5, 3.2, 12.3, 0.45, "◆ 物件が手に来たら", size=15, color=BRICK, bold=True)
    L.box(s, 0.7, 3.7, 12.0, 0.9,
          "・PDF → Drive『_査定受け箱』へ入れる（or 私に渡す）\n"
          "・URL/住所 → 私に「査定して」と言う（その場で判定＋デッキ）", size=14, color=INK, ls=1.3)
    L.box(s, 0.5, 4.85, 12.3, 0.45, "◆ 良い物件だった時", size=15, color=GREEN, bold=True)
    L.box(s, 0.7, 5.35, 12.0, 1.0,
          "🟢なら融資資料デッキが自動で出来ている → レントロール等を添えて銀行へ持参。\n"
          "『最後まで全自動』＝判定で終わらず、銀行に出せる資料まで勝手に進む。", size=14, color=INK, ls=1.3)
    L.footer(s)

    # 8 結び
    s = L.add_slide(prs)
    L.bar(s, 0, 2.4, 13.333, 0.08, BRICK)
    L.box(s, 0.9, 2.7, 11.5, 1.3,
          "物件を渡すだけ。良い物件は、銀行資料まで勝手に出来る。\nあとは『出す』だけ。",
          size=22, color=INK, ls=1.3)
    L.box(s, 0.95, 5.4, 11.5, 0.7, "KHD AI査定エンジン ／ バイセル構想をTTP→自社で完成・無人化", size=14, color=GRAY)

    out = Path.home() / "01_honbu_docs_automation" / "out_screener" / "KHD_AI査定エンジン_マニュアル.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


if __name__ == "__main__":
    print("✅", build())

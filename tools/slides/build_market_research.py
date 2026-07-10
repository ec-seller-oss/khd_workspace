"""
クリニックLINE導入 市場調査（他社の販売手法／KHDの独自性・弱み）。
クリーム白×レンガ赤。出力: market_research.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"クリニックLINE導入 市場調査 ｜ KHD",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,valc in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(valc); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.09); cell.margin_right=Inches(0.06); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=(TEALBG if is_hi else (CARD if ri%2==1 else BG))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or (ci==0) or is_hi
                    r.font.color.rgb=(WHT if ri==0 else (TEALD if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A))))
    return tb

# ════ 1 表紙 ════
s=sl()
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"MARKET RESEARCH",Inches(0.9),Inches(1.6),Inches(11),Inches(0.4),sz=15,bold=True,col=RED)
t(s,"クリニックLINE導入 市場調査",Inches(0.88),Inches(2.15),Inches(11.6),Inches(0.9),sz=34,bold=True,col=INK)
t(s,"他社はどう売り、我々は何で勝つか",Inches(0.9),Inches(3.05),Inches(11.6),Inches(0.6),sz=22,bold=True,col=RED)
t(s,"LINE導入を売る各社の販売手法・価格を整理し、KHDだけができること／逆にできないこと（弱み）まで網羅。",
  Inches(0.92),Inches(3.95),Inches(11.4),Inches(0.6),sz=13.5,col=GRY,line_sp=1.25)
offers=[("他社の初期","20〜250万円"),("他社の月額","5〜40万円"),("KHD","0円＋成果報酬")]
ox,ow,og=Inches(0.9),Inches(3.7),Inches(0.2)
for i,(lab,val) in enumerate(offers):
    cx=ox+(ow+og)*i; isk=(i==2)
    bx(s,cx,Inches(4.85),ow,Inches(1.35),REDBG if isk else CARD,line=CARDLN,lw=1.0); bx(s,cx,Inches(4.85),ow,Inches(0.06),RED)
    t(s,lab,cx,Inches(5.02),ow,Inches(0.35),sz=12,col=GRY,align=PP_ALIGN.CENTER)
    t(s,val,cx,Inches(5.4),ow,Inches(0.6),sz=22,bold=True,col=RED,align=PP_ALIGN.CENTER)
bx(s,Inches(0.9),Inches(6.55),Inches(11.5),Pt(1.2),LN)
t(s,"出典：Lステップ構築代行各社・クリニック向けLINE支援各社の公開価格（2025-2026）",Inches(0.9),Inches(6.65),Inches(11),Inches(0.4),sz=10,col=GRY)

# ════ 2 市場のプレイヤー4類型 ════
s=sl(); ft(s)
hdr(s,"WHO SELLS WHAT","市場のプレイヤーは大きく4類型","誰が・何を・いくらで売っているか")
rows=[
 ("類型","代表例","売り方","価格帯（目安）"),
 ("Lステップ構築代行","正規代理店／フリーランス","シナリオ・リッチメニュー構築＋月額運用代行","初期20〜250万＋月5〜40万"),
 ("クリニック特化 LINE/予約","カルー／メディカル革命byGMO／CLIUS／ヒーロー","HP制作・診療予約システムにLINEを付帯","初期8.8万〜＋月3,300円〜"),
 ("MA・拡張ツール本体","Lステップ／エルメ／Liny／MicoCloud","ツール月額課金＋販売パートナー（代理店）","月額課金（数千〜数万）"),
 ("広告・運用代理店","LINE広告代理店","広告→友だち追加→ステップ配信で育成","月10〜30万＋広告費"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(4.3),[Inches(2.6),Inches(3.5),Inches(3.6),Inches(2.53)],hi_col=None,sz=11,header_sz=12)
t(s,"→ 共通点：「ツールを売る」「高い初期費用＋月額固定を取る」モデル。効果が出なくても費用は発生。",Inches(0.55),Inches(6.4),Inches(12.2),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 3 他社の販売手法 ════
s=sl(); ft(s)
hdr(s,"HOW THEY SELL","他社の「売り方」5パターン","ここを知ると、我々の刺し方が決まる")
ways=[("高額初期＋月額固定","初期20〜250万・月5〜40万。効果に関わらず固定収益（院はリスクを負う）"),
      ("無料セミナー集客","「LINEで集客」セミナー→構築・ツール契約へ誘導（情報→高額提案）"),
      ("代理店・パートナー制度","Liny等がツールの販売代理店を募集→販路を面で広げる"),
      ("フリーランス副業市場","「LINE構築代行で稼ぐ」副業層が格安〜中価格で構築を量産"),
      ("自社もLINEで実演","広告→友だち追加→ステップ配信で「自分たちが結果を出して」見せる")]
cw,ch,gx,gy=Inches(6.0),Inches(1.42),Inches(0.45),Inches(0.3); x0,y0=Inches(0.55),Inches(2.0)
for i,(ti,ds) in enumerate(ways[:4]):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.28),cy+Inches(0.12),cw-Inches(0.5),Inches(0.4),sz=14,bold=True,col=INK)
    t(s,ds,cx+Inches(0.3),cy+Inches(0.55),cw-Inches(0.55),Inches(0.8),sz=11,col=GRY,line_sp=1.1)
cy=y0+(ch+gy)*2
bx(s,Inches(0.55),cy,Inches(12.45),ch,GRYBG,line=CARDLN,lw=1.0); bx(s,Inches(0.55),cy,Inches(0.1),ch,RED)
t(s,ways[4][0],Inches(0.83),cy+Inches(0.12),Inches(4),Inches(0.4),sz=14,bold=True,col=INK)
t(s,ways[4][1],Inches(0.85),cy+Inches(0.55),Inches(12),Inches(0.6),sz=11,col=GRY)
t(s,"※ ほとんどが「ツール販売」と「固定費」。「中立で・成果報酬で・自分で作る」プレイヤーは少ない＝そこが空白。",Inches(0.55),Inches(6.65),Inches(12.4),Inches(0.4),sz=11.5,bold=True,col=REDD)

# ════ 4 我々しかできない ════
s=sl(); ft(s)
hdr(s,"ONLY KHD","我々しかできない（独自性）","他社の前提＝高額ツール販売を、ぜんぶ外している")
adv=[("初期0円・月額0円","OSS（LINE Harness）×Cloudflare無料枠。他社の初期20-250万＋月額固定がゼロ"),
     ("AIで全構築・全運用","Claude Codeで構築〜運用。京橋は実質1日で本番化。人を増やさず回せる"),
     ("中立（ツールを売らない）","特定ツールの代理店でない→院に最適な組合せを選べる"),
     ("院インサイダー視点","福井＝事務長／医療コンサル。経営課題の翻訳・現場スタッフの感情に配慮"),
     ("成果報酬が可能","固定費0だから「効果が出た分だけ」が成立。院のリスクを消せる"),
     ("横展開パイプライン","1院で型化→次院は固有値の差し替えだけ。量産できる")]
cw,ch,gx,gy=Inches(6.0),Inches(1.32),Inches(0.45),Inches(0.22); x0,y0=Inches(0.55),Inches(1.95)
for i,(ti,ds) in enumerate(adv):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,TEALBG,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,TEALD)
    t(s,ti,cx+Inches(0.28),cy+Inches(0.12),cw-Inches(0.5),Inches(0.4),sz=13.5,bold=True,col=TEALD)
    t(s,ds,cx+Inches(0.3),cy+Inches(0.55),cw-Inches(0.55),Inches(0.72),sz=10.5,col=INK,line_sp=1.08)

# ════ 5 我々にはできない（弱み）＋補い方 ════
s=sl(); ft(s)
hdr(s,"NOT KHD (yet)","逆に、我々にはできない（弱み）と その補い方","正直に弱みを置き、潰し方をセットで持つ")
rows=[
 ("我々にはできない／弱い","補い方（戦い方）"),
 ("大量の導入実績・ブランド信頼","小さく始めて実績を積む（京橋が実証1号）。事例で語る"),
 ("24時間サポート・専任CS体制","少数精鋭の伴走で対応。規模拡大時に体制を作る"),
 ("レセコン・電子カルテの正式API連携","メーカー（中央ビジコム等）と住み分け＝「器」は任せ「中身」を担う"),
 ("大規模な営業網・代理店網","福井ルート・紹介で「質」重視。数で勝負しない"),
 ("法務・医療広告ガイドラインの専門体制","専門家と連携・公開GL順守。誇大表現を避ける"),
 ("OSSの保守責任・障害復旧","自己保守＋影響範囲を限定（まず1機能から）。重要連携は冗長化"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(4.5),[Inches(5.4),Inches(6.83)],hi_col=None,sz=11.5,header_sz=12.5)
t(s,"→ 弱みは「実績・体制・正式連携」。「小さく実績／メーカーと住み分け／成果で信頼」で埋める。",Inches(0.55),Inches(6.6),Inches(12.2),Inches(0.4),sz=12,bold=True,col=REDD)

# ════ 6 ポジショニング／勝ち筋 ════
s=sl(); ft(s)
hdr(s,"POSITIONING","勝ち筋 ── 空白は「中立 × 低コスト × 中身を担う」","他社が「ツールを高く売る」なら、我々は「中身と経営を担う中立の伴走者」")
# 2x2 マップ
mx,my,mw,mh=Inches(0.8),Inches(2.1),Inches(7.6),Inches(4.3)
bx(s,mx,my,mw,mh,WHT,line=CARDLN,lw=1.0)
bx(s,mx,my+mh/2,mw,Pt(1.2),LN); bx(s,mx+mw/2,my,Pt(1.2),mh,LN)
t(s,"低コスト・低リスク（院）",mx,my-Inches(0.05),mw,Inches(0.3),sz=10.5,bold=True,col=GRY,align=PP_ALIGN.CENTER)
t(s,"高コスト・固定費",mx,my+mh-Inches(0.02),mw,Inches(0.3),sz=10.5,bold=True,col=GRY,align=PP_ALIGN.CENTER)
t(s,"ツールを売る",mx-Inches(0.1),my+mh/2-Inches(0.18),Inches(1.6),Inches(0.3),sz=10.5,bold=True,col=GRY)
t(s,"中身・経営を担う",mx+mw-Inches(1.7),my+mh/2-Inches(0.18),Inches(1.7),Inches(0.3),sz=10.5,bold=True,col=GRY,align=PP_ALIGN.RIGHT)
# KHD（右上＝中身担う×低コスト）
bx(s,mx+mw*0.55,my+Inches(0.5),Inches(2.6),Inches(1.0),TEALBG,line=TEALD,lw=1.5)
t(s,"KHD",mx+mw*0.55,my+Inches(0.6),Inches(2.6),Inches(0.4),sz=15,bold=True,col=TEALD,align=PP_ALIGN.CENTER)
t(s,"中立・0円・成果報酬\n×AI×経営伴走",mx+mw*0.55,my+Inches(1.0),Inches(2.6),Inches(0.5),sz=10,col=TEALD,align=PP_ALIGN.CENTER,line_sp=1.0)
# 他社（左下＝ツール販売×高コスト）
bx(s,mx+Inches(0.4),my+mh-Inches(1.5),Inches(2.7),Inches(1.0),GRYBG,line=CARDLN,lw=1.0)
t(s,"既存の構築代行・ツール各社",mx+Inches(0.4),my+mh-Inches(1.4),Inches(2.7),Inches(0.5),sz=11,bold=True,col=GRY,align=PP_ALIGN.CENTER,line_sp=1.0)
t(s,"高額初期＋月額固定",mx+Inches(0.4),my+mh-Inches(0.85),Inches(2.7),Inches(0.4),sz=10,col=GRY,align=PP_ALIGN.CENTER)
# 右：勝ち筋
rx=Inches(8.7)
t(s,"勝ち筋（3つ）",rx,Inches(2.1),Inches(4.2),Inches(0.4),sz=15,bold=True,col=RED)
for i,(ti,ds) in enumerate([("リスクを消す","固定費0×成果報酬で、院の「払って効果ゼロ」を排除"),("中立で最適化","ツールを売らないから、院に合う組合せを選べる"),("経営に接続","LINEで終わらず、承継・テナント等の本丸へつなぐ")]):
    cy=Inches(2.7)+Inches(1.25)*i
    bx(s,rx,cy,Inches(4.2),Inches(1.1),CARD,line=CARDLN,lw=1.0); bx(s,rx,cy,Inches(0.08),Inches(1.1),RED)
    t(s,ti,rx+Inches(0.25),cy+Inches(0.1),Inches(3.8),Inches(0.4),sz=13,bold=True,col=REDD)
    t(s,ds,rx+Inches(0.27),cy+Inches(0.5),Inches(3.8),Inches(0.55),sz=10.5,col=GRY,line_sp=1.05)

prs.save("market_research.pptx")
print("saved market_research.pptx slides:",len(prs.slides._sldIdLst))

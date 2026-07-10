"""
YouTube 医療チャンネル EP01「クリニック開業、9割が立地で失敗する」 画面サポート用スライド
build_slides_minimal.py のテーマ（ダーク×ゴールド）を流用。
李牧の指摘反映: 診療圏に具体例 / AI活用の一言 / ②を厚く。
出力: medical_ep01_2026.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG  = RGBColor(0x0D, 0x1B, 0x2A)
ACC = RGBColor(0xE8, 0xA8, 0x00)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
LGR = RGBColor(0xCC, 0xD6, 0xE0)
MID = RGBColor(0x1A, 0x2E, 0x44)
BLU = RGBColor(0x4A, 0x9E, 0xCB)
GRN = RGBColor(0x4C, 0xC2, 0x8C)
RED = RGBColor(0xFF, 0x77, 0x77)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=20, bold=False, col=WHT,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
        r.font.italic = italic
    return tb


def bx(slide, x, y, w, h, col, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = col
    if lw:
        s.line.color.rgb = col
    else:
        s.line.fill.background()
    return s


def hdr(slide, main, sub=""):
    bx(slide, 0, 0, W, Inches(1.1), MID)
    t(slide, main, Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), sz=24, bold=True)
    if sub:
        t(slide, sub, Inches(0.6), Inches(0.78), Inches(12), Inches(0.3), sz=11, col=LGR)


def ft(slide):
    bx(slide, 0, H-Inches(0.38), W, Inches(0.38), MID)
    t(slide, "クリニック開業支援  |  KHD × 医療メディア",
      Inches(0.5), H-Inches(0.35), Inches(10), Inches(0.32), sz=9, col=LGR)


def ai_chip(slide, text, x, y, w):
    bx(slide, x, y, w, Inches(0.7), RGBColor(0x10, 0x24, 0x38))
    bx(slide, x, y, Inches(0.08), Inches(0.7), GRN)
    t(slide, text, x+Inches(0.22), y+Inches(0.12), w-Inches(0.35), Inches(0.5),
      sz=13, bold=True, col=GRN)


# ── 01 カバー ──────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "クリニック開業・承継 専門メディア",
  Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.5), sz=15, col=ACC, italic=True)
t(s, "クリニック開業、\n9割が立地で失敗する",
  Inches(0.9), Inches(1.45), Inches(11.8), Inches(2.2), sz=46, bold=True)
t(s, "診療圏とテナントの見極め方 ── 開業医が最初に見るべき「立地の数字」3つ",
  Inches(0.9), Inches(3.85), Inches(11.5), Inches(0.55), sz=16, col=LGR)
bx(s, Inches(0.9), Inches(4.6), Inches(10), Inches(0.04), ACC)
t(s, "菊池ホールディングス（KHD）  |  @khd_medical01",
  Inches(0.9), Inches(4.75), Inches(10), Inches(0.45), sz=14, col=LGR)

# ── 02 結論 ────────────────────────────────
s = sl(); ft(s)
hdr(s, "結論")
t(s, "開業は、立地で9割が決まる。",
  Inches(0.7), Inches(1.7), Inches(12), Inches(1.0), sz=40, bold=True, col=ACC)
t(s, "内装でも、設備でも、先生の腕でもなく ── まず立地。",
  Inches(0.7), Inches(3.0), Inches(12), Inches(0.6), sz=20, col=WHT)
bx(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(1.5), MID)
bx(s, Inches(0.7), Inches(4.0), Inches(0.08), Inches(1.5), ACC)
t(s, "探すのは「良い場所」ではない。\n「自分の科目に合う場所」を見抜く。",
  Inches(0.95), Inches(4.25), Inches(11.6), Inches(1.1), sz=22, bold=True, italic=True)

# ── 03 3つの視点 ───────────────────────────
s = sl(); ft(s)
hdr(s, "立地を見抜く 3つの視点", "感覚で決めない。数字と事実で見る。")
items = [("01", "診療圏", "推計患者数で\n損益分岐を超えるか"),
         ("02", "科目との相性", "良い物件は存在しない。\n科目で正解が変わる"),
         ("03", "テナント", "坪単価で選ばない。\n視認性・動線・競合の隙間")]
cw = Inches(4.0)
for i, (no, ti, di) in enumerate(items):
    cx = Inches(0.55)+i*(cw+Inches(0.14))
    bx(s, cx, Inches(1.55), cw, Inches(4.3), MID)
    bx(s, cx, Inches(1.55), cw, Inches(0.06), ACC)
    t(s, no, cx+Inches(0.25), Inches(1.8), cw-Inches(0.4), Inches(1.0), sz=44, bold=True, col=ACC)
    t(s, ti, cx+Inches(0.25), Inches(2.95), cw-Inches(0.4), Inches(0.7), sz=20, bold=True)
    t(s, di, cx+Inches(0.25), Inches(3.8), cw-Inches(0.4), Inches(1.8), sz=14, col=LGR)

# ── 04 視点① 診療圏（具体例＋AI） ──────────
s = sl(); ft(s)
hdr(s, "視点 1  /  診療圏 = 推計患者数", "最初に出すべき、たった1つの数字")
bx(s, Inches(0.7), Inches(1.45), Inches(12.0), Inches(0.95), MID)
t(s, "半径◯km以内の人口  ×  科目の受療率  ÷  競合の数  =  推計患者数",
  Inches(0.95), Inches(1.7), Inches(11.6), Inches(0.5), sz=20, bold=True, col=ACC)
t(s, "【具体例】内科 / 半径2km・人口5万人 × 受療率6% ÷ 競合4院 ≒ 1日あたりの推計患者数",
  Inches(0.7), Inches(2.7), Inches(12), Inches(0.5), sz=15, col=WHT)
t(s, "→ この推計が開業後の損益分岐を超えるか。これが最初の関門。",
  Inches(0.7), Inches(3.3), Inches(12), Inches(0.5), sz=15, col=LGR)
t(s, "内装やこだわりは、この数字をクリアしてからの話。\n順番を逆にすると「立派な箱に患者が来ない」最悪が起きる。",
  Inches(0.7), Inches(4.0), Inches(12), Inches(1.0), sz=16, col=WHT)
ai_chip(s, "💡 この診療圏分析、今はAIで数分。感覚に頼らず、誰でも数字で判断できる時代。",
        Inches(0.7), Inches(5.3), Inches(12.0))

# ── 05 視点② 科目と立地の相性 ──────────────
s = sl(); ft(s)
hdr(s, "視点 2  /  科目と立地の相性", "同じ立地でも、科目で正解が変わる")
bx(s, Inches(0.6), Inches(1.5), Inches(5.95), Inches(3.6), MID)
bx(s, Inches(0.6), Inches(1.5), Inches(5.95), Inches(0.06), BLU)
t(s, "内科・小児科", Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.6), sz=22, bold=True, col=BLU)
t(s, "住宅地の「生活動線」", Inches(0.85), Inches(2.5), Inches(5.5), Inches(0.5), sz=18, bold=True)
t(s, "毎日の買い物や通勤で\n前を通る場所が強い。",
  Inches(0.85), Inches(3.2), Inches(5.5), Inches(1.5), sz=15, col=LGR)
bx(s, Inches(6.78), Inches(1.5), Inches(5.95), Inches(3.6), MID)
bx(s, Inches(6.78), Inches(1.5), Inches(5.95), Inches(0.06), ACC)
t(s, "美容・自由診療", Inches(7.03), Inches(1.7), Inches(5.5), Inches(0.6), sz=22, bold=True, col=ACC)
t(s, "ターミナル駅の「視認性」", Inches(7.03), Inches(2.5), Inches(5.5), Inches(0.5), sz=18, bold=True)
t(s, "わざわざ来てもらう前提。\n目立つ場所が効く。",
  Inches(7.03), Inches(3.2), Inches(5.5), Inches(1.5), sz=15, col=LGR)
t(s, "「良い物件」は存在しない。あるのは「その科目に合う物件」だけ。",
  Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.6), sz=18, bold=True, col=ACC, align=PP_ALIGN.CENTER)

# ── 06 視点③ テナント ──────────────────────
s = sl(); ft(s)
hdr(s, "視点 3  /  テナントは坪単価で選ばない")
t(s, "効くのは ── 視認性・動線・競合との距離。",
  Inches(0.7), Inches(1.55), Inches(12), Inches(0.6), sz=22, bold=True, col=ACC)
checks = ["1階で人目につくか", "駅やスーパーの導線上にあるか", "近くの競合は土日をやっているか"]
for i, x in enumerate(checks):
    bx(s, Inches(0.7), Inches(2.5)+i*Inches(0.7), Inches(0.14), Inches(0.5), ACC)
    t(s, x, Inches(1.0), Inches(2.5)+i*Inches(0.7), Inches(11.5), Inches(0.5), sz=17)
bx(s, Inches(0.7), Inches(4.9), Inches(12.0), Inches(1.4), MID)
bx(s, Inches(0.7), Inches(4.9), Inches(0.08), Inches(1.4), RED)
t(s, "家賃が安くても患者が来ない立地は、いちばん高くつく。\n競合が手を抜く時間帯・曜日 ── そこがあなたの患者になる。",
  Inches(0.95), Inches(5.15), Inches(11.6), Inches(1.0), sz=17, bold=True)

# ── 07 まとめ + CTA ────────────────────────
s = sl(); ft(s)
hdr(s, "まとめ", "立地は、この3つを「数字と事実」で見る")
items = ["①  診療圏の推計患者数", "②  科目との相性", "③  テナントの視認性と競合の隙間"]
for i, x in enumerate(items):
    t(s, x, Inches(0.8), Inches(1.55)+i*Inches(0.75), Inches(11.5), Inches(0.6), sz=22, bold=True)
t(s, "感覚で決めない。これだけで失敗確率はぐっと下がる。",
  Inches(0.8), Inches(4.0), Inches(12), Inches(0.6), sz=17, col=LGR)
bx(s, Inches(0.7), Inches(4.8), Inches(12.0), Inches(1.5), MID)
bx(s, Inches(0.7), Inches(4.8), Inches(0.08), Inches(1.5), ACC)
t(s, "診療圏調査から事業計画、テナント選びまで一緒に組み立てています。\n「うちの科目だとどの立地が合う？」 ── 概要欄リンクから気軽にご相談を。",
  Inches(0.95), Inches(5.0), Inches(11.6), Inches(1.1), sz=16, col=WHT)

prs.save("medical_ep01_2026.pptx")
print("✅ medical_ep01_2026.pptx を生成しました（全7枚）")

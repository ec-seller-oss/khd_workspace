# -*- coding: utf-8 -*-
"""
オンライン診療 3社比較 ＋ 月曜デモ手順（京橋・内科／初診から想定）1枚資料
KHD配色 クリーム白#F9F6EF × レンガ赤#AA2E26
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); WHT=RGBColor(0xFF,0xFF,0xFF)
GRN=RGBColor(0x2E,0x7D,0x32); REDBG=RGBColor(0xF4,0xE4,0xE2)
FONT="Hiragino Sans"
W=Inches(13.333); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
s=prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb=BG

def bx(x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.shadow.inherit=False
    if col is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=col
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp
def t(x,y,w,h,text,sz=11,b=False,c=INK,al=PP_ALIGN.LEFT,an=MSO_ANCHOR.TOP,ls=None):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=an
    tf.margin_left=Emu(0); tf.margin_right=Emu(0); tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=al
        if ls: p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.name=FONT; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
    return tb

# ヘッダー
bx(Inches(0.5),Inches(0.4),Pt(4),Inches(1.0),RED)
t(Inches(0.9),Inches(0.4),Inches(11.5),Inches(0.32),"ONLINE CONSULTATION ｜ 京橋クリニック（内科・初診からを想定）",11,True,RED)
t(Inches(0.88),Inches(0.72),Inches(11.6),Inches(0.5),"オンライン診療 ── 3社比較と、月曜デモの手順",22,True,INK)
bx(Inches(0.9),Inches(1.32),Inches(11.5),Pt(3),RED)

# ===== 左：比較表（native table）=====
rows=[
 ["","CLINICS（メドレー）","curon（MICIN）","YaDoc"],
 ["iPad","◎","◎ Web版可","◎"],
 ["費用","月1.5万＋決済3.45%","初期0・月0／決済4%","月3.3万"],
 ["Medicom","別立てになりがち","公式が連携を案内","公式連携(カルテにボタン)"],
 ["将来のAI拡張","◎ 一気通貫","○","△"],
 ["補助金","認定実績あり","要確認","要確認"],
]
tx,ty,tw,th=Inches(0.6),Inches(1.55),Inches(7.5),Inches(3.0)
tbl=s.shapes.add_table(len(rows),4,tx,ty,tw,th).table
tbl.first_row=False; tbl.horz_banding=False
tbl.columns[0].width=Inches(1.5); tbl.columns[1].width=Inches(2.2); tbl.columns[2].width=Inches(2.0); tbl.columns[3].width=Inches(1.8)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        c=tbl.cell(ri,ci); c.text=val; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        c.margin_left=Inches(0.06); c.margin_right=Inches(0.04); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
        c.fill.solid()
        hi = (ci==2)  # curon列を強調
        if ri==0: c.fill.fore_color.rgb = REDD if hi else RED
        else: c.fill.fore_color.rgb = REDBG if hi else (CARD if ri%2 else BG)
        for p in c.text_frame.paragraphs:
            p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name=FONT; r.font.size=Pt(10.5 if ri==0 else 10)
                r.font.bold=(ri==0) or (ci==0) or hi
                r.font.color.rgb = WHT if ri==0 else (REDD if hi else INK)

# ===== 右上：初診からの注意 =====
RX,RW=Inches(8.35),Inches(4.45)
bx(RX,Inches(1.55),RW,Inches(1.7),CARD,line=CARDLN,lw=1.0)
bx(RX,Inches(1.55),RW,Inches(0.06),RED)
t(RX+Inches(0.2),Inches(1.66),RW-Inches(0.4),Inches(0.3),"「初診から」やる場合の前提",12,True,RED)
t(RX+Inches(0.22),Inches(2.0),RW-Inches(0.45),Inches(1.2),
  "・2026年4月施行の改訂指針に準拠\n・原則かかりつけ医。初診も、既往/服薬/アレルギー等を把握でき医師が可と判断すれば可\n・初診は原則 顔写真付き本人確認（マイナ/免許等）",
  9.8,False,INK,ls=1.18)

# ===== 右下：二択の薦め方 =====
bx(RX,Inches(3.4),RW,Inches(2.45),WHT,line=CARDLN,lw=1.0)
bx(RX,Inches(3.4),RW,Inches(0.06),RED)
t(RX+Inches(0.2),Inches(3.5),RW-Inches(0.4),Inches(0.3),"先生への薦め方（正直な二択）",12,True,RED)
t(RX+Inches(0.22),Inches(3.85),RW-Inches(0.45),Inches(0.95),
  "A. 将来AIまで一気通貫を取る → CLINICS\n   ただし電子カルテをMedicomから移す覚悟が要る",9.8,True,INK,ls=1.15)
t(RX+Inches(0.22),Inches(4.78),RW-Inches(0.45),Inches(1.0),
  "B. Medicomを生かし低リスクで即始める → curon\n   月0円・公式連携。まず試すのに最適\n   ※連携でもカルテ自動転記までは要確認(二重入力)",9.8,True,REDD,ls=1.15)

# ===== 下部：月曜デモ手順（curon・iPad）=====
by=Inches(5.95)
bx(Inches(0.6),by,Inches(12.2),Inches(1.05),REDBG)
bx(Inches(0.6),by,Inches(0.1),Inches(1.05),RED)
t(Inches(0.85),by+Inches(0.1),Inches(12.0),Inches(0.3),"月曜デモ手順（curon を iPad で・初期0円で試せる）",12,True,REDD)
t(Inches(0.85),by+Inches(0.45),Inches(12.0),Inches(0.55),
  "① 患者目線：予約 → Web問診 → ビデオ診療 → クレカ決済 の流れをiPadで実演　② 医師目線：予約一覧→ビデオ診察→処方/薬局配送（医師側iPad完結は要実機確認）　③ 締め：初診運用（本人確認・かかりつけ）の設計を一緒に詰める",
  9.8,False,INK,ls=1.15)

t(Inches(0.6),Inches(7.16),Inches(11),Inches(0.3),"テナントアシスト・ウイン株式会社 ｜ 菊池 研太　※数値・連携は商談前に各社へ要確認",8.5,False,GRY)

import shutil
prs.save("online_shinryo_compare.pptx")
_dst="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/京橋_オンライン診療比較と月曜デモ.pptx"
shutil.copy("online_shinryo_compare.pptx",_dst)
print("SAVED:",_dst)

"""
KHD 2デッキ生成（クリーム白×レンガ赤）
 ① KHD_オペレーティングモデル＆作業分解マニュアル（人=打席 / Claude=それ以外）
 ② 営業の型＝羽鳥×無敗営業ノウハウ（キラー3問→握る→決め手／90日道場）
SSoT: 2026-06-24 セッションの結論。出力: ops_model.pptx / eigyo_kata.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)

def newprs():
    p=Presentation(); p.slide_width=W; p.slide_height=H; return p

def sl(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=col; r.font.name=FONT
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)

def ftmaker(label):
    def ft(slide):
        bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
        t(slide,label,Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
    return ft

def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.08)
            cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04); cell.fill.solid()
            is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=REDBG if is_hi else (CARD if ri%2==1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or is_hi or (ci==0)
                    if ri==0: r.font.color.rgb=WHT
                    elif is_hi: r.font.color.rgb=RED
                    elif ci==0: r.font.color.rgb=INK
                    else: r.font.color.rgb=RGBColor(0x3A,0x3A,0x3A)
    return tb

def cover(s,eng,l1,l2red,sub,byline):
    bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
    t(s,eng,Inches(0.9),Inches(1.5),Inches(11.5),Inches(0.45),sz=15,bold=True,col=RED)
    t(s,l1,Inches(0.88),Inches(2.1),Inches(11.6),Inches(0.95),sz=38,bold=True,col=INK)
    t(s,l2red,Inches(0.88),Inches(2.95),Inches(11.6),Inches(0.95),sz=38,bold=True,col=RED)
    t(s,sub,Inches(0.9),Inches(4.1),Inches(11.4),Inches(1.0),sz=14,col=GRY,line_sp=1.3)
    bx(s,Inches(0.9),Inches(6.6),Inches(11.5),Pt(1.2),LINE)
    t(s,byline,Inches(0.9),Inches(6.72),Inches(11),Inches(0.4),sz=13,bold=True,col=INK)

def cards2x2(s,items,y0=Inches(1.85),ch=Inches(2.4),gy=Inches(0.4),bodysz=12):
    cw,gx=Inches(6.0),Inches(0.45); x0=Inches(0.55)
    by=Emu(int(ch*0.40)); bh=Emu(int(ch*0.55))
    for i,(no,ti,body) in enumerate(items):
        cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
        bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.12),ch,RED)
        t(s,no,cx+Inches(0.3),cy+Inches(0.18),Inches(1.1),Inches(0.55),sz=24,bold=True,col=RED)
        t(s,ti,cx+Inches(1.15),cy+Inches(0.24),cw-Inches(1.35),Inches(0.6),sz=15,bold=True,col=INK)
        t(s,body,cx+Inches(0.34),cy+by,cw-Inches(0.6),bh,sz=bodysz,col=GRY,line_sp=1.12)

# ════════════════════════════════════════════════════════════
#  DECK ① オペレーティングモデル＆作業分解マニュアル
# ════════════════════════════════════════════════════════════
prs=newprs(); ft=ftmaker("KHD オペレーティングモデル  ｜  2026-06")

# S1 表紙
s=sl(prs)
cover(s,"KHD OPERATING MODEL",
      "人は「打席」だけ。",
      "残りは全部、AIが持つ。",
      "経営の最小単位＝会って・聞いて・詰めて・信頼を作る。\nそれ以外（記録/集計/資料/分析/可視化）はClaudeが自動化する。",
      "菊池 研太  ｜  KHD  ｜  作業分解マニュアル")

# S2 大原則
s=sl(prs); ft(s)
hdr(s,"THE PRINCIPLE","経営の最小単位 ── 人にしかできない核は1つだけ","ここに時間を集約する。空いた時間は調査士と家族へ")
bx(s,Inches(0.8),Inches(1.95),Inches(11.7),Inches(1.15),REDBG); bx(s,Inches(0.8),Inches(1.95),Inches(0.12),Inches(1.15),RED)
t(s,"人にしかできない核 ＝「会って・聞いて・詰めて・信頼を作る」（＝相手のYesを作る打席）",Inches(1.1),Inches(2.12),Inches(11.1),Inches(0.5),sz=16,bold=True,col=REDD)
t(s,"関係・判断・存在＝自動化できない。だから人はここに全集中。それ以外は全部、機械でいい。",Inches(1.1),Inches(2.62),Inches(11.1),Inches(0.4),sz=12.5,col=INK)
cards2x2(s,[
 ("◯","人＝打席・判断・信頼","会う/電話/ヒアリング/詰める/クロージング・GO/見送り・指値・採用・GIVE"),
 ("AI","Claude＝それ以外全部","記録(02転記)/集計(週次KPI)/資料作成/分析(決め手)/リマインド/可視化"),
 ("⚠","当たり前度を高める","営業は才能でなく型×量。反復でキラー3問が「呼吸」になれば打席の負荷↓"),
 ("⛔","逃げの罠","「自動化の作り込み」自体が新しい謎の業務。空いた時間は打席と調査士へ"),
], y0=Inches(3.3), ch=Inches(1.7), gy=Inches(0.26), bodysz=11.5)

# S3 役割分担表
s=sl(prs); ft(s)
hdr(s,"DIVISION OF LABOR","作業分解 ── 人がやる / Claudeが回す","迷ったらこの表。「誰のYesに繋がるか」言えない作業は、やめる")
rows=[
 ("領域","人（菊池）だけがやる","Claudeが回す（自動/代行）"),
 ("営業","会う・聞く・詰める・クロージング","文面ドラフト・送信前3秒チェック・追客リマインド"),
 ("意思決定","GO/見送り・指値・採用・撤退","数字の集計・シナリオ比較・資金繰り可視化"),
 ("記録","結果の「決め手」を一言","対話を02へ全転記・案件メタ補完"),
 ("分析","次の打ち手を決める","勝因/敗因の集計・成約率・ボトルネック抽出"),
 ("段取り","相手との約束を握る","スケジュール整理・朝晩ブリーフ・週次棚卸し"),
 ("人的資本","調査士・家族（死守）","学習リマインド・進捗可視化"),
]
light_table(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(4.6),
            [Inches(2.0),Inches(5.3),Inches(4.93)],hi_col=1,sz=12.5,header_sz=13)

# S4 1日の流れ
s=sl(prs); ft(s)
hdr(s,"DAILY FLOW","1日の流れ ── 朝・日中・夜で何をするか","朝の5-7時は調査士の聖域(不可侵)。人は打席、Claudeは裏で回す")
rows=[
 ("時間帯","人がやること","Claudeがやること"),
 ("朝","調査士(5-7聖域)→当日の打席を1つ決める","カレンダー→02へ当日予定登録／「今日の打席は誰?」提示"),
 ("日中","打席に立つ(会う・聞く・詰める)","対話を全部02へ報告転記／文面ドラフト／資料作成"),
 ("夜","結果と「決め手」を一言／明日の方針","実績化(実時間)・KPI更新・明日の繰越を準備"),
 ("週次(月)","先週レビュー→今週の打席を選ぶ","05_週次KPIを自動再構築(本部別h・成約率・繰越)"),
]
light_table(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(3.9),
            [Inches(2.0),Inches(5.3),Inches(4.93)],hi_col=1,sz=12.5,header_sz=13)
bx(s,Inches(0.8),Inches(6.1),Inches(11.7),Inches(0.7),REDBG); bx(s,Inches(0.8),Inches(6.1),Inches(0.1),Inches(0.7),RED)
t(s,"朝の第一声は「今日の打席は誰？」── 作業でなく「会う相手」から1日を組む。",Inches(1.05),Inches(6.22),Inches(11.4),Inches(0.5),sz=14,bold=True,col=REDD)

# S5 作業マニュアル（手順）
s=sl(prs); ft(s)
hdr(s,"MANUAL","各作業マニュアル ── 迷わず回す手順","人は依頼するだけ。実作業はClaudeがこの手順で巻き取る")
mans=[
 ("①朝ブリーフ＆02登録","Claude：日付確認→カレンダー全網羅取得→02へ当日「予定」登録(本部/案件番号/営業直結)→『カネ直結2-3件＋今日の打席は誰か』を提示。人：打席を1つ決める。"),
 ("②日中＝打席＋報告転記","人：会って聞いて詰める。Claude：対話は全業務報告→02へ行追加/更新。結果・次アクション(誰いつ何)・決定事項を1行で残す。"),
 ("③夜の終業＆実績化","Claude：先に記録を読む→当日を実績化(カレンダーの実時間→実開始/終了→所要が自動算出)→結果に「決め手」記録→明日の下準備を提示。"),
 ("④週次棚卸し(月曜)","Claude：runAllで05_週次KPI再構築。先週レビュー(本部別h/営業比率/成約率)＋今週やること=繰越を自動集約。人：優先S/A/Bを付けるだけ。"),
]
y=Inches(1.85)
for ti,body in mans:
    bx(s,Inches(0.55),y,Inches(12.23),Inches(1.12),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),y,Inches(0.1),Inches(1.12),RED)
    t(s,ti,Inches(0.8),y+Inches(0.12),Inches(11.7),Inches(0.4),sz=15,bold=True,col=RED)
    t(s,body,Inches(0.8),y+Inches(0.5),Inches(11.7),Inches(0.6),sz=11.5,col=INK,line_sp=1.12)
    y=y+Inches(1.22)

# S6 締め
s=sl(prs)
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"ONE LINE",Inches(0.9),Inches(1.7),Inches(11),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"あなたは「会って聞いて詰める」だけやればいい。",Inches(0.9),Inches(2.4),Inches(11.7),Inches(0.9),sz=30,bold=True,col=INK,line_sp=1.15)
t(s,"それが一番難しく、一番価値があり、反復で当たり前になる。",Inches(0.9),Inches(3.5),Inches(11.7),Inches(0.7),sz=24,bold=True,col=RED)
t(s,"残りは全部、こっちが持つ。",Inches(0.9),Inches(4.4),Inches(11.7),Inches(0.7),sz=24,bold=True,col=INK)
t(s,"空いた時間は、調査士と家族へ。",Inches(0.92),Inches(5.5),Inches(11),Inches(0.5),sz=14,col=GRY)
prs.save("ops_model.pptx")
print("saved ops_model.pptx", len(prs.slides._sldIdLst))

# ════════════════════════════════════════════════════════════
#  DECK ② 営業の型＝羽鳥×無敗営業
# ════════════════════════════════════════════════════════════
prs=newprs(); ft=ftmaker("営業の型  ｜  羽鳥 × 無敗営業（高橋浩一）  ｜  2026-06")

# S1 表紙
s=sl(prs)
cover(s,"THE SALES PLAYBOOK",
      "聞く＝最速の",
      "クロージング。",
      "Yesは取りに行かない。相手の意思決定を整える。\n信頼が先、収益は結果（中核信条）。",
      "菊池 研太  ｜  羽鳥式 × 無敗営業 ノウハウ")

# S2 大原則
s=sl(prs); ft(s)
hdr(s,"MINDSET","売れない営業の正体 ── Yesを急ぐ癖","焦って相手の疑問をすっ飛ばす→刺さらない→もっと焦る。この悪循環を断つ")
bx(s,Inches(0.8),Inches(1.95),Inches(11.7),Inches(1.15),REDBG); bx(s,Inches(0.8),Inches(1.95),Inches(0.12),Inches(1.15),RED)
t(s,"Yesを「成果の単位」にするのをやめる。単位は「相手の疑問を1つ潰したか」。",Inches(1.1),Inches(2.12),Inches(11.1),Inches(0.5),sz=15,bold=True,col=REDD)
t(s,"羽鳥が年収を最初に聞くのは「急ぎ」でなく「逆算」。勘所(制約)を先に握るから提案が外れない。",Inches(1.1),Inches(2.62),Inches(11.1),Inches(0.4),sz=12,col=INK)
cards2x2(s,[
 ("呪文","急いだら唱える","「Yesは結果。今やるのは相手の疑問を1つ潰すこと」"),
 ("順番","聞く→黙る→出す","質問したら3秒黙る(被せない)。沈黙が本音を出す"),
 ("KPI","Yes数を追わない","「キラー3問聞けたか/何割喋らせたか」を追う"),
 ("信条","信頼が先、収益は結果","急いで取りに行くほど信頼が減りYesが遠のく"),
], y0=Inches(3.3), ch=Inches(1.7), gy=Inches(0.26), bodysz=11.5)

# S3 キラー3問
s=sl(prs); ft(s)
hdr(s,"STEP 1 — ASK","キラー3問 ── 聞き切るまでクロージングしない","この3つを「顧客マスター」に貯める。これが「自分の顧客」の土台")
rows=[
 ("キラー質問","握れるもの","聞き方の例"),
 ("① 予算","どこまで出せるか(提案の上限)","差し支えなければ、ご予算の上限は…"),
 ("② 期限","どれだけ急ぐか(進め方)","いつまでに、が理想ですか？"),
 ("③ なぜ(こだわり)","本当のニーズ・決め手","なぜこの物件(商品)が気になりました？"),
]
light_table(s,rows,Inches(0.8),Inches(1.95),Inches(11.7),Inches(3.0),
            [Inches(3.0),Inches(4.7),Inches(4.0)],hi_col=0,sz=14,header_sz=13)
bx(s,Inches(0.8),Inches(5.2),Inches(11.7),Inches(1.5),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.8),Inches(5.2),Inches(0.1),Inches(1.5),RED)
t(s,"★最重要・抜けがち＝③こだわり",Inches(1.05),Inches(5.34),Inches(11.3),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"条件(予算/エリア/広さ)が合っても「キッチンが壁付けだから嫌」で齟齬る。\n申し込むかどうかは「気に入り角度の一致」で決まる。生の声を最初に聞く。",Inches(1.05),Inches(5.74),Inches(11.3),Inches(0.8),sz=12,col=INK,line_sp=1.2)

# S4 質問の3段ギア
s=sl(prs); ft(s)
hdr(s,"TECHNIQUE","質問の3段ギア ── 答えやすく、深く聞く","無敗営業の質問術。焦る時ほど「枕詞」で1枚クッションを置く")
steps=[
 ("枕詞(クッション)","警戒を解く","「差し支えなければ」「お答えいただける範囲で」\n→聞きにくい予算・年収を最初に聞ける"),
 ("深掘り質問","本音を広げる","「と言いますと？」「具体的には？」\n→相手に喋らせる。オープンで背景を取る"),
 ("特定質問","選択肢で絞る","「初期費を抑える／家賃を保つ、どっち？」\n→相手が答えやすく、こちらは条件を把握"),
]
cw,gx,x0,y0=Inches(3.95),Inches(0.24),Inches(0.55),Inches(1.95); CARDH=Inches(3.7)
for i,(ti,sub,body) in enumerate(steps):
    cx=x0+(cw+gx)*i
    bx(s,cx,y0,cw,CARDH,CARD,line=CARDLN,lw=1.0); bx(s,cx,y0,cw,Inches(0.75),RED)
    t(s,"GEAR "+str(i+1),cx,y0+Inches(0.1),cw,Inches(0.4),sz=15,bold=True,col=WHT,align=PP_ALIGN.CENTER)
    t(s,sub,cx,y0+Inches(0.48),cw,Inches(0.28),sz=10.5,col=RGBColor(0xF2,0xD8,0xD6),align=PP_ALIGN.CENTER)
    t(s,ti,cx+Inches(0.2),y0+Inches(0.95),cw-Inches(0.4),Inches(0.6),sz=17,bold=True,col=INK,align=PP_ALIGN.CENTER)
    t(s,body,cx+Inches(0.3),y0+Inches(1.7),cw-Inches(0.6),Inches(1.8),sz=12.5,col=GRY,line_sp=1.25)

# S5 握る（持ち帰り対策）
s=sl(prs); ft(s)
hdr(s,"STEP 2 — GRIP","「検討します」と言われてからが勝負 ── 握る","その場で決まらない時こそ、意思決定プロセスを握り、次の約束を取る")
cards2x2(s,[
 ("①","いつ","いつまでに決めますか？(意思決定の期限)"),
 ("②","誰が","誰が決めますか？(決裁者・配偶者・本部)"),
 ("③","何を基準に","何が決め手になりますか？(判断軸)"),
 ("✓","次の約束を握る","その場で次アクション(誰・いつ・何)を確定。持ち帰りを放置しない"),
])

# S6 決め手を記録（無敗データ）
s=sl(prs); ft(s)
hdr(s,"STEP 3 — RECORD","決め手を記録する ── あなただけの「無敗データ」","勝っても負けても「最後の決め手は？」を残す＝自社実績DB")
bx(s,Inches(0.8),Inches(2.0),Inches(11.7),Inches(1.3),REDBG); bx(s,Inches(0.8),Inches(2.0),Inches(0.12),Inches(1.3),RED)
t(s,"02作業DBの「結果」に →  決め手:◯◯（なぜYes / なぜNo）",Inches(1.1),Inches(2.2),Inches(11.1),Inches(0.5),sz=16,bold=True,col=REDD)
t(s,"これを貯めると勝ち筋/負け筋が見える。ルート別の成約率も可視化＝どこを太らせるか分かる。",Inches(1.1),Inches(2.78),Inches(11.1),Inches(0.5),sz=12.5,col=INK)
rows=[
 ("貯める情報","置き場","狙い"),
 ("予算・期限・こだわり","顧客マスター(H###)","人の属性=1人1行で蓄積"),
 ("結果・決め手・次アクション","02作業DB","時系列=いつ誰に何をしたか"),
 ("ルート別 成約率／自分起点比率","週次KPI","進化メーター(自分の顧客が増えてるか)"),
]
light_table(s,rows,Inches(0.8),Inches(3.5),Inches(11.7),Inches(2.6),
            [Inches(4.6),Inches(3.6),Inches(3.5)],hi_col=2,sz=12.5,header_sz=12.5)

# S7 90日道場
s=sl(prs); ft(s)
hdr(s,"90-DAY DOJO","営業道場 90日設計 ── 一点突破","聞く力は反復でしか矯正できない。道場を1つに絞り、毎日打席に立つ")
rows=[
 ("時期","道場","目的"),
 ("7月","賃貸追客(反響・クラウドミル・羽鳥師事)","キラー3問を「無意識」でできるまで反復"),
 ("8月〜","オーロラ(ケアマネ営業・採用・利用者)","固めた型をセンターピンに移植(公庫実行も8月)"),
 ("9-10月","発信(実戦ネタで)＋紹介設計","自分の蛇口を持つ ※調査士10/18は朝聖域で並行"),
]
light_table(s,rows,Inches(0.55),Inches(1.95),Inches(12.23),Inches(3.0),
            [Inches(2.0),Inches(5.5),Inches(4.73)],hi_col=0,sz=13,header_sz=13)
bx(s,Inches(0.8),Inches(5.3),Inches(11.7),Inches(1.3),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.8),Inches(5.3),Inches(0.1),Inches(1.3),RED)
t(s,"賃貸は「道場」であって「事業の柱」ではない",Inches(1.05),Inches(5.44),Inches(11.3),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"道場として自分でやる=◎(応用が効く基礎体力)。業務委託で深入り=×(薄利・センターピンを食う)。型ができたら8月オーロラへ。",Inches(1.05),Inches(5.84),Inches(11.3),Inches(0.6),sz=12,col=INK,line_sp=1.15)

# S(追加A) 心理学トリガー（羽鳥/無敗営業の実技）
s=sl(prs); ft(s)
hdr(s,"PSYCHOLOGY","心理学トリガー ── 押さずに、相手が動く8つ","羽鳥/無敗営業の実技。1チャット・1商談に最低1つ仕込む")
rows=[
 ("トリガー","使い方（実例）"),
 ("返報性","先にGIVE(有益情報/提案/紹介)→「この人に返したい」を生む(羽鳥:紹介はGIVE先行)"),
 ("損失回避","「今キャンセル無料で先に「抑える」だけ→無くなる前に」。申込ハードルを下げて先に確保"),
 ("一貫性","「内見しますよね」前提で小さなYesを積む→大きなYesへ"),
 ("希少性/締切","退去・期限を握る=「今動く理由」を相手の中に作る"),
 ("ラベリング","「ここ「重要」なのでお時間いいですか?」=重要性のラベルで耳を傾けさせる"),
 ("アンカリング","初期費用は「家賃×5倍が目安」と先に基準を置く→不安を数字で扱う"),
 ("ザイオンス","即レス＋中間報告で接触回数を稼ぐ=淡々レスでなく「提案」で返す"),
 ("沈黙","質問したら3秒黙る。沈黙が相手の本音を引き出す"),
]
light_table(s,rows,Inches(0.55),Inches(1.85),Inches(12.23),Inches(4.85),
            [Inches(2.6),Inches(9.63)],hi_col=0,sz=12,header_sz=13)

# S(追加B) 羽鳥の実トーク＆やってはいけない8
s=sl(prs); ft(s)
hdr(s,"REAL TALK","羽鳥の実トーク＆やってはいけない8 ── 実記録から","左=効いた言い回し(実記録)／右=送信前に避ける型(LINE43本の実証)")
# 左：効いた実トーク
bx(s,Inches(0.55),Inches(1.85),Inches(6.0),Inches(4.75),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(1.85),Inches(6.0),Inches(0.06),RED)
t(s,"◯ 効いた実トーク（羽鳥/ロープレ）",Inches(0.8),Inches(2.0),Inches(5.6),Inches(0.4),sz=14,bold=True,col=RED)
for i,ln in enumerate([
 "・「ここ重要なのでお時間いいですか?」",
 "・「キャンセル無料・お金は一切かからないので先に抑えましょう」",
 "・「なぜこの物件が気に入りました?」(こだわりを先に)",
 "・「初期費を抑える/家賃を保つ、どっちが近い?」",
 "・おとり物件は正直に説明→信頼(SUUMOは募集終了が大半)",
 "・ADより「顧客が住みたい家」を出す=顧客満足優先",
 "・実務はスピード優先で代行に回す(申込/FAX/内見依頼)",
]):
    t(s,ln,Inches(0.8),Inches(2.5)+Inches(0.55)*i,Inches(5.55),Inches(0.55),sz=11.5,col=INK,line_sp=1.05)
# 右：やってはいけない8ブロッカー
bx(s,Inches(6.78),Inches(1.85),Inches(6.0),Inches(4.75),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(1.85),Inches(6.0),Inches(0.06),RED)
t(s,"⛔ 送信前に避ける8ブロッカー",Inches(7.0),Inches(2.0),Inches(5.6),Inches(0.4),sz=14,bold=True,col=RED)
for i,ln in enumerate([
 "①長文 ②聞く前に提案 ③確認を詰めすぎ",
 "④構想倒れ(語るだけ) ⑤進捗共有せず個別逃げ",
 "⑥絵文字/自分語り過多 ⑦過剰謝罪 ⑧ビッグマウス",
 "── 一語「黙る/削る」で直す ──",
 "送信前3秒：①提案乗ったか ②本音を聞く質問は",
 "③数×パターンか ④心理学を1つ ⑤GIVE先行か",
 "⑥迷ったら相談(初期は必須)",
]):
    col=REDD if ln.startswith("──") else INK
    t(s,ln,Inches(7.0),Inches(2.5)+Inches(0.55)*i,Inches(5.55),Inches(0.55),sz=11.5,bold=ln.startswith("──"),col=col,line_sp=1.05)

# S(追加C) 毎日セルフチェック＆先行KPI（羽鳥不在でも回す）
s=sl(prs); ft(s)
hdr(s,"DAILY CHECK","毎日セルフチェック＆先行KPI ── 師匠がいなくても回す","羽鳥の役割はAIが代替。毎日この問いを自動で投げ、02と週次KPIに記録")
# 左：夜4問
bx(s,Inches(0.55),Inches(1.85),Inches(6.0),Inches(2.5),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.55),Inches(1.85),Inches(0.1),Inches(2.5),RED)
t(s,"🌙 夜の4問（毎晩・AIが聞く）",Inches(0.8),Inches(2.0),Inches(5.6),Inches(0.4),sz=14,bold=True,col=RED)
for i,ln in enumerate([
 "1. 今日「提案」を何件した?(0なら赤)",
 "2. 相手の本音を1つでも引き出せたか?",
 "3. 先にGIVEした相手はいるか?",
 "4. 相談した案件は?(相談ゼロも赤)",
]):
    t(s,ln,Inches(0.8),Inches(2.5)+Inches(0.42)*i,Inches(5.5),Inches(0.42),sz=12.5,col=INK)
# 右：先行KPI
bx(s,Inches(6.78),Inches(1.85),Inches(6.0),Inches(2.5),CARD,line=CARDLN,lw=1.0); bx(s,Inches(6.78),Inches(1.85),Inches(0.1),Inches(2.5),RED)
t(s,"📈 先行指標(自分で増やせる=死守)",Inches(7.0),Inches(2.0),Inches(5.6),Inches(0.4),sz=14,bold=True,col=RED)
for i,ln in enumerate([
 "・提案数 ← 最重要(追客で何件提案したか)",
 "・ヒアリング深度(本音/条件を引き出した数)",
 "・GIVE数(先に与えた回数)",
 "→ 内見→申込→成約・紹介発生(結果指標)",
]):
    t(s,ln,Inches(7.0),Inches(2.5)+Inches(0.42)*i,Inches(5.5),Inches(0.42),sz=12.5,col=INK)
# 下：連動バンド
bx(s,Inches(0.55),Inches(4.6),Inches(12.23),Inches(2.0),REDBG); bx(s,Inches(0.55),Inches(4.6),Inches(0.1),Inches(2.0),RED)
t(s,"02・顧客マスターと連動（毎日見れる/分析できる）",Inches(0.85),Inches(4.74),Inches(11.8),Inches(0.4),sz=14,bold=True,col=REDD)
for i,ln in enumerate([
 "・送信前：AIが営業文に「提案乗ったか/本音質問/心理学1つ/GIVE」を3秒チェック(羽鳥代替)",
 "・日中：対話=報告→02へ。結果に「決め手」、次アクションに「誰・いつ・何」",
 "・顧客マスター：予算/期限/こだわり＋効いた技を蓄積(1人1行)",
 "・週次KPI(05)：提案数・GIVE数・相談数・ルート別成約率を集計＝先行指標が右肩上がりか可視化",
]):
    t(s,ln,Inches(0.9),Inches(5.18)+Inches(0.35)*i,Inches(11.7),Inches(0.35),sz=11.5,col=INK)

# S8 締め
s=sl(prs)
bx(s,Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(s,"ONE LINE",Inches(0.9),Inches(1.6),Inches(11),Inches(0.4),sz=14,bold=True,col=RED)
t(s,"聞く → 詰める → 貯める。",Inches(0.9),Inches(2.3),Inches(11.7),Inches(0.9),sz=34,bold=True,col=INK)
t(s,"あとは「当たり前度」を高めるだけ。",Inches(0.9),Inches(3.4),Inches(11.7),Inches(0.8),sz=28,bold=True,col=RED)
t(s,"予算・期限・なぜ を聞く → 決まらなければ いつ・誰が・何を を握る → 決め手を記録。\nYesは追わない。相手の意思決定を整えれば、Yesは結果としてついてくる。",
  Inches(0.92),Inches(4.5),Inches(11.5),Inches(1.2),sz=14,col=GRY,line_sp=1.35)
prs.save("eigyo_kata.pptx")
print("saved eigyo_kata.pptx", len(prs.slides._sldIdLst))

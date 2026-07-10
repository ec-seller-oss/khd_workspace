# -*- coding: utf-8 -*-
"""
KHD オウンドメディア構築の流れ サマリーデッキ（イメージ写真入り）
デザインシステム: クリーム白#F9F6EF × レンガ赤#AA2E26 / ヒラギノ角ゴシック
出力: KHD_ownedmedia_flow.pptx ＋ 連動図 ownedmedia_hubspoke.png
"""
import os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = "/Users/kikuchikenta/01_honbu_docs_automation"
CREAM=(249,246,239); BRICK=(170,46,38); DARK=(43,43,43); MUTED=(138,129,122)
WHITE=(255,255,255); LINE=(224,216,204); SOFT=(245,224,220)
W6="/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc"
W3="/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc"

def f(b,s): return ImageFont.truetype(W6 if b else W3, s)

# ========== 連動図(ハブ&スポーク) ==========
def build_diagram(path):
    Wd,Ht=1600,900
    img=Image.new("RGB",(Wd,Ht),CREAM); d=ImageDraw.Draw(img)
    d.rectangle([0,0,14,Ht],fill=BRICK)
    def box(x,y,w,h,fill,outline,ow=2,r=18):
        d.rounded_rectangle([x,y,x+w,y+h],radius=r,fill=fill,outline=outline,width=ow)
    def ctext(cx,cy,s,font,fill):
        d.text((cx,cy),s,font=font,fill=fill,anchor="mm")
    def arrow(x,y0,y1):
        d.line([x,y0,x,y1-18],fill=BRICK,width=8)
        d.polygon([(x,y1),(x-16,y1-22),(x+16,y1-22)],fill=BRICK)
    # 1ソース
    box(640,40,320,70,WHITE,BRICK,3); ctext(800,75,"週1台本（1ソース）",f(True,30),DARK)
    arrow(800,112,168); ctext(940,140,"AIチームが展開",f(True,22),BRICK)
    # 4 SNS
    labels=["YouTube","X","note","Instagram"]
    bw=300; gap=30; total=bw*4+gap*3; x0=(Wd-total)//2; y=180
    for i,l in enumerate(labels):
        x=x0+i*(bw+gap); box(x,y,bw,90,WHITE,LINE,2); ctext(x+bw/2,y+45,l,f(True,30),DARK)
    # arrows down to HP
    hp_y=360
    for i in range(4):
        x=x0+i*bw+bw/2+i*gap
        d.line([x,y+90,800,hp_y-6],fill=LINE,width=2)
    arrow(800,290,hp_y)
    # HP hub
    box(540,hp_y,520,120,SOFT,BRICK,4,r=22)
    ctext(800,hp_y+44,"HP（顔）＝全SNSの集約ハブ",f(True,34),BRICK)
    ctext(800,hp_y+86,"gocinc.jp風の構成／FV・3ブロック・2段CTA",f(False,24),DARK)
    arrow(800,hp_y+120,hp_y+185)
    ctext(960,hp_y+150,"2段CTA",f(True,22),BRICK)
    # outputs
    oy=hp_y+185
    box(360,oy,520,150,WHITE,LINE,2)
    ctext(620,oy+45,"①医師",f(True,30),BRICK)
    ctext(620,oy+92,"診療圏ミニ診断（無料）",f(True,26),DARK)
    ctext(620,oy+126,"→ 医療コンサル/承継/テナント",f(False,22),MUTED)
    box(920,oy,320,150,WHITE,LINE,2)
    ctext(1080,oy+45,"②仲間",f(True,30),BRICK)
    ctext(1080,oy+92,"スモビジサロン",f(True,26),DARK)
    ctext(1080,oy+126,"（箱）",f(False,22),MUTED)
    d.text((70,Ht-40),"※ 1ソースを各SNSへ展開し、全部HPに集約→2段CTAで収益へ。回すのはAI、人は判断のみ。",font=f(False,22),fill=MUTED)
    img.save(path,"PNG"); print("saved",path)

# ========== デッキ ==========
def rgb(t): return RGBColor(*t)
def add_bg(slide, prs):
    s=slide.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb=rgb(CREAM); s.line.fill.background()
    s.shadow.inherit=False
    # left accent
    a=slide.shapes.add_shape(1,0,0,Inches(0.14),prs.slide_height)
    a.fill.solid(); a.fill.fore_color.rgb=rgb(BRICK); a.line.fill.background(); a.shadow.inherit=False
    return s

def tb(slide,x,y,w,h,text,size,color,bold=True,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    box=slide.shapes.add_textbox(x,y,w,h); tf=box.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    lines=text.split("\n")
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.name="ヒラギノ角ゴシック W6" if bold else "ヒラギノ角ゴシック W3"
        r.font.color.rgb=rgb(color)
    return box

def bullets(slide,x,y,w,h,items,size=18,color=DARK):
    box=slide.shapes.add_textbox(x,y,w,h); tf=box.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(8); r=p.add_run(); r.text="・"+it
        r.font.size=Pt(size); r.font.name="ヒラギノ角ゴシック W3"; r.font.color.rgb=rgb(color)
    return box

def build_deck(path, diagram):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    SW,SH=prs.slide_width,prs.slide_height

    # S1 表紙
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.9),Inches(2.4),Inches(11.5),Inches(1.6),"KHD オウンドメディア構築の流れ",40,DARK)
    tb(s,Inches(0.9),Inches(3.9),Inches(11.5),Inches(0.8),"医療 × 不動産 × AI ／ 1ソース→各SNS→HP集約 の連動設計",22,BRICK,bold=False)
    tb(s,Inches(0.9),Inches(6.5),Inches(11.5),Inches(0.6),"2026-06-05 時点サマリー",16,MUTED,bold=False)

    # S2 全体像 + 連動図
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.5),Inches(0.35),Inches(12),Inches(0.7),"① 全体像：オウンドメディアOS（連動図）",26,DARK)
    s.shapes.add_picture(diagram,Inches(2.1),Inches(1.25),height=Inches(5.6))

    # S3 顔HP
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.5),Inches(0.35),Inches(12),Inches(0.7),"② 顔＝HP（最も集客効果が高い『受け皿』）",26,DARK)
    bullets(s,Inches(0.7),Inches(1.4),Inches(11.8),Inches(5),[
        "全SNSが集まる常設の受け皿。ゴール逆算では一番大事＝ここを先に立てる。",
        "参考HP＝gocinc.jp（ヒーロー→実績→信頼→Founder→SNS実績→Contact）をKHDに写像。",
        "中身は完成済（実行A）：FVコピー／3ブロック（共感→解決→信頼）／2段CTA。",
        "出口①医師＝診療圏ミニ診断（無料フォーム・設問も作成済）／②仲間＝スモビジサロン。",
        "残り＝site_v3への流し込み（菊池の手番）＋Googleフォーム作成。← 次の最優先。",
    ],size=19)

    # S4 組織
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.5),Inches(0.35),Inches(12),Inches(0.7),"③ 運用組織（回すのはAI、人は判断のみ）",26,DARK)
    bullets(s,Inches(0.7),Inches(1.4),Inches(11.8),Inches(5),[
        "SNSチーム：企画→原稿→制作→投稿→振り返り（外部マニュアルをKHD実装）。",
        "YouTube制作ワークフロー：編集本体は人、前後（タイトル/サムネ/字幕/切り抜き）をAI。",
        "役割：Director=菊池／Creator=AI＋レビュー／Analyst=AI＋秘書（1人運用→将来5人）。",
        "週次フロー：月=企画→火=原稿→水=制作→木=投稿→金=振り返り（午後のメディア枠）。",
        "中核信条：売り込まずGIVE。誇大・煽りは人が必ず止める。",
    ],size=19)

    # S5 器DB
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.5),Inches(0.35),Inches(12),Inches(0.7),"④ 器（毎週回るNotion DB・実体化済）",26,DARK)
    bullets(s,Inches(0.7),Inches(1.4),Inches(11.8),Inches(5),[
        "📋 コンテンツ企画・投稿カレンダーDB：テーマ/媒体/優先度/CTA/状態/公開予定日。",
        "  ビュー3つ＝テーブル／パイプライン（状態別ボード）／投稿カレンダー。",
        "📊 週次振り返りログDB：表示/クリック/ミニ診断申込/サロン問合せ→勝ち3・改善3・来週企画。",
        "初期投入済：第1ソース（進行中）＋今週のX企画10本（未着手）＝弾が装填された状態。",
        "全部🕸️全体連動図ページからリンクで辿れる（Notion検索で繋がる）。",
    ],size=19)

    # S6 コンテンツ実例（イメージ写真）
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.5),Inches(0.35),Inches(12),Inches(0.7),"⑤ コンテンツ実例＝1ソースが図3点・4媒体に（イメージ）",24,DARK)
    imgs=[("shinryoken_demo_01_stacking.png","図1：受療率の積み上げ→約21人/日"),
          ("shinryoken_demo_02_competition.png","図2：競合按分 21→4〜6人/日"),
          ("shinryoken_demo_03_3axis.png","図A：立地を数字で見る3軸")]
    x=Inches(0.45); w=Inches(4.05)
    for fn,cap in imgs:
        p=os.path.join(BASE,fn)
        if os.path.exists(p):
            s.shapes.add_picture(p,x,Inches(1.5),width=w)
        tb(s,x,Inches(4.05),w,Inches(0.8),cap,13,DARK,bold=False,align=PP_ALIGN.CENTER)
        x=Emu(int(x)+int(w)+Inches(0.18))
    tb(s,Inches(0.7),Inches(5.1),Inches(12),Inches(1.5),
       "第1ソース「診療圏調査」＝Xスレ5連／note長文／YouTube絵コンテ／デモ画像2点。\n第2ソース「腕より立地」＝X5連／note800字／YouTube台本／HP一言＋図A。媒体ごとに別ネタを作らない＝工数1/4。",
       16,MUTED,bold=False)

    # S7 現在地・次の一手
    s=prs.slides.add_slide(blank); add_bg(s,prs)
    tb(s,Inches(0.5),Inches(0.35),Inches(12),Inches(0.7),"⑥ 現在地と次の一手",26,DARK)
    tb(s,Inches(0.7),Inches(1.3),Inches(11.8),Inches(0.5),"■ できた（設計→組織→器→中身まで通電）",18,BRICK)
    bullets(s,Inches(0.9),Inches(1.85),Inches(11.5),Inches(1.7),[
        "連動図／SNSチーム／YouTubeワークフロー／HPコピー／ミニ診断設問／DB2本",
        "第1・第2ソースの素材（テキスト4媒体＋図3点）",
    ],size=17)
    tb(s,Inches(0.7),Inches(3.7),Inches(11.8),Inches(0.5),"■ 残り（受け皿＝最優先）",18,BRICK)
    bullets(s,Inches(0.9),Inches(4.25),Inches(11.5),Inches(1.7),[
        "HP（site_v3）流し込み＋Googleフォーム作成＝送客の受け皿を立てる",
        "X下書き保存（ログイン要）／図の手動添付／福井OK／公開",
    ],size=17)
    tb(s,Inches(0.7),Inches(6.1),Inches(11.8),Inches(0.9),
       "→ 次の最優先＝あなたの逆算どおり「HPの受け皿」を立てる（フォーム作成＋流し込み手順）。",18,DARK)

    prs.save(path); print("saved",path)

if __name__=="__main__":
    dia=os.path.join(BASE,"ownedmedia_hubspoke.png")
    build_diagram(dia)
    build_deck(os.path.join(BASE,"KHD_ownedmedia_flow.pptx"), dia)

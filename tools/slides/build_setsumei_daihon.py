"""
医療事務・看護師さん向け 説明台本（提案資料に沿った司会スクリプト）。
冒頭の喋り出し・前提を最重視。全員女性・仕事喪失不安に配慮。出力: setsumei_daihon.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LN=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF); TEALBG=RGBColor(0xE1,0xF5,0xEE); TEALD=RGBColor(0x0F,0x6E,0x56)
FONT="Hiragino Sans"; W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]
def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s
def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb
def bx(slide,x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    s=slide.shapes.add_shape(shape,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s
def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=22,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.46),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)
def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LN)
    t(slide,"医療事務・看護師さん向け 説明台本（京橋クリニック 公式LINE）",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)
def block(slide,x,y,w,h,label,spoken,cue):
    bx(slide,x,y,w,h,CARD,line=CARDLN,lw=1.0); bx(slide,x,y,Inches(0.1),h,RED)
    t(slide,label,x+Inches(0.25),y+Inches(0.12),w-Inches(0.45),Inches(0.32),sz=11.5,bold=True,col=REDD)
    t(slide,spoken,x+Inches(0.27),y+Inches(0.46),w-Inches(0.5),h-Inches(0.9),sz=12.5,col=INK,line_sp=1.22)
    if cue: t(slide,cue,x+Inches(0.27),y+h-Inches(0.4),w-Inches(0.5),Inches(0.34),sz=10,col=GRY)

# ════ 1 表紙＋使い方 ════
s=sl()
bx(s,Inches(0.5),Inches(0.5),Pt(4),H-Inches(1.0),RED)
t(s,"医療事務・看護師さん向け",Inches(0.9),Inches(1.05),Inches(11),Inches(0.4),sz=15,bold=True,col=RED)
t(s,"説明台本（このまま読めばOK）",Inches(0.88),Inches(1.55),Inches(11.7),Inches(0.8),sz=29,bold=True,col=INK)
t(s,"提案資料（公式LINE）の流れに沿った、司会スクリプトです。",Inches(0.92),Inches(2.5),Inches(11.4),Inches(0.4),sz=13,col=GRY)
bx(s,Inches(0.9),Inches(3.2),Inches(11.5),Inches(1.35),CARD,line=CARDLN,lw=1.0); bx(s,Inches(0.9),Inches(3.2),Inches(0.1),Inches(1.35),RED)
t(s,"使い方",Inches(1.15),Inches(3.32),Inches(11),Inches(0.35),sz=13,bold=True,col=REDD)
t(s,"・ふつうの文＝そのまま読むセリフ　／　（ ）＝間・トーンの指示　／　太字＝必ず言う言葉\n・困ったら、この3つだけ守れば崩れません → 売り込まない／ダメ出しをもらう／皆さんが楽になる",
  Inches(1.15),Inches(3.7),Inches(11),Inches(0.8),sz=12,col=INK,line_sp=1.3)
bx(s,Inches(0.9),Inches(4.85),Inches(11.5),Inches(1.05),REDBG); bx(s,Inches(0.9),Inches(4.85),Inches(0.1),Inches(1.05),RED)
t(s,"今日の相手は、全員女性の事務・看護師さん。「あなたの仕事は大事」「仕事を奪わない」が伝わるトーンで、ゆっくり。",
  Inches(1.2),Inches(4.85),Inches(11),Inches(1.05),sz=13,bold=True,col=REDD,anchor=MSO_ANCHOR.MIDDLE)

# ════ 2 冒頭の喋り出し ════
s=sl(); ft(s)
hdr(s,"① 冒頭の喋り出し（ここが最重要）","最初の30秒で「警戒」を解く","売り込みでない・仕事を奪わない・あなたの声が欲しい、を先に伝える")
block(s,Inches(0.55),Inches(2.0),Inches(12.23),Inches(1.95),
 "つかみ（やわらかく・ゆっくり）",
 "本日はお時間ありがとうございます。最初に、一つだけ。\n今日は「売り込み」でも、皆さんのお仕事を機械に置き換える話でも、ありません。むしろ逆で——皆さんが少し楽になって、患者さんに向き合う時間を増やすための話です。",
 "（目を見て、ゆっくり。ここで場の空気がやわらぐ。急がない）")
block(s,Inches(0.55),Inches(4.1),Inches(12.23),Inches(1.95),
 "ダメ出しをお願いする（主役は皆さん）",
 "そしてもう一つ。今日いちばん欲しいのは、皆さんの「ダメ出し」です。\n受付をされている皆さんが、患者さんの動きを一番ご存知ですから。「これは要らない」「これは違う」を、どうぞ遠慮なく言ってください。",
 "（頷きを待ってから、本題へ。沈黙を怖がらない）")

# ════ 3 4つの前提 ════
s=sl(); ft(s)
hdr(s,"② 本題の前に置く「4つの前提」","この4つを先に握ると、皆さんが安心して聞ける","本題に入る前に、ゆっくり1つずつ")
pre=[("①「受付の代わり」ではありません","これは“受付の代わり”ではなく、面倒な一次対応を引き受ける「助手」です。皆さんの仕事を奪うものではありません。"),
     ("② まず1つだけ、無料で","全部を一度にやりません。まず1つだけ、無料で小さく試します。ご負担はかけません。"),
     ("③ 当院は「当日の順番受付制」のまま","予約制に変える話ではありません。今のやり方は、そのままです。"),
     ("④ お薬・症状の相談は、必ず人が","むずかしいご相談に自動で答えることはしません。医療の判断は、これまでどおり先生とスタッフが行います。")]
cw,ch,gx,gy=Inches(6.0),Inches(1.7),Inches(0.23),Inches(0.3); x0,y0=Inches(0.55),Inches(2.05)
for i,(ti,ds) in enumerate(pre):
    cx=x0+(cw+gx)*(i%2); cy=y0+(ch+gy)*(i//2)
    bx(s,cx,cy,cw,ch,CARD,line=CARDLN,lw=1.0); bx(s,cx,cy,Inches(0.1),ch,RED)
    t(s,ti,cx+Inches(0.25),cy+Inches(0.18),cw-Inches(0.45),Inches(0.5),sz=13.5,bold=True,col=REDD)
    t(s,ds,cx+Inches(0.27),cy+Inches(0.72),cw-Inches(0.5),Inches(0.9),sz=11.5,col=INK,line_sp=1.25)

# ════ 4 本題の流れ 前半 ════
s=sl(); ft(s)
hdr(s,"③ 本題の流れ（スライドに沿って）── 前半","現状 → 比較 → 患者さんの流れ","共感を取りながら。押し付けない")
block(s,Inches(0.55),Inches(1.95),Inches(12.23),Inches(1.25),
 "【現状の4つの負担】のスライドで",
 "まずは、皆さんが普段感じている負担です。鳴り止まない電話、紙からカルテへの転記、待ち時間の問い合わせ、月末の忙しさ。——どれも“あるある”ですよね？",
 "（断定せず、問いかけて頷きをもらう）")
block(s,Inches(0.55),Inches(3.35),Inches(12.23),Inches(1.15),
 "【よくある運用との比較】のスライドで",
 "他のクリニックでは、こういう所をLINEで楽にしています。当院に当てはめると、こう変わります、というのがこの表です。",
 "（他院の事例＝安心材料として淡々と）")
block(s,Inches(0.55),Inches(4.65),Inches(12.23),Inches(1.25),
 "【患者さんの流れ】のスライドで",
 "患者さんの動きで言うと——問い合わせはLINEが自動でお返事、順番が近づいたら通知、来院前にスマホで問診。受付に来た時には、もうスムーズ、という流れです。",
 "（受付が楽になる絵を、具体的に）")

# ════ 5 本題の流れ 後半 ════
s=sl(); ft(s)
hdr(s,"③ 本題の流れ ── 後半","メリット → 実物 → 効果の測り方","「楽になる」「評価でない」を必ず添える")
block(s,Inches(0.55),Inches(1.95),Inches(12.23),Inches(1.25),
 "【皆さまのメリット】のスライドで",
 "事務・看護の皆さんには、電話と転記が減ります。先生には、残業と書類が減ります。どちらか一方でなく、クリニック全体が少し楽になります。",
 "（“人を減らす”は絶対に言わない。“楽になる”で通す）")
block(s,Inches(0.55),Inches(3.35),Inches(12.23),Inches(1.15),
 "【実際の画面】のスライドで",
 "これは作っただけの絵ではなく、もう実際に動いている本物です。あとでスマホでも触ってみてください。",
 "（実機を一緒に触る時間をつくる）")
block(s,Inches(0.55),Inches(4.65),Inches(12.23),Inches(1.25),
 "【効果の測り方】のスライドで",
 "やりっぱなしにせず、電話の本数や残業を、導入の前後で一緒に測ります。これは皆さんを評価するためではなく、本当に楽になったかを確かめるためです。",
 "（“評価ではない”を必ず言い切る）")

# ════ 6 ダメ出し → はじめ方 → 締め ════
s=sl(); ft(s)
hdr(s,"④ ダメ出しをもらう → はじめ方 → 締め","ここが本番。皆さんに喋ってもらう","沈黙を怖がらず、待つ")
block(s,Inches(0.55),Inches(1.95),Inches(12.23),Inches(1.45),
 "【一緒に良くする＝ダメ出し】のスライドで",
 "ここからは、ぜひ皆さんの声を聞かせてください。他院で“これいいな”と思った仕組み、この文面や順番で伝わりにくい所、現場の困りごと——なんでも。「これは無理そう」というダメ出しが、一番ありがたいです。",
 "（問いかけたら黙る。相手が話し終わるまで待つ）")
block(s,Inches(0.55),Inches(3.55),Inches(12.23),Inches(1.0),
 "【はじめ方】のスライドで",
 "やるとしても、まず1つだけ、無料で。効果が出た所から、少しずつ。合わなければ、いつでも見直せます。",
 "（負担ゼロ・後戻りできる、を強調）")
block(s,Inches(0.55),Inches(4.7),Inches(12.23),Inches(1.2),
 "締めの言葉",
 "今日いただいたダメ出しを反映して、もう一度お持ちします。皆さんと一緒に作っていけたら嬉しいです。本日は、ありがとうございました。",
 "（“一緒に作る”で終わる。売り込みで終わらない）")

# ════ 7 NGワード＆配慮（裏カンペ） ════
s=sl(); ft(s)
hdr(s,"⑤ 言ってはいけない言葉（裏カンペ）","不安を煽る言葉は厳禁。言い換えで通す","全員女性・立場と言葉遣いに配慮")
bx(s,Inches(0.55),Inches(2.0),Inches(6.0),Inches(3.6),REDBG,line=RED,lw=1.0); bx(s,Inches(0.55),Inches(2.0),Inches(0.1),Inches(3.6),RED)
t(s,"✗ 言ってはいけない",Inches(0.8),Inches(2.12),Inches(5.5),Inches(0.4),sz=14,bold=True,col=REDD)
t(s,"✗「人件費削減」「効率化で人を減らせる」\n✗「仕事が減る／無くなる」（不安を煽る）\n✗ 難しいIT用語（API・システム連携・OSS 等）\n✗ 未確定の数字を断言（“絶対◯%減る”）\n✗ 先生にだけ向けて話す（事務・看護を置き去りにしない）",
  Inches(0.82),Inches(2.6),Inches(5.5),Inches(2.9),sz=12,col=INK,line_sp=1.5)
bx(s,Inches(6.78),Inches(2.0),Inches(6.0),Inches(3.6),TEALBG,line=TEALD,lw=1.0); bx(s,Inches(6.78),Inches(2.0),Inches(0.1),Inches(3.6),TEALD)
t(s,"○ こう言い換える",Inches(7.03),Inches(2.12),Inches(5.5),Inches(0.4),sz=14,bold=True,col=TEALD)
t(s,"○「皆さんが患者さんに向き合う時間が増える」\n○「面倒な一次対応を引き受ける“助手”」\n○「楽になったか、一緒に確かめる」\n○「まず1つだけ、無料で。いつでも見直せる」\n○ 事務・看護の皆さんに、まず先に問いかける",
  Inches(7.0),Inches(2.6),Inches(5.5),Inches(2.9),sz=12,col=INK,line_sp=1.5)
t(s,"※ 困ったら「あなたの仕事は大事です」が伝わるかを基準に、一言ずつ選ぶ。",Inches(0.55),Inches(5.8),Inches(12),Inches(0.35),sz=11.5,bold=True,col=REDD)

prs.save("setsumei_daihon.pptx")
print("saved setsumei_daihon.pptx slides:",len(prs.slides._sldIdLst))

"""
軽量版：テキスト重視のシンプルなPPTX生成
背景色・ボックスを最小限にしてファイルサイズを削減
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG   = RGBColor(0x0D, 0x1B, 0x2A)
ACC  = RGBColor(0xE8, 0xA8, 0x00)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
LGR  = RGBColor(0xCC, 0xD6, 0xE0)
MID  = RGBColor(0x1A, 0x2E, 0x44)
BLU  = RGBColor(0x4A, 0x9E, 0xCB)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def t(slide, text, x, y, w, h, sz=20, bold=False, col=WHT,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
        r.font.italic = italic
    return tb

def bx(slide, x, y, w, h, col, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = col
    if lw:
        s.line.color.rgb = col
    else:
        s.line.fill.background()
    return s

def hdr(slide, main, sub=""):
    bx(slide, 0, 0, W, Inches(1.1), MID)
    t(slide, main, Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), sz=24, bold=True)
    if sub:
        t(slide, sub, Inches(0.6), Inches(0.78), Inches(12), Inches(0.3), sz=11, col=LGR)

def ft(slide):
    bx(slide, 0, H-Inches(0.38), W, Inches(0.38), MID)
    t(slide, "医療専門不動産  |  ビジネスモデル",
      Inches(0.5), H-Inches(0.35), Inches(10), Inches(0.32), sz=9, col=LGR)

def ep_box(slide, text, y):
    bx(slide, Inches(0.55), y, Inches(0.08), Inches(0.85), ACC)
    t(slide, text, Inches(0.8), y+Inches(0.08),
      Inches(12.0), Inches(0.72), sz=11, col=WHT)

# ──────────────────────────────────────────
# SLIDE 01 カバー
# ──────────────────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "「物件を追うな、情報の非対称性を狩れ」",
  Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.55), sz=17, col=ACC, italic=True)
t(s, "医療専門不動産のビジネスモデル",
  Inches(0.9), Inches(1.55), Inches(11.5), Inches(1.2), sz=40, bold=True)
t(s, "一般不動産屋との事業モデルの比較、そして「30年単位で勝つ」ための設計思想。",
  Inches(0.9), Inches(3.0), Inches(11), Inches(0.55), sz=15, col=LGR)
bx(s, Inches(0.9), Inches(3.75), Inches(10), Inches(0.04), ACC)
t(s, "チームてっかん  |  菊池　　宅建業免許：東京都知事 (4)　　2026年",
  Inches(0.9), Inches(3.9), Inches(10), Inches(0.45), sz=14, col=LGR)

# ──────────────────────────────────────────
# SLIDE 02 AGENDA
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "AGENDA  /  本日の構成")
items = [
    ("01  事業モデル",
     "なぜ「医療×不動産」は長期で勝てるのか。／情報の非対称性という最強の堀"),
    ("02  営業ノウハウ",
     "再現できる「型」になるまでの実装。／ハンター式ヒアリングトーク含む"),
    ("03  成約事例",
     "1件の入口が、どう収益として重なるか。／6,000万粗利事例の全解剖"),
    ("04  事業展開",
     "チームとSNSで「ダムを作る」次のステージ。"),
]
cw = Inches(3.05)
for i, (ti, di) in enumerate(items):
    cx = Inches(0.5)+i*(cw+Inches(0.14))
    bx(s, cx, Inches(1.35), cw, Inches(4.7), MID)
    bx(s, cx, Inches(1.35), cw, Inches(0.06), ACC)
    t(s, ti, cx+Inches(0.15), Inches(1.5), cw-Inches(0.25), Inches(0.55), sz=16, bold=True, col=ACC)
    t(s, di, cx+Inches(0.15), Inches(2.15), cw-Inches(0.25), Inches(3.7), sz=12, col=WHT)

# ──────────────────────────────────────────
# SLIDE 03 SPEAKER
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "SPEAKER  /  自己紹介", "「カバン持ちから始まった、医療専門不動産」")
t(s, "大叔父の医業コンサル法人を承継中。宅建業免許：東京都知事 (4)",
  Inches(0.7), Inches(1.25), Inches(12), Inches(0.42), sz=14, col=LGR)
ep_box(s, "【原点エピソード】大叔父（医業コンサル第一人者）の後ろを歩きながら医師の診察室に入り続けた。"
       "気づいたのは「医師は不動産の話を、信頼できる身内にしか相談しない」こと。"
       "情報は看板ではなく、信頼の深さに従って流れる──これが業態の本質。", Inches(1.75))

takes = [("01  医療業界の構造的な「強さ」","なぜ値下げ競争に巻き込まれないのか。"),
         ("02  30年単位の顧客LTV設計","開業→拡張→機材→承継→売却まで関わる仕組み。"),
         ("03  横のつながり＝チームビルディング","士業・医療機器・保険・金融を束ねる動き方。"),
         ("04  明日から動ける、シンプルな原則","派手な近道はない。1件ずつ、目の前を大切にする。")]
for i, (ti, di) in enumerate(takes):
    rx = Inches(0.7)+(i%2)*Inches(6.2)
    ry = Inches(3.5)+(i//2)*Inches(1.0)
    t(s, ti, rx, ry, Inches(5.8), Inches(0.42), sz=13, bold=True, col=BLU)
    t(s, di, rx+Inches(0.2), ry+Inches(0.42), Inches(5.6), Inches(0.42), sz=11, col=LGR)

# ──────────────────────────────────────────
# SLIDE 04 この場の約束
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "今日お伝えするのは「不動産の売り方」ではない。")
nots = ["✕  物件の見つけ方", "✕  SUUMO・レインズの活用法", "✕  仲介手数料の交渉術"]
yess = ["✓  情報の源泉（門番）にどう近づくか", "✓  「悩み」から数千万の収益を生む設計", "✓  値下げ競争に入らず30年勝ち続ける構造"]
bx(s, Inches(0.6), Inches(1.35), Inches(5.7), Inches(2.4), MID)
t(s, "NOT TEACHING", Inches(0.8), Inches(1.45), Inches(5.3), Inches(0.38), sz=12, col=LGR)
for i, x in enumerate(nots):
    t(s, x, Inches(0.8), Inches(1.95)+i*Inches(0.6), Inches(5.3), Inches(0.52), sz=14, col=RGBColor(0xFF,0x66,0x66))
bx(s, Inches(6.9), Inches(1.35), Inches(6.0), Inches(2.4), MID)
bx(s, Inches(6.9), Inches(1.35), Inches(6.0), Inches(0.06), ACC)
t(s, "TEACHING", Inches(7.1), Inches(1.45), Inches(5.6), Inches(0.38), sz=12, bold=True, col=ACC)
for i, x in enumerate(yess):
    t(s, x, Inches(7.1), Inches(1.95)+i*Inches(0.6), Inches(5.6), Inches(0.52), sz=13, col=WHT)
bx(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.4), MID)
bx(s, Inches(0.6), Inches(4.0), Inches(0.08), Inches(1.4), ACC)
t(s, "「物件を追っているうちは、仲介手数料（3%）の奴隷から抜け出せない。\n"
     " 10〜30%の粗利は、ネットの向こうではなく『人の悩みの深さ』の中にある。」",
  Inches(0.85), Inches(4.1), Inches(11.6), Inches(1.2), sz=16, col=ACC, italic=True)

# ──────────────────────────────────────────
# SLIDE 05 SECTION 01 DIVIDER
# ──────────────────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "SECTION", Inches(0.9), Inches(1.2), Inches(5), Inches(0.6), sz=20, col=LGR)
t(s, "01", Inches(0.9), Inches(1.8), Inches(5), Inches(1.8), sz=90, bold=True, col=ACC)
t(s, "事業モデル", Inches(0.9), Inches(3.7), Inches(10), Inches(0.75), sz=32, bold=True)
t(s, "なぜ「医療×不動産」は、長期で勝てるのか。\n情報の非対称性・業界構造・長期事業の設計、の4枚で。",
  Inches(0.9), Inches(4.55), Inches(11), Inches(1.0), sz=15, col=LGR)

# ──────────────────────────────────────────
# SLIDE 06 会社の成り立ち
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "01  /  会社の成り立ちと実績",
    "「医療コンサル不動産」は、銀行の現場発注で生まれた。後付けではなく、現場で組み上げた業態。")
steps = [("01 証券会社時代","金融の原則を身につける"),
         ("02 NECメディカル","医療業界の構造に気付く"),
         ("03 住友銀行との連携","「貸せるがコンサルがいない」"),
         ("04 事業化・会社設立","型を文書化、事業として確立")]
sw = Inches(2.85)
for i, (ti, di) in enumerate(steps):
    sx = Inches(0.5)+i*(sw+Inches(0.18))
    bx(s, sx, Inches(1.3), sw, Inches(1.65), MID)
    bx(s, sx, Inches(1.3), sw, Inches(0.06), ACC)
    t(s, ti, sx+Inches(0.15), Inches(1.38), sw-Inches(0.2), Inches(0.5), sz=14, bold=True, col=ACC)
    t(s, di, sx+Inches(0.15), Inches(1.9), sw-Inches(0.2), Inches(0.9), sz=12, col=LGR)

ep_box(s, "【エピソード：住友銀行の担当者の一言】「菊池さん、うちは医師に貸したい。でも事業計画書を"
       "正しく読めるコンサルがいないから融資の判断ができない」──銀行は資金を持つが目利きができない。"
       "私は人脈と知見を持つが融資力が弱い。「タッグを組めば最強だ」翌日から連携が始まった。", Inches(3.15))

stats = [("140件超","開業支援"), ("20件超","事業承継"), ("17件超","経営支援・救済")]
for i, (v, l) in enumerate(stats):
    sx = Inches(0.5)+i*Inches(4.25)
    bx(s, sx, Inches(4.25), Inches(4.0), Inches(1.1), MID)
    t(s, v, sx+Inches(0.2), Inches(4.32), Inches(3.6), Inches(0.55), sz=28, bold=True, col=ACC)
    t(s, l, sx+Inches(0.2), Inches(4.87), Inches(3.6), Inches(0.38), sz=13, col=LGR)

# ──────────────────────────────────────────
# SLIDE 07 業界の違い
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "02  /  業界の違いと構造的優位",
    "医療不動産は「売買」ではなく「コンサルティング」。真の強みは業界横断の「チームビルディング」にある。")
bx(s, Inches(0.5), Inches(1.25), Inches(5.8), Inches(2.3), RGBColor(0x15,0x1B,0x2E))
t(s, "一般不動産　単発の「売買仲介」",
  Inches(0.7), Inches(1.35), Inches(5.4), Inches(0.42), sz=14, bold=True, col=LGR)
for i, x in enumerate(["情報はオープン、利幅は薄い","毎回ゼロから集客","1件決まれば、次は振り出し"]):
    t(s, f"・{x}", Inches(0.7), Inches(1.85)+i*Inches(0.55), Inches(5.4), Inches(0.48), sz=13, col=LGR)
bx(s, Inches(6.9), Inches(1.25), Inches(5.9), Inches(2.3), MID)
bx(s, Inches(6.9), Inches(1.25), Inches(5.9), Inches(0.06), ACC)
t(s, "医療不動産　「コンサル」＋「チームビルディング」",
  Inches(7.1), Inches(1.35), Inches(5.5), Inches(0.42), sz=14, bold=True, col=WHT)
for i, x in enumerate(["開業→移転→拡張→機材→承継→売却と連鎖","信頼は次の紹介を呼ぶ","利益が長期で複利的に積み上がる"]):
    t(s, f"・{x}", Inches(7.1), Inches(1.85)+i*Inches(0.55), Inches(5.5), Inches(0.48), sz=13, col=WHT)
t(s, "VS", Inches(6.0), Inches(2.0), Inches(0.85), Inches(0.6), sz=18, bold=True, col=ACC, align=PP_ALIGN.CENTER)
reasons = [("①  情報の非対称性","医師は本業多忙で不動産情報に疎い。比較サイトの安値競争に陥らず、信頼で選ばれる。"),
           ("②  ライフサイクルの長さ","1人の医師に30年関わり続けられる。開業→拡張→機材→承継→引退後の自宅売却まで。"),
           ("③  チームビルディング","士業・医療機器・薬品卸・保険・設計・金融──業界横断チームを束ねる位置に立つ。")]
rw = Inches(3.95)
for i, (ti, di) in enumerate(reasons):
    rx = Inches(0.5)+i*(rw+Inches(0.14))
    bx(s, rx, Inches(3.75), rw, Inches(2.0), MID)
    bx(s, rx, Inches(3.75), rw, Inches(0.06), ACC)
    t(s, ti, rx+Inches(0.15), Inches(3.83), rw-Inches(0.25), Inches(0.42), sz=13, bold=True, col=ACC)
    t(s, di, rx+Inches(0.15), Inches(4.32), rw-Inches(0.25), Inches(1.3), sz=11, col=WHT)

# ──────────────────────────────────────────
# SLIDE 08 情報の非対称性（NEW）
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "情報の非対称性  ──  核心スライド  【NEW】",
    "良い情報は「不動産屋」には来ない。「信頼できる身内」にしか落ちない。")
flow = ["医師の本音", "門番\n（税理士・弁護士\n・銀行員）", "私たち\n（コンサルとして\n認知された場合のみ）", "不動産を\n「処方箋」として\n活用"]
fw = Inches(2.75)
for i, f in enumerate(flow):
    fx = Inches(0.5)+i*(fw+Inches(0.45))
    bx(s, fx, Inches(1.3), fw, Inches(2.0), MID)
    bx(s, fx, Inches(1.3), fw, Inches(0.06), ACC)
    t(s, f, fx+Inches(0.15), Inches(1.45), fw-Inches(0.25), Inches(1.7), sz=13, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    if i < 3:
        t(s, "→", fx+fw+Inches(0.07), Inches(2.05), Inches(0.38), Inches(0.5), sz=18, col=ACC)

bx(s, Inches(0.5), Inches(3.5), Inches(5.7), Inches(1.35), RGBColor(0x15,0x15,0x2E))
t(s, "一般流通物件", Inches(0.7), Inches(3.6), Inches(5.3), Inches(0.38), sz=14, bold=True, col=LGR)
t(s, "仲介手数料3%　全員が同じ情報を持つ", Inches(0.7), Inches(4.05), Inches(5.3), Inches(0.65), sz=13, col=LGR)
bx(s, Inches(7.0), Inches(3.5), Inches(5.8), Inches(1.35), MID)
bx(s, Inches(7.0), Inches(3.5), Inches(5.8), Inches(0.06), ACC)
t(s, "非公開・難件", Inches(7.2), Inches(3.6), Inches(5.4), Inches(0.38), sz=14, bold=True, col=ACC)
t(s, "粗利10〜30%　情報を持つのは1人だけ", Inches(7.2), Inches(4.05), Inches(5.4), Inches(0.65), sz=13, col=WHT)

ep_box(s, "【エピソード】大叔父のカバン持ち中、医師が「来年クリニックを移転しようと思っているが、"
       "どこに相談すればいいか分からなくて……」と言った。不動産屋としてではなく「連れ」として入ったからこそ聞けた話。"
       "情報は看板ではなく信頼の深さに従って流れる。", Inches(5.05))

# ──────────────────────────────────────────
# SLIDE 09 長期事業モデル
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "03  /  長期事業モデルの作り方",
    "1件の相談を、生涯にわたる収益ストリームに育てる。複利で積み上がる事業設計。")
streams = [("0年目","相談・受付","紹介・無料"),
           ("0–1年目","診療圏調査","コンサルfee"),
           ("1年目","事業計画・物件","仲介手数料"),
           ("1年〜継続","機材・備品","卸売・継続収益"),
           ("5–30年目","拡張・承継・売却","再仲介・買取再販")]
sw2 = Inches(2.38)
for i, (yr, ti, rev) in enumerate(streams):
    sx = Inches(0.45)+i*(sw2+Inches(0.12))
    bx(s, sx, Inches(1.3), sw2, Inches(2.1), MID)
    bx(s, sx, Inches(1.3), sw2, Inches(0.06), ACC)
    t(s, yr, sx+Inches(0.12), Inches(1.38), sw2-Inches(0.2), Inches(0.38), sz=11, col=LGR)
    t(s, ti, sx+Inches(0.12), Inches(1.8), sw2-Inches(0.2), Inches(0.55), sz=13, bold=True, col=WHT)
    t(s, rev, sx+Inches(0.12), Inches(2.42), sw2-Inches(0.2), Inches(0.8), sz=11, col=ACC)
t(s, "1人の医師から、生涯で  数千万〜数億円規模  の累積取引が積み上がる構造。",
  Inches(0.5), Inches(3.6), Inches(12.2), Inches(0.42), sz=15, bold=True, col=ACC)
principles = [("01 利他のマインド","目先の手数料を追わず伴走"),
              ("02 ハブとして繋ぐ","全部やらない。最初に電話される存在に"),
              ("03 双方向ネットワーク","「もらう」関係を脱し「紹介する側」にも"),
              ("04 30年の領域選び","ライフサイクルが長い業界でLTVを伸ばす"),
              ("05 ノウハウの「型」化","属人性を抜き、再現性を残す")]
pw = Inches(2.38)
for i, (ti, di) in enumerate(principles):
    px = Inches(0.45)+i*(pw+Inches(0.12))
    bx(s, px, Inches(4.2), pw, Inches(1.95), MID)
    t(s, ti, px+Inches(0.12), Inches(4.3), pw-Inches(0.2), Inches(0.42), sz=12, bold=True, col=BLU)
    t(s, di, px+Inches(0.12), Inches(4.78), pw-Inches(0.2), Inches(1.2), sz=11, col=LGR)

# ──────────────────────────────────────────
# SLIDE 10 SECTION 02
# ──────────────────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "SECTION", Inches(0.9), Inches(1.2), Inches(5), Inches(0.6), sz=20, col=LGR)
t(s, "02", Inches(0.9), Inches(1.8), Inches(5), Inches(1.8), sz=90, bold=True, col=ACC)
t(s, "営業ノウハウ", Inches(0.9), Inches(3.7), Inches(10), Inches(0.75), sz=32, bold=True)
t(s, "再現できる「型」になるまで、現場で何をやってきたか。\n5ステップ──士業との関係・相談ハブ化・診療圏調査・事業計画書・ハンター式ヒアリング。",
  Inches(0.9), Inches(4.55), Inches(11), Inches(1.1), sz=15, col=LGR)

# ──────────────────────────────────────────
# SLIDE 11 STEP01 士業との関係構築
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "STEP 01  /  士業との関係構築",
    "「お願いします」をやめた瞬間、紹介が来始めた。5年通った、泥臭い5年。")
bx(s, Inches(0.5), Inches(1.25), Inches(5.8), Inches(2.3), RGBColor(0x15,0x15,0x2E))
t(s, "BEFORE  「お願いします」の営業",
  Inches(0.7), Inches(1.33), Inches(5.4), Inches(0.4), sz=13, bold=True, col=LGR)
for i, x in enumerate(["士業事務所を回り案件を「もらう」立場","他業者と並列、紹介は来ない","条件の悪い案件だけ回ってくる"]):
    t(s, f"・{x}", Inches(0.7), Inches(1.83)+i*Inches(0.58), Inches(5.4), Inches(0.5), sz=12, col=LGR)
bx(s, Inches(6.85), Inches(1.25), Inches(5.95), Inches(2.3), MID)
bx(s, Inches(6.85), Inches(1.25), Inches(5.95), Inches(0.06), ACC)
t(s, "AFTER  「ご紹介します」の営業",
  Inches(7.05), Inches(1.33), Inches(5.5), Inches(0.4), sz=13, bold=True, col=ACC)
for i, x in enumerate(["医師に「税理士いますか？◯◯先生を紹介します」","士業に頼られる位置に、案件は双方向に","コミュニティ内で「あの会社」と固有名詞で呼ばれる"]):
    t(s, f"・{x}", Inches(7.05), Inches(1.83)+i*Inches(0.58), Inches(5.5), Inches(0.5), sz=12, col=WHT)
t(s, "VS", Inches(6.05), Inches(2.0), Inches(0.8), Inches(0.6), sz=18, bold=True, col=ACC, align=PP_ALIGN.CENTER)
ep_box(s, "【転換点エピソード】税理士の先生への5回目の挨拶。「先週、信頼できる税理士を探している医師の先生から相談がありました。"
       "ぜひ先生をご紹介したいのですが」──その瞬間表情が変わった。1ヶ月後「クリニックの移転で相談したい先生がいる」と電話が来た。"
       "最初に投げることで、初めて返ってくる。", Inches(3.72))
reality = [("断られても顔を出し続ける","5回、10回通ってようやく覚えられる"),
           ("「便利屋」上等で動く","「いつでも対応してくれる」が入口"),
           ("勉強会・葬式まで出る","業務外で会う回数が信頼の上限"),
           ("「先に紹介」を実装する","もらう前に投げる。義理で返ってくる")]
rw2 = Inches(2.95)
for i, (ti, di) in enumerate(reality):
    rx = Inches(0.5)+i*(rw2+Inches(0.12))
    bx(s, rx, Inches(5.0), rw2, Inches(1.05), MID)
    t(s, ti, rx+Inches(0.12), Inches(5.08), rw2-Inches(0.2), Inches(0.38), sz=12, bold=True, col=BLU)
    t(s, di, rx+Inches(0.12), Inches(5.5), rw2-Inches(0.2), Inches(0.45), sz=10, col=LGR)

# ──────────────────────────────────────────
# SLIDE 12 STEP02 相談プラットフォーム化
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "STEP 02  /  相談プラットフォーム化",
    "横のつながりが「集客の入口」になる。「最初に電話される存在」になる。")
bx(s, Inches(5.4), Inches(2.65), Inches(2.5), Inches(0.9), ACC)
t(s, "医療チームのHUB", Inches(5.4), Inches(2.68), Inches(2.5), Inches(0.85),
  sz=13, bold=True, col=BG, align=PP_ALIGN.CENTER)
partners = [("医師・クリニック","起点", Inches(5.6), Inches(1.35)),
            ("会計士・税理士","提携", Inches(0.6), Inches(1.95)),
            ("弁護士・司法書士","提携", Inches(0.6), Inches(3.0)),
            ("医療機器メーカー","自社対応", Inches(0.6), Inches(4.05)),
            ("薬品卸・販売企業","提携", Inches(0.6), Inches(5.1)),
            ("保険代理店","提携", Inches(9.55), Inches(1.95)),
            ("設計建築・内装","提携", Inches(9.55), Inches(3.0)),
            ("ハウスメーカー","提携", Inches(9.55), Inches(4.05)),
            ("金融機関","提携", Inches(9.55), Inches(5.1)),
            ("人材紹介","提携", Inches(5.6), Inches(5.7))]
for nm, ro, px, py in partners:
    bx(s, px, py, Inches(2.7), Inches(0.65), MID)
    t(s, nm, px+Inches(0.1), py+Inches(0.05), Inches(2.0), Inches(0.35), sz=12, bold=True)
    t(s, ro, px+Inches(0.1), py+Inches(0.38), Inches(1.5), Inches(0.25), sz=9, col=ACC)
t(s, "自分で全部やる必要はない。「最初に電話される存在」になれば、案件は向こうから来る。",
  Inches(0.6), Inches(6.62), Inches(12.2), Inches(0.38), sz=13, bold=True, col=ACC, italic=True)

# ──────────────────────────────────────────
# SLIDE 13 STEP03 診療圏調査
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "STEP 03  /  診療圏調査",
    "数字で「ここなら勝てる」を示す。サンプル：T整形外科リウマチ科（川崎市宮前区／鷺沼駅圏）")
ps = [("01 診療圏マップ確定","地形・道路・公共交通でエリア確定"),
      ("02 競合施設の調査","競合7施設をリスト化"),
      ("03 診療圏人口集計","総人口48,201人"),
      ("04 受療率の適用","厚労省データで最適化"),
      ("05 自院推定外来算定","(人口×受療率)÷(競合数+1)")]
pw2 = Inches(2.38)
for i, (ti, di) in enumerate(ps):
    px = Inches(0.45)+i*(pw2+Inches(0.12))
    bx(s, px, Inches(1.28), pw2, Inches(1.8), MID)
    bx(s, px, Inches(1.28), pw2, Inches(0.06), ACC)
    t(s, ti, px+Inches(0.12), Inches(1.35), pw2-Inches(0.2), Inches(0.55), sz=12, bold=True, col=ACC)
    t(s, di, px+Inches(0.12), Inches(1.95), pw2-Inches(0.2), Inches(0.95), sz=11, col=LGR)
out_items = [("対象エリア","鷺沼駅/半径2km"), ("診療圏人口","48,201人"),
             ("65歳以上比率","23.8%"), ("競合施設","7施設"),
             ("競合総外来","523人/日"), ("自院推定（昼）","133人/日"), ("自院推定（夜）","123人/日")]
ow = Inches(1.68)
for i, (la, va) in enumerate(out_items):
    ox = Inches(0.45)+i*(ow+Inches(0.1))
    bx(s, ox, Inches(3.28), ow, Inches(1.3), MID)
    t(s, va, ox+Inches(0.08), Inches(3.35), ow-Inches(0.12), Inches(0.62), sz=17, bold=True, col=ACC, align=PP_ALIGN.CENTER)
    t(s, la, ox+Inches(0.08), Inches(3.97), ow-Inches(0.12), Inches(0.45), sz=9, col=LGR, align=PP_ALIGN.CENTER)
t(s, "テンプレに見えるが、過去140件超の現場データに裏打ちされた一品物。AI時代でも、ここは経験で差が出る。",
  Inches(0.45), Inches(4.73), Inches(12.2), Inches(0.35), sz=11, col=LGR, italic=True)
t(s, "【現場の声】「こんな資料を出してくれた業者は初めてです。先生は本当に私のことを考えてくれているんですね」──数字は信頼を作る。",
  Inches(0.45), Inches(5.18), Inches(12.2), Inches(0.55), sz=12, col=WHT, italic=True)

# ──────────────────────────────────────────
# SLIDE 14 STEP04 事業計画書
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "STEP 04  /  事業計画書",
    "銀行が通す事業計画書──融資が通れば、すべてが動く。必要資金2.3億円 / 借入2億円")
fmt = [("①資金調達","自己資金/借入/利率/返済"), ("②概算総事業費","建設・不動産・諸費用・運転"),
       ("③施設概要","法人・住所・診療科・床面積"), ("④減価償却","建物・設備・医療機器"),
       ("⑤人件費","常勤・パート・福利・賞与"), ("⑥収支算定","外来×単価×日数/経費率"),
       ("⑦年度別損益","5期分・税引後・返済財源"), ("⑧借入返済明細","元金均等の年次CF")]
fw2 = Inches(2.85)
for i, (ti, di) in enumerate(fmt[:4]):
    fx = Inches(0.5)+i*(fw2+Inches(0.1))
    bx(s, fx, Inches(1.28), fw2, Inches(0.95), MID)
    t(s, ti, fx+Inches(0.12), Inches(1.35), fw2-Inches(0.2), Inches(0.35), sz=12, bold=True, col=BLU)
    t(s, di, fx+Inches(0.12), Inches(1.73), fw2-Inches(0.2), Inches(0.4), sz=10, col=LGR)
for i, (ti, di) in enumerate(fmt[4:]):
    fx = Inches(0.5)+i*(fw2+Inches(0.1))
    bx(s, fx, Inches(2.33), fw2, Inches(0.95), MID)
    t(s, ti, fx+Inches(0.12), Inches(2.4), fw2-Inches(0.2), Inches(0.35), sz=12, bold=True, col=BLU)
    t(s, di, fx+Inches(0.12), Inches(2.78), fw2-Inches(0.2), Inches(0.4), sz=10, col=LGR)
years = [("R8（初年度）","+2,970"), ("R9","+7,962"), ("R10","+12,966"), ("R11","+18,618"), ("R12","+25,822")]
yw = Inches(2.38)
for i, (yr, va) in enumerate(years):
    yx = Inches(0.5)+i*(yw+Inches(0.12))
    bx(s, yx, Inches(3.5), yw, Inches(1.15), MID)
    bx(s, yx, Inches(3.5), yw, Inches(0.06), ACC)
    t(s, yr, yx+Inches(0.12), Inches(3.57), yw-Inches(0.2), Inches(0.35), sz=10, col=LGR)
    t(s, f"{va} 千円", yx+Inches(0.12), Inches(3.95), yw-Inches(0.2), Inches(0.55), sz=18, bold=True, col=ACC)
t(s, "※ 5年で借入残高1.0億円・返済原資が剰余に転じる",
  Inches(0.5), Inches(4.8), Inches(12), Inches(0.35), sz=11, col=LGR)
t(s, "「この計画書を持ってくる業者の案件なら、うちは前向きに検討します」──融資担当者を動かす資料が作れる。これが最大の差。",
  Inches(0.5), Inches(5.3), Inches(12.2), Inches(0.55), sz=12, col=WHT, italic=True)

# ──────────────────────────────────────────
# SLIDE 15 STEP05 ハンター式ヒアリング（NEW）
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "STEP 05  /  ハンター式ヒアリングトーク  【NEW】",
    "「不動産屋に話さない本音」を引き出す3ステップ。入り口は絶対に「不動産」の話をしない。")
steps_h = [
    ("STEP①  業界特有の「苦しみ」への共感",
     "「先生、最近は診療報酬の改定で、現場を回すだけでも手出しが増えてるんじゃないですか？」",
     "→ 「この人は自分の業界を知っている」というインサイダー認定を得る。"),
    ("STEP②  最悪のシナリオの言語化",
     "「もし今月、メインバンクからの融資が止まったら、どの資産から手をつけますか？」\n「後継者の方とは、万が一の際の出口について具体的に話したことはありますか？」",
     "→ 相手が「言えない悩み」を持っているかを確認する。"),
    ("STEP③  処方箋の提示（出口）",
     "「売ることが目的ではなく、先生の法人の『再起動』が目的です。\n私が泥を被って、1ヶ月以内に現金を用意するスキームを組みます。」",
     "→ 「守るための現金化」として提案。「業者」ではなく「救済者」として認識される。"),
]
sh = Inches(1.75)
for i, (ti, qu, no) in enumerate(steps_h):
    bx(s, Inches(0.5), Inches(1.28)+i*(sh+Inches(0.12)), Inches(12.3), sh, MID)
    bx(s, Inches(0.5), Inches(1.28)+i*(sh+Inches(0.12)), Inches(0.08), sh, ACC)
    t(s, ti, Inches(0.75), Inches(1.35)+i*(sh+Inches(0.12)), Inches(3.0), sh-Inches(0.2), sz=12, bold=True, col=ACC)
    t(s, qu, Inches(3.9), Inches(1.35)+i*(sh+Inches(0.12)), Inches(5.3), sh-Inches(0.2), sz=12, col=WHT, italic=True)
    t(s, no, Inches(9.35), Inches(1.35)+i*(sh+Inches(0.12)), Inches(3.3), sh-Inches(0.2), sz=10, col=LGR)
t(s, "相手が「この人は同じ側の人間だ」と感じた瞬間から、本当の営業が始まる。",
  Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), sz=12, col=ACC, italic=True)

# ──────────────────────────────────────────
# SLIDE 16 SECTION 03
# ──────────────────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "SECTION", Inches(0.9), Inches(1.2), Inches(5), Inches(0.6), sz=20, col=LGR)
t(s, "03", Inches(0.9), Inches(1.8), Inches(5), Inches(1.8), sz=90, bold=True, col=ACC)
t(s, "成約事例", Inches(0.9), Inches(3.7), Inches(10), Inches(0.75), sz=32, bold=True)
t(s, "1件の入口が、どう収益として重なっていくか。\n6,000万粗利の全解剖と、対照的な2件──「単発で終わらない」事業の実態を。",
  Inches(0.9), Inches(4.55), Inches(11), Inches(1.0), sz=15, col=LGR)

# ──────────────────────────────────────────
# SLIDE 17 CASE01 弁護士ルート
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "CASE 01  /  弁護士ルート  ──  横のつながりが直接効いた事例",
    "「医療を分かる人いる？」──弁護士からの1本の電話。一般不動産屋が断った居抜き案件。")
c1 = [("01 入口","弁護士から居抜きの紹介","電源・水回り・フロア用途を評価できず一般屋で断られた案件が、医療分かる窓口として回ってきた。"),
      ("02 なぜ我々に","「医療だけ」と思われていない","事務所内で「人柄含めて信頼できる窓口」と認知。不動産が絡む案件はカジュアルに投げてもらえる位置。"),
      ("03 アプローチ","自社買取＋再販","売り急ぎで市場に出すと足元を見られる。買取で時間を作り、医療向けにバリューアップして次の使い手へ。"),
      ("04 メッセージ","専門特化＝逆に間口が広がる","医療を切り口にすると、一般不動産の窓口にもなれる。コンサルできる軸を1本持つことが効く。")]
cw3 = Inches(2.95)
for i, (n, ti, di) in enumerate(c1):
    cx = Inches(0.5)+i*(cw3+Inches(0.12))
    bx(s, cx, Inches(1.25), cw3, Inches(2.6), MID)
    bx(s, cx, Inches(1.25), cw3, Inches(0.06), ACC)
    t(s, n, cx+Inches(0.12), Inches(1.32), cw3-Inches(0.2), Inches(0.35), sz=10, col=LGR)
    t(s, ti, cx+Inches(0.12), Inches(1.7), cw3-Inches(0.2), Inches(0.5), sz=13, bold=True, col=WHT)
    t(s, di, cx+Inches(0.12), Inches(2.25), cw3-Inches(0.2), Inches(1.45), sz=11, col=LGR)
cp_active = ["物件仲介手数料", "買取・再販差益", "士業・税理士紹介"]
cp_all = ["診療圏調査fee","事業計画書fee","物件仲介手数料","買取・再販差益","機材・リース","士業・税理士紹介","顧問料","承継・M&A"]
cpw = Inches(1.48)
for i, lab in enumerate(cp_all):
    cx = Inches(0.5)+i*(cpw+Inches(0.06))
    col = ACC if lab in cp_active else RGBColor(0x30,0x40,0x55)
    bx(s, cx, Inches(4.05), cpw, Inches(0.7), col)
    t(s, lab, cx+Inches(0.05), Inches(4.12), cpw-Inches(0.08), Inches(0.55), sz=9, col=BG if lab in cp_active else LGR)
ep_box(s, "【弁護士からの電話】「医療を分かる人いる？顧問先のクリニックが資金繰りで詰まっていて、"
       "一般の不動産屋が全員断ってくる」──断られた理由はすぐわかった。医療機器の残存価値、"
       "電気容量、フロア用途変更費用。翌朝8時に現地へ。当日夕方に買取価格を提示した。", Inches(5.0))

# ──────────────────────────────────────────
# SLIDE 18 CASE02 危機対応
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "CASE 02  /  危機対応  ──  長期顧客化の典型例",
    "「閉院 vs 救済」の岐路で、医師の人生に伴走した数年。事務長の使い込みで資金破綻寸前。")
c2 = [("01 入口","事務長の使い込みで資金破綻寸前","「人」の問題で経営が傾いた典型例。チームを設計してきたから相談が来る。"),
      ("02 普通なら","閉院→物件売却で「終わり」","医師は廃業、設備は二束三文、患者は離散。一般不動産屋ならここまで。"),
      ("03 アプローチ","自宅売却＋チームから人材派遣","資金は不動産売却で確保。チームから事務長を送り込み、診療を止めずに再建。"),
      ("04 結果","売却益→派遣→顧問料の連結","1件の不動産案件が、派遣収入と顧問料の継続収益に変わった。")]
for i, (n, ti, di) in enumerate(c2):
    cx = Inches(0.5)+i*(cw3+Inches(0.12))
    bx(s, cx, Inches(1.25), cw3, Inches(2.4), MID)
    bx(s, cx, Inches(1.25), cw3, Inches(0.06), ACC)
    t(s, n, cx+Inches(0.12), Inches(1.32), cw3-Inches(0.2), Inches(0.35), sz=10, col=LGR)
    t(s, ti, cx+Inches(0.12), Inches(1.7), cw3-Inches(0.2), Inches(0.5), sz=13, bold=True, col=WHT)
    t(s, di, cx+Inches(0.12), Inches(2.25), cw3-Inches(0.2), Inches(1.25), sz=11, col=LGR)
rev = [("Y1","自宅売却仲介＋利益","資金繰り解消"),("Y1","事務長の派遣収入","チームから人材送り売上化"),
       ("継続","顧問料（毎月）","経営支援の継続収入"),("Y2+","クリニック移転・機材","立て直し後の案件"),
       ("Y5+","事業承継の伴走","後継候補接続・関係継続")]
rw3 = Inches(2.38)
for i, (yr, ti, di) in enumerate(rev):
    rx = Inches(0.5)+i*(rw3+Inches(0.12))
    bx(s, rx, Inches(3.85), rw3, Inches(1.45), MID)
    bx(s, rx, Inches(3.85), rw3, Inches(0.06), ACC)
    t(s, yr, rx+Inches(0.12), Inches(3.92), rw3-Inches(0.2), Inches(0.32), sz=12, bold=True, col=ACC)
    t(s, ti, rx+Inches(0.12), Inches(4.27), rw3-Inches(0.2), Inches(0.42), sz=12, bold=True, col=WHT)
    t(s, di, rx+Inches(0.12), Inches(4.72), rw3-Inches(0.2), Inches(0.48), sz=10, col=LGR)
t(s, "学び：「人」の問題まで踏み込めると、不動産1件＋派遣＋顧問料で収益が重なる。チームがあるからこそ動ける。",
  Inches(0.5), Inches(5.48), Inches(12.2), Inches(0.38), sz=12, col=LGR, italic=True)
t(s, "→ 次のスライド：6,000万が生まれた瞬間──この構造を徹底解剖",
  Inches(0.5), Inches(5.95), Inches(12.2), Inches(0.38), sz=13, col=ACC, bold=True)

# ──────────────────────────────────────────
# SLIDE 19 6,000万が生まれた瞬間（NEW）
# ──────────────────────────────────────────
s = sl(); ft(s)
bx(s, 0, 0, W, Inches(1.1), RGBColor(0x1A,0x2E,0x1A))
t(s, "6,000万が生まれた瞬間  ──  破産寸前→資産現金化→再起動  【NEW】",
  Inches(0.6), Inches(0.18), Inches(12), Inches(0.65), sz=21, bold=True, col=ACC)
t(s, "医療法人救済の全記録。半年で完結。",
  Inches(0.6), Inches(0.78), Inches(12), Inches(0.3), sz=12, col=LGR)
timeline = [("0ヶ月目  接触","大叔父（医業コンサル第一人者）のカバン持ちとして医療法人の経営相談に同席。不動産業者としてではなく「信頼できる連れ」として入室。"),
            ("3〜5ヶ月目  発掘","「事務長が資金を横領」という極秘情報が門番（大叔父）を通じて届く。資金ショートまであと2週間。"),
            ("6ヶ月目  信頼の借用","「大叔父が信頼している人」というお墨付きで院長に面談。警戒心はゼロ。「任せます」をその場でもらった。"),
            ("決済期  スピードが命","ノンバンク×プロジェクト融資を駆使。自宅＋医院不動産を48時間で査定・提示。「閉院」を「再起動」に変えた。")]
tw2 = Inches(2.95)
for i, (ti, di) in enumerate(timeline):
    tx = Inches(0.5)+i*(tw2+Inches(0.12))
    bx(s, tx, Inches(1.28), tw2, Inches(2.9), MID)
    bx(s, tx, Inches(1.28), tw2, Inches(0.06), ACC)
    t(s, ti, tx+Inches(0.12), Inches(1.35), tw2-Inches(0.2), Inches(0.55), sz=12, bold=True, col=ACC)
    t(s, di, tx+Inches(0.12), Inches(1.95), tw2-Inches(0.2), Inches(2.1), sz=11, col=WHT)
bx(s, Inches(0.5), Inches(4.38), Inches(12.3), Inches(0.82), RGBColor(0x1A,0x3A,0x1A))
t(s, "結果：自社粗利  約6,000万円",
  Inches(0.7), Inches(4.45), Inches(6), Inches(0.42), sz=22, bold=True, col=ACC)
t(s, "院長の感想：「あなたがいなかったら終わっていた」",
  Inches(0.7), Inches(4.88), Inches(12), Inches(0.28), sz=13, col=WHT, italic=True)
ep_box(s, "【院長の一言】「最初は半信半疑でした。でもあなたは逃げなかった。48時間で現金を用意してくれた。"
       "うちの病院は終わりかけていたのに、今ここで診療を続けられている。本当にありがとうございます」\n"
       "──不動産は手段に過ぎない。誰かの人生の危機に、最後まで伴走できるかどうか。それだけです。", Inches(5.42))

# ──────────────────────────────────────────
# SLIDE 20 PATTERN CATALOG
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "PATTERN CATALOG  /  案件のバラエティ",
    "チームを束ねている結果、こうした案件パターンが日常的に発生する──「単発仲介」では捉えられない世界。")
pats = [("P01","新規開業の物件選定","診療圏調査→候補抽出→賃貸契約まで士業ハブで完結。"),
        ("P02","クリニック移転・拡張","現物件の出口設計と新物件選定を同時に。"),
        ("P03","医療法人化に伴う本部移転","登記要件と運営動線を両立。会計士・行政書士と並走。"),
        ("P04","機材入替・リース紹介","提携機材商社へ。設備更新タイミングで再接触。"),
        ("P05","内装リニューアル","提携内装業者へ繋ぎ、紹介マージン。動線改善も。"),
        ("P06","事業承継・M&A","後継候補との接続、銀行調整、設備譲渡。1案件で複数収益。"),
        ("P07","廃業・救済（債権者対応）","閉院ではなく事業継続の道を提示。弁護士・会計士と数年並走。"),
        ("P08","多店舗展開支援","2院目・3院目。組織設計と物件取得を同時伴走。")]
pw3 = Inches(2.95)
for i, (n, ti, di) in enumerate(pats):
    px = Inches(0.45)+(i%4)*(pw3+Inches(0.12))
    py = Inches(1.28)+(i//4)*Inches(2.35)
    bx(s, px, py, pw3, Inches(2.15), MID)
    bx(s, px, py, pw3, Inches(0.06), ACC)
    t(s, n, px+Inches(0.12), py+Inches(0.1), pw3, Inches(0.38), sz=15, bold=True, col=ACC)
    t(s, ti, px+Inches(0.12), py+Inches(0.52), pw3-Inches(0.2), Inches(0.5), sz=12, bold=True, col=WHT)
    t(s, di, px+Inches(0.12), py+Inches(1.08), pw3-Inches(0.2), Inches(0.95), sz=10, col=LGR)
t(s, "年140件超の中で、「単発で終わる案件」はほとんど無い。どのパターンも、次のパターンの入口になっている。",
  Inches(0.45), Inches(6.9), Inches(12.2), Inches(0.38), sz=11, col=LGR, italic=True)

# ──────────────────────────────────────────
# SLIDE 21 SECTION 04（NEW）
# ──────────────────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
t(s, "SECTION", Inches(0.9), Inches(1.2), Inches(5), Inches(0.6), sz=20, col=LGR)
t(s, "04", Inches(0.9), Inches(1.8), Inches(5), Inches(1.8), sz=90, bold=True, col=ACC)
t(s, "事業展開", Inches(0.9), Inches(3.7), Inches(10), Inches(0.75), sz=32, bold=True)
t(s, "チームとSNSで「信頼のダムを作る」次のステージ。\n1人でやる限界を超えるための設計。",
  Inches(0.9), Inches(4.55), Inches(11), Inches(1.0), sz=15, col=LGR)

# ──────────────────────────────────────────
# SLIDE 22 次のステージ（NEW）
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "次のステージ  /  「信頼のダムを作る」  【NEW】",
    "あなたが構築すべきは「不動産を売る仕組み」ではなく、「誰にも言えない悩みが自動的に集まってくるダム」。")
t(s, "入り口はコンサル、出口はハンター。",
  Inches(0.6), Inches(1.28), Inches(12), Inches(0.42), sz=18, bold=True)
sns = [("① 恐怖と救済\n（エピソード発信）",
        "「今月ショートする医療法人。事務長は逃亡。残されたのは動かない不動産。そこで私がやったこと──」\n→ 特定層が震えるリアルな実話を投稿。"),
       ("② 情報の非対称性\n（専門家発信）",
        "「不動産業者に相談すると叩かれるだけの物件も、医療コンサルの視点を加えると価値が3倍になる理由」\n→ 専門知識を公開し権威を確立。"),
       ("③ 逃げない覚悟\n（人間性発信）",
        "「カバン持ちから学んだ、最後の一人まで向き合う泥臭さ」\n→ 顧客が最も恐れる「業者の逃げ」を否定する。")]
sw3 = Inches(3.95)
for i, (ti, di) in enumerate(sns):
    sx = Inches(0.5)+i*(sw3+Inches(0.12))
    bx(s, sx, Inches(1.9), sw3, Inches(2.8), MID)
    bx(s, sx, Inches(1.9), sw3, Inches(0.06), ACC)
    t(s, ti, sx+Inches(0.15), Inches(1.98), sw3-Inches(0.25), Inches(0.7), sz=13, bold=True, col=ACC)
    t(s, di, sx+Inches(0.15), Inches(2.72), sw3-Inches(0.25), Inches(1.85), sz=11, col=WHT)
teams = ["各業界の「太客の隣」に潜り込む\n（医療・福祉・士業など）",
         "難件を拾ってくる\n（入り口コンサル）",
         "出口の刈り取りを\n指揮・バックアップ",
         "粗利を分配し\n実績としてコンテンツ化"]
tw3 = Inches(2.95)
for i, t_ in enumerate(teams):
    tx = Inches(0.5)+i*(tw3+Inches(0.12))
    bx(s, tx, Inches(4.9), tw3, Inches(1.0), RGBColor(0x1A,0x35,0x50))
    t(s, f"Step {i+1}", tx+Inches(0.12), Inches(4.97), tw3-Inches(0.2), Inches(0.3), sz=10, col=ACC)
    t(s, t_, tx+Inches(0.12), Inches(5.3), tw3-Inches(0.2), Inches(0.55), sz=11, col=WHT)
t(s, "「不動産屋を名乗っているうちは、仲介手数料の枠から出られない。」",
  Inches(0.5), Inches(6.1), Inches(12.2), Inches(0.35), sz=12, col=ACC, italic=True, bold=True)
t(s, "「一生営業したくないなら、一生モノの『信頼のダム』を今すぐ掘れ。」",
  Inches(0.5), Inches(6.5), Inches(12.2), Inches(0.35), sz=12, col=ACC, italic=True, bold=True)

# ──────────────────────────────────────────
# SLIDE 23 CLOSING
# ──────────────────────────────────────────
s = sl(); ft(s)
hdr(s, "CLOSING  /  明日から動ける、ノウハウの結論",
    "派手な近道はない。1件ずつ、目の前を大切にする。")
closes = [("フットワーク軽く","連絡は即返信。スピードで信頼を作る。"),
          ("顔を出し続けろ","1回では覚えられない。3回、5回、10回。義理が信頼に化ける。"),
          ("足を使え","現地・士業事務所・医師の診察室。机では分からないことが現場にある。"),
          ("頭を使え","30年スパンで考える。1件の手数料ではなく生涯LTVから逆算する。"),
          ("金を使え","情報・人脈・道具に投資する。ケチると、機会のほうが先に去る。"),
          ("まず1件取れ","完璧を待たずに動く。1件目を「事故なく」完遂すれば次の100件の入口になる。")]
cw4 = Inches(3.95)
for i, (ti, di) in enumerate(closes):
    cx = Inches(0.45)+(i%3)*(cw4+Inches(0.12))
    cy = Inches(1.28)+(i//3)*Inches(1.32)
    bx(s, cx, cy, cw4, Inches(1.18), MID)
    bx(s, cx, cy, cw4, Inches(0.06), ACC)
    t(s, ti, cx+Inches(0.15), cy+Inches(0.12), cw4-Inches(0.25), Inches(0.38), sz=14, bold=True, col=ACC)
    t(s, di, cx+Inches(0.15), cy+Inches(0.55), cw4-Inches(0.25), Inches(0.55), sz=11, col=WHT)
ep_box(s, "【エピソード：最初の1件】正直、何もわからなかった。でも逃げなかった。深夜まで資料を直した。融資担当者に電話をかけ続けた。"
       "その1件が完遂できた日、医師からこう言われた。「次の先生も紹介しますよ」──それが2件目の入口だった。"
       "完璧を待っていたら、その1件も2件目もなかった。", Inches(4.08))
t(s, "そして──「医師を分かる業者」になれ。士業も医師も、結局そこを見ている。",
  Inches(0.45), Inches(5.6), Inches(12.2), Inches(0.42), sz=15, bold=True, col=ACC)

# ──────────────────────────────────────────
# SLIDE 24 エンドスライド
# ──────────────────────────────────────────
s = sl()
bx(s, 0, 0, Inches(0.5), H, ACC)
bx(s, 0, H-Inches(1.1), W, Inches(1.1), MID)
bx(s, Inches(0.9), Inches(3.45), Inches(10), Inches(0.04), ACC)
t(s, "ご清聴ありがとうございました。",
  Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.55), sz=19, col=LGR)
t(s, "1人の医師に、\n30年寄り添う。",
  Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.7), sz=44, bold=True)
t(s, "これが、医療専門不動産のビジネスモデル。",
  Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.52), sz=19, col=LGR)
t(s, "チームてっかん  |  菊池",
  Inches(0.9), Inches(5.0), Inches(7), Inches(0.48), sz=19, bold=True)
t(s, "ご相談・面談、お気軽にどうぞ。",
  Inches(0.9), Inches(5.55), Inches(7), Inches(0.4), sz=13, col=LGR)

out = "/Users/kikuchikenta/01_honbu_docs_automation/medical_realestate_seminar_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
import os; print(f"Size: {os.path.getsize(out)/1024:.1f} KB")

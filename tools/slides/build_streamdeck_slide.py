# -*- coding: utf-8 -*-
"""
Stream Deck イメージスライド（京橋クリニック向け）
物理ボタン15キーのモックアップ＋御院での使い方。KHD配色。月曜デモ用イメージ。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); INK=RGBColor(0x1A,0x1A,0x1A)
GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF); CARD=RGBColor(0xF1,0xEC,0xE1)
CARDLN=RGBColor(0xE1,0xDA,0xCB); WHT=RGBColor(0xFF,0xFF,0xFF); GRN=RGBColor(0x2E,0x7D,0x32)
DARK=RGBColor(0x23,0x24,0x28); KEY=RGBColor(0x2E,0x6E,0x8E); KEY2=RGBColor(0x3A,0x7A,0x52); KEYR=RGBColor(0xB5,0x3A,0x32)
FONT="Hiragino Sans"
W=Inches(13.333); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
s=prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb=BG

def bx(x,y,w,h,col,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.shadow.inherit=False
    if col is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=col
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp
def t(x,y,w,h,text,sz=12,b=False,c=INK,al=PP_ALIGN.LEFT,an=MSO_ANCHOR.TOP,ls=None):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=an
    tf.margin_left=Emu(0); tf.margin_right=Emu(0); tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=al
        if ls: p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.name=FONT; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
    return tb

# ヘッダー
bx(Inches(0.5),Inches(0.45),Pt(4),H-Inches(0.9),RED)
t(Inches(0.9),Inches(0.45),Inches(11.5),Inches(0.35),"STREAM DECK ｜ 入力をワンタッチに",11,True,RED)
t(Inches(0.88),Inches(0.78),Inches(11.5),Inches(0.55),"ボタン1つで、定型文・紹介状・予約案内を一瞬で入力",23,True,INK)
bx(Inches(0.9),Inches(1.42),Inches(11.5),Pt(3),RED)
t(Inches(0.92),Inches(1.55),Inches(11.5),Inches(0.3),"医療事務の「最初の文字入力」と「紹介状の二度手間」を、押すだけに。電子カルテはそのまま使えます。",11.5,False,GRY)

# ===== 左：Stream Deckデバイス・モックアップ（15キー 5x3）=====
dev_x, dev_y, dev_w, dev_h = Inches(0.95), Inches(2.15), Inches(6.0), Inches(4.0)
bx(dev_x, dev_y, dev_w, dev_h, DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
t(dev_x, dev_y+Inches(0.12), dev_w, Inches(0.3), "Stream Deck（Elgato）", 11, True, WHT, al=PP_ALIGN.CENTER)
# キー配置
keys=[
 ("再診\nお呼び出し",KEY),("発熱\n外来案内",KEYR),("予約\n空き確認",KEY),("処方\n説明定型",KEY2),("会計\n案内",KEY),
 ("紹介状\nひな型",KEYR),("診断書\nひな型",KEYR),("検査\n説明",KEY2),("次回\n予約案内",KEY),("LINE\n登録案内",KEY2),
 ("よくある\n質問①",KEY),("よくある\n質問②",KEY),("住所\n定型",KEY2),("自費\n料金案内",KEY),("受付\n締め時間",KEYR),
]
gx0, gy0 = dev_x+Inches(0.3), dev_y+Inches(0.6)
kw, kh, gap = Inches(1.0), Inches(0.92), Inches(0.12)
for idx,(lab,col) in enumerate(keys):
    rr=idx//5; cc=idx%5
    kx=gx0+(kw+gap)*cc; ky=gy0+(kh+gap)*rr
    bx(kx,ky,kw,kh,col,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t(kx,ky,kw,kh,lab,8.5,True,WHT,al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE,ls=0.95)
t(dev_x, dev_y+dev_h+Inches(0.08), dev_w, Inches(0.3), "※ボタンの中身（定型文）は御院に合わせて自由に設定できます", 9.5, False, GRY, al=PP_ALIGN.CENTER)

# ===== 右：御院での使い方＋効果 =====
RX, RW = Inches(7.25), Inches(5.2)
bx(RX,Inches(2.15),RW,Inches(2.55),CARD,line=CARDLN,lw=1.0)
bx(RX,Inches(2.15),RW,Inches(0.06),RED)
t(RX+Inches(0.25),Inches(2.3),RW-Inches(0.5),Inches(0.3),"御院での使い方（例）",13,True,RED)
uses=[
 ("①","紙問診の決まり文句や再診の呼び出しを、1ボタンで入力"),
 ("②","紹介状・診断書のひな型を一発呼び出し→二度手間を解消"),
 ("③","「次回予約は…」「LINE登録は…」の案内文もワンタッチ"),
 ("④","電子カルテ（Medicom）はそのまま。キー操作で文字を挿入するだけ"),
]
yy=Inches(2.7)
for n,d in uses:
    t(RX+Inches(0.28),yy,Inches(0.4),Inches(0.5),n,12,True,RED)
    t(RX+Inches(0.65),yy,RW-Inches(0.95),Inches(0.6),d,10.5,False,INK,ls=1.05)
    yy=yy+Inches(0.47)

bx(RX,Inches(4.9),RW,Inches(1.25),WHT,line=CARDLN,lw=1.0)
t(RX+Inches(0.25),Inches(5.02),RW-Inches(0.5),Inches(0.3),"なぜ最初の一手に向くか",12,True,RED)
t(RX+Inches(0.28),Inches(5.36),RW-Inches(0.55),Inches(0.7),
  "・低コスト（機器代のみ）・電子カルテ改修不要・すぐ試せる\n・効果が目に見える＝現場が「楽になった」を体感しやすい",10,False,INK,ls=1.2)

# 下バンド：月曜デモ
by=Inches(6.35)
bx(Inches(0.95),by,Inches(11.5),Inches(0.7),RED)
t(Inches(1.2),by+Inches(0.08),Inches(11.0),Inches(0.55),
  "月曜デモ予定：実機を持ち込み、御院の定型文を入れて「押すだけ入力」をその場で体感いただきます（宮崎が事前検証）",12,True,WHT,an=MSO_ANCHOR.MIDDLE)
t(Inches(0.95),Inches(7.12),Inches(11),Inches(0.3),"テナントアシスト・ウイン株式会社 ｜ 菊池 研太",9,False,GRY)

import shutil
prs.save("streamdeck_image.pptx")
_dst="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/260526_AI医療コンサル/京橋_StreamDeckイメージ.pptx"
shutil.copy("streamdeck_image.pptx",_dst)
print("SAVED:",_dst)

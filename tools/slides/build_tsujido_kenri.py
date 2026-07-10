"""
辻堂東海岸3-3-20 近隣 権利関係マップ
KHDデザイン（クリーム白×レンガ赤）。公図＋取得謄本を統合し、登記名義 vs 現地表札・私道共有者・位置を整理。
出力: tsujido_kenri.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0xF9, 0xF6, 0xEF)
RED    = RGBColor(0xAA, 0x2E, 0x26)
REDD   = RGBColor(0x8C, 0x24, 0x1D)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GRY    = RGBColor(0x6E, 0x6E, 0x6E)
LINE   = RGBColor(0xDA, 0xD6, 0xCF)
CARD   = RGBColor(0xF1, 0xEC, 0xE1)
CARDLN = RGBColor(0xE1, 0xDA, 0xCB)
REDBG  = RGBColor(0xF4, 0xE4, 0xE2)
GRYBG  = RGBColor(0xEC, 0xE8, 0xDF)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Hiragino Sans"
W = Inches(13.33); H = Inches(7.5)
KOUZU = "/tmp/tsujido_kouzu.png"

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s


def t(slide, text, x, y, w, h, sz=18, bold=False, col=INK, align=PP_ALIGN.LEFT,
      italic=False, anchor=MSO_ANCHOR.TOP, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp: p.line_spacing = line_sp
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col; r.font.name = FONT
    return tb


def bx(slide, x, y, w, h, col, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def hdr(slide, eyebrow, main, sub=""):
    t(slide, eyebrow, Inches(0.6), Inches(0.4), Inches(12), Inches(0.4), sz=13, bold=True, col=RED)
    bx(slide, Inches(0.62), Inches(0.78), Inches(1.7), Pt(3), RED)
    t(slide, main, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.55), sz=22, bold=True, col=INK)
    if sub:
        t(slide, sub, Inches(0.62), Inches(1.44), Inches(12.1), Inches(0.3), sz=11.5, col=GRY)


def ft(slide):
    bx(slide, Inches(0.5), H-Inches(0.5), Inches(12.33), Pt(1.2), LINE)
    t(slide, "辻堂東海岸3-3-20  ｜  近隣 権利関係マップ（2026-06-17時点）", Inches(0.5), H-Inches(0.42), Inches(10), Inches(0.32), sz=9, col=GRY)


def light_table(slide, rows, x, y, w, h, col_w, hi_col=None, sz=12, header_sz=12):
    n, m = len(rows), len(rows[0])
    tb = slide.shapes.add_table(n, m, x, y, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w): tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci); cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            is_hi = (hi_col is not None and ci == hi_col)
            if ri == 0: cell.fill.fore_color.rgb = REDD if is_hi else RED
            else: cell.fill.fore_color.rgb = REDBG if is_hi else (CARD if ri % 2 == 1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_sz if ri == 0 else sz)
                    r.font.bold = (ri == 0) or is_hi or (ci == 0)
                    if ri == 0: r.font.color.rgb = WHT
                    elif is_hi: r.font.color.rgb = RED
                    elif ci == 0: r.font.color.rgb = INK
                    else: r.font.color.rgb = RGBColor(0x3A,0x3A,0x3A)
    return tb


def framed_img(slide, path, x, y, frame_w, frame_h, cap=""):
    bx(slide, x, y, frame_w, frame_h, CARD, line=CARDLN, lw=1.0)
    bx(slide, x, y, frame_w, Inches(0.06), RED)
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih
    pad = Inches(0.18)
    avail_w = frame_w - pad*2; avail_h = frame_h - pad*2 - (Inches(0.3) if cap else Inches(0))
    if avail_w/avail_h > ar:
        ph = avail_h; pw = Emu(int(ph*ar))
    else:
        pw = avail_w; ph = Emu(int(pw/ar))
    px = x + Emu(int((frame_w-pw)/2)); py = y + pad
    slide.shapes.add_picture(path, px, py, width=pw, height=ph)
    if cap:
        t(slide, cap, x, y+frame_h-Inches(0.32), frame_w, Inches(0.3), sz=9.5, col=GRY, align=PP_ALIGN.CENTER)


# ════════ SLIDE 1 — 表紙 ════════
s = sl()
bx(s, Inches(0.5), Inches(0.45), Pt(4), H-Inches(0.9), RED)
t(s, "TSUJIDO ｜ 近隣 権利関係マップ", Inches(0.9), Inches(1.5), Inches(7.6), Inches(0.45), sz=15, bold=True, col=RED)
t(s, "辻堂東海岸3-3-20", Inches(0.88), Inches(2.05), Inches(7.7), Inches(0.9), sz=38, bold=True, col=INK)
t(s, "お隣は誰か？ 権利関係マップ", Inches(0.88), Inches(2.85), Inches(7.7), Inches(0.9), sz=34, bold=True, col=RED)
t(s, "公図＋取得謄本を突き合わせ、隣人の「登記名義（同意を取る相手）」と\n「現地の表札（会う人）」を整理。前面道路・私道の共有者まで一枚に。",
  Inches(0.9), Inches(3.85), Inches(7.4), Inches(0.9), sz=13, col=GRY, line_sp=1.25)
# 3カード
ox, oy, ow, og = Inches(0.9), Inches(5.0), Inches(2.42), Inches(0.14)
cards = [("取得した謄本", "6筆"), ("お隣の登記名義", "別人"), ("私道(共有)", "無し")]
for i, (lab, val) in enumerate(cards):
    cx = ox + (ow + og) * i
    bx(s, cx, oy, ow, Inches(1.25), CARD, line=CARDLN, lw=1.0)
    bx(s, cx, oy, ow, Inches(0.06), RED)
    t(s, lab, cx, oy+Inches(0.16), ow, Inches(0.35), sz=11.5, col=GRY, align=PP_ALIGN.CENTER)
    t(s, val, cx, oy+Inches(0.5), ow, Inches(0.7), sz=27, bold=True, col=RED, align=PP_ALIGN.CENTER)
bx(s, Inches(0.9), Inches(6.62), Inches(7.5), Pt(1.2), LINE)
t(s, "テナントアシスト・ウイン  ｜  菊池 研太  ｜  2026-06-17", Inches(0.9), Inches(6.74), Inches(7.5), Inches(0.4), sz=12, bold=True, col=INK)
# 右：公図
framed_img(s, KOUZU, Inches(8.95), Inches(1.4), Inches(3.95), Inches(5.55), cap="公図（7239-245ほか・2025/6/24取得）")

# ════════ SLIDE 2 — 結論 ════════
s = sl(); ft(s)
hdr(s, "KEY POINT", "結論：名義人＝同意を取る相手／表札＝現地で会う人（別人）",
    "ここを取り違えると交渉が空回りする。福井さんも曖昧なので、まずこの3点を握る")
pts = [
    ("関口さん＝表札（居住者）", "福井さんが現地・南側で会った人の表札が「関口」。だが南側の登記名義（449=柏木ら／738=有川／739=宮崎）に関口はいない＝居住者の可能性。"),
    ("「藤本」という別登記は無い", "738/739を取得して表札と一致するか照合した際のラベル。登記名義はあくまで有川・宮崎で、藤本という名義は存在しない。"),
    ("同意を取る相手＝登記名義人", "南側ブロック後退（駐車解消）の同意は、登記名義人＝7239-449の柏木ら4名 全員から取る必要がある。関口さんは関係づくりに有効だが、ハンコを持つ人ではない。"),
]
cy = Inches(1.95)
for i, (ti, body) in enumerate(pts):
    bx(s, Inches(0.55), cy, Inches(12.23), Inches(1.5), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.55), cy, Inches(0.12), Inches(1.5), RED)
    t(s, str(i+1), Inches(0.8), cy+Inches(0.32), Inches(0.8), Inches(0.8), sz=30, bold=True, col=RED)
    t(s, ti, Inches(1.75), cy+Inches(0.2), Inches(10.7), Inches(0.5), sz=17, bold=True, col=INK)
    t(s, body, Inches(1.77), cy+Inches(0.72), Inches(10.8), Inches(0.7), sz=12.5, col=GRY, line_sp=1.15)
    cy = cy + Inches(1.7)

# ════════ SLIDE 3 — 位置関係（公図＋並び） ════════
s = sl(); ft(s)
hdr(s, "POSITION", "位置関係 ── 公図で見る並び", "住居表示の並び（3-3-19→24）と地番・名義を対応。★が対象地")
framed_img(s, KOUZU, Inches(0.55), Inches(1.8), Inches(5.0), Inches(5.0), cap="公図（辻堂東海岸三丁目7239番台）")
# 右：並びダイアグラム
seq = [
    ("3-3-19", "7239-246", "櫻井 照雄", "北隣", False),
    ("3-3-20", "7239-245+985", "松井・大石ほか", "★対象地（売主=松井）", True),
    ("3-3-21", "7239-739", "宮崎 茂", "隣接・きれい", False),
    ("3-3-22", "7239-738", "有川 龍", "隣接・⚠仮差押+抵当", False),
    ("3-3-23", "7239-641", "小池 元壽", "0.43㎡の残地（私道でない）", False),
    ("3-3-24", "7239-449", "柏木ら4名", "南隣＝本丸（4名同意要）", False),
]
dx, dy, dw, dh, dg = Inches(5.95), Inches(1.85), Inches(6.85), Inches(0.66), Inches(0.1)
for i, (juu, chi, nm, memo, hi) in enumerate(seq):
    cy = dy + (dh+dg)*i
    bx(s, dx, cy, dw, dh, REDBG if hi else CARD, line=RED if hi else CARDLN, lw=1.5 if hi else 1.0)
    bx(s, dx, cy, Inches(0.1), dh, RED)
    t(s, juu, dx+Inches(0.22), cy+Inches(0.06), Inches(1.25), dh-Inches(0.12), sz=15, bold=True, col=RED, anchor=MSO_ANCHOR.MIDDLE)
    t(s, chi, dx+Inches(1.5), cy+Inches(0.06), Inches(1.7), dh-Inches(0.12), sz=11.5, col=GRY, anchor=MSO_ANCHOR.MIDDLE)
    t(s, nm, dx+Inches(3.15), cy+Inches(0.06), Inches(1.85), dh-Inches(0.12), sz=14, bold=True, col=INK, anchor=MSO_ANCHOR.MIDDLE)
    t(s, memo, dx+Inches(4.95), cy+Inches(0.06), Inches(1.8), dh-Inches(0.12), sz=10, col=GRY, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.0)
t(s, "※「関口」は南側の現地表札（居住者）で、上記いずれの登記名義とも不一致＝別人。",
  dx, dy+(dh+dg)*6+Inches(0.02), dw, Inches(0.4), sz=10.5, bold=True, col=REDD)

# ════════ SLIDE 4 — 権利関係テーブル ════════
s = sl(); ft(s)
hdr(s, "RIGHTS", "権利関係 一覧 ── 取得謄本ベース（登記名義＝同意を取る相手）",
    "「現地の表札」列は福井さんが会った人。名義と表札が別なら、まず名義人を押さえる")
rows = [
    ("地番", "住居表示", "登記名義人（同意を取る相手）", "持分", "現地の表札", "メモ"),
    ("7239-245+985", "3-3-20", "松井弘枝・敏行／大石3名", "—", "—", "対象地（売主=松井）"),
    ("7239-246", "3-3-19", "櫻井 照雄", "単独", "—", "北隣・筆界確認済"),
    ("7239-739", "3-3-21", "宮崎 茂", "単独", "—", "S57取得・権利関係きれい"),
    ("7239-738", "3-3-22", "有川 龍", "単独", "—", "⚠仮差押(R7/2/4)＋抵当多数"),
    ("7239-449", "3-3-24", "柏木・菊池英一・清水2名", "1/4等", "関口（居住者）", "南隣＝本丸・4名全員同意要"),
    ("7239-641", "3-3-23", "小池 元壽", "単独", "—", "738/739の親番・0.43㎡の残地"),
]
light_table(s, rows, Inches(0.45), Inches(1.85), Inches(12.43), Inches(4.6),
            [Inches(1.85), Inches(1.25), Inches(3.55), Inches(1.0), Inches(1.85), Inches(2.93)],
            hi_col=2, sz=11.5, header_sz=12)
t(s, "★南側ブロック後退の同意権者＝7239-449の柏木ら4名全員。関口さんは表札（居住者）で同意権者ではない。",
  Inches(0.45), Inches(6.6), Inches(12.4), Inches(0.35), sz=11.5, bold=True, col=REDD)

# ════════ SLIDE 4.5 — 関口は誰か／駐車交渉の意味 ════════
s = sl(); ft(s)
hdr(s, "KEY QUESTION", "関口さんは誰か／駐車交渉に意味はあるか", "近隣6筆の登記名義を全照合した結論。福井さんの「デメリット」質問への実質的な答え")
bx(s, Inches(0.55), Inches(1.8), Inches(12.23), Inches(1.3), CARD, line=CARDLN, lw=1.0)
bx(s, Inches(0.55), Inches(1.8), Inches(0.12), Inches(1.3), RED)
t(s, "関口さん＝南側隣地（7239-449・3-3-24）の居住者。登記の所有者ではない。", Inches(0.85), Inches(1.95), Inches(11.6), Inches(0.4), sz=15, bold=True, col=RED)
t(s, "近隣6筆の名義を全照合→関口は所有者に不在。その土地7239-449の名義は柏木洋子・菊池英一・清水信敬・清水澄子の4名共有。",
  Inches(0.87), Inches(2.38), Inches(11.6), Inches(0.6), sz=12.5, col=INK, line_sp=1.2)
judg = [
    ("△", "関口さん単独では駐車は空かない", "名義人ではない＝ブロック後退の同意権が無い。関口さんとの交渉だけでは前進しない。"),
    ("○", "ただし名義人4名への「入口」にはなる", "449世帯の住人で、福井さんが既に関係構築済。名義人4名へ当たる糸口として活かせる。"),
    ("✗", "本丸＝4名全員の同意（難度高）", "共有4名の全員合意が必須。1人反対で不成立。長期化すれば金利(10/30満了)で粗利が溶ける。"),
]
cy = Inches(3.35)
for mk, ti, body in judg:
    bx(s, Inches(0.55), cy, Inches(12.23), Inches(0.85), CARD, line=CARDLN, lw=1.0)
    t(s, mk, Inches(0.68), cy+Inches(0.18), Inches(0.8), Inches(0.55), sz=22, bold=True, col=RED, align=PP_ALIGN.CENTER)
    t(s, ti, Inches(1.6), cy+Inches(0.08), Inches(10.9), Inches(0.4), sz=14, bold=True, col=INK)
    t(s, body, Inches(1.62), cy+Inches(0.46), Inches(10.9), Inches(0.35), sz=11, col=GRY)
    cy = cy + Inches(0.98)
bx(s, Inches(0.55), Inches(6.42), Inches(12.23), Inches(0.55), REDBG)
bx(s, Inches(0.55), Inches(6.42), Inches(0.1), Inches(0.55), RED)
t(s, "結論：「隣に空けてもらう」で駐車確保は現実性が低い（相手4名・全員同意・他人）。本線は荻原3,400万の業者買取。交渉はやるなら期限付き、ダメなら荻原。",
  Inches(0.8), Inches(6.5), Inches(11.9), Inches(0.4), sz=11.5, bold=True, col=REDD)

# ════════ SLIDE 5 — 前面道路・私道 ════════
s = sl(); ft(s)
hdr(s, "ROAD", "前面道路は私道か？ ── 結論：共有の私道は無し", "親番7239-641まで取得して確認。738/739/641すべて単独所有で、共有名義の道路は存在しない")
facts = [
    ("①  738/739/641すべて単独所有だった", "7239-738＝有川龍／7239-739＝宮崎茂／親番7239-641＝小池元壽（3-3-23）。いずれも単独で、共有名義の土地は無い。"),
    ("②  7239-641は0.43㎡の極小残地", "208㎡から738・739を分筆した残り＝わずか0.43㎡の宅地。道路として機能する共有地ではなかった。"),
    ("③  結論：共有の私道は存在しない", "「道路の共有者」に当たる共有私道は無し。前面道路は公道（または共有でない通路）。＝この論点はクローズ。"),
]
cy = Inches(2.0)
for i, (ti, body) in enumerate(facts):
    bx(s, Inches(0.55), cy, Inches(12.23), Inches(1.35), CARD, line=CARDLN, lw=1.0)
    bx(s, Inches(0.55), cy, Inches(0.12), Inches(1.35), RED)
    t(s, ti, Inches(0.85), cy+Inches(0.18), Inches(11.6), Inches(0.5), sz=16, bold=True, col=INK)
    t(s, body, Inches(0.87), cy+Inches(0.68), Inches(11.7), Inches(0.6), sz=12.5, col=GRY, line_sp=1.15)
    cy = cy + Inches(1.52)

# ════════ SLIDE 6 — 次アクション＋航空写真枠 ════════
s = sl(); ft(s)
hdr(s, "NEXT", "次の一手 ＆ 現地 航空写真", "名義と表札の橋渡しを確定し、私道を押さえる")
# 左：航空写真 貼付枠
bx(s, Inches(0.55), Inches(1.9), Inches(6.1), Inches(4.55), GRYBG, line=CARDLN, lw=1.0)
bx(s, Inches(0.55), Inches(1.9), Inches(6.1), Inches(0.06), RED)
t(s, "Googleマップ 航空写真（貼付枠）", Inches(0.55), Inches(3.7), Inches(6.1), Inches(0.4), sz=14, bold=True, col=RED, align=PP_ALIGN.CENTER)
t(s, "辻堂東海岸3-3-20 を航空写真でスクショ→ここに貼付\n（細い住宅街道路に面した密集地・前面道路の確認用）",
  Inches(0.75), Inches(4.15), Inches(5.7), Inches(0.8), sz=11, col=GRY, align=PP_ALIGN.CENTER, line_sp=1.2)
# 右：次アクション
steps = [
    ("1", "福井さんに確認", "「関口さんは449の名義4名（柏木ら）とどういう関係か＝家族か単なる居住者か」を一言。名義⇔表札の橋渡しが最短で埋まる。"),
    ("2", "荻原3,400万を本線に", "駐車交渉は4名全員同意で難度高。確実な出口＝荻原の業者買取を本線に、交渉は期限付きの上振れ狙いに位置づけ。"),
    ("3", "一覧をDBへ反映", "確定後、この権利関係マップを物件管理DB／案件シートへ反映し、交渉相手リストを確定。"),
]
dx, dy = Inches(6.95), Inches(1.95)
for i, (no, ti, body) in enumerate(steps):
    cy = dy + Inches(1.5)*i
    bx(s, dx, cy, Inches(5.83), Inches(1.35), CARD, line=CARDLN, lw=1.0)
    bx(s, dx, cy, Inches(0.7), Inches(1.35), RED)
    t(s, no, dx, cy+Inches(0.35), Inches(0.7), Inches(0.6), sz=26, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    t(s, ti, dx+Inches(0.9), cy+Inches(0.16), Inches(4.8), Inches(0.4), sz=15, bold=True, col=RED)
    t(s, body, dx+Inches(0.92), cy+Inches(0.6), Inches(4.8), Inches(0.7), sz=11, col=GRY, line_sp=1.12)

prs.save("tsujido_kenri.pptx")
print("saved tsujido_kenri.pptx  /  slides:", len(prs.slides._sldIdLst))

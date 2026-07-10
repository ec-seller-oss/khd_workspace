"""
KHD 運用マニュアル「人 × AI の作業分担 ＋ 毎日評価」 全7枚
クリーム白×レンガ赤(KHD標準)。SSoT: .company/secretary/notes/2026-06-28-claude-daikou-catalog-and-juushi.md
出力: unei_manual.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF9,0xF6,0xEF); RED=RGBColor(0xAA,0x2E,0x26); REDD=RGBColor(0x8C,0x24,0x1D)
INK=RGBColor(0x1A,0x1A,0x1A); GRY=RGBColor(0x6E,0x6E,0x6E); LINE=RGBColor(0xDA,0xD6,0xCF)
CARD=RGBColor(0xF1,0xEC,0xE1); CARDLN=RGBColor(0xE1,0xDA,0xCB); REDBG=RGBColor(0xF4,0xE4,0xE2)
GRYBG=RGBColor(0xEC,0xE8,0xDF); WHT=RGBColor(0xFF,0xFF,0xFF)
FONT="Hiragino Sans"
W=Inches(13.33); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def sl():
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG; return s

def t(slide,text,x,y,w,h,sz=18,bold=False,col=INK,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP,line_sp=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_sp: p.line_spacing=line_sp
        r=p.add_run(); r.text=line; r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=col; r.font.name=FONT
    return tb

def bx(slide,x,y,w,h,col,line=None,lw=1.0):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    if col is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=col
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=line; s.line.width=Pt(lw)
    s.shadow.inherit=False; return s

def hdr(slide,eyebrow,main,sub=""):
    t(slide,eyebrow,Inches(0.6),Inches(0.4),Inches(12),Inches(0.4),sz=13,bold=True,col=RED)
    bx(slide,Inches(0.62),Inches(0.78),Inches(1.7),Pt(3),RED)
    t(slide,main,Inches(0.6),Inches(0.9),Inches(12.1),Inches(0.55),sz=23,bold=True,col=INK)
    if sub: t(slide,sub,Inches(0.62),Inches(1.44),Inches(12.1),Inches(0.3),sz=11.5,col=GRY)

def ft(slide):
    bx(slide,Inches(0.5),H-Inches(0.5),Inches(12.33),Pt(1.2),LINE)
    t(slide,"KHD 運用マニュアル  ｜  人 × AI 作業分担 ＋ 毎日評価",Inches(0.5),H-Inches(0.42),Inches(10),Inches(0.32),sz=9,col=GRY)

def light_table(slide,rows,x,y,w,h,col_w,hi_col=None,sz=12,header_sz=12):
    n,m=len(rows),len(rows[0]); tb=slide.shapes.add_table(n,m,x,y,w,h).table
    tb.first_row=False; tb.horz_banding=False
    for ci,cw in enumerate(col_w): tb.columns[ci].width=cw
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            cell=tb.cell(ri,ci); cell.text=str(val); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.fill.solid(); is_hi=(hi_col is not None and ci==hi_col)
            if ri==0: cell.fill.fore_color.rgb=REDD if is_hi else RED
            else: cell.fill.fore_color.rgb=REDBG if is_hi else (CARD if ri%2==1 else BG)
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name=FONT; r.font.size=Pt(header_sz if ri==0 else sz)
                    r.font.bold=(ri==0) or is_hi or (ci==0)
                    r.font.color.rgb=WHT if ri==0 else (RED if is_hi else (INK if ci==0 else RGBColor(0x3A,0x3A,0x3A)))
    return tb

def redlabel(slide,text,x,y,w,h,sz=12):
    bx(slide,x,y,w,h,RED); t(slide,text,x,y,w,h,sz=sz,bold=True,col=WHT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def card(slide,x,y,w,h,title,body,tcol=RED):
    bx(slide,x,y,w,h,CARD,CARDLN,1.0); bx(slide,x,y,w,Inches(0.06),tcol)
    t(slide,title,x+Inches(0.18),y+Inches(0.16),w-Inches(0.36),Inches(0.4),sz=14,bold=True,col=tcol)
    t(slide,body,x+Inches(0.18),y+Inches(0.62),w-Inches(0.36),h-Inches(0.7),sz=11.5,col=INK,line_sp=1.15)

# ════════ SLIDE 1 — 表紙 ════════
s=sl()
bx(s,Inches(0.6),Inches(1.5),Pt(5),Inches(3.0),RED)
t(s,"KHD OPERATING MANUAL — HUMAN × AI",Inches(0.9),Inches(1.5),Inches(11),Inches(0.5),sz=14,bold=True,col=RED)
t(s,"人 と AI の作業分担\n＋ 毎日評価",Inches(0.88),Inches(2.05),Inches(11.5),Inches(2.0),sz=44,bold=True,col=INK,line_sp=1.05)
bx(s,Inches(0.92),Inches(4.35),Inches(7.2),Pt(1.4),LINE)
t(s,"北極星 ── 人を笑わせ \"ファン\" にする人徳 × 一流の経営者 ＆ かっこいいパパ",Inches(0.9),Inches(4.5),Inches(11.6),Inches(0.6),sz=15,bold=True,col=REDD)
t(s,"打席は人が立つ。内務はAIが巻き取る。毎日、数字で採点する。",Inches(0.92),Inches(5.15),Inches(11.6),Inches(0.4),sz=12.5,col=GRY)
t(s,"2026-06-28  ｜  菊池研太 / KHD",Inches(0.92),Inches(6.5),Inches(8),Inches(0.4),sz=11,col=GRY)

# ════════ SLIDE 2 — 北極星 ════════
s=sl(); hdr(s,"NORTH STAR / 目指す姿","面白くて人間力ある経営者は、仕事とリアルの両輪で作る","片方では作れない。リアルの厚みが、トークと人間味の源泉。")
card(s,Inches(0.6),Inches(1.95),Inches(5.95),Inches(4.6),"［ 仕事 ］売り込まず、信頼の対価で稼ぐ",
     "・打席に立つ量＝トーク／人間力の筋トレ\n・本音を聞く（感情ギャップ・相手目線）\n・GIVE先行・売り込まない＝信頼の対価\n・提案数 × パターン（淡々レスのAIにならない）\n・AIで情報ギャップを埋め、人は感情に全振り\n・専門性（調査士・診療圏）＝信頼の武器\n・数字で自己承認（他人評価で消耗しない）",RED)
card(s,Inches(6.78),Inches(1.95),Inches(5.95),Inches(4.6),"［ リアル ］人間の厚みを貯める",
     "・家族時間（妻・葵斗）＝幸福の土台・先回りGIVE\n・健康・睡眠・余白＝余裕＝人間味（焦らない／平和）\n・趣味・教養（サッカー／読書／投資）＝話題の引き出し\n・人脈・異業種交流＝縁・人間力\n・感謝の発信（妻・仲間）＝関係の土台\n・ブレない軸・自己変革＝経営者の芯",REDD)
ft(s)

# ════════ SLIDE 3 — 時間配分 ════════
s=sl(); hdr(s,"TIME ALLOCATION / 1日の時間配分","朝・日中・夜の3帯で、調査士と打席を先に置く","内務は固定枠を取らない＝AI代行＋リモート隙間で。")
light_table(s,[
 ["活動","朝","日中","夜","1日計"],
 ["調査士（聖域）","3","1","1","5h"],
 ["打席（営業・対面）","2","2","2","6h 狙い"],
 ["内務","—","隙間","—","AI代行＋隙間"],
],Inches(0.6),Inches(2.05),Inches(8.4),Inches(2.4),
 [Inches(2.9),Inches(1.1),Inches(1.3),Inches(1.1),Inches(2.0)],hi_col=4,sz=14,header_sz=13)
card(s,Inches(9.3),Inches(2.05),Inches(3.45),Inches(2.4),"運用思想",
     "・調査士＝朝3h聖域（電話/チャット禁止）\n・打席は機会が少なければ無理に埋めず、空きは内務へ\n・内務はクロード代行＋移動中の隙間で",RED)
t(s,"狙い：脳の主役を「打席」と「調査士」に空ける。内務で焦って手を動かさない。",
  Inches(0.6),Inches(4.8),Inches(12),Inches(0.5),sz=13,bold=True,col=REDD)
ft(s)

# ════════ SLIDE 4 — 人がやること ════════
s=sl(); hdr(s,"HUMAN / 人がやること","👤 菊池は「打席」だけ ── 人にしかできないこと","ここに時間とエネルギーを全振りする。")
items=[("人と会う","対面・電話・面談で打席に立つ"),
       ("本音を聞く","条件の奥の感情・恐怖・プライドを引き出す"),
       ("決断する","買付・撤退・価格・GO/No-Go"),
       ("人間関係を築く","親友・先生・業者・家族（GIVE先行）"),
       ("調査士を学ぶ","朝3h聖域。全事業の仕入・信頼・出口の土台"),
       ("最終チェック","送信・契約・実印は人の目で")]
x0=Inches(0.6); y0=Inches(2.0); cw=Inches(3.95); ch=Inches(2.1); gx=Inches(0.18); gy=Inches(0.2)
for i,(ti,bo) in enumerate(items):
    cx=x0+(cw+gx)*(i%3); cy=y0+(ch+gy)*(i//3)
    card(s,cx,cy,cw,ch,ti,bo,RED)
ft(s)

# ════════ SLIDE 5 — AIがやること（カタログ） ════════
s=sl(); hdr(s,"AI / クロード代行カタログ","🤖 内務はAIに丸投げ ── 隙間で「A1やって」と投げるだけ","リモートOK。番号で指示。")
light_table(s,[
 ["#","区分","代行タスク"],
 ["A1","お金","カード/銀行明細DL → MFインポートCSV"],
 ["A2","お金","EC粗利集計・販売管理表の分析"],
 ["A3","お金","資金繰り/PL/BS更新・古田土ダッシュ"],
 ["B1","物件","マイソク→収支試算（玉川式KPI判定）"],
 ["B2","物件","レインズ/ポータルで出口価格収集・市場調査"],
 ["B3","物件","物件マスターDB・在庫表 更新"],
 ["C1","医療","診療圏調査レポート作成"],
 ["C2","医療","補助金/制度リサーチ"],
 ["C3","医療","提案書・事業計画スライド"],
 ["D1","営業","追客/LINE/メール下書き（送信前チェック付）"],
 ["D2","営業","ザオラル（眠った関係の再点火文）"],
 ["D3","営業","顧客マスター/02DB更新・KPI記録"],
 ["E1","仕組","スプシ設計・GAS"],
 ["E2","仕組","Notion記録・整理"],
 ["E3","仕組","朝ブリーフ・日報・週次レビュー"],
],Inches(0.6),Inches(1.85),Inches(12.1),Inches(5.1),
 [Inches(0.9),Inches(1.5),Inches(9.7)],sz=11.5,header_sz=12)
ft(s)

# ════════ SLIDE 6 — 毎日評価運用 ════════
s=sl(); hdr(s,"DAILY REVIEW / 毎日評価","コード × 02_作業DB を紐付けて、毎日 数字で採点する","新しい列は作らない＝既存の「作業カテゴリ」を転用。")
card(s,Inches(0.6),Inches(1.95),Inches(5.95),Inches(2.6),"仕組み（列を増やさない）",
     "・人/AI判別＝作業カテゴリで：\n　営業打席・学習・家族・会議連絡・移動 → 人\n　仕込み資料・内務事務・その他 → AI\n・カタログ番号は内容の頭に [A1] タグ\n・「相談」列は本来の用途のまま温存",RED)
light_table(s,[
 ["毎日の4KPI","目標"],
 ["打席比率（営業直結の時間）","60% 以上"],
 ["調査士の時間","5h"],
 ["AI代行 件数","計測・最大化"],
 ["顧客KPI 記録漏れ","0 件"],
],Inches(6.78),Inches(1.95),Inches(5.95),Inches(2.5),
 [Inches(4.0),Inches(1.95)],hi_col=1,sz=12.5,header_sz=12.5)
t(s,"朝ブリーフで提示 → 夜の日報で採点。人がどれだけ打席に立てたか／AIがどれだけ内務を巻き取ったかが、毎日見える。",
  Inches(0.6),Inches(4.75),Inches(12.1),Inches(0.6),sz=13,bold=True,col=REDD,line_sp=1.1)
ft(s)

# ════════ SLIDE 7 — 決意 ════════
s=sl()
bx(s,Inches(0.6),Inches(1.8),Pt(5),Inches(3.4),RED)
t(s,"DECISION / 決意",Inches(0.9),Inches(1.8),Inches(11),Inches(0.5),sz=14,bold=True,col=RED)
t(s,"一流かつ、かっこいい大人を\n本気で目指す。",Inches(0.88),Inches(2.4),Inches(11.6),Inches(1.8),sz=40,bold=True,col=INK,line_sp=1.08)
t(s,"人を笑わせ、ファンにする人徳を作り続ける。\n経営者として、パパとして。",Inches(0.92),Inches(4.5),Inches(11.6),Inches(1.0),sz=17,bold=True,col=REDD,line_sp=1.2)
t(s,"打席は人が立つ。内務はAIが巻き取る。毎日、数字で前へ。",Inches(0.92),Inches(6.4),Inches(11),Inches(0.4),sz=12,col=GRY)

prs.save("unei_manual.pptx")
print("saved unei_manual.pptx /", len(prs.slides.__iter__.__self__._sldIdLst), "slides")

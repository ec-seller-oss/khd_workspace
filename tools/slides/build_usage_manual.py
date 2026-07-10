# -*- coding: utf-8 -*-
# KHD 経営ダッシュボード 使い方マニュアル（妻も使える・模式図）2026-06-03
# デザイン：クリーム白#F9F6EF × レンガ赤#AA2E26（スライド標準）
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM=RGBColor(0xF9,0xF6,0xEF); BRICK=RGBColor(0xAA,0x2E,0x26); INK=RGBColor(0x33,0x33,0x33)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0xCD,0xE9,0xD6); YEL=RGBColor(0xFF,0xF4,0xD6)
GRAY=RGBColor(0xEC,0xEC,0xEC); SUBF=RGBColor(0xF0,0xE2,0xDF)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]

def slide(bg=CREAM):
    s=prs.slides.add_slide(blank)
    r=s.shapes.add_shape(1,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def box(s,x,y,w,h,fill,line=None):
    b=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h)); b.fill.solid(); b.fill.fore_color.rgb=fill
    if line: b.line.color.rgb=line; b.line.width=Pt(1)
    else: b.line.fill.background()
    b.shadow.inherit=False; return b
def txt(s,x,y,w,h,text,size,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    lines=text.split("\n")
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=sp
        r=p.add_run(); r.text=ln; f=r.font; f.size=Pt(size); f.bold=bold; f.color.rgb=color; f.name="Meiryo"
    return tb
def titlebar(s,t,sub=None):
    box(s,0,0,13.333,1.15,BRICK)
    txt(s,0.5,0.12,12.3,0.95,t,26,WHITE,True,anchor=MSO_ANCHOR.MIDDLE)
    if sub: txt(s,0.5,0.82,12.3,0.35,sub,12,RGBColor(0xF3,0xD9,0xD6))

# S1 表紙
s=slide(BRICK)
txt(s,0,2.4,13.333,1.2,"KHD 経営ダッシュボード 使い方",40,WHITE,True,PP_ALIGN.CENTER)
txt(s,0,3.7,13.333,0.7,"妻と一緒に見る、1枚の「家計 ＋ 事業」",20,RGBColor(0xF3,0xD9,0xD6),False,PP_ALIGN.CENTER)
txt(s,0,6.5,13.333,0.5,"2026-06　／　毎日は①を開くだけ・押すボタンは週1回だけ",13,RGBColor(0xF3,0xD9,0xD6),False,PP_ALIGN.CENTER)

# S2 全体像
s=slide(); titlebar(s,"全体像 ── このスプシは何か","「見る①〜⑦」＋「〔元〕データ」。答えるのは3つの問い")
box(s,0.5,1.5,6.0,3.4,WHITE,SUBF); txt(s,0.7,1.6,5.6,0.5,"👀 見るタブ（①〜⑦）",16,BRICK,True)
txt(s,0.8,2.2,5.6,2.6,"① 司令塔（毎日ここ）\n② 時間×金×家族\n③ 資産・負債（BS）\n④ 資金繰り（現金）\n⑤ 借入\n⑥ 使い方\n⑦ 損益（PL）",15,INK,sp=1.15)
box(s,6.8,1.5,6.0,3.4,GRAY); txt(s,7.0,1.6,5.6,0.5,"🗄 〔元〕データ／🔐機密（触らない）",16,INK,True)
txt(s,7.1,2.2,5.6,2.6,"〔元〕売上見込み・諸経費・借入返済・税金\n〔元〕未来会計図表・経費削減損切り\n→ ①〜⑦が自動で参照する“源”\n\n🔐 パス・ID（機密）",14,INK,sp=1.15)
box(s,0.5,5.2,12.3,1.7,GREEN); txt(s,0.7,5.3,12.0,1.5,"3つの問い：\n① 儲かるか？ → ⑦損益（PL・損益分岐）\n② お金が回るか？ → ④資金繰り（現金・ランウェイ）\n③ 家族の時間は？ → ②時間×金×家族（目的）",15,INK,True,sp=1.1)

# S3 毎日
s=slide(); titlebar(s,"【毎日】① 司令塔を開くだけ","数字は自動。押すボタンも入力も無し。眺めるだけ")
cards=[("純資産","2,809万","過去＝今ある資産"),("現預金","752万","すぐ使えるお金"),
       ("ランウェイ","11.9ヶ月 🟢","あと何ヶ月もつか"),("営業直結%","30%（目標60%）","時間が営業に向いてるか")]
for i,(t,v,d) in enumerate(cards):
    x=0.5+i*3.1; box(s,x,1.6,2.9,2.2,WHITE,SUBF)
    txt(s,x+0.15,1.75,2.6,0.5,t,14,BRICK,True); txt(s,x+0.15,2.3,2.6,0.7,v,22,INK,True)
    txt(s,x+0.15,3.15,2.6,0.6,d,11,INK)
box(s,0.5,4.2,12.3,2.5,YEL)
txt(s,0.7,4.35,12.0,2.2,"見方：左から「過去（資産）→ 現在（現金・ランウェイ）→ 未来（谷）→ 行動（営業%）」の4軸が1画面。\n\n🟢攻めOK＝6ヶ月以上 ／ 🟡注意＝3〜6 ／ 🔴守り＝3ヶ月未満。\n営業直結%が60%を割ってたら「内務を減らして営業へ」のサイン。\n\n★毎日はこれを開いて眺めるだけ。入力も計算も自動。",15,INK,sp=1.15)

# S4 週次
s=slide(); titlebar(s,"【週次・月曜】② 時間×金×家族","メニュー『📊統合ダッシュボード→実績hをカレンダーから更新』を1回押す")
box(s,0.5,1.5,12.3,0.5,BRICK);
for i,h in enumerate(["事業","実績h","期待粗利","円/h","家族◯ヶ月分","判断"]):
    txt(s,0.6+i*2.05,1.52,2.0,0.45,h,12,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
rows=[("不動産(栄町)","23h","171万","42,750","2.7ヶ月","続ける"),
      ("医療コンサル","10h","33万","5,500","0.5ヶ月","増やす"),
      ("EC","2h","13万","3,375","0.2ヶ月","維持"),
      ("内務","35h","—","—","—","減らす"),
      ("家族(目的)","29h","—","—","—","死守")]
for r,row in enumerate(rows):
    y=2.0+r*0.62; bg=GREEN if "家族" in row[0] else WHITE
    box(s,0.5,y,12.3,0.6,bg,SUBF)
    for i,c in enumerate(row): txt(s,0.6+i*2.05,y+0.05,2.0,0.5,c,12,INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.5,5.3,12.3,1.6,"見方：1時間がいくら生むか（円/h）が高い事業に時間を寄せる。\n営業直結%・本命メディア時間が薄くないか毎週点検 → 翌週の配分を決める。\n★押すのはこの『実績h更新』だけ（週1）。",14,INK,True,sp=1.15)

# S5 月初
s=slide(); titlebar(s,"【月初・5分】残高だけ手で更新 → あとは自動","BSの残高を写すと、①司令塔・④資金繰り・⑦損益が全部追従")
steps=[("1","MoneyForwardを開く","各口座の残高を確認"),
       ("2","③BSの残高列を更新","銀行・証券の数字を写すだけ（5分）"),
       ("3","①④⑦が自動で最新に","純資産・ランウェイ・谷が更新"),
       ("4","解約・損切りレビュー","⑥の月次チェック")]
for i,(n,t,d) in enumerate(steps):
    y=1.6+i*1.25; box(s,0.6,y,0.9,0.9,BRICK); txt(s,0.6,y+0.1,0.9,0.7,n,30,WHITE,True,PP_ALIGN.CENTER)
    box(s,1.7,y,11.0,0.9,WHITE,SUBF); txt(s,1.9,y+0.08,10.6,0.45,t,16,BRICK,True); txt(s,1.9,y+0.5,10.6,0.35,d,12,INK)
txt(s,0.6,6.7,12.0,0.6,"※信金2口座(大東京/法人TB)は2ヶ月毎に窓口/電話で取得 → その時だけ手更新。",12,INK)

# S6 妻と見る核心
s=slide(GREEN); titlebar(s,"💑 妻と見る核心 ── お金を「家族の時間」に翻訳","②本部マトリクスの下のサマリー")
box(s,1.0,1.7,11.3,2.0,WHITE,BRICK)
txt(s,1.2,1.9,11.0,1.7,"今月の期待粗利 合計  =  2,205,000円\n＝ 家族 3.5ヶ月分の暮らしを確保",24,BRICK,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE,1.2)
txt(s,1.0,4.0,11.3,2.8,"だから「今は家族の時間を少し割いてでも、この事業に集中する」が腹落ちする。\n\n・栄町(不動産)＝1時間で¥42,750＝家族2.7ヶ月分 → 6月はここに集中\n・医療コンサル＝継続性◎ → 並行で仕込む\n\n稼ぐ目的は“家族の時間を増やすこと”。家族時間が削れてないか毎月確認。",16,INK,sp=1.2)

# S7 解約ルール
s=slide(); titlebar(s,"💸 解約・損切りルール（カード年会費の垂れ流し防止）","契約する時に“やめる期限”を決める")
txt(s,0.6,1.5,12.1,2.2,"① 新しい固定費（カード年会費・サブスク）を契約する時、\n   「いくらの売上・効果で元が取れるか＝損益分岐」を出す。\n\n② その達成期限＝“解約判断期限”を 🔐パス・ID の出口/期限欄に書く。\n\n③ 期限が来て損益分岐を超えてなければ解約。\n   （カードは年会費が落ちる前にレビュー＝手遅れ防止）",17,INK,sp=1.25)
box(s,0.6,4.6,12.1,2.0,YEL); txt(s,0.8,4.75,11.7,1.8,"★ポイント・マイルは有効期限が近いとアラート（期限アラート列）。\nB1ポイントを1/1で失効させた失敗を二度と繰り返さない。\n解約期限と有効期限は“同じ仕組み”で先に手を打つ。",15,INK,True,sp=1.15)

# S8 データ置き場
s=slide(); titlebar(s,"🗄 データ置き場（〔元〕）と 🔐機密 ── 触らない","ここは①〜⑦が自動で参照する“源”。普段は開かない")
box(s,0.5,1.5,12.3,2.6,GRAY); txt(s,0.7,1.6,12.0,0.5,"〔元〕データ（自動で参照される源）",16,INK,True)
txt(s,0.8,2.2,12.0,1.8,"〔元〕売上見込み … ④資金繰りの入金へ（シナリオ切替で自動）\n〔元〕借入返済・借入条件 … ⑤借入・④へ（証票ベース）\n〔元〕諸経費・税金・未来会計図表 … 損益・予測の源\n〔元〕経費削減・損切り … 妻の打ち手",14,INK,sp=1.2)
box(s,0.5,4.4,12.3,1.3,SUBF); txt(s,0.7,4.55,12.0,1.1,"🔐 パス・ID … ログイン情報（機密）。開始日=A列、出口・期限=同じ列。\n各資産の残高確認リンクは③BSのG列から飛べる。",14,INK,sp=1.15)
txt(s,0.5,6.0,12.3,0.8,"★〔元〕は基本いじらない。月初に銀行帳票が来たら該当口座フォルダに保存（ルート直置き禁止）。",13,BRICK,True)

# S9 困った時
s=slide(); titlebar(s,"🆘 困った時 ── 帳票の取り方・タスクの見方","")
txt(s,0.6,1.5,12.1,2.6,"【銀行帳票（返済予定表）の取得】\n・朝日 … 朝日ビジネスポータル（電子交付）でDL ※PWは1111\n・公庫 … 公庫ダイレクト（ネット） 0120-154-505\n・大東京/東京ベイ（信金） … 2ヶ月毎に窓口・融資アポと同時\n→ 取ったら⑤借入を証票どおりに直す（私に渡せばやります）",16,INK,sp=1.25)
box(s,0.6,4.5,12.1,2.1,WHITE,SUBF); txt(s,0.8,4.65,11.7,1.9,"【やること（タスク）】Google Tasks「マイタスク」に全部・期限つき。\n・営業直結（福井送信/物件確認）＝今週\n・資産系（PL再設計/ポイント自動化）＝来週\n・信金/公庫の帳票取得＝7月\n毎朝：カレンダー＋Tasksで今日やることを確認。",15,INK,sp=1.2)

# S10 まとめ
s=slide(BRICK)
txt(s,0,1.0,13.333,0.9,"まとめ ── これだけ覚えればOK",30,WHITE,True,PP_ALIGN.CENTER)
items=[("毎日","① 司令塔を開いて眺めるだけ（自動・入力なし）"),
       ("週1回","メニュー『実績hをカレンダーから更新』を押す"),
       ("月初5分","③BSの残高を手で更新 → ①④⑦が自動で最新"),
       ("妻と","②で『家族◯ヶ月分』を見て、来月の時間配分を決める"),
       ("ルール","契約時に解約期限を決める／タスクは全部期限つき")]
for i,(k,v) in enumerate(items):
    y=2.1+i*0.95; box(s,1.2,y,2.3,0.78,YEL); txt(s,1.2,y+0.12,2.3,0.55,k,17,BRICK,True,PP_ALIGN.CENTER)
    box(s,3.7,y,8.4,0.78,CREAM); txt(s,3.9,y+0.1,8.1,0.6,v,15,INK,anchor=MSO_ANCHOR.MIDDLE)

out="/Users/kikuchikenta/Library/CloudStorage/GoogleDrive-ec-seller@kikuchi-hd.net/マイドライブ/KHD_経営ダッシュボード_使い方.pptx"
prs.save(out); print("保存:",out); print("slides:",len(prs.slides._sldIdLst))
